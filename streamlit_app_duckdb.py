from __future__ import annotations

import io
import json
import re
from pathlib import Path

import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st


# =========================
# Optional OpenAI Support
# =========================

OPENAI_AVAILABLE = False
client = None

try:
    from openai import OpenAI

    if "OPENAI_API_KEY" in st.secrets:
        client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])
        OPENAI_AVAILABLE = True

except Exception:
    OPENAI_AVAILABLE = False
    client = None


def retrieve_market_context_with_openai(
    context_package: dict,
    market_context_query: str,
    model: str = "gpt-4.1-mini",
) -> str:
    """Controlled market / sector context retrieval using OpenAI web search.

    This is separate from final commentary generation so the app can show users
    exactly what market context is being used before synthesis.
    """

    if not OPENAI_AVAILABLE or client is None:
        return (
            "Market context retrieval unavailable. Confirm that `openai` is in requirements.txt "
            "and `OPENAI_API_KEY` is configured in Streamlit Secrets."
        )

    retrieval_prompt = {
        "task": "Retrieve and summarize public market context for muni commentary.",
        "strict_rules": [
            "Focus on public market context only.",
            "Do not invent issuer-specific explanations.",
            "Separate rates/Treasury context, muni market context, sector context, and issuer-specific public headlines if any.",
            "If relevant public context is not found, say so clearly.",
            "Keep the output concise and evidence-oriented.",
        ],
        "dashboard_context": {
            "issuer": context_package.get("issuer"),
            "sector": context_package.get("sector"),
            "bucket": context_package.get("bucket"),
            "benchmark": context_package.get("benchmark"),
            "period": context_package.get("period"),
            "signals": context_package.get("signals", {}),
        },
        "search_focus": market_context_query,
        "preferred_context_categories": [
            "Treasury curve / rates movement",
            "municipal bond market tone",
            "municipal fund flows",
            "sector-specific public news",
            "issuer-specific public news only if clearly available",
        ],
        "output_format": {
            "Rates / Treasury Context": "2-4 bullets",
            "Muni Market Context": "2-4 bullets",
            "Sector / Issuer Context": "2-4 bullets",
            "Relevance to Dashboard Signals": "2-4 bullets",
            "Caveats": "1-2 bullets",
        },
    }

    try:
        response = client.responses.create(
            model=model,
            tools=[{"type": "web_search_preview"}],
            input=[
                {
                    "role": "system",
                    "content": (
                        "You are a market context retrieval assistant for a municipal bond analytics dashboard. "
                        "Retrieve public context conservatively. Do not overstate causality."
                    ),
                },
                {
                    "role": "user",
                    "content": json.dumps(retrieval_prompt, indent=2, default=str),
                },
            ],
            temperature=0.15,
            max_output_tokens=900,
        )

        return response.output_text

    except Exception as e:
        return f"Market Context Retrieval Error: {str(e)}"


def generate_ai_market_commentary(
    context_package: dict,
    manual_market_context: str = "",
    retrieved_market_context: str = "",
    use_web_search: bool = False,
    market_context_query: str = "",
    model: str = "gpt-4.1-mini",
) -> str:
    """Generate evidence-linked institutional commentary.

    The model should only synthesize the analytics and market context provided.
    Web search is available as an optional fallback, but preferred workflow is:
    dashboard signals -> controlled retrieval -> review context -> commentary.
    """

    if not OPENAI_AVAILABLE or client is None:
        return (
            "AI commentary unavailable. Confirm that `openai` is in requirements.txt "
            "and `OPENAI_API_KEY` is configured in Streamlit Secrets."
        )

    system_prompt = """
You are an institutional municipal bond market strategist.

Write concise, evidence-linked secondary-market commentary for a muni trading / public finance team.

Rules:
- Use ONLY the provided dashboard analytics and provided/retrieved market context.
- Do NOT invent issuer-specific news, ratings actions, trades, or market events.
- Do NOT claim causality unless the provided market context supports it.
- If market context is missing or weak, say that context is limited.
- Separate data-backed observations from interpretation.
- Keep tone professional, like buy-side or broker-dealer strategy commentary.
- Mention that signals are screening indicators, not investment recommendations.
- Prefer 4 sections:
  1) Market Commentary
  2) Why This May Be Happening
  3) Risks / Caveats
  4) Evidence Used
"""

    user_payload = {
        "task": "Generate institutional municipal secondary-market commentary.",
        "manual_market_context": manual_market_context,
        "retrieved_market_context": retrieved_market_context,
        "market_context_query": market_context_query,
        "dashboard_context_package": context_package,
        "requested_output_format": {
            "Market Commentary": "2-4 concise bullet points",
            "Why This May Be Happening": "2-4 concise bullet points tied to evidence",
            "Risks / Caveats": "1-3 bullets noting weak evidence, liquidity/data limitations, or alternative explanations",
            "Evidence Used": "bullet list of exact dashboard signals and market-context items used",
        },
    }

    try:
        tools = [{"type": "web_search_preview"}] if use_web_search else []

        response = client.responses.create(
            model=model,
            tools=tools,
            input=[
                {"role": "system", "content": system_prompt},
                {
                    "role": "user",
                    "content": json.dumps(user_payload, indent=2, default=str),
                },
            ],
            temperature=0.25,
            max_output_tokens=1100,
        )

        return response.output_text

    except Exception as e:
        return f"AI Commentary Error: {str(e)}"



from data_utils import (
    build_issuer_master,
    merge_market_data,
    read_uploaded_file,
    standardize_bonds,
    standardize_issuer_mapping,
    standardize_mmd,
    standardize_trades,
)


st.set_page_config(page_title="Municipal Secondary Market Dashboard Generator", layout="wide")
st.title("Municipal Secondary Market Dashboard Generator")
st.caption("Bring your own bond master and trade-history exports. Generate issuer-level relative value and liquidity analytics.")

st.markdown(
    """
<style>
/* Overall page polish */
.block-container {
    padding-top: 2.2rem;
    padding-bottom: 3rem;
    max-width: 1500px;
}

h1, h2, h3 {
    letter-spacing: -0.02em;
}

section[data-testid="stSidebar"] {
    min-width: 330px !important;
}

div[data-testid="stMetric"] {
    background: #ffffff;
    border: 1px solid #e6e8ef;
    border-radius: 16px;
    padding: 18px 18px 14px 18px;
    box-shadow: 0 1px 3px rgba(15, 23, 42, 0.05);
}

.clean-card {
    background: #ffffff;
    border: 1px solid #e6e8ef;
    border-radius: 18px;
    padding: 18px 20px;
    min-height: 124px;
    box-shadow: 0 1px 3px rgba(15, 23, 42, 0.05);
}

.clean-card-label {
    font-size: 0.86rem;
    font-weight: 700;
    color: #64748b;
    margin-bottom: 8px;
    text-transform: uppercase;
    letter-spacing: 0.03em;
}

.clean-card-value-large {
    font-size: 1.65rem;
    font-weight: 720;
    line-height: 1.15;
    color: #111827;
    overflow-wrap: anywhere;
}

.clean-card-value-small {
    font-size: 1.32rem;
    font-weight: 720;
    line-height: 1.2;
    color: #111827;
    overflow-wrap: anywhere;
}

.clean-card-note {
    font-size: 0.82rem;
    color: #94a3b8;
    margin-top: 8px;
}

.nav-card {
    background: #f8fafc;
    border: 1px solid #e2e8f0;
    border-radius: 16px;
    padding: 14px 16px;
    margin: 10px 0 18px 0;
}

.nav-card a {
    text-decoration: none;
    color: #334155;
    font-size: 0.92rem;
}

.nav-card a:hover {
    color: #0f172a;
    text-decoration: underline;
}

.sidebar-nav-small {
    font-size: 0.88rem;
    line-height: 1.55;
}

/* Keep dataframes/charts visually lighter */
div[data-testid="stDataFrame"] {
    border-radius: 14px;
    overflow: hidden;
}
</style>
""",
    unsafe_allow_html=True,
)


def section_anchor(anchor_id: str, title: str, level: int = 2):
    """Create a stable HTML anchor plus a Streamlit header/subheader."""
    st.markdown(f"<a id='{anchor_id}'></a>", unsafe_allow_html=True)
    if level == 1:
        st.title(title)
    elif level == 2:
        st.header(title)
    else:
        st.subheader(title)


def clean_metric_card(label: str, value: object, size: str = "large", note: str | None = None):
    """Compact custom metric card that gives long text more room than st.metric."""
    value_class = "clean-card-value-large" if size == "large" else "clean-card-value-small"
    safe_value = "—" if value is None else str(value)
    note_html = f"<div class='clean-card-note'>{note}</div>" if note else ""
    st.markdown(
        f"""
<div class="clean-card">
  <div class="clean-card-label">{label}</div>
  <div class="{value_class}">{safe_value}</div>
  {note_html}
</div>
""",
        unsafe_allow_html=True,
    )


def section_directory():
    """Compact workflow map for the main page.

    The full jump list lives in the sidebar. This main-page version is intentionally
    concise so it does not crowd the dashboard before users upload data.
    """
    with st.expander("Dashboard workflow map", expanded=False):
        st.markdown(
            """
<div class="nav-card">
<b>How to read this dashboard</b><br><br>

<b>1. Data readiness</b><br>
<a href="#file-readiness">File Readiness</a> ·
<a href="#data-quality-scorecard">Data Quality Scorecard</a> ·
<a href="#executive-snapshot">Executive Snapshot</a><br><br>

<b>2. Benchmark & spread framework</b><br>
<a href="#yield-relative-value">Yield / RV Trend</a> ·
<a href="#issuer-curve">Issuer Curve</a> ·
<a href="#spread-level">Spread Level</a> ·
<a href="#spread-attribution">Spread Attribution</a><br><br>

<b>3. Relative value signals</b><br>
<a href="#peer-rv">Peer RV</a> ·
<a href="#cross-issuer-rv">Cross-Issuer RV</a> ·
<a href="#historical-spread">Historical Percentile</a> ·
<a href="#recommendation-engine">Rule-Based Narrative</a> ·
<a href="#ai-commentary-studio">AI Commentary Studio</a><br><br>

<b>4. Risk, flow & opportunity screening</b><br>
<a href="#curve-shape">Curve Shape</a> ·
<a href="#scenario-shock">Scenario Shock</a> ·
<a href="#dealer-proxy">Dealer Proxy</a> ·
<a href="#security-screener">Security Screener</a> ·
<a href="#watchlist">Watchlist</a><br><br>

<b>5. Outputs, methodology & raw detail</b><br>
<a href="#spread-movement">Spread Movement</a> ·
<a href="#cusip-drilldown">CUSIP Drilldown</a> ·
<a href="#report-export-center">Report Export Center</a> ·
<a href="#export-summary">Export Summary</a> ·
<a href="#admin-methodology">Admin Methodology</a> ·
<a href="#version-changelog">Version / Change Log</a> ·
<a href="#downloads">Downloads</a>
</div>
""",
            unsafe_allow_html=True,
        )


with st.expander("Instructions", expanded=False):
    st.markdown(
        """
<div style='font-size:15px; color:black; line-height:1.4;'>

<h5 style='margin-bottom:4px;'>Step 1: Upload Required Files</h5>

<div style='padding-left:18px;'>

<b>1. Bond File</b>

<ul style='margin-top:2px; margin-bottom:6px;'>
<li>Information can be found from Munipro</li>
<li>Row 1 must contain column headers</li>
<li>Actual data should begin from Row 2</li>
<li>Multiple issuers’ bond data should be combined into the same file</li>
</ul>

<b>Minimum Required Columns:</b><br>Issuer, Cusip, Maturity<br><br><b>Recommended Columns:</b><br>Type, Lien, Election, Series, Secondary Credit, Term, Par Amount, Outstanding Amount, Coupon, Call Date, Call Price, Fed Tax, AMT

<div style='height:10px;'></div>

<b>2. Trade History File(s)</b>

<ul style='margin-top:2px; margin-bottom:6px;'>
<li>Information can be extracted from Munipro</li>
<li>Row 1 must contain column headers</li>
<li>Actual data should begin from Row 2</li>
<li>Trade files should be uploaded separately</li>
</ul>

<b>Minimum Required Columns:</b><br>CUSIP9, Trade Date, Yield<br><br><b>Recommended Columns:</b><br>Trade Date/Time, Description, Maturity Date, Settlement Date, Coupon, Price, Trade Amount, Calculation Date, Calculation Price, Index, Index Rate, Spread, Trade Type, Ratings M/S/F

<div style='height:8px;'></div>

<b>Important:</b><br>
CUSIP9 in Trade Files must match Cusip in Bond File.

<div style='height:10px;'></div>

<b>3. Optional Files</b>

<ul style='margin-top:2px; margin-bottom:2px;'>
<li>Issuer / Sector Mapping File</li>
<li>MMD Curve File</li>
</ul>

</div>

<h5 style='margin-top:10px; margin-bottom:4px;'>Step 2: Automatic Issuer Detection</h5>

<div style='padding-left:18px;'>
The dashboard automatically detects issuer names from uploaded datasets.
</div>

<h5 style='margin-top:10px; margin-bottom:4px;'>Step 3: Select Uploaded Issuer</h5>

<div style='padding-left:18px;'>

<ul style='margin-top:2px; margin-bottom:6px;'>
<li>Select one of the detected issuers</li>
<li>Apply optional filters:
    <ul style='margin-top:2px; margin-bottom:2px;'>
        <li>Maturity Bucket</li>
        <li>Time Window</li>
        <li>Relative Value Comparison</li>
    </ul>
</li>
</ul>

</div>

</div>
""",
        unsafe_allow_html=True

    )

# -----------------------------------------------------------------------------
# Team-readiness validation layer
# -----------------------------------------------------------------------------
# Goal: keep the dashboard usable even when files come from different people or
# slightly different Munipro exports. We separate fields into:
#   1) REQUIRED: the app needs these to run.
#   2) RECOMMENDED: the app can run without them, but analytics become weaker.
#   3) OPTIONAL: nice-to-have reference fields.

COLUMN_ALIASES: dict[str, list[str]] = {
    # Shared identifiers
    "cusip": ["cusip", "cusip9", "cusip 9", "cusip_9", "security id", "security_id"],
    "issuer": ["issuer", "issuer name", "issuer_name", "obligor", "borrower"],
    "sector": ["sector", "industry", "sector name", "sector_name"],
    "primary_type": ["type", "primary type", "primary_type", "bond type"],
    # Bond master fields
    "lien": ["lien"],
    "election": ["election"],
    "series": ["series"],
    "secondary_credit": ["secondary credit", "secondary_credit", "credit", "credit enhancement"],
    "term": ["term"],
    "maturity": ["maturity", "maturity date", "maturity_date", "mty"],
    "par_amount": ["par amount", "par_amount", "par", "amount issued"],
    "outstanding_amount": ["outstanding amount", "outstanding_amount", "amount outstanding", "current amount outstanding"],
    "coupon": ["coupon", "coupon rate", "coupon_rate", "cpn"],
    "call_date": ["call date", "call_date", "first call date", "first_call_date"],
    "call_price": ["call price", "call_price"],
    "fed_tax": ["fed tax", "fed_tax", "tax status", "tax_status"],
    "amt": ["amt", "alternative minimum tax"],
    "rating": ["rating", "ratings", "ratings m/s/f", "ratings_m_s_f", "moody/s&p/fitch"],
    # Trade fields
    "trade_datetime": ["trade date/time", "td & time", "td time", "trade datetime", "trade_datetime", "datetime"],
    "trade_date": ["trade date", "td & time", "td time", "trade_date", "date", "transaction date"],
    "settlement_date": ["settlement date", "settlement_date", "settle date"],
    "description": ["description", "security description", "bond description"],
    "maturity_trade": ["maturity date", "maturity_date", "maturity", "mty"],
    "yield": ["yield", "yield to worst", "ytw", "yield_to_worst", "yield to maturity", "ytm"],
    "price": ["price", "trade price", "execution price"],
    "trade_amount": ["trade amount", "trade_amount", "par traded", "par amount", "amount", "quantity", "qty", "qty (m)"],
    "calculation_date": ["calculation date", "calculation_date"],
    "calculation_price": ["calculation price", "calculation_price"],
    "index": ["index", "benchmark"],
    "index_rate": ["index rate", "index_rate", "benchmark rate"],
    "spread": ["spread", "g spread", "z spread", "spread to benchmark"],
    "trade_type": ["trade type", "tde type", "trade_type", "side", "buy/sell"],
}

BOND_REQUIRED = ["cusip", "issuer", "maturity"]
BOND_RECOMMENDED = [
    "coupon", "outstanding_amount", "call_date", "call_price", "sector", "secondary_credit", "fed_tax", "amt"
]
BOND_OPTIONAL = ["primary_type", "lien", "election", "series", "term", "par_amount", "rating"]

TRADE_REQUIRED = ["cusip", "trade_date", "yield"]
TRADE_RECOMMENDED = ["price", "trade_amount", "trade_type", "spread", "settlement_date", "rating"]
TRADE_OPTIONAL = ["trade_datetime", "description", "maturity_trade", "calculation_date", "calculation_price", "index", "index_rate"]

MMD_REQUIRED = ["date"]
MMD_RECOMMENDED = ["1Y", "2Y", "5Y", "10Y", "20Y", "30Y"]
CURVE_TEMPLATE_COLUMNS = [
    "date", "5Y", "10Y", "20Y", "30Y",
    "AA+_5Y", "AA+_10Y", "AA+_20Y", "AA+_30Y",
    "AA_5Y", "AA_10Y", "AA_20Y", "AA_30Y",
    "AA-_5Y", "AA-_10Y", "AA-_20Y", "AA-_30Y",
    "A_5Y", "A_10Y", "A_20Y", "A_30Y",
    "BBB_5Y", "BBB_10Y", "BBB_20Y", "BBB_30Y",
]


# -----------------------------------------------------------------------------
# Benchmark curve assumptions
# -----------------------------------------------------------------------------
# MMD is treated as the AAA municipal benchmark curve. Non-AAA curves are
# approximated by adding transparent, maturity-adjusted credit spread assumptions
# to the selected MMD tenor. Units are percentage points, not basis points:
#   0.10 = 10 bps.
# These assumptions are intentionally visible in the app so the team can review,
# override, or replace them with paid/internal curve data later.

RATING_SPREADS: dict[str, dict[str, float]] = {
    "AAA": {"5Y": 0.00, "10Y": 0.00, "20Y": 0.00, "30Y": 0.00},
    "AA+": {"5Y": 0.08, "10Y": 0.10, "20Y": 0.12, "30Y": 0.15},
    "AA": {"5Y": 0.10, "10Y": 0.14, "20Y": 0.17, "30Y": 0.20},
    "AA-": {"5Y": 0.14, "10Y": 0.18, "20Y": 0.22, "30Y": 0.28},
    "A+": {"5Y": 0.22, "10Y": 0.28, "20Y": 0.35, "30Y": 0.42},
    "A": {"5Y": 0.30, "10Y": 0.38, "20Y": 0.48, "30Y": 0.58},
    "A-": {"5Y": 0.42, "10Y": 0.55, "20Y": 0.68, "30Y": 0.82},
    "BBB": {"5Y": 0.60, "10Y": 0.80, "20Y": 1.00, "30Y": 1.20},
}

# ----------------------------------------------------------------------------‑
# Professional maturity bucket labels
# ----------------------------------------------------------------------------‑
# The original prototype used tenor-style labels (10Y / 20Y / 30Y) for broad
# maturity ranges. For presentation clarity, the app now uses descriptive curve
# sector labels instead.
#
# Short          = <= 7Y
# Intermediate  = 7Y to 15Y
# Long          = 15Y to 25Y
# Extended Long = 25Y+
#
# The benchmark tenor mapping still uses the closest MMD tenor for each bucket.

MATURITY_BUCKET_ORDER = ["Short", "Intermediate", "Long", "Extended Long"]
MATURITY_BUCKET_OPTIONS = ["All"] + MATURITY_BUCKET_ORDER

# Backward compatibility: data_utils or older uploaded processed data may still
# create legacy labels. Normalize them immediately after processing.
MATURITY_BUCKET_RENAME = {
    "10Y": "Intermediate",
    "20Y": "Long",
    "30Y": "Extended Long",
}

MMD_BUCKET_MAP = {
    "Short": "5Y",
    "Intermediate": "10Y",
    "Long": "20Y",
    "Extended Long": "30Y",
    "All": "10Y",
}
BENCHMARK_RATINGS = list(RATING_SPREADS.keys())


def _curve_column_key(name: object) -> str:
    """Normalize curve column names for flexible matching.

    Examples that should match the same idea:
    - AA 10Y, AA_10Y, AA-10Y
    - AAA 5Y, MMD 5Y, 5Y
    - AA+ 20Y, AA Plus 20Y
    """
    text = str(name).strip().lower()
    text = text.replace("+", " plus ").replace("-", " minus ")
    keep = []
    for ch in text:
        keep.append(ch if ch.isalnum() else " ")
    return " ".join("".join(keep).split()).replace(" ", "")


def _rating_key(rating: str) -> str:
    return _curve_column_key(rating)


def find_uploaded_benchmark_column(mmd_df: pd.DataFrame, tenor: str, rating: str) -> str | None:
    """Find an explicitly uploaded benchmark column for rating + tenor.

    Priority logic:
    1. Exact user-provided curve columns, e.g. AA_10Y / AA 10Y / AA Curve 10Y.
    2. For AAA only, also allow MMD/vanilla tenor columns, e.g. 10Y.

    This lets users upload vendor/internal AA/A/BBB curves. If they do not,
    the app falls back to MMD AAA + transparent spread assumptions.
    """
    normalized = {_curve_column_key(c): c for c in mmd_df.columns}
    r = _rating_key(rating)
    t = _curve_column_key(tenor)

    candidates = [
        f"{r}{t}",
        f"{r}curve{t}",
        f"{r}yield{t}",
        f"{r}muni{t}",
        f"{r}mmd{t}",
        f"{t}{r}",
    ]

    if rating == "AAA":
        candidates.extend([
            t,
            f"mmd{t}",
            f"mmdaaa{t}",
            f"aaammd{t}",
            f"aaacurve{t}",
        ])

    for key in candidates:
        if key in normalized:
            return normalized[key]
    return None


def get_benchmark_curve(mmd_plot: pd.DataFrame, tenor: str, rating: str) -> tuple[pd.Series, dict] | tuple[None, dict]:
    """Return benchmark yield and metadata.

    Priority:
    - Use explicitly uploaded rating curve column when available.
    - Otherwise use uploaded AAA/MMD tenor column + visible rating-spread assumption.
    """
    explicit_col = find_uploaded_benchmark_column(mmd_plot, tenor, rating)
    if explicit_col is not None:
        return pd.to_numeric(mmd_plot[explicit_col], errors="coerce"), {
            "benchmark_source": "Uploaded curve",
            "source_column": explicit_col,
            "rating_spread_bps": 0.0,
        }

    base_col = find_uploaded_benchmark_column(mmd_plot, tenor, "AAA")
    if base_col is None:
        return None, {
            "benchmark_source": "Unavailable",
            "source_column": None,
            "rating_spread_bps": pd.NA,
        }

    base_curve = pd.to_numeric(mmd_plot[base_col], errors="coerce")
    spread_adjustment = RATING_SPREADS.get(rating, RATING_SPREADS["AAA"]).get(tenor, 0.00)
    return base_curve + spread_adjustment, {
        "benchmark_source": "Modeled from MMD + spread assumption" if rating != "AAA" else "Uploaded MMD / AAA curve",
        "source_column": base_col,
        "rating_spread_bps": spread_adjustment * 100,
    }


def benchmark_curve_from_mmd(mmd_plot: pd.DataFrame, mmd_col: str, rating: str) -> pd.Series:
    """Backward-compatible wrapper used by older chart blocks."""
    curve, _meta = get_benchmark_curve(mmd_plot, mmd_col, rating)
    if curve is None:
        return pd.Series([pd.NA] * len(mmd_plot), index=mmd_plot.index, dtype="float")
    return curve


def rating_spread_table() -> pd.DataFrame:
    """User-facing spread assumption table in both percentage points and bps."""
    rows = []
    for rating, tenors in RATING_SPREADS.items():
        row = {"Rating": rating}
        for tenor, spread_pct in tenors.items():
            row[f"{tenor} Spread"] = spread_pct
            row[f"{tenor} Spread (bps)"] = round(spread_pct * 100, 1)
        rows.append(row)
    return pd.DataFrame(rows)


def _detect_mmd_date_column(mmd_df: pd.DataFrame) -> str | None:
    """Find the MMD date column across common naming variants."""
    if "Date" in mmd_df.columns:
        return "Date"
    if "date" in mmd_df.columns:
        return "date"
    return None


def make_benchmark_long(mmd_df: pd.DataFrame, rating: str) -> pd.DataFrame:
    """Convert MMD wide curve data into long benchmark data by maturity bucket.

    Output columns:
    - trade_date: normalized MMD date
    - maturity_bucket: Short / Intermediate / Long / Extended Long
    - benchmark_rating
    - mmd_tenor
    - benchmark_yield
    - rating_spread_bps
    """
    if mmd_df.empty:
        return pd.DataFrame()

    date_col = _detect_mmd_date_column(mmd_df)
    if date_col is None:
        return pd.DataFrame()

    frames = []
    mmd_base = mmd_df.copy()
    mmd_base[date_col] = pd.to_datetime(mmd_base[date_col], errors="coerce")
    mmd_base = mmd_base.dropna(subset=[date_col])

    for bucket, tenor in MMD_BUCKET_MAP.items():
        if bucket == "All":
            continue
        benchmark_yield, meta = get_benchmark_curve(mmd_base, tenor, rating)
        if benchmark_yield is None:
            continue
        frames.append(
            pd.DataFrame(
                {
                    "trade_date": mmd_base[date_col].dt.normalize(),
                    "maturity_bucket": bucket,
                    "benchmark_rating": rating,
                    "mmd_tenor": tenor,
                    "benchmark_yield": benchmark_yield,
                    "rating_spread_bps": meta.get("rating_spread_bps"),
                    "benchmark_source": meta.get("benchmark_source"),
                    "source_column": meta.get("source_column"),
                }
            )
        )

    return pd.concat(frames, ignore_index=True) if frames else pd.DataFrame()


def build_spread_observations(
    market_df: pd.DataFrame,
    mmd_df: pd.DataFrame,
    issuer: str,
    rating: str,
) -> pd.DataFrame:
    """Build daily issuer spread observations by maturity bucket.

    Spread is calculated in basis points:
    (average issuer trade yield - synthetic benchmark yield) * 100.
    """
    required_cols = {"issuer", "trade_date", "maturity_bucket", "yield"}
    if market_df.empty or mmd_df.empty or not required_cols.issubset(set(market_df.columns)):
        return pd.DataFrame()

    issuer_df = market_df[market_df["issuer"] == issuer].copy()
    issuer_df = issuer_df[issuer_df["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)]
    if issuer_df.empty:
        return pd.DataFrame()

    issuer_df["trade_date"] = pd.to_datetime(issuer_df["trade_date"], errors="coerce").dt.normalize()
    issuer_df["yield"] = pd.to_numeric(issuer_df["yield"], errors="coerce")
    issuer_df = issuer_df.dropna(subset=["trade_date", "yield", "maturity_bucket"])

    daily_issuer = (
        issuer_df.groupby(["trade_date", "maturity_bucket"], as_index=False)
        .agg(
            avg_yield=("yield", "mean"),
            trade_count=("yield", "count"),
            total_trade_amount=("trade_amount", "sum") if "trade_amount" in issuer_df.columns else ("yield", "count"),
        )
    )

    benchmark_long = make_benchmark_long(mmd_df, rating)
    if benchmark_long.empty:
        return pd.DataFrame()

    spread_obs = daily_issuer.merge(
        benchmark_long,
        on=["trade_date", "maturity_bucket"],
        how="inner",
    )
    if spread_obs.empty:
        return pd.DataFrame()

    spread_obs["spread_to_benchmark_bps"] = (
        spread_obs["avg_yield"] - spread_obs["benchmark_yield"]
    ) * 100
    return spread_obs.sort_values(["maturity_bucket", "trade_date"])


def build_spread_movement_heatmap_data(
    spread_obs: pd.DataFrame,
    windows: dict[str, int] | None = None,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    """Return heatmap matrix and audit table for spread movement.

    For each maturity bucket and lookback window:
    Spread movement = latest available spread - historical spread at/before target date.

    Positive value means widening; negative value means tightening.
    """
    if windows is None:
        windows = {"1W": 7, "1M": 30, "3M": 90, "6M": 180, "1Y": 365}

    maturity_order = MATURITY_BUCKET_ORDER
    matrix = pd.DataFrame(index=maturity_order, columns=list(windows.keys()), dtype="float")
    audit_rows = []

    if spread_obs.empty:
        return matrix, pd.DataFrame(audit_rows)

    obs = spread_obs.copy()
    obs["trade_date"] = pd.to_datetime(obs["trade_date"], errors="coerce").dt.normalize()
    obs = obs.dropna(subset=["trade_date", "spread_to_benchmark_bps"])

    for bucket in maturity_order:
        bucket_obs = obs[obs["maturity_bucket"] == bucket].sort_values("trade_date")
        if bucket_obs.empty:
            continue

        latest_row = bucket_obs.iloc[-1]
        latest_date = latest_row["trade_date"]
        latest_spread = latest_row["spread_to_benchmark_bps"]

        for label, days in windows.items():
            target_date = latest_date - pd.Timedelta(days=days)
            historical_candidates = bucket_obs[bucket_obs["trade_date"] <= target_date]
            if historical_candidates.empty:
                audit_rows.append(
                    {
                        "maturity_bucket": bucket,
                        "window": label,
                        "latest_date": latest_date,
                        "latest_spread_bps": latest_spread,
                        "target_date": target_date,
                        "historical_date": pd.NaT,
                        "historical_spread_bps": pd.NA,
                        "spread_movement_bps": pd.NA,
                        "note": "No historical observation at or before target date",
                    }
                )
                continue

            historical_row = historical_candidates.iloc[-1]
            historical_date = historical_row["trade_date"]
            historical_spread = historical_row["spread_to_benchmark_bps"]
            movement = latest_spread - historical_spread
            matrix.loc[bucket, label] = movement
            audit_rows.append(
                {
                    "maturity_bucket": bucket,
                    "window": label,
                    "latest_date": latest_date,
                    "latest_spread_bps": latest_spread,
                    "target_date": target_date,
                    "historical_date": historical_date,
                    "historical_spread_bps": historical_spread,
                    "spread_movement_bps": movement,
                    "note": "Positive = widening; negative = tightening",
                }
            )

    return matrix, pd.DataFrame(audit_rows)


def build_spread_level_data(
    market_df: pd.DataFrame,
    mmd_df: pd.DataFrame,
    issuer: str,
    ratings: list[str],
) -> tuple[pd.DataFrame, pd.DataFrame]:
    """Return current spread level matrix and audit table.

    Matrix rows are maturity buckets; columns are benchmark ratings.
    Each cell is the latest available issuer spread to that benchmark, in bps:
        (Average Issuer Trade Yield - Synthetic Benchmark Yield) * 100

    This is different from spread movement. Spread level answers "where is
    the issuer trading now?" Movement answers "how much did it change?"
    """
    maturity_order = MATURITY_BUCKET_ORDER
    clean_ratings = [r for r in ratings if r in BENCHMARK_RATINGS]
    matrix = pd.DataFrame(index=maturity_order, columns=clean_ratings, dtype="float")
    audit_rows: list[dict] = []

    if not clean_ratings or market_df.empty or mmd_df.empty:
        return matrix, pd.DataFrame(audit_rows)

    for rating in clean_ratings:
        spread_obs = build_spread_observations(
            market_df=market_df,
            mmd_df=mmd_df,
            issuer=issuer,
            rating=rating,
        )
        if spread_obs.empty:
            continue

        spread_obs = spread_obs.copy()
        spread_obs["trade_date"] = pd.to_datetime(spread_obs["trade_date"], errors="coerce").dt.normalize()
        spread_obs = spread_obs.dropna(subset=["trade_date", "spread_to_benchmark_bps"])

        for bucket in maturity_order:
            bucket_obs = spread_obs[spread_obs["maturity_bucket"] == bucket].sort_values("trade_date")
            if bucket_obs.empty:
                audit_rows.append(
                    {
                        "maturity_bucket": bucket,
                        "benchmark_rating": rating,
                        "latest_date": pd.NaT,
                        "avg_yield": pd.NA,
                        "benchmark_yield": pd.NA,
                        "spread_to_benchmark_bps": pd.NA,
                        "mmd_tenor": MMD_BUCKET_MAP.get(bucket),
                        "rating_spread_bps": RATING_SPREADS.get(rating, RATING_SPREADS["AAA"]).get(MMD_BUCKET_MAP.get(bucket, "10Y"), 0.00) * 100,
                        "benchmark_source": "No matching benchmark/date",
                        "source_column": pd.NA,
                        "trade_count": pd.NA,
                        "total_trade_amount": pd.NA,
                        "note": "No overlapping issuer trade and benchmark observation",
                    }
                )
                continue

            latest = bucket_obs.iloc[-1]
            spread_level = latest["spread_to_benchmark_bps"]
            matrix.loc[bucket, rating] = spread_level
            audit_rows.append(
                {
                    "maturity_bucket": bucket,
                    "benchmark_rating": rating,
                    "latest_date": latest["trade_date"],
                    "avg_yield": latest.get("avg_yield"),
                    "benchmark_yield": latest.get("benchmark_yield"),
                    "spread_to_benchmark_bps": spread_level,
                    "mmd_tenor": latest.get("mmd_tenor"),
                    "rating_spread_bps": latest.get("rating_spread_bps"),
                    "benchmark_source": latest.get("benchmark_source"),
                    "source_column": latest.get("source_column"),
                    "trade_count": latest.get("trade_count"),
                    "total_trade_amount": latest.get("total_trade_amount"),
                    "note": "Latest available spread observation for maturity bucket and benchmark",
                }
            )

    return matrix, pd.DataFrame(audit_rows)


def build_issuer_curve_snapshot(
    market_df: pd.DataFrame,
    mmd_df: pd.DataFrame,
    issuer: str,
    ratings: list[str],
    as_of_date: pd.Timestamp,
    lookback_days: int,
    aggregation_method: str,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    """Build issuer yield curve vs benchmark curves by maturity bucket.

    This is a cross-sectional curve snapshot, not a time-series chart.

    Issuer curve logic:
    - Average Last N Days: average uploaded trade yield by maturity bucket over the
      lookback window ending on the selected as-of date.
    - Latest Trade Per Bucket: latest available trade observation at or before the
      selected as-of date for each maturity bucket.

    Benchmark curve logic:
    - Use uploaded rating curve columns when available.
    - Otherwise use MMD/AAA + visible rating spread assumptions.
    - For each bucket/rating, use the latest benchmark observation at or before
      the selected as-of date.
    """
    maturity_order = MATURITY_BUCKET_ORDER
    clean_ratings = [r for r in ratings if r in BENCHMARK_RATINGS]

    if market_df.empty or mmd_df.empty or not clean_ratings:
        return pd.DataFrame(), pd.DataFrame()

    required_cols = {"issuer", "trade_date", "maturity_bucket", "yield"}
    if not required_cols.issubset(set(market_df.columns)):
        return pd.DataFrame(), pd.DataFrame()

    as_of_date = pd.to_datetime(as_of_date).normalize()
    issuer_df = market_df[market_df["issuer"] == issuer].copy()
    issuer_df = issuer_df[issuer_df["maturity_bucket"].isin(maturity_order)]
    issuer_df["trade_date"] = pd.to_datetime(issuer_df["trade_date"], errors="coerce").dt.normalize()
    issuer_df["yield"] = pd.to_numeric(issuer_df["yield"], errors="coerce")
    issuer_df = issuer_df.dropna(subset=["trade_date", "yield", "maturity_bucket"])
    issuer_df = issuer_df[issuer_df["trade_date"] <= as_of_date]
    if issuer_df.empty:
        return pd.DataFrame(), pd.DataFrame()

    if aggregation_method == "Latest trade per bucket":
        latest_rows = (
            issuer_df.sort_values(["maturity_bucket", "trade_date"])
            .groupby("maturity_bucket", as_index=False)
            .tail(1)
        )
        issuer_curve = latest_rows[["maturity_bucket", "trade_date", "yield"]].rename(
            columns={"trade_date": "issuer_observation_date", "yield": "issuer_yield"}
        )
        counts = issuer_df.groupby("maturity_bucket", as_index=False).agg(trade_count=("yield", "count"))
        issuer_curve = issuer_curve.merge(counts, on="maturity_bucket", how="left")
        issuer_curve["aggregation_method"] = aggregation_method
        issuer_curve["lookback_start"] = pd.NaT
        issuer_curve["lookback_end"] = as_of_date
    else:
        lookback_start = as_of_date - pd.Timedelta(days=int(lookback_days))
        window_df = issuer_df[(issuer_df["trade_date"] >= lookback_start) & (issuer_df["trade_date"] <= as_of_date)].copy()
        if window_df.empty:
            return pd.DataFrame(), pd.DataFrame()
        agg_dict = {
            "issuer_yield": ("yield", "mean"),
            "trade_count": ("yield", "count"),
            "issuer_observation_date": ("trade_date", "max"),
        }
        if "trade_amount" in window_df.columns:
            agg_dict["total_trade_amount"] = ("trade_amount", "sum")
        issuer_curve = window_df.groupby("maturity_bucket", as_index=False).agg(**agg_dict)
        issuer_curve["aggregation_method"] = f"Average last {lookback_days} days"
        issuer_curve["lookback_start"] = lookback_start
        issuer_curve["lookback_end"] = as_of_date

    # Preserve intuitive curve order.
    issuer_curve["maturity_bucket"] = pd.Categorical(
        issuer_curve["maturity_bucket"], categories=maturity_order, ordered=True
    )
    issuer_curve = issuer_curve.sort_values("maturity_bucket")

    date_col = _detect_mmd_date_column(mmd_df)
    if date_col is None:
        return pd.DataFrame(), pd.DataFrame()

    mmd_base = mmd_df.copy()
    mmd_base[date_col] = pd.to_datetime(mmd_base[date_col], errors="coerce").dt.normalize()
    mmd_base = mmd_base.dropna(subset=[date_col])
    mmd_base = mmd_base[mmd_base[date_col] <= as_of_date]
    if mmd_base.empty:
        return pd.DataFrame(), pd.DataFrame()

    rows = []
    for rating in clean_ratings:
        for bucket in maturity_order:
            tenor = MMD_BUCKET_MAP.get(bucket, "10Y")
            y, meta = get_benchmark_curve(mmd_base, tenor, rating)
            if y is None:
                continue
            tmp = mmd_base[[date_col]].copy()
            tmp["benchmark_yield"] = pd.to_numeric(y, errors="coerce")
            tmp = tmp.dropna(subset=["benchmark_yield"])
            if tmp.empty:
                continue
            latest_bench = tmp.iloc[-1]
            rows.append(
                {
                    "maturity_bucket": bucket,
                    "benchmark_rating": rating,
                    "benchmark_date": latest_bench[date_col],
                    "benchmark_yield": latest_bench["benchmark_yield"],
                    "mmd_tenor": tenor,
                    "benchmark_source": meta.get("benchmark_source"),
                    "source_column": meta.get("source_column"),
                    "rating_spread_bps": meta.get("rating_spread_bps"),
                }
            )

    benchmark_curve = pd.DataFrame(rows)
    if benchmark_curve.empty:
        return pd.DataFrame(), pd.DataFrame()

    curve_data = issuer_curve.merge(benchmark_curve, on="maturity_bucket", how="inner")
    if curve_data.empty:
        return pd.DataFrame(), pd.DataFrame()

    curve_data["spread_to_benchmark_bps"] = (
        curve_data["issuer_yield"] - curve_data["benchmark_yield"]
    ) * 100

    # Long format for one clean Plotly line chart.
    issuer_line = issuer_curve[["maturity_bucket", "issuer_yield", "trade_count", "issuer_observation_date"]].copy()
    issuer_line = issuer_line.rename(columns={"issuer_yield": "yield_value"})
    issuer_line["curve"] = f"{issuer} issuer curve"
    issuer_line["curve_type"] = "Issuer"

    benchmark_line = benchmark_curve.rename(columns={"benchmark_yield": "yield_value"}).copy()
    benchmark_line["curve"] = benchmark_line["benchmark_rating"].astype(str) + " benchmark curve"
    benchmark_line["curve_type"] = "Benchmark"
    benchmark_line["trade_count"] = pd.NA
    benchmark_line["issuer_observation_date"] = pd.NaT

    plot_df = pd.concat(
        [
            issuer_line[["maturity_bucket", "yield_value", "curve", "curve_type", "trade_count", "issuer_observation_date"]],
            benchmark_line[["maturity_bucket", "yield_value", "curve", "curve_type", "trade_count", "issuer_observation_date"]],
        ],
        ignore_index=True,
    )
    plot_df["maturity_bucket"] = pd.Categorical(plot_df["maturity_bucket"], categories=maturity_order, ordered=True)
    plot_df = plot_df.sort_values(["curve_type", "curve", "maturity_bucket"])

    return plot_df, curve_data


def _normalize_col_name(name: object) -> str:
    """Normalize external column names so Munipro/Excel variants can be detected."""
    text = str(name).strip().lower()
    for ch in ["_", "-", "/", "\\", "\n", "\t"]:
        text = text.replace(ch, " ")
    return " ".join(text.split())


def _find_column(df: pd.DataFrame, canonical_name: str) -> str | None:
    """Return the actual uploaded column matching a canonical internal field."""
    normalized_columns = {_normalize_col_name(c): c for c in df.columns}
    aliases = COLUMN_ALIASES.get(canonical_name, [canonical_name])
    for alias in aliases:
        hit = normalized_columns.get(_normalize_col_name(alias))
        if hit is not None:
            return hit
    return None


def build_column_mapping(df: pd.DataFrame, expected_fields: list[str]) -> dict[str, str | None]:
    return {field: _find_column(df, field) for field in expected_fields}


def validate_dataset(
    df: pd.DataFrame,
    dataset_name: str,
    required_fields: list[str],
    recommended_fields: list[str],
    optional_fields: list[str] | None = None,
) -> dict:
    """Create a file-readiness report without blocking on non-critical fields."""
    optional_fields = optional_fields or []
    all_fields = required_fields + recommended_fields + optional_fields
    mapping = build_column_mapping(df, all_fields)

    missing_required = [field for field in required_fields if mapping.get(field) is None]
    missing_recommended = [field for field in recommended_fields if mapping.get(field) is None]
    detected_required = [field for field in required_fields if mapping.get(field) is not None]
    detected_recommended = [field for field in recommended_fields if mapping.get(field) is not None]

    return {
        "dataset": dataset_name,
        "can_run": len(missing_required) == 0,
        "row_count": len(df),
        "column_count": len(df.columns),
        "mapping": mapping,
        "missing_required": missing_required,
        "missing_recommended": missing_recommended,
        "detected_required": detected_required,
        "detected_recommended": detected_recommended,
    }


def validate_basic_values(df: pd.DataFrame, mapping: dict[str, str | None], dataset_type: str) -> list[str]:
    """Soft data-quality checks. These generate warnings instead of killing the app."""
    warnings: list[str] = []

    cusip_col = mapping.get("cusip")
    if cusip_col and cusip_col in df.columns:
        blank_cusips = df[cusip_col].isna().sum() + (df[cusip_col].astype(str).str.strip() == "").sum()
        if blank_cusips:
            warnings.append(f"{blank_cusips:,} row(s) have blank CUSIP values.")

    date_field = "maturity" if dataset_type == "bond" else "trade_date"
    date_col = mapping.get(date_field)
    if date_col and date_col in df.columns:
        parsed = pd.to_datetime(df[date_col], errors="coerce")
        bad_dates = parsed.isna().sum()
        if bad_dates:
            warnings.append(f"{bad_dates:,} row(s) have invalid or blank {date_field} values.")

    yield_col = mapping.get("yield")
    if yield_col and yield_col in df.columns:
        parsed_yield = pd.to_numeric(df[yield_col], errors="coerce")
        bad_yields = parsed_yield.isna().sum()
        extreme_yields = ((parsed_yield < -5) | (parsed_yield > 30)).sum()
        if bad_yields:
            warnings.append(f"{bad_yields:,} row(s) have non-numeric yield values.")
        if extreme_yields:
            warnings.append(f"{extreme_yields:,} row(s) have yield values outside the expected -5% to 30% range.")

    amount_col = mapping.get("trade_amount") or mapping.get("outstanding_amount")
    if amount_col and amount_col in df.columns:
        parsed_amount = pd.to_numeric(df[amount_col], errors="coerce")
        negative_amounts = (parsed_amount < 0).sum()
        if negative_amounts:
            warnings.append(f"{negative_amounts:,} row(s) have negative amount values.")

    return warnings


def display_validation_report(title: str, report: dict, warnings: list[str] | None = None):
    """Render a user-facing readiness card in Streamlit."""
    warnings = warnings or []
    status_icon = "✅" if report["can_run"] else "❌"
    with st.expander(f"{status_icon} {title} readiness check", expanded=not report["can_run"]):
        st.caption(f"Rows: {report['row_count']:,} · Columns: {report['column_count']:,}")

        c1, c2, c3 = st.columns(3)
        c1.metric("Required detected", f"{len(report['detected_required'])}/{len(report['detected_required']) + len(report['missing_required'])}")
        c2.metric("Recommended detected", f"{len(report['detected_recommended'])}/{len(report['detected_recommended']) + len(report['missing_recommended'])}")
        c3.metric("Ready to run", "Yes" if report["can_run"] else "No")

        if report["missing_required"]:
            st.error("Missing required fields: " + ", ".join(report["missing_required"]))
        if report["missing_recommended"]:
            st.warning("Missing recommended fields: " + ", ".join(report["missing_recommended"]))
        if warnings:
            for warning in warnings:
                st.warning(warning)

        mapping_rows = [
            {"Internal Field": key, "Uploaded Column Detected": value or "—"}
            for key, value in report["mapping"].items()
        ]
        st.dataframe(pd.DataFrame(mapping_rows), use_container_width=True, hide_index=True)



def display_trade_centric_metadata_review(
    title: str,
    report: dict,
    warnings: list[str] | None = None,
    issuer_can_be_derived: bool = False,
    maturity_can_be_derived: bool = False,
):
    """Render bond/security metadata as an informational review, not a blocker.

    In the flexible combined-trade architecture, the trade file is the required
    source. Bond-master fields such as issuer, sector, coupon, and outstanding
    amount are useful metadata, but they are not required for the dashboard to
    run. This avoids showing a scary red readiness failure when the app can
    infer issuer/maturity from Security Description, Mty, or reference files.
    """
    warnings = warnings or []
    missing_required = list(report.get("missing_required", []))
    adjusted_missing = []
    for field in missing_required:
        if field == "issuer" and issuer_can_be_derived:
            continue
        if field == "maturity" and maturity_can_be_derived:
            continue
        adjusted_missing.append(field)

    metadata_status = "✅" if not adjusted_missing else "🟡"
    with st.expander(f"{metadata_status} {title}", expanded=False):
        st.caption(f"Rows: {report['row_count']:,} · Columns: {report['column_count']:,}")

        c1, c2, c3 = st.columns(3)
        c1.metric(
            "Bond fields detected",
            f"{len(report['detected_required'])}/{len(report['detected_required']) + len(report['missing_required'])}",
        )
        c2.metric(
            "Recommended metadata detected",
            f"{len(report['detected_recommended'])}/{len(report['detected_recommended']) + len(report['missing_recommended'])}",
        )
        c3.metric("Dashboard mode", "Trade-centric")

        st.info(
            "Running in flexible trade-centric mode. Issuer and bond metadata are "
            "inferred from Security Description, maturity fields, and reference files when available."
        )

        if issuer_can_be_derived and "issuer" in missing_required:
            st.success("Issuer column is missing, but issuer names can be inferred from Security Description / reference mapping.")
        if maturity_can_be_derived and "maturity" in missing_required:
            st.success("Maturity column is missing, but maturity can be inferred from available maturity-style columns.")
        if adjusted_missing:
            st.warning(
                "Some bond/security metadata is still unavailable: " + ", ".join(adjusted_missing) +
                ". Related sections will show N/A or use fallback logic."
            )
        if report.get("missing_recommended"):
            st.caption("Optional metadata not found: " + ", ".join(report["missing_recommended"]))
        if warnings:
            with st.expander("Metadata warnings", expanded=False):
                for warning in warnings:
                    st.warning(warning)

        mapping_rows = [
            {"Internal Field": key, "Uploaded Column Detected": value or "—"}
            for key, value in report["mapping"].items()
        ]
        st.dataframe(pd.DataFrame(mapping_rows), use_container_width=True, hide_index=True)


def template_download_button(columns: list[str], label: str, filename: str):
    template = pd.DataFrame(columns=columns)
    st.download_button(
        label=label,
        data=template.to_csv(index=False).encode("utf-8"),
        file_name=filename,
        mime="text/csv",
    )


@st.cache_data(show_spinner="Processing uploaded data...")
def process_uploads(
    bond_bytes: bytes,
    bond_name: str,
    trade_payloads: list[tuple[str, bytes]],
    issuer_mapping_payload: tuple[str, bytes] | None,
    mmd_payload: tuple[str, bytes] | None,
):
    bond_file = io.BytesIO(bond_bytes)
    raw_bonds = read_uploaded_file(bond_file, bond_name)
    bonds_df = standardize_bonds(raw_bonds)

    issuer_mapping_df = pd.DataFrame()
    if issuer_mapping_payload is not None:
        name, payload = issuer_mapping_payload
        raw_mapping = read_uploaded_file(io.BytesIO(payload), name)
        issuer_mapping_df = standardize_issuer_mapping(raw_mapping)

    issuer_master = build_issuer_master(bonds_df, issuer_mapping_df)

    trade_frames = []
    failed_files = []
    for trade_name, trade_bytes in trade_payloads:
        try:
            raw_trade = read_uploaded_file(io.BytesIO(trade_bytes), trade_name)
            trade_frames.append(standardize_trades(raw_trade, source_file=trade_name))
        except Exception as exc:
            failed_files.append((trade_name, str(exc)))

    trades_df = pd.concat(trade_frames, ignore_index=True) if trade_frames else pd.DataFrame()

    # Data Health metric: remove exact duplicate trade rows before analytics.
    before_dedup = len(trades_df)
    trades_df = trades_df.drop_duplicates().reset_index(drop=True)
    duplicates_removed = before_dedup - len(trades_df)

    market_df = merge_market_data(bonds_df, trades_df, issuer_master)

    mmd_df = pd.DataFrame()
    if mmd_payload is not None:
        name, payload = mmd_payload
        raw_mmd = read_uploaded_file(io.BytesIO(payload), name)
        mmd_df = standardize_mmd(raw_mmd)

    return bonds_df, trades_df, issuer_master, market_df, mmd_df, failed_files, duplicates_removed


def dataframe_download_button(df: pd.DataFrame, label: str, filename: str):
    if df.empty:
        return
    csv = df.to_csv(index=False).encode("utf-8")
    st.download_button(label=label, data=csv, file_name=filename, mime="text/csv")


# -----------------------------------------------------------------------------
# Local / Desktop Data Source Mode
# -----------------------------------------------------------------------------
# This DuckDB-oriented version no longer uses Streamlit file uploaders.
# Data stays outside GitHub by default.
#
# Expected local files:
#   data/processed/Trade_Output_Sample.csv
#   data/processed/issuers.csv        optional but recommended
#   data/processed/mmd.csv            optional but recommended
#
# Trade_Output_Sample.csv is treated as the combined source file containing both
# bond/security metadata and trade-history fields. It is standardized once as a
# bond master and once as a trade file, then merged by CUSIP through the existing
# analytics pipeline.
# -----------------------------------------------------------------------------

BASE_DIR = Path(__file__).resolve().parent
DEFAULT_REFERENCE_DATA_DIR = BASE_DIR / "data" / "processed"
DEFAULT_TRADE_OUTPUT_PATH = DEFAULT_REFERENCE_DATA_DIR / "Trade_Output_Sample.csv"
DEFAULT_DUCKDB_PATH = BASE_DIR / "data" / "muni_market.duckdb"


def _first_existing_file(folder: Path, candidates: list[str]) -> Path | None:
    """Return the first existing file from a list of candidate filenames."""
    for name in candidates:
        candidate = folder / name
        if candidate.exists():
            return candidate
    return None


def _safe_read_local_file(path: Path) -> pd.DataFrame:
    """Read a local CSV/XLSX/XLS file with the existing file reader."""
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")
    with path.open("rb") as f:
        return read_uploaded_file(io.BytesIO(f.read()), path.name)


def _pick_col(df: pd.DataFrame, canonical_name: str) -> str | None:
    """Find a column in a flexible Munipro / trade-output file."""
    return _find_column(df, canonical_name)


def _series_or_na(df: pd.DataFrame, col: str | None, default=pd.NA) -> pd.Series:
    """Return a column series if available; otherwise return an NA series."""
    if col is not None and col in df.columns:
        return df[col]
    return pd.Series([default] * len(df), index=df.index)


def _derive_issuer_from_description(desc: object) -> str:
    """Best-effort issuer extraction from Munipro-style Security Description.

    This is intentionally conservative: if no clear delimiter is found, it keeps
    the full description. An optional issuers.csv can later overwrite/clean the
    issuer and sector mapping.
    """
    text = "" if pd.isna(desc) else str(desc).strip()
    if not text:
        return "Unknown"
    upper = text.upper()
    cut_words = [
        " CAP APPREC", " CAP APPRECIATION", " CAP APPN", " CAP APP", " CAB",
        " ELECTION", " REV", " REF", " SER", " SERIES", " BOND", " GO ",
        " TAXABLE", " UNLTD", " LTD", " 20", " 19",
    ]
    cut_positions = [upper.find(w) for w in cut_words if upper.find(w) > 8]
    if cut_positions:
        text = text[:min(cut_positions)].strip(" ,-–—")
    return text or "Unknown"


# -----------------------------------------------------------------------------
# Issuer normalization / reference-data layer
# -----------------------------------------------------------------------------
# Combined trade exports often include issuer-like text inside Security Description.
# Without normalization, the same issuer can appear as multiple dropdown options
# because one bond description says "CAP APPN", another says "SER A", etc.
# This layer creates:
#   raw_issuer       = the original extracted issuer text
#   issuer_key       = a cleaned matching key
#   issuer           = standardized display name, using issuers.csv when available
#   issuer_mapped    = whether issuers.csv supplied a standard name
# It makes the dashboard behave more like a real reference-data system.

ISSUER_NOISE_PATTERNS = [
    # Capital-appreciation / CAB language should never define issuer identity.
    r"\bCAP(?:ITAL)?\s+APP(?:REC|RECIATION|RECIATN|N)?\b.*$",
    r"\bCAP\s+APP(?:REC|N)?\b.*$",
    r"\bCAP\s+A\b.*$",
    r"\bCAB(?:S)?\b.*$",
    r"\bZERO\s+CPN\b.*$",
    r"\bCURRENT\s+INT(?:EREST)?\b.*$",

    # Series / refunding / security-structure descriptors.
    r"\bREF(?:UNDING)?\b.*$",
    r"\bREV(?:ENUE)?\b.*$",
    r"\bGO\b.*$",
    r"\bG\s+O\b.*$",
    r"\bBDS?\b.*$",
    r"\bBONDS?\b.*$",
    r"\bSER(?:IES)?\s+[A-Z0-9\-]+\b.*$",
    r"\bSER(?:IES)?\b.*$",
    r"\bISSUE\b.*$",
    r"\bCERT(?:IFICATE)?S?\s+OF\s+PART(?:ICIPATION)?\b.*$",
    r"\bCOP(?:S)?\b.*$",
    r"\bLEASE\s+REV(?:ENUE)?\b.*$",

    # Election / tax / measure / dating language.
    r"\bELECTION\b.*$",
    r"\bMEASURE\b.*$",
    r"\bTAXABLE\b.*$",
    r"\bTAX\s+EXEMPT\b.*$",
    r"\bUNLTD\s+TAX\b.*$",
    r"\bLTD\s+TAX\b.*$",
    r"\bDATED\b.*$",
    r"\bDUE\b.*$",

    # Dates / years usually indicate bond series or maturity metadata.
    r"\b\d{1,2}/\d{1,2}/\d{2,4}\b.*$",
    r"\b(?:19|20)\d{2}\b.*$",
]

ISSUER_ENTITY_TERMINALS = [
    "SCHOOL DISTRICT",
    "HIGH SCHOOL DISTRICT",
    "UNIFIED SCHOOL DISTRICT",
    "UNION SCHOOL DISTRICT",
    "COMMUNITY COLLEGE DISTRICT",
    "COLLEGE DISTRICT",
    "WATER DISTRICT",
    "IRRIGATION DISTRICT",
    "SANITATION DISTRICT",
    "UTILITY DISTRICT",
    "DISTRICT",
    "AUTHORITY",
    "FINANCE AUTHORITY",
    "FINANCING AUTHORITY",
    "CITY",
    "COUNTY",
    "STATE",
]

COMMON_ISSUER_ABBREVIATIONS = {
    "CALIF": "CA",
    "CALIFORNIA": "CA",
    "CNTY": "COUNTY",
    "CTY": "CITY",
    "SCH": "SCHOOL",
    "DIST": "DISTRICT",
    "CMNTY": "COMMUNITY",
    "COMM": "COMMUNITY",
    "COLL": "COLLEGE",
    "JUNIOR": "JR",
    "UN": "UNION",
    "UNI": "UNIFIED",
    "UNIF": "UNIFIED",
    "AUTH": "AUTHORITY",
    "DEV": "DEVELOPMENT",
    "PUB": "PUBLIC",
    "FIN": "FINANCE",
    "CORP": "CORPORATION",
}


def normalize_issuer_key(value: object) -> str:
    """Create a stable issuer matching key from messy trade descriptions."""
    if pd.isna(value):
        return "UNKNOWN"

    text = str(value).upper().strip()
    if not text:
        return "UNKNOWN"

    # Remove punctuation and normalize separators before pattern stripping.
    text = re.sub(r"[\u2010-\u2015]", "-", text)
    text = re.sub(r"[,;:\(\)\[\]\{\}\.]", " ", text)
    text = re.sub(r"\s+", " ", text).strip()

    # Remove bond-structure suffixes that should not define issuer identity.
    for pattern in ISSUER_NOISE_PATTERNS:
        text = re.sub(pattern, "", text).strip()

    # Remove obvious coupon/date fragments if they slipped into the extracted text.
    text = re.sub(r"\b\d+(?:\.\d+)?\s*%?\b", "", text)
    text = re.sub(r"\s+", " ", text).strip(" -_/,")

    if not text:
        return "UNKNOWN"

    tokens = [COMMON_ISSUER_ABBREVIATIONS.get(tok, tok) for tok in text.split()]
    text = " ".join(tokens)
    text = re.sub(r"\s+", " ", text).strip()

    # Final canonicalization pass: after abbreviation expansion, remove any
    # remaining bond-structure suffixes and trim fragments after recognized
    # municipal entity endings. This collapses examples such as:
    #   ACALANES CA UNION HIGH SCHOOL DISTRICT CAP APPREC
    #   ACALANES CA UNION HIGH SCHOOL DISTRICT SERIES 2016
    # into one issuer key.
    for pattern in ISSUER_NOISE_PATTERNS:
        text = re.sub(pattern, "", text).strip()

    for terminal in sorted(ISSUER_ENTITY_TERMINALS, key=len, reverse=True):
        idx = text.find(terminal)
        if idx >= 0:
            end = idx + len(terminal)
            tail = text[end:].strip()
            if not tail or re.search(r"\b(CAP|APP|APPREC|APPRECIATION|CAB|SER|SERIES|REF|REFUNDING|REV|BOND|BONDS|BDS|ELECTION|MEASURE|TAX|DATED|DUE|\d{4})\b", tail):
                text = text[:end].strip()
                break

    text = re.sub(r"\s+", " ", text).strip(" -_/,")
    return text or "UNKNOWN"


def display_issuer_from_key(key: object) -> str:
    """Readable fallback when issuers.csv does not map an issuer_key."""
    if pd.isna(key):
        return "Unknown"
    key = str(key).strip()
    if not key or key == "UNKNOWN":
        return "Unknown"
    return key.title().replace(" Ca ", " CA ").replace(" Usd", " USD")


def build_issuer_reference_table(issuer_mapping_df: pd.DataFrame) -> pd.DataFrame:
    """Build a flexible issuer reference table from issuers.csv / issuer_mapping.csv.

    Supported columns, case-insensitive:
    - issuer_key: optional cleaned key
    - raw_issuer: optional raw matching text
    - standard_issuer / issuer / issuer_name: display name
    - sector
    - primary_type
    """
    if issuer_mapping_df is None or issuer_mapping_df.empty:
        return pd.DataFrame(columns=["issuer_key", "standard_issuer", "sector", "primary_type"])

    ref = issuer_mapping_df.copy()

    key_col = _find_column(ref, "issuer_key")
    raw_col = _find_column(ref, "raw_issuer")
    standard_col = (
        _find_column(ref, "standard_issuer")
        or _find_column(ref, "issuer")
        or _find_column(ref, "issuer_name")
    )
    sector_col = _find_column(ref, "sector")
    primary_col = _find_column(ref, "primary_type")

    if key_col:
        issuer_key = ref[key_col].apply(normalize_issuer_key)
    elif raw_col:
        issuer_key = ref[raw_col].apply(normalize_issuer_key)
    elif standard_col:
        issuer_key = ref[standard_col].apply(normalize_issuer_key)
    else:
        return pd.DataFrame(columns=["issuer_key", "standard_issuer", "sector", "primary_type"])

    out = pd.DataFrame({"issuer_key": issuer_key})
    out["standard_issuer"] = (
        ref[standard_col].astype(str).str.strip() if standard_col else out["issuer_key"].apply(display_issuer_from_key)
    )
    out["sector"] = ref[sector_col].astype(str).str.strip() if sector_col else pd.NA
    out["primary_type"] = ref[primary_col].astype(str).str.strip() if primary_col else pd.NA

    out = out.replace({"nan": pd.NA, "None": pd.NA, "": pd.NA})
    out = out.dropna(subset=["issuer_key"])
    out = out[out["issuer_key"] != "UNKNOWN"]
    out = out.drop_duplicates("issuer_key", keep="first")
    return out.reset_index(drop=True)


def apply_manual_issuer_override_columns(
    df: pd.DataFrame,
    issuer_reference: pd.DataFrame | None = None,
) -> pd.DataFrame:
    """Apply persistent manual issuer overrides from issuers.csv.

    This is intentionally stronger than ordinary fuzzy normalization. If
    issuers.csv contains a row like:
        raw_issuer / issuer_key -> standard_issuer
    then every matching row is forced to that standard issuer for dropdowns,
    analytics, and sector grouping. This is the human-in-the-loop reference
    data override layer.
    """
    if df is None or df.empty or issuer_reference is None or issuer_reference.empty:
        return df

    out = df.copy()
    ref = issuer_reference.copy()
    if "issuer_key" not in ref.columns or "standard_issuer" not in ref.columns:
        return out

    ref = ref.dropna(subset=["issuer_key", "standard_issuer"]).copy()
    if ref.empty:
        return out

    # Normalize keys one more time to avoid stale/manual typo variants.
    ref["issuer_key"] = ref["issuer_key"].apply(normalize_issuer_key)
    key_to_standard = dict(zip(ref["issuer_key"], ref["standard_issuer"]))
    key_to_sector = dict(zip(ref["issuer_key"], ref["sector"])) if "sector" in ref.columns else {}
    key_to_primary = dict(zip(ref["issuer_key"], ref["primary_type"])) if "primary_type" in ref.columns else {}

    if "issuer_key" not in out.columns:
        if "raw_issuer" in out.columns:
            out["issuer_key"] = out["raw_issuer"].apply(normalize_issuer_key)
        elif "issuer" in out.columns:
            out["issuer_key"] = out["issuer"].apply(normalize_issuer_key)
        elif "description" in out.columns:
            out["issuer_key"] = out["description"].apply(_derive_issuer_from_description).apply(normalize_issuer_key)

    if "issuer_key" not in out.columns:
        return out

    mapped_standard = out["issuer_key"].map(key_to_standard)
    if "issuer" not in out.columns:
        out["issuer"] = out["issuer_key"].apply(display_issuer_from_key)
    out["issuer"] = mapped_standard.combine_first(out["issuer"])
    out["issuer_mapped"] = mapped_standard.notna() | out.get("issuer_mapped", False)

    if key_to_sector:
        mapped_sector = out["issuer_key"].map(key_to_sector)
        if "sector" not in out.columns:
            out["sector"] = mapped_sector
        else:
            out["sector"] = mapped_sector.combine_first(out["sector"])

    if key_to_primary:
        mapped_primary = out["issuer_key"].map(key_to_primary)
        if "primary_type" not in out.columns:
            out["primary_type"] = mapped_primary
        else:
            out["primary_type"] = mapped_primary.combine_first(out["primary_type"])

    return out


def apply_issuer_normalization_layer(
    df: pd.DataFrame,
    issuer_reference: pd.DataFrame | None = None,
) -> pd.DataFrame:
    """Normalize issuer names and optionally map them to standard reference data."""
    if df is None or df.empty:
        return df

    out = df.copy()

    if "issuer" in out.columns:
        raw = out["issuer"].astype(str).replace({"nan": "", "None": ""}).str.strip()
    elif "description" in out.columns:
        raw = out["description"].apply(_derive_issuer_from_description)
    else:
        raw = pd.Series(["Unknown"] * len(out), index=out.index)

    # Preserve original extracted name for audit/review.
    if "raw_issuer" not in out.columns:
        out["raw_issuer"] = raw

    out["issuer_key"] = raw.apply(normalize_issuer_key)
    out["issuer_clean_fallback"] = out["issuer_key"].apply(display_issuer_from_key)

    if issuer_reference is not None and not issuer_reference.empty:
        ref = issuer_reference.copy()
        ref = ref.loc[:, ~ref.columns.duplicated()].copy()
        out = out.merge(ref, on="issuer_key", how="left", suffixes=("", "_ref"))
        out["issuer_mapped"] = out["standard_issuer"].notna()
        out["issuer"] = out["standard_issuer"].combine_first(out["issuer_clean_fallback"])

        if "sector_ref" in out.columns:
            if "sector" not in out.columns:
                out["sector"] = out["sector_ref"]
            else:
                out["sector"] = out["sector_ref"].combine_first(out["sector"])
            out = out.drop(columns=["sector_ref"])

        if "primary_type_ref" in out.columns:
            if "primary_type" not in out.columns:
                out["primary_type"] = out["primary_type_ref"]
            else:
                out["primary_type"] = out["primary_type_ref"].combine_first(out["primary_type"])
            out = out.drop(columns=["primary_type_ref"])
    else:
        out["issuer_mapped"] = False
        out["issuer"] = out["issuer_clean_fallback"]

    if "sector" not in out.columns:
        out["sector"] = "Unknown"
    out["sector"] = out["sector"].fillna("Unknown").replace({"nan": "Unknown", "": "Unknown"})

    if "primary_type" not in out.columns:
        out["primary_type"] = "Unknown"
    out["primary_type"] = out["primary_type"].fillna("Unknown").replace({"nan": "Unknown", "": "Unknown"})

    return out


def build_unmapped_issuer_review_table(market_df: pd.DataFrame) -> pd.DataFrame:
    """Summarize normalized but unmapped issuers for reference-data cleanup."""
    if market_df is None or market_df.empty or "issuer_key" not in market_df.columns:
        return pd.DataFrame()

    review = market_df.copy()
    if "issuer_mapped" in review.columns:
        review = review[~review["issuer_mapped"].fillna(False)]

    agg_map = {
        "trade_count": ("cusip", "count") if "cusip" in review.columns else ("issuer_key", "count"),
    }
    if "trade_amount" in review.columns:
        agg_map["total_trade_amount"] = ("trade_amount", "sum")
    if "raw_issuer" in review.columns:
        agg_map["sample_raw_issuer"] = ("raw_issuer", "first")
    if "issuer" in review.columns:
        agg_map["suggested_issuer"] = ("issuer", "first")
    if "sector" in review.columns:
        agg_map["current_sector"] = ("sector", "first")

    out = review.groupby("issuer_key", as_index=False).agg(**agg_map)
    if "total_trade_amount" in out.columns:
        out = out.sort_values(["trade_count", "total_trade_amount"], ascending=False)
    else:
        out = out.sort_values("trade_count", ascending=False)
    return out.reset_index(drop=True)



def _unique_text_values(df: pd.DataFrame, column: str) -> list[str]:
    """Return clean unique text values from a possibly duplicated column name."""
    if df is None or df.empty or column not in df.columns:
        return []
    series_or_df = df.loc[:, column]
    if isinstance(series_or_df, pd.DataFrame):
        series = series_or_df.iloc[:, 0]
    else:
        series = series_or_df
    values = (
        series.dropna()
        .astype(str)
        .str.strip()
        .replace({"": pd.NA, "nan": pd.NA, "None": pd.NA, "Unknown": pd.NA, "UNKNOWN": pd.NA})
        .dropna()
        .unique()
        .tolist()
    )
    return sorted(set(values))


def _load_existing_issuer_mapping(mapping_path: Path) -> pd.DataFrame:
    """Load issuers.csv if it exists, otherwise create the expected schema."""
    base_cols = ["raw_issuer", "issuer_key", "standard_issuer", "sector", "primary_type"]
    if mapping_path.exists():
        try:
            existing = pd.read_csv(mapping_path)
        except Exception:
            existing = pd.DataFrame(columns=base_cols)
    else:
        existing = pd.DataFrame(columns=base_cols)

    for col in base_cols:
        if col not in existing.columns:
            existing[col] = pd.NA
    return existing


def save_manual_issuer_merge_mapping(
    mapping_path: Path,
    selected_names: list[str],
    canonical_issuer: str,
    sector: str = "Unknown",
    primary_type: str = "Unknown",
) -> pd.DataFrame:
    """Append manual issuer merge rows to issuers.csv.

    The saved file becomes a reference-data table. On the next rerun, the issuer
    normalization layer maps every selected raw/detected name to the canonical
    issuer selected by the user.
    """
    mapping_path = Path(mapping_path)
    mapping_path.parent.mkdir(parents=True, exist_ok=True)

    canonical_issuer = str(canonical_issuer).strip()
    if not canonical_issuer:
        raise ValueError("Canonical issuer name cannot be blank.")
    if not selected_names:
        raise ValueError("Select at least one issuer name to merge.")

    sector = str(sector).strip() or "Unknown"
    primary_type = str(primary_type).strip() or "Unknown"

    new_rows = []
    for raw_name in selected_names:
        raw_name = str(raw_name).strip()
        if not raw_name:
            continue
        new_rows.append(
            {
                "raw_issuer": raw_name,
                "issuer_key": normalize_issuer_key(raw_name),
                "standard_issuer": canonical_issuer,
                "sector": sector,
                "primary_type": primary_type,
            }
        )

    if not new_rows:
        raise ValueError("No valid issuer names were selected.")

    existing = _load_existing_issuer_mapping(mapping_path)
    updated = pd.concat([existing, pd.DataFrame(new_rows)], ignore_index=True)

    # Keep the latest manual decision for each normalized issuer key.
    updated["issuer_key"] = updated["issuer_key"].apply(normalize_issuer_key)
    updated = updated.replace({"nan": pd.NA, "None": pd.NA, "": pd.NA})
    updated = updated.dropna(subset=["issuer_key"])
    updated = updated[updated["issuer_key"] != "UNKNOWN"]
    updated = updated.drop_duplicates(subset=["issuer_key"], keep="last")
    updated = updated.sort_values(["standard_issuer", "issuer_key"]).reset_index(drop=True)
    updated.to_csv(mapping_path, index=False)
    return updated


def upsert_single_issuer_reference(
    reference_data_dir: Path,
    issuer: str,
    sector: str = "Unknown",
    primary_type: str = "Unknown",
) -> pd.DataFrame:
    """Persist a sector / primary-type correction for one selected issuer.

    This is the lightweight human-in-the-loop override used from the sidebar.
    It writes to data/processed/issuers.csv so the correction survives reruns
    and becomes part of the reference dataset.
    """
    mapping_path = Path(reference_data_dir) / "issuers.csv"
    return save_manual_issuer_merge_mapping(
        mapping_path=mapping_path,
        selected_names=[issuer],
        canonical_issuer=issuer,
        sector=sector,
        primary_type=primary_type,
    )


def infer_selected_issuer_sector(
    issuer: str,
    market_df: pd.DataFrame,
    issuer_master: pd.DataFrame,
) -> str:
    """Return the best available sector for the selected issuer."""
    selected_sector = "Unknown"
    if market_df is not None and not market_df.empty and "issuer" in market_df.columns and "sector" in market_df.columns:
        sector_values = (
            market_df.loc[market_df["issuer"].astype(str) == str(issuer), "sector"]
            .dropna()
            .astype(str)
            .str.strip()
        )
        sector_values = [x for x in sector_values.unique().tolist() if x and x.lower() != "nan"]
        if sector_values:
            selected_sector = sector_values[0]
    elif issuer_master is not None and not issuer_master.empty and "issuer" in issuer_master.columns and "sector" in issuer_master.columns:
        sector_values = (
            issuer_master.loc[issuer_master["issuer"].astype(str) == str(issuer), "sector"]
            .dropna()
            .astype(str)
            .str.strip()
        )
        sector_values = [x for x in sector_values.unique().tolist() if x and x.lower() != "nan"]
        if sector_values:
            selected_sector = sector_values[0]
    return selected_sector or "Unknown"


def render_manual_issuer_merge_tool(
    market_df: pd.DataFrame,
    reference_data_dir: Path,
):
    """Render a Streamlit UI for manually merging detected issuer names.

    This solves the unavoidable issuer-normalization problem: automated parsing
    is helpful, but a human reviewer should be able to override entity mapping
    and persist the result into issuers.csv.
    """
    mapping_path = Path(reference_data_dir) / "issuers.csv"

    with st.expander("Manual Issuer Merge Tool", expanded=False):
        st.markdown(
            """
Use this tool when several detected issuer names actually represent the same issuer.  
Select the messy/raw names, enter one canonical display name, and save. The mapping is written to `data/processed/issuers.csv` and will apply after the app reruns.
            """
        )

        if market_df is None or market_df.empty:
            st.info("No market data is loaded yet.")
            return

        candidate_values = []
        for col in ["raw_issuer", "issuer", "issuer_key"]:
            candidate_values.extend(_unique_text_values(market_df, col))
        candidate_values = sorted(set(candidate_values))

        if not candidate_values:
            st.info("No issuer candidates are available for manual merging.")
            return

        selected_names = st.multiselect(
            "Select issuer names to merge",
            candidate_values,
            help="Choose every raw/detected issuer label that should collapse into one standard issuer.",
        )

        suggested_name = ""
        if selected_names:
            suggested_name = display_issuer_from_key(normalize_issuer_key(selected_names[0]))

        canonical_issuer = st.text_input(
            "Canonical issuer name",
            value=suggested_name,
            placeholder="e.g. Acalanes Union High School District",
            help="This is the clean issuer name that will appear in the issuer dropdown.",
        )

        sector_values = ["Unknown"] + _unique_text_values(market_df, "sector")
        sector_values = list(dict.fromkeys(sector_values))
        sector = st.selectbox(
            "Sector",
            sector_values,
            index=0,
            help="Optional, but recommended for sector filtering and peer comparison.",
        )

        primary_type_values = ["Unknown"] + _unique_text_values(market_df, "primary_type")
        primary_type_values = list(dict.fromkeys(primary_type_values))
        primary_type = st.selectbox(
            "Primary Type",
            primary_type_values,
            index=0,
            help="Optional reference field. Leave Unknown if you do not need it yet.",
        )

        if selected_names:
            preview = pd.DataFrame(
                {
                    "selected_name": selected_names,
                    "issuer_key_to_save": [normalize_issuer_key(x) for x in selected_names],
                    "standard_issuer_to_save": canonical_issuer,
                    "sector_to_save": sector,
                    "primary_type_to_save": primary_type,
                }
            )
            st.caption("Preview of mapping rows to be saved")
            st.dataframe(preview, use_container_width=True, hide_index=True)

        save_col, file_col = st.columns([1, 2])
        with save_col:
            save_clicked = st.button("Save merge mapping", type="primary")
        with file_col:
            st.caption(f"Mapping file: `{mapping_path}`")

        if save_clicked:
            try:
                updated_mapping = save_manual_issuer_merge_mapping(
                    mapping_path=mapping_path,
                    selected_names=selected_names,
                    canonical_issuer=canonical_issuer,
                    sector=sector,
                    primary_type=primary_type,
                )
                st.success(
                    f"Saved {len(selected_names)} selected name(s) into issuers.csv. "
                    "The app will rerun and apply the merge now."
                )
                st.cache_data.clear()
                st.rerun()
            except Exception as exc:
                st.error(f"Could not save issuer merge mapping: {exc}")

        if mapping_path.exists():
            with st.expander("Current issuer mapping file preview", expanded=False):
                try:
                    current_mapping = pd.read_csv(mapping_path)
                    st.dataframe(current_mapping.tail(300), use_container_width=True, hide_index=True)
                    dataframe_download_button(current_mapping, "Download current issuers.csv", "issuers.csv")
                except Exception as exc:
                    st.warning(f"Could not preview issuers.csv: {exc}")


def _assign_maturity_bucket_from_years(years: float | int | None) -> str:
    if pd.isna(years):
        return "Unknown"
    if years <= 7:
        return "Short"
    if years <= 15:
        return "Intermediate"
    if years <= 25:
        return "Long"
    return "Extended Long"


def _add_required_optional_columns(df: pd.DataFrame) -> pd.DataFrame:
    """Make downstream analytics schema-flexible.

    Older dashboard sections expect bond-master fields such as coupon_bond,
    maturity_bond, and outstanding_amount. A combined trade-output file may not
    have them. This function creates safe placeholder / fallback columns so those
    sections degrade gracefully instead of crashing.
    """
    df = df.copy()

    if "coupon_bond" not in df.columns:
        if "coupon" in df.columns:
            df["coupon_bond"] = df["coupon"]
        elif "coupon_trade" in df.columns:
            df["coupon_bond"] = df["coupon_trade"]
        else:
            df["coupon_bond"] = pd.NA

    if "coupon_trade" not in df.columns:
        if "coupon" in df.columns:
            df["coupon_trade"] = df["coupon"]
        elif "coupon_bond" in df.columns:
            df["coupon_trade"] = df["coupon_bond"]
        else:
            df["coupon_trade"] = pd.NA

    if "maturity_bond" not in df.columns:
        if "maturity" in df.columns:
            df["maturity_bond"] = df["maturity"]
        elif "maturity_trade" in df.columns:
            df["maturity_bond"] = df["maturity_trade"]
        else:
            df["maturity_bond"] = pd.NaT

    if "maturity_trade" not in df.columns:
        if "maturity" in df.columns:
            df["maturity_trade"] = df["maturity"]
        elif "maturity_bond" in df.columns:
            df["maturity_trade"] = df["maturity_bond"]
        else:
            df["maturity_trade"] = pd.NaT

    if "outstanding_amount" not in df.columns:
        # Do not pretend trade amount is true outstanding. Use NA so turnover
        # calculations show as unavailable instead of being misleading.
        df["outstanding_amount"] = pd.NA

    if "par_amount" not in df.columns:
        df["par_amount"] = pd.NA

    if "trade_amount" not in df.columns:
        df["trade_amount"] = pd.NA

    if "sector" not in df.columns:
        df["sector"] = "Unknown"

    if "primary_type" not in df.columns:
        df["primary_type"] = "Unknown"

    if "maturity_bucket" not in df.columns:
        mat = pd.to_datetime(df.get("maturity_bond", pd.NaT), errors="coerce")
        ref_date = pd.to_datetime(df.get("trade_date", pd.Timestamp.today()), errors="coerce")
        if isinstance(ref_date, pd.Series):
            years = (mat - ref_date).dt.days / 365.25
        else:
            years = (mat - pd.Timestamp.today()).dt.days / 365.25
        df["maturity_bucket"] = years.apply(_assign_maturity_bucket_from_years)

    df["maturity_bucket"] = df["maturity_bucket"].replace(MATURITY_BUCKET_RENAME)
    return df


def build_combined_trade_output_frames(raw_df: pd.DataFrame, source_file: str) -> tuple[pd.DataFrame, pd.DataFrame]:
    """Build flexible bonds_df and trades_df from one combined trade-output file.

    This replaces the older assumption that users always have separate bond and
    trade files. Missing bond-master fields are allowed and handled as optional.
    """
    df = raw_df.copy()

    cusip_col = _pick_col(df, "cusip")
    trade_dt_col = _pick_col(df, "trade_datetime") or _pick_col(df, "trade_date")
    maturity_col = _pick_col(df, "maturity") or _pick_col(df, "maturity_trade")
    ytw_col = _pick_col(df, "yield")
    qty_col = _pick_col(df, "trade_amount")
    desc_col = _pick_col(df, "description")
    coupon_col = _pick_col(df, "coupon")
    price_col = _pick_col(df, "price")
    trade_type_col = _pick_col(df, "trade_type")
    issuer_col = _pick_col(df, "issuer")
    sector_col = _pick_col(df, "sector")

    if cusip_col is None:
        raise ValueError("Combined trade output must include a CUSIP/Cusip column.")
    if trade_dt_col is None:
        raise ValueError("Combined trade output must include TD & Time or Trade Date.")
    if ytw_col is None:
        raise ValueError("Combined trade output must include YTW/Yield.")

    trade_dt = pd.to_datetime(_series_or_na(df, trade_dt_col), errors="coerce")
    maturity = pd.to_datetime(_series_or_na(df, maturity_col), errors="coerce") if maturity_col else pd.Series(pd.NaT, index=df.index)
    desc = _series_or_na(df, desc_col, "")
    issuer = _series_or_na(df, issuer_col) if issuer_col else desc.apply(_derive_issuer_from_description)
    sector = _series_or_na(df, sector_col, "Unknown") if sector_col else pd.Series(["Unknown"] * len(df), index=df.index)
    coupon = pd.to_numeric(_series_or_na(df, coupon_col), errors="coerce") if coupon_col else pd.Series(pd.NA, index=df.index)
    ytw = pd.to_numeric(_series_or_na(df, ytw_col), errors="coerce")
    qty = pd.to_numeric(_series_or_na(df, qty_col), errors="coerce") if qty_col else pd.Series(pd.NA, index=df.index)
    price = pd.to_numeric(_series_or_na(df, price_col), errors="coerce") if price_col else pd.Series(pd.NA, index=df.index)

    rating_parts = []
    for rating_col in ["M", "S", "F", "Rating", "Ratings"]:
        if rating_col in df.columns:
            rating_parts.append(df[rating_col].astype(str).replace("nan", ""))
    ratings = rating_parts[0] if rating_parts else pd.Series(pd.NA, index=df.index)
    if len(rating_parts) > 1:
        ratings = rating_parts[0]
        for part in rating_parts[1:]:
            ratings = ratings.str.cat(part, sep="/")

    years = (maturity - trade_dt).dt.days / 365.25
    maturity_bucket = years.apply(_assign_maturity_bucket_from_years)

    trades_df = pd.DataFrame({
        "trade_datetime": trade_dt,
        "trade_date": trade_dt.dt.normalize(),
        "cusip": _series_or_na(df, cusip_col).astype(str).str.strip().str.upper(),
        "description": desc.astype(str),
        "maturity_trade": maturity,
        "coupon_trade": coupon,
        "yield": ytw,
        "price": price,
        "trade_amount": qty,
        "trade_type": _series_or_na(df, trade_type_col, pd.NA),
        "ratings_m_s_f": ratings,
        "source_file": source_file,
        "issuer": issuer.astype(str).str.strip(),
        "sector": sector.astype(str).str.strip(),
        "maturity_bucket": maturity_bucket,
    })

    bonds_df = pd.DataFrame({
        "issuer": issuer.astype(str).str.strip(),
        "sector": sector.astype(str).str.strip(),
        "primary_type": "Unknown",
        "cusip": _series_or_na(df, cusip_col).astype(str).str.strip().str.upper(),
        "maturity": maturity,
        "coupon": coupon,
        "outstanding_amount": pd.NA,
        "par_amount": pd.NA,
        "call_date": pd.NaT,
        "call_price": pd.NA,
        "fed_tax": pd.NA,
        "amt": pd.NA,
        "rating": ratings,
        "maturity_bucket": maturity_bucket,
    }).drop_duplicates(subset=["cusip"], keep="first")

    trades_df = _add_required_optional_columns(trades_df)
    bonds_df = _add_required_optional_columns(bonds_df)
    return bonds_df.reset_index(drop=True), trades_df.reset_index(drop=True)


@st.cache_data(show_spinner="Loading project data...")
def process_local_desktop_data(
    trade_output_path_str: str,
    reference_data_dir_str: str,
    duckdb_path_str: str,
):
    """Load project data and build dashboard-ready DataFrames.

    This keeps user data outside GitHub while preserving the existing analytics
    logic downstream.
    """
    trade_output_path = Path(trade_output_path_str).expanduser()
    reference_data_dir = Path(reference_data_dir_str).expanduser()
    duckdb_path = Path(duckdb_path_str).expanduser()

    failed_files: list[tuple[str, str]] = []

    # Optional: create / refresh a local DuckDB table from the combined CSV.
    # The dashboard still uses pandas DataFrames after standardization because
    # the existing analytics code expects DataFrames.
    try:
        import duckdb

        duckdb_path.parent.mkdir(parents=True, exist_ok=True)
        con = duckdb.connect(str(duckdb_path))
        con.execute(
            "CREATE OR REPLACE TABLE trade_output AS "
            "SELECT * FROM read_csv_auto(?)",
            [str(trade_output_path)],
        )
        con.close()
    except Exception as exc:
        # Do not block the dashboard if DuckDB refresh fails; the file can still
        # be loaded directly through the existing reader.
        failed_files.append(("DuckDB refresh", str(exc)))

    # Main combined source file.
    # Flexible combined-trade architecture: one Munipro-style trade output can
    # contain both trade fields and security metadata. Missing bond-master
    # columns are treated as optional and filled with safe placeholders.
    raw_trade_output = _safe_read_local_file(trade_output_path)
    bonds_df, trades_df = build_combined_trade_output_frames(
        raw_trade_output,
        source_file=trade_output_path.name,
    )

    before_dedup = len(trades_df)
    trades_df = trades_df.drop_duplicates().reset_index(drop=True)
    duplicates_removed = before_dedup - len(trades_df)

    # Optional issuer / sector reference file.
    issuer_mapping_df = pd.DataFrame()
    issuer_mapping_path = _first_existing_file(
        reference_data_dir,
        [
            "issuers.csv",
            "issuer.csv",
            "issuer_mapping.csv",
            "Issuer_Mapping.csv",
            "issuers.xlsx",
            "issuer_mapping.xlsx",
        ],
    )
    if issuer_mapping_path is not None:
        try:
            # Keep the raw issuer mapping schema instead of forcing the old
            # two-file uploader schema. This supports manual override columns
            # such as raw_issuer, issuer_key, standard_issuer, sector, and
            # primary_type. The flexible reference builder below handles all
            # common column variants.
            issuer_mapping_df = _safe_read_local_file(issuer_mapping_path)
        except Exception as exc:
            failed_files.append((issuer_mapping_path.name, str(exc)))

    # Build a reference-data normalization layer from issuers.csv / issuer_mapping.csv.
    # This prevents small Security Description differences from splitting the same
    # issuer into many dropdown entries.
    issuer_reference = build_issuer_reference_table(issuer_mapping_df)

    # First normalize the raw bond/trade issuer labels before merge.
    bonds_df = apply_issuer_normalization_layer(bonds_df, issuer_reference)
    trades_df = apply_issuer_normalization_layer(trades_df, issuer_reference)

    issuer_master = build_issuer_master(bonds_df, issuer_mapping_df)
    try:
        market_df = merge_market_data(bonds_df, trades_df, issuer_master)
    except Exception:
        # Fallback for trade-centric files where the old two-file merge schema
        # is too strict. Preserve all trade rows and attach security metadata by CUSIP.
        bond_merge_cols = [c for c in bonds_df.columns if c not in {"issuer", "sector", "maturity_bucket"}]
        bond_merge = bonds_df[["cusip"] + [c for c in bond_merge_cols if c != "cusip"]].drop_duplicates("cusip")
        market_df = trades_df.merge(bond_merge, on="cusip", how="left", suffixes=("_trade", "_bond"))

    # Normalize again after merge so the final displayed issuer is always the
    # standard_issuer from issuers.csv when available, otherwise a cleaned fallback.
    market_df = apply_issuer_normalization_layer(market_df, issuer_reference)

    # Strong human-in-the-loop override layer. This makes manually saved merge
    # decisions apply to every downstream dataframe, including the issuer
    # dropdown, peer ranking, and sector analytics.
    bonds_df = apply_manual_issuer_override_columns(bonds_df, issuer_reference)
    trades_df = apply_manual_issuer_override_columns(trades_df, issuer_reference)
    market_df = apply_manual_issuer_override_columns(market_df, issuer_reference)

    bonds_df = _add_required_optional_columns(bonds_df)
    trades_df = _add_required_optional_columns(trades_df)
    market_df = _add_required_optional_columns(market_df)

    # Rebuild issuer_master from the normalized output used by the dashboard.
    issuer_master_cols = [c for c in ["issuer", "sector", "primary_type", "issuer_key", "raw_issuer", "issuer_mapped"] if c in market_df.columns]
    issuer_master = (
        market_df[issuer_master_cols]
        .drop_duplicates(subset=["issuer"])
        .sort_values("issuer")
        .reset_index(drop=True)
        if issuer_master_cols and "issuer" in issuer_master_cols
        else issuer_master
    )

    # Optional MMD / benchmark curve file.
    mmd_df = pd.DataFrame()
    mmd_path = _first_existing_file(
        reference_data_dir,
        [
            "mmd.csv",
            "MMD.csv",
            "mmd_curve.csv",
            "MMD_Curve.csv",
            "benchmark_curve.csv",
            "mmd.xlsx",
            "mmd_curve.xlsx",
        ],
    )
    if mmd_path is not None:
        try:
            raw_mmd = _safe_read_local_file(mmd_path)
            mmd_df = standardize_mmd(raw_mmd)
        except Exception as exc:
            failed_files.append((mmd_path.name, str(exc)))

    data_source_summary = {
        "trade_output_path": str(trade_output_path),
        "reference_data_dir": str(reference_data_dir),
        "duckdb_path": str(duckdb_path),
        "issuer_mapping_path": str(issuer_mapping_path) if issuer_mapping_path else "Not found",
        "mmd_path": str(mmd_path) if mmd_path else "Not found",
        "issuer_reference_rows": int(len(issuer_reference)) if 'issuer_reference' in locals() else 0,
        "normalized_issuer_count": int(market_df["issuer"].nunique()) if "issuer" in market_df.columns else 0,
        "mapped_issuer_rows": int(market_df.get("issuer_mapped", pd.Series(dtype=bool)).fillna(False).sum()) if "issuer_mapped" in market_df.columns else 0,
    }

    return bonds_df, trades_df, issuer_master, market_df, mmd_df, failed_files, duplicates_removed, data_source_summary



# Sidebar Project Data Source controls are rendered at the bottom of the sidebar
# after issuer selection. Values are initialized here so the data pipeline can
# run before those controls are visually displayed.
if "trade_output_path_input" not in st.session_state:
    st.session_state["trade_output_path_input"] = str(DEFAULT_TRADE_OUTPUT_PATH)
if "reference_data_dir_input" not in st.session_state:
    st.session_state["reference_data_dir_input"] = str(DEFAULT_REFERENCE_DATA_DIR)
if "duckdb_path_input" not in st.session_state:
    st.session_state["duckdb_path_input"] = str(DEFAULT_DUCKDB_PATH)

trade_output_path_input = st.session_state["trade_output_path_input"]
reference_data_dir_input = st.session_state["reference_data_dir_input"]
duckdb_path_input = st.session_state["duckdb_path_input"]

# -----------------------------------------------------------------------------
# Local-readiness gate
# -----------------------------------------------------------------------------
trade_output_path = Path(trade_output_path_input).expanduser()
reference_data_dir = Path(reference_data_dir_input).expanduser()

section_anchor("file-readiness", "Project Data Readiness Check")

if not trade_output_path.exists():
    st.error(f"Required file not found: `{trade_output_path}`")
    st.info("Put `Trade_Output_Sample.csv` in `data/processed/`, or adjust the path in the sidebar.")
    st.stop()

if not reference_data_dir.exists():
    st.warning(
        f"Reference folder not found: `{reference_data_dir}`. The dashboard can still run, "
        "but issuer mapping and MMD comparison may be unavailable."
    )

raw_combined_preview = _safe_read_local_file(trade_output_path)

# In the combined trade-output architecture, the trade source is the only
# true blocker. Bond/security metadata is informational because issuer,
# maturity bucket, coupon, and other fields can often be inferred from
# Security Description, Mty, Cpn, and optional reference files.
bond_report = validate_dataset(
    raw_combined_preview,
    trade_output_path.name,
    BOND_REQUIRED,
    BOND_RECOMMENDED,
    BOND_OPTIONAL,
)
bond_warnings = validate_basic_values(raw_combined_preview, bond_report["mapping"], dataset_type="bond")

trade_report = validate_dataset(
    raw_combined_preview,
    trade_output_path.name,
    TRADE_REQUIRED,
    TRADE_RECOMMENDED,
    TRADE_OPTIONAL,
)
trade_warnings = validate_basic_values(raw_combined_preview, trade_report["mapping"], dataset_type="trade")

# Trade source readiness remains required.
display_validation_report("Combined File as Trade Source", trade_report, trade_warnings)

if not trade_report["can_run"]:
    st.error(
        "The project-data dashboard cannot run yet because the combined file is missing required trade fields. "
        "At minimum, it needs CUSIP, TD & Time/Trade Date, and YTW/Yield."
    )
    st.stop()

# Bond/security metadata review is informational only.
# A missing issuer column is acceptable if Security Description exists because
# the flexible pipeline can derive issuer names from the description or clean
# them with issuers.csv.
issuer_can_be_derived = bool(
    bond_report["mapping"].get("issuer")
    or trade_report["mapping"].get("description")
    or _find_column(raw_combined_preview, "description")
)
maturity_can_be_derived = bool(
    bond_report["mapping"].get("maturity")
    or trade_report["mapping"].get("maturity_trade")
    or _find_column(raw_combined_preview, "maturity_trade")
    or _find_column(raw_combined_preview, "maturity")
)

display_trade_centric_metadata_review(
    "Bond / Security Metadata Review",
    bond_report,
    bond_warnings,
    issuer_can_be_derived=issuer_can_be_derived,
    maturity_can_be_derived=maturity_can_be_derived,
)

st.success(
    "Running in flexible trade-centric mode. Issuer and bond metadata are inferred "
    "from Security Description / reference files when available."
)

with st.expander("Methodology: project data mode", expanded=False):
    st.markdown(
        """
- The dashboard no longer uses Streamlit upload boxes.
- `Trade_Output_Sample.csv` is treated as the unified local source for bond/security metadata and trade history.
- `issuers.csv` / `issuer_mapping.csv` and `mmd.csv` / `mmd_curve.csv` are optional reference files read from `Intern_Muni_Data`.
- The app refreshes a local DuckDB table named `trade_output`, while the existing analytics engine continues to use standardized DataFrames downstream.
- Real data stays outside GitHub by default.
        """
    )

(
    bonds_df,
    trades_df,
    issuer_master,
    market_df,
    mmd_df,
    failed_files,
    duplicates_removed,
    data_source_summary,
) = process_local_desktop_data(
    trade_output_path_str=str(trade_output_path_input),
    reference_data_dir_str=str(reference_data_dir_input),
    duckdb_path_str=str(duckdb_path_input),
)

# Normalize legacy tenor-style bucket labels into presentation-friendly curve sectors.
# This keeps the dashboard compatible with existing data_utils output while making
# the user-facing terminology clearer.
for _df_name, _df in [("bonds_df", bonds_df), ("trades_df", trades_df), ("market_df", market_df)]:
    if isinstance(_df, pd.DataFrame):
        _fixed = _add_required_optional_columns(_df)
        if _df_name == "bonds_df":
            bonds_df = _fixed
        elif _df_name == "trades_df":
            trades_df = _fixed
        else:
            market_df = _fixed

if failed_files:
    with st.warning("Some local reference files could not be processed."):
        st.write(failed_files)

with st.expander("Local data source summary", expanded=False):
    st.json(data_source_summary)

# Reference-data QA: show issuer strings that were normalized but not mapped by issuers.csv.
# This is designed as a maintenance queue: users can copy suggested rows into
# issuers.csv over time to make the issuer universe cleaner.
unmapped_issuer_review = build_unmapped_issuer_review_table(market_df)
with st.expander("Unmapped issuer review table", expanded=False):
    st.caption(
        "Rows below were normalized from Security Description but did not match issuers.csv. "
        "Add important names to issuers.csv with issuer_key, standard_issuer, sector, and primary_type to improve the dropdown."
    )
    if unmapped_issuer_review.empty:
        st.success("All issuer keys are mapped, or no issuer_key field is available.")
    else:
        st.dataframe(unmapped_issuer_review.head(500), use_container_width=True, hide_index=True)
        dataframe_download_button(
            unmapped_issuer_review,
            "Download unmapped issuer review CSV",
            "unmapped_issuer_review.csv",
        )


# Manual reference-data override: let users merge messy issuer labels into one
# canonical issuer without editing CSV files by hand.
render_manual_issuer_merge_tool(market_df, reference_data_dir)

if bonds_df.empty:
    st.error("No usable bond/security rows found. Check that Trade_Output_Sample.csv includes CUSIP, issuer, and maturity fields.")
    st.stop()

if market_df.empty:
    st.error("No usable trade rows found. Check that Trade_Output_Sample.csv includes CUSIP, trade date, and yield fields.")
    st.stop()

uploaded_issuers = sorted(market_df["issuer"].dropna().astype(str).unique().tolist())

if not uploaded_issuers:
    st.error("No issuer names were detected from the local combined file. Please check the issuer field or issuer mapping file.")
    st.stop()

st.success(
    f"Processed {len(bonds_df):,} bond/security rows and {len(market_df):,} merged trade rows "
    f"from project data. Detected {len(uploaded_issuers):,} issuer(s)."
)


# -----------------------------------------------------------------------------
# Main-page guidance and data health
# -----------------------------------------------------------------------------
section_anchor("dashboard-contents", "Contents")
st.markdown(
    """
<div class="nav-card">
<b>Data & Setup</b><br>
<a href="#file-readiness">1. Project Data Readiness</a><br>
<a href="#data-quality-scorecard">2. Data Quality Scorecard</a><br>
<a href="#executive-snapshot">3. Executive Snapshot</a><br><br>

<b>Benchmark / Spread Framework</b><br>
<a href="#yield-relative-value">4. Yield & Relative Value</a><br>
<a href="#issuer-curve">5. Issuer Curve vs Benchmark</a><br>
<a href="#spread-level">6. Current Spread Level</a><br>
<a href="#spread-attribution">7. Spread Attribution</a><br><br>

<b>Relative Value Signals</b><br>
<a href="#market-narrative">8. Market Narrative & Opportunity Map</a><br>
<a href="#peer-rv">9. Peer RV Comparison</a><br>
<a href="#cross-issuer-rv">10. Cross-Issuer RV Analytics</a><br>
<a href="#historical-spread">11. Historical Spread Percentile</a><br>
<a href="#recommendation-engine">12. Rule-Based Recommendation</a><br>
<a href="#ai-commentary-studio">13. AI Commentary Studio</a><br><br>

<b>Risk / Flow / Screening</b><br>
<a href="#curve-shape">14. Curve Shape Analytics</a><br>
<a href="#scenario-shock">15. Scenario Shock Analysis</a><br>
<a href="#dealer-proxy">16. Dealer Behavior Proxy</a><br>
<a href="#security-screener">17. Security Screener</a><br>
<a href="#watchlist">18. Watchlist / Saved Candidates</a><br><br>

<b>Bond-Level Drilldown</b><br>
<a href="#spread-movement">19. Spread Movement</a><br>
<a href="#cusip-drilldown">20. CUSIP Opportunity Drilldown</a><br>
<a href="#rv-positioning">21. RV Positioning Map</a><br>
<a href="#liquidity">22. Liquidity Analysis</a><br><br>

<b>Reference / Admin / Outputs</b><br>
<a href="#bond-master">23. Bond Master</a><br>
<a href="#trade-detail">24. Trade Detail</a><br>
<a href="#report-export-center">25. Report Export Center</a><br>
<a href="#export-summary">26. Export Summary</a><br>
<a href="#admin-methodology">27. Admin Methodology</a><br>
<a href="#version-changelog">28. Version / Change Log</a><br>
<a href="#downloads">29. Downloads</a>
</div>
""",
    unsafe_allow_html=True,
)

section_anchor("data-health-overview", "Data Health")
if not market_df.empty and "trade_date" in market_df.columns:
    trade_dates = pd.to_datetime(market_df["trade_date"], errors="coerce").dropna()
    if not trade_dates.empty:
        earliest_trade = trade_dates.min()
        latest_trade = trade_dates.max()
        data_coverage_text = f"{earliest_trade:%Y-%m-%d} → {latest_trade:%Y-%m-%d}"
    else:
        data_coverage_text = "No valid trade dates detected"
else:
    data_coverage_text = "No trade data loaded"

total_rows = len(market_df)
if total_rows > 0 and "cusip" in market_df.columns:
    bond_cusips = set(bonds_df["cusip"].dropna().astype(str).str.upper()) if "cusip" in bonds_df.columns else set()
    trade_cusips = market_df["cusip"].dropna().astype(str).str.upper()
    matched_cusips_count = trade_cusips.isin(bond_cusips).sum() if bond_cusips else 0
    match_rate = matched_cusips_count / total_rows * 100
else:
    matched_cusips_count = 0
    match_rate = 0

missing_issuers = market_df["issuer"].isna().sum() if "issuer" in market_df.columns else total_rows
missing_issuer_rate = missing_issuers / total_rows * 100 if total_rows > 0 else 0

health_cols = st.columns(5)
health_cols[0].metric("Data Coverage", data_coverage_text)
health_cols[1].metric("Trades Loaded", f"{len(market_df):,}")
health_cols[2].metric("CUSIP Match Rate", f"{match_rate:.1f}%")
health_cols[3].metric("Missing Issuers", f"{missing_issuers:,}")
health_cols[4].metric("Duplicate Trades Removed", f"{duplicates_removed:,}")

with st.expander("Data Health methodology", expanded=False):
    st.markdown(
        """
- **Data Coverage** uses the earliest and latest valid trade dates after standardization.
- **Trades Loaded** counts merged trade rows available for analytics.
- **CUSIP Match Rate** is the share of merged trade rows whose CUSIP appears in the uploaded bond master.
- **Missing Issuers** counts rows without an issuer after the bond/trade merge and issuer-mapping logic.
- **Duplicate Trades Removed** counts exact duplicate standardized trade rows removed before analytics.
        """
    )

# -----------------------------------------------------------------------------
# Sidebar issuer-first workflow
# -----------------------------------------------------------------------------
with st.sidebar:
    st.header("Select From Uploaded Issuers")
    selected_issuer = st.selectbox(
        "Issuer detected from uploaded files",
        uploaded_issuers,
        help="This list is generated from the processed project files."
    )

    selected_sector_sidebar = infer_selected_issuer_sector(selected_issuer, market_df, issuer_master)
    available_sectors = ["Unknown"] + _unique_text_values(market_df, "sector")
    default_sector_options = [
        "Education",
        "School District",
        "Utilities",
        "Water / Sewer",
        "Transportation",
        "Airport",
        "Healthcare",
        "Housing",
        "General Government",
        "Local Government",
        "Public Finance Authority",
        "Other",
    ]
    sector_options = list(dict.fromkeys(available_sectors + default_sector_options))
    if selected_sector_sidebar not in sector_options:
        sector_options.insert(1, selected_sector_sidebar)

    with st.expander("Issuer Sector Override", expanded=(selected_sector_sidebar == "Unknown")):
        st.caption("Use this when the issuer sector is Unknown or classified incorrectly. The correction is saved to issuers.csv.")
        sector_choice = st.selectbox(
            "Issuer Sector",
            sector_options,
            index=sector_options.index(selected_sector_sidebar) if selected_sector_sidebar in sector_options else 0,
        )
        custom_sector = st.text_input(
            "Custom Sector",
            value="" if sector_choice != "Other" else selected_sector_sidebar,
            placeholder="Enter a sector if Other is selected",
        )
        primary_type_choice = st.text_input(
            "Primary Type",
            value="Unknown",
            help="Optional reference field. Leave Unknown if you do not need it yet.",
        )
        final_sector_choice = custom_sector.strip() if sector_choice == "Other" and custom_sector.strip() else sector_choice

        if st.button("Save Sector to issuers.csv", type="primary"):
            try:
                upsert_single_issuer_reference(
                    reference_data_dir=reference_data_dir,
                    issuer=selected_issuer,
                    sector=final_sector_choice,
                    primary_type=primary_type_choice,
                )
                st.success(f"Saved: {selected_issuer} → {final_sector_choice}")
                st.cache_data.clear()
                st.rerun()
            except Exception as exc:
                st.error(f"Could not save issuer sector override: {exc}")

    # -----------------------------------------------------------------------------
    # Maturity Bucket Methodology
    # -----------------------------------------------------------------------------
    with st.expander("Maturity Bucket Methodology", expanded=False):
        st.markdown(
            """
<style>
.bucket-table {
    width: 100%;
    border-collapse: collapse;
    font-size: 0.86rem;
}

.bucket-table th,
.bucket-table td {
    border: 1px solid #d9dce3;
    padding: 8px 10px;
    vertical-align: middle;
    text-align: left;
}

.bucket-table th {
    font-weight: 700;
    background-color: #f6f7fb;
}

.bucket-col {
    width: 26%;
    white-space: nowrap;
    font-weight: 600;
}

.years-col {
    width: 24%;
    white-space: nowrap;
}

.interp-col {
    width: 50%;
}
</style>

### Bucket Definitions

<table class="bucket-table">
  <tr>
    <th>Bucket</th>
    <th>Years to Maturity</th>
    <th>Interpretation</th>
  </tr>
  <tr>
    <td class="bucket-col">All</td>
    <td class="years-col">All maturities</td>
    <td class="interp-col">Full uploaded trade universe</td>
  </tr>
  <tr>
    <td class="bucket-col">Short</td>
    <td class="years-col">≤ 7Y</td>
    <td class="interp-col">Front-end / lower-duration bonds</td>
  </tr>
  <tr>
    <td class="bucket-col">Intermediate</td>
    <td class="years-col">7–15Y</td>
    <td class="interp-col">Intermediate curve sector</td>
  </tr>
  <tr>
    <td class="bucket-col">Long</td>
    <td class="years-col">15–25Y</td>
    <td class="interp-col">Long-duration municipal sector</td>
  </tr>
  <tr>
    <td class="bucket-col">Extended Long</td>
    <td class="years-col">25Y+</td>
    <td class="interp-col">Long-end institutional duration sector</td>
  </tr>
</table>

<br>

### Why This Matters

These maturity buckets are used to:

- Compare issuer spreads across the municipal curve
- Analyze relative value positioning
- Evaluate duration sensitivity
- Review liquidity and trading activity by curve sector
- Standardize secondary-market analytics

### Important Notes

- Bucket labels represent practical curve sectors, not exact maturity points.
- The benchmark mapping uses the closest MMD tenor: Short → 5Y, Intermediate → 10Y, Long → 20Y, Extended Long → 30Y.
- The framework supports relative-value, curve, and liquidity analysis.
            """,
            unsafe_allow_html=True,
        )

    maturity_bucket = st.selectbox(
        "Maturity Bucket",
        MATURITY_BUCKET_OPTIONS,
        help="""
Bucket Definitions

• All = All available maturities
• Short = ≤ 7 Years
• Intermediate = 7–15 Years
• Long = 15–25 Years
• Extended Long = 25+ Years

Used for:
- Relative Value Analysis
- Yield Curve Positioning
- Secondary Market Trend Analysis

Buckets represent maturity ranges,
NOT exact maturity tenors.
"""
    )

    time_window = st.selectbox(
        "Time Window",
        ["All", "1Y", "3Y", "5Y"],
        help="""
Historical time range used for trend analysis.

• 1Y = Last 1 Year
• 3Y = Last 3 Years
• 5Y = Last 5 Years
• All = Entire uploaded dataset

Used for:
- Yield trend analysis
- Spread movement review
- Historical relative value comparison
"""
    )

    show_raw_tables = st.checkbox(
        "Show Raw Tables",
        value=False,
        help="""
Display underlying trade-level and bond-level data tables.

Useful for:
- Audit review
- Data validation
- Trade-level investigation
- CUSIP drilldowns
"""
    )

    st.markdown("---")
    st.header("Project Data Source")
    st.caption("This version reads project data files from data/processed instead of using upload boxes.")

    trade_output_path_input = st.text_input(
        "Combined Trade Output CSV",
        key="trade_output_path_input",
        help="Default: data/processed/Trade_Output_Sample.csv. This combined file can contain both security metadata and trade-history fields.",
    )

    reference_data_dir_input = st.text_input(
        "Reference Data Folder",
        key="reference_data_dir_input",
        help="Default: data/processed. Put issuers.csv and mmd.csv here.",
    )

    duckdb_path_input = st.text_input(
        "DuckDB Database Path",
        key="duckdb_path_input",
        help="The app refreshes a local DuckDB table named trade_output from your combined CSV.",
    )

    st.caption("Tip: for large/company data, keep the same folder structure locally or on a shared drive and do not commit real data to GitHub.")

    with st.expander("Expected local files"):
        st.markdown(
            """
**Required**
- `data/processed/Trade_Output_Sample.csv`

**Optional reference files** inside `data/processed/`
- `issuers.csv` or `issuer_mapping.csv`
- `mmd.csv` or `mmd_curve.csv`

The combined trade-output file is used as both the bond/security source and the trade-history source.
            """
        )

    with st.expander("Version / Change Log", expanded=False):
        st.markdown(
            """
**Current Version:** `v1.2-issuer-first-sector-override`

Recent additions:
- issuer selection moved to the top of the sidebar
- project data source controls moved to the bottom of the sidebar
- contents and data health moved into the main dashboard body
- issuer sector override saves corrections to issuers.csv
            """
        )

issuer_bonds = bonds_df[bonds_df["issuer"] == selected_issuer].copy()
issuer_trades = market_df[market_df["issuer"] == selected_issuer].copy()

selected_sector = "Unknown"
if "sector" in market_df.columns:
    sector_values = issuer_trades["sector"].dropna().astype(str).unique().tolist()
    if sector_values:
        selected_sector = sector_values[0]
elif "sector" in issuer_master.columns:
    sector_values = issuer_master.loc[issuer_master["issuer"] == selected_issuer, "sector"].dropna().astype(str).unique().tolist()
    if sector_values:
        selected_sector = sector_values[0]

if not issuer_trades.empty and maturity_bucket != "All":
    issuer_trades = issuer_trades[issuer_trades["maturity_bucket"] == maturity_bucket].copy()

if not issuer_trades.empty and time_window != "All":
    latest_date = issuer_trades["trade_date"].max()
    years = {"1Y": 1, "3Y": 3, "5Y": 5}[time_window]
    issuer_trades = issuer_trades[issuer_trades["trade_date"] >= latest_date - pd.DateOffset(years=years)].copy()


section_anchor("data-quality-scorecard", "Data Quality Scorecard")
with st.expander("Methodology: data quality scorecard", expanded=False):
    st.markdown(
        """
This section evaluates whether the uploaded data is reliable enough for secondary-market analytics.

**Scorecard components:**

- **CUSIP match rate**: share of merged trade rows whose CUSIP appears in the uploaded bond master.
- **Valid trade date rate**: share of rows with parseable trade dates.
- **Known maturity bucket rate**: share of rows assigned to Short / 10Y / 20Y / 30Y.
- **Positive trade amount rate**: share of rows with positive par/trade amount.
- **Issuer coverage rate**: share of rows with an issuer after merge / mapping.
- **Duplicate rows removed**: exact duplicate standardized trade rows removed before analytics.

The score is a practical reliability indicator, not a guarantee of correctness.
        """
    )

dq_rows = []
dq_total = len(market_df)

def dq_pct(numer, denom):
    return (numer / denom * 100) if denom and denom > 0 else 0

if dq_total > 0:
    bond_cusips = set(bonds_df["cusip"].dropna().astype(str).str.upper()) if "cusip" in bonds_df.columns else set()
    trade_cusips = market_df["cusip"].dropna().astype(str).str.upper() if "cusip" in market_df.columns else pd.Series(dtype=str)
    cusip_match_rate = dq_pct(trade_cusips.isin(bond_cusips).sum(), dq_total) if bond_cusips else 0

    valid_trade_dates = pd.to_datetime(market_df["trade_date"], errors="coerce").notna().sum() if "trade_date" in market_df.columns else 0
    valid_trade_date_rate = dq_pct(valid_trade_dates, dq_total)

    valid_buckets = MATURITY_BUCKET_ORDER
    known_bucket_rate = dq_pct(market_df["maturity_bucket"].isin(valid_buckets).sum(), dq_total) if "maturity_bucket" in market_df.columns else 0

    positive_amount_rate = (
        dq_pct((pd.to_numeric(market_df["trade_amount"], errors="coerce") > 0).sum(), dq_total)
        if "trade_amount" in market_df.columns else 0
    )

    issuer_coverage_rate = dq_pct(market_df["issuer"].notna().sum(), dq_total) if "issuer" in market_df.columns else 0

    data_quality_score = (
        0.30 * cusip_match_rate
        + 0.20 * valid_trade_date_rate
        + 0.20 * known_bucket_rate
        + 0.15 * issuer_coverage_rate
        + 0.15 * positive_amount_rate
    )

    dq_rows = [
        {"Metric": "CUSIP Match Rate", "Value": cusip_match_rate, "Weight": "30%", "Interpretation": "Trade CUSIPs found in bond master"},
        {"Metric": "Valid Trade Date Rate", "Value": valid_trade_date_rate, "Weight": "20%", "Interpretation": "Rows with parseable trade dates"},
        {"Metric": "Known Maturity Bucket Rate", "Value": known_bucket_rate, "Weight": "20%", "Interpretation": "Rows mapped to Short / Intermediate / Long / Extended Long"},
        {"Metric": "Issuer Coverage Rate", "Value": issuer_coverage_rate, "Weight": "15%", "Interpretation": "Rows with issuer after merge/mapping"},
        {"Metric": "Positive Trade Amount Rate", "Value": positive_amount_rate, "Weight": "15%", "Interpretation": "Rows with positive trade amount"},
    ]

    dq1, dq2, dq3, dq4 = st.columns(4)
    dq1.metric("Data Quality Score", f"{data_quality_score:.1f}/100")
    dq2.metric("CUSIP Match Rate", f"{cusip_match_rate:.1f}%")
    dq3.metric("Known Bucket Rate", f"{known_bucket_rate:.1f}%")
    dq4.metric("Duplicates Removed", f"{duplicates_removed:,}")

    if data_quality_score >= 90:
        st.success("Data quality looks strong for dashboard-level analytics.")
    elif data_quality_score >= 75:
        st.warning("Data quality is usable, but some analytics may be affected by missing fields.")
    else:
        st.error("Data quality is weak. Review missing CUSIPs, maturity dates, issuer mapping, and trade amount fields before relying on analytics.")

    dq_display = pd.DataFrame(dq_rows)
    dq_display["Value"] = dq_display["Value"].map(lambda x: f"{x:.1f}%")
    st.dataframe(dq_display, use_container_width=True, hide_index=True)

    with st.expander("Data quality issue drilldown", expanded=False):
        issue_cols = []
        issue_df = market_df.copy()
        if "cusip" in issue_df.columns:
            issue_df["cusip_matches_bond_master"] = issue_df["cusip"].astype(str).str.upper().isin(bond_cusips)
            issue_cols.append("cusip_matches_bond_master")
        if "maturity_bucket" in issue_df.columns:
            issue_df["known_maturity_bucket"] = issue_df["maturity_bucket"].isin(valid_buckets)
            issue_cols.append("known_maturity_bucket")
        if "trade_amount" in issue_df.columns:
            issue_df["positive_trade_amount"] = pd.to_numeric(issue_df["trade_amount"], errors="coerce") > 0
            issue_cols.append("positive_trade_amount")
        if "issuer" in issue_df.columns:
            issue_df["issuer_present"] = issue_df["issuer"].notna()
            issue_cols.append("issuer_present")

        display_issue_cols = [
            c for c in ["issuer", "cusip", "trade_date", "maturity_bucket", "yield", "trade_amount"] + issue_cols
            if c in issue_df.columns
        ]
        st.dataframe(issue_df[display_issue_cols].head(5000), use_container_width=True, hide_index=True)
else:
    st.info("Data quality scorecard will appear after market data is processed.")


section_anchor("executive-snapshot", "Executive Snapshot")

latest_trade_display = (
    issuer_trades["trade_date"].max().strftime("%Y-%m-%d")
    if not issuer_trades.empty
    else "No trades"
)

# Custom cards give long sector/issuer names enough horizontal room, while keeping numeric fields quieter.
snap_col1, snap_col2, snap_col3, snap_col4, snap_col5 = st.columns([1.55, 2.15, 0.75, 0.9, 1.1])
with snap_col1:
    clean_metric_card("Sector", selected_sector, size="large")
with snap_col2:
    clean_metric_card("Issuer", selected_issuer, size="large")
with snap_col3:
    clean_metric_card("Bonds", f"{len(issuer_bonds):,}", size="small")
with snap_col4:
    clean_metric_card("Trades", f"{len(issuer_trades):,}", size="small")
with snap_col5:
    clean_metric_card("Latest Trade", latest_trade_display, size="small")

section_anchor("yield-relative-value", "Yield Trend / Relative Value Comparison")
with st.expander("Methodology: benchmark curve framework", expanded=False):
    st.markdown(
        """
This section groups uploaded trade rows by **trade date** and **issuer**, then plots average observed trade yield.

**Benchmark logic:**

- **AAA Curve = uploaded MMD / AAA curve.**
- **If users upload explicit AA+/AA/AA-/A+/A/A-/BBB curve columns, the app uses those directly.**
- **If explicit non-AAA curves are missing, the app falls back to MMD + transparent rating-spread assumptions.**
- Spread assumptions are **maturity-adjusted**. For example, the 30Y AA spread can be wider than the 5Y AA spread.
- Units in the code are percentage points: `0.10 = 10 bps`.
- This is an internal analytical benchmark, not a live Bloomberg/BVAL/ICE curve. Replace the assumptions with firm-approved or vendor curves when available.
        """
    )
    st.dataframe(rating_spread_table(), use_container_width=True, hide_index=True)

issuer_choices = uploaded_issuers
default_compare = [selected_issuer] if selected_issuer in issuer_choices else issuer_choices[:1]
compare_issuers = st.multiselect("Compare Issuers", issuer_choices, default=default_compare)
compare_bucket = st.selectbox("Comparison Maturity Bucket", MATURITY_BUCKET_OPTIONS, key="compare_bucket")
benchmark_ratings = st.multiselect(
    "Benchmark Curve(s)",
    BENCHMARK_RATINGS,
    default=["AAA", "AA"],
    help="Priority: uploaded rating curve columns first; otherwise MMD/AAA plus the visible rating-spread assumptions above.",
)
show_spread_to_benchmark = st.checkbox(
    "Show issuer spread to selected benchmark",
    value=True,
    help="Calculates average issuer yield minus selected benchmark curve for dates where both are available.",
)

chart_df = market_df[market_df["issuer"].isin(compare_issuers)].copy()
if compare_bucket != "All":
    chart_df = chart_df[chart_df["maturity_bucket"] == compare_bucket].copy()

if chart_df.empty:
    st.warning("No trade data found for selected comparison filters.")
else:
    date_min = chart_df["trade_date"].min().date()
    date_max = chart_df["trade_date"].max().date()
    selected_dates = st.date_input("Trade Date Range", value=(date_min, date_max), min_value=date_min, max_value=date_max)
    if isinstance(selected_dates, tuple) and len(selected_dates) == 2:
        start_date, end_date = selected_dates
        chart_df = chart_df[(chart_df["trade_date"].dt.date >= start_date) & (chart_df["trade_date"].dt.date <= end_date)].copy()

    daily = (
        chart_df.groupby(["trade_date", "issuer"], as_index=False)
        .agg(avg_yield=("yield", "mean"), trade_count=("yield", "count"), total_trade_amount=("trade_amount", "sum"))
    )
    fig = px.line(
        daily.sort_values("trade_date"),
        x="trade_date",
        y="avg_yield",
        color="issuer",
        markers=True,
        hover_data=["trade_count", "total_trade_amount"],
        title="Average Trade Yield by Issuer",
    )

    benchmark_daily = pd.DataFrame()
    benchmark_ready = False
    if not mmd_df.empty and benchmark_ratings:
        date_col = _detect_mmd_date_column(mmd_df)
        mmd_col = MMD_BUCKET_MAP.get(compare_bucket, "10Y")
        if date_col:
            mmd_plot = mmd_df.copy()
            mmd_plot[date_col] = pd.to_datetime(mmd_plot[date_col], errors="coerce")
            mmd_plot = mmd_plot.dropna(subset=[date_col])
            if isinstance(selected_dates, tuple) and len(selected_dates) == 2:
                mmd_plot = mmd_plot[(mmd_plot[date_col].dt.date >= start_date) & (mmd_plot[date_col].dt.date <= end_date)]

            benchmark_frames = []
            unavailable_ratings = []
            for rating in benchmark_ratings:
                y, meta = get_benchmark_curve(mmd_plot, mmd_col, rating)
                if y is None:
                    unavailable_ratings.append(rating)
                    continue
                fig.add_scatter(
                    x=mmd_plot[date_col],
                    y=y,
                    mode="lines",
                    name=f"{rating} Curve ({mmd_col})",
                )
                benchmark_frames.append(
                    pd.DataFrame({
                        "trade_date": mmd_plot[date_col].dt.normalize(),
                        "benchmark_rating": rating,
                        "benchmark_yield": y,
                        "mmd_tenor": mmd_col,
                        "rating_spread_bps": meta.get("rating_spread_bps"),
                        "benchmark_source": meta.get("benchmark_source"),
                        "source_column": meta.get("source_column"),
                    })
                )
            benchmark_daily = pd.concat(benchmark_frames, ignore_index=True) if benchmark_frames else pd.DataFrame()
            benchmark_ready = not benchmark_daily.empty
            if unavailable_ratings:
                st.warning(
                    "Some benchmark curves could not be built because neither an uploaded curve column nor a usable AAA/MMD base tenor was found: "
                    + ", ".join(unavailable_ratings)
                )
        else:
            st.warning("Benchmark curves could not be plotted because the curve file does not contain a usable date column.")

    fig.update_layout(xaxis_title="Trade Date", yaxis_title="Yield (%)", hovermode="x unified")
    st.plotly_chart(fig, use_container_width=True)

    if show_spread_to_benchmark and benchmark_ready and not daily.empty:
        spread_base = daily.copy()
        spread_base["trade_date"] = pd.to_datetime(spread_base["trade_date"], errors="coerce").dt.normalize()
        spread_to_benchmark = spread_base.merge(benchmark_daily, on="trade_date", how="inner")
        if spread_to_benchmark.empty:
            st.info("No overlapping dates were found between issuer trades and the selected benchmark curve.")
        else:
            spread_to_benchmark["spread_to_benchmark_bps"] = (
                spread_to_benchmark["avg_yield"] - spread_to_benchmark["benchmark_yield"]
            ) * 100
            spread_fig = px.line(
                spread_to_benchmark.sort_values("trade_date"),
                x="trade_date",
                y="spread_to_benchmark_bps",
                color="issuer",
                line_dash="benchmark_rating",
                markers=True,
                hover_data=["benchmark_rating", "mmd_tenor", "benchmark_source", "source_column", "rating_spread_bps", "trade_count", "total_trade_amount"],
                title="Issuer Spread to Selected Benchmark Curve(s)",
            )
            spread_fig.update_layout(xaxis_title="Trade Date", yaxis_title="Spread to Benchmark (bps)", hovermode="x unified")
            st.plotly_chart(spread_fig, use_container_width=True)

            with st.expander("Spread-to-benchmark calculation details", expanded=False):
                st.markdown(
                    """
For each issuer/date/rating benchmark:

`Spread to Benchmark (bps) = (Average Issuer Trade Yield - Synthetic Benchmark Yield) × 100`

Where:

`Benchmark Yield = uploaded rating curve if available; otherwise MMD/AAA Tenor Yield + Rating Spread Assumption`
                    """
                )
                st.dataframe(
                    spread_to_benchmark[[
                        "trade_date", "issuer", "benchmark_rating", "mmd_tenor", "avg_yield",
                        "benchmark_yield", "benchmark_source", "source_column", "rating_spread_bps", "spread_to_benchmark_bps",
                        "trade_count", "total_trade_amount",
                    ]].sort_values(["trade_date", "issuer", "benchmark_rating"], ascending=[False, True, True]).head(1000),
                    use_container_width=True,
                    hide_index=True,
                )
    elif show_spread_to_benchmark and mmd_df.empty:
        st.info("Upload an MMD curve file to enable AAA/AA/A/BBB benchmark curves and spread-to-benchmark analytics.")


section_anchor("issuer-curve", "Issuer Curve vs Benchmark Curve")
with st.expander("Methodology: issuer curve vs benchmark curve", expanded=False):
    st.markdown(
        """
This chart shows a **cross-sectional yield curve** by maturity bucket, rather than a time-series trend.

**Issuer curve logic:**

- The issuer curve is built from uploaded trade yields by maturity bucket: **Short / 10Y / 20Y / 30Y**.
- Default aggregation uses **average yield over the latest selected window** ending on the curve date. This reduces noise from sparse municipal trading.
- You can also use **latest trade per bucket** when you want the most recent observation in each maturity bucket.

**Benchmark curve logic:**

- The benchmark curve uses uploaded rating curve columns when available.
- If an uploaded AA/A/BBB curve is missing, the app falls back to **MMD/AAA + transparent rating-spread assumptions**.
- The benchmark value uses the latest available curve observation at or before the selected curve date.

**How to read it:**

- If the issuer line is above the benchmark line, the issuer trades cheaper / wider for that bucket.
- If the issuer line is below the benchmark line, the issuer trades richer / tighter for that bucket.
- The accompanying table shows exact yields and spreads in basis points.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD / benchmark curve file to enable Issuer Curve vs Benchmark Curve analysis.")
else:
    selected_issuer_dates = pd.to_datetime(
        market_df.loc[market_df["issuer"] == selected_issuer, "trade_date"], errors="coerce"
    ).dropna()

    if selected_issuer_dates.empty:
        st.warning("No valid trade dates were found for the selected issuer, so the issuer curve cannot be built.")
    else:
        curve_min_date = selected_issuer_dates.min().date()
        curve_max_date = selected_issuer_dates.max().date()

        curve_ctrl1, curve_ctrl2, curve_ctrl3 = st.columns([1, 1, 1.4])
        with curve_ctrl1:
            curve_as_of_date = st.date_input(
                "Curve Date",
                value=curve_max_date,
                min_value=curve_min_date,
                max_value=curve_max_date,
                key="issuer_curve_as_of_date",
                help="The issuer and benchmark curves use observations available at or before this date.",
            )
        with curve_ctrl2:
            curve_aggregation = st.selectbox(
                "Issuer Curve Aggregation",
                ["Average last N days", "Latest trade per bucket"],
                index=0,
                key="issuer_curve_aggregation",
            )
        with curve_ctrl3:
            curve_benchmark_ratings = st.multiselect(
                "Benchmark Curve(s) for Curve Chart",
                BENCHMARK_RATINGS,
                default=[r for r in ["AAA", "AA"] if r in BENCHMARK_RATINGS],
                key="issuer_curve_benchmark_ratings",
                help="Priority: uploaded rating curve columns first; otherwise MMD/AAA plus the visible spread assumptions.",
            )

        curve_lookback_days = 30
        if curve_aggregation == "Average last N days":
            curve_lookback_days = st.select_slider(
                "Lookback Window for Issuer Curve",
                options=[7, 14, 30, 60, 90, 180],
                value=30,
                format_func=lambda x: f"{x} days",
                key="issuer_curve_lookback_days",
                help="Municipal trades can be sparse, so averaging over a window usually gives a more stable curve than using one day only.",
            )

        if not curve_benchmark_ratings:
            st.info("Select at least one benchmark curve to display the issuer curve comparison.")
        else:
            issuer_curve_plot_df, issuer_curve_audit = build_issuer_curve_snapshot(
                market_df=market_df,
                mmd_df=mmd_df,
                issuer=selected_issuer,
                ratings=curve_benchmark_ratings,
                as_of_date=pd.Timestamp(curve_as_of_date),
                lookback_days=curve_lookback_days,
                aggregation_method=curve_aggregation,
            )

            if issuer_curve_plot_df.empty or issuer_curve_audit.empty:
                st.warning(
                    "No overlapping issuer trades and benchmark curve observations were found for this curve setup. "
                    "Try a longer lookback window, a different curve date, or check that the benchmark file has usable 5Y/10Y/20Y/30Y columns."
                )
            else:
                curve_fig = px.line(
                    issuer_curve_plot_df,
                    x="maturity_bucket",
                    y="yield_value",
                    color="curve",
                    markers=True,
                    hover_data=["curve_type", "trade_count", "issuer_observation_date"],
                    title=f"{selected_issuer} Issuer Curve vs Benchmark Curve(s)",
                    labels={
                        "maturity_bucket": "Maturity Bucket",
                        "yield_value": "Yield (%)",
                        "curve": "Curve",
                    },
                )
                curve_fig.update_layout(hovermode="x unified", height=500)
                st.plotly_chart(curve_fig, use_container_width=True)

                table_cols = [
                    "maturity_bucket", "benchmark_rating", "issuer_yield", "benchmark_yield",
                    "spread_to_benchmark_bps", "trade_count", "issuer_observation_date", "benchmark_date",
                    "mmd_tenor", "benchmark_source", "source_column", "rating_spread_bps",
                    "aggregation_method", "lookback_start", "lookback_end",
                ]
                curve_table = issuer_curve_audit[[c for c in table_cols if c in issuer_curve_audit.columns]].copy()
                for c in ["issuer_yield", "benchmark_yield", "spread_to_benchmark_bps", "rating_spread_bps"]:
                    if c in curve_table.columns:
                        curve_table[c] = pd.to_numeric(curve_table[c], errors="coerce").round(2)

                st.subheader("Curve Spread Table")
                st.dataframe(curve_table, use_container_width=True, hide_index=True)

                primary_curve_rating = curve_benchmark_ratings[0]
                primary_rows = issuer_curve_audit[issuer_curve_audit["benchmark_rating"] == primary_curve_rating].copy()
                primary_rows = primary_rows.dropna(subset=["spread_to_benchmark_bps"])
                if not primary_rows.empty:
                    cheap_row = primary_rows.loc[primary_rows["spread_to_benchmark_bps"].idxmax()]
                    rich_row = primary_rows.loc[primary_rows["spread_to_benchmark_bps"].idxmin()]
                    st.info(
                        f"Curve read-through vs {primary_curve_rating}: "
                        f"{cheap_row['maturity_bucket']} is the widest bucket at {cheap_row['spread_to_benchmark_bps']:+.1f} bp, "
                        f"while {rich_row['maturity_bucket']} is the tightest bucket at {rich_row['spread_to_benchmark_bps']:+.1f} bp."
                    )



section_anchor("curve-shape", "Curve Shape Analytics")
with st.expander("Methodology: curve shape analytics", expanded=False):
    st.markdown(
        """
This section turns the issuer curve into **curve mathematics**, similar to what rates / muni desks monitor.

**Metrics:**

- **5s10s Slope** = 10Y yield − Short/5Y yield
- **10s30s Slope** = 30Y yield − 10Y yield
- **5s30s Slope** = 30Y yield − Short/5Y yield
- **5s10s30s Butterfly** = 10Y yield − average(Short/5Y yield, 30Y yield)

**How to read it:**

- Higher positive slope = steeper curve.
- Lower or negative slope = flatter / inverted curve shape.
- Positive butterfly = 10Y “belly” is high/cheap versus wings.
- Negative butterfly = 10Y “belly” is low/rich versus wings.

The issuer curve uses uploaded trade data over the selected lookback window. The benchmark curve uses uploaded rating curves when available; otherwise it falls back to MMD/AAA plus transparent rating-spread assumptions.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable curve shape analytics.")
else:
    cs_col1, cs_col2, cs_col3 = st.columns([1, 1, 1])
    with cs_col1:
        cs_rating = st.selectbox(
            "Curve Shape Benchmark",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="curve_shape_benchmark",
        )
    with cs_col2:
        cs_lookback = st.selectbox(
            "Issuer Curve Lookback",
            [7, 30, 60, 90, 180],
            index=1,
            format_func=lambda x: f"Latest {x} days",
            key="curve_shape_lookback",
        )
    with cs_col3:
        cs_curve_basis = st.selectbox(
            "Curve Basis",
            ["Yield Curve", "Spread Curve"],
            index=0,
            key="curve_shape_basis",
            help="Yield Curve uses issuer yields. Spread Curve uses issuer spread to selected benchmark.",
        )

    cs_base = market_df[market_df["issuer"] == selected_issuer].copy()
    cs_base["trade_date"] = pd.to_datetime(cs_base["trade_date"], errors="coerce").dt.normalize()
    cs_base["yield"] = pd.to_numeric(cs_base["yield"], errors="coerce")
    cs_base = cs_base.dropna(subset=["trade_date", "yield", "maturity_bucket"])
    cs_base = cs_base[cs_base["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)].copy()

    if cs_base.empty:
        st.warning("No usable issuer trade rows were available for curve shape analytics.")
    else:
        cs_latest_date = cs_base["trade_date"].max()
        cs_start_date = cs_latest_date - pd.Timedelta(days=int(cs_lookback))
        cs_window = cs_base[cs_base["trade_date"] >= cs_start_date].copy()

        if cs_window.empty:
            st.warning("No issuer trades were found inside the selected lookback window.")
        else:
            issuer_curve = (
                cs_window.groupby("maturity_bucket", as_index=False)
                .agg(
                    issuer_yield=("yield", "mean"),
                    trade_count=("yield", "count"),
                    total_trade_amount=("trade_amount", "sum") if "trade_amount" in cs_window.columns else ("yield", "count"),
                    latest_trade=("trade_date", "max"),
                )
            )

            date_col = _detect_mmd_date_column(mmd_df)
            if date_col is None:
                st.warning("Curve shape analytics cannot build benchmark curve because the curve file has no usable date column.")
            else:
                cs_mmd = mmd_df.copy()
                cs_mmd[date_col] = pd.to_datetime(cs_mmd[date_col], errors="coerce")
                cs_mmd = cs_mmd.dropna(subset=[date_col])
                cs_mmd = cs_mmd[cs_mmd[date_col].dt.normalize() <= cs_latest_date].sort_values(date_col)

                if cs_mmd.empty:
                    st.warning("No benchmark curve observation was available on or before the latest issuer trade date.")
                else:
                    cs_latest_mmd = cs_mmd.iloc[[-1]].copy()
                    cs_benchmark_date = cs_latest_mmd[date_col].iloc[0]

                    bench_rows = []
                    for bucket in MATURITY_BUCKET_ORDER:
                        tenor = MMD_BUCKET_MAP.get(bucket, "10Y")
                        y, meta = get_benchmark_curve(cs_latest_mmd, tenor, cs_rating)
                        if y is not None and pd.notna(y.iloc[0]):
                            bench_rows.append(
                                {
                                    "maturity_bucket": bucket,
                                    "mmd_tenor": tenor,
                                    "benchmark_yield": float(y.iloc[0]),
                                    "benchmark_source": meta.get("benchmark_source"),
                                    "source_column": meta.get("source_column"),
                                    "rating_spread_bps": meta.get("rating_spread_bps"),
                                }
                            )
                    bench_curve = pd.DataFrame(bench_rows)

                    if bench_curve.empty:
                        st.warning("Selected benchmark curve could not be built for curve shape analytics.")
                    else:
                        curve_shape_df = issuer_curve.merge(bench_curve, on="maturity_bucket", how="outer")
                        curve_shape_df["spread_to_benchmark_bps"] = (
                            curve_shape_df["issuer_yield"] - curve_shape_df["benchmark_yield"]
                        ) * 100

                        maturity_order = MATURITY_BUCKET_ORDER
                        curve_shape_df["maturity_bucket"] = pd.Categorical(
                            curve_shape_df["maturity_bucket"],
                            categories=maturity_order,
                            ordered=True,
                        )
                        curve_shape_df = curve_shape_df.sort_values("maturity_bucket")

                        metric_col = "issuer_yield" if cs_curve_basis == "Yield Curve" else "spread_to_benchmark_bps"
                        metric_label = "Issuer Yield (%)" if cs_curve_basis == "Yield Curve" else f"Spread to {cs_rating} (bps)"

                        curve_plot = curve_shape_df.dropna(subset=[metric_col]).copy()
                        if curve_plot.empty:
                            st.warning("Not enough curve points were available to calculate curve shape metrics.")
                        else:
                            fig_curve_shape = px.line(
                                curve_plot,
                                x="maturity_bucket",
                                y=metric_col,
                                markers=True,
                                hover_data=[
                                    c for c in [
                                        "issuer_yield", "benchmark_yield", "spread_to_benchmark_bps",
                                        "trade_count", "total_trade_amount", "latest_trade",
                                        "benchmark_source", "source_column"
                                    ] if c in curve_plot.columns
                                ],
                                title=f"{selected_issuer} {cs_curve_basis} Shape",
                                labels={
                                    "maturity_bucket": "Maturity Bucket",
                                    metric_col: metric_label,
                                },
                            )
                            fig_curve_shape.update_layout(height=450, hovermode="x unified")
                            st.plotly_chart(fig_curve_shape, use_container_width=True)

                            curve_values = (
                                curve_shape_df.set_index("maturity_bucket")[metric_col]
                                .astype(float)
                                .to_dict()
                            )

                            def get_curve_value(bucket: str):
                                value = curve_values.get(bucket)
                                return value if pd.notna(value) else pd.NA

                            v_short = get_curve_value("Short")
                            v_10 = get_curve_value("10Y")
                            v_20 = get_curve_value("20Y")
                            v_30 = get_curve_value("30Y")

                            # -------------------------
                            # Dynamic curve diagnostics
                            # -------------------------
                            available_points = []
                            missing_points = []

                            for bucket_name, bucket_value in {
                                "Short": v_short,
                                "10Y": v_10,
                                "20Y": v_20,
                                "30Y": v_30,
                            }.items():
                                if pd.notna(bucket_value):
                                    available_points.append(bucket_name)
                                else:
                                    missing_points.append(bucket_name)

                            diag_col1, diag_col2 = st.columns(2)

                            with diag_col1:
                                st.success(
                                    "Available Curve Points:\n\n"
                                    + ", ".join(available_points)
                                    if available_points
                                    else "No usable curve points detected."
                                )

                            with diag_col2:
                                if missing_points:
                                    st.warning(
                                        "Missing Curve Points:\n\n"
                                        + ", ".join(missing_points)
                                    )
                                else:
                                    st.success("All core curve points detected.")

                            metrics_rows = []
                            analytics_status = []

                            # 5s10s
                            if pd.notna(v_short) and pd.notna(v_10):
                                metrics_rows.append({
                                    "Metric": "5s10s Slope",
                                    "Value": v_10 - v_short,
                                })
                                analytics_status.append({
                                    "Analytics": "5s10s Slope",
                                    "Status": "Available",
                                    "Requirement": "Short + 10Y",
                                })
                            else:
                                analytics_status.append({
                                    "Analytics": "5s10s Slope",
                                    "Status": "Missing Required Points",
                                    "Requirement": "Short + 10Y",
                                })

                            # 10s30s
                            if pd.notna(v_10) and pd.notna(v_30):
                                metrics_rows.append({
                                    "Metric": "10s30s Slope",
                                    "Value": v_30 - v_10,
                                })
                                analytics_status.append({
                                    "Analytics": "10s30s Slope",
                                    "Status": "Available",
                                    "Requirement": "10Y + 30Y",
                                })
                            else:
                                analytics_status.append({
                                    "Analytics": "10s30s Slope",
                                    "Status": "Missing Required Points",
                                    "Requirement": "10Y + 30Y",
                                })

                            # 5s30s
                            if pd.notna(v_short) and pd.notna(v_30):
                                metrics_rows.append({
                                    "Metric": "5s30s Slope",
                                    "Value": v_30 - v_short,
                                })
                                analytics_status.append({
                                    "Analytics": "5s30s Slope",
                                    "Status": "Available",
                                    "Requirement": "Short + 30Y",
                                })
                            else:
                                analytics_status.append({
                                    "Analytics": "5s30s Slope",
                                    "Status": "Missing Required Points",
                                    "Requirement": "Short + 30Y",
                                })

                            # Butterfly
                            if pd.notna(v_short) and pd.notna(v_10) and pd.notna(v_30):
                                metrics_rows.append({
                                    "Metric": "5s10s30s Butterfly",
                                    "Value": v_10 - ((v_short + v_30) / 2),
                                })
                                analytics_status.append({
                                    "Analytics": "5s10s30s Butterfly",
                                    "Status": "Available",
                                    "Requirement": "Short + 10Y + 30Y",
                                })
                            else:
                                analytics_status.append({
                                    "Analytics": "5s10s30s Butterfly",
                                    "Status": "Missing Required Points",
                                    "Requirement": "Short + 10Y + 30Y",
                                })

                            # Steepness score
                            if pd.notna(v_short) and pd.notna(v_10) and pd.notna(v_20) and pd.notna(v_30):
                                metrics_rows.append({
                                    "Metric": "Steepness Score",
                                    "Value": (
                                        ((v_10 - v_short)
                                        + (v_30 - v_10)
                                        + (v_30 - v_short)) / 3
                                    ),
                                })
                                analytics_status.append({
                                    "Analytics": "Steepness Score",
                                    "Status": "Available",
                                    "Requirement": "Short + 10Y + 20Y + 30Y",
                                })
                            else:
                                analytics_status.append({
                                    "Analytics": "Steepness Score",
                                    "Status": "Missing Required Points",
                                    "Requirement": "Short + 10Y + 20Y + 30Y",
                                })

                            analytics_status_df = pd.DataFrame(analytics_status)

                            with st.expander("Curve Analytics Availability", expanded=False):
                                st.dataframe(
                                    analytics_status_df,
                                    use_container_width=True,
                                    hide_index=True,
                                )

                            if not metrics_rows:
                                st.info("At least two compatible curve points are needed to calculate curve shape metrics.")
                            else:
                                metrics_df = pd.DataFrame(metrics_rows)
                                unit = "bp" if cs_curve_basis == "Spread Curve" else "%"
                                metrics_df["Display"] = metrics_df["Value"].map(
                                    lambda x: f"{x:+.1f} {unit}" if cs_curve_basis == "Spread Curve" else f"{x:+.2f}%"
                                )

                                mcols = st.columns(min(4, len(metrics_df)))
                                for idx, (_, row) in enumerate(metrics_df.head(4).iterrows()):
                                    mcols[idx % len(mcols)].metric(row["Metric"], row["Display"])

                                st.dataframe(metrics_df, use_container_width=True, hide_index=True)

                                # Read-through
                                slope_1030 = metrics_df.loc[metrics_df["Metric"] == "10s30s Slope", "Value"]
                                butterfly = metrics_df.loc[metrics_df["Metric"] == "5s10s30s Butterfly", "Value"]

                                notes = []
                                if not slope_1030.empty:
                                    s_val = float(slope_1030.iloc[0])
                                    if s_val > (20 if cs_curve_basis == "Spread Curve" else 0.20):
                                        notes.append("long end screens steep versus the 10Y point")
                                    elif s_val < (-10 if cs_curve_basis == "Spread Curve" else -0.10):
                                        notes.append("long end screens flat/inverted versus the 10Y point")
                                    else:
                                        notes.append("10s30s slope appears relatively contained")

                                if not butterfly.empty:
                                    b_val = float(butterfly.iloc[0])
                                    if b_val > (10 if cs_curve_basis == "Spread Curve" else 0.10):
                                        notes.append("10Y belly appears cheap/high versus wings")
                                    elif b_val < (-10 if cs_curve_basis == "Spread Curve" else -0.10):
                                        notes.append("10Y belly appears rich/low versus wings")

                                if notes:
                                    st.info("Curve read-through: " + "; ".join(notes) + ".")

                                with st.expander("Curve shape audit table", expanded=False):
                                    audit_cols = [
                                        "maturity_bucket",
                                        "issuer_yield",
                                        "benchmark_yield",
                                        "spread_to_benchmark_bps",
                                        "trade_count",
                                        "total_trade_amount",
                                        "latest_trade",
                                        "mmd_tenor",
                                        "benchmark_source",
                                        "source_column",
                                        "rating_spread_bps",
                                    ]
                                    audit_curve = curve_shape_df[[c for c in audit_cols if c in curve_shape_df.columns]].copy()
                                    for c in ["issuer_yield", "benchmark_yield", "spread_to_benchmark_bps", "rating_spread_bps"]:
                                        if c in audit_curve.columns:
                                            audit_curve[c] = pd.to_numeric(audit_curve[c], errors="coerce").round(2)
                                    st.caption(
                                        f"Issuer lookback: latest {cs_lookback} days. "
                                        f"Benchmark date: {cs_benchmark_date.strftime('%Y-%m-%d')}."
                                    )
                                    st.dataframe(audit_curve, use_container_width=True, hide_index=True)


section_anchor("spread-level", "Current Spread Level Framework")
with st.expander("Methodology: current spread level", expanded=False):
    st.markdown(
        """
This section shows where the selected issuer is trading **now** versus transparent benchmark curves.

**Calculation:**

`Current Spread Level = (Average Issuer Trade Yield - Benchmark Yield) × 100`

Where:

`Benchmark Yield = uploaded rating curve if available; otherwise MMD/AAA Curve + Rating Spread Assumption`

**How to read it:**

- **Positive spread**: issuer yield is above the selected benchmark curve; the issuer/bucket is trading cheaper than that benchmark.
- **Negative spread**: issuer yield is below the selected benchmark curve; the issuer/bucket is trading richer than that benchmark.
- Rows are maturity buckets. Columns are benchmark curves.
- This is a **level** view, not a movement view. Level answers: *is it cheap or rich right now?* Movement answers: *did it widen or tighten recently?*
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD curve file to enable current spread level analytics.")
else:
    level_col1, level_col2 = st.columns([1, 2])
    with level_col1:
        level_ratings = st.multiselect(
            "Spread Level Benchmark Curves",
            BENCHMARK_RATINGS,
            default=[r for r in ["AAA", "AA", "A", "BBB"] if r in BENCHMARK_RATINGS],
            help="Priority: uploaded rating curve columns first; otherwise MMD/AAA plus the visible rating-spread assumptions.",
        )
    with level_col2:
        st.caption(
            "Cells show latest available issuer spread to each benchmark curve, in basis points. "
            "Higher positive values generally indicate cheaper relative value versus that benchmark."
        )

    if not level_ratings:
        st.info("Select at least one benchmark curve to display current spread levels.")
    else:
        level_matrix, level_audit = build_spread_level_data(
            market_df=market_df,
            mmd_df=mmd_df,
            issuer=selected_issuer,
            ratings=level_ratings,
        )
        if level_matrix.isna().all().all():
            st.warning(
                "No overlapping issuer trade dates and benchmark dates were found for current spread levels. "
                "Check that the curve file has a Date column plus either 5Y/10Y/20Y/30Y base columns or explicit rating curve columns such as AA_10Y, and that trade dates overlap with the curve history."
            )
        else:
            level_text = level_matrix.map(lambda x: "" if pd.isna(x) else f"{x:+.1f} bp")

            # 1) Spread level curve: one line per selected benchmark rating.
            curve_df = level_matrix.reset_index().rename(columns={"index": "maturity_bucket"})
            curve_long = curve_df.melt(
                id_vars="maturity_bucket",
                var_name="benchmark_rating",
                value_name="spread_to_benchmark_bps",
            ).dropna(subset=["spread_to_benchmark_bps"])
            curve_long["maturity_bucket"] = pd.Categorical(
                curve_long["maturity_bucket"],
                categories=MATURITY_BUCKET_ORDER,
                ordered=True,
            )
            curve_long = curve_long.sort_values(["benchmark_rating", "maturity_bucket"])

            st.subheader("1. Current Spread Curve")
            level_curve_fig = px.line(
                curve_long,
                x="maturity_bucket",
                y="spread_to_benchmark_bps",
                color="benchmark_rating",
                markers=True,
                title=f"{selected_issuer} Current Spread Curve vs Selected Benchmarks",
                labels={
                    "maturity_bucket": "Maturity Bucket",
                    "spread_to_benchmark_bps": "Spread to Benchmark (bps)",
                    "benchmark_rating": "Benchmark Curve",
                },
            )
            level_curve_fig.add_hline(y=0, line_dash="dash", opacity=0.5)
            level_curve_fig.update_layout(hovermode="x unified")
            st.plotly_chart(level_curve_fig, use_container_width=True)

            # 2) Spread level heatmap: maturity bucket x benchmark rating.
            st.subheader("2. Current Spread Level Heatmap")
            level_heatmap_fig = px.imshow(
                level_matrix.astype(float),
                x=level_matrix.columns,
                y=level_matrix.index,
                color_continuous_scale=["#1a9850", "#f7f7f7", "#d73027"],
                color_continuous_midpoint=0,
                aspect="auto",
                title=f"{selected_issuer} Current Spread Level vs Benchmark Curves",
                labels={"x": "Benchmark Curve", "y": "Maturity Bucket", "color": "Current Spread (bps)"},
            )
            level_heatmap_fig.update_traces(
                text=level_text.values,
                texttemplate="%{text}",
                hovertemplate="Maturity=%{y}<br>Benchmark=%{x}<br>Spread=%{z:.1f} bp<extra></extra>",
            )
            level_heatmap_fig.update_layout(height=420)
            st.plotly_chart(level_heatmap_fig, use_container_width=True)

            # 3) Quick signal: identify the cheapest bucket vs the first selected benchmark.
            primary_rating = level_ratings[0]
            if primary_rating in level_matrix.columns and level_matrix[primary_rating].notna().any():
                cheapest_bucket = level_matrix[primary_rating].astype(float).idxmax()
                cheapest_spread = level_matrix.loc[cheapest_bucket, primary_rating]
                richest_bucket = level_matrix[primary_rating].astype(float).idxmin()
                richest_spread = level_matrix.loc[richest_bucket, primary_rating]
                st.info(
                    f"Relative value read-through vs {primary_rating}: "
                    f"{cheapest_bucket} appears cheapest at {cheapest_spread:+.1f} bp, "
                    f"while {richest_bucket} appears richest at {richest_spread:+.1f} bp."
                )

            with st.expander("Current spread level audit table", expanded=False):
                display_cols = [
                    "maturity_bucket", "benchmark_rating", "latest_date", "avg_yield", "benchmark_yield",
                    "spread_to_benchmark_bps", "mmd_tenor", "benchmark_source", "source_column", "rating_spread_bps", "trade_count",
                    "total_trade_amount", "note",
                ]
                audit_display = level_audit[[c for c in display_cols if c in level_audit.columns]].copy()
                for c in ["avg_yield", "benchmark_yield", "spread_to_benchmark_bps", "rating_spread_bps"]:
                    if c in audit_display.columns:
                        audit_display[c] = pd.to_numeric(audit_display[c], errors="coerce").round(2)
                st.dataframe(audit_display, use_container_width=True, hide_index=True)


section_anchor("spread-attribution", "Spread Attribution Waterfall")
with st.expander("Methodology: spread attribution waterfall", expanded=False):
    st.markdown(
        """
This section decomposes the selected issuer's spread versus the **AAA/MMD curve** into transparent, reviewable components.

**Framework:**

`Issuer Spread vs AAA = Rating Premium + Liquidity Premium + Callable Adjustment + Residual / Issuer-Specific Premium`

**Important notes:**

- This is a **modeled attribution**, not a vendor curve or investment recommendation.
- **Rating Premium** uses the visible maturity-adjusted rating-spread assumptions already shown in the benchmark methodology.
- **Liquidity Premium** is estimated from the issuer/bucket liquidity score. Less liquid buckets receive a larger modeled premium.
- **Callable Adjustment** is a simple proxy based on whether bonds in the bucket appear callable.
- **Residual / Issuer-Specific Premium** captures what remains after the modeled components. This can reflect credit, sector, structure, supply/demand, data noise, or model misspecification.
- The purpose is pitchbook-style explanation and screening, not final pricing.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD curve file to enable spread attribution waterfall analytics.")
else:
    wf_col1, wf_col2, wf_col3 = st.columns([1, 1, 1])
    with wf_col1:
        wf_bucket = st.selectbox(
            "Waterfall Maturity Bucket",
            MATURITY_BUCKET_ORDER,
            index=1,
            help="The issuer spread will be attributed for this maturity bucket.",
        )
    with wf_col2:
        wf_rating = st.selectbox(
            "Modeled Rating Premium",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AA") if "AA" in BENCHMARK_RATINGS else 0,
            help="Used only for the rating-premium component. The total spread is measured versus AAA/MMD.",
        )
    with wf_col3:
        wf_lookback_days = st.selectbox(
            "Issuer Yield Lookback",
            [7, 30, 60, 90, 180],
            index=1,
            format_func=lambda x: f"Latest {x} days",
            help="Average issuer yield is calculated from trades in this lookback window to reduce muni trading noise.",
        )

    wf_tenor = MMD_BUCKET_MAP.get(wf_bucket, "10Y")
    date_col = _detect_mmd_date_column(mmd_df)

    if date_col is None:
        st.warning("Waterfall cannot run because the MMD/curve file does not contain a usable date column.")
    else:
        issuer_bucket_trades = market_df[
            (market_df["issuer"] == selected_issuer)
            & (market_df["maturity_bucket"] == wf_bucket)
        ].copy()

        if issuer_bucket_trades.empty:
            st.warning(f"No {wf_bucket} trade rows were found for {selected_issuer}.")
        else:
            issuer_bucket_trades["trade_date"] = pd.to_datetime(issuer_bucket_trades["trade_date"], errors="coerce")
            issuer_bucket_trades["yield"] = pd.to_numeric(issuer_bucket_trades["yield"], errors="coerce")
            issuer_bucket_trades = issuer_bucket_trades.dropna(subset=["trade_date", "yield"])

            if issuer_bucket_trades.empty:
                st.warning("Waterfall cannot run because no valid trade dates/yields remain after cleaning.")
            else:
                wf_latest_trade_date = issuer_bucket_trades["trade_date"].max().normalize()
                wf_start_date = wf_latest_trade_date - pd.Timedelta(days=int(wf_lookback_days))
                wf_window_trades = issuer_bucket_trades[issuer_bucket_trades["trade_date"] >= wf_start_date].copy()

                if wf_window_trades.empty:
                    st.warning("No trades were found inside the selected lookback window.")
                else:
                    wf_avg_issuer_yield = wf_window_trades["yield"].mean()
                    wf_trade_count = len(wf_window_trades)

                    if "trade_amount" in wf_window_trades.columns:
                        wf_total_trade_amount = pd.to_numeric(
                            wf_window_trades["trade_amount"], errors="coerce"
                        ).fillna(0).sum()
                    else:
                        wf_total_trade_amount = 0.0

                    wf_mmd = mmd_df.copy()
                    wf_mmd[date_col] = pd.to_datetime(wf_mmd[date_col], errors="coerce")
                    wf_mmd = wf_mmd.dropna(subset=[date_col])
                    wf_mmd = wf_mmd[wf_mmd[date_col].dt.normalize() <= wf_latest_trade_date].copy()

                    if wf_mmd.empty:
                        st.warning("No benchmark curve observations were available on or before the latest issuer trade date.")
                    else:
                        wf_mmd = wf_mmd.sort_values(date_col)
                        wf_latest_mmd = wf_mmd.iloc[[-1]].copy()
                        wf_benchmark_date = wf_latest_mmd[date_col].iloc[0]

                        wf_aaa_yield_series, wf_aaa_meta = get_benchmark_curve(wf_latest_mmd, wf_tenor, "AAA")
                        if wf_aaa_yield_series is None or pd.isna(wf_aaa_yield_series.iloc[0]):
                            st.warning(f"AAA/MMD {wf_tenor} curve could not be built for the waterfall.")
                        else:
                            wf_aaa_yield = float(wf_aaa_yield_series.iloc[0])
                            wf_total_spread_bps = (wf_avg_issuer_yield - wf_aaa_yield) * 100

                            wf_rating_premium_bps = (
                                RATING_SPREADS.get(wf_rating, RATING_SPREADS["AAA"]).get(wf_tenor, 0.00) * 100
                            )

                            wf_liq_source = wf_window_trades.copy()
                            wf_liq_source["trade_month"] = wf_liq_source["trade_date"].dt.to_period("M").astype(str)
                            wf_today = pd.Timestamp.today().normalize()

                            if "trade_amount" not in wf_liq_source.columns:
                                wf_liq_source["trade_amount"] = 0.0

                            wf_liq_by_cusip = (
                                wf_liq_source.groupby("cusip", dropna=False)
                                .agg(
                                    trade_count=("trade_date", "count"),
                                    latest_trade=("trade_date", "max"),
                                    active_months=("trade_month", "nunique"),
                                    total_trade_amount=("trade_amount", "sum"),
                                )
                                .reset_index()
                            )

                            if wf_liq_by_cusip.empty:
                                wf_avg_liquidity_score = pd.NA
                                wf_liquidity_premium_bps = 10.0
                                wf_liquidity_note = "No CUSIP-level liquidity score available; default proxy used."
                            else:
                                wf_liq_by_cusip["days_since_last_trade"] = (
                                    wf_today - wf_liq_by_cusip["latest_trade"]
                                ).dt.days
                                wf_liq_by_cusip["recent_90d_trades"] = wf_liq_by_cusip["trade_count"]
                                wf_liq_by_cusip["liquidity_score"] = (
                                    wf_liq_by_cusip["trade_count"].rank(pct=True) * 35
                                    + wf_liq_by_cusip["total_trade_amount"].rank(pct=True) * 25
                                    + wf_liq_by_cusip["recent_90d_trades"].rank(pct=True) * 25
                                    + (1 - wf_liq_by_cusip["days_since_last_trade"].rank(pct=True)) * 15
                                )
                                wf_avg_liquidity_score = wf_liq_by_cusip["liquidity_score"].mean()

                                if pd.isna(wf_avg_liquidity_score):
                                    wf_liquidity_premium_bps = 10.0
                                    wf_liquidity_note = "Liquidity score unavailable; default proxy used."
                                elif wf_avg_liquidity_score < 45:
                                    wf_liquidity_premium_bps = 15.0
                                    wf_liquidity_note = "Low liquidity bucket proxy."
                                elif wf_avg_liquidity_score < 75:
                                    wf_liquidity_premium_bps = 7.5
                                    wf_liquidity_note = "Medium liquidity bucket proxy."
                                else:
                                    wf_liquidity_premium_bps = 2.5
                                    wf_liquidity_note = "High liquidity bucket proxy."

                            callable_cols = [c for c in ["call_date", "call_date_bond"] if c in wf_window_trades.columns]
                            wf_callable_share = 0.0
                            if callable_cols:
                                call_col = callable_cols[0]
                                parsed_calls = pd.to_datetime(wf_window_trades[call_col], errors="coerce")
                                wf_callable_share = parsed_calls.notna().mean()
                            wf_callable_adjustment_bps = 5.0 if wf_callable_share >= 0.50 else 0.0

                            wf_residual_bps = (
                                wf_total_spread_bps
                                - wf_rating_premium_bps
                                - wf_liquidity_premium_bps
                                - wf_callable_adjustment_bps
                            )

                            waterfall_df = pd.DataFrame(
                                {
                                    "Component": [
                                        "AAA / MMD Base",
                                        "Rating Premium",
                                        "Liquidity Premium",
                                        "Callable Adjustment",
                                        "Residual / Issuer-Specific Premium",
                                        "Implied Issuer Yield",
                                    ],
                                    "Value": [
                                        wf_aaa_yield,
                                        wf_rating_premium_bps / 100,
                                        wf_liquidity_premium_bps / 100,
                                        wf_callable_adjustment_bps / 100,
                                        wf_residual_bps / 100,
                                        wf_avg_issuer_yield,
                                    ],
                                    "Display": [
                                        f"{wf_aaa_yield:.2f}%",
                                        f"{wf_rating_premium_bps:+.1f} bp",
                                        f"{wf_liquidity_premium_bps:+.1f} bp",
                                        f"{wf_callable_adjustment_bps:+.1f} bp",
                                        f"{wf_residual_bps:+.1f} bp",
                                        f"{wf_avg_issuer_yield:.2f}%",
                                    ],
                                }
                            )

                            wf_fig = go.Figure(
                                go.Waterfall(
                                    name="Spread Attribution",
                                    orientation="v",
                                    measure=["absolute", "relative", "relative", "relative", "relative", "total"],
                                    x=waterfall_df["Component"],
                                    y=waterfall_df["Value"],
                                    text=waterfall_df["Display"],
                                    textposition="outside",
                                    connector={"line": {"width": 1}},
                                )
                            )
                            wf_fig.update_layout(
                                title=f"{selected_issuer} Spread Attribution Waterfall ({wf_bucket}, vs AAA/MMD)",
                                yaxis_title="Yield / Spread Contribution (%)",
                                height=540,
                                showlegend=False,
                            )
                            st.plotly_chart(wf_fig, use_container_width=True)

                            wf_metric1, wf_metric2, wf_metric3, wf_metric4 = st.columns(4)
                            wf_metric1.metric("Issuer Yield", f"{wf_avg_issuer_yield:.2f}%")
                            wf_metric2.metric("AAA / MMD Yield", f"{wf_aaa_yield:.2f}%")
                            wf_metric3.metric("Total Spread vs AAA", f"{wf_total_spread_bps:+.1f} bp")
                            wf_metric4.metric("Residual Premium", f"{wf_residual_bps:+.1f} bp")

                            if wf_residual_bps > 15:
                                st.info(
                                    f"Read-through: after modeled rating, liquidity, and callable components, "
                                    f"{selected_issuer}'s {wf_bucket} bucket still shows a positive residual premium "
                                    f"of {wf_residual_bps:+.1f} bp. This may indicate issuer-specific cheapness, "
                                    f"sector/supply pressure, data noise, or a component assumption that should be reviewed."
                                )
                            elif wf_residual_bps < -15:
                                st.info(
                                    f"Read-through: the modeled components exceed the observed spread by "
                                    f"{abs(wf_residual_bps):.1f} bp. This may indicate rich trading, stronger demand, "
                                    f"or overly conservative component assumptions."
                                )
                            else:
                                st.info(
                                    "Read-through: modeled components broadly explain the observed spread versus AAA/MMD. "
                                    "Residual premium is relatively modest."
                                )

                            with st.expander("Waterfall calculation audit table", expanded=False):
                                audit_df = pd.DataFrame(
                                    [
                                        {"Metric": "Selected issuer", "Value": selected_issuer},
                                        {"Metric": "Maturity bucket", "Value": wf_bucket},
                                        {"Metric": "MMD tenor", "Value": wf_tenor},
                                        {"Metric": "Latest issuer trade date", "Value": wf_latest_trade_date.strftime("%Y-%m-%d")},
                                        {"Metric": "Benchmark curve date", "Value": wf_benchmark_date.strftime("%Y-%m-%d")},
                                        {"Metric": "Issuer yield lookback", "Value": f"{wf_lookback_days} days"},
                                        {"Metric": "Trade count in lookback", "Value": f"{wf_trade_count:,}"},
                                        {"Metric": "Total trade amount in lookback", "Value": f"{wf_total_trade_amount:,.0f}"},
                                        {"Metric": "Average issuer yield", "Value": f"{wf_avg_issuer_yield:.4f}%"},
                                        {"Metric": "AAA/MMD yield", "Value": f"{wf_aaa_yield:.4f}%"},
                                        {"Metric": "Total spread vs AAA", "Value": f"{wf_total_spread_bps:+.2f} bp"},
                                        {"Metric": "Rating premium assumption", "Value": f"{wf_rating} / {wf_tenor}: {wf_rating_premium_bps:+.2f} bp"},
                                        {"Metric": "Average liquidity score", "Value": "" if pd.isna(wf_avg_liquidity_score) else f"{wf_avg_liquidity_score:.2f}"},
                                        {"Metric": "Liquidity premium proxy", "Value": f"{wf_liquidity_premium_bps:+.2f} bp — {wf_liquidity_note}"},
                                        {"Metric": "Callable share proxy", "Value": f"{wf_callable_share:.1%}"},
                                        {"Metric": "Callable adjustment proxy", "Value": f"{wf_callable_adjustment_bps:+.2f} bp"},
                                        {"Metric": "Residual / issuer-specific premium", "Value": f"{wf_residual_bps:+.2f} bp"},
                                        {"Metric": "Benchmark source", "Value": wf_aaa_meta.get("benchmark_source")},
                                        {"Metric": "Benchmark source column", "Value": wf_aaa_meta.get("source_column")},
                                    ]
                                )
                                st.dataframe(audit_df, use_container_width=True, hide_index=True)



section_anchor("market-narrative", "Market Narrative & Opportunity Map")
with st.expander("Methodology: market narrative and opportunity map", expanded=False):
    st.markdown(
        """
This section combines recent trading behavior with relative-value positioning.

**Trading Activity Timeline**

- Aggregates selected issuer trades by day.
- Shows daily trade count and total par traded.
- Adds a daily average yield line.
- Generates simple narrative events when activity, volume, or yield moves are unusually high relative to the issuer's recent history.

**Rich / Cheap Quadrant**

- Uses security-level observations.
- `x = Liquidity Score`
- `y = Spread to Benchmark` when benchmark data is available; otherwise `y = Average Yield`
- Vertical and horizontal median lines divide the map into four desk-style zones:
    - **Cheap + Liquid**: potential buy candidate
    - **Cheap + Illiquid**: opportunistic / liquidity premium required
    - **Rich + Liquid**: benchmark-like / monitor
    - **Rich + Illiquid**: low priority / avoid
        """
    )

mn_tab1, mn_tab2 = st.tabs(["Trading Activity Timeline", "Rich / Cheap Quadrant"])

with mn_tab1:
    if issuer_trades.empty:
        st.warning("No trade rows found for the selected issuer and filters.")
    else:
        timeline_df = issuer_trades.copy()
        timeline_df["trade_date"] = pd.to_datetime(timeline_df["trade_date"], errors="coerce")
        timeline_df["yield"] = pd.to_numeric(timeline_df["yield"], errors="coerce")
        if "trade_amount" in timeline_df.columns:
            timeline_df["trade_amount"] = pd.to_numeric(timeline_df["trade_amount"], errors="coerce").fillna(0)
        else:
            timeline_df["trade_amount"] = 0.0

        timeline_df = timeline_df.dropna(subset=["trade_date"])
        if timeline_df.empty:
            st.warning("Timeline cannot be generated because no valid trade dates were found.")
        else:
            timeline_daily = (
                timeline_df.groupby(timeline_df["trade_date"].dt.normalize(), as_index=False)
                .agg(
                    trade_date=("trade_date", "first"),
                    trade_count=("cusip", "count") if "cusip" in timeline_df.columns else ("yield", "count"),
                    total_trade_amount=("trade_amount", "sum"),
                    avg_yield=("yield", "mean"),
                )
                .sort_values("trade_date")
            )
            timeline_daily["yield_change_bp"] = timeline_daily["avg_yield"].diff() * 100

            lookback_options = {
                "Latest 30D": 30,
                "Latest 60D": 60,
                "Latest 90D": 90,
                "All": None,
            }
            timeline_window_label = st.selectbox(
                "Timeline Window",
                list(lookback_options.keys()),
                index=2,
                key="timeline_window",
            )
            timeline_days = lookback_options[timeline_window_label]
            timeline_plot = timeline_daily.copy()
            if timeline_days is not None and not timeline_plot.empty:
                cutoff = timeline_plot["trade_date"].max() - pd.Timedelta(days=timeline_days)
                timeline_plot = timeline_plot[timeline_plot["trade_date"] >= cutoff].copy()

            if timeline_plot.empty:
                st.info("No timeline observations are available for the selected window.")
            else:
                tl_fig = px.bar(
                    timeline_plot,
                    x="trade_date",
                    y="trade_count",
                    hover_data={
                        "total_trade_amount": ":,.0f",
                        "avg_yield": ":.2f",
                        "yield_change_bp": ":.1f",
                    },
                    title=f"{selected_issuer} Trading Activity Timeline",
                    labels={
                        "trade_date": "Trade Date",
                        "trade_count": "Trade Count",
                        "total_trade_amount": "Total Trade Amount",
                        "avg_yield": "Average Yield",
                        "yield_change_bp": "Daily Yield Change (bp)",
                    },
                )
                tl_fig.add_scatter(
                    x=timeline_plot["trade_date"],
                    y=timeline_plot["avg_yield"],
                    mode="lines+markers",
                    name="Average Yield",
                    yaxis="y2",
                )
                tl_fig.update_layout(
                    height=500,
                    yaxis=dict(title="Trade Count"),
                    yaxis2=dict(title="Average Yield (%)", overlaying="y", side="right"),
                    hovermode="x unified",
                )
                st.plotly_chart(tl_fig, use_container_width=True)

                amount_fig = px.bar(
                    timeline_plot,
                    x="trade_date",
                    y="total_trade_amount",
                    title="Daily Total Par Traded",
                    labels={
                        "trade_date": "Trade Date",
                        "total_trade_amount": "Total Trade Amount",
                    },
                    hover_data={"trade_count": ":,.0f", "avg_yield": ":.2f"},
                )
                amount_fig.update_layout(height=380)
                st.plotly_chart(amount_fig, use_container_width=True)

                # Simple narrative event detection
                event_rows = []
                tc_mean = timeline_daily["trade_count"].mean()
                tc_std = timeline_daily["trade_count"].std()
                amt_mean = timeline_daily["total_trade_amount"].mean()
                amt_std = timeline_daily["total_trade_amount"].std()
                yield_abs_threshold = 10.0

                for _, row in timeline_plot.iterrows():
                    notes = []
                    if pd.notna(tc_std) and tc_std > 0 and row["trade_count"] >= tc_mean + tc_std:
                        notes.append("Trade count above recent norm")
                    if pd.notna(amt_std) and amt_std > 0 and row["total_trade_amount"] >= amt_mean + amt_std:
                        notes.append("Heavy par traded")
                    if pd.notna(row.get("yield_change_bp")) and abs(row["yield_change_bp"]) >= yield_abs_threshold:
                        direction = "moved higher" if row["yield_change_bp"] > 0 else "moved lower"
                        notes.append(f"Average yield {direction} by {row['yield_change_bp']:+.1f} bp")
                    if notes:
                        event_rows.append(
                            {
                                "Date": row["trade_date"].strftime("%Y-%m-%d"),
                                "Narrative Signal": "; ".join(notes),
                                "Trade Count": int(row["trade_count"]),
                                "Total Trade Amount": row["total_trade_amount"],
                                "Average Yield": row["avg_yield"],
                                "Yield Change (bp)": row.get("yield_change_bp"),
                            }
                        )

                if event_rows:
                    st.subheader("Narrative Signals")
                    events_df = pd.DataFrame(event_rows)
                    st.dataframe(events_df, use_container_width=True, hide_index=True)
                else:
                    st.info("No unusually large activity/volume/yield-move events were detected in the selected timeline window.")

with mn_tab2:
    if "liq" not in locals() or liq.empty:
        st.info("The quadrant map will be available after liquidity metrics are calculated for the selected issuer.")
    else:
        quadrant_df = liq.copy()

        # Choose Y-axis: spread to benchmark when available, else avg yield.
        quadrant_y_col = "avg_yield"
        quadrant_y_label = "Average Yield (%)"
        quadrant_source_note = "Using Average Yield because benchmark spread is not available inside this liquidity section."

        if not mmd_df.empty:
            try:
                q_rating = st.selectbox(
                    "Quadrant Benchmark Curve",
                    BENCHMARK_RATINGS,
                    index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
                    key="quadrant_benchmark_rating",
                )
                q_spread_obs = build_spread_observations(
                    market_df=market_df,
                    mmd_df=mmd_df,
                    issuer=selected_issuer,
                    rating=q_rating,
                )
                if not q_spread_obs.empty and "cusip" in issuer_trades.columns:
                    # Approximate CUSIP-level benchmark spread by merging the latest bucket-level spread to each CUSIP bucket.
                    latest_bucket_spread = (
                        q_spread_obs.sort_values("trade_date")
                        .groupby("maturity_bucket", as_index=False)
                        .tail(1)[["maturity_bucket", "spread_to_benchmark_bps"]]
                    )
                    if "maturity" in quadrant_df.columns:
                        pass
                    if "maturity_bucket" not in quadrant_df.columns and "maturity" in quadrant_df.columns:
                        # If liquidity table does not retain bucket, infer from maturity when possible.
                        maturity_tmp = pd.to_datetime(quadrant_df["maturity"], errors="coerce")
                        years_tmp = (maturity_tmp - pd.Timestamp.today()).dt.days / 365.25
                        quadrant_df["maturity_bucket"] = pd.cut(
                            years_tmp,
                            bins=[-float("inf"), 7, 15, 25, float("inf")],
                            labels=MATURITY_BUCKET_ORDER,
                        ).astype("string")
                    if "maturity_bucket" in quadrant_df.columns:
                        quadrant_df = quadrant_df.merge(latest_bucket_spread, on="maturity_bucket", how="left")
                        if quadrant_df["spread_to_benchmark_bps"].notna().any():
                            quadrant_y_col = "spread_to_benchmark_bps"
                            quadrant_y_label = f"Spread to {q_rating} Benchmark (bps)"
                            quadrant_source_note = (
                                "Using latest available bucket-level spread to benchmark. "
                                "Higher values generally indicate cheaper relative value."
                            )
            except Exception as exc:
                st.warning(f"Benchmark spread overlay was unavailable, so the quadrant uses average yield. Details: {exc}")

        required_cols = ["liquidity_score", quadrant_y_col]
        for col in required_cols:
            if col in quadrant_df.columns:
                quadrant_df[col] = pd.to_numeric(quadrant_df[col], errors="coerce")
        quadrant_df = quadrant_df.dropna(subset=[c for c in required_cols if c in quadrant_df.columns])

        if quadrant_df.empty:
            st.warning("No usable observations were available for the rich/cheap quadrant.")
        else:
            valid_buckets = MATURITY_BUCKET_ORDER
            if "maturity_bucket" not in quadrant_df.columns:
                quadrant_df["maturity_bucket"] = "Unknown"
            quadrant_df["maturity_bucket"] = quadrant_df["maturity_bucket"].astype(str)

            if "total_trade_amount" not in quadrant_df.columns:
                quadrant_df["total_trade_amount"] = 1
            quadrant_df["total_trade_amount"] = pd.to_numeric(
                quadrant_df["total_trade_amount"], errors="coerce"
            ).fillna(0).clip(lower=0)
            if quadrant_df["total_trade_amount"].sum() <= 0:
                quadrant_df["point_size"] = 10
                q_size_col = "point_size"
            else:
                q_size_col = "total_trade_amount"

            q_median_liq = quadrant_df["liquidity_score"].median()
            q_median_y = quadrant_df[quadrant_y_col].median()

            quadrant_df["Quadrant"] = "Unclassified"
            quadrant_df.loc[
                (quadrant_df["liquidity_score"] >= q_median_liq) & (quadrant_df[quadrant_y_col] >= q_median_y),
                "Quadrant",
            ] = "Cheap + Liquid"
            quadrant_df.loc[
                (quadrant_df["liquidity_score"] < q_median_liq) & (quadrant_df[quadrant_y_col] >= q_median_y),
                "Quadrant",
            ] = "Cheap + Illiquid"
            quadrant_df.loc[
                (quadrant_df["liquidity_score"] >= q_median_liq) & (quadrant_df[quadrant_y_col] < q_median_y),
                "Quadrant",
            ] = "Rich + Liquid"
            quadrant_df.loc[
                (quadrant_df["liquidity_score"] < q_median_liq) & (quadrant_df[quadrant_y_col] < q_median_y),
                "Quadrant",
            ] = "Rich + Illiquid"

            st.caption(quadrant_source_note)

            q_fig = px.scatter(
                quadrant_df,
                x="liquidity_score",
                y=quadrant_y_col,
                size=q_size_col,
                size_max=38,
                color="Quadrant",
                hover_name="cusip" if "cusip" in quadrant_df.columns else None,
                hover_data=[c for c in [
                    "maturity_bucket", "liquidity_tier", "trade_count", "recent_90d_trades",
                    "days_since_last_trade", "avg_yield", "spread_to_benchmark_bps",
                    "total_trade_amount", "outstanding_amount"
                ] if c in quadrant_df.columns],
                title=f"{selected_issuer} Rich / Cheap Liquidity Quadrant",
                labels={
                    "liquidity_score": "Liquidity Score",
                    quadrant_y_col: quadrant_y_label,
                    "total_trade_amount": "Total Trade Amount",
                },
            )
            q_fig.add_vline(x=q_median_liq, line_dash="dash", opacity=0.45)
            q_fig.add_hline(y=q_median_y, line_dash="dash", opacity=0.45)
            q_fig.add_annotation(
                x=0.98, y=0.98, xref="paper", yref="paper",
                text="Cheap + Liquid<br>Buy candidate",
                showarrow=False, align="right",
                bgcolor="rgba(255,255,255,0.75)",
            )
            q_fig.add_annotation(
                x=0.02, y=0.98, xref="paper", yref="paper",
                text="Cheap + Illiquid<br>Opportunistic",
                showarrow=False, align="left",
                bgcolor="rgba(255,255,255,0.75)",
            )
            q_fig.add_annotation(
                x=0.98, y=0.02, xref="paper", yref="paper",
                text="Rich + Liquid<br>Benchmark / monitor",
                showarrow=False, align="right",
                bgcolor="rgba(255,255,255,0.75)",
            )
            q_fig.add_annotation(
                x=0.02, y=0.02, xref="paper", yref="paper",
                text="Rich + Illiquid<br>Low priority",
                showarrow=False, align="left",
                bgcolor="rgba(255,255,255,0.75)",
            )
            q_fig.update_layout(height=560, hovermode="closest")
            st.plotly_chart(q_fig, use_container_width=True)

            q_summary = (
                quadrant_df.groupby("Quadrant", as_index=False)
                .agg(
                    cusip_count=("cusip", "count") if "cusip" in quadrant_df.columns else ("liquidity_score", "count"),
                    avg_liquidity_score=("liquidity_score", "mean"),
                    avg_y_axis=(quadrant_y_col, "mean"),
                    total_trade_amount=("total_trade_amount", "sum") if "total_trade_amount" in quadrant_df.columns else ("liquidity_score", "count"),
                )
            )
            st.dataframe(q_summary, use_container_width=True, hide_index=True)

            with st.expander("Quadrant security-level table", expanded=False):
                display_cols = [
                    "Quadrant", "cusip", "maturity_bucket", "liquidity_score", quadrant_y_col,
                    "avg_yield", "trade_count", "recent_90d_trades", "days_since_last_trade",
                    "total_trade_amount", "outstanding_amount", "liquidity_tier",
                ]
                st.dataframe(
                    quadrant_df[[c for c in display_cols if c in quadrant_df.columns]]
                    .sort_values(["Quadrant", "liquidity_score"], ascending=[True, False])
                    .head(5000),
                    use_container_width=True,
                    hide_index=True,
                )



section_anchor("peer-rv", "Peer Relative Value Comparison")
with st.expander("Methodology: peer relative value comparison", expanded=False):
    st.markdown(
        """
This module is intentionally **optional**. It only becomes fully useful when the uploaded dataset contains multiple issuers.

**Purpose:**

- Compare the selected issuer against uploaded peers.
- Prefer same-sector peers when available.
- Allow manual cross-sector comparison when needed, with a warning that sector differences may dominate issuer-level relative value.

**Core calculation:**

`Issuer Spread to Benchmark = (Average Issuer Yield - Benchmark Yield) × 100`

**Default framework:**

- Peer universe = uploaded issuers in the same sector as the selected issuer.
- Benchmark = user-selected AAA / AA / A / BBB curve.
- Window = recent average, because municipal trading can be sparse.
- If fewer than two issuers are uploaded, the module shows an information message instead of failing.
        """
    )

if len(uploaded_issuers) < 2:
    st.info(
        "Peer comparison is unavailable. Upload bond/trade data for at least two issuers "
        "to compare relative value across peers."
    )
elif mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable peer spread comparison.")
else:
    peer_base = market_df.copy()
    peer_base["issuer"] = peer_base["issuer"].astype(str)

    # Determine same-sector peer universe from uploaded data.
    selected_sector_for_peers = selected_sector if selected_sector and selected_sector != "Unknown" else None

    if selected_sector_for_peers and "sector" in peer_base.columns:
        same_sector_issuers = sorted(
            peer_base.loc[
                peer_base["sector"].astype(str) == str(selected_sector_for_peers),
                "issuer",
            ].dropna().astype(str).unique().tolist()
        )
    else:
        same_sector_issuers = []

    peer_col1, peer_col2, peer_col3 = st.columns([1.2, 1, 1])
    with peer_col1:
        peer_mode = st.radio(
            "Peer Universe",
            ["Same-sector uploaded peers", "Manual selection"],
            index=0 if len(same_sector_issuers) >= 2 else 1,
            horizontal=False,
            help="Same-sector comparison is usually cleaner. Manual selection is useful when sector data is missing or when the team wants a custom comp set.",
        )
    with peer_col2:
        peer_rating = st.selectbox(
            "Peer Benchmark Curve",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="peer_benchmark_rating",
        )
    with peer_col3:
        peer_window_label = st.selectbox(
            "Peer Lookback Window",
            ["Latest 30D", "Latest 60D", "Latest 90D", "All"],
            index=0,
            key="peer_lookback_window",
        )

    if peer_mode == "Same-sector uploaded peers":
        if len(same_sector_issuers) < 2:
            st.info(
                "No same-sector peer set with at least two issuers was detected. "
                "Use Manual selection to compare across uploaded issuers, but interpret cross-sector comparisons carefully."
            )
            peer_issuers = [selected_issuer]
        else:
            st.success(
                f"Detected {len(same_sector_issuers):,} uploaded issuer(s) in sector: {selected_sector_for_peers}."
            )
            default_peer_issuers = same_sector_issuers
            peer_issuers = st.multiselect(
                "Same-sector peers",
                same_sector_issuers,
                default=default_peer_issuers,
                key="same_sector_peer_issuers",
            )
    else:
        st.warning(
            "Manual peer selection can compare issuers across sectors. Cross-sector spreads may reflect sector risk, "
            "not just issuer-level relative value."
        )
        default_manual = [selected_issuer]
        # add up to 3 additional issuers by default
        for issuer in uploaded_issuers:
            if issuer != selected_issuer and len(default_manual) < 4:
                default_manual.append(issuer)
        peer_issuers = st.multiselect(
            "Manual peers from uploaded issuers",
            uploaded_issuers,
            default=default_manual,
            key="manual_peer_issuers",
        )

    peer_days = {"Latest 30D": 30, "Latest 60D": 60, "Latest 90D": 90, "All": None}[peer_window_label]

    if len(peer_issuers) < 2:
        st.info("Select at least two issuers to generate peer comparison charts.")
    else:
        peer_work = peer_base[peer_base["issuer"].isin(peer_issuers)].copy()
        peer_work["trade_date"] = pd.to_datetime(peer_work["trade_date"], errors="coerce").dt.normalize()
        peer_work["yield"] = pd.to_numeric(peer_work["yield"], errors="coerce")
        if "trade_amount" in peer_work.columns:
            peer_work["trade_amount"] = pd.to_numeric(peer_work["trade_amount"], errors="coerce").fillna(0)
        else:
            peer_work["trade_amount"] = 0.0

        peer_work = peer_work.dropna(subset=["trade_date", "yield", "issuer", "maturity_bucket"])
        peer_work = peer_work[peer_work["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)].copy()

        if peer_work.empty:
            st.warning("No usable peer trade rows were found after cleaning.")
        else:
            latest_peer_date = peer_work["trade_date"].max()
            if peer_days is not None:
                peer_cutoff = latest_peer_date - pd.Timedelta(days=peer_days)
                peer_work = peer_work[peer_work["trade_date"] >= peer_cutoff].copy()

            if peer_work.empty:
                st.info("No peer observations remain inside the selected lookback window.")
            else:
                # Build issuer/bucket average yields over lookback.
                peer_summary = (
                    peer_work.groupby(["issuer", "maturity_bucket"], as_index=False)
                    .agg(
                        avg_yield=("yield", "mean"),
                        trade_count=("yield", "count"),
                        latest_trade=("trade_date", "max"),
                        total_trade_amount=("trade_amount", "sum"),
                    )
                )

                # Build benchmark curves and use latest benchmark observation at/before latest peer date.
                peer_date_col = _detect_mmd_date_column(mmd_df)
                if peer_date_col is None:
                    st.warning("Peer comparison cannot run because benchmark curve file has no usable date column.")
                else:
                    peer_mmd = mmd_df.copy()
                    peer_mmd[peer_date_col] = pd.to_datetime(peer_mmd[peer_date_col], errors="coerce")
                    peer_mmd = peer_mmd.dropna(subset=[peer_date_col])
                    peer_mmd = peer_mmd[peer_mmd[peer_date_col].dt.normalize() <= latest_peer_date].sort_values(peer_date_col)

                    if peer_mmd.empty:
                        st.warning("No benchmark curve observation was available on or before the latest peer trade date.")
                    else:
                        peer_latest_mmd = peer_mmd.iloc[[-1]].copy()
                        peer_benchmark_date = peer_latest_mmd[peer_date_col].iloc[0]

                        bench_rows = []
                        for bucket in MATURITY_BUCKET_ORDER:
                            tenor = MMD_BUCKET_MAP.get(bucket, "10Y")
                            y, meta = get_benchmark_curve(peer_latest_mmd, tenor, peer_rating)
                            if y is not None and pd.notna(y.iloc[0]):
                                bench_rows.append(
                                    {
                                        "maturity_bucket": bucket,
                                        "mmd_tenor": tenor,
                                        "benchmark_yield": float(y.iloc[0]),
                                        "benchmark_source": meta.get("benchmark_source"),
                                        "source_column": meta.get("source_column"),
                                        "rating_spread_bps": meta.get("rating_spread_bps"),
                                    }
                                )
                        bench_df = pd.DataFrame(bench_rows)

                        if bench_df.empty:
                            st.warning("Selected benchmark curve could not be built for any maturity bucket.")
                        else:
                            peer_summary = peer_summary.merge(bench_df, on="maturity_bucket", how="left")
                            peer_summary["spread_to_benchmark_bps"] = (
                                peer_summary["avg_yield"] - peer_summary["benchmark_yield"]
                            ) * 100
                            peer_summary = peer_summary.dropna(subset=["spread_to_benchmark_bps"])

                            if peer_summary.empty:
                                st.info("No overlapping peer observations and benchmark tenors were available.")
                            else:
                                maturity_order = MATURITY_BUCKET_ORDER
                                peer_summary["maturity_bucket"] = pd.Categorical(
                                    peer_summary["maturity_bucket"],
                                    categories=maturity_order,
                                    ordered=True,
                                )
                                peer_summary = peer_summary.sort_values(["issuer", "maturity_bucket"])

                                st.subheader("1. Peer Spread Curve Comparison")
                                peer_curve_fig = px.line(
                                    peer_summary,
                                    x="maturity_bucket",
                                    y="spread_to_benchmark_bps",
                                    color="issuer",
                                    markers=True,
                                    hover_data=[
                                        "avg_yield",
                                        "benchmark_yield",
                                        "trade_count",
                                        "total_trade_amount",
                                        "latest_trade",
                                        "benchmark_source",
                                        "source_column",
                                    ],
                                    title=f"Peer Spread Curve Comparison vs {peer_rating} Benchmark",
                                    labels={
                                        "maturity_bucket": "Maturity Bucket",
                                        "spread_to_benchmark_bps": "Spread to Benchmark (bps)",
                                        "issuer": "Issuer",
                                    },
                                )
                                peer_curve_fig.add_hline(y=0, line_dash="dash", opacity=0.45)
                                peer_curve_fig.update_layout(height=520, hovermode="x unified")
                                st.plotly_chart(peer_curve_fig, use_container_width=True)

                                st.subheader("2. Peer Spread Heatmap")
                                peer_matrix = peer_summary.pivot_table(
                                    index="issuer",
                                    columns="maturity_bucket",
                                    values="spread_to_benchmark_bps",
                                    aggfunc="mean",
                                    observed=False,
                                ).reindex(columns=maturity_order)
                                peer_text = peer_matrix.map(lambda x: "" if pd.isna(x) else f"{x:+.1f} bp")
                                peer_heatmap_fig = px.imshow(
                                    peer_matrix.astype(float),
                                    x=peer_matrix.columns.astype(str),
                                    y=peer_matrix.index.astype(str),
                                    color_continuous_scale=["#1a9850", "#f7f7f7", "#d73027"],
                                    color_continuous_midpoint=0,
                                    aspect="auto",
                                    title=f"Peer Spread Heatmap vs {peer_rating}",
                                    labels={
                                        "x": "Maturity Bucket",
                                        "y": "Issuer",
                                        "color": "Spread (bps)",
                                    },
                                )
                                peer_heatmap_fig.update_traces(
                                    text=peer_text.values,
                                    texttemplate="%{text}",
                                    hovertemplate="Issuer=%{y}<br>Bucket=%{x}<br>Spread=%{z:.1f} bp<extra></extra>",
                                )
                                peer_heatmap_fig.update_layout(height=max(380, 70 * len(peer_matrix.index)))
                                st.plotly_chart(peer_heatmap_fig, use_container_width=True)

                                st.subheader("3. Peer Ranking Table")

                                # -------------------------------------------------------------
                                # Peer ranking methodology
                                # -------------------------------------------------------------
                                # Average spread is a simple arithmetic average across maturity
                                # buckets. Weighted average spread gives more influence to buckets
                                # with larger uploaded trade amount / par volume.
                                #
                                # Spread-to-benchmark is calculated as:
                                #   (Issuer average trade yield - benchmark yield) * 100
                                # where yields are in percentage-point terms and the output is bps.
                                # -------------------------------------------------------------

                                def _weighted_avg_spread_by_amount(group: pd.DataFrame) -> float:
                                    spreads = pd.to_numeric(group["spread_to_benchmark_bps"], errors="coerce")
                                    weights = pd.to_numeric(group["total_trade_amount"], errors="coerce").fillna(0)
                                    valid = spreads.notna() & weights.gt(0)
                                    if valid.any():
                                        return float((spreads[valid] * weights[valid]).sum() / weights[valid].sum())
                                    return float(spreads.mean())

                                weighted_spread = (
                                    peer_summary.groupby("issuer")
                                    .apply(_weighted_avg_spread_by_amount)
                                    .reset_index(name="weighted_avg_spread_bps")
                                )

                                weight_basis = (
                                    peer_summary.groupby("issuer", as_index=False)
                                    .agg(weighted_trade_amount_used=("total_trade_amount", "sum"))
                                )
                                weight_basis["weighting_method"] = weight_basis["weighted_trade_amount_used"].apply(
                                    lambda x: "Trade-amount weighted" if pd.notna(x) and x > 0 else "Unweighted fallback"
                                )

                                peer_rank = (
                                    peer_summary.groupby("issuer", as_index=False)
                                    .agg(
                                        avg_spread_bps=("spread_to_benchmark_bps", "mean"),
                                        max_spread_bps=("spread_to_benchmark_bps", "max"),
                                        trade_count=("trade_count", "sum"),
                                        total_trade_amount=("total_trade_amount", "sum"),
                                        latest_trade=("latest_trade", "max"),
                                    )
                                    .merge(weighted_spread, on="issuer", how="left")
                                    .merge(weight_basis[["issuer", "weighting_method"]], on="issuer", how="left")
                                )

                                peer_rank["rank_by_weighted_spread"] = (
                                    peer_rank["weighted_avg_spread_bps"]
                                    .rank(method="dense", ascending=False)
                                    .astype(int)
                                )
                                peer_rank["rank_by_avg_spread"] = (
                                    peer_rank["avg_spread_bps"]
                                    .rank(method="dense", ascending=False)
                                    .astype(int)
                                )
                                peer_rank = peer_rank.sort_values(
                                    ["rank_by_weighted_spread", "rank_by_avg_spread", "issuer"],
                                    ascending=[True, True, True],
                                )

                                with st.expander("Peer ranking methodology", expanded=False):
                                    st.markdown(f"""
### How the Peer Ranking Is Calculated

This table ranks issuers by **weighted average spread-to-benchmark**, measured in basis points.

**Spread-to-benchmark formula:**

`Spread (bps) = (Issuer Average Trade Yield - Benchmark Yield) × 100`

**Weighted average spread formula:**

`Weighted Avg Spread = Σ(Spread × Trade Amount) / Σ(Trade Amount)`

### Current Ranking Logic

- Benchmark rating selected: **{peer_rating}**
- Lookback window selected: **{peer_window_label}**
- Maturity buckets included: selected peer maturity scope shown in this section
- Issuers are ranked from **widest** to **tightest** based on `weighted_avg_spread_bps`
- If trade amount is unavailable or zero, the dashboard falls back to the simple average spread

### Interpretation

- **Positive spread**: issuer is trading wider than the selected benchmark
- **Negative spread**: issuer is trading tighter than the selected benchmark
- **Higher rank**: wider screened spread, potentially more yield compensation versus benchmark

### Important Caveats

This is a **relative-value screening tool**, not a standalone trading recommendation. Rankings can be affected by:

- Trade size and odd-lot effects
- Liquidity differences across CUSIPs
- Callable structures and coupon differences
- Credit quality differences not fully captured by benchmark rating
- Limited observations in the selected time window
- Stale or uneven trading activity across issuers

Use this ranking as a first-pass screen, then review the CUSIP-level drilldown and liquidity metrics before drawing trading conclusions.
""")

                                st.dataframe(
                                    peer_rank[
                                        [
                                            "rank_by_weighted_spread",
                                            "rank_by_avg_spread",
                                            "issuer",
                                            "weighted_avg_spread_bps",
                                            "avg_spread_bps",
                                            "max_spread_bps",
                                            "trade_count",
                                            "total_trade_amount",
                                            "weighting_method",
                                            "latest_trade",
                                        ]
                                    ],
                                    use_container_width=True,
                                    hide_index=True,
                                )

                                # Read-through
                                if not peer_rank.empty:
                                    widest = peer_rank.iloc[0]
                                    tightest = peer_rank.iloc[-1]
                                    st.info(
                                        f"Peer read-through: {widest['issuer']} screens widest on a trade-amount weighted basis "
                                        f"at {widest['weighted_avg_spread_bps']:+.1f} bp versus {peer_rating}, while "
                                        f"{tightest['issuer']} screens tightest at {tightest['weighted_avg_spread_bps']:+.1f} bp. "
                                        f"The simple average spread is also shown for comparison. Use this as a screening signal "
                                        f"and review CUSIP-level liquidity before drawing trading conclusions."
                                    )

                                with st.expander("Peer comparison audit table", expanded=False):
                                    audit_cols = [
                                        "issuer",
                                        "maturity_bucket",
                                        "avg_yield",
                                        "benchmark_yield",
                                        "spread_to_benchmark_bps",
                                        "trade_count",
                                        "total_trade_amount",
                                        "latest_trade",
                                        "mmd_tenor",
                                        "benchmark_source",
                                        "source_column",
                                        "rating_spread_bps",
                                    ]
                                    audit_df = peer_summary[[c for c in audit_cols if c in peer_summary.columns]].copy()
                                    for c in ["avg_yield", "benchmark_yield", "spread_to_benchmark_bps", "rating_spread_bps"]:
                                        if c in audit_df.columns:
                                            audit_df[c] = pd.to_numeric(audit_df[c], errors="coerce").round(2)
                                    st.caption(
                                        f"Benchmark date used: {peer_benchmark_date.strftime('%Y-%m-%d')}. "
                                        f"Lookback window: {peer_window_label}."
                                    )
                                    st.dataframe(audit_df, use_container_width=True, hide_index=True)




section_anchor("cross-issuer-rv", "Cross-Issuer Relative Value Analytics")
with st.expander("Methodology: cross-issuer relative value analytics", expanded=False):
    st.markdown(
        """
This section upgrades peer comparison from **visual comparison** into a systematic **issuer-bucket ranking framework**.

**Core purpose:**

- Identify which issuer / maturity bucket screens cheap or rich versus the uploaded peer group.
- Convert peer spreads into **peer gaps**, **z-scores**, and **relative value scores**.
- Keep the module optional: it only becomes meaningful when at least two issuers are uploaded.

**Core calculations:**

`Issuer Spread = (Average Issuer Yield - Benchmark Yield) × 100`

`Peer Gap = Issuer Spread - Peer Median Spread within the same maturity bucket`

`Peer Z-Score = (Issuer Spread - Bucket Mean Spread) / Bucket Spread Std`

**Simplified RV Score:**

`RV Score = 45% Spread Percentile + 35% Liquidity Percentile + 20% Trade Activity Percentile`

**How to read it:**

- Positive peer gap: issuer/bucket is wider than peers, potentially cheaper.
- Negative peer gap: issuer/bucket is tighter than peers, potentially richer.
- Higher RV score: screens cheaper while retaining better liquidity/trading support.
- This is a screening tool, not a final trade recommendation.
        """
    )

if len(uploaded_issuers) < 2:
    st.info(
        "Cross-issuer RV analytics is unavailable. Upload at least two issuers to compare issuer-bucket relative value."
    )
elif mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable cross-issuer spread analytics.")
else:
    xrv_col1, xrv_col2, xrv_col3, xrv_col4 = st.columns([1.2, 1, 1, 1])
    with xrv_col1:
        xrv_mode = st.radio(
            "Cross-Issuer Universe",
            ["Same-sector uploaded issuers", "Manual issuer set"],
            index=0 if "sector" in market_df.columns and selected_sector != "Unknown" else 1,
            key="xrv_universe_mode",
            help="Same-sector is cleaner. Manual set is useful when sector data is missing or when the desk wants a custom comp set.",
        )
    with xrv_col2:
        xrv_rating = st.selectbox(
            "X-Issuer Benchmark",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="xrv_benchmark",
        )
    with xrv_col3:
        xrv_window = st.selectbox(
            "X-Issuer Lookback",
            ["Latest 30D", "Latest 60D", "Latest 90D", "All"],
            index=1,
            key="xrv_window",
        )
    with xrv_col4:
        xrv_min_trades = st.number_input(
            "Min Trades / Bucket",
            min_value=1,
            max_value=100,
            value=1,
            step=1,
            key="xrv_min_trades",
        )

    if xrv_mode == "Same-sector uploaded issuers":
        if "sector" not in market_df.columns or selected_sector == "Unknown":
            st.info("Same-sector universe is unavailable because sector data is missing or unknown. Use Manual issuer set.")
            xrv_issuers = [selected_issuer]
        else:
            xrv_issuers = sorted(
                market_df.loc[
                    market_df["sector"].astype(str) == str(selected_sector),
                    "issuer",
                ].dropna().astype(str).unique().tolist()
            )
            if len(xrv_issuers) < 2:
                st.info(
                    "Fewer than two same-sector issuers were detected. Use Manual issuer set or upload more peer data."
                )
    else:
        default_xrv = [selected_issuer]
        for issuer in uploaded_issuers:
            if issuer != selected_issuer and len(default_xrv) < 5:
                default_xrv.append(issuer)
        xrv_issuers = st.multiselect(
            "Manual cross-issuer set",
            uploaded_issuers,
            default=default_xrv,
            key="xrv_manual_issuers",
        )
        st.warning(
            "Manual cross-issuer comparison may include different sectors. Interpret peer gaps carefully because sector risk can dominate issuer-level RV."
        )

    if len(xrv_issuers) < 2:
        st.info("Select or upload at least two issuers to generate cross-issuer RV analytics.")
    else:
        xrv_days = {"Latest 30D": 30, "Latest 60D": 60, "Latest 90D": 90, "All": None}[xrv_window]

        xrv_df = market_df[market_df["issuer"].astype(str).isin([str(x) for x in xrv_issuers])].copy()
        xrv_df["trade_date"] = pd.to_datetime(xrv_df["trade_date"], errors="coerce").dt.normalize()
        xrv_df["yield"] = pd.to_numeric(xrv_df["yield"], errors="coerce")
        if "trade_amount" in xrv_df.columns:
            xrv_df["trade_amount"] = pd.to_numeric(xrv_df["trade_amount"], errors="coerce").fillna(0)
        else:
            xrv_df["trade_amount"] = 0.0

        xrv_df = xrv_df.dropna(subset=["trade_date", "yield", "issuer", "maturity_bucket"])
        xrv_df = xrv_df[xrv_df["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)].copy()

        if xrv_df.empty:
            st.warning("No usable cross-issuer observations remained after cleaning.")
        else:
            xrv_latest_date = xrv_df["trade_date"].max()
            if xrv_days is not None:
                xrv_cutoff = xrv_latest_date - pd.Timedelta(days=xrv_days)
                xrv_df = xrv_df[xrv_df["trade_date"] >= xrv_cutoff].copy()

            if xrv_df.empty:
                st.info("No cross-issuer observations remain inside the selected lookback window.")
            else:
                xrv_summary = (
                    xrv_df.groupby(["issuer", "maturity_bucket"], as_index=False)
                    .agg(
                        avg_yield=("yield", "mean"),
                        trade_count=("yield", "count"),
                        total_trade_amount=("trade_amount", "sum"),
                        latest_trade=("trade_date", "max"),
                        sector=("sector", "first") if "sector" in xrv_df.columns else ("issuer", "first"),
                    )
                )
                xrv_summary = xrv_summary[xrv_summary["trade_count"] >= xrv_min_trades].copy()

                if xrv_summary.empty:
                    st.info("No issuer-bucket observations met the minimum trade filter.")
                else:
                    date_col = _detect_mmd_date_column(mmd_df)
                    if date_col is None:
                        st.warning("Cross-issuer RV cannot run because the benchmark file has no usable date column.")
                    else:
                        xrv_mmd = mmd_df.copy()
                        xrv_mmd[date_col] = pd.to_datetime(xrv_mmd[date_col], errors="coerce")
                        xrv_mmd = xrv_mmd.dropna(subset=[date_col])
                        xrv_mmd = xrv_mmd[xrv_mmd[date_col].dt.normalize() <= xrv_latest_date].sort_values(date_col)

                        if xrv_mmd.empty:
                            st.warning("No benchmark curve observation was available on or before the latest cross-issuer trade date.")
                        else:
                            xrv_latest_mmd = xrv_mmd.iloc[[-1]].copy()
                            xrv_benchmark_date = xrv_latest_mmd[date_col].iloc[0]

                            bench_rows = []
                            for bucket in MATURITY_BUCKET_ORDER:
                                tenor = MMD_BUCKET_MAP.get(bucket, "10Y")
                                y, meta = get_benchmark_curve(xrv_latest_mmd, tenor, xrv_rating)
                                if y is not None and pd.notna(y.iloc[0]):
                                    bench_rows.append(
                                        {
                                            "maturity_bucket": bucket,
                                            "mmd_tenor": tenor,
                                            "benchmark_yield": float(y.iloc[0]),
                                            "benchmark_source": meta.get("benchmark_source"),
                                            "source_column": meta.get("source_column"),
                                            "rating_spread_bps": meta.get("rating_spread_bps"),
                                        }
                                    )
                            xrv_bench = pd.DataFrame(bench_rows)

                            if xrv_bench.empty:
                                st.warning("Selected benchmark curve could not be built for cross-issuer RV.")
                            else:
                                xrv_summary = xrv_summary.merge(xrv_bench, on="maturity_bucket", how="left")
                                xrv_summary["spread_to_benchmark_bps"] = (
                                    xrv_summary["avg_yield"] - xrv_summary["benchmark_yield"]
                                ) * 100
                                xrv_summary = xrv_summary.dropna(subset=["spread_to_benchmark_bps"])

                                if xrv_summary.empty:
                                    st.info("No issuer-bucket observations had usable benchmark spreads.")
                                else:
                                    # Peer-relative metrics by maturity bucket.
                                    xrv_summary["bucket_peer_median_bps"] = (
                                        xrv_summary.groupby("maturity_bucket")["spread_to_benchmark_bps"].transform("median")
                                    )
                                    xrv_summary["bucket_peer_mean_bps"] = (
                                        xrv_summary.groupby("maturity_bucket")["spread_to_benchmark_bps"].transform("mean")
                                    )
                                    xrv_summary["bucket_peer_std_bps"] = (
                                        xrv_summary.groupby("maturity_bucket")["spread_to_benchmark_bps"].transform("std")
                                    )
                                    xrv_summary["peer_gap_bps"] = (
                                        xrv_summary["spread_to_benchmark_bps"] - xrv_summary["bucket_peer_median_bps"]
                                    )
                                    xrv_summary["peer_z_score"] = (
                                        (xrv_summary["spread_to_benchmark_bps"] - xrv_summary["bucket_peer_mean_bps"])
                                        / xrv_summary["bucket_peer_std_bps"].replace({0: pd.NA})
                                    )

                                    # Liquidity proxy within universe.
                                    today_xrv = pd.Timestamp.today().normalize()
                                    xrv_summary["days_since_last_trade"] = (today_xrv - xrv_summary["latest_trade"]).dt.days
                                    xrv_summary["liquidity_score"] = (
                                        xrv_summary["trade_count"].rank(pct=True) * 35
                                        + xrv_summary["total_trade_amount"].rank(pct=True) * 35
                                        + (1 - xrv_summary["days_since_last_trade"].rank(pct=True)) * 30
                                    )

                                    xrv_summary["spread_percentile"] = (
                                        xrv_summary.groupby("maturity_bucket")["spread_to_benchmark_bps"]
                                        .rank(pct=True)
                                    )
                                    xrv_summary["liquidity_percentile"] = xrv_summary["liquidity_score"].rank(pct=True)
                                    xrv_summary["trade_activity_percentile"] = xrv_summary["trade_count"].rank(pct=True)
                                    xrv_summary["x_issuer_rv_score"] = (
                                        xrv_summary["spread_percentile"] * 45
                                        + xrv_summary["liquidity_percentile"] * 35
                                        + xrv_summary["trade_activity_percentile"] * 20
                                    )

                                    def classify_xrv(row):
                                        if row["peer_gap_bps"] >= 10 and row["liquidity_score"] >= 60:
                                            return "Cheap + Liquid"
                                        if row["peer_gap_bps"] >= 10:
                                            return "Cheap / Needs Liquidity Check"
                                        if row["peer_gap_bps"] <= -10:
                                            return "Rich vs Peers"
                                        return "In Line"

                                    xrv_summary["x_issuer_signal"] = xrv_summary.apply(classify_xrv, axis=1)

                                    st.subheader("1. Peer Gap Matrix")
                                    maturity_order = MATURITY_BUCKET_ORDER
                                    gap_matrix = xrv_summary.pivot_table(
                                        index="issuer",
                                        columns="maturity_bucket",
                                        values="peer_gap_bps",
                                        aggfunc="mean",
                                        observed=False,
                                    ).reindex(columns=maturity_order)
                                    gap_text = gap_matrix.map(lambda x: "" if pd.isna(x) else f"{x:+.1f} bp")
                                    gap_fig = px.imshow(
                                        gap_matrix.astype(float),
                                        x=gap_matrix.columns.astype(str),
                                        y=gap_matrix.index.astype(str),
                                        color_continuous_scale=["#1a9850", "#f7f7f7", "#d73027"],
                                        color_continuous_midpoint=0,
                                        aspect="auto",
                                        title=f"Peer Gap Matrix vs {xrv_rating} Benchmark",
                                        labels={
                                            "x": "Maturity Bucket",
                                            "y": "Issuer",
                                            "color": "Peer Gap (bps)",
                                        },
                                    )
                                    gap_fig.update_traces(
                                        text=gap_text.values,
                                        texttemplate="%{text}",
                                        hovertemplate="Issuer=%{y}<br>Bucket=%{x}<br>Peer Gap=%{z:.1f} bp<extra></extra>",
                                    )
                                    gap_fig.update_layout(height=max(390, 72 * len(gap_matrix.index)))
                                    st.plotly_chart(gap_fig, use_container_width=True)

                                    st.subheader("2. Cross-Issuer RV Ranking")
                                    ranking = xrv_summary.sort_values("x_issuer_rv_score", ascending=False, na_position="last").copy()
                                    rank_cols = [
                                        "issuer",
                                        "sector",
                                        "maturity_bucket",
                                        "spread_to_benchmark_bps",
                                        "bucket_peer_median_bps",
                                        "peer_gap_bps",
                                        "peer_z_score",
                                        "liquidity_score",
                                        "trade_count",
                                        "total_trade_amount",
                                        "days_since_last_trade",
                                        "x_issuer_rv_score",
                                        "x_issuer_signal",
                                    ]
                                    ranking_display = ranking[[c for c in rank_cols if c in ranking.columns]].copy()
                                    for c in [
                                        "spread_to_benchmark_bps",
                                        "bucket_peer_median_bps",
                                        "peer_gap_bps",
                                        "peer_z_score",
                                        "liquidity_score",
                                        "x_issuer_rv_score",
                                    ]:
                                        if c in ranking_display.columns:
                                            ranking_display[c] = pd.to_numeric(ranking_display[c], errors="coerce").round(2)

                                    st.dataframe(ranking_display, use_container_width=True, hide_index=True, height=430)

                                    st.subheader("3. Cross-Issuer Opportunity Map")
                                    xrv_scatter = px.scatter(
                                        ranking,
                                        x="liquidity_score",
                                        y="peer_gap_bps",
                                        size="total_trade_amount",
                                        size_max=38,
                                        color="maturity_bucket",
                                        symbol="x_issuer_signal",
                                        hover_name="issuer",
                                        hover_data=[
                                            "sector",
                                            "maturity_bucket",
                                            "spread_to_benchmark_bps",
                                            "bucket_peer_median_bps",
                                            "peer_z_score",
                                            "trade_count",
                                            "total_trade_amount",
                                            "x_issuer_rv_score",
                                            "x_issuer_signal",
                                        ],
                                        title="Cross-Issuer Relative Value Opportunity Map",
                                        labels={
                                            "liquidity_score": "Liquidity Score",
                                            "peer_gap_bps": "Peer Gap (bps)",
                                            "maturity_bucket": "Maturity Bucket",
                                        },
                                    )
                                    xrv_scatter.add_hline(y=0, line_dash="dash", opacity=0.45)
                                    xrv_scatter.add_vline(x=60, line_dash="dash", opacity=0.35)
                                    xrv_scatter.update_layout(height=540, hovermode="closest")
                                    st.plotly_chart(xrv_scatter, use_container_width=True)

                                    if not ranking.empty:
                                        top = ranking.iloc[0]
                                        st.info(
                                            f"Cross-issuer read-through: {top['issuer']} / {top['maturity_bucket']} screens highest by RV score. "
                                            f"It is {top['peer_gap_bps']:+.1f} bp versus the bucket peer median, "
                                            f"with liquidity score {top['liquidity_score']:.1f} and signal: {top['x_issuer_signal']}."
                                        )

                                    with st.expander("Cross-issuer RV audit table", expanded=False):
                                        audit_cols = [
                                            "issuer",
                                            "sector",
                                            "maturity_bucket",
                                            "avg_yield",
                                            "benchmark_yield",
                                            "spread_to_benchmark_bps",
                                            "bucket_peer_median_bps",
                                            "peer_gap_bps",
                                            "peer_z_score",
                                            "trade_count",
                                            "total_trade_amount",
                                            "latest_trade",
                                            "mmd_tenor",
                                            "benchmark_source",
                                            "source_column",
                                            "rating_spread_bps",
                                        ]
                                        audit_xrv = xrv_summary[[c for c in audit_cols if c in xrv_summary.columns]].copy()
                                        for c in [
                                            "avg_yield",
                                            "benchmark_yield",
                                            "spread_to_benchmark_bps",
                                            "bucket_peer_median_bps",
                                            "peer_gap_bps",
                                            "peer_z_score",
                                            "rating_spread_bps",
                                        ]:
                                            if c in audit_xrv.columns:
                                                audit_xrv[c] = pd.to_numeric(audit_xrv[c], errors="coerce").round(2)
                                        st.caption(
                                            f"Benchmark date used: {xrv_benchmark_date.strftime('%Y-%m-%d')}. "
                                            f"Lookback window: {xrv_window}."
                                        )
                                        st.dataframe(audit_xrv, use_container_width=True, hide_index=True)


section_anchor("historical-spread", "Historical Spread Range & Percentile")
with st.expander("Methodology: historical spread range and percentile", expanded=False):
    st.markdown(
        """
This section compares the selected issuer's current spread against its own historical spread range.

**Core question:**

> Is the current spread historically wide, tight, or normal?

**Calculation:**

`Spread to Benchmark = (Average Issuer Yield - Benchmark Yield) × 100`

The module then calculates:

- **Current Spread**: latest available spread observation in the selected bucket.
- **Historical Median**: median spread over the selected history window.
- **Percentile**: where the current spread ranks versus the historical spread distribution.
- **Z-Score**: current spread distance from historical mean in standard deviation units.

**Signal guide:**

- **> 90th percentile**: very wide / potentially cheap
- **75th–90th percentile**: wide
- **25th–75th percentile**: normal range
- **10th–25th percentile**: tight
- **< 10th percentile**: very tight / potentially rich

Municipal trading can be sparse, so this is a screening signal rather than a final trade recommendation.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable historical spread percentile analysis.")
else:
    hist_col1, hist_col2, hist_col3, hist_col4 = st.columns([1, 1, 1, 1])
    with hist_col1:
        hist_bucket = st.selectbox(
            "Historical Maturity Bucket",
            MATURITY_BUCKET_ORDER,
            index=1,
            key="hist_spread_bucket",
        )
    with hist_col2:
        hist_rating = st.selectbox(
            "Historical Benchmark Curve",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="hist_spread_benchmark",
        )
    with hist_col3:
        hist_window_label = st.selectbox(
            "History Window",
            ["1M", "3M", "6M", "1Y", "All"],
            index=3,
            key="hist_spread_window",
        )
    with hist_col4:
        hist_smoothing = st.selectbox(
            "Smoothing",
            ["Daily", "7D average", "30D average"],
            index=0,
            key="hist_spread_smoothing",
        )

    hist_window_days = {"1M": 30, "3M": 90, "6M": 180, "1Y": 365, "All": None}[hist_window_label]

    hist_obs = build_spread_observations(
        market_df=market_df,
        mmd_df=mmd_df,
        issuer=selected_issuer,
        rating=hist_rating,
    )

    if hist_obs.empty:
        st.warning(
            "No overlapping issuer trade dates and benchmark dates were found for historical spread analysis. "
            "Check that the curve file date range overlaps with trade dates."
        )
    else:
        hist_obs = hist_obs.copy()
        hist_obs["trade_date"] = pd.to_datetime(hist_obs["trade_date"], errors="coerce").dt.normalize()
        hist_obs = hist_obs[
            (hist_obs["maturity_bucket"] == hist_bucket)
            & hist_obs["spread_to_benchmark_bps"].notna()
        ].sort_values("trade_date")

        if hist_obs.empty:
            st.warning(f"No historical spread observations found for {hist_bucket}.")
        else:
            latest_hist_date = hist_obs["trade_date"].max()
            if hist_window_days is not None:
                hist_cutoff = latest_hist_date - pd.Timedelta(days=hist_window_days)
                hist_plot = hist_obs[hist_obs["trade_date"] >= hist_cutoff].copy()
            else:
                hist_plot = hist_obs.copy()

            if hist_plot.empty or len(hist_plot) < 2:
                st.info("Not enough historical spread observations to calculate a percentile for the selected window.")
            else:
                # Aggregate to one observation per date first.
                hist_daily = (
                    hist_plot.groupby("trade_date", as_index=False)
                    .agg(
                        spread_to_benchmark_bps=("spread_to_benchmark_bps", "mean"),
                        avg_yield=("avg_yield", "mean"),
                        benchmark_yield=("benchmark_yield", "mean"),
                        trade_count=("trade_count", "sum"),
                        total_trade_amount=("total_trade_amount", "sum"),
                        benchmark_source=("benchmark_source", "first"),
                        source_column=("source_column", "first"),
                    )
                    .sort_values("trade_date")
                )

                if hist_smoothing == "7D average":
                    hist_daily["display_spread_bps"] = (
                        hist_daily["spread_to_benchmark_bps"].rolling(window=7, min_periods=1).mean()
                    )
                    display_label = "7D Average Spread"
                elif hist_smoothing == "30D average":
                    hist_daily["display_spread_bps"] = (
                        hist_daily["spread_to_benchmark_bps"].rolling(window=30, min_periods=1).mean()
                    )
                    display_label = "30D Average Spread"
                else:
                    hist_daily["display_spread_bps"] = hist_daily["spread_to_benchmark_bps"]
                    display_label = "Daily Spread"

                current_spread = float(hist_daily["display_spread_bps"].iloc[-1])
                hist_values = pd.to_numeric(hist_daily["display_spread_bps"], errors="coerce").dropna()

                hist_median = float(hist_values.median())
                hist_mean = float(hist_values.mean())
                hist_std = float(hist_values.std()) if len(hist_values) > 1 else 0.0
                hist_p25 = float(hist_values.quantile(0.25))
                hist_p75 = float(hist_values.quantile(0.75))
                hist_p10 = float(hist_values.quantile(0.10))
                hist_p90 = float(hist_values.quantile(0.90))
                hist_min = float(hist_values.min())
                hist_max = float(hist_values.max())

                percentile = float((hist_values <= current_spread).mean() * 100)
                z_score = (current_spread - hist_mean) / hist_std if hist_std and hist_std > 0 else pd.NA

                # Signal label
                if percentile >= 90:
                    signal = "Very Wide / Potentially Cheap"
                elif percentile >= 75:
                    signal = "Wide"
                elif percentile >= 25:
                    signal = "Normal Range"
                elif percentile >= 10:
                    signal = "Tight"
                else:
                    signal = "Very Tight / Potentially Rich"

                card1, card2, card3, card4 = st.columns(4)
                card1.metric("Current Spread", f"{current_spread:+.1f} bp")
                card2.metric("Historical Percentile", f"{percentile:.0f}th")
                card3.metric("Historical Median", f"{hist_median:+.1f} bp")
                card4.metric("Z-Score", "N/A" if pd.isna(z_score) else f"{z_score:+.2f}")

                st.info(
                    f"Read-through: {selected_issuer}'s {hist_bucket} spread to {hist_rating} is currently "
                    f"at the {percentile:.0f}th percentile versus the selected {hist_window_label} history. "
                    f"Signal: **{signal}**."
                )

                hist_fig = px.line(
                    hist_daily,
                    x="trade_date",
                    y="display_spread_bps",
                    markers=True,
                    hover_data={
                        "spread_to_benchmark_bps": ":.1f",
                        "avg_yield": ":.2f",
                        "benchmark_yield": ":.2f",
                        "trade_count": ":,.0f",
                        "total_trade_amount": ":,.0f",
                        "benchmark_source": True,
                        "source_column": True,
                    },
                    title=f"{selected_issuer} {hist_bucket} Spread to {hist_rating} Benchmark Over Time",
                    labels={
                        "trade_date": "Trade Date",
                        "display_spread_bps": f"{display_label} (bps)",
                        "spread_to_benchmark_bps": "Raw Daily Spread (bps)",
                    },
                )
                hist_fig.add_hline(y=current_spread, line_dash="solid", opacity=0.65, annotation_text="Current")
                hist_fig.add_hline(y=hist_median, line_dash="dash", opacity=0.65, annotation_text="Median")
                hist_fig.add_hline(y=hist_p25, line_dash="dot", opacity=0.45, annotation_text="25th")
                hist_fig.add_hline(y=hist_p75, line_dash="dot", opacity=0.45, annotation_text="75th")
                hist_fig.update_layout(height=520, hovermode="x unified")
                st.plotly_chart(hist_fig, use_container_width=True)

                dist_col1, dist_col2 = st.columns([1.2, 1])
                with dist_col1:
                    hist_dist_fig = px.histogram(
                        hist_daily,
                        x="display_spread_bps",
                        nbins=25,
                        title=f"Historical Spread Distribution ({hist_window_label})",
                        labels={"display_spread_bps": f"{display_label} (bps)", "count": "Observation Count"},
                    )
                    hist_dist_fig.add_vline(x=current_spread, line_dash="solid", annotation_text="Current")
                    hist_dist_fig.add_vline(x=hist_median, line_dash="dash", annotation_text="Median")
                    hist_dist_fig.update_layout(height=430)
                    st.plotly_chart(hist_dist_fig, use_container_width=True)

                with dist_col2:
                    pct_table = pd.DataFrame(
                        [
                            {"Statistic": "Minimum", "Spread (bps)": hist_min},
                            {"Statistic": "10th Percentile", "Spread (bps)": hist_p10},
                            {"Statistic": "25th Percentile", "Spread (bps)": hist_p25},
                            {"Statistic": "Median", "Spread (bps)": hist_median},
                            {"Statistic": "75th Percentile", "Spread (bps)": hist_p75},
                            {"Statistic": "90th Percentile", "Spread (bps)": hist_p90},
                            {"Statistic": "Maximum", "Spread (bps)": hist_max},
                            {"Statistic": "Current", "Spread (bps)": current_spread},
                        ]
                    )
                    pct_table["Spread (bps)"] = pct_table["Spread (bps)"].round(2)
                    st.dataframe(pct_table, use_container_width=True, hide_index=True)

                with st.expander("Historical spread audit table", expanded=False):
                    audit_cols = [
                        "trade_date",
                        "spread_to_benchmark_bps",
                        "display_spread_bps",
                        "avg_yield",
                        "benchmark_yield",
                        "trade_count",
                        "total_trade_amount",
                        "benchmark_source",
                        "source_column",
                    ]
                    audit_hist = hist_daily[[c for c in audit_cols if c in hist_daily.columns]].copy()
                    for c in ["spread_to_benchmark_bps", "display_spread_bps", "avg_yield", "benchmark_yield"]:
                        if c in audit_hist.columns:
                            audit_hist[c] = pd.to_numeric(audit_hist[c], errors="coerce").round(2)
                    st.dataframe(
                        audit_hist.sort_values("trade_date", ascending=False).head(1000),
                        use_container_width=True,
                        hide_index=True,
                    )



section_anchor("dealer-proxy", "Bid / Ask & Dealer Behavior Proxy")
with st.expander("Methodology: dealer behavior proxy", expanded=False):
    st.markdown(
        """
This section is intentionally **data-dependent**. It becomes most useful when uploaded trade files include a buy/sell indicator, dealer side, customer side, or trade type.

**Realistic implementation:**

- If a usable `trade_type` / side field exists, the app classifies trades into:
    - **Customer Buy / Buy**
    - **Customer Sell / Sell**
    - **Other / Unknown**
- It estimates a simple imbalance proxy:

`Selling Imbalance = (Sell Amount - Buy Amount) / (Sell Amount + Buy Amount)`

**How to read it:**

- Positive value: more selling pressure / potential dealer inventory pressure.
- Negative value: more buying demand.
- Near zero: balanced flow.

**Important limitation:**

This is a **proxy**, not true dealer inventory. True dealer behavior would require richer MSRB/dealer-side fields, bid-wanted data, or inventory data.
        """
    )

dealer_trade_type_col = None
for candidate_col in ["trade_type", "side", "buy_sell", "customer_side", "dealer_side"]:
    if candidate_col in market_df.columns:
        dealer_trade_type_col = candidate_col
        break

if dealer_trade_type_col is None:
    st.info(
        "Dealer behavior proxy is unavailable because no trade side / trade type field was detected. "
        "Upload trade data with fields such as Trade Type, Side, Customer Buy/Sell, or Dealer Side to enable this module."
    )
else:
    dealer_col1, dealer_col2, dealer_col3 = st.columns([1, 1, 1])
    with dealer_col1:
        dealer_scope = st.selectbox(
            "Dealer Proxy Scope",
            ["Selected issuer", "All uploaded issuers"],
            index=0,
            key="dealer_proxy_scope",
        )
    with dealer_col2:
        dealer_bucket = st.selectbox(
            "Dealer Proxy Maturity Bucket",
            MATURITY_BUCKET_OPTIONS,
            index=0,
            key="dealer_proxy_bucket",
        )
    with dealer_col3:
        dealer_window = st.selectbox(
            "Dealer Proxy Window",
            ["Latest 30D", "Latest 60D", "Latest 90D", "All"],
            index=2,
            key="dealer_proxy_window",
        )

    dealer_df = market_df.copy()
    dealer_df["trade_date"] = pd.to_datetime(dealer_df["trade_date"], errors="coerce").dt.normalize()
    dealer_df["yield"] = pd.to_numeric(dealer_df["yield"], errors="coerce")
    if "trade_amount" in dealer_df.columns:
        dealer_df["trade_amount"] = pd.to_numeric(dealer_df["trade_amount"], errors="coerce").fillna(0)
    else:
        dealer_df["trade_amount"] = 0.0

    dealer_df = dealer_df.dropna(subset=["trade_date"])
    if dealer_scope == "Selected issuer":
        dealer_df = dealer_df[dealer_df["issuer"] == selected_issuer].copy()
    if dealer_bucket != "All" and "maturity_bucket" in dealer_df.columns:
        dealer_df = dealer_df[dealer_df["maturity_bucket"] == dealer_bucket].copy()

    dealer_days = {"Latest 30D": 30, "Latest 60D": 60, "Latest 90D": 90, "All": None}[dealer_window]
    if not dealer_df.empty and dealer_days is not None:
        dealer_latest = dealer_df["trade_date"].max()
        dealer_df = dealer_df[dealer_df["trade_date"] >= dealer_latest - pd.Timedelta(days=dealer_days)].copy()

    if dealer_df.empty:
        st.warning("No trades remain for the selected dealer-proxy filters.")
    else:
        def classify_trade_side(value: object) -> str:
            text_val = str(value).strip().lower()
            if not text_val or text_val in {"nan", "none"}:
                return "Unknown"
            # Common variants from trade exports. These are intentionally broad.
            if any(token in text_val for token in ["sell", "sold", "sld", "customer sell", "cust sell", "cs"]):
                return "Sell"
            if any(token in text_val for token in ["buy", "bought", "purchase", "customer buy", "cust buy", "cb"]):
                return "Buy"
            if text_val in {"s"}:
                return "Sell"
            if text_val in {"b"}:
                return "Buy"
            return "Other / Unknown"

        dealer_df["flow_side"] = dealer_df[dealer_trade_type_col].map(classify_trade_side)

        flow_summary = (
            dealer_df.groupby("flow_side", as_index=False)
            .agg(
                trade_count=("trade_date", "count"),
                total_trade_amount=("trade_amount", "sum"),
                avg_yield=("yield", "mean"),
            )
        )

        buy_amount = flow_summary.loc[flow_summary["flow_side"] == "Buy", "total_trade_amount"].sum()
        sell_amount = flow_summary.loc[flow_summary["flow_side"] == "Sell", "total_trade_amount"].sum()
        buy_count = flow_summary.loc[flow_summary["flow_side"] == "Buy", "trade_count"].sum()
        sell_count = flow_summary.loc[flow_summary["flow_side"] == "Sell", "trade_count"].sum()
        denom_amount = buy_amount + sell_amount
        denom_count = buy_count + sell_count
        amount_imbalance = (sell_amount - buy_amount) / denom_amount if denom_amount > 0 else pd.NA
        count_imbalance = (sell_count - buy_count) / denom_count if denom_count > 0 else pd.NA

        dp1, dp2, dp3, dp4 = st.columns(4)
        dp1.metric("Buy Amount", f"{buy_amount:,.0f}")
        dp2.metric("Sell Amount", f"{sell_amount:,.0f}")
        dp3.metric("Amount Imbalance", "N/A" if pd.isna(amount_imbalance) else f"{amount_imbalance:+.1%}")
        dp4.metric("Count Imbalance", "N/A" if pd.isna(count_imbalance) else f"{count_imbalance:+.1%}")

        if pd.notna(amount_imbalance):
            if amount_imbalance >= 0.25:
                st.info("Read-through: flow appears sell-heavy, which may indicate customer selling pressure or dealer balance-sheet pressure.")
            elif amount_imbalance <= -0.25:
                st.info("Read-through: flow appears buy-heavy, which may indicate stronger demand or dealer distribution.")
            else:
                st.info("Read-through: buy/sell flow appears relatively balanced in the selected window.")
        else:
            st.info("Read-through: buy/sell imbalance cannot be calculated because buy/sell amount data is unavailable or zero.")

        flow_fig = px.bar(
            flow_summary,
            x="flow_side",
            y="total_trade_amount",
            hover_data={"trade_count": ":,.0f", "avg_yield": ":.2f"},
            title="Trade Amount by Flow Side",
            labels={
                "flow_side": "Flow Side",
                "total_trade_amount": "Total Trade Amount",
                "trade_count": "Trade Count",
                "avg_yield": "Average Yield",
            },
        )
        flow_fig.update_layout(height=430)
        st.plotly_chart(flow_fig, use_container_width=True)

        dealer_daily = (
            dealer_df.groupby(["trade_date", "flow_side"], as_index=False)
            .agg(
                trade_count=("trade_date", "count"),
                total_trade_amount=("trade_amount", "sum"),
            )
        )
        daily_pivot = (
            dealer_daily.pivot_table(
                index="trade_date",
                columns="flow_side",
                values="total_trade_amount",
                aggfunc="sum",
                fill_value=0,
                observed=False,
            )
            .reset_index()
        )
        if "Buy" not in daily_pivot.columns:
            daily_pivot["Buy"] = 0
        if "Sell" not in daily_pivot.columns:
            daily_pivot["Sell"] = 0
        daily_pivot["net_selling_pressure"] = daily_pivot["Sell"] - daily_pivot["Buy"]

        pressure_fig = px.bar(
            daily_pivot,
            x="trade_date",
            y="net_selling_pressure",
            title="Daily Net Selling Pressure Proxy",
            labels={
                "trade_date": "Trade Date",
                "net_selling_pressure": "Sell Amount - Buy Amount",
            },
        )
        pressure_fig.add_hline(y=0, line_dash="dash", opacity=0.45)
        pressure_fig.update_layout(height=430)
        st.plotly_chart(pressure_fig, use_container_width=True)

        with st.expander("Dealer proxy audit table", expanded=False):
            audit_cols = [
                "trade_date",
                "issuer",
                "cusip",
                "maturity_bucket",
                dealer_trade_type_col,
                "flow_side",
                "yield",
                "price",
                "trade_amount",
                "trade_type",
            ]

            # Remove missing columns and duplicate column names while preserving order.
            # This matters when dealer_trade_type_col itself is "trade_type".
            audit_cols = list(dict.fromkeys([c for c in audit_cols if c in dealer_df.columns]))

            dealer_audit_display = dealer_df[audit_cols].copy()
            dealer_audit_display = dealer_audit_display.loc[:, ~dealer_audit_display.columns.duplicated()].copy()

            st.caption(f"Side classification source column: `{dealer_trade_type_col}`.")
            st.dataframe(
                dealer_audit_display.sort_values("trade_date", ascending=False).head(5000),
                use_container_width=True,
                hide_index=True,
            )



section_anchor("security-screener", "Security Screener — Top Relative Value Candidates")
with st.expander("Methodology: security screener", expanded=False):
    st.markdown(
        """
This section turns the dashboard into a practical **find me bonds** workflow.

**Goal:**

Screen uploaded bonds/trades for securities that are both relatively cheap and sufficiently liquid.

**Core fields used:**

- Sector / issuer / maturity bucket
- Spread to benchmark
- Liquidity score
- Trade count
- Total trade amount
- Days since last trade

**Core spread calculation:**

`Spread to Benchmark = (Average CUSIP Yield - Benchmark Yield) × 100`

**Important limitation:**

This is a screening tool. It does not replace credit review, call analysis, structure review, tax status review, or PM/trader judgment.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable security screener spread calculations.")
else:
    screen_col1, screen_col2, screen_col3, screen_col4 = st.columns([1, 1, 1, 1])
    with screen_col1:
        screen_sector_options = ["All"]
        if "sector" in market_df.columns:
            screen_sector_options += sorted(market_df["sector"].dropna().astype(str).unique().tolist())
        screen_sector = st.selectbox("Screener Sector", screen_sector_options, index=0, key="screen_sector")
    with screen_col2:
        screen_bucket = st.selectbox(
            "Screener Maturity Bucket",
            MATURITY_BUCKET_OPTIONS,
            index=0,
            key="screen_bucket",
        )
    with screen_col3:
        screen_rating = st.selectbox(
            "Screener Benchmark",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="screen_rating",
        )
    with screen_col4:
        screen_window = st.selectbox(
            "Screener Lookback",
            ["Latest 30D", "Latest 60D", "Latest 90D", "All"],
            index=2,
            key="screen_window",
        )

    filt_col1, filt_col2, filt_col3, filt_col4 = st.columns([1, 1, 1, 1])
    with filt_col1:
        min_spread = st.number_input("Min Spread (bps)", value=40.0, step=5.0, key="min_screen_spread")
    with filt_col2:
        min_liquidity = st.number_input("Min Liquidity Score", value=50.0, min_value=0.0, max_value=100.0, step=5.0, key="min_screen_liq")
    with filt_col3:
        min_trades_screen = st.number_input("Min Trades", value=2, min_value=1, max_value=100, step=1, key="min_screen_trades")
    with filt_col4:
        min_trade_amount = st.number_input("Min Total Trade Amount", value=0.0, step=100000.0, key="min_screen_amount")

    screener_df = market_df.copy()
    screener_df["trade_date"] = pd.to_datetime(screener_df["trade_date"], errors="coerce").dt.normalize()
    screener_df["yield"] = pd.to_numeric(screener_df["yield"], errors="coerce")
    if "trade_amount" in screener_df.columns:
        screener_df["trade_amount"] = pd.to_numeric(screener_df["trade_amount"], errors="coerce").fillna(0)
    else:
        screener_df["trade_amount"] = 0.0
    if "price" in screener_df.columns:
        screener_df["price"] = pd.to_numeric(screener_df["price"], errors="coerce")
    else:
        screener_df["price"] = pd.NA

    screener_df = screener_df.dropna(subset=["trade_date", "yield", "cusip"])
    if screen_sector != "All" and "sector" in screener_df.columns:
        screener_df = screener_df[screener_df["sector"].astype(str) == str(screen_sector)].copy()
    if screen_bucket != "All" and "maturity_bucket" in screener_df.columns:
        screener_df = screener_df[screener_df["maturity_bucket"] == screen_bucket].copy()

    screen_days = {"Latest 30D": 30, "Latest 60D": 60, "Latest 90D": 90, "All": None}[screen_window]
    if not screener_df.empty and screen_days is not None:
        latest_screen_date = screener_df["trade_date"].max()
        screener_df = screener_df[screener_df["trade_date"] >= latest_screen_date - pd.Timedelta(days=screen_days)].copy()

    if screener_df.empty:
        st.warning("No securities remain after the selected screener universe filters.")
    else:
        screen_summary = (
            screener_df.groupby("cusip", dropna=False)
            .agg(
                issuer=("issuer", "first"),
                sector=("sector", "first") if "sector" in screener_df.columns else ("issuer", "first"),
                maturity_bucket=("maturity_bucket", "first") if "maturity_bucket" in screener_df.columns else ("issuer", "first"),
                maturity=("maturity_bond", "first") if "maturity_bond" in screener_df.columns else ("trade_date", "max"),
                coupon=("coupon_bond", "first") if "coupon_bond" in screener_df.columns else ("yield", "count"),
                call_date=("call_date", "first") if "call_date" in screener_df.columns else ("trade_date", "max"),
                avg_yield=("yield", "mean"),
                latest_yield=("yield", "last"),
                avg_price=("price", "mean"),
                trade_count=("trade_date", "count"),
                latest_trade=("trade_date", "max"),
                first_trade=("trade_date", "min"),
                total_trade_amount=("trade_amount", "sum"),
                avg_trade_amount=("trade_amount", "mean"),
                outstanding_amount=("outstanding_amount", "first") if "outstanding_amount" in screener_df.columns else ("trade_amount", "sum"),
            )
            .reset_index()
        )

        # Latest benchmark by maturity bucket.
        screen_date_col = _detect_mmd_date_column(mmd_df)
        if screen_date_col is None:
            st.warning("Security screener cannot calculate spreads because the benchmark file has no usable date column.")
        else:
            screen_mmd = mmd_df.copy()
            screen_mmd[screen_date_col] = pd.to_datetime(screen_mmd[screen_date_col], errors="coerce")
            screen_mmd = screen_mmd.dropna(subset=[screen_date_col])
            latest_trade_for_screen = screener_df["trade_date"].max()
            screen_mmd = screen_mmd[screen_mmd[screen_date_col].dt.normalize() <= latest_trade_for_screen].sort_values(screen_date_col)

            if screen_mmd.empty:
                st.warning("No benchmark curve observation was available on or before the latest screener trade date.")
            else:
                screen_latest_mmd = screen_mmd.iloc[[-1]].copy()
                screen_benchmark_date = screen_latest_mmd[screen_date_col].iloc[0]

                bench_rows = []
                for bucket in MATURITY_BUCKET_ORDER:
                    tenor = MMD_BUCKET_MAP.get(bucket, "10Y")
                    y, meta = get_benchmark_curve(screen_latest_mmd, tenor, screen_rating)
                    if y is not None and pd.notna(y.iloc[0]):
                        bench_rows.append(
                            {
                                "maturity_bucket": bucket,
                                "mmd_tenor": tenor,
                                "benchmark_yield": float(y.iloc[0]),
                                "benchmark_source": meta.get("benchmark_source"),
                                "source_column": meta.get("source_column"),
                                "rating_spread_bps": meta.get("rating_spread_bps"),
                            }
                        )
                screen_bench = pd.DataFrame(bench_rows)

                if screen_bench.empty:
                    st.warning("Selected benchmark curve could not be built for screener.")
                else:
                    screen_summary["maturity_bucket"] = screen_summary["maturity_bucket"].astype(str)
                    screen_summary = screen_summary.merge(screen_bench, on="maturity_bucket", how="left")
                    screen_summary["spread_to_benchmark_bps"] = (
                        screen_summary["avg_yield"] - screen_summary["benchmark_yield"]
                    ) * 100

                    # Liquidity score proxy.
                    today_screen = pd.Timestamp.today().normalize()
                    screen_summary["days_since_last_trade"] = (today_screen - screen_summary["latest_trade"]).dt.days
                    screen_summary["liquidity_score"] = (
                        screen_summary["trade_count"].rank(pct=True) * 35
                        + screen_summary["total_trade_amount"].rank(pct=True) * 35
                        + (1 - screen_summary["days_since_last_trade"].rank(pct=True)) * 30
                    )
                    screen_summary["turnover_ratio"] = (
                        screen_summary["total_trade_amount"]
                        / pd.to_numeric(screen_summary["outstanding_amount"], errors="coerce").replace({0: pd.NA})
                    )

                    # Apply screen filters.
                    candidates = screen_summary[
                        (screen_summary["spread_to_benchmark_bps"] >= min_spread)
                        & (screen_summary["liquidity_score"] >= min_liquidity)
                        & (screen_summary["trade_count"] >= min_trades_screen)
                        & (screen_summary["total_trade_amount"] >= min_trade_amount)
                    ].copy()

                    candidates["rv_score"] = (
                        candidates["spread_to_benchmark_bps"].rank(pct=True) * 45
                        + candidates["liquidity_score"].rank(pct=True) * 35
                        + candidates["trade_count"].rank(pct=True) * 10
                        + candidates["total_trade_amount"].rank(pct=True) * 10
                    )

                    candidates = candidates.sort_values("rv_score", ascending=False, na_position="last")

                    s1, s2, s3, s4 = st.columns(4)
                    s1.metric("Candidates", f"{len(candidates):,}")
                    s2.metric("Benchmark", f"{screen_rating}")
                    s3.metric("Benchmark Date", screen_benchmark_date.strftime("%Y-%m-%d"))
                    s4.metric("Universe CUSIPs", f"{len(screen_summary):,}")

                    if candidates.empty:
                        st.info(
                            "No securities met the current screener filters. Try lowering minimum spread, liquidity score, or trade count."
                        )
                    else:
                        top_candidate = candidates.iloc[0]
                        st.info(
                            f"Top candidate by RV score: CUSIP {top_candidate['cusip']} "
                            f"({top_candidate['issuer']}) screens at {top_candidate['spread_to_benchmark_bps']:+.1f} bp "
                            f"to {screen_rating}, liquidity score {top_candidate['liquidity_score']:.1f}, "
                            f"and {int(top_candidate['trade_count'])} trades in the selected window."
                        )

                        display_cols = [
                            "rv_score",
                            "cusip",
                            "issuer",
                            "sector",
                            "maturity_bucket",
                            "maturity",
                            "coupon",
                            "call_date",
                            "avg_yield",
                            "benchmark_yield",
                            "spread_to_benchmark_bps",
                            "liquidity_score",
                            "trade_count",
                            "total_trade_amount",
                            "avg_trade_amount",
                            "days_since_last_trade",
                            "avg_price",
                            "outstanding_amount",
                            "turnover_ratio",
                            "benchmark_source",
                            "source_column",
                        ]
                        display_candidates = candidates[[c for c in display_cols if c in candidates.columns]].copy()
                        for c in ["rv_score", "avg_yield", "benchmark_yield", "spread_to_benchmark_bps", "liquidity_score", "avg_price", "turnover_ratio"]:
                            if c in display_candidates.columns:
                                display_candidates[c] = pd.to_numeric(display_candidates[c], errors="coerce").round(2)

                        st.dataframe(
                            display_candidates.head(1000),
                            use_container_width=True,
                            hide_index=True,
                            height=480,
                        )

                        screener_fig = px.scatter(
                            candidates,
                            x="liquidity_score",
                            y="spread_to_benchmark_bps",
                            size="total_trade_amount",
                            size_max=38,
                            color="maturity_bucket",
                            hover_name="cusip",
                            hover_data=[
                                c for c in [
                                    "issuer",
                                    "sector",
                                    "avg_yield",
                                    "benchmark_yield",
                                    "trade_count",
                                    "days_since_last_trade",
                                    "rv_score",
                                ] if c in candidates.columns
                            ],
                            title="Top Relative Value Candidates",
                            labels={
                                "liquidity_score": "Liquidity Score",
                                "spread_to_benchmark_bps": "Spread to Benchmark (bps)",
                                "maturity_bucket": "Maturity Bucket",
                            },
                        )
                        screener_fig.add_vline(x=min_liquidity, line_dash="dash", opacity=0.45)
                        screener_fig.add_hline(y=min_spread, line_dash="dash", opacity=0.45)
                        screener_fig.update_layout(height=520, hovermode="closest")
                        st.plotly_chart(screener_fig, use_container_width=True)

                        csv_candidates = candidates.to_csv(index=False).encode("utf-8")
                        st.download_button(
                            label="Download Top Relative Value Candidates CSV",
                            data=csv_candidates,
                            file_name="top_relative_value_candidates.csv",
                            mime="text/csv",
                        )

                    with st.expander("Screener universe audit table", expanded=False):
                        audit_cols = [
                            "cusip",
                            "issuer",
                            "sector",
                            "maturity_bucket",
                            "avg_yield",
                            "benchmark_yield",
                            "spread_to_benchmark_bps",
                            "liquidity_score",
                            "trade_count",
                            "total_trade_amount",
                            "latest_trade",
                            "benchmark_source",
                            "source_column",
                        ]
                        audit_screen = screen_summary[[c for c in audit_cols if c in screen_summary.columns]].copy()
                        for c in ["avg_yield", "benchmark_yield", "spread_to_benchmark_bps", "liquidity_score"]:
                            if c in audit_screen.columns:
                                audit_screen[c] = pd.to_numeric(audit_screen[c], errors="coerce").round(2)
                        st.dataframe(audit_screen.head(5000), use_container_width=True, hide_index=True)




section_anchor("watchlist", "Watchlist / Saved Candidates")
with st.expander("Methodology: watchlist / saved candidates", expanded=False):
    st.markdown(
        """
This section lets users save CUSIPs for later review during the current session.

**Why it matters:**

A trading workflow often moves from screening → shortlist → detailed review. The watchlist keeps promising CUSIPs visible without forcing users to re-filter every time.

**Current implementation:**

- Uses Streamlit session state, so it persists while the app session is active.
- Users can add CUSIPs from the selected issuer or all uploaded data.
- Users can download the watchlist as CSV.
- This is not a database-backed permanent watchlist yet; that would be a future production feature.
        """
    )

if "watchlist_cusips" not in st.session_state:
    st.session_state["watchlist_cusips"] = []

watch_col1, watch_col2 = st.columns([1, 1])
with watch_col1:
    watch_scope = st.selectbox(
        "Watchlist Add Scope",
        ["Selected issuer", "All uploaded issuers"],
        index=0,
        key="watchlist_scope",
    )
with watch_col2:
    watch_universe = market_df.copy()
    if watch_scope == "Selected issuer":
        watch_universe = watch_universe[watch_universe["issuer"] == selected_issuer].copy()
    watch_options = sorted(watch_universe["cusip"].dropna().astype(str).unique().tolist()) if "cusip" in watch_universe.columns else []
    cusips_to_add = st.multiselect(
        "Add CUSIPs to Watchlist",
        watch_options,
        default=[],
        key="watchlist_add_cusips",
    )

add_col, clear_col = st.columns([1, 1])
with add_col:
    if st.button("Add selected CUSIPs", key="add_watchlist_cusips"):
        current = set(st.session_state.get("watchlist_cusips", []))
        current.update(cusips_to_add)
        st.session_state["watchlist_cusips"] = sorted(current)
        st.success(f"Added {len(cusips_to_add):,} CUSIP(s) to watchlist.")
with clear_col:
    if st.button("Clear watchlist", key="clear_watchlist_cusips"):
        st.session_state["watchlist_cusips"] = []
        st.info("Watchlist cleared.")

watchlist_cusips = st.session_state.get("watchlist_cusips", [])
if not watchlist_cusips:
    st.info("No CUSIPs saved yet. Add candidates from the selector above.")
else:
    watchlist_df = market_df[market_df["cusip"].astype(str).isin(watchlist_cusips)].copy()
    if watchlist_df.empty:
        st.warning("Saved CUSIPs were not found in the current uploaded dataset.")
    else:
        watchlist_df["trade_date"] = pd.to_datetime(watchlist_df["trade_date"], errors="coerce")
        if "trade_amount" in watchlist_df.columns:
            watchlist_df["trade_amount"] = pd.to_numeric(watchlist_df["trade_amount"], errors="coerce").fillna(0)
        else:
            watchlist_df["trade_amount"] = 0.0
        watchlist_summary = (
            watchlist_df.groupby("cusip", dropna=False)
            .agg(
                issuer=("issuer", "first"),
                sector=("sector", "first") if "sector" in watchlist_df.columns else ("issuer", "first"),
                maturity_bucket=("maturity_bucket", "first") if "maturity_bucket" in watchlist_df.columns else ("issuer", "first"),
                maturity=("maturity_bond", "first") if "maturity_bond" in watchlist_df.columns else ("trade_date", "max"),
                coupon=("coupon_bond", "first") if "coupon_bond" in watchlist_df.columns else ("trade_date", "count"),
                avg_yield=("yield", "mean"),
                avg_price=("price", "mean") if "price" in watchlist_df.columns else ("yield", "mean"),
                trade_count=("trade_date", "count"),
                total_trade_amount=("trade_amount", "sum"),
                latest_trade=("trade_date", "max"),
            )
            .reset_index()
        )
        st.metric("Saved CUSIPs", f"{len(watchlist_summary):,}")
        st.dataframe(watchlist_summary, use_container_width=True, hide_index=True)

        st.download_button(
            label="Download Watchlist CSV",
            data=watchlist_summary.to_csv(index=False).encode("utf-8"),
            file_name="watchlist_saved_candidates.csv",
            mime="text/csv",
        )


section_anchor("recommendation-engine", "Trade Recommendation Narrative Engine")
with st.expander("Methodology: rule-based recommendation narrative", expanded=False):
    st.markdown(
        """
This section generates a **rule-driven market commentary** from the analytics already shown in the dashboard.

It is intentionally not an AI black box. Each phrase is triggered by a transparent data rule.

**Signals used when available:**

- **Spread movement:** whether the selected maturity bucket widened/tightened over the chosen lookback window.
- **Historical percentile:** whether the current spread is wide/tight versus its own history.
- **Liquidity:** whether the bucket remains tradable based on trade count, amount traded, and recency.
- **Peer comparison:** whether the selected issuer screens wide/tight versus uploaded peers.
- **Dealer flow proxy:** whether observed trade type/side suggests buy-heavy or sell-heavy activity.

**Example rule mapping:**

| Data trigger | Narrative phrase |
|---|---|
| Spread change > +15 bp | widened materially |
| Historical percentile > 90th | near the upper end of its historical range |
| Liquidity score > 70 | while maintaining above-average liquidity |
| Wider than peer median by > 10 bp | screens wide versus uploaded peers |
| Sell imbalance > 25% | flow appears sell-heavy |

**Important limitation:**

This is a screening commentary, not an investment recommendation. It should be reviewed alongside credit fundamentals, call structure, tax status, and actual executable market levels.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable rule-based recommendation commentary.")
else:
    rec_col1, rec_col2, rec_col3 = st.columns([1, 1, 1])
    with rec_col1:
        rec_bucket = st.selectbox(
            "Narrative Maturity Bucket",
            MATURITY_BUCKET_ORDER,
            index=3,
            key="rec_bucket",
        )
    with rec_col2:
        rec_rating = st.selectbox(
            "Narrative Benchmark",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="rec_rating",
        )
    with rec_col3:
        rec_window_label = st.selectbox(
            "Narrative Lookback",
            ["1W", "1M", "3M", "6M", "1Y"],
            index=1,
            key="rec_window",
        )

    rec_window_days = {"1W": 7, "1M": 30, "3M": 90, "6M": 180, "1Y": 365}[rec_window_label]

    # -----------------------------
    # Signal 1: spread movement + current spread
    # -----------------------------
    rec_obs = build_spread_observations(
        market_df=market_df,
        mmd_df=mmd_df,
        issuer=selected_issuer,
        rating=rec_rating,
    )

    narrative_lines = []
    evidence_rows = []

    current_spread = pd.NA
    spread_change = pd.NA
    historical_percentile = pd.NA
    liquidity_score = pd.NA
    peer_gap = pd.NA
    amount_imbalance = pd.NA

    if rec_obs.empty:
        st.warning("No overlapping issuer trade dates and benchmark dates were found for narrative generation.")
    else:
        rec_obs = rec_obs.copy()
        rec_obs["trade_date"] = pd.to_datetime(rec_obs["trade_date"], errors="coerce").dt.normalize()
        rec_obs = rec_obs[
            (rec_obs["maturity_bucket"] == rec_bucket)
            & rec_obs["spread_to_benchmark_bps"].notna()
        ].sort_values("trade_date")

        if rec_obs.empty:
            st.warning(f"No spread observations found for {selected_issuer} / {rec_bucket}.")
        else:
            latest_row = rec_obs.iloc[-1]
            latest_date = latest_row["trade_date"]
            current_spread = float(latest_row["spread_to_benchmark_bps"])
            target_date = latest_date - pd.Timedelta(days=rec_window_days)
            hist_candidates = rec_obs[rec_obs["trade_date"] <= target_date]

            if not hist_candidates.empty:
                hist_row = hist_candidates.iloc[-1]
                historical_spread = float(hist_row["spread_to_benchmark_bps"])
                spread_change = current_spread - historical_spread
            else:
                historical_spread = pd.NA

            # Historical percentile uses up to 1Y of history when available.
            hist_1y = rec_obs[rec_obs["trade_date"] >= latest_date - pd.Timedelta(days=365)].copy()
            hist_values = pd.to_numeric(hist_1y["spread_to_benchmark_bps"], errors="coerce").dropna()
            if len(hist_values) >= 2:
                historical_percentile = float((hist_values <= current_spread).mean() * 100)

            evidence_rows.append({
                "Signal": "Current spread",
                "Value": f"{current_spread:+.1f} bp",
                "Rule / Source": f"Latest {rec_bucket} spread to {rec_rating}",
            })
            if pd.notna(spread_change):
                evidence_rows.append({
                    "Signal": f"{rec_window_label} spread movement",
                    "Value": f"{spread_change:+.1f} bp",
                    "Rule / Source": "Latest spread minus historical spread at/before lookback target date",
                })
            if pd.notna(historical_percentile):
                evidence_rows.append({
                    "Signal": "Historical percentile",
                    "Value": f"{historical_percentile:.0f}th",
                    "Rule / Source": "Current spread percentile versus latest 1Y spread observations",
                })

            # Movement phrase.
            if pd.notna(spread_change):
                if spread_change >= 15:
                    movement_phrase = f"{rec_bucket} widened materially by {spread_change:+.1f} bp over {rec_window_label}"
                elif spread_change >= 5:
                    movement_phrase = f"{rec_bucket} widened modestly by {spread_change:+.1f} bp over {rec_window_label}"
                elif spread_change <= -15:
                    movement_phrase = f"{rec_bucket} tightened materially by {spread_change:+.1f} bp over {rec_window_label}"
                elif spread_change <= -5:
                    movement_phrase = f"{rec_bucket} tightened modestly by {spread_change:+.1f} bp over {rec_window_label}"
                else:
                    movement_phrase = f"{rec_bucket} was broadly stable over {rec_window_label}"
                narrative_lines.append(f"{selected_issuer} {movement_phrase} versus the {rec_rating} benchmark.")
            else:
                narrative_lines.append(
                    f"{selected_issuer} {rec_bucket} currently screens at {current_spread:+.1f} bp versus the {rec_rating} benchmark."
                )

            # Historical phrase.
            if pd.notna(historical_percentile):
                if historical_percentile >= 90:
                    narrative_lines.append(
                        f"Current spread sits near the upper end of its recent historical range ({historical_percentile:.0f}th percentile)."
                    )
                elif historical_percentile >= 75:
                    narrative_lines.append(
                        f"Current spread screens wide versus recent history ({historical_percentile:.0f}th percentile)."
                    )
                elif historical_percentile <= 10:
                    narrative_lines.append(
                        f"Current spread sits near the tight end of its recent historical range ({historical_percentile:.0f}th percentile)."
                    )
                elif historical_percentile <= 25:
                    narrative_lines.append(
                        f"Current spread screens tight versus recent history ({historical_percentile:.0f}th percentile)."
                    )
                else:
                    narrative_lines.append(
                        f"Current spread is close to its normal recent historical range ({historical_percentile:.0f}th percentile)."
                    )

    # -----------------------------
    # Signal 2: liquidity proxy
    # -----------------------------
    rec_trades = market_df[
        (market_df["issuer"] == selected_issuer)
        & (market_df["maturity_bucket"] == rec_bucket)
    ].copy()

    if not rec_trades.empty:
        rec_trades["trade_date"] = pd.to_datetime(rec_trades["trade_date"], errors="coerce").dt.normalize()
        if "trade_amount" in rec_trades.columns:
            rec_trades["trade_amount"] = pd.to_numeric(rec_trades["trade_amount"], errors="coerce").fillna(0)
        else:
            rec_trades["trade_amount"] = 0.0

        rec_trades = rec_trades.dropna(subset=["trade_date"])
        if not rec_trades.empty:
            latest_trade_date = rec_trades["trade_date"].max()
            rec_trade_window = rec_trades[rec_trades["trade_date"] >= latest_trade_date - pd.Timedelta(days=rec_window_days)].copy()

            if not rec_trade_window.empty:
                trade_count = len(rec_trade_window)
                total_trade_amount = rec_trade_window["trade_amount"].sum()
                days_since_last = (pd.Timestamp.today().normalize() - latest_trade_date).days

                # Simple bounded liquidity proxy.
                trade_count_score = min(trade_count / 10, 1) * 35
                amount_score = min(total_trade_amount / 5_000_000, 1) * 35
                recency_score = max(0, 1 - min(days_since_last / 180, 1)) * 30
                liquidity_score = trade_count_score + amount_score + recency_score

                evidence_rows.append({
                    "Signal": "Liquidity score",
                    "Value": f"{liquidity_score:.1f}",
                    "Rule / Source": "35% trade count + 35% par traded + 30% recency proxy",
                })
                evidence_rows.append({
                    "Signal": "Window trade count",
                    "Value": f"{trade_count:,}",
                    "Rule / Source": f"Trades in selected {rec_window_label} window",
                })
                evidence_rows.append({
                    "Signal": "Window par traded",
                    "Value": f"{total_trade_amount:,.0f}",
                    "Rule / Source": f"Total par traded in selected {rec_window_label} window",
                })

                if liquidity_score >= 70:
                    narrative_lines.append("The bucket maintains above-average liquidity based on recent trade count, par traded, and recency.")
                elif liquidity_score >= 45:
                    narrative_lines.append("Liquidity appears moderate; execution quality should still be checked at the CUSIP level.")
                else:
                    narrative_lines.append("Liquidity appears limited, so any apparent cheapness may include a meaningful liquidity premium.")

    # -----------------------------
    # Signal 3: peer comparison proxy
    # -----------------------------
    try:
        if "sector" in market_df.columns and selected_sector and selected_sector != "Unknown":
            peer_universe = market_df[
                market_df["sector"].astype(str) == str(selected_sector)
            ].copy()
        else:
            peer_universe = market_df.copy()

        peer_universe = peer_universe[
            (peer_universe["issuer"].astype(str) != str(selected_issuer))
            & (peer_universe["maturity_bucket"] == rec_bucket)
        ].copy()

        if not peer_universe.empty and pd.notna(current_spread):
            peer_universe["trade_date"] = pd.to_datetime(peer_universe["trade_date"], errors="coerce").dt.normalize()
            peer_universe["yield"] = pd.to_numeric(peer_universe["yield"], errors="coerce")
            peer_universe = peer_universe.dropna(subset=["trade_date", "yield"])
            if not peer_universe.empty:
                latest_peer_date = peer_universe["trade_date"].max()
                peer_universe = peer_universe[
                    peer_universe["trade_date"] >= latest_peer_date - pd.Timedelta(days=rec_window_days)
                ].copy()

                if not peer_universe.empty:
                    date_col = _detect_mmd_date_column(mmd_df)
                    if date_col is not None:
                        peer_mmd = mmd_df.copy()
                        peer_mmd[date_col] = pd.to_datetime(peer_mmd[date_col], errors="coerce")
                        peer_mmd = peer_mmd.dropna(subset=[date_col])
                        peer_mmd = peer_mmd[peer_mmd[date_col].dt.normalize() <= latest_peer_date].sort_values(date_col)
                        if not peer_mmd.empty:
                            peer_latest_mmd = peer_mmd.iloc[[-1]].copy()
                            peer_tenor = MMD_BUCKET_MAP.get(rec_bucket, "10Y")
                            peer_bench, _peer_meta = get_benchmark_curve(peer_latest_mmd, peer_tenor, rec_rating)
                            if peer_bench is not None and pd.notna(peer_bench.iloc[0]):
                                peer_benchmark_yield = float(peer_bench.iloc[0])
                                peer_summary_for_narrative = (
                                    peer_universe.groupby("issuer", as_index=False)
                                    .agg(avg_yield=("yield", "mean"), trade_count=("yield", "count"))
                                )
                                peer_summary_for_narrative["peer_spread_bps"] = (
                                    peer_summary_for_narrative["avg_yield"] - peer_benchmark_yield
                                ) * 100
                                peer_median = peer_summary_for_narrative["peer_spread_bps"].median()
                                peer_gap = current_spread - peer_median

                                evidence_rows.append({
                                    "Signal": "Peer gap",
                                    "Value": f"{peer_gap:+.1f} bp",
                                    "Rule / Source": f"Selected issuer spread minus uploaded peer median in {rec_bucket}",
                                })

                                if peer_gap >= 10:
                                    narrative_lines.append(
                                        f"The issuer screens {peer_gap:+.1f} bp wide to the uploaded peer median in the same bucket."
                                    )
                                elif peer_gap <= -10:
                                    narrative_lines.append(
                                        f"The issuer screens {abs(peer_gap):.1f} bp tight to the uploaded peer median in the same bucket."
                                    )
                                else:
                                    narrative_lines.append(
                                        "The issuer screens broadly in line with the uploaded peer median."
                                    )
    except Exception as exc:
        evidence_rows.append({
            "Signal": "Peer comparison",
            "Value": "Unavailable",
            "Rule / Source": f"Peer narrative skipped: {exc}",
        })

    # -----------------------------
    # Signal 4: dealer flow proxy
    # -----------------------------
    dealer_col = None
    for candidate_col in ["trade_type", "side", "buy_sell", "customer_side", "dealer_side"]:
        if candidate_col in market_df.columns:
            dealer_col = candidate_col
            break

    if dealer_col is not None:
        try:
            flow_df = market_df[
                (market_df["issuer"] == selected_issuer)
                & (market_df["maturity_bucket"] == rec_bucket)
            ].copy()
            flow_df["trade_date"] = pd.to_datetime(flow_df["trade_date"], errors="coerce").dt.normalize()
            if "trade_amount" in flow_df.columns:
                flow_df["trade_amount"] = pd.to_numeric(flow_df["trade_amount"], errors="coerce").fillna(0)
            else:
                flow_df["trade_amount"] = 0.0
            flow_df = flow_df.dropna(subset=["trade_date"])
            if not flow_df.empty:
                latest_flow_date = flow_df["trade_date"].max()
                flow_df = flow_df[flow_df["trade_date"] >= latest_flow_date - pd.Timedelta(days=rec_window_days)].copy()

            def rec_classify_side(value: object) -> str:
                t = str(value).strip().lower()
                if any(token in t for token in ["sell", "sold", "sld", "customer sell", "cust sell", "cs"]):
                    return "Sell"
                if any(token in t for token in ["buy", "bought", "purchase", "customer buy", "cust buy", "cb"]):
                    return "Buy"
                if t == "s":
                    return "Sell"
                if t == "b":
                    return "Buy"
                return "Other"

            if not flow_df.empty:
                flow_df["flow_side"] = flow_df[dealer_col].map(rec_classify_side)
                buy_amt = flow_df.loc[flow_df["flow_side"] == "Buy", "trade_amount"].sum()
                sell_amt = flow_df.loc[flow_df["flow_side"] == "Sell", "trade_amount"].sum()
                denom = buy_amt + sell_amt
                if denom > 0:
                    amount_imbalance = (sell_amt - buy_amt) / denom
                    evidence_rows.append({
                        "Signal": "Dealer / flow imbalance",
                        "Value": f"{amount_imbalance:+.1%}",
                        "Rule / Source": f"(Sell amount - Buy amount) / total classified amount from `{dealer_col}`",
                    })
                    if amount_imbalance >= 0.25:
                        narrative_lines.append("Classified flow appears sell-heavy, which may indicate customer selling pressure.")
                    elif amount_imbalance <= -0.25:
                        narrative_lines.append("Classified flow appears buy-heavy, which may indicate stronger demand.")
        except Exception as exc:
            evidence_rows.append({
                "Signal": "Dealer / flow imbalance",
                "Value": "Unavailable",
                "Rule / Source": f"Flow narrative skipped: {exc}",
            })

    # -----------------------------
    # Recommendation label
    # -----------------------------
    score = 0
    if pd.notna(spread_change) and spread_change >= 15:
        score += 1
    if pd.notna(historical_percentile) and historical_percentile >= 75:
        score += 1
    if pd.notna(liquidity_score) and liquidity_score >= 60:
        score += 1
    if pd.notna(peer_gap) and peer_gap >= 10:
        score += 1
    if pd.notna(amount_imbalance) and amount_imbalance >= 0.25:
        score += 0.5

    if score >= 3:
        recommendation_label = "Potential Relative Value Candidate"
    elif score >= 2:
        recommendation_label = "Watchlist Candidate"
    elif score <= 0.5 and pd.notna(historical_percentile) and historical_percentile <= 25:
        recommendation_label = "Potentially Rich / Lower Priority"
    else:
        recommendation_label = "Neutral / Needs More Evidence"

    rec_m1, rec_m2, rec_m3, rec_m4 = st.columns(4)
    rec_m1.metric("Narrative Signal", recommendation_label)
    rec_m2.metric("Current Spread", "N/A" if pd.isna(current_spread) else f"{current_spread:+.1f} bp")
    rec_m3.metric("Spread Movement", "N/A" if pd.isna(spread_change) else f"{spread_change:+.1f} bp")
    rec_m4.metric("Historical Percentile", "N/A" if pd.isna(historical_percentile) else f"{historical_percentile:.0f}th")

    if narrative_lines:
        st.markdown("### Generated Commentary")
        commentary_text = " ".join(narrative_lines)
        st.info(commentary_text)
    else:
        st.info("Not enough data was available to generate a recommendation narrative for the selected inputs.")

    if evidence_rows:
        st.markdown("### Evidence Trail")
        evidence_df = pd.DataFrame(evidence_rows)
        st.dataframe(evidence_df, use_container_width=True, hide_index=True)

    with st.expander("Rule thresholds used in this narrative", expanded=False):
        rule_df = pd.DataFrame(
            [
                {"Rule": "Material widening", "Threshold": "Spread movement >= +15 bp"},
                {"Rule": "Material tightening", "Threshold": "Spread movement <= -15 bp"},
                {"Rule": "Historically wide", "Threshold": "Historical percentile >= 75th"},
                {"Rule": "Very historically wide", "Threshold": "Historical percentile >= 90th"},
                {"Rule": "Above-average liquidity", "Threshold": "Liquidity score >= 70"},
                {"Rule": "Moderate liquidity", "Threshold": "Liquidity score 45–70"},
                {"Rule": "Wide versus peers", "Threshold": "Peer gap >= +10 bp"},
                {"Rule": "Tight versus peers", "Threshold": "Peer gap <= -10 bp"},
                {"Rule": "Sell-heavy flow", "Threshold": "Flow imbalance >= +25%"},
                {"Rule": "Buy-heavy flow", "Threshold": "Flow imbalance <= -25%"},
            ]
        )
        st.dataframe(rule_df, use_container_width=True, hide_index=True)




section_anchor("ai-commentary-studio", "AI Commentary Studio")
with st.expander("Methodology: AI commentary studio", expanded=False):
    st.markdown(
        """
This section adds an **AI writing layer** on top of the dashboard analytics.

**Design principle:**

The AI should not invent the analysis. The dashboard first creates a structured evidence package from spreads, liquidity, historical percentile, peer gap, curve shape, and selected trade activity. The AI then converts that package into polished institutional commentary.

**Recommended use:**

- Use **Dashboard Analytics Only** when you want controlled commentary from internal numbers.
- Add **manual market context** when you already know what happened during the period.
- Enable **web search** only when you want current public market/news context. Review the output carefully.

**Important limitation:**

This commentary is for market discussion and screening. It is not an investment recommendation, credit opinion, or substitute for trader/PM judgment.
        """
    )

ai_col1, ai_col2, ai_col3, ai_col4 = st.columns([1, 1, 1, 1])
with ai_col1:
    ai_bucket = st.selectbox(
        "AI Commentary Bucket",
        MATURITY_BUCKET_ORDER,
        index=3,
        key="ai_commentary_bucket",
    )
with ai_col2:
    ai_rating = st.selectbox(
        "AI Benchmark",
        BENCHMARK_RATINGS,
        index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
        key="ai_commentary_rating",
    )
with ai_col3:
    ai_period = st.selectbox(
        "AI Commentary Period",
        ["1W", "1M", "3M", "6M", "1Y"],
        index=1,
        key="ai_commentary_period",
    )
with ai_col4:
    ai_model = st.selectbox(
        "AI Model",
        ["gpt-4.1-mini", "gpt-4o-mini"],
        index=0,
        key="ai_commentary_model",
    )

ai_period_days = {"1W": 7, "1M": 30, "3M": 90, "6M": 180, "1Y": 365}[ai_period]

st.markdown("### Controlled Market Context Retrieval")
st.caption(
    "Recommended workflow: first build dashboard signals, then retrieve/review public market context, "
    "then generate final commentary. This keeps the AI evidence-linked instead of free-form."
)

context_col1, context_col2 = st.columns([1.2, 1])
with context_col1:
    market_context_query = st.text_area(
        "Market / sector context search focus",
        value=(
            f"Municipal bond market context for {selected_sector} sector and {selected_issuer}; "
            f"focus on Treasury curve, muni market tone, fund flows, and sector headlines over the last {ai_period_days} days."
        ),
        height=110,
        key="ai_market_context_query",
        help="This tells the AI web search what public context to retrieve. Avoid confidential internal details.",
    )
with context_col2:
    ai_context_mode = st.radio(
        "AI Context Mode",
        [
            "Dashboard analytics only",
            "Manual context only",
            "Retrieve market context first",
            "Manual + retrieved context",
        ],
        index=2,
        key="ai_context_mode",
    )

manual_market_context = st.text_area(
    "Optional manual market / sector context",
    placeholder=(
        "Example: Treasury curve steepened during the period; long-duration munis underperformed; "
        "utility sector saw weaker tone after fund outflows..."
    ),
    height=120,
    key="manual_ai_market_context",
)

direct_web_search_in_commentary = st.checkbox(
    "Allow direct web search during final commentary generation",
    value=False,
    key="direct_web_search_in_commentary",
    help=(
        "Usually keep this off. Preferred workflow is: Retrieve Market Context → review it → Generate Commentary. "
        "Turn this on only if you want the final commentary call to do additional web search."
    ),
)

# -----------------------------
# Build structured context package
# -----------------------------
ai_context = {
    "issuer": selected_issuer,
    "sector": selected_sector,
    "bucket": ai_bucket,
    "benchmark": ai_rating,
    "period": ai_period,
    "analytics_as_of": pd.Timestamp.now().strftime("%Y-%m-%d %H:%M"),
    "executive_snapshot": {
        "bonds": len(issuer_bonds),
        "trades_current_filter": len(issuer_trades),
        "latest_trade": issuer_trades["trade_date"].max().strftime("%Y-%m-%d") if not issuer_trades.empty else None,
    },
    "signals": {},
    "market_context_request": {
        "sector": selected_sector,
        "issuer": selected_issuer,
        "period_days": ai_period_days,
        "requested_public_context": [
            "Treasury curve / rates move",
            "municipal market tone",
            "fund flows if available",
            "sector-specific headlines",
            "issuer-specific public news only if available",
        ],
    },
}

# Spread movement and historical percentile.
try:
    ai_obs = build_spread_observations(
        market_df=market_df,
        mmd_df=mmd_df,
        issuer=selected_issuer,
        rating=ai_rating,
    )
    ai_obs = ai_obs.copy()
    ai_obs["trade_date"] = pd.to_datetime(ai_obs["trade_date"], errors="coerce").dt.normalize()
    ai_obs = ai_obs[
        (ai_obs["maturity_bucket"] == ai_bucket)
        & ai_obs["spread_to_benchmark_bps"].notna()
    ].sort_values("trade_date")

    if not ai_obs.empty:
        ai_latest = ai_obs.iloc[-1]
        ai_latest_date = ai_latest["trade_date"]
        ai_current_spread = float(ai_latest["spread_to_benchmark_bps"])
        ai_target_date = ai_latest_date - pd.Timedelta(days=ai_period_days)
        ai_hist_candidates = ai_obs[ai_obs["trade_date"] <= ai_target_date]

        if not ai_hist_candidates.empty:
            ai_hist_row = ai_hist_candidates.iloc[-1]
            ai_spread_change = ai_current_spread - float(ai_hist_row["spread_to_benchmark_bps"])
        else:
            ai_spread_change = None

        ai_1y = ai_obs[ai_obs["trade_date"] >= ai_latest_date - pd.Timedelta(days=365)]
        ai_values = pd.to_numeric(ai_1y["spread_to_benchmark_bps"], errors="coerce").dropna()
        ai_percentile = float((ai_values <= ai_current_spread).mean() * 100) if len(ai_values) >= 2 else None

        ai_context["signals"]["spread"] = {
            "current_spread_bps": round(ai_current_spread, 2),
            "spread_change_bps": None if ai_spread_change is None else round(ai_spread_change, 2),
            "historical_percentile_1y": None if ai_percentile is None else round(ai_percentile, 1),
            "latest_spread_date": ai_latest_date.strftime("%Y-%m-%d"),
            "benchmark_source": ai_latest.get("benchmark_source"),
            "benchmark_column": ai_latest.get("source_column"),
        }
except Exception as exc:
    ai_context["signals"]["spread_error"] = str(exc)

# Liquidity proxy.
try:
    ai_trades = market_df[
        (market_df["issuer"] == selected_issuer)
        & (market_df["maturity_bucket"] == ai_bucket)
    ].copy()
    ai_trades["trade_date"] = pd.to_datetime(ai_trades["trade_date"], errors="coerce").dt.normalize()
    if "trade_amount" in ai_trades.columns:
        ai_trades["trade_amount"] = pd.to_numeric(ai_trades["trade_amount"], errors="coerce").fillna(0)
    else:
        ai_trades["trade_amount"] = 0.0
    ai_trades = ai_trades.dropna(subset=["trade_date"])
    if not ai_trades.empty:
        ai_latest_trade = ai_trades["trade_date"].max()
        ai_trade_window = ai_trades[ai_trades["trade_date"] >= ai_latest_trade - pd.Timedelta(days=ai_period_days)].copy()
        ai_trade_count = len(ai_trade_window)
        ai_total_amount = float(ai_trade_window["trade_amount"].sum())
        ai_days_since_last = int((pd.Timestamp.today().normalize() - ai_latest_trade).days)
        ai_liq_score = (
            min(ai_trade_count / 10, 1) * 35
            + min(ai_total_amount / 5_000_000, 1) * 35
            + max(0, 1 - min(ai_days_since_last / 180, 1)) * 30
        )

        ai_context["signals"]["liquidity"] = {
            "liquidity_score_proxy": round(ai_liq_score, 1),
            "trade_count": int(ai_trade_count),
            "total_trade_amount": round(ai_total_amount, 0),
            "days_since_last_trade": ai_days_since_last,
        }
except Exception as exc:
    ai_context["signals"]["liquidity_error"] = str(exc)

# Peer gap proxy.
try:
    if "sector" in market_df.columns and selected_sector and selected_sector != "Unknown":
        ai_peer_universe = market_df[
            (market_df["sector"].astype(str) == str(selected_sector))
            & (market_df["issuer"].astype(str) != str(selected_issuer))
            & (market_df["maturity_bucket"] == ai_bucket)
        ].copy()
    else:
        ai_peer_universe = market_df[
            (market_df["issuer"].astype(str) != str(selected_issuer))
            & (market_df["maturity_bucket"] == ai_bucket)
        ].copy()

    if not ai_peer_universe.empty and "spread" in ai_context["signals"]:
        ai_peer_universe["trade_date"] = pd.to_datetime(ai_peer_universe["trade_date"], errors="coerce").dt.normalize()
        ai_peer_universe["yield"] = pd.to_numeric(ai_peer_universe["yield"], errors="coerce")
        ai_peer_universe = ai_peer_universe.dropna(subset=["trade_date", "yield"])
        if not ai_peer_universe.empty:
            ai_peer_latest = ai_peer_universe["trade_date"].max()
            ai_peer_universe = ai_peer_universe[
                ai_peer_universe["trade_date"] >= ai_peer_latest - pd.Timedelta(days=ai_period_days)
            ].copy()

            ai_date_col = _detect_mmd_date_column(mmd_df)
            if ai_date_col is not None and not ai_peer_universe.empty:
                ai_peer_mmd = mmd_df.copy()
                ai_peer_mmd[ai_date_col] = pd.to_datetime(ai_peer_mmd[ai_date_col], errors="coerce")
                ai_peer_mmd = ai_peer_mmd.dropna(subset=[ai_date_col])
                ai_peer_mmd = ai_peer_mmd[
                    ai_peer_mmd[ai_date_col].dt.normalize() <= ai_peer_latest
                ].sort_values(ai_date_col)

                if not ai_peer_mmd.empty:
                    ai_peer_tenor = MMD_BUCKET_MAP.get(ai_bucket, "10Y")
                    ai_peer_bench, _ = get_benchmark_curve(ai_peer_mmd.iloc[[-1]], ai_peer_tenor, ai_rating)
                    if ai_peer_bench is not None and pd.notna(ai_peer_bench.iloc[0]):
                        ai_peer_benchmark_yield = float(ai_peer_bench.iloc[0])
                        ai_peer_summary = (
                            ai_peer_universe.groupby("issuer", as_index=False)
                            .agg(avg_yield=("yield", "mean"), trade_count=("yield", "count"))
                        )
                        ai_peer_summary["peer_spread_bps"] = (
                            ai_peer_summary["avg_yield"] - ai_peer_benchmark_yield
                        ) * 100
                        ai_peer_median = float(ai_peer_summary["peer_spread_bps"].median())
                        ai_peer_gap = ai_context["signals"]["spread"]["current_spread_bps"] - ai_peer_median

                        ai_context["signals"]["peer_comparison"] = {
                            "peer_median_spread_bps": round(ai_peer_median, 2),
                            "peer_gap_bps": round(ai_peer_gap, 2),
                            "peer_count": int(len(ai_peer_summary)),
                        }
except Exception as exc:
    ai_context["signals"]["peer_error"] = str(exc)

# Curve shape snapshot if enough points.
try:
    ai_curve = market_df[market_df["issuer"] == selected_issuer].copy()
    ai_curve["trade_date"] = pd.to_datetime(ai_curve["trade_date"], errors="coerce").dt.normalize()
    ai_curve["yield"] = pd.to_numeric(ai_curve["yield"], errors="coerce")
    ai_curve = ai_curve.dropna(subset=["trade_date", "yield", "maturity_bucket"])
    ai_curve = ai_curve[ai_curve["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)]
    if not ai_curve.empty:
        ai_curve_latest = ai_curve["trade_date"].max()
        ai_curve_window = ai_curve[ai_curve["trade_date"] >= ai_curve_latest - pd.Timedelta(days=30)]
        ai_curve_summary = ai_curve_window.groupby("maturity_bucket")["yield"].mean().to_dict()
        ai_context["signals"]["curve_shape"] = {
            "short_yield": ai_curve_summary.get("Short"),
            "ten_yield": ai_curve_summary.get("10Y"),
            "twenty_yield": ai_curve_summary.get("20Y"),
            "thirty_yield": ai_curve_summary.get("30Y"),
        }
        if "10Y" in ai_curve_summary and "30Y" in ai_curve_summary:
            ai_context["signals"]["curve_shape"]["10s30s_slope_pct"] = round(
                ai_curve_summary["30Y"] - ai_curve_summary["10Y"], 4
            )
except Exception as exc:
    ai_context["signals"]["curve_error"] = str(exc)

st.subheader("AI Context Package")
with st.expander("Review structured evidence before sending to AI", expanded=False):
    st.json(ai_context)

# -----------------------------
# Controlled retrieval + synthesis workflow
# -----------------------------
retrieve_enabled = ai_context_mode in ["Retrieve market context first", "Manual + retrieved context"]
manual_enabled = ai_context_mode in ["Manual context only", "Manual + retrieved context"]

if retrieve_enabled or direct_web_search_in_commentary:
    st.warning(
        "Public web context may be retrieved. Do not include confidential, proprietary, or client-sensitive information "
        "in the market context query or manual context field."
    )

action_col1, action_col2 = st.columns([1, 1])

with action_col1:
    if st.button("Retrieve Market / Sector Context", key="retrieve_ai_market_context", disabled=not retrieve_enabled):
        with st.spinner("Retrieving controlled public market context..."):
            retrieved_context = retrieve_market_context_with_openai(
                context_package=ai_context,
                market_context_query=market_context_query,
                model=ai_model,
            )
            st.session_state["latest_retrieved_market_context"] = retrieved_context

with action_col2:
    if st.button("Generate AI Institutional Commentary", key="generate_ai_institutional_commentary"):
        retrieved_context_for_commentary = (
            st.session_state.get("latest_retrieved_market_context", "")
            if ai_context_mode in ["Retrieve market context first", "Manual + retrieved context"]
            else ""
        )
        manual_context_for_commentary = manual_market_context if manual_enabled else ""

        with st.spinner("Generating evidence-linked commentary..."):
            commentary = generate_ai_market_commentary(
                context_package=ai_context,
                manual_market_context=manual_context_for_commentary,
                retrieved_market_context=retrieved_context_for_commentary,
                use_web_search=direct_web_search_in_commentary,
                market_context_query=market_context_query,
                model=ai_model,
            )
            st.session_state["latest_ai_commentary"] = commentary

if "latest_retrieved_market_context" in st.session_state:
    st.subheader("Retrieved Market / Sector Context")
    st.markdown(st.session_state["latest_retrieved_market_context"])

    st.download_button(
        label="Download Retrieved Market Context Markdown",
        data=st.session_state["latest_retrieved_market_context"].encode("utf-8"),
        file_name=f"{selected_issuer}_retrieved_market_context.md".replace(" ", "_"),
        mime="text/markdown",
    )

if "latest_ai_commentary" in st.session_state:
    st.subheader("Generated Institutional Commentary")
    st.markdown(st.session_state["latest_ai_commentary"])

    st.download_button(
        label="Download AI Commentary Markdown",
        data=st.session_state["latest_ai_commentary"].encode("utf-8"),
        file_name=f"{selected_issuer}_ai_market_commentary.md".replace(" ", "_"),
        mime="text/markdown",
    )

with st.expander("Recommended AI architecture for this dashboard", expanded=False):
    st.markdown(
        """
**Best practice: use a centralized AI Commentary Studio, not one AI button per section.**

Why:
1. Lower API cost.
2. Less duplicated commentary.
3. Lower risk of conflicting explanations.
4. Easier review by analysts / managers.
5. Cleaner evidence trail.

Recommended workflow:
1. Dashboard computes deterministic analytics.
2. AI Context Package captures the evidence.
3. Controlled Market Context Retrieval gathers public market/sector context.
4. User reviews retrieved context.
5. Final AI commentary synthesizes data + context into institutional language.
        """
    )

with st.expander("Where should AI commentary live?", expanded=False):
    st.markdown(
        """
I recommend **not** putting an AI button inside every single section.

A cleaner institutional structure is:

1. **Section-level rule commentary** stays deterministic and auditable.
2. **AI Commentary Studio** synthesizes across multiple sections.
3. Specific section AI can be added later only for the highest-value areas:
   - Historical Spread Percentile
   - Peer / Cross-Issuer RV
   - Security Screener
   - Recommendation Narrative

This keeps cost lower, avoids repeated/conflicting narratives, and makes the output easier for a team to review.
        """
    )


section_anchor("scenario-shock", "Scenario Shock Analysis")
with st.expander("Methodology: scenario shock analysis", expanded=False):
    st.markdown(
        """
This section estimates how the selected issuer or uploaded securities may react under simple interest-rate shock scenarios.

**Purpose:**

- Estimate approximate price impact under rate shocks.
- Identify which maturity buckets or CUSIPs are most exposed to parallel moves, steepening, or flattening.
- Provide a first-pass risk lens for secondary trading and pitchbook discussion.

**Version 1 approximation:**

This is a **duration-proxy model**, not a full bond pricing engine.

Core formula:

`Approximate Price Impact ≈ -Duration × Yield Shock`

Where:

- Duration is proxied by maturity bucket unless a duration field is uploaded.
- Yield shock is expressed in decimal form. Example: `+25 bp = +0.0025`.
- Price impact is an approximate percentage price move.

**Default proxy durations:**

| Bucket | Proxy Duration |
|---|---:|
| Short | 2.0 |
| 10Y | 8.0 |
| 20Y | 13.0 |
| 30Y | 18.0 |

**Important limitations:**

- This does not model callable optionality, convexity, OAS, amortization, tax effects, or full cash flows.
- Callable / premium bonds may behave differently from this simple duration approximation.
- Treat this as a screening and risk-discussion tool, not a final valuation model.
        """
    )

DURATION_PROXY = {
    "Short": 2.0,
    "10Y": 8.0,
    "20Y": 13.0,
    "30Y": 18.0,
}

SHOCK_SCENARIOS_BPS = {
    "+25bp Parallel": {"Short": 25, "10Y": 25, "20Y": 25, "30Y": 25},
    "+50bp Parallel": {"Short": 50, "10Y": 50, "20Y": 50, "30Y": 50},
    "-25bp Parallel": {"Short": -25, "10Y": -25, "20Y": -25, "30Y": -25},
    "Bear Steepening": {"Short": 5, "10Y": 15, "20Y": 25, "30Y": 35},
    "Bull Flattening": {"Short": -25, "10Y": -20, "20Y": -10, "30Y": -5},
    "Front-End Selloff": {"Short": 35, "10Y": 20, "20Y": 10, "30Y": 5},
}

shock_col1, shock_col2, shock_col3 = st.columns([1, 1, 1])
with shock_col1:
    shock_scope = st.selectbox(
        "Shock Scope",
        ["Selected issuer", "All uploaded issuers"],
        index=0,
        key="shock_scope",
    )
with shock_col2:
    shock_scenario = st.selectbox(
        "Rate Shock Scenario",
        list(SHOCK_SCENARIOS_BPS.keys()) + ["Custom"],
        index=0,
        key="shock_scenario",
    )
with shock_col3:
    shock_view = st.selectbox(
        "Shock View",
        ["Maturity bucket summary", "CUSIP-level detail"],
        index=0,
        key="shock_view",
    )

if shock_scenario == "Custom":
    custom_col1, custom_col2, custom_col3, custom_col4 = st.columns(4)
    with custom_col1:
        shock_short = st.number_input("Short Shock (bp)", value=25.0, step=5.0, key="shock_short")
    with custom_col2:
        shock_10 = st.number_input("10Y Shock (bp)", value=25.0, step=5.0, key="shock_10")
    with custom_col3:
        shock_20 = st.number_input("20Y Shock (bp)", value=25.0, step=5.0, key="shock_20")
    with custom_col4:
        shock_30 = st.number_input("30Y Shock (bp)", value=25.0, step=5.0, key="shock_30")
    selected_shocks = {
        "Short": float(shock_short),
        "10Y": float(shock_10),
        "20Y": float(shock_20),
        "30Y": float(shock_30),
    }
else:
    selected_shocks = SHOCK_SCENARIOS_BPS[shock_scenario]

shock_base = market_df.copy()
if shock_scope == "Selected issuer":
    shock_base = shock_base[shock_base["issuer"] == selected_issuer].copy()

if shock_base.empty:
    st.warning("No trade rows are available for the selected shock scope.")
else:
    shock_base["trade_date"] = pd.to_datetime(shock_base["trade_date"], errors="coerce").dt.normalize()
    shock_base["yield"] = pd.to_numeric(shock_base["yield"], errors="coerce")
    if "price" in shock_base.columns:
        shock_base["price"] = pd.to_numeric(shock_base["price"], errors="coerce")
    else:
        shock_base["price"] = pd.NA
    if "trade_amount" in shock_base.columns:
        shock_base["trade_amount"] = pd.to_numeric(shock_base["trade_amount"], errors="coerce").fillna(0)
    else:
        shock_base["trade_amount"] = 0.0

    shock_base = shock_base.dropna(subset=["trade_date", "yield", "maturity_bucket"])
    shock_base = shock_base[shock_base["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)].copy()

    if shock_base.empty:
        st.warning("No usable rows with maturity buckets were available for scenario shock analysis.")
    else:
        # Build bucket summary from latest available selected-scope data.
        latest_shock_date = shock_base["trade_date"].max()
        shock_recent = shock_base[shock_base["trade_date"] >= latest_shock_date - pd.Timedelta(days=90)].copy()
        if shock_recent.empty:
            shock_recent = shock_base.copy()

        bucket_summary = (
            shock_recent.groupby("maturity_bucket", as_index=False)
            .agg(
                avg_yield=("yield", "mean"),
                avg_price=("price", "mean"),
                trade_count=("yield", "count"),
                total_trade_amount=("trade_amount", "sum"),
                latest_trade=("trade_date", "max"),
            )
        )
        bucket_summary["proxy_duration"] = bucket_summary["maturity_bucket"].map(DURATION_PROXY)
        bucket_summary["shock_bps"] = bucket_summary["maturity_bucket"].map(selected_shocks)
        bucket_summary["shock_decimal"] = bucket_summary["shock_bps"] / 10000
        bucket_summary["approx_price_impact_pct"] = -bucket_summary["proxy_duration"] * bucket_summary["shock_decimal"] * 100
        bucket_summary["shocked_yield"] = bucket_summary["avg_yield"] + (bucket_summary["shock_bps"] / 100)
        bucket_summary["impact_direction"] = bucket_summary["approx_price_impact_pct"].map(
            lambda x: "Price Down" if x < 0 else "Price Up" if x > 0 else "Flat"
        )

        shock_m1, shock_m2, shock_m3, shock_m4 = st.columns(4)
        worst_bucket_row = bucket_summary.sort_values("approx_price_impact_pct").iloc[0]
        best_bucket_row = bucket_summary.sort_values("approx_price_impact_pct", ascending=False).iloc[0]
        weighted_impact = (
            (bucket_summary["approx_price_impact_pct"] * bucket_summary["total_trade_amount"]).sum()
            / bucket_summary["total_trade_amount"].sum()
            if bucket_summary["total_trade_amount"].sum() > 0
            else bucket_summary["approx_price_impact_pct"].mean()
        )

        shock_m1.metric("Scenario", shock_scenario)
        shock_m2.metric("Worst Bucket", f"{worst_bucket_row['maturity_bucket']}")
        shock_m3.metric("Worst Approx Impact", f"{worst_bucket_row['approx_price_impact_pct']:+.2f}%")
        shock_m4.metric("Weighted Impact", f"{weighted_impact:+.2f}%")

        st.info(
            f"Read-through: under **{shock_scenario}**, the most rate-sensitive bucket is "
            f"{worst_bucket_row['maturity_bucket']} with an approximate price impact of "
            f"{worst_bucket_row['approx_price_impact_pct']:+.2f}%. "
            f"This is based on proxy duration and should be treated as a first-pass risk estimate."
        )

        st.subheader("1. Shock Impact by Maturity Bucket")
        shock_bar = px.bar(
            bucket_summary,
            x="maturity_bucket",
            y="approx_price_impact_pct",
            hover_data={
                "avg_yield": ":.2f",
                "shocked_yield": ":.2f",
                "proxy_duration": ":.1f",
                "shock_bps": ":.0f",
                "trade_count": ":,.0f",
                "total_trade_amount": ":,.0f",
            },
            title=f"Approximate Price Impact by Bucket — {shock_scenario}",
            labels={
                "maturity_bucket": "Maturity Bucket",
                "approx_price_impact_pct": "Approximate Price Impact (%)",
                "avg_yield": "Current Avg Yield",
                "shocked_yield": "Shocked Yield",
                "proxy_duration": "Proxy Duration",
                "shock_bps": "Shock (bp)",
            },
        )
        shock_bar.add_hline(y=0, line_dash="dash", opacity=0.45)
        shock_bar.update_layout(height=440)
        st.plotly_chart(shock_bar, use_container_width=True)

        st.subheader("2. Shock Summary Table")
        bucket_display = bucket_summary.copy()
        for col in ["avg_yield", "avg_price", "proxy_duration", "shocked_yield", "approx_price_impact_pct"]:
            if col in bucket_display.columns:
                bucket_display[col] = pd.to_numeric(bucket_display[col], errors="coerce").round(2)
        st.dataframe(
            bucket_display[
                [
                    "maturity_bucket",
                    "avg_yield",
                    "shocked_yield",
                    "shock_bps",
                    "proxy_duration",
                    "approx_price_impact_pct",
                    "trade_count",
                    "total_trade_amount",
                    "latest_trade",
                    "impact_direction",
                ]
            ],
            use_container_width=True,
            hide_index=True,
        )

        # CUSIP-level shock detail.
        if shock_view == "CUSIP-level detail":
            st.subheader("3. CUSIP-Level Shock Detail")
            cusip_shock = (
                shock_recent.groupby("cusip", dropna=False)
                .agg(
                    issuer=("issuer", "first"),
                    sector=("sector", "first") if "sector" in shock_recent.columns else ("issuer", "first"),
                    maturity_bucket=("maturity_bucket", "first"),
                    maturity=("maturity_bond", "first") if "maturity_bond" in shock_recent.columns else ("trade_date", "max"),
                    coupon=("coupon_bond", "first") if "coupon_bond" in shock_recent.columns else ("yield", "count"),
                    avg_yield=("yield", "mean"),
                    avg_price=("price", "mean"),
                    trade_count=("yield", "count"),
                    total_trade_amount=("trade_amount", "sum"),
                    latest_trade=("trade_date", "max"),
                )
                .reset_index()
            )
            cusip_shock["proxy_duration"] = cusip_shock["maturity_bucket"].map(DURATION_PROXY)
            cusip_shock["shock_bps"] = cusip_shock["maturity_bucket"].map(selected_shocks)
            cusip_shock["shock_decimal"] = cusip_shock["shock_bps"] / 10000
            cusip_shock["approx_price_impact_pct"] = -cusip_shock["proxy_duration"] * cusip_shock["shock_decimal"] * 100
            cusip_shock["shocked_yield"] = cusip_shock["avg_yield"] + (cusip_shock["shock_bps"] / 100)
            cusip_shock = cusip_shock.sort_values("approx_price_impact_pct")

            detail_cols = [
                "cusip",
                "issuer",
                "sector",
                "maturity_bucket",
                "maturity",
                "coupon",
                "avg_yield",
                "shocked_yield",
                "shock_bps",
                "proxy_duration",
                "approx_price_impact_pct",
                "trade_count",
                "total_trade_amount",
                "latest_trade",
            ]
            detail_display = cusip_shock[[c for c in detail_cols if c in cusip_shock.columns]].copy()
            for col in ["avg_yield", "shocked_yield", "proxy_duration", "approx_price_impact_pct"]:
                if col in detail_display.columns:
                    detail_display[col] = pd.to_numeric(detail_display[col], errors="coerce").round(2)

            st.dataframe(detail_display.head(5000), use_container_width=True, hide_index=True, height=480)

            shock_scatter = px.scatter(
                cusip_shock,
                x="proxy_duration",
                y="approx_price_impact_pct",
                size="total_trade_amount",
                color="maturity_bucket",
                hover_name="cusip",
                hover_data=[
                    c for c in ["issuer", "avg_yield", "shocked_yield", "shock_bps", "trade_count", "latest_trade"]
                    if c in cusip_shock.columns
                ],
                title="CUSIP Shock Exposure Map",
                labels={
                    "proxy_duration": "Proxy Duration",
                    "approx_price_impact_pct": "Approx Price Impact (%)",
                    "maturity_bucket": "Maturity Bucket",
                },
            )
            shock_scatter.add_hline(y=0, line_dash="dash", opacity=0.45)
            shock_scatter.update_layout(height=500)
            st.plotly_chart(shock_scatter, use_container_width=True)

        with st.expander("Scenario shock assumptions and audit", expanded=False):
            shock_assumption_df = pd.DataFrame(
                [
                    {
                        "Maturity Bucket": bucket,
                        "Shock (bp)": selected_shocks.get(bucket),
                        "Proxy Duration": DURATION_PROXY.get(bucket),
                        "Formula": "Approx Price Impact ≈ -Duration × Shock",
                    }
                    for bucket in MATURITY_BUCKET_ORDER
                ]
            )
            st.dataframe(shock_assumption_df, use_container_width=True, hide_index=True)

            st.download_button(
                label="Download Scenario Shock Results CSV",
                data=bucket_summary.to_csv(index=False).encode("utf-8"),
                file_name="scenario_shock_results.csv",
                mime="text/csv",
            )


section_anchor("spread-movement", "Spread Movement Heatmap")
with st.expander("Methodology: spread movement heatmap", expanded=False):
    st.markdown(
        """
This heatmap shows whether the selected issuer has become richer or cheaper versus the selected benchmark curve.

**Calculation:**

`Issuer Spread = (Average Issuer Trade Yield - Benchmark Yield) × 100`

`Spread Movement = Latest Available Issuer Spread - Historical Issuer Spread`

**How to read it:**

- **Positive / red = widening**: issuer spread increased versus the benchmark; the issuer/bucket became cheaper or underperformed.
- **Negative / green = tightening**: issuer spread decreased versus the benchmark; the issuer/bucket became richer or outperformed.
- Rows are maturity buckets. Columns are lookback windows.
- Because municipal bonds can trade sparsely, the historical value uses the latest available observation at or before the lookback target date.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD curve file to enable the spread movement heatmap.")
else:
    heatmap_col1, heatmap_col2 = st.columns([1, 2])
    with heatmap_col1:
        heatmap_rating = st.selectbox(
            "Heatmap Benchmark Curve",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            help="Priority: uploaded rating curve columns first; otherwise MMD/AAA plus the visible rating-spread assumptions.",
        )
    with heatmap_col2:
        st.caption(
            "Cells show change in spread, in basis points, from the latest available observation to each lookback window."
        )

    heatmap_spread_obs = build_spread_observations(
        market_df=market_df,
        mmd_df=mmd_df,
        issuer=selected_issuer,
        rating=heatmap_rating,
    )

    if heatmap_spread_obs.empty:
        st.warning(
            "No overlapping issuer trade dates and benchmark dates were found for the heatmap. "
            "Check that the curve file has a Date column plus either 5Y/10Y/20Y/30Y base columns or explicit rating curve columns such as AA_10Y, and that trade dates overlap with the curve history."
        )
    else:
        heatmap_matrix, heatmap_audit = build_spread_movement_heatmap_data(heatmap_spread_obs)
        if heatmap_matrix.isna().all().all():
            st.info("Not enough historical spread observations to calculate movement across the selected windows yet.")
        else:
            heatmap_text = heatmap_matrix.map(lambda x: "" if pd.isna(x) else f"{x:+.1f} bp")
            heatmap_fig = px.imshow(
                heatmap_matrix.astype(float),
                x=heatmap_matrix.columns,
                y=heatmap_matrix.index,
                color_continuous_scale=["#1a9850", "#f7f7f7", "#d73027"],
                color_continuous_midpoint=0,
                aspect="auto",
                title=f"{selected_issuer} Spread Movement vs {heatmap_rating} Curve",
                labels={"x": "Lookback Window", "y": "Maturity Bucket", "color": "Spread Movement (bps)"},
            )
            heatmap_fig.update_traces(text=heatmap_text.values, texttemplate="%{text}", hovertemplate="Maturity=%{y}<br>Window=%{x}<br>Movement=%{z:.1f} bp<extra></extra>")
            heatmap_fig.update_layout(height=420)
            st.plotly_chart(heatmap_fig, use_container_width=True)

            latest_obs_date = heatmap_spread_obs["trade_date"].max()
            st.caption(
                f"Latest available spread observation used: {latest_obs_date.strftime('%Y-%m-%d')}. "
                "Positive values indicate spread widening; negative values indicate spread tightening."
            )

            with st.expander("Heatmap calculation audit table", expanded=False):
                display_cols = [
                    "maturity_bucket", "window", "latest_date", "latest_spread_bps", "target_date",
                    "historical_date", "historical_spread_bps", "spread_movement_bps", "note",
                ]
                audit_display = heatmap_audit[[c for c in display_cols if c in heatmap_audit.columns]].copy()
                for c in ["latest_spread_bps", "historical_spread_bps", "spread_movement_bps"]:
                    if c in audit_display.columns:
                        audit_display[c] = pd.to_numeric(audit_display[c], errors="coerce").round(2)
                st.dataframe(audit_display, use_container_width=True, hide_index=True)



section_anchor("cusip-drilldown", "CUSIP Opportunity Drilldown")
with st.expander("Methodology: CUSIP opportunity drilldown", expanded=False):
    st.markdown(
        """
This section moves from issuer-level signals into **specific bond-level candidates**.

**Purpose:**

- Identify which CUSIPs are driving a maturity bucket's relative-value signal.
- Compare current CUSIP-level spread, recent movement, trade count, and liquidity.
- Help the team move from: *"30Y widened"* to *"which 30Y bonds should we look at?"*

**Calculation overview:**

- Current CUSIP yield is calculated over the selected lookback window.
- Current spread is calculated as:

`CUSIP Spread = (Average CUSIP Yield - Benchmark Yield) × 100`

- Historical spread uses the most recent observation at or before the lookback target date.
- Spread change is:

`Spread Change = Current Spread - Historical Spread`

Positive spread change means widening; negative spread change means tightening.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD curve file to enable CUSIP-level spread drilldown.")
elif issuer_trades.empty:
    st.warning("No trade rows found for the selected issuer and filters.")
else:
    dd_col1, dd_col2, dd_col3, dd_col4 = st.columns([1, 1, 1, 1])
    with dd_col1:
        dd_bucket = st.selectbox(
            "Drilldown Maturity Bucket",
            MATURITY_BUCKET_ORDER,
            index=3,
            help="Focus the drilldown on one maturity bucket.",
        )
    with dd_col2:
        dd_rating = st.selectbox(
            "Drilldown Benchmark",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            help="Priority: uploaded benchmark curve first; otherwise modeled from MMD + spread assumptions.",
        )
    with dd_col3:
        dd_lookback_label = st.selectbox(
            "Movement Lookback",
            ["1W", "1M", "3M", "6M", "1Y"],
            index=1,
        )
    with dd_col4:
        dd_min_trades = st.number_input(
            "Minimum Trades",
            min_value=1,
            max_value=50,
            value=1,
            step=1,
            help="Filter out CUSIPs with fewer trades in the current lookback window.",
        )

    dd_window_days = {"1W": 7, "1M": 30, "3M": 90, "6M": 180, "1Y": 365}[dd_lookback_label]
    dd_tenor = MMD_BUCKET_MAP.get(dd_bucket, "10Y")
    dd_date_col = _detect_mmd_date_column(mmd_df)

    if dd_date_col is None:
        st.warning("CUSIP drilldown cannot run because the benchmark file does not contain a usable date column.")
    else:
        dd_base = market_df[
            (market_df["issuer"] == selected_issuer)
            & (market_df["maturity_bucket"] == dd_bucket)
        ].copy()

        if dd_base.empty:
            st.warning(f"No {dd_bucket} CUSIP-level trade rows were found for {selected_issuer}.")
        else:
            dd_base["trade_date"] = pd.to_datetime(dd_base["trade_date"], errors="coerce").dt.normalize()
            dd_base["yield"] = pd.to_numeric(dd_base["yield"], errors="coerce")
            if "trade_amount" in dd_base.columns:
                dd_base["trade_amount"] = pd.to_numeric(dd_base["trade_amount"], errors="coerce").fillna(0)
            else:
                dd_base["trade_amount"] = 0.0
            if "price" in dd_base.columns:
                dd_base["price"] = pd.to_numeric(dd_base["price"], errors="coerce")
            else:
                dd_base["price"] = pd.NA

            dd_base = dd_base.dropna(subset=["trade_date", "yield", "cusip"])

            if dd_base.empty:
                st.warning("CUSIP drilldown cannot run because no valid CUSIP/date/yield rows remain after cleaning.")
            else:
                dd_latest_date = dd_base["trade_date"].max()
                dd_current_start = dd_latest_date - pd.Timedelta(days=dd_window_days)
                dd_hist_target = dd_current_start

                dd_current = dd_base[dd_base["trade_date"] >= dd_current_start].copy()
                if dd_current.empty:
                    st.warning("No CUSIP trades were found inside the selected current lookback window.")
                else:
                    # Benchmark curve on or before latest issuer trade date.
                    dd_mmd = mmd_df.copy()
                    dd_mmd[dd_date_col] = pd.to_datetime(dd_mmd[dd_date_col], errors="coerce")
                    dd_mmd = dd_mmd.dropna(subset=[dd_date_col])
                    dd_mmd = dd_mmd[dd_mmd[dd_date_col].dt.normalize() <= dd_latest_date].sort_values(dd_date_col)

                    if dd_mmd.empty:
                        st.warning("No benchmark curve observation was available on or before the latest issuer trade date.")
                    else:
                        dd_latest_mmd = dd_mmd.iloc[[-1]].copy()
                        dd_benchmark_date = dd_latest_mmd[dd_date_col].iloc[0]
                        dd_benchmark_yield_series, dd_meta = get_benchmark_curve(dd_latest_mmd, dd_tenor, dd_rating)

                        if dd_benchmark_yield_series is None or pd.isna(dd_benchmark_yield_series.iloc[0]):
                            st.warning(f"{dd_rating} {dd_tenor} benchmark could not be built for this drilldown.")
                        else:
                            dd_benchmark_yield = float(dd_benchmark_yield_series.iloc[0])

                            current_summary = (
                                dd_current.groupby("cusip", dropna=False)
                                .agg(
                                    current_avg_yield=("yield", "mean"),
                                    latest_yield=("yield", "last"),
                                    avg_price=("price", "mean"),
                                    trade_count=("trade_date", "count"),
                                    latest_trade=("trade_date", "max"),
                                    first_trade=("trade_date", "min"),
                                    total_trade_amount=("trade_amount", "sum"),
                                    avg_trade_amount=("trade_amount", "mean"),
                                    maturity=("maturity_bond", "first") if "maturity_bond" in dd_current.columns else ("trade_date", "max"),
                                    coupon=("coupon_bond", "first") if "coupon_bond" in dd_current.columns else ("yield", "count"),
                                    call_date=("call_date", "first") if "call_date" in dd_current.columns else ("trade_date", "max"),
                                    call_price=("call_price", "first") if "call_price" in dd_current.columns else ("yield", "count"),
                                    outstanding_amount=("outstanding_amount", "first") if "outstanding_amount" in dd_current.columns else ("trade_amount", "sum"),
                                )
                                .reset_index()
                            )

                            current_summary["current_spread_bps"] = (
                                current_summary["current_avg_yield"] - dd_benchmark_yield
                            ) * 100

                            # Historical CUSIP spread at or before the lookback target date.
                            hist_candidates = dd_base[dd_base["trade_date"] <= dd_hist_target].copy()
                            if not hist_candidates.empty:
                                hist_rows = (
                                    hist_candidates.sort_values("trade_date")
                                    .groupby("cusip", as_index=False)
                                    .tail(1)[["cusip", "trade_date", "yield", "price", "trade_amount"]]
                                    .rename(
                                        columns={
                                            "trade_date": "historical_trade_date",
                                            "yield": "historical_yield",
                                            "price": "historical_price",
                                            "trade_amount": "historical_trade_amount",
                                        }
                                    )
                                )

                                dd_hist_mmd = mmd_df.copy()
                                dd_hist_mmd[dd_date_col] = pd.to_datetime(dd_hist_mmd[dd_date_col], errors="coerce")
                                dd_hist_mmd = dd_hist_mmd.dropna(subset=[dd_date_col])
                                dd_hist_mmd = dd_hist_mmd[dd_hist_mmd[dd_date_col].dt.normalize() <= dd_hist_target].sort_values(dd_date_col)

                                if not dd_hist_mmd.empty:
                                    dd_hist_latest_mmd = dd_hist_mmd.iloc[[-1]].copy()
                                    dd_hist_benchmark_yield_series, dd_hist_meta = get_benchmark_curve(
                                        dd_hist_latest_mmd, dd_tenor, dd_rating
                                    )
                                    if dd_hist_benchmark_yield_series is not None and pd.notna(dd_hist_benchmark_yield_series.iloc[0]):
                                        dd_hist_benchmark_yield = float(dd_hist_benchmark_yield_series.iloc[0])
                                        hist_rows["historical_spread_bps"] = (
                                            hist_rows["historical_yield"] - dd_hist_benchmark_yield
                                        ) * 100
                                    else:
                                        hist_rows["historical_spread_bps"] = pd.NA
                                else:
                                    hist_rows["historical_spread_bps"] = pd.NA
                            else:
                                hist_rows = pd.DataFrame(columns=["cusip", "historical_trade_date", "historical_yield", "historical_price", "historical_trade_amount", "historical_spread_bps"])

                            dd_opps = current_summary.merge(hist_rows, on="cusip", how="left")
                            dd_opps["spread_change_bps"] = dd_opps["current_spread_bps"] - dd_opps["historical_spread_bps"]
                            dd_opps["yield_change_bps"] = (dd_opps["current_avg_yield"] - dd_opps["historical_yield"]) * 100

                            # Liquidity score proxy for current window.
                            dd_today = pd.Timestamp.today().normalize()
                            dd_opps["days_since_last_trade"] = (dd_today - dd_opps["latest_trade"]).dt.days
                            dd_opps["liquidity_score"] = (
                                dd_opps["trade_count"].rank(pct=True) * 40
                                + dd_opps["total_trade_amount"].rank(pct=True) * 35
                                + (1 - dd_opps["days_since_last_trade"].rank(pct=True)) * 25
                            )
                            dd_opps["liquidity_tier"] = pd.cut(
                                dd_opps["liquidity_score"],
                                bins=[-1, 45, 75, 101],
                                labels=["Low", "Medium", "High"],
                            ).astype(str)

                            dd_opps = dd_opps[dd_opps["trade_count"] >= dd_min_trades].copy()

                            if dd_opps.empty:
                                st.info("No CUSIPs met the selected minimum trade filter.")
                            else:
                                dd_sort_options = {
                                    "Current Spread": "current_spread_bps",
                                    "Spread Change": "spread_change_bps",
                                    "Liquidity Score": "liquidity_score",
                                    "Trade Count": "trade_count",
                                    "Total Trade Amount": "total_trade_amount",
                                }
                                dd_sort_label = st.selectbox(
                                    "Sort Opportunities By",
                                    list(dd_sort_options.keys()),
                                    index=0,
                                    key="dd_sort_opportunities",
                                )
                                dd_sort_col = dd_sort_options[dd_sort_label]
                                dd_opps = dd_opps.sort_values(dd_sort_col, ascending=False, na_position="last")

                                summary_c1, summary_c2, summary_c3, summary_c4 = st.columns(4)
                                summary_c1.metric("CUSIPs Found", f"{len(dd_opps):,}")
                                summary_c2.metric("Bucket", dd_bucket)
                                summary_c3.metric("Total Par Traded", f"{dd_opps['total_trade_amount'].sum():,.0f}")
                                summary_c4.metric("Benchmark", f"{dd_rating} {dd_tenor}")

                                top_row = dd_opps.iloc[0]
                                spread_change_text = ""
                                if pd.notna(top_row.get("spread_change_bps")):
                                    spread_change_text = f"{top_row.get('spread_change_bps'):+.1f} bp spread change, "

                                st.info(
                                    f"Top read-through by {dd_sort_label}: CUSIP {top_row['cusip']} shows "
                                    f"{top_row['current_spread_bps']:+.1f} bp current spread to {dd_rating}, "
                                    f"{spread_change_text}"
                                    f"{int(top_row['trade_count'])} trades, and {top_row['liquidity_tier']} liquidity in the selected window."
                                )

                                display_cols = [
                                    "cusip", "coupon", "maturity", "call_date", "call_price",
                                    "current_avg_yield", "current_spread_bps", "spread_change_bps",
                                    "yield_change_bps", "trade_count", "total_trade_amount", "avg_trade_amount",
                                    "latest_trade", "historical_trade_date", "historical_spread_bps",
                                    "avg_price", "historical_price", "liquidity_score", "liquidity_tier",
                                    "outstanding_amount",
                                ]
                                dd_display = dd_opps[[c for c in display_cols if c in dd_opps.columns]].copy()
                                for col in ["current_avg_yield", "current_spread_bps", "spread_change_bps", "yield_change_bps", "avg_price", "historical_price", "liquidity_score"]:
                                    if col in dd_display.columns:
                                        dd_display[col] = pd.to_numeric(dd_display[col], errors="coerce").round(2)

                                st.subheader("CUSIP Opportunity Table")
                                st.dataframe(dd_display, use_container_width=True, hide_index=True, height=420)

                                st.subheader("Security Detail")
                                selected_cusip = st.selectbox(
                                    "Select CUSIP for detail",
                                    dd_opps["cusip"].astype(str).tolist(),
                                    index=0,
                                    key="selected_cusip_drilldown",
                                )

                                sec_trades = dd_base[dd_base["cusip"].astype(str) == str(selected_cusip)].copy()
                                sec_trades = sec_trades.sort_values("trade_date")
                                if sec_trades.empty:
                                    st.warning("No trade rows found for the selected CUSIP.")
                                else:
                                    sec_daily = (
                                        sec_trades.groupby("trade_date", as_index=False)
                                        .agg(
                                            avg_yield=("yield", "mean"),
                                            trade_count=("yield", "count"),
                                            total_trade_amount=("trade_amount", "sum"),
                                            avg_price=("price", "mean"),
                                        )
                                    )

                                    # Build benchmark series for selected security dates.
                                    bench_long = make_benchmark_long(mmd_df, dd_rating)
                                    if not bench_long.empty:
                                        sec_daily = sec_daily.merge(
                                            bench_long[bench_long["maturity_bucket"] == dd_bucket][["trade_date", "benchmark_yield", "benchmark_source", "source_column"]],
                                            on="trade_date",
                                            how="left",
                                        )
                                        sec_daily["spread_to_benchmark_bps"] = (
                                            sec_daily["avg_yield"] - sec_daily["benchmark_yield"]
                                        ) * 100

                                    detail_col1, detail_col2 = st.columns(2)
                                    with detail_col1:
                                        sec_yield_fig = px.line(
                                            sec_daily,
                                            x="trade_date",
                                            y="avg_yield",
                                            markers=True,
                                            hover_data=["trade_count", "total_trade_amount", "avg_price"],
                                            title=f"{selected_cusip} Yield History",
                                            labels={"trade_date": "Trade Date", "avg_yield": "Average Yield (%)"},
                                        )
                                        sec_yield_fig.update_layout(height=380)
                                        st.plotly_chart(sec_yield_fig, use_container_width=True)

                                    with detail_col2:
                                        if "spread_to_benchmark_bps" in sec_daily.columns and sec_daily["spread_to_benchmark_bps"].notna().any():
                                            sec_spread_fig = px.line(
                                                sec_daily,
                                                x="trade_date",
                                                y="spread_to_benchmark_bps",
                                                markers=True,
                                                hover_data=["trade_count", "total_trade_amount", "benchmark_source", "source_column"],
                                                title=f"{selected_cusip} Spread to {dd_rating} Benchmark",
                                                labels={"trade_date": "Trade Date", "spread_to_benchmark_bps": "Spread (bps)"},
                                            )
                                            sec_spread_fig.update_layout(height=380)
                                            st.plotly_chart(sec_spread_fig, use_container_width=True)
                                        else:
                                            sec_amt_fig = px.bar(
                                                sec_daily,
                                                x="trade_date",
                                                y="total_trade_amount",
                                                hover_data=["trade_count", "avg_yield", "avg_price"],
                                                title=f"{selected_cusip} Trade Amount History",
                                                labels={"trade_date": "Trade Date", "total_trade_amount": "Total Trade Amount"},
                                            )
                                            sec_amt_fig.update_layout(height=380)
                                            st.plotly_chart(sec_amt_fig, use_container_width=True)

                                    with st.expander("Latest trades for selected CUSIP", expanded=False):
                                        latest_trade_cols = [
                                            "trade_datetime", "trade_date", "cusip", "description", "maturity_trade",
                                            "maturity_bond", "coupon_trade", "coupon_bond", "yield", "price",
                                            "trade_amount", "spread", "trade_type", "ratings_m_s_f",
                                        ]
                                        st.dataframe(
                                            sec_trades[[c for c in latest_trade_cols if c in sec_trades.columns]]
                                            .sort_values("trade_date", ascending=False)
                                            .head(500),
                                            use_container_width=True,
                                            hide_index=True,
                                        )

                            with st.expander("Drilldown benchmark/audit details", expanded=False):
                                st.markdown(
                                    f"""
- Latest CUSIP trade date used: **{dd_latest_date.strftime('%Y-%m-%d')}**
- Current window start: **{dd_current_start.strftime('%Y-%m-%d')}**
- Historical target date: **{dd_hist_target.strftime('%Y-%m-%d')}**
- Benchmark date: **{dd_benchmark_date.strftime('%Y-%m-%d')}**
- Benchmark source: **{dd_meta.get('benchmark_source')}**
- Source column: **{dd_meta.get('source_column')}**
- Benchmark yield: **{dd_benchmark_yield:.4f}%**
                                    """
                                )


section_anchor("rv-positioning", "Relative Value Positioning Map")
with st.expander("Methodology: relative value positioning map", expanded=False):
    st.markdown(
        """
This scatter plot maps individual CUSIPs by **tradability** and **relative value**.

**Default interpretation:**

- **X-axis = Liquidity Score**: higher means more actively traded, larger traded amount, more recent activity, and less staleness.
- **Y-axis = Spread to Benchmark**: higher means the bond is trading cheaper versus the selected benchmark curve.
- **Bubble size = Total Trade Amount**: larger dots indicate more secondary-market trading volume.
- **Color = Maturity Bucket**: Short / Intermediate / Long / Extended Long.

**Quadrants:**

- **Upper-right:** cheap and liquid; often the first area to investigate.
- **Upper-left:** cheap but illiquid; may require a liquidity premium.
- **Lower-right:** liquid but rich; useful benchmark-like bonds.
- **Lower-left:** illiquid and rich; usually less attractive from a relative-value screen.

This is a **screening view**, not an investment recommendation. It helps analysts identify bonds worth deeper review.
        """
    )

if issuer_trades.empty:
    st.warning("No trade rows found for this issuer and filter.")
else:
    rv_controls = st.columns([1, 1, 1, 1])
    with rv_controls[0]:
        rv_benchmark_rating = st.selectbox(
            "RV Benchmark Curve",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="rv_benchmark_rating",
            help="Used only when Y-axis is spread to benchmark. Uploaded curve columns are used first; otherwise MMD + assumption spread.",
        )
    with rv_controls[1]:
        rv_y_axis = st.selectbox(
            "Y-axis",
            ["Spread to Benchmark (bps)", "Average Yield (%)"],
            index=0,
            key="rv_y_axis",
        )
    with rv_controls[2]:
        rv_size_by = st.selectbox(
            "Bubble size",
            ["Total Trade Amount", "Outstanding Amount", "Trade Count"],
            index=0,
            key="rv_size_by",
        )
    with rv_controls[3]:
        rv_min_trades = st.number_input(
            "Minimum Trades",
            min_value=1,
            max_value=100,
            value=1,
            step=1,
            key="rv_min_trades",
        )

    rv_base = issuer_trades.copy()
    rv_base["trade_date"] = pd.to_datetime(rv_base["trade_date"], errors="coerce").dt.normalize()
    rv_base["yield"] = pd.to_numeric(rv_base["yield"], errors="coerce")
    if "trade_amount" in rv_base.columns:
        rv_base["trade_amount"] = pd.to_numeric(rv_base["trade_amount"], errors="coerce")
    else:
        rv_base["trade_amount"] = pd.NA
    if "price" in rv_base.columns:
        rv_base["price"] = pd.to_numeric(rv_base["price"], errors="coerce")
    else:
        rv_base["price"] = pd.NA

    rv_base = rv_base.dropna(subset=["cusip", "trade_date", "yield"])

    if rv_base.empty:
        st.warning("No usable CUSIP-level trade rows are available for the positioning map.")
    else:
        today_rv = pd.Timestamp.today().normalize()
        rv_base["trade_month"] = rv_base["trade_date"].dt.to_period("M").astype(str)
        rv_summary = (
            rv_base.groupby("cusip", dropna=False)
            .agg(
                avg_yield=("yield", "mean"),
                latest_yield=("yield", "last"),
                avg_price=("price", "mean"),
                trade_count=("trade_date", "count"),
                first_trade=("trade_date", "min"),
                latest_trade=("trade_date", "max"),
                active_months=("trade_month", "nunique"),
                total_trade_amount=("trade_amount", "sum"),
                avg_trade_amount=("trade_amount", "mean"),
                maturity_bucket=("maturity_bucket", "first"),
                maturity=("maturity_bond", "first"),
                coupon=("coupon_bond", "first"),
                outstanding_amount=("outstanding_amount", "first"),
                description=("description", "first") if "description" in rv_base.columns else ("yield", "count"),
            )
            .reset_index()
        )
        rv_summary["days_since_last_trade"] = (today_rv - rv_summary["latest_trade"]).dt.days
        rv_summary["trading_period_days"] = (rv_summary["latest_trade"] - rv_summary["first_trade"]).dt.days.clip(lower=1)
        rv_summary["avg_days_between_trades"] = rv_summary["trading_period_days"] / rv_summary["trade_count"].clip(lower=1)
        rv_summary["avg_trades_per_month"] = rv_summary["trade_count"] / rv_summary["active_months"].clip(lower=1)

        recent_cutoff_rv = today_rv - pd.DateOffset(days=90)
        rv_recent = (
            rv_base[rv_base["trade_date"] >= recent_cutoff_rv]
            .groupby("cusip")
            .agg(recent_90d_trades=("trade_date", "count"))
            .reset_index()
        )
        rv_summary = rv_summary.merge(rv_recent, on="cusip", how="left")
        rv_summary["recent_90d_trades"] = rv_summary["recent_90d_trades"].fillna(0).astype(int)

        for numeric_col in ["total_trade_amount", "outstanding_amount", "avg_trade_amount"]:
            if numeric_col in rv_summary.columns:
                rv_summary[numeric_col] = pd.to_numeric(rv_summary[numeric_col], errors="coerce")
        rv_summary["turnover_ratio"] = rv_summary["total_trade_amount"] / rv_summary["outstanding_amount"].replace({0: pd.NA})
        rv_summary["liquidity_score"] = (
            rv_summary["trade_count"].rank(pct=True) * 35
            + rv_summary["total_trade_amount"].fillna(0).rank(pct=True) * 25
            + rv_summary["recent_90d_trades"].rank(pct=True) * 25
            + (1 - rv_summary["days_since_last_trade"].rank(pct=True)) * 15
        )
        rv_summary["liquidity_tier"] = pd.cut(
            rv_summary["liquidity_score"],
            bins=[-1, 45, 75, 101],
            labels=["Low Liquidity", "Medium Liquidity", "High Liquidity"],
        ).astype(str)
        rv_summary.loc[rv_summary["days_since_last_trade"] > 365, "liquidity_tier"] = "Stale"

        rv_summary = rv_summary[rv_summary["trade_count"] >= rv_min_trades].copy()

        # Add benchmark spread at each CUSIP's latest trade date and maturity bucket.
        if rv_y_axis == "Spread to Benchmark (bps)":
            if mmd_df.empty:
                st.info("Upload an MMD / benchmark curve file to use Spread to Benchmark. Showing Average Yield instead.")
                rv_y_axis_col = "avg_yield"
                rv_y_axis_label = "Average Yield (%)"
            else:
                benchmark_long_rv = make_benchmark_long(mmd_df, rv_benchmark_rating)
                if benchmark_long_rv.empty:
                    st.info("No usable benchmark curve was found for the selected rating. Showing Average Yield instead.")
                    rv_y_axis_col = "avg_yield"
                    rv_y_axis_label = "Average Yield (%)"
                else:
                    benchmark_long_rv = benchmark_long_rv.sort_values(["maturity_bucket", "trade_date"])
                    merge_frames = []
                    for bucket in MATURITY_BUCKET_ORDER:
                        left = rv_summary[rv_summary["maturity_bucket"] == bucket].sort_values("latest_trade")
                        right = benchmark_long_rv[benchmark_long_rv["maturity_bucket"] == bucket].sort_values("trade_date")
                        if left.empty or right.empty:
                            continue
                        merged_bucket = pd.merge_asof(
                            left,
                            right,
                            left_on="latest_trade",
                            right_on="trade_date",
                            direction="backward",
                            tolerance=pd.Timedelta(days=14),
                        )
                        merge_frames.append(merged_bucket)
                    if merge_frames:
                        rv_summary = pd.concat(merge_frames, ignore_index=True)
                        rv_summary["spread_to_benchmark_bps"] = (
                            rv_summary["avg_yield"] - rv_summary["benchmark_yield"]
                        ) * 100
                        rv_y_axis_col = "spread_to_benchmark_bps"
                        rv_y_axis_label = "Spread to Benchmark (bps)"
                    else:
                        st.info("No overlapping CUSIP latest-trade dates and benchmark dates were found. Showing Average Yield instead.")
                        rv_y_axis_col = "avg_yield"
                        rv_y_axis_label = "Average Yield (%)"
        else:
            rv_y_axis_col = "avg_yield"
            rv_y_axis_label = "Average Yield (%)"

        rv_summary = rv_summary.dropna(subset=["liquidity_score", rv_y_axis_col])

        if rv_summary.empty:
            st.warning("No CUSIPs meet the selected filters for the positioning map.")
        else:
            size_map = {
                "Total Trade Amount": "total_trade_amount",
                "Outstanding Amount": "outstanding_amount",
                "Trade Count": "trade_count",
            }
            size_col = size_map.get(rv_size_by, "total_trade_amount")

            # Defensive plotting layer -------------------------------------------------
            # Plotly scatter is sensitive to missing/non-numeric/negative values in
            # size, x, and y columns. Muni exports often have blank outstanding amount,
            # missing trade amount, missing maturity dates, or unmatched benchmark values.
            #
            # We handle this in two steps:
            #   1) Clean numeric plotting inputs so the chart does not crash.
            #   2) Split known maturity buckets from unknown maturity buckets so
            #      "Unknown" does not dominate or pollute the main positioning map.
            rv_plot = rv_summary.copy()

            for numeric_col in ["liquidity_score", rv_y_axis_col, size_col]:
                if numeric_col in rv_plot.columns:
                    rv_plot[numeric_col] = pd.to_numeric(rv_plot[numeric_col], errors="coerce")
                    rv_plot[numeric_col] = rv_plot[numeric_col].replace([float("inf"), -float("inf")], pd.NA)

            required_plot_cols = ["liquidity_score", rv_y_axis_col]
            rv_plot = rv_plot.dropna(subset=[c for c in required_plot_cols if c in rv_plot.columns])

            # Resolve maturity bucket from common merge variants. If the bucket still
            # cannot be determined, keep the row for audit but exclude it from the
            # main scatter chart.
            valid_buckets = MATURITY_BUCKET_ORDER

            if "maturity_bucket" not in rv_plot.columns:
                possible_bucket_cols = [
                    "maturity_bucket_x",
                    "maturity_bucket_y",
                    "maturity_bucket_trade",
                    "maturity_bucket_bond",
                ]
                found_bucket_col = next((c for c in possible_bucket_cols if c in rv_plot.columns), None)
                if found_bucket_col:
                    rv_plot["maturity_bucket"] = rv_plot[found_bucket_col]
                else:
                    rv_plot["maturity_bucket"] = pd.NA

            rv_plot["maturity_bucket"] = rv_plot["maturity_bucket"].astype("string")

            if "cusip" not in rv_plot.columns:
                rv_plot["cusip"] = rv_plot.index.astype(str)
            else:
                rv_plot["cusip"] = rv_plot["cusip"].fillna("Unknown").astype(str)

            rv_known = rv_plot[rv_plot["maturity_bucket"].isin(valid_buckets)].copy()
            rv_unknown = rv_plot[~rv_plot["maturity_bucket"].isin(valid_buckets)].copy()

            # Clean the bubble-size column only on the known-bucket plotting set.
            if size_col not in rv_known.columns:
                rv_known["point_size"] = 10
                size_col = "point_size"
            else:
                rv_known[size_col] = pd.to_numeric(rv_known[size_col], errors="coerce")
                rv_known[size_col] = rv_known[size_col].replace([float("inf"), -float("inf")], pd.NA)
                rv_known[size_col] = rv_known[size_col].fillna(0).clip(lower=0)

                if rv_known[size_col].sum() <= 0:
                    rv_known["point_size"] = 10
                    size_col = "point_size"

            hover_cols = [
                "cusip", "maturity_bucket", "maturity", "coupon", "avg_yield", "avg_price",
                "trade_count", "recent_90d_trades", "days_since_last_trade", "total_trade_amount",
                "outstanding_amount", "turnover_ratio", "liquidity_tier",
            ]
            if "spread_to_benchmark_bps" in rv_known.columns:
                hover_cols.extend(["spread_to_benchmark_bps", "benchmark_yield", "benchmark_source", "source_column"])
            hover_cols = [c for c in hover_cols if c in rv_known.columns]

            if rv_plot.empty:
                st.warning(
                    "No valid observations remain after cleaning the positioning-map inputs. "
                    "Try lowering the minimum trade filter or using Average Yield instead of Spread to Benchmark."
                )
                rv_summary = rv_plot
                median_liquidity = pd.NA
                median_y = pd.NA

            elif rv_known.empty:
                st.warning(
                    "No bonds with known maturity buckets were available for the main positioning map. "
                    "Unknown-maturity bonds are listed below for audit."
                )
                rv_summary = rv_known
                median_liquidity = pd.NA
                median_y = pd.NA

            else:
                try:
                    rv_fig = px.scatter(
                        rv_known,
                        x="liquidity_score",
                        y=rv_y_axis_col,
                        size=size_col,
                        size_max=38,
                        color="maturity_bucket",
                        category_orders={"maturity_bucket": valid_buckets},
                        hover_name="cusip",
                        hover_data=hover_cols,
                        title=f"{selected_issuer} Relative Value Positioning Map",
                        labels={
                            "liquidity_score": "Liquidity Score",
                            rv_y_axis_col: rv_y_axis_label,
                            "maturity_bucket": "Maturity Bucket",
                            size_col: rv_size_by if size_col != "point_size" else "Fixed Point Size",
                        },
                    )
                    median_liquidity = rv_known["liquidity_score"].median()
                    median_y = rv_known[rv_y_axis_col].median()
                    if pd.notna(median_liquidity):
                        rv_fig.add_vline(x=median_liquidity, line_dash="dash", opacity=0.45)
                    if pd.notna(median_y):
                        rv_fig.add_hline(y=median_y, line_dash="dash", opacity=0.45)
                    rv_fig.update_layout(height=560, hovermode="closest")
                    st.plotly_chart(rv_fig, use_container_width=True)
                except Exception as exc:
                    st.warning(
                        "The positioning map could not be plotted because the scatter inputs were not usable. "
                        f"The cleaned known-maturity data table is shown below for review. Error: {exc}"
                    )
                    st.dataframe(rv_known.head(1000), use_container_width=True, hide_index=True)
                    median_liquidity = rv_known["liquidity_score"].median() if "liquidity_score" in rv_known.columns else pd.NA
                    median_y = rv_known[rv_y_axis_col].median() if rv_y_axis_col in rv_known.columns else pd.NA

                # Use the cleaned known-bucket plotting data for quadrant/read-through logic.
                rv_summary = rv_known

            # Unknown maturity bucket audit ------------------------------------------
            # These rows are not bad data; they are simply excluded from the main map
            # because the maturity bucket could not be determined from the uploaded
            # bond/trade data. Keeping them visible makes the dashboard transparent
            # without letting Unknown dominate the legend.
            if not rv_unknown.empty:
                with st.expander(
                    f"Unknown maturity bucket bonds excluded from main map ({len(rv_unknown):,})",
                    expanded=False,
                ):
                    st.caption(
                        "These CUSIPs were excluded from the main positioning map because their maturity bucket "
                        "could not be determined from the uploaded bond/trade data. They are retained here for audit."
                    )
                    unknown_display_cols = [
                        "cusip",
                        "avg_yield",
                        "spread_to_benchmark_bps",
                        "liquidity_score",
                        "trade_count",
                        "recent_90d_trades",
                        "days_since_last_trade",
                        "total_trade_amount",
                        "outstanding_amount",
                    ]
                    unknown_existing_cols = [c for c in unknown_display_cols if c in rv_unknown.columns]
                    st.dataframe(
                        rv_unknown[unknown_existing_cols].head(5000),
                        use_container_width=True,
                        hide_index=True,
                    )

            if (
                rv_y_axis_col == "spread_to_benchmark_bps"
                and not rv_summary.empty
                and pd.notna(median_liquidity)
                and pd.notna(median_y)
            ):
                candidates = rv_summary[
                    (rv_summary["liquidity_score"] >= median_liquidity)
                    & (rv_summary["spread_to_benchmark_bps"] >= median_y)
                ].sort_values(["spread_to_benchmark_bps", "liquidity_score"], ascending=False)
                if not candidates.empty:
                    top = candidates.iloc[0]
                    st.info(
                        f"Positioning read-through: {top['cusip']} screens as relatively cheap and liquid "
                        f"at {top['spread_to_benchmark_bps']:+.1f} bp versus {rv_benchmark_rating}, "
                        f"with a liquidity score of {top['liquidity_score']:.1f}."
                    )

            with st.expander("Positioning map data table", expanded=False):
                display_cols = [
                    "cusip", "maturity_bucket", "liquidity_score", "liquidity_tier", rv_y_axis_col,
                    "avg_yield", "benchmark_yield", "benchmark_source", "source_column", "trade_count",
                    "recent_90d_trades", "days_since_last_trade", "total_trade_amount", "outstanding_amount",
                    "turnover_ratio", "maturity", "coupon", "avg_price",
                ]
                # Keep only existing columns and remove duplicates.
                # In flexible trade-centric mode, rv_y_axis_col can equal a column
                # already listed above (for example, liquidity_score). Duplicate
                # column names make rv_display[c] return a DataFrame instead of a
                # Series, which breaks pd.to_numeric().
                display_cols = [c for c in display_cols if c in rv_summary.columns]
                display_cols = list(dict.fromkeys(display_cols))

                rv_display = rv_summary[display_cols].copy()
                rv_display = rv_display.loc[:, ~rv_display.columns.duplicated()].copy()

                numeric_cols = [
                    "liquidity_score",
                    rv_y_axis_col,
                    "avg_yield",
                    "benchmark_yield",
                    "turnover_ratio",
                    "avg_price",
                ]
                numeric_cols = list(dict.fromkeys([c for c in numeric_cols if c in rv_display.columns]))

                for c in numeric_cols:
                    try:
                        rv_display[c] = pd.to_numeric(rv_display[c], errors="coerce").round(2)
                    except Exception:
                        # Leave problematic optional columns unchanged instead of
                        # stopping the full dashboard.
                        pass

                sort_cols = [c for c in [rv_y_axis_col, "liquidity_score"] if c in rv_display.columns]
                if sort_cols:
                    rv_display = rv_display.sort_values(sort_cols, ascending=False)

                st.dataframe(
                    rv_display,
                    use_container_width=True,
                    hide_index=True,
                    height=420,
                )

section_anchor("liquidity", "Liquidity / Trading Frequency Analysis")
with st.expander("Methodology", expanded=False):
    st.write("Liquidity score is a transparent ranking measure: 35% trade count, 25% total trade amount, 25% recent 90-day trades, and 15% recency. It is a screening metric, not a credit rating or valuation recommendation.")
if issuer_trades.empty:
    st.warning("No trade rows found for this issuer and filter.")
else:
    today = pd.Timestamp.today().normalize()
    liq_base = issuer_trades.copy()
    liq_base["trade_month"] = liq_base["trade_date"].dt.to_period("M").astype(str)
    liq = (
        liq_base.groupby("cusip", dropna=False)
        .agg(
            trade_count=("trade_date", "count"),
            first_trade=("trade_date", "min"),
            latest_trade=("trade_date", "max"),
            active_months=("trade_month", "nunique"),
            avg_yield=("yield", "mean"),
            min_yield=("yield", "min"),
            max_yield=("yield", "max"),
            avg_price=("price", "mean"),
            total_trade_amount=("trade_amount", "sum"),
            avg_trade_amount=("trade_amount", "mean"),
            median_trade_amount=("trade_amount", "median"),
            maturity=("maturity_bond", "first"),
            coupon=("coupon_bond", "first"),
            outstanding_amount=("outstanding_amount", "first"),
        )
        .reset_index()
    )
    liq["days_since_last_trade"] = (today - liq["latest_trade"]).dt.days
    liq["trading_period_days"] = (liq["latest_trade"] - liq["first_trade"]).dt.days.clip(lower=1)
    liq["avg_days_between_trades"] = liq["trading_period_days"] / liq["trade_count"].clip(lower=1)
    liq["avg_trades_per_month"] = liq["trade_count"] / liq["active_months"].clip(lower=1)
    recent_cutoff = today - pd.DateOffset(days=90)
    recent = liq_base[liq_base["trade_date"] >= recent_cutoff].groupby("cusip").agg(recent_90d_trades=("trade_date", "count")).reset_index()
    liq = liq.merge(recent, on="cusip", how="left")
    liq["recent_90d_trades"] = liq["recent_90d_trades"].fillna(0).astype(int)
    liq["yield_range"] = liq["max_yield"] - liq["min_yield"]
    liq["turnover_ratio"] = liq["total_trade_amount"] / liq["outstanding_amount"].replace({0: pd.NA})
    liq["liquidity_score"] = (
        liq["trade_count"].rank(pct=True) * 35
        + liq["total_trade_amount"].rank(pct=True) * 25
        + liq["recent_90d_trades"].rank(pct=True) * 25
        + (1 - liq["days_since_last_trade"].rank(pct=True)) * 15
    )
    liq["liquidity_tier"] = pd.cut(
        liq["liquidity_score"], bins=[-1, 45, 75, 101], labels=["Low Liquidity", "Medium Liquidity", "High Liquidity"]
    ).astype(str)
    liq.loc[liq["days_since_last_trade"] > 365, "liquidity_tier"] = "Stale"
    liq = liq.sort_values(["liquidity_score", "trade_count", "total_trade_amount"], ascending=False)

    monthly = liq_base.groupby("trade_month", as_index=False).agg(trade_count=("trade_date", "count"), total_trade_amount=("trade_amount", "sum"), avg_yield=("yield", "mean"))
    st.subheader("1. Market Activity Over Time")
    st.plotly_chart(px.line(monthly, x="trade_month", y="trade_count", markers=True, title="Monthly Trade Count"), use_container_width=True)

    st.subheader("2. Trade Size Distribution")
    with st.expander("Methodology: trade size distribution", expanded=False):
        st.markdown(
            """
This chart groups trades by par/trade amount to show whether activity is primarily retail-sized, institutional-sized, or block-oriented.

**Default buckets:**

- **< $100k**: odd-lot / retail-sized activity
- **$100k–$250k**: small institutional or advisor-sized activity
- **$250k–$1mm**: institutional-sized activity
- **$1mm+**: block trade / larger institutional flow

This is useful because trade count alone can overstate liquidity when most activity comes from small trades.
            """
        )

    if "trade_amount" not in liq_base.columns:
        st.info("Trade size distribution is unavailable because trade_amount is missing from the uploaded trade data.")
    else:
        trade_size_df = liq_base.copy()
        trade_size_df["trade_amount"] = pd.to_numeric(trade_size_df["trade_amount"], errors="coerce")
        trade_size_df = trade_size_df.dropna(subset=["trade_amount"])
        trade_size_df = trade_size_df[trade_size_df["trade_amount"] > 0]

        if trade_size_df.empty:
            st.info("Trade size distribution is unavailable because no positive trade_amount values were found.")
        else:
            trade_size_bins = [0, 100_000, 250_000, 1_000_000, float("inf")]
            trade_size_labels = ["< $100k", "$100k–$250k", "$250k–$1mm", "$1mm+"]
            trade_size_df["trade_size_bucket"] = pd.cut(
                trade_size_df["trade_amount"],
                bins=trade_size_bins,
                labels=trade_size_labels,
                include_lowest=True,
                right=False,
            )

            size_summary = (
                trade_size_df.groupby("trade_size_bucket", observed=False)
                .agg(
                    trade_count=("trade_amount", "count"),
                    total_trade_amount=("trade_amount", "sum"),
                    avg_trade_amount=("trade_amount", "mean"),
                    median_trade_amount=("trade_amount", "median"),
                )
                .reset_index()
            )
            size_summary["trade_size_bucket"] = size_summary["trade_size_bucket"].astype(str)
            size_summary["trade_count_share"] = size_summary["trade_count"] / size_summary["trade_count"].sum()
            size_summary["amount_share"] = size_summary["total_trade_amount"] / size_summary["total_trade_amount"].sum()

            size_fig = px.bar(
                size_summary,
                x="trade_size_bucket",
                y="trade_count",
                hover_data={
                    "total_trade_amount": ":,.0f",
                    "avg_trade_amount": ":,.0f",
                    "median_trade_amount": ":,.0f",
                    "trade_count_share": ":.1%",
                    "amount_share": ":.1%",
                },
                title="Trade Count by Size Bucket",
                labels={
                    "trade_size_bucket": "Trade Size Bucket",
                    "trade_count": "Number of Trades",
                    "total_trade_amount": "Total Trade Amount",
                    "trade_count_share": "Share of Trades",
                    "amount_share": "Share of Par Traded",
                },
            )
            size_fig.update_layout(height=430)
            st.plotly_chart(size_fig, use_container_width=True)

            amount_fig = px.bar(
                size_summary,
                x="trade_size_bucket",
                y="total_trade_amount",
                hover_data={
                    "trade_count": ":,.0f",
                    "avg_trade_amount": ":,.0f",
                    "median_trade_amount": ":,.0f",
                    "trade_count_share": ":.1%",
                    "amount_share": ":.1%",
                },
                title="Total Par Traded by Size Bucket",
                labels={
                    "trade_size_bucket": "Trade Size Bucket",
                    "total_trade_amount": "Total Trade Amount",
                    "trade_count": "Number of Trades",
                    "trade_count_share": "Share of Trades",
                    "amount_share": "Share of Par Traded",
                },
            )
            amount_fig.update_layout(height=430)
            st.plotly_chart(amount_fig, use_container_width=True)

            retail_trade_share = size_summary.loc[
                size_summary["trade_size_bucket"] == "< $100k", "trade_count_share"
            ]
            block_amount_share = size_summary.loc[
                size_summary["trade_size_bucket"] == "$1mm+", "amount_share"
            ]

            retail_trade_share_val = float(retail_trade_share.iloc[0]) if not retail_trade_share.empty else 0.0
            block_amount_share_val = float(block_amount_share.iloc[0]) if not block_amount_share.empty else 0.0

            if retail_trade_share_val >= 0.60 and block_amount_share_val < 0.25:
                st.info(
                    f"Read-through: trading activity appears retail / odd-lot heavy. "
                    f"< $100k trades account for {retail_trade_share_val:.1%} of trades, "
                    f"while $1mm+ blocks account for {block_amount_share_val:.1%} of par traded."
                )
            elif block_amount_share_val >= 0.50:
                st.info(
                    f"Read-through: activity appears institutionally active. "
                    f"$1mm+ blocks account for {block_amount_share_val:.1%} of par traded."
                )
            else:
                st.info(
                    f"Read-through: trade activity is mixed across retail-sized and institutional-sized buckets. "
                    f"< $100k trades account for {retail_trade_share_val:.1%} of trades; "
                    f"$1mm+ blocks account for {block_amount_share_val:.1%} of par traded."
                )

            with st.expander("Trade size distribution table", expanded=False):
                table_display = size_summary.copy()
                for pct_col in ["trade_count_share", "amount_share"]:
                    table_display[pct_col] = table_display[pct_col].map(lambda x: f"{x:.1%}" if pd.notna(x) else "")
                for amt_col in ["total_trade_amount", "avg_trade_amount", "median_trade_amount"]:
                    table_display[amt_col] = pd.to_numeric(table_display[amt_col], errors="coerce").round(0)
                st.dataframe(table_display, use_container_width=True, hide_index=True)

    st.subheader("3. Most Frequently Traded CUSIPs")
    st.plotly_chart(px.bar(liq.head(25), x="cusip", y="trade_count", color="liquidity_tier", title="Top 25 Most Frequently Traded CUSIPs"), use_container_width=True)

    st.subheader("4. Trade Recency / Staleness")
    st.plotly_chart(px.histogram(liq, x="days_since_last_trade", nbins=30, color="liquidity_tier", title="Distribution of Days Since Last Trade"), use_container_width=True)

    st.subheader("5. Liquidity Ranking Table")
    display_cols = [
        "cusip", "liquidity_tier", "liquidity_score", "trade_count", "recent_90d_trades", "active_months",
        "avg_trades_per_month", "avg_days_between_trades", "days_since_last_trade", "first_trade", "latest_trade",
        "avg_yield", "yield_range", "avg_price", "total_trade_amount", "avg_trade_amount", "turnover_ratio",
        "maturity", "coupon", "outstanding_amount",
    ]
    st.dataframe(liq[[c for c in display_cols if c in liq.columns]], use_container_width=True, height=500)

section_anchor("bond-master", "Bond Master / Security Reference")
bond_cols = ["issuer", "sector", "primary_type", "election", "series", "cusip", "secondary_credit", "term", "maturity", "par_amount", "outstanding_amount", "coupon", "call_date", "call_price", "fed_tax", "amt"]
st.dataframe(issuer_bonds[[c for c in bond_cols if c in issuer_bonds.columns]].sort_values(["maturity", "cusip"]), use_container_width=True)

section_anchor("trade-detail", "Underlying Trade Detail")
trade_cols = ["trade_datetime", "cusip", "description", "maturity_trade", "maturity_bond", "maturity_bucket", "coupon_trade", "yield", "price", "trade_amount", "spread", "trade_type", "ratings_m_s_f"]
st.dataframe(issuer_trades[[c for c in trade_cols if c in issuer_trades.columns]].sort_values("trade_datetime", ascending=False).head(20000), use_container_width=True)



section_anchor("report-export-center", "Report Export Center")
with st.expander("Methodology: report export center", expanded=False):
    st.markdown(
        """
This section creates exportable reporting packages from the current dashboard state.

**Export options:**

- **Interactive HTML report:** includes selected charts, summary metrics, methodology notes, and key tables.
- **PDF summary:** creates a lightweight PDF using `reportlab` when available. For a full visual PDF, download the HTML report and use browser **Print → Save as PDF**.
- **PowerPoint slides:** creates a simple presentation using `python-pptx` when available.
- **Chart HTML bundle:** exports selected Plotly charts as standalone HTML files inside a ZIP.
- **Chart data CSV bundle:** exports the underlying data used to generate selected charts.

**Important limitation:**

A Streamlit app cannot reliably export the exact live browser page, all expanded/collapsed states, and all interactive widgets into a perfect PDF/PPTX without a browser automation service. This module therefore exports a clean, reproducible report built from the uploaded data and current issuer selection.
        """
    )

# -----------------------------
# Build exportable chart/data objects
# -----------------------------
export_chart_items = []
export_data_items = {}

def add_export_chart(name: str, fig, data: pd.DataFrame | None = None):
    """Collect charts for HTML/ZIP/PPT export without breaking if one chart fails."""
    try:
        export_chart_items.append((name, fig))
        if data is not None and not data.empty:
            export_data_items[name] = data.copy()
    except Exception:
        pass

# 1) Yield trend chart
try:
    export_yield_df = market_df[market_df["issuer"] == selected_issuer].copy()
    export_yield_df["trade_date"] = pd.to_datetime(export_yield_df["trade_date"], errors="coerce")
    export_yield_df["yield"] = pd.to_numeric(export_yield_df["yield"], errors="coerce")
    export_yield_df = export_yield_df.dropna(subset=["trade_date", "yield"])
    if not export_yield_df.empty:
        export_yield_daily = (
            export_yield_df.groupby("trade_date", as_index=False)
            .agg(avg_yield=("yield", "mean"), trade_count=("yield", "count"), total_trade_amount=("trade_amount", "sum") if "trade_amount" in export_yield_df.columns else ("yield", "count"))
            .sort_values("trade_date")
        )
        export_yield_fig = px.line(
            export_yield_daily,
            x="trade_date",
            y="avg_yield",
            markers=True,
            title=f"{selected_issuer} Average Trade Yield",
            labels={"trade_date": "Trade Date", "avg_yield": "Average Yield (%)"},
        )
        add_export_chart("yield_trend", export_yield_fig, export_yield_daily)
except Exception:
    pass

# 2) Issuer curve chart
try:
    export_curve_df = market_df[market_df["issuer"] == selected_issuer].copy()
    export_curve_df["trade_date"] = pd.to_datetime(export_curve_df["trade_date"], errors="coerce").dt.normalize()
    export_curve_df["yield"] = pd.to_numeric(export_curve_df["yield"], errors="coerce")
    export_curve_df = export_curve_df.dropna(subset=["trade_date", "yield", "maturity_bucket"])
    export_curve_df = export_curve_df[export_curve_df["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)].copy()
    if not export_curve_df.empty:
        latest_curve_date = export_curve_df["trade_date"].max()
        export_curve_df = export_curve_df[export_curve_df["trade_date"] >= latest_curve_date - pd.Timedelta(days=30)]
        export_curve_summary = (
            export_curve_df.groupby("maturity_bucket", as_index=False)
            .agg(avg_yield=("yield", "mean"), trade_count=("yield", "count"))
        )
        export_curve_summary["maturity_bucket"] = pd.Categorical(
            export_curve_summary["maturity_bucket"],
            categories=MATURITY_BUCKET_ORDER,
            ordered=True,
        )
        export_curve_summary = export_curve_summary.sort_values("maturity_bucket")
        export_curve_fig = px.line(
            export_curve_summary,
            x="maturity_bucket",
            y="avg_yield",
            markers=True,
            title=f"{selected_issuer} Issuer Curve — Latest 30D Average",
            labels={"maturity_bucket": "Maturity Bucket", "avg_yield": "Average Yield (%)"},
        )
        add_export_chart("issuer_curve_latest_30d", export_curve_fig, export_curve_summary)
except Exception:
    pass

# 3) Current spread level heatmap
try:
    if not mmd_df.empty:
        export_level_matrix, export_level_audit = build_spread_level_data(
            market_df=market_df,
            mmd_df=mmd_df,
            issuer=selected_issuer,
            ratings=["AAA", "AA", "A", "BBB"],
        )
        if not export_level_matrix.empty and not export_level_matrix.isna().all().all():
            export_level_text = export_level_matrix.map(lambda x: "" if pd.isna(x) else f"{x:+.1f} bp")
            export_level_fig = px.imshow(
                export_level_matrix.astype(float),
                x=export_level_matrix.columns,
                y=export_level_matrix.index,
                color_continuous_scale=["#1a9850", "#f7f7f7", "#d73027"],
                color_continuous_midpoint=0,
                aspect="auto",
                title=f"{selected_issuer} Current Spread Level",
                labels={"x": "Benchmark Curve", "y": "Maturity Bucket", "color": "Spread (bps)"},
            )
            export_level_fig.update_traces(text=export_level_text.values, texttemplate="%{text}")
            add_export_chart("current_spread_level_heatmap", export_level_fig, export_level_audit)
except Exception:
    pass

# 4) Liquidity monthly activity chart
try:
    export_liq_df = market_df[market_df["issuer"] == selected_issuer].copy()
    export_liq_df["trade_date"] = pd.to_datetime(export_liq_df["trade_date"], errors="coerce")
    export_liq_df = export_liq_df.dropna(subset=["trade_date"])
    if "trade_amount" in export_liq_df.columns:
        export_liq_df["trade_amount"] = pd.to_numeric(export_liq_df["trade_amount"], errors="coerce").fillna(0)
    else:
        export_liq_df["trade_amount"] = 0.0
    if not export_liq_df.empty:
        export_liq_df["trade_month"] = export_liq_df["trade_date"].dt.to_period("M").astype(str)
        export_monthly = (
            export_liq_df.groupby("trade_month", as_index=False)
            .agg(trade_count=("trade_date", "count"), total_trade_amount=("trade_amount", "sum"))
        )
        export_monthly_fig = px.line(
            export_monthly,
            x="trade_month",
            y="trade_count",
            markers=True,
            title=f"{selected_issuer} Monthly Trade Count",
            labels={"trade_month": "Trade Month", "trade_count": "Trade Count"},
        )
        add_export_chart("monthly_trade_count", export_monthly_fig, export_monthly)
except Exception:
    pass

# -----------------------------
# Export controls
# -----------------------------
export_options = st.multiselect(
    "Select charts to include",
    [name for name, _fig in export_chart_items],
    default=[name for name, _fig in export_chart_items],
    help="These are reconstructed export charts based on current selected issuer and uploaded data.",
)

selected_export_charts = [(name, fig) for name, fig in export_chart_items if name in export_options]

export_meta = {
    "Generated": pd.Timestamp.now().strftime("%Y-%m-%d %H:%M"),
    "Selected Issuer": selected_issuer,
    "Sector": selected_sector,
    "Bonds": f"{len(issuer_bonds):,}",
    "Trades in Current Filter": f"{len(issuer_trades):,}",
    "Latest Trade": issuer_trades["trade_date"].max().strftime("%Y-%m-%d") if not issuer_trades.empty else "No trades",
}

report_html_parts = [
    "<html><head><meta charset='utf-8'><title>Municipal Secondary Market Dashboard Report</title>",
    "<style>body{font-family:Arial,sans-serif;margin:32px;color:#111827;} h1,h2{color:#111827;} table{border-collapse:collapse;width:100%;margin:16px 0;} td,th{border:1px solid #e5e7eb;padding:8px;text-align:left;} .note{color:#64748b;font-size:13px;}</style>",
    "</head><body>",
    "<h1>Municipal Secondary Market Dashboard Report</h1>",
    "<h2>Executive Summary</h2>",
    "<table>",
]
for k, v in export_meta.items():
    report_html_parts.append(f"<tr><th>{k}</th><td>{v}</td></tr>")
report_html_parts.extend([
    "</table>",
    "<p class='note'>Benchmark curves use uploaded rating curves when available; otherwise the app falls back to MMD/AAA plus transparent spread assumptions. Screening outputs are not investment recommendations.</p>",
])

for name, fig in selected_export_charts:
    report_html_parts.append(f"<h2>{name.replace('_', ' ').title()}</h2>")
    report_html_parts.append(fig.to_html(full_html=False, include_plotlyjs="cdn"))

report_html_parts.append("</body></html>")
full_report_html = "\n".join(report_html_parts)

export_col1, export_col2, export_col3 = st.columns(3)

with export_col1:
    st.download_button(
        label="Download Interactive HTML Report",
        data=full_report_html.encode("utf-8"),
        file_name=f"{selected_issuer}_dashboard_report.html".replace(" ", "_"),
        mime="text/html",
        help="Open this file in a browser. For a visual PDF, use browser Print → Save as PDF.",
    )

with export_col2:
    # Lightweight PDF summary via reportlab, if installed.
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib import colors

        pdf_buffer = io.BytesIO()
        doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = [
            Paragraph("Municipal Secondary Market Dashboard Summary", styles["Title"]),
            Spacer(1, 12),
        ]
        meta_table = Table([[k, str(v)] for k, v in export_meta.items()])
        meta_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.lightgrey),
            ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        story.append(meta_table)
        story.append(Spacer(1, 14))
        story.append(Paragraph("Included Charts", styles["Heading2"]))
        for name, _fig in selected_export_charts:
            story.append(Paragraph(f"• {name.replace('_', ' ').title()}", styles["BodyText"]))
        story.append(Spacer(1, 14))
        story.append(Paragraph("Methodology Note", styles["Heading2"]))
        story.append(Paragraph(
            "This PDF is a lightweight summary. For interactive charts and fuller visual output, use the HTML report and browser Print → Save as PDF.",
            styles["BodyText"],
        ))
        doc.build(story)
        pdf_bytes = pdf_buffer.getvalue()

        st.download_button(
            label="Download PDF Summary",
            data=pdf_bytes,
            file_name=f"{selected_issuer}_dashboard_summary.pdf".replace(" ", "_"),
            mime="application/pdf",
        )
    except Exception:
        st.info("PDF summary export requires `reportlab`. Add `reportlab` to requirements.txt, or download HTML and print/save as PDF.")

with export_col3:
    # PPTX summary via python-pptx, if installed.
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Municipal Secondary Market Dashboard"
        slide.placeholders[1].text = f"{selected_issuer} | {selected_sector} | Generated {export_meta['Generated']}"

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive Snapshot"
        body = slide.placeholders[1].text_frame
        body.clear()
        for k, v in export_meta.items():
            p = body.add_paragraph()
            p.text = f"{k}: {v}"
            p.font.size = Pt(18)

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Included Dashboard Charts"
        body = slide.placeholders[1].text_frame
        body.clear()
        for name, _fig in selected_export_charts:
            p = body.add_paragraph()
            p.text = name.replace("_", " ").title()
            p.level = 0
            p.font.size = Pt(18)

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Methodology Notes"
        body = slide.placeholders[1].text_frame
        body.clear()
        notes = [
            "Benchmark curves use uploaded rating curves when available.",
            "Fallback curves use MMD/AAA plus visible rating-spread assumptions.",
            "Liquidity and RV scores are screening tools, not trade recommendations.",
            "Scenario shock uses duration proxies, not full cash-flow pricing.",
        ]
        for note in notes:
            p = body.add_paragraph()
            p.text = note
            p.font.size = Pt(16)

        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_bytes = pptx_buffer.getvalue()

        st.download_button(
            label="Download PowerPoint Slides",
            data=pptx_bytes,
            file_name=f"{selected_issuer}_dashboard_slides.pptx".replace(" ", "_"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception:
        st.info("PowerPoint export requires `python-pptx`. Add `python-pptx` to requirements.txt to enable slide export.")

# Chart HTML bundle and chart data bundle
bundle_col1, bundle_col2 = st.columns(2)

with bundle_col1:
    try:
        import zipfile
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", compression=zipfile.ZIP_DEFLATED) as zf:
            for name, fig in selected_export_charts:
                zf.writestr(f"{name}.html", fig.to_html(full_html=True, include_plotlyjs="cdn"))
        st.download_button(
            label="Download Chart HTML Bundle",
            data=zip_buffer.getvalue(),
            file_name=f"{selected_issuer}_chart_html_bundle.zip".replace(" ", "_"),
            mime="application/zip",
        )
    except Exception as exc:
        st.info(f"Chart bundle export unavailable: {exc}")

with bundle_col2:
    try:
        import zipfile
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", compression=zipfile.ZIP_DEFLATED) as zf:
            for name, df in export_data_items.items():
                if name in export_options and df is not None and not df.empty:
                    zf.writestr(f"{name}_data.csv", df.to_csv(index=False))
        st.download_button(
            label="Download Chart Data CSV Bundle",
            data=zip_buffer.getvalue(),
            file_name=f"{selected_issuer}_chart_data_bundle.zip".replace(" ", "_"),
            mime="application/zip",
        )
    except Exception as exc:
        st.info(f"Chart data bundle export unavailable: {exc}")

with st.expander("How to export the full live webpage as PDF", expanded=False):
    st.markdown(
        """
For the exact live Streamlit page:

1. Open the dashboard in your browser.
2. Expand the sections you want included.
3. Press **Cmd+P** on Mac or **Ctrl+P** on Windows.
4. Choose **Save as PDF**.
5. Set scale to 70–85% if charts are too wide.

For a cleaner report with reproducible charts, use **Download Interactive HTML Report** above.
        """
    )


section_anchor("export-summary", "Export Summary Package")
with st.expander("Methodology: export summary package", expanded=False):
    st.markdown(
        """
This section generates a lightweight export package that can be copied into internal updates, pitchbook drafts, or meeting notes.

**Current implementation:**

- Generates a Markdown summary and an HTML summary.
- Uses the selected issuer, selected sector, counts, latest trade date, and available dashboard outputs.
- For PDF export, open the HTML file in a browser and print/save as PDF.
- For PowerPoint, copy the HTML/Markdown summary into your deck and insert key charts from the dashboard.

This avoids adding fragile report-generation dependencies while keeping the workflow practical for internal use.
        """
    )

latest_trade_text = issuer_trades["trade_date"].max().strftime("%Y-%m-%d") if not issuer_trades.empty else "No trades"
summary_timestamp = pd.Timestamp.now().strftime("%Y-%m-%d %H:%M")
summary_lines = [
    f"# Municipal Secondary Market Dashboard Summary",
    "",
    f"**Generated:** {summary_timestamp}",
    f"**Selected Issuer:** {selected_issuer}",
    f"**Sector:** {selected_sector}",
    f"**Bonds:** {len(issuer_bonds):,}",
    f"**Trades in Current Filter:** {len(issuer_trades):,}",
    f"**Latest Trade:** {latest_trade_text}",
    "",
    "## Key Dashboard Modules",
    "- Yield Trend / Relative Value Comparison",
    "- Issuer Curve vs Benchmark Curve",
    "- Current Spread Level Framework",
    "- Peer and Cross-Issuer Relative Value",
    "- Historical Spread Percentile",
    "- Security Screener",
    "- Recommendation Narrative",
    "- Scenario Shock Analysis",
    "- CUSIP Opportunity Drilldown",
    "",
    "## Notes",
    "- Benchmark curves use uploaded rating curves when available; otherwise the dashboard falls back to MMD/AAA plus visible spread assumptions.",
    "- Liquidity, RV score, dealer proxy, and scenario shock outputs are screening tools, not final trade recommendations.",
]

if "data_quality_score" in locals():
    summary_lines.extend([
        "",
        "## Data Quality",
        f"- Data Quality Score: {data_quality_score:.1f}/100",
        f"- CUSIP Match Rate: {cusip_match_rate:.1f}%",
        f"- Known Maturity Bucket Rate: {known_bucket_rate:.1f}%",
        f"- Duplicates Removed: {duplicates_removed:,}",
    ])

summary_md = "\n".join(summary_lines)
summary_html = summary_md.replace("\n", "<br>")

export_c1, export_c2 = st.columns(2)
with export_c1:
    st.download_button(
        label="Download Markdown Summary",
        data=summary_md.encode("utf-8"),
        file_name=f"{selected_issuer}_dashboard_summary.md".replace(" ", "_"),
        mime="text/markdown",
    )
with export_c2:
    st.download_button(
        label="Download HTML Summary",
        data=f"<html><body>{summary_html}</body></html>".encode("utf-8"),
        file_name=f"{selected_issuer}_dashboard_summary.html".replace(" ", "_"),
        mime="text/html",
    )

with st.expander("Preview export summary", expanded=True):
    st.markdown(summary_md)

section_anchor("admin-methodology", "Admin Methodology Page")
st.markdown(
    """
This page centralizes the assumptions used throughout the dashboard so the tool is easier to hand off and maintain.

### Benchmark Curves
- AAA = uploaded MMD / AAA curve.
- Non-AAA curves use uploaded rating-specific columns when available.
- If not available, non-AAA curves are modeled as MMD/AAA + visible rating-spread assumptions.

### Liquidity Score
Uses trade count, total trade amount, recent activity, and recency. It is a screening score, not a credit rating.

### Relative Value Score
Combines spread percentile, liquidity percentile, and trade activity percentile to identify cheap + tradable candidates.

### Dealer Behavior Proxy
Only enabled when trade side/type data exists. It estimates buy/sell imbalance but does not represent true dealer inventory.

### Scenario Shock
Uses duration proxies by maturity bucket. It does not model full cash flows, convexity, OAS, tax status, or callable optionality.

### Recommendation Narrative
Rule-based and explainable. Each phrase is triggered by spread movement, historical percentile, liquidity, peer gap, or flow proxy thresholds.
"""
)

with st.expander("Rating spread assumptions", expanded=False):
    st.dataframe(rating_spread_table(), use_container_width=True, hide_index=True)

with st.expander("Duration proxy assumptions", expanded=False):
    duration_proxy_df = pd.DataFrame(
        [
            {"Maturity Bucket": "Short", "Proxy Duration": 2.0},
            {"Maturity Bucket": "Intermediate", "Proxy Duration": 8.0},
            {"Maturity Bucket": "Long", "Proxy Duration": 13.0},
            {"Maturity Bucket": "Extended Long", "Proxy Duration": 18.0},
        ]
    )
    st.dataframe(duration_proxy_df, use_container_width=True, hide_index=True)

section_anchor("version-changelog", "Version / Change Log")
version_rows = [
    {"Version": "v1.0-team-ready", "Change": "Stabilized data validation, benchmark framework, relative value analytics, and team-readiness modules."},
    {"Version": "v1.1", "Change": "Added Cross-Issuer RV Analytics, Scenario Shock, Recommendation Narrative, and CUSIP Drilldown."},
    {"Version": "v1.2", "Change": "Added Data Quality Scorecard, Export Summary Package, Admin Methodology Page, and Watchlist."},
]
st.dataframe(pd.DataFrame(version_rows), use_container_width=True, hide_index=True)
st.caption("Update this changelog whenever the team changes methodology, assumptions, or major modules.")


section_anchor("downloads", "Download Outputs")
d1, d2, d3 = st.columns(3)
with d1:
    dataframe_download_button(market_df, "Download Merged Market Data CSV", "merged_market_data.csv")
with d2:
    dataframe_download_button(issuer_master, "Download Issuer Master CSV", "issuer_master.csv")
with d3:
    dataframe_download_button(bonds_df, "Download Cleaned Bonds CSV", "cleaned_bonds.csv")

if show_raw_tables:
    st.header("Raw / Processed Tables")
    st.subheader("Issuer Master")
    st.dataframe(issuer_master, use_container_width=True)
    st.subheader("All Bonds")
    st.dataframe(bonds_df, use_container_width=True)
    st.subheader("All Trades")
    st.dataframe(trades_df.head(20000), use_container_width=True)
    st.subheader("Merged Market Data")
    st.dataframe(market_df.head(20000), use_container_width=True)
