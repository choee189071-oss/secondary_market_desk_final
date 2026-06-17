from __future__ import annotations

import html
import io
import re
import zipfile

import pandas as pd

from engine.methodology import methodology_trust_layers
from reports.reviewer_handoff import reviewer_handoff_markdown


def _make_unique_columns(columns) -> list[str]:
    counts: dict[str, int] = {}
    out = []
    for col in columns:
        base = str(col)
        count = counts.get(base, 0)
        out.append(base if count == 0 else f"{base}.{count}")
        counts[base] = count + 1
    return out


def prepare_report_dataframe(df: pd.DataFrame, max_rows: int | None = None) -> pd.DataFrame:
    if df is None:
        return pd.DataFrame()
    if not isinstance(df, pd.DataFrame):
        try:
            df = pd.DataFrame(df)
        except Exception:
            return pd.DataFrame({"value": [str(df)]})

    out = df.copy()
    out.columns = _make_unique_columns(out.columns)
    if max_rows is not None and len(out) > max_rows:
        out = out.head(int(max_rows)).copy()
    for col in out.columns:
        try:
            if pd.api.types.is_datetime64_any_dtype(out[col]):
                out[col] = pd.to_datetime(out[col], errors="coerce").dt.strftime("%m/%d/%Y")
            elif "date" in str(col).lower() or "maturity" in str(col).lower():
                converted = pd.to_datetime(out[col], errors="coerce")
                if len(converted) == 0 or converted.notna().mean() >= 0.6:
                    out[col] = converted.dt.strftime("%m/%d/%Y")
        except Exception:
            pass
    return out


def _fmt_bps(x, digits: int = 1) -> str:
    try:
        if pd.isna(x):
            return "N/A"
        return f"{float(x):+.{digits}f} bp"
    except Exception:
        return "N/A"


def _fmt_num(x, digits: int = 1) -> str:
    try:
        if pd.isna(x):
            return "N/A"
        return f"{float(x):,.{digits}f}"
    except Exception:
        return "N/A"


def _html_escape(value: object) -> str:
    return html.escape("" if value is None else str(value), quote=True)


def focused_report_value(value: object, empty: str = "N/A") -> str:
    if value is None:
        return empty
    try:
        if pd.isna(value):
            return empty
    except Exception:
        pass
    return str(value)


def focused_report_markdown_table(df: pd.DataFrame, max_rows: int = 20) -> str:
    """Dependency-free Markdown table writer; pandas.to_markdown needs tabulate."""
    if df is None or df.empty:
        return "_No rows available._"
    display_df = prepare_report_dataframe(df, max_rows=max_rows)
    if display_df.empty:
        return "_No rows available._"

    def clean_cell(value: object) -> str:
        text = focused_report_value(value, "")
        text = text.replace("\n", " ").replace("|", "\\|")
        return text

    cols = [str(c) for c in display_df.columns]
    lines = [
        "| " + " | ".join(cols) + " |",
        "| " + " | ".join(["---"] * len(cols)) + " |",
    ]
    for _, row in display_df.iterrows():
        lines.append("| " + " | ".join(clean_cell(row.get(c)) for c in display_df.columns) + " |")
    if len(df) > len(display_df):
        lines.append(f"\n_Showing first {len(display_df):,} of {len(df):,} rows._")
    return "\n".join(lines)


def focused_report_html_table(df: pd.DataFrame, max_rows: int = 30) -> str:
    if df is None or df.empty:
        return "<p class='muted'>No rows available.</p>"
    display_df = prepare_report_dataframe(df, max_rows=max_rows)
    if display_df.empty:
        return "<p class='muted'>No rows available.</p>"
    note = ""
    if len(df) > len(display_df):
        note = f"<p class='muted'>Showing first {len(display_df):,} of {len(df):,} rows.</p>"
    return display_df.to_html(index=False, classes="report-table", border=0, escape=True) + note


def focused_report_filename(label: str, suffix: str) -> str:
    safe_label = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(label)).strip("_") or "secondary_market_report"
    return f"{safe_label}_{suffix}"


def _context_trust_layers(context: dict) -> dict[str, pd.DataFrame]:
    existing = context.get("methodology_trust_layers")
    if isinstance(existing, dict):
        return existing
    return methodology_trust_layers(
        benchmark_source_mode=str(context.get("benchmark_source_mode", "")),
        benchmark_priority=str(context.get("benchmark_priority", "")),
        benchmark_conflict_policy=str(context.get("benchmark_conflict_policy", "")),
    )


def focused_methodology_appendix(
    mmd_df: pd.DataFrame,
    benchmark_source_mode: str,
    benchmark_priority: str,
    benchmark_conflict_policy: str,
) -> pd.DataFrame:
    return pd.DataFrame(
        [
            {
                "Section": "Benchmark",
                "Policy": "Uploaded MMD role",
                "Current treatment": "Uploaded MMD is treated as the AAA benchmark curve when external MMD fallback is active.",
                "Audit evidence": f"{len(mmd_df) if isinstance(mmd_df, pd.DataFrame) else 0:,} MMD row(s) loaded.",
            },
            {
                "Section": "Benchmark",
                "Policy": "Active benchmark source",
                "Current treatment": benchmark_source_mode,
                "Audit evidence": benchmark_priority,
            },
            {
                "Section": "Benchmark",
                "Policy": "Conflict policy",
                "Current treatment": benchmark_conflict_policy,
                "Audit evidence": "Trade-sheet index rates and uploaded MMD are not blended in the same run.",
            },
            {
                "Section": "Spread",
                "Policy": "Spread calculation",
                "Current treatment": "Issuer yield minus active benchmark yield, shown in basis points.",
                "Audit evidence": "Uses trade-sheet spread when present; otherwise derives from yield and index_rate.",
            },
            {
                "Section": "Peer grouping",
                "Policy": "Ratings fallback",
                "Current treatment": "Ratings can guide peer grouping when available; missing ratings fall back to sector and maturity.",
                "Audit evidence": "Rating effects remain separate from the benchmark spread.",
            },
            {
                "Section": "Attribution",
                "Policy": "Callable, liquidity, and sector effects",
                "Current treatment": "Displayed as separate attribution layers rather than embedded into benchmark spread.",
                "Audit evidence": "This preserves benchmark explainability.",
            },
            {
                "Section": "Liquidity",
                "Policy": "Liquidity score",
                "Current treatment": "Screening score based on trade count, total par amount, and recency.",
                "Audit evidence": "Higher score means more observable trading support in the uploaded file.",
            },
            {
                "Section": "RV",
                "Policy": "RV score",
                "Current treatment": "Screening score combining spread rank and liquidity rank.",
                "Audit evidence": "Output is a shortlist aid, not an investment recommendation.",
            },
            {
                "Section": "Watchlist",
                "Policy": "Saved candidates",
                "Current treatment": "Saved CUSIPs and notes are stored in the active Streamlit session.",
                "Audit evidence": "Download watchlist CSV/Markdown or report bundle before clearing the session.",
            },
            {
                "Section": "Export",
                "Policy": "Report reproducibility",
                "Current treatment": "Reports are regenerated from the current uploaded data, issuer selection, and saved watchlist.",
                "Audit evidence": "Interactive Streamlit widget state is summarized, not pixel-copied.",
            },
        ]
    )


def focused_core_chart_explanations(selected_issuer: str, selected_sector: str, benchmark_source_mode: str) -> pd.DataFrame:
    return pd.DataFrame(
        [
            {
                "Chart": "Spread Trend",
                "Question answered": "Is the issuer trading wider or tighter over time?",
                "How to read": f"Primary line tracks {selected_issuer} median spread by selected frequency; reference lines compare against {selected_sector or 'sector'} and uploaded-universe medians.",
                "Methodology note": f"Spread source follows active benchmark policy: {benchmark_source_mode}.",
            },
            {
                "Chart": "Volume & Activity",
                "Question answered": "Is observed trading support strong enough to trust the signal?",
                "How to read": "Bars show par traded; lower panel shows trade count and selected issuer share of uploaded activity.",
                "Methodology note": "Volume supports liquidity confidence but does not change benchmark spread.",
            },
            {
                "Chart": "Issuer Curve",
                "Question answered": "Where does the issuer curve sit by maturity?",
                "How to read": "Issuer average yield by maturity is compared with sector, uploaded universe, and optional MMD/rating benchmark curves.",
                "Methodology note": "External MMD is treated as AAA; rating assumptions are disclosed in the appendix.",
            },
            {
                "Chart": "CUSIP Trade Path",
                "Question answered": "What happened to the selected bond through time?",
                "How to read": "Spread, yield/price, and par panels show path, magnitude, and support for a saved candidate.",
                "Methodology note": "CUSIP path uses only uploaded trades for that CUSIP.",
            },
        ]
    )


def focused_report_markdown(context: dict, include_watchlist: bool = True, include_methodology: bool = True) -> str:
    lines = [
        f"# {context['title']}",
        "",
        f"Generated: {context['generated']}",
        f"Prepared for: {context.get('prepared_for') or 'Internal desk review'}",
        "",
    ]
    if context.get("analyst_note"):
        lines.extend(["## Analyst Note", "", str(context["analyst_note"]), ""])

    lines.extend(
        [
            "## Desk Snapshot",
            "",
            focused_report_markdown_table(context["metrics"], max_rows=25),
            "",
            "## Analyst Takeaway",
            "",
        ]
    )
    lines.extend([f"- {bullet}" for bullet in context.get("takeaway_bullets", [])])
    lines.extend(["", "## Top Opportunities", "", focused_report_markdown_table(context["top_opportunities"], max_rows=10), ""])

    if include_watchlist:
        lines.extend(["## Saved Watchlist", "", focused_report_markdown_table(context["saved_watchlist"], max_rows=30), ""])

    lines.extend(["## Core Chart Guide", "", focused_report_markdown_table(context["chart_explanations"], max_rows=10), ""])
    lines.extend(["## Methodology Warnings", "", focused_report_markdown_table(context["warning_rows"], max_rows=10), ""])

    if include_methodology:
        lines.extend(["## Methodology Trust Layer", ""])
        for layer_name, layer_df in _context_trust_layers(context).items():
            lines.extend([f"### {layer_name}", "", focused_report_markdown_table(layer_df, max_rows=30), ""])
        lines.extend(["## Methodology Appendix", "", focused_report_markdown_table(context["methodology"], max_rows=30), ""])

    lines.extend(
        [
            "## Important Limitation",
            "",
            "RV and liquidity scores are screening tools generated from the uploaded data. They are not investment recommendations.",
        ]
    )
    return "\n".join(lines)


def focused_report_html(context: dict, include_watchlist: bool = True, include_methodology: bool = True) -> str:
    note_html = ""
    if context.get("analyst_note"):
        note_html = f"<section><h2>Analyst Note</h2><p>{_html_escape(context['analyst_note'])}</p></section>"
    takeaway = "".join(f"<li>{_html_escape(bullet)}</li>" for bullet in context.get("takeaway_bullets", []))
    watchlist_section = ""
    if include_watchlist:
        watchlist_section = f"""
<section>
  <h2>Saved Watchlist</h2>
  {focused_report_html_table(context["saved_watchlist"], max_rows=30)}
</section>
"""
    methodology_section = ""
    if include_methodology:
        trust_html = "".join(
            f"""
<section>
  <h2>Methodology Trust Layer: {_html_escape(layer_name)}</h2>
  {focused_report_html_table(layer_df, max_rows=30)}
</section>
"""
            for layer_name, layer_df in _context_trust_layers(context).items()
        )
        methodology_section = f"""
{trust_html}
<section>
  <h2>Methodology Appendix</h2>
  {focused_report_html_table(context["methodology"], max_rows=30)}
</section>
"""
    return f"""
<!doctype html>
<html>
<head>
<meta charset="utf-8">
<title>{_html_escape(context['title'])}</title>
<style>
  body {{ font-family: Arial, sans-serif; color: #111827; margin: 34px; line-height: 1.45; }}
  h1 {{ font-size: 30px; margin-bottom: 4px; }}
  h2 {{ font-size: 20px; margin-top: 28px; border-bottom: 1px solid #dbe3ef; padding-bottom: 8px; }}
  .muted {{ color: #64748b; font-size: 13px; }}
  .report-table {{ border-collapse: collapse; width: 100%; margin: 12px 0 8px; font-size: 13px; }}
  .report-table th, .report-table td {{ border: 1px solid #dbe3ef; padding: 8px 10px; text-align: left; vertical-align: top; }}
  .report-table th {{ background: #f5f7fb; color: #334155; }}
  .callout {{ border: 1px solid #b7d4ce; background: #eef8f5; padding: 12px 14px; border-radius: 8px; }}
  @media print {{ body {{ margin: 18mm; }} .report-table {{ font-size: 11px; }} }}
</style>
</head>
<body>
<h1>{_html_escape(context['title'])}</h1>
<p class="muted">Generated {context['generated']} | Prepared for {_html_escape(context.get('prepared_for') or 'Internal desk review')}</p>
{note_html}
<section>
  <h2>Desk Snapshot</h2>
  {focused_report_html_table(context["metrics"], max_rows=25)}
</section>
<section>
  <h2>Analyst Takeaway</h2>
  <div class="callout"><ul>{takeaway}</ul></div>
</section>
<section>
  <h2>Top Opportunities</h2>
  {focused_report_html_table(context["top_opportunities"], max_rows=10)}
</section>
{watchlist_section}
<section>
  <h2>Core Chart Guide</h2>
  {focused_report_html_table(context["chart_explanations"], max_rows=10)}
</section>
<section>
  <h2>Methodology Warnings</h2>
  {focused_report_html_table(context["warning_rows"], max_rows=10)}
</section>
{methodology_section}
<section>
  <h2>Important Limitation</h2>
  <p>RV and liquidity scores are screening tools generated from the uploaded data. They are not investment recommendations.</p>
</section>
</body>
</html>
"""


def focused_report_bundle_bytes(context: dict, markdown: str, html_report: str) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("report.md", markdown)
        zf.writestr("report.html", html_report)
        zf.writestr("reviewer_handoff.md", reviewer_handoff_markdown(str(context.get("issuer", "Selected issuer"))))
        zf.writestr("desk_snapshot.csv", context["metrics"].to_csv(index=False))
        zf.writestr("top_opportunities.csv", context["top_opportunities"].to_csv(index=False))
        zf.writestr("saved_watchlist.csv", context["saved_watchlist"].to_csv(index=False))
        zf.writestr("chart_guide.csv", context["chart_explanations"].to_csv(index=False))
        zf.writestr("methodology_appendix.csv", context["methodology"].to_csv(index=False))
        zf.writestr("methodology_warnings.csv", context["warning_rows"].to_csv(index=False))
        for layer_name, layer_df in _context_trust_layers(context).items():
            safe_name = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(layer_name)).strip("_").lower()
            zf.writestr(f"methodology_trust_{safe_name}.csv", layer_df.to_csv(index=False))
    return buffer.getvalue()


def focused_report_pdf_bytes(context: dict) -> tuple[bytes | None, str | None]:
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
    except Exception as exc:
        return None, str(exc)

    def small_table(df: pd.DataFrame, max_rows: int = 10) -> Table:
        if df is None or df.empty:
            table = Table([["No rows available."]])
            table.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.25, colors.lightgrey), ("FONTSIZE", (0, 0), (-1, -1), 8)]))
            return table
        display = prepare_report_dataframe(df, max_rows=max_rows)
        data = [list(display.columns)] + display.astype(str).values.tolist()
        table = Table(data, repeatRows=1)
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.lightgrey),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("FONTSIZE", (0, 0), (-1, -1), 7),
                ]
            )
        )
        return table

    try:
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=28, leftMargin=28, topMargin=28, bottomMargin=28)
        styles = getSampleStyleSheet()
        story = [
            Paragraph(_html_escape(context["title"]), styles["Title"]),
            Paragraph(f"Generated {context['generated']}", styles["Normal"]),
            Spacer(1, 10),
            Paragraph("Desk Snapshot", styles["Heading2"]),
            small_table(context["metrics"], max_rows=25),
            Spacer(1, 12),
            Paragraph("Analyst Takeaway", styles["Heading2"]),
        ]
        for bullet in context.get("takeaway_bullets", []):
            story.append(Paragraph(f"- {_html_escape(bullet)}", styles["BodyText"]))
        story.extend(
            [
                Spacer(1, 12),
                Paragraph("Top Opportunities", styles["Heading2"]),
                small_table(context["top_opportunities"], max_rows=6),
                Spacer(1, 12),
                Paragraph("Saved Watchlist", styles["Heading2"]),
                small_table(context["saved_watchlist"], max_rows=8) if not context["saved_watchlist"].empty else Paragraph("No saved watchlist rows.", styles["BodyText"]),
                Spacer(1, 12),
                Paragraph("Methodology Trust Layer", styles["Heading2"]),
                small_table(_context_trust_layers(context)["Benchmark Policy"], max_rows=6),
                Spacer(1, 12),
                Paragraph("Methodology Appendix", styles["Heading2"]),
                small_table(context["methodology"], max_rows=10),
                Spacer(1, 10),
                Paragraph("RV and liquidity scores are screening tools, not investment recommendations.", styles["Italic"]),
            ]
        )
        doc.build(story)
        return buffer.getvalue(), None
    except Exception as exc:
        return None, str(exc)


def focused_report_pptx_bytes(context: dict) -> tuple[bytes | None, str | None]:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception as exc:
        return None, str(exc)

    def add_bullets(slide, items: list[str], font_size: int = 16):
        body = slide.placeholders[1].text_frame
        body.clear()
        for item in items:
            p = body.add_paragraph()
            p.text = str(item)
            p.font.size = Pt(font_size)

    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = context["title"]
        slide.placeholders[1].text = f"{context['issuer']} | {context['sector']} | {context['generated']}"

        metrics = {row["Metric"]: row["Value"] for _, row in context["metrics"].iterrows()}
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Desk Snapshot"
        add_bullets(
            slide,
            [
                f"Date range: {metrics.get('Trade date range', 'N/A')}",
                f"Trade rows: {metrics.get('Trade rows', 'N/A')} | CUSIPs: {metrics.get('CUSIPs', 'N/A')}",
                f"Median spread: {metrics.get('Median spread', 'N/A')}",
                f"Top candidate: {metrics.get('Top candidate', 'N/A')}",
                f"Benchmark source: {metrics.get('Benchmark source', 'N/A')}",
            ],
        )

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Analyst Takeaway"
        add_bullets(slide, context.get("takeaway_bullets", [])[:5])

        top_rows = []
        for _, row in context["top_opportunities"].head(5).iterrows():
            top_rows.append(
                f"{row.get('cusip', 'N/A')}: {row.get('signal', 'Monitor')}, spread {_fmt_bps(row.get('current_spread_bps'))}, liquidity {_fmt_num(row.get('liquidity_score'))}"
            )
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Top Opportunities"
        add_bullets(slide, top_rows or ["No top opportunities available."])

        watch_rows = []
        for _, row in context["saved_watchlist"].head(5).iterrows():
            note = str(row.get("note", "") or "").strip()
            status = row.get("status", "Review")
            next_step = str(row.get("next_step", "") or "").strip()
            parts = [str(row.get("signal", "Monitor")), f"status {status}"]
            if next_step:
                parts.append(f"next: {next_step}")
            if note:
                parts.append(note)
            watch_rows.append(f"{row.get('cusip', 'N/A')}: " + " - ".join(parts))
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Saved Watchlist"
        add_bullets(slide, watch_rows or ["No saved watchlist rows."])

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Methodology Notes"
        add_bullets(
            slide,
            [
                "Uploaded MMD is AAA benchmark when external MMD fallback is active.",
                "Trade-sheet Index / Index Rate is not blended with external MMD.",
                "Callable, liquidity, sector, and rating effects remain separate from benchmark spread.",
                "RV and liquidity scores are screening tools, not investment recommendations.",
            ],
            font_size=15,
        )

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue(), None
    except Exception as exc:
        return None, str(exc)
