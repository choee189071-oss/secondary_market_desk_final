from __future__ import annotations

import html
import io
import json
import re

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import streamlit as st

from app_state import (
    ENABLE_REPORT_EXPORTS,
    FULL_DASHBOARD_LABEL,
    LARGE_TABLE_COL_THRESHOLD,
    LARGE_TABLE_ROW_THRESHOLD,
    MAX_HEATMAP_ROWS,
    MAX_TABLE_ROWS,
    PERFORMANCE_MODE,
    SHOW_FULL_RAW_TABLES,
    TABLE_PREVIEW_ROWS,
    WORKFLOW_LABELS,
    WORKFLOW_STEPS,
)
from engine.normalize import (
    coerce_maturity_label,
    ensure_model_columns,
    resolve_model_col,
    safe_melt_by_maturity,
)
from engine.benchmark import (
    BENCHMARK_RATINGS,
    MATURITY_BUCKET_OPTIONS,
    MATURITY_BUCKET_ORDER,
    MATURITY_BUCKET_RENAME,
    MAX_MATURITY_YEAR,
    MMD_BUCKET_MAP,
    RATING_SPREADS,
    build_issuer_curve_snapshot,
    build_spread_level_data,
    build_spread_movement_ladder_data,
    build_spread_observations,
    detect_mmd_date_column as _detect_mmd_date_column,
    get_benchmark_curve,
    make_benchmark_long,
    maturity_year_sort_key,
    observed_maturity_years,
    rating_spread_table,
)
from engine.load_data import process_uploads
from engine.scoring import (
    add_workflow_spread_bps as _add_workflow_spread_bps,
    build_workflow_cusip_summary as _build_workflow_cusip_summary,
    focused_summary_with_peer_gaps as _focused_summary_with_peer_gaps,
    focused_trade_side as _focused_trade_side,
    workflow_date_range_text as _workflow_date_range_text,
)
from engine.validation import (
    BOND_OPTIONAL,
    BOND_RECOMMENDED,
    BOND_REQUIRED,
    CURVE_TEMPLATE_COLUMNS,
    MMD_RECOMMENDED,
    MMD_REQUIRED,
    TRADE_OPTIONAL,
    TRADE_RECOMMENDED,
    TRADE_REQUIRED,
    validate_basic_values,
    validate_dataset,
)
from reports.export_center import (
    focused_core_chart_explanations as _focused_core_chart_explanations,
    focused_methodology_appendix as _focused_methodology_appendix,
)
from ui.export_center import render_focused_export_methodology


# =========================
# Optional OpenAI Support
# =========================

OPENAI_AVAILABLE = False
client = None

try:
    from openai import OpenAI

    if "OPENAI_API_KEY" in st.secrets:
        client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])
        OPENAI_AVAILABLE = True

except Exception:
    OPENAI_AVAILABLE = False
    client = None


def retrieve_market_context_with_openai(
    context_package: dict,
    market_context_query: str,
    model: str = "gpt-4.1-mini",
) -> str:
    """Controlled market / sector context retrieval using OpenAI web search.

    This is separate from final commentary generation so the app can show users
    exactly what market context is being used before synthesis.
    """

    if not OPENAI_AVAILABLE or client is None:
        return (
            "Market context retrieval unavailable. Confirm that `openai` is in requirements.txt "
            "and `OPENAI_API_KEY` is configured in Streamlit Secrets."
        )

    retrieval_prompt = {
        "task": "Retrieve and summarize public market context for muni commentary.",
        "strict_rules": [
            "Focus on public market context only.",
            "Do not invent issuer-specific explanations.",
            "Separate rates/Treasury context, muni market context, sector context, and issuer-specific public headlines if any.",
            "If relevant public context is not found, say so clearly.",
            "Keep the output concise and evidence-oriented.",
        ],
        "dashboard_context": {
            "issuer": context_package.get("issuer"),
            "sector": context_package.get("sector"),
            "bucket": context_package.get("bucket"),
            "benchmark": context_package.get("benchmark"),
            "period": context_package.get("period"),
            "signals": context_package.get("signals", {}),
        },
        "search_focus": market_context_query,
        "preferred_context_categories": [
            "Treasury curve / rates movement",
            "municipal bond market tone",
            "municipal fund flows",
            "sector-specific public news",
            "issuer-specific public news only if clearly available",
        ],
        "output_format": {
            "Rates / Treasury Context": "2-4 bullets",
            "Muni Market Context": "2-4 bullets",
            "Sector / Issuer Context": "2-4 bullets",
            "Relevance to Dashboard Signals": "2-4 bullets",
            "Caveats": "1-2 bullets",
        },
    }

    try:
        response = client.responses.create(
            model=model,
            tools=[{"type": "web_search_preview"}],
            input=[
                {
                    "role": "system",
                    "content": (
                        "You are a market context retrieval assistant for a municipal bond analytics dashboard. "
                        "Retrieve public context conservatively. Do not overstate causality."
                    ),
                },
                {
                    "role": "user",
                    "content": json.dumps(retrieval_prompt, indent=2, default=str),
                },
            ],
            temperature=0.15,
            max_output_tokens=900,
        )

        return response.output_text

    except Exception as e:
        return f"Market Context Retrieval Error: {str(e)}"


def generate_ai_market_commentary(
    context_package: dict,
    manual_market_context: str = "",
    retrieved_market_context: str = "",
    use_web_search: bool = False,
    market_context_query: str = "",
    model: str = "gpt-4.1-mini",
) -> str:
    """Generate evidence-linked institutional commentary.

    The model should only synthesize the analytics and market context provided.
    Web search is available as an optional fallback, but preferred workflow is:
    dashboard signals -> controlled retrieval -> review context -> commentary.
    """

    if not OPENAI_AVAILABLE or client is None:
        return (
            "AI commentary unavailable. Confirm that `openai` is in requirements.txt "
            "and `OPENAI_API_KEY` is configured in Streamlit Secrets."
        )

    system_prompt = """
You are an institutional municipal bond market strategist.

Write concise, evidence-linked secondary-market commentary for a muni trading / public finance team.

Rules:
- Use ONLY the provided dashboard analytics and provided/retrieved market context.
- Do NOT invent issuer-specific news, ratings actions, trades, or market events.
- Do NOT claim causality unless the provided market context supports it.
- If market context is missing or weak, say that context is limited.
- Separate data-backed observations from interpretation.
- Keep tone professional, like buy-side or broker-dealer strategy commentary.
- Mention that signals are screening indicators, not investment recommendations.
- Prefer 4 sections:
  1) Market Commentary
  2) Why This May Be Happening
  3) Risks / Caveats
  4) Evidence Used
"""

    user_payload = {
        "task": "Generate institutional municipal secondary-market commentary.",
        "manual_market_context": manual_market_context,
        "retrieved_market_context": retrieved_market_context,
        "market_context_query": market_context_query,
        "dashboard_context_package": context_package,
        "requested_output_format": {
            "Market Commentary": "2-4 concise bullet points",
            "Why This May Be Happening": "2-4 concise bullet points tied to evidence",
            "Risks / Caveats": "1-3 bullets noting weak evidence, liquidity/data limitations, or alternative explanations",
            "Evidence Used": "bullet list of exact dashboard signals and market-context items used",
        },
    }

    try:
        tools = [{"type": "web_search_preview"}] if use_web_search else []

        response = client.responses.create(
            model=model,
            tools=tools,
            input=[
                {"role": "system", "content": system_prompt},
                {
                    "role": "user",
                    "content": json.dumps(user_payload, indent=2, default=str),
                },
            ],
            temperature=0.25,
            max_output_tokens=1100,
        )

        return response.output_text

    except Exception as e:
        return f"AI Commentary Error: {str(e)}"


def generate_ai_section_readthrough(
    section_title: str,
    python_quote: str,
    evidence: list[str] | None = None,
    model: str = "gpt-4.1-mini",
) -> str:
    """Lightweight AI polish layer for one section.

    The numeric facts are produced by Python first. The AI only rewrites them into
    senior-analyst slide language and must not invent new numbers or causality.
    """
    if not OPENAI_AVAILABLE or client is None:
        return (
            "AI section read-through unavailable. Confirm that `openai` is in requirements.txt "
            "and `OPENAI_API_KEY` is configured in Streamlit Secrets."
        )

    payload = {
        "section_title": section_title,
        "python_generated_quote": python_quote,
        "calculation_evidence": evidence or [],
        "output_rules": [
            "Return 2-3 concise bullets only.",
            "Use only the quote and evidence provided.",
            "Do not invent market news, causality, ratings changes, or additional numbers.",
            "Write in senior municipal fixed-income analyst style.",
            "Make the language slide-ready for a public finance / trading audience.",
            "If the evidence is thin, explicitly phrase the conclusion as a screening signal."
        ],
    }

    try:
        response = client.responses.create(
            model=model,
            input=[
                {
                    "role": "system",
                    "content": (
                        "You are a senior municipal bond analyst. Convert Python-generated dashboard facts "
                        "into concise, evidence-linked slide bullets. Do not add unsupported facts."
                    ),
                },
                {"role": "user", "content": json.dumps(payload, indent=2, default=str)},
            ],
            temperature=0.2,
            max_output_tokens=350,
        )
        return response.output_text
    except Exception as e:
        return f"AI Section Read-through Error: {str(e)}"


def answer_dashboard_question_with_ai(
    context_package: dict,
    user_question: str,
    model: str = "gpt-4.1-mini",
) -> str:
    """Single centralized analyst copilot for dashboard questions.

    This is intentionally centralized rather than repeated under every chart.
    """
    if not OPENAI_AVAILABLE or client is None:
        return (
            "AI dashboard assistant unavailable. Confirm that `openai` is in requirements.txt "
            "and `OPENAI_API_KEY` is configured in Streamlit Secrets."
        )

    payload = {
        "dashboard_context_package": context_package,
        "user_question": user_question,
        "rules": [
            "Answer using only the structured dashboard context package.",
            "Do not invent issuer news, market events, ratings actions, or recommendations.",
            "Separate data-backed observations from interpretation.",
            "If the context package is insufficient, say what additional data is needed.",
            "Keep the answer concise and desk-oriented."
        ],
    }
    try:
        response = client.responses.create(
            model=model,
            input=[
                {
                    "role": "system",
                    "content": (
                        "You are a municipal secondary-market analyst copilot embedded in a dashboard. "
                        "You help users interpret dashboard signals without inventing unsupported facts."
                    ),
                },
                {"role": "user", "content": json.dumps(payload, indent=2, default=str)},
            ],
            temperature=0.2,
            max_output_tokens=700,
        )
        return response.output_text
    except Exception as e:
        return f"AI Dashboard Assistant Error: {str(e)}"




from data_utils import read_uploaded_file


st.set_page_config(page_title="Municipal Secondary Market Dashboard Generator", layout="wide")
st.title("Municipal Secondary Market Dashboard Generator")
st.caption("Bring your own MuniPro trade-history exports. Generate issuer-level relative value and liquidity analytics; bond reference data is optional enrichment.")

st.markdown(
    """
<style>
/* Overall page polish */
.block-container {
    padding-top: 2.2rem;
    padding-bottom: 3rem;
    max-width: 1500px;
}

h1, h2, h3 {
    letter-spacing: -0.02em;
}

section[data-testid="stSidebar"] {
    min-width: 330px !important;
}

div[data-testid="stMetric"] {
    background: #ffffff;
    border: 1px solid #e6e8ef;
    border-radius: 16px;
    padding: 18px 18px 14px 18px;
    box-shadow: 0 1px 3px rgba(15, 23, 42, 0.05);
}

.clean-card {
    background: #ffffff;
    border: 1px solid #e6e8ef;
    border-radius: 18px;
    padding: 18px 20px;
    min-height: 124px;
    box-shadow: 0 1px 3px rgba(15, 23, 42, 0.05);
}

.clean-card-label {
    font-size: 0.86rem;
    font-weight: 700;
    color: #64748b;
    margin-bottom: 8px;
    text-transform: uppercase;
    letter-spacing: 0.03em;
}

.clean-card-value-large {
    font-size: 1.65rem;
    font-weight: 720;
    line-height: 1.15;
    color: #111827;
    overflow-wrap: anywhere;
}

.clean-card-value-small {
    font-size: 1.32rem;
    font-weight: 720;
    line-height: 1.2;
    color: #111827;
    overflow-wrap: anywhere;
}

.clean-card-note {
    font-size: 0.82rem;
    color: #94a3b8;
    margin-top: 8px;
}

.nav-card {
    background: #f8fafc;
    border: 1px solid #e2e8f0;
    border-radius: 16px;
    padding: 14px 16px;
    margin: 10px 0 18px 0;
}

.nav-card a {
    text-decoration: none;
    color: #334155;
    font-size: 0.92rem;
}

.nav-card a:hover {
    color: #0f172a;
    text-decoration: underline;
}

.workflow-grid {
    display: grid;
    grid-template-columns: repeat(6, minmax(0, 1fr));
    gap: 12px;
    margin: 18px 0 22px 0;
}

.workflow-step {
    background: #ffffff;
    border: 1px solid #dbe3ee;
    border-radius: 14px;
    padding: 14px 14px 12px 14px;
    min-height: 112px;
}

.workflow-step-active {
    border-color: #277568;
    background: #eef8f5;
    box-shadow: inset 0 0 0 1px rgba(39, 117, 104, 0.18);
}

.workflow-step-num {
    color: #64748b;
    font-size: 0.78rem;
    font-weight: 750;
    margin-bottom: 8px;
}

.workflow-step-title {
    color: #111827;
    font-size: 0.98rem;
    font-weight: 780;
    line-height: 1.18;
    margin-bottom: 8px;
}

.workflow-step-note {
    color: #64748b;
    font-size: 0.84rem;
    line-height: 1.28;
}

.focus-band {
    background: #f8fafc;
    border: 1px solid #dbe3ee;
    border-radius: 14px;
    padding: 16px 18px;
    margin: 12px 0 18px 0;
}

.methodology-note {
    background: #fffdf4;
    border: 1px solid #f5d36b;
    border-radius: 12px;
    padding: 14px 16px;
    margin: 10px 0;
    color: #4f3b05;
}

.file-card-grid,
.status-card-grid {
    display: grid;
    grid-template-columns: repeat(4, minmax(0, 1fr));
    gap: 12px;
    margin: 12px 0 20px 0;
}

.status-card-grid {
    grid-template-columns: repeat(3, minmax(0, 1fr));
}

.file-card,
.status-card,
.ready-card {
    background: #ffffff;
    border: 1px solid #dbe3ee;
    border-left: 6px solid #94a3b8;
    border-radius: 14px;
    padding: 14px 16px;
    min-height: 116px;
}

.ready-card {
    margin: 12px 0 18px 0;
    min-height: 0;
}

.status-good {
    border-left-color: #15803d;
    background: #f0fdf4;
}

.status-warn {
    border-left-color: #ca8a04;
    background: #fffbeb;
}

.status-bad {
    border-left-color: #b91c1c;
    background: #fef2f2;
}

.status-neutral {
    border-left-color: #64748b;
    background: #f8fafc;
}

.card-kicker {
    color: #64748b;
    font-size: 0.78rem;
    font-weight: 780;
    text-transform: uppercase;
    letter-spacing: 0.03em;
    margin-bottom: 7px;
}

.card-title {
    color: #111827;
    font-size: 1.02rem;
    font-weight: 780;
    line-height: 1.18;
    margin-bottom: 7px;
    overflow-wrap: anywhere;
}

.card-value {
    color: #1f2937;
    font-size: 0.94rem;
    font-weight: 650;
    line-height: 1.35;
    margin-bottom: 7px;
    overflow-wrap: anywhere;
}

.card-detail {
    color: #64748b;
    font-size: 0.85rem;
    line-height: 1.35;
    overflow-wrap: anywhere;
}

.status-pill {
    display: inline-block;
    border-radius: 999px;
    padding: 3px 10px;
    font-size: 0.76rem;
    font-weight: 780;
    margin-bottom: 8px;
}

.status-good .status-pill {
    background: #dcfce7;
    color: #166534;
}

.status-warn .status-pill {
    background: #fef3c7;
    color: #92400e;
}

.status-bad .status-pill {
    background: #fee2e2;
    color: #991b1b;
}

.status-neutral .status-pill {
    background: #e2e8f0;
    color: #334155;
}

@media (max-width: 1200px) {
    .workflow-grid {
        grid-template-columns: repeat(3, minmax(0, 1fr));
    }

    .file-card-grid,
    .status-card-grid {
        grid-template-columns: repeat(2, minmax(0, 1fr));
    }
}

@media (max-width: 760px) {
    .workflow-grid {
        grid-template-columns: repeat(1, minmax(0, 1fr));
    }

    .file-card-grid,
    .status-card-grid {
        grid-template-columns: repeat(1, minmax(0, 1fr));
    }
}

.sidebar-nav-small {
    font-size: 0.88rem;
    line-height: 1.55;
}

/* Keep dataframes/charts visually lighter */
div[data-testid="stDataFrame"] {
    border-radius: 14px;
    overflow: hidden;
}
</style>
""",

    unsafe_allow_html=True,
)


def safe_chart_df(df: pd.DataFrame, required_concepts: list[str], section_name: str = "chart") -> pd.DataFrame:
    """Prepare chart dataframe and warn, not crash, if required concepts are absent."""
    out = ensure_model_columns(df)
    missing = [c for c in required_concepts if resolve_model_col(out, c, required=False) is None]
    if missing:
        st.warning(f"{section_name}: missing required data fields: {', '.join(missing)}. This section was skipped safely.")
        return pd.DataFrame()
    return out


def _make_unique_columns(columns) -> list[str]:
    """Return unique, human-readable column names for Streamlit/Arrow display."""
    seen: dict[str, int] = {}
    unique_cols: list[str] = []
    for col in columns:
        base = str(col)
        if base not in seen:
            seen[base] = 0
            unique_cols.append(base)
        else:
            seen[base] += 1
            unique_cols.append(f"{base}_{seen[base]}")
    return unique_cols


def prepare_display_dataframe(df: pd.DataFrame, max_rows: int | None = TABLE_PREVIEW_ROWS) -> pd.DataFrame:
    """Prepare a dataframe for safe Streamlit display.

    - Removes duplicate column-name ambiguity by suffixing duplicates.
    - Converts dates to MM/DD/YYYY strings for readability.
    - Limits preview rows to keep the web app responsive.
    """
    if df is None:
        return pd.DataFrame()
    if not isinstance(df, pd.DataFrame):
        try:
            df = pd.DataFrame(df)
        except Exception:
            return pd.DataFrame({"value": [str(df)]})

    out = df.copy()
    out.columns = _make_unique_columns(out.columns)

    if max_rows is not None and len(out) > max_rows:
        out = out.head(int(max_rows)).copy()

    for col in out.columns:
        try:
            if pd.api.types.is_datetime64_any_dtype(out[col]):
                out[col] = pd.to_datetime(out[col], errors="coerce").dt.strftime("%m/%d/%Y")
            elif "date" in str(col).lower() or "maturity" in str(col).lower():
                converted = pd.to_datetime(out[col], errors="coerce")
                if len(converted) == 0 or converted.notna().mean() >= 0.6:
                    out[col] = converted.dt.strftime("%m/%d/%Y")
        except Exception:
            # If a column has mixed objects that do not convert cleanly, leave it as-is.
            pass
    return out


def safe_dataframe(
    df: pd.DataFrame,
    *args,
    expander_label: str | None = None,
    expanded: bool = False,
    max_rows: int | None = TABLE_PREVIEW_ROWS,
    auto_collapse: bool = True,
    top_rows: int = 10,
    **kwargs,
):
    """Render dataframes safely for Streamlit Cloud.

    This wrapper solves three common production issues:
    1) duplicate column names that break Arrow serialization;
    2) large audit/detail tables dominating the page;
    3) deprecated Streamlit `use_container_width` warnings flooding logs.

    Large tables show a compact Top-N preview first, with a collapsed preview table
    behind an expander. The full underlying analytics still use the full dataset.
    """
    # Backward compatibility: older calls may still pass use_container_width.
    if "use_container_width" in kwargs and "width" not in kwargs:
        kwargs["width"] = "stretch" if kwargs.pop("use_container_width") else "content"
    else:
        kwargs.pop("use_container_width", None)
    kwargs.setdefault("width", "stretch")

    # Respect the sidebar display row limit when the caller did not specify a tighter one.
    effective_max_rows = max_rows
    try:
        if effective_max_rows == TABLE_PREVIEW_ROWS:
            effective_max_rows = min(TABLE_PREVIEW_ROWS, int(MAX_TABLE_ROWS))
    except Exception:
        pass

    display_df = prepare_display_dataframe(df, max_rows=effective_max_rows)
    row_count = len(df) if isinstance(df, pd.DataFrame) else len(display_df)
    col_count = len(display_df.columns)
    is_large = row_count >= LARGE_TABLE_ROW_THRESHOLD or col_count >= LARGE_TABLE_COL_THRESHOLD

    if expander_label is None:
        expander_label = f"View data table ({row_count:,} rows × {col_count:,} cols)"
        if effective_max_rows is not None and row_count > effective_max_rows:
            expander_label += f" — preview capped at {effective_max_rows:,} rows"

    if auto_collapse and is_large:
        preview_rows = min(max(int(top_rows), 1), len(display_df))
        st.caption(f"Showing top {preview_rows:,} rows. Expand for a larger preview.")
        st.dataframe(display_df.head(preview_rows), *args, **kwargs)
        with st.expander(expander_label, expanded=expanded):
            if effective_max_rows is not None and row_count > effective_max_rows:
                st.caption(f"Large-table protection: showing first {effective_max_rows:,} of {row_count:,} rows.")
            return st.dataframe(display_df, *args, **kwargs)

    return st.dataframe(display_df, *args, **kwargs)


def safe_plotly_chart(fig, *args, **kwargs):
    """Central wrapper for Plotly charts and deprecated Streamlit kwargs."""
    if "use_container_width" in kwargs and "width" not in kwargs:
        kwargs["width"] = "stretch" if kwargs.pop("use_container_width") else "content"
    else:
        kwargs.pop("use_container_width", None)
    kwargs.setdefault("width", "stretch")
    try:
        if fig is not None:
            fig.update_layout(uirevision="keep")
    except Exception:
        pass
    return st.plotly_chart(fig, *args, **kwargs)


# -----------------------------------------------------------------------------
# Analyst read-through engine: numeric, slide-ready commentary from model factors
# -----------------------------------------------------------------------------
def _fmt_bps(x, digits: int = 1) -> str:
    try:
        if pd.isna(x):
            return "N/A"
        return f"{float(x):+.{digits}f} bp"
    except Exception:
        return "N/A"


def _fmt_num(x, digits: int = 1) -> str:
    try:
        if pd.isna(x):
            return "N/A"
        return f"{float(x):,.{digits}f}"
    except Exception:
        return "N/A"


def _fmt_pct(x, digits: int = 1) -> str:
    try:
        if pd.isna(x):
            return "N/A"
        return f"{float(x):.{digits}f}%"
    except Exception:
        return "N/A"


def _fmt_mm(x, digits: int = 1) -> str:
    try:
        if pd.isna(x):
            return "N/A"
        return f"${float(x) / 1_000_000:,.{digits}f}M"
    except Exception:
        return "N/A"


def _fmt_month(x) -> str:
    try:
        return pd.to_datetime(x).strftime("%b %Y")
    except Exception:
        return str(x)


def _fmt_date(x) -> str:
    try:
        return pd.to_datetime(x).strftime("%m/%d/%Y")
    except Exception:
        return str(x)


def _first_existing_col(df: pd.DataFrame, candidates: list[str]) -> str | None:
    if not isinstance(df, pd.DataFrame):
        return None
    for c in candidates:
        if c in df.columns:
            return c
    return None


def _render_slide_quote(title: str, quote: str, evidence: list[str] | None = None, expanded: bool = False):
    """Render a compact slide-ready quote plus optional calculation evidence and AI polish."""
    if not quote:
        return
    st.markdown(f"**Analyst read-through — {title}**")
    st.info(f"Slide-ready quote: {quote}")
    if evidence:
        with st.expander("Evidence / calculation details", expanded=expanded):
            for item in evidence:
                st.markdown(f"- {item}")

    # Optional per-section AI read-through.
    # Python remains the source of truth for numbers; AI only polishes the narrative.
    safe_key = re.sub(r"[^a-zA-Z0-9_]+", "_", str(title)).strip("_").lower()[:60]
    with st.expander("AI read-through polish", expanded=False):
        st.caption(
            "Uses the Python-generated quote and evidence above. The AI is instructed not to invent numbers or causes."
        )
        ai_model_for_section = st.selectbox(
            "Section AI model",
            ["gpt-4.1-mini", "gpt-4.1", "gpt-4o-mini"],
            index=0,
            key=f"section_ai_model_{safe_key}",
        )
        if st.button("Generate polished section bullets", key=f"section_ai_button_{safe_key}"):
            with st.spinner("Generating section read-through..."):
                st.session_state[f"section_ai_output_{safe_key}"] = generate_ai_section_readthrough(
                    section_title=title,
                    python_quote=quote,
                    evidence=evidence or [],
                    model=ai_model_for_section,
                )
        if f"section_ai_output_{safe_key}" in st.session_state:
            st.markdown(st.session_state[f"section_ai_output_{safe_key}"])


def render_spread_trend_readthrough(df: pd.DataFrame, primary_issuer: str, compare_issuers: list[str] | None = None):
    """Spread trend commentary built from daily median spread factors."""
    if df is None or df.empty:
        return
    date_col = _first_existing_col(df, ["trade_date", "date"])
    issuer_col = _first_existing_col(df, ["issuer", "line_item", "volume_group"])
    spread_col = _first_existing_col(df, ["spread_bps", "spread_to_benchmark_bps", "spread"])
    if not all([date_col, issuer_col, spread_col]):
        return
    d = df.copy()
    d[date_col] = pd.to_datetime(d[date_col], errors="coerce")
    d[spread_col] = pd.to_numeric(d[spread_col], errors="coerce")
    if spread_col == "spread" and d[spread_col].abs().median(skipna=True) < 5:
        d[spread_col] = d[spread_col] * 100
    d = d.dropna(subset=[date_col, issuer_col, spread_col])
    if d.empty:
        return
    daily = d.groupby([pd.Grouper(key=date_col, freq="D"), issuer_col], as_index=False).agg(
        spread_bps=(spread_col, "median"), trade_count=(spread_col, "count")
    ).sort_values(date_col)
    primary = daily[daily[issuer_col].astype(str) == str(primary_issuer)]
    if primary.empty:
        return
    start, end = primary.iloc[0], primary.iloc[-1]
    chg = end["spread_bps"] - start["spread_bps"]
    max_row = primary.loc[primary["spread_bps"].idxmax()]
    min_row = primary.loc[primary["spread_bps"].idxmin()]
    direction = "widened" if chg > 0 else "tightened" if chg < 0 else "was largely unchanged"
    peer_text = ""
    peers = [x for x in (compare_issuers or []) if str(x) != str(primary_issuer)]
    evidence = [
        f"{primary_issuer}: {_fmt_bps(start['spread_bps'])} on {_fmt_date(start[date_col])} to {_fmt_bps(end['spread_bps'])} on {_fmt_date(end[date_col])}.",
        f"Peak / trough in selected window: {_fmt_bps(max_row['spread_bps'])} on {_fmt_date(max_row[date_col])}; {_fmt_bps(min_row['spread_bps'])} on {_fmt_date(min_row[date_col])}.",
    ]
    if peers:
        peer_df = daily[daily[issuer_col].astype(str).isin([str(p) for p in peers])]
        if not peer_df.empty:
            peer_moves = []
            for p, g in peer_df.groupby(issuer_col):
                g = g.sort_values(date_col)
                if len(g) >= 2:
                    peer_moves.append(g.iloc[-1]["spread_bps"] - g.iloc[0]["spread_bps"])
            if peer_moves:
                peer_med = float(np.nanmedian(peer_moves))
                peer_text = f" versus a peer median move of {_fmt_bps(peer_med)}"
                evidence.append(f"Peer median move across selected comparison issuers: {_fmt_bps(peer_med)}.")
    quote = (
        f"{primary_issuer} spreads {direction} by {_fmt_bps(chg).replace('+','')} over the selected window, "
        f"moving from {_fmt_bps(start['spread_bps'])} to {_fmt_bps(end['spread_bps'])}{peer_text}."
    )
    _render_slide_quote("spread trend", quote, evidence)


def render_volume_readthrough(df: pd.DataFrame, primary_issuer: str):
    """Trading volume commentary from monthly par/volume and primary issuer share."""
    if df is None or df.empty:
        return
    date_col = _first_existing_col(df, ["trade_date", "date", "month"])
    issuer_col = _first_existing_col(df, ["issuer", "volume_group"])
    amt_col = _first_existing_col(df, ["trade_amount", "monthly_volume", "total_trade_amount"])
    if not all([date_col, issuer_col, amt_col]):
        return
    d = df.copy()
    d[date_col] = pd.to_datetime(d[date_col], errors="coerce")
    d[amt_col] = pd.to_numeric(d[amt_col], errors="coerce")
    d = d.dropna(subset=[date_col, issuer_col, amt_col])
    if d.empty:
        return
    d["_month"] = d[date_col].dt.to_period("M").dt.to_timestamp()
    monthly = d.groupby(["_month", issuer_col], as_index=False).agg(volume=(amt_col, "sum"), trades=(amt_col, "count"))
    totals = monthly.groupby("_month", as_index=False).agg(total_volume=("volume", "sum"))
    primary = monthly[monthly[issuer_col].astype(str) == str(primary_issuer)].groupby("_month", as_index=False).agg(primary_volume=("volume", "sum"))
    share = totals.merge(primary, on="_month", how="left").fillna({"primary_volume": 0})
    share["primary_share"] = np.where(share["total_volume"] > 0, share["primary_volume"] / share["total_volume"] * 100, np.nan)
    if share.empty:
        return
    peak_share = share.loc[share["primary_share"].idxmax()]
    avg_share = share["primary_share"].mean()
    latest = share.sort_values("_month").iloc[-1]
    peak_vol = share.loc[share["primary_volume"].idxmax()]
    quote = (
        f"{primary_issuer} represented an average {_fmt_pct(avg_share)} of selected secondary-market volume, "
        f"peaking at {_fmt_pct(peak_share['primary_share'])} in {_fmt_month(peak_share['_month'])}; "
        f"latest-month volume was {_fmt_mm(latest['primary_volume'])}."
    )
    evidence = [
        f"Peak issuer share: {_fmt_pct(peak_share['primary_share'])} in {_fmt_month(peak_share['_month'])}.",
        f"Peak issuer volume: {_fmt_mm(peak_vol['primary_volume'])} in {_fmt_month(peak_vol['_month'])}.",
        f"Latest month: {_fmt_mm(latest['primary_volume'])} and {_fmt_pct(latest['primary_share'])} of selected volume.",
    ]
    _render_slide_quote("trading volume", quote, evidence)


def render_monthly_activity_readthrough(monthly: pd.DataFrame):
    if monthly is None or monthly.empty or "trade_count" not in monthly.columns:
        return
    date_col = _first_existing_col(monthly, ["trade_month", "month", "_month"])
    if not date_col:
        return
    d = monthly.copy()
    d[date_col] = pd.to_datetime(d[date_col], errors="coerce")
    d["trade_count"] = pd.to_numeric(d["trade_count"], errors="coerce")
    d = d.dropna(subset=[date_col, "trade_count"]).sort_values(date_col)
    if d.empty:
        return
    latest = d.iloc[-1]
    avg = d["trade_count"].mean()
    peak = d.loc[d["trade_count"].idxmax()]
    q = f"Trading activity ended at {latest['trade_count']:,.0f} trades in {_fmt_month(latest[date_col])}, versus an average of {avg:,.0f}; the peak month was {_fmt_month(peak[date_col])} with {peak['trade_count']:,.0f} trades."
    _render_slide_quote("market activity", q, [f"Average monthly trade count: {avg:,.0f}.", f"Peak month: {_fmt_month(peak[date_col])}, {peak['trade_count']:,.0f} trades."])


def render_trade_size_readthrough(size_summary: pd.DataFrame):
    if size_summary is None or size_summary.empty:
        return
    if not {"trade_size_bucket", "trade_count_share", "amount_share"}.issubset(size_summary.columns):
        return
    d = size_summary.copy()
    d["trade_count_share"] = pd.to_numeric(d["trade_count_share"], errors="coerce") * 100
    d["amount_share"] = pd.to_numeric(d["amount_share"], errors="coerce") * 100
    count_leader = d.loc[d["trade_count_share"].idxmax()]
    amt_leader = d.loc[d["amount_share"].idxmax()]
    quote = (
        f"Trade count is concentrated in {count_leader['trade_size_bucket']} trades ({_fmt_pct(count_leader['trade_count_share'])} of tickets), "
        f"while par traded is concentrated in {amt_leader['trade_size_bucket']} blocks ({_fmt_pct(amt_leader['amount_share'])} of volume)."
    )
    _render_slide_quote("trade-size mix", quote, [
        f"Largest share by ticket count: {count_leader['trade_size_bucket']} at {_fmt_pct(count_leader['trade_count_share'])}.",
        f"Largest share by par amount: {amt_leader['trade_size_bucket']} at {_fmt_pct(amt_leader['amount_share'])}.",
    ])


def render_liquidity_readthrough(liq: pd.DataFrame):
    if liq is None or liq.empty:
        return
    d = liq.copy()
    tier_col = _first_existing_col(d, ["liquidity_tier"])
    score_col = _first_existing_col(d, ["liquidity_score"])
    stale_col = _first_existing_col(d, ["days_since_last_trade"])
    cusip_col = _first_existing_col(d, ["cusip"])
    if score_col:
        d[score_col] = pd.to_numeric(d[score_col], errors="coerce")
    if stale_col:
        d[stale_col] = pd.to_numeric(d[stale_col], errors="coerce")
    top = d.sort_values(score_col, ascending=False).iloc[0] if score_col and d[score_col].notna().any() else None
    high_share = None
    if tier_col:
        high_share = (d[tier_col].astype(str).str.contains("High", case=False, na=False).mean() * 100)
    stale_share = (d[stale_col].gt(30).mean() * 100) if stale_col else None
    parts = []
    evidence = []
    if top is not None and cusip_col:
        parts.append(f"the most liquid CUSIP is {top[cusip_col]} with a liquidity score of {_fmt_num(top[score_col])}")
        evidence.append(f"Top liquidity score: {top[cusip_col]} at {_fmt_num(top[score_col])}.")
    if high_share is not None:
        parts.append(f"{_fmt_pct(high_share)} of securities screen as high liquidity")
        evidence.append(f"High-liquidity share: {_fmt_pct(high_share)}.")
    if stale_share is not None:
        parts.append(f"{_fmt_pct(stale_share)} have not traded in more than 30 days")
        evidence.append(f"Staleness share >30 days: {_fmt_pct(stale_share)}.")
    if parts:
        quote = "Liquidity screen indicates " + "; ".join(parts) + "."
        _render_slide_quote("liquidity", quote, evidence)


def render_ladder_readthrough(df: pd.DataFrame, value_col: str, label_col: str | None = None, title: str = "ranking"):
    if df is None or df.empty or value_col not in df.columns:
        return
    d = df.copy()
    d[value_col] = pd.to_numeric(d[value_col], errors="coerce")
    d = d.dropna(subset=[value_col])
    if d.empty:
        return
    if label_col is None or label_col not in d.columns:
        label_col = _first_existing_col(d, ["security_label", "security_bucket", "maturity_bucket", "maturity_zone", "cusip", "issuer"])
    if label_col is None:
        return
    wide = d.loc[d[value_col].idxmax()]
    rich = d.loc[d[value_col].idxmin()]
    quote = (
        f"The widest/richest dispersion in the {title} is {wide[label_col]} at {_fmt_bps(wide[value_col])}, "
        f"versus {rich[label_col]} at {_fmt_bps(rich[value_col])}."
    )
    _render_slide_quote(title, quote, [
        f"Widest / cheapest point: {wide[label_col]}, {_fmt_bps(wide[value_col])}.",
        f"Richest / tightest point: {rich[label_col]}, {_fmt_bps(rich[value_col])}.",
    ])


def render_security_detail_readthrough(sec_daily: pd.DataFrame, cusip: str, benchmark_label: str = "benchmark"):
    if sec_daily is None or sec_daily.empty:
        return
    date_col = _first_existing_col(sec_daily, ["trade_date", "date"])
    y_col = _first_existing_col(sec_daily, ["avg_yield", "yield"])
    s_col = _first_existing_col(sec_daily, ["spread_to_benchmark_bps", "current_spread_bps"])
    if not date_col:
        return
    d = sec_daily.copy()
    d[date_col] = pd.to_datetime(d[date_col], errors="coerce")
    d = d.dropna(subset=[date_col]).sort_values(date_col)
    if len(d) < 2:
        return
    evidence = []
    components = []
    if y_col:
        d[y_col] = pd.to_numeric(d[y_col], errors="coerce")
        dy = d[y_col].iloc[-1] - d[y_col].iloc[0]
        components.append(f"yield moved {_fmt_bps(dy * 100).replace(' bp',' bps')}")
        evidence.append(f"Yield: {_fmt_num(d[y_col].iloc[0], 2)}% to {_fmt_num(d[y_col].iloc[-1], 2)}%.")
    if s_col:
        d[s_col] = pd.to_numeric(d[s_col], errors="coerce")
        ds = d[s_col].iloc[-1] - d[s_col].iloc[0]
        components.append(f"spread to {benchmark_label} moved {_fmt_bps(ds).replace(' bp',' bps')}")
        evidence.append(f"Spread: {_fmt_bps(d[s_col].iloc[0])} to {_fmt_bps(d[s_col].iloc[-1])}.")
    if components:
        q = f"CUSIP {cusip} traded from {_fmt_date(d[date_col].iloc[0])} to {_fmt_date(d[date_col].iloc[-1])}; " + " and ".join(components) + "."
        _render_slide_quote("CUSIP detail", q, evidence)


def render_analyst_pack(market_df: pd.DataFrame, selected_issuer: str):
    """A compact end-of-report commentary pack generated from model factors, not chart pixels."""
    st.markdown("### Automated Analyst Commentary Pack")
    st.caption("Python-generated, factor-based bullets designed to be copied into slides. These use uploaded trade / benchmark fields rather than image-level interpretation.")
    try:
        if market_df is not None and not market_df.empty:
            tmp = market_df.copy()
            if "issuer" in tmp.columns:
                tmp = tmp[tmp["issuer"].astype(str) == str(selected_issuer)]
            if not tmp.empty:
                if "spread" in tmp.columns:
                    tmp["_spread_bps"] = pd.to_numeric(tmp["spread"], errors="coerce") * 100
                elif {"yield", "index_rate"}.issubset(tmp.columns):
                    tmp["_spread_bps"] = (pd.to_numeric(tmp["yield"], errors="coerce") - pd.to_numeric(tmp["index_rate"], errors="coerce")) * 100
                if {"trade_date", "_spread_bps"}.issubset(tmp.columns):
                    render_spread_trend_readthrough(tmp.rename(columns={"_spread_bps":"spread_bps"}), selected_issuer, [])
                if {"trade_date", "trade_amount", "issuer"}.issubset(market_df.columns):
                    render_volume_readthrough(market_df, selected_issuer)
    except Exception as e:
        st.caption(f"Commentary pack skipped safely: {e}")


def compact_ladder_table_for_display(table: pd.DataFrame, max_rows: int | None = None) -> pd.DataFrame:
    """Limit ladder rows to the most informative non-empty maturity rows.

    Rows are ranked by maximum absolute movement so the chart stays readable.
    """
    if table is None or table.empty:
        return table
    out = table.dropna(how="all").copy()
    if out.empty:
        return out
    max_rows = max_rows or MAX_HEATMAP_ROWS
    if PERFORMANCE_MODE and len(out) > max_rows:
        score = out.abs().max(axis=1).sort_values(ascending=False)
        keep = score.head(max_rows).index.tolist()
        out = out.loc[keep]
        out = out.loc[sorted(out.index, key=maturity_year_sort_key)]
    return out


def maturity_zone_label(value: object) -> str:
    """Collapse annual maturity years into desk-readable curve sectors."""
    y = maturity_year_sort_key(value)
    if y == 9999:
        return "Unknown"
    if y <= 3:
        return "1-3Y"
    if y <= 7:
        return "4-7Y"
    if y <= 12:
        return "8-12Y"
    if y <= 20:
        return "13-20Y"
    return "21Y+"


MATURITY_ZONE_ORDER = ["1-3Y", "4-7Y", "8-12Y", "13-20Y", "21Y+"]


def maturity_display_order(values: object) -> list[str]:
    """Return a stable maturity display order for annual buckets or maturity zones.

    This prevents charts from going blank when a table has been aggregated from
    1Y/2Y/... into 1-3Y/4-7Y/... zones. Previously some charts forced
    MATURITY_BUCKET_ORDER only, which converted zone labels into NaN categories.
    """
    try:
        vals = pd.Series(list(values)).dropna().astype(str).unique().tolist()
    except Exception:
        vals = []

    zone_vals = [z for z in MATURITY_ZONE_ORDER if z in vals]
    annual_vals = [v for v in vals if re.fullmatch(r"\d{1,2}Y", str(v).strip().upper())]
    annual_vals = sorted(annual_vals, key=maturity_year_sort_key)
    other_vals = sorted([v for v in vals if v not in set(zone_vals + annual_vals)])
    return zone_vals + annual_vals + other_vals


def sanitize_curve_long_for_plot(
    df: pd.DataFrame,
    x_col: str = "maturity_bucket",
    y_col: str = "spread_to_benchmark_bps",
    color_col: str = "benchmark_rating",
) -> pd.DataFrame:
    """Coerce curve data into a plot-safe long dataframe.

    The dashboard accepts uploaded data with evolving schemas. This function
    makes the current spread curve defensive by:
    - resolving maturity columns through the central data model;
    - coercing spread values to numeric;
    - dropping rows with missing x/y values;
    - preserving maturity-zone labels instead of forcing them into annual buckets.
    """
    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        return pd.DataFrame()

    out = ensure_model_columns(df).copy()
    out = out.loc[:, ~out.columns.duplicated()].copy()

    maturity_col = x_col if x_col in out.columns else resolve_model_col(out, "maturity_bucket", required=False)
    spread_col = y_col if y_col in out.columns else resolve_model_col(out, "spread_bps", required=False)

    if maturity_col is None or spread_col is None:
        return pd.DataFrame()

    if maturity_col != x_col:
        out[x_col] = out[maturity_col]
    if spread_col != y_col:
        out[y_col] = out[spread_col]

    if color_col not in out.columns:
        out[color_col] = "Benchmark"

    out[x_col] = out[x_col].apply(coerce_maturity_label).astype("string")
    out[y_col] = pd.to_numeric(out[y_col], errors="coerce")
    out = out.dropna(subset=[x_col, y_col, color_col]).copy()

    order = maturity_display_order(out[x_col].tolist())
    if order:
        out[x_col] = pd.Categorical(out[x_col].astype(str), categories=order, ordered=True)
        out = out.sort_values([color_col, x_col])

    return out


def curve_data_audit(df: pd.DataFrame, required_cols: list[str] | None = None) -> pd.DataFrame:
    """Small diagnostic table for empty charts, shown only inside an expander."""
    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        return pd.DataFrame({"check": ["input_rows"], "value": [0]})
    required_cols = required_cols or ["maturity_bucket", "benchmark_rating", "spread_to_benchmark_bps"]
    rows = [
        {"check": "input_rows", "value": len(df)},
        {"check": "columns", "value": ", ".join(map(str, df.columns[:25])) + ("..." if len(df.columns) > 25 else "")},
    ]
    for col in required_cols:
        if col in df.columns:
            rows.append({"check": f"{col}_missing", "value": int(df[col].isna().sum())})
            rows.append({"check": f"{col}_non_missing", "value": int(df[col].notna().sum())})
        else:
            rows.append({"check": f"{col}_present", "value": False})
    return pd.DataFrame(rows)


def aggregate_maturity_rows_for_display(table: pd.DataFrame, agg: str = "median") -> pd.DataFrame:
    """Aggregate ladder rows from 1Y..40Y into readable maturity zones."""
    if table is None or table.empty:
        return table
    out = table.dropna(how="all").copy()
    if out.empty:
        return out
    out["__maturity_zone__"] = [maturity_zone_label(idx) for idx in out.index]
    grouped = out.groupby("__maturity_zone__").median(numeric_only=True) if agg == "median" else out.groupby("__maturity_zone__").mean(numeric_only=True)
    grouped = grouped.reindex([z for z in MATURITY_ZONE_ORDER if z in grouped.index])
    return grouped


def aggregate_maturity_columns_for_display(table: pd.DataFrame, agg: str = "median") -> pd.DataFrame:
    """Aggregate ladder columns from 1Y..40Y into readable maturity zones."""
    if table is None or table.empty:
        return table
    out = table.dropna(how="all").copy()
    if out.empty:
        return out
    zone_map = {col: maturity_zone_label(col) for col in out.columns}
    frames = []
    for zone in MATURITY_ZONE_ORDER:
        cols = [c for c, z in zone_map.items() if z == zone]
        if not cols:
            continue
        numeric = out[cols].apply(pd.to_numeric, errors="coerce")
        frames.append(numeric.median(axis=1).rename(zone) if agg == "median" else numeric.mean(axis=1).rename(zone))
    return pd.concat(frames, axis=1) if frames else out


def ranked_bar_chart(
    df: pd.DataFrame,
    value_col: str,
    label_col: str,
    title: str,
    x_title: str,
    top_n: int = 15,
    color_col: str | None = None,
    hover_cols: list[str] | None = None,
):
    """Desk-friendly horizontal bar chart for top ranked opportunities."""
    if df is None or df.empty or value_col not in df.columns or label_col not in df.columns:
        return None
    plot_df = df.copy()
    plot_df[value_col] = pd.to_numeric(plot_df[value_col], errors="coerce")
    plot_df = plot_df.dropna(subset=[value_col])
    if plot_df.empty:
        return None
    plot_df = plot_df.sort_values(value_col, ascending=False).head(top_n)
    plot_df[label_col] = plot_df[label_col].astype(str)
    fig = px.bar(
        plot_df.sort_values(value_col, ascending=True),
        x=value_col,
        y=label_col,
        orientation="h",
        color=color_col if color_col in plot_df.columns else None,
        hover_data=[c for c in (hover_cols or []) if c in plot_df.columns],
        title=title,
        labels={value_col: x_title, label_col: "Security / Bucket"},
    )
    fig.add_vline(x=0, line_dash="dash", opacity=0.35)
    fig.update_layout(height=max(360, 28 * len(plot_df) + 140), showlegend=bool(color_col and color_col in plot_df.columns))
    return fig


def add_security_label(df: pd.DataFrame, label_col: str = "security_label") -> pd.DataFrame:
    """Create a compact, human-readable label for CUSIP/security rows."""
    out = df.copy()
    if "cusip" in out.columns:
        out[label_col] = out["cusip"].astype(str)
    elif "issuer" in out.columns and "maturity_bucket" in out.columns:
        out[label_col] = out["issuer"].astype(str) + " " + out["maturity_bucket"].astype(str)
    elif "issuer" in out.columns:
        out[label_col] = out["issuer"].astype(str)
    else:
        out[label_col] = out.index.astype(str)
    if "maturity_bucket" in out.columns and "cusip" in out.columns:
        out[label_col] = out[label_col] + " (" + out["maturity_bucket"].astype(str) + ")"
    return out


def section_anchor(anchor_id: str, title: str, level: int = 2):
    """Create a stable HTML anchor plus a Streamlit header/subheader."""
    st.markdown(f"<a id='{anchor_id}'></a>", unsafe_allow_html=True)
    if level == 1:
        st.title(title)
    elif level == 2:
        st.header(title)
    else:
        st.subheader(title)


def clean_metric_card(label: str, value: object, size: str = "large", note: str | None = None):
    """Compact custom metric card that gives long text more room than st.metric."""
    value_class = "clean-card-value-large" if size == "large" else "clean-card-value-small"
    safe_value = "—" if value is None else str(value)
    note_html = f"<div class='clean-card-note'>{note}</div>" if note else ""
    st.markdown(
        f"""
<div class="clean-card">
  <div class="clean-card-label">{label}</div>
  <div class="{value_class}">{safe_value}</div>
  {note_html}
</div>
""",
        unsafe_allow_html=True,
    )


def _html_escape(value: object) -> str:
    return html.escape("" if value is None else str(value), quote=True)


def _status_label(status: str) -> str:
    return {
        "good": "Green",
        "warn": "Yellow",
        "bad": "Red",
        "neutral": "Info",
    }.get(status, "Info")


def _render_card_grid(cards: list[dict], grid_class: str):
    parts = [f"<div class='{grid_class}'>"]
    for card in cards:
        status = card.get("status", "neutral")
        parts.append(
            f"""
<div class="{card.get('class_name', 'status-card')} status-{_html_escape(status)}">
  <div class="status-pill">{_html_escape(_status_label(status))}</div>
  <div class="card-kicker">{_html_escape(card.get('kicker', ''))}</div>
  <div class="card-title">{_html_escape(card.get('title', ''))}</div>
  <div class="card-value">{_html_escape(card.get('value', ''))}</div>
  <div class="card-detail">{_html_escape(card.get('detail', ''))}</div>
</div>
"""
        )
    parts.append("</div>")
    st.markdown("".join(parts), unsafe_allow_html=True)


def render_upload_file_cards(
    trade_file_names: list[str],
    bond_file_name: str | None,
    issuer_mapping_file_name: str | None,
    mmd_file_name: str | None,
    use_external_mmd_fallback: bool,
):
    """Show selected upload files as role-based cards before the heavier audit."""
    trade_count = len(trade_file_names)
    cards = [
        {
            "class_name": "file-card",
            "status": "good" if trade_count else "bad",
            "kicker": "Required",
            "title": "Trade files",
            "value": f"{trade_count:,} selected" if trade_count else "No trade file selected",
            "detail": ", ".join(trade_file_names[:3]) + (" ..." if trade_count > 3 else "") if trade_count else "Upload at least one MuniPro trade-history export.",
        },
        {
            "class_name": "file-card",
            "status": "good" if bond_file_name else "neutral",
            "kicker": "Optional",
            "title": "Bond reference",
            "value": bond_file_name or "Not uploaded",
            "detail": "Adds call, tax, lien, and static security metadata when available.",
        },
        {
            "class_name": "file-card",
            "status": "good" if issuer_mapping_file_name else "neutral",
            "kicker": "Optional",
            "title": "Issuer mapping",
            "value": issuer_mapping_file_name or "Not uploaded",
            "detail": "Adds persistent issuer/sector labels instead of manual overrides.",
        },
        {
            "class_name": "file-card",
            "status": "good" if mmd_file_name else ("warn" if use_external_mmd_fallback else "neutral"),
            "kicker": "Optional benchmark",
            "title": "MMD / AAA curve",
            "value": mmd_file_name or ("Fallback enabled; no file selected" if use_external_mmd_fallback else "Fallback disabled"),
            "detail": "Uploaded MMD is treated as the AAA curve only when it is the active fallback benchmark.",
        },
    ]
    _render_card_grid(cards, "file-card-grid")


def _nonnull_rate(df: pd.DataFrame, col: str) -> float | None:
    if df.empty or col not in df.columns or len(df) == 0:
        return None
    series = df[col]
    if series.dtype == object:
        valid = series.notna() & (series.astype(str).str.strip() != "")
    else:
        valid = series.notna()
    return float(valid.mean() * 100)


def _numeric_rate(df: pd.DataFrame, col: str) -> float | None:
    if df.empty or col not in df.columns or len(df) == 0:
        return None
    return float(pd.to_numeric(df[col], errors="coerce").notna().mean() * 100)


def _rate_status(rate: float | None, good_threshold: float = 95, warn_threshold: float = 80) -> str:
    if rate is None:
        return "bad"
    if rate >= good_threshold:
        return "good"
    if rate >= warn_threshold:
        return "warn"
    return "bad"


def build_upload_audit_cards(
    trade_reports: list[dict],
    market_df: pd.DataFrame,
    benchmark_source_mode: str,
    use_external_mmd_fallback: bool,
    mmd_file_provided: bool,
) -> tuple[list[dict], dict]:
    """Return red/yellow/green audit cards and a compact readiness summary."""
    required_ok = bool(trade_reports) and all(report.get("can_run") for report in trade_reports)
    required_missing = sorted(
        {
            missing
            for report in trade_reports
            for missing in report.get("missing_required", [])
        }
    )

    trade_date_rate = _nonnull_rate(market_df, "trade_date")
    dates = pd.to_datetime(market_df["trade_date"], errors="coerce").dropna() if "trade_date" in market_df.columns else pd.Series(dtype="datetime64[ns]")
    if dates.empty:
        date_status = "bad"
        date_value = "No valid trade dates"
        date_detail = "Trade date is required for time-series charts and snapshot period filters."
    else:
        unique_dates = dates.dt.normalize().nunique()
        date_status = "good" if trade_date_rate and trade_date_rate >= 95 and unique_dates >= 2 else "warn"
        date_value = f"{dates.min():%m/%d/%Y} - {dates.max():%m/%d/%Y}"
        date_detail = f"{trade_date_rate:.1f}% valid date rows across {unique_dates:,} unique date(s)."

    cusip_rate = _nonnull_rate(market_df, "cusip")
    cusip_status = _rate_status(cusip_rate, good_threshold=95, warn_threshold=80)

    yield_rate = _numeric_rate(market_df, "yield")
    index_rate = _numeric_rate(market_df, "index_rate")
    spread_rate = _numeric_rate(market_df, "spread")
    yield_status = _rate_status(yield_rate, good_threshold=90, warn_threshold=70)

    best_benchmark_input_rate = max([x for x in [index_rate, spread_rate] if x is not None], default=0)
    if benchmark_source_mode in {"Trade Sheet Index / Index Rate", "Uploaded MMD fallback"}:
        benchmark_status = "good"
    elif use_external_mmd_fallback and not mmd_file_provided:
        benchmark_status = "warn"
    else:
        benchmark_status = "bad"

    if benchmark_source_mode == "Trade Sheet Index / Index Rate":
        benchmark_value = "Trade sheet Index / Index Rate"
        benchmark_detail = "Primary benchmark is active. External MMD is not mixed into this run."
    elif benchmark_source_mode == "Uploaded MMD fallback":
        benchmark_value = "Uploaded MMD fallback"
        benchmark_detail = "Uploaded MMD is active as the AAA benchmark because trade index data was unavailable."
    else:
        benchmark_value = "No active benchmark"
        benchmark_detail = "Yield/liquidity analytics can run, but spread-to-benchmark views are degraded."

    if yield_status == "bad":
        spread_input_status = "bad"
    elif best_benchmark_input_rate >= 70:
        spread_input_status = "good"
    elif benchmark_status == "good":
        spread_input_status = "warn"
    else:
        spread_input_status = "bad"

    cards = [
        {
            "status": "good" if required_ok else "bad",
            "kicker": "Required fields",
            "title": "Minimum schema",
            "value": "Pass" if required_ok else "Blocking issue",
            "detail": "All required trade fields were detected." if required_ok else "Missing: " + ", ".join(required_missing),
        },
        {
            "status": date_status,
            "kicker": "Date coverage",
            "title": "Trade date window",
            "value": date_value,
            "detail": date_detail,
        },
        {
            "status": cusip_status,
            "kicker": "CUSIP quality",
            "title": "Valid CUSIP rate",
            "value": "N/A" if cusip_rate is None else f"{cusip_rate:.1f}%",
            "detail": "CUSIP-level drilldown and watchlist depend on this field.",
        },
        {
            "status": benchmark_status,
            "kicker": "Benchmark source",
            "title": "Active curve policy",
            "value": benchmark_value,
            "detail": benchmark_detail,
        },
        {
            "status": yield_status,
            "kicker": "Yield availability",
            "title": "Numeric yield rows",
            "value": "N/A" if yield_rate is None else f"{yield_rate:.1f}%",
            "detail": "Yield is required for spread, curve, and RV calculations.",
        },
        {
            "status": spread_input_status,
            "kicker": "Spread inputs",
            "title": "Index Rate / Spread",
            "value": f"Index {index_rate or 0:.1f}% / Spread {spread_rate or 0:.1f}%",
            "detail": "Trade Index Rate is preferred; uploaded MMD can fill the benchmark role only as fallback.",
        },
    ]

    blocking = [c for c in cards if c["status"] == "bad" and c["kicker"] in {"Required fields", "Date coverage", "Yield availability"}]
    benchmark_missing = benchmark_status == "bad"
    if blocking:
        ready_status = "bad"
        ready_value = "Not ready"
        next_step = "Stay on Upload / Data Audit and fix blocking fields before analysis."
    elif benchmark_missing:
        ready_status = "warn"
        ready_value = "Ready for yield/liquidity only"
        next_step = "Go to Desk Snapshot for basic review, but add Index Rate or enable MMD fallback before relying on spread/RV outputs."
    elif any(c["status"] == "warn" for c in cards):
        ready_status = "warn"
        ready_value = "Ready with warnings"
        next_step = "Go to Desk Snapshot, then review warnings before using CUSIP/RV outputs."
    else:
        ready_status = "good"
        ready_value = "Ready to analyze"
        next_step = "Next step: open Desk Snapshot."

    readiness = {
        "status": ready_status,
        "value": ready_value,
        "next_step": next_step,
        "bad_count": sum(1 for c in cards if c["status"] == "bad"),
        "warn_count": sum(1 for c in cards if c["status"] == "warn"),
    }
    return cards, readiness


def render_ready_to_analyze_card(readiness: dict):
    status = readiness.get("status", "neutral")
    st.markdown(
        f"""
<div class="ready-card status-{_html_escape(status)}">
  <div class="status-pill">{_html_escape(_status_label(status))}</div>
  <div class="card-kicker">Ready to Analyze</div>
  <div class="card-title">{_html_escape(readiness.get('value', 'Review upload'))}</div>
  <div class="card-value">{_html_escape(readiness.get('next_step', 'Review the audit cards above.'))}</div>
  <div class="card-detail">Warnings: {_html_escape(readiness.get('warn_count', 0))} | Blocking issues: {_html_escape(readiness.get('bad_count', 0))}</div>
</div>
""",
        unsafe_allow_html=True,
    )


def render_workflow_header(active_label: str, files_loaded: int = 0, issuers_loaded: int = 0):
    """Render the six-step workstation flow as a compact visual map."""
    html_parts = ["<div class='workflow-grid'>"]
    for idx, step in enumerate(WORKFLOW_STEPS, start=1):
        active_class = " workflow-step-active" if step["label"] == active_label else ""
        if idx == 1 and files_loaded:
            note = f"{files_loaded:,} file(s); {issuers_loaded:,} issuer(s)."
        elif idx == 2 and issuers_loaded:
            note = "Ready after issuer selection."
        else:
            note = step["note"]
        html_parts.append(
            f"""
<div class='workflow-step{active_class}'>
  <div class='workflow-step-num'>{idx:02d}</div>
  <div class='workflow-step-title'>{step['title']}</div>
  <div class='workflow-step-note'>{note}</div>
</div>
"""
        )
    html_parts.append("</div>")
    st.markdown("".join(html_parts), unsafe_allow_html=True)


def _focused_watchlist_records() -> dict:
    """Return mutable watchlist records, migrating older list-based session state."""
    if "focused_watchlist_records" not in st.session_state:
        records = {}
        for cusip in st.session_state.get("focused_watchlist", []):
            records[str(cusip)] = {
                "cusip": str(cusip),
                "issuer": "",
                "signal": "",
                "note": "",
                "source": "Migrated",
                "added_at": pd.Timestamp.now().strftime("%Y-%m-%d %H:%M"),
            }
        st.session_state["focused_watchlist_records"] = records
    return st.session_state["focused_watchlist_records"]


def _upsert_focused_watchlist(cusip: object, issuer: str, source: str, row: pd.Series | dict | None = None, note: str = ""):
    records = _focused_watchlist_records()
    key = str(cusip)
    existing = records.get(key, {})
    row_dict = row.to_dict() if isinstance(row, pd.Series) else (row or {})
    records[key] = {
        "cusip": key,
        "issuer": issuer or existing.get("issuer", ""),
        "signal": row_dict.get("signal", existing.get("signal", "")),
        "maturity_bucket": row_dict.get("maturity_bucket", existing.get("maturity_bucket", "")),
        "current_spread_bps": row_dict.get("current_spread_bps", existing.get("current_spread_bps", pd.NA)),
        "peer_median_gap_bps": row_dict.get("peer_median_gap_bps", existing.get("peer_median_gap_bps", pd.NA)),
        "liquidity_score": row_dict.get("liquidity_score", existing.get("liquidity_score", pd.NA)),
        "rv_score": row_dict.get("rv_score", existing.get("rv_score", pd.NA)),
        "trade_count": row_dict.get("trade_count", existing.get("trade_count", pd.NA)),
        "total_trade_amount": row_dict.get("total_trade_amount", existing.get("total_trade_amount", pd.NA)),
        "latest_trade": row_dict.get("latest_trade", existing.get("latest_trade", pd.NA)),
        "note": note if note else existing.get("note", ""),
        "source": source or existing.get("source", ""),
        "added_at": existing.get("added_at") or pd.Timestamp.now().strftime("%Y-%m-%d %H:%M"),
        "updated_at": pd.Timestamp.now().strftime("%Y-%m-%d %H:%M"),
    }
    st.session_state["focused_watchlist"] = sorted(records.keys())


def _focused_watchlist_dataframe(summary: pd.DataFrame | None = None) -> pd.DataFrame:
    records = _focused_watchlist_records()
    if not records:
        return pd.DataFrame()
    out = pd.DataFrame(records.values())
    if summary is not None and not summary.empty and "cusip" in summary.columns:
        refresh_cols = [
            "cusip", "signal", "maturity_bucket", "current_spread_bps", "peer_median_gap_bps", "liquidity_score",
            "rv_score", "trade_count", "total_trade_amount", "latest_trade",
        ]
        current = summary[[c for c in refresh_cols if c in summary.columns]].copy()
        current["cusip"] = current["cusip"].astype(str)
        out = out.merge(current, on="cusip", how="left", suffixes=("", "_current"))
        for col in [c for c in current.columns if c != "cusip"]:
            current_col = f"{col}_current"
            if current_col not in out.columns:
                continue
            if col in out.columns:
                out[col] = out[current_col].combine_first(out[col])
            else:
                out[col] = out[current_col]
            out = out.drop(columns=[current_col])
    return out


def _focused_watchlist_markdown(saved_df: pd.DataFrame, issuer: str) -> str:
    if saved_df.empty:
        return f"# {issuer} Watchlist\n\nNo saved candidates."
    lines = [
        f"# {issuer} Watchlist",
        "",
        f"Generated: {pd.Timestamp.now():%Y-%m-%d %H:%M}",
        "",
    ]
    for _, row in saved_df.iterrows():
        lines.extend(
            [
                f"## {row.get('cusip', 'N/A')}",
                f"- Signal: {row.get('signal', 'N/A')}",
                f"- Maturity bucket: {row.get('maturity_bucket', 'N/A')}",
                f"- Spread: {_fmt_bps(row.get('current_spread_bps'))}",
                f"- Peer median gap: {_fmt_bps(row.get('peer_median_gap_bps'))}",
                f"- Liquidity score: {_fmt_num(row.get('liquidity_score'))}",
                f"- RV score: {_fmt_num(row.get('rv_score'))}",
                f"- Note: {row.get('note', '') or 'N/A'}",
                "",
            ]
        )
    return "\n".join(lines)


def _build_focused_report_context(
    report_title: str,
    prepared_for: str,
    analyst_note: str,
    selected_issuer: str,
    selected_sector: str,
    market_df: pd.DataFrame,
    issuer_trades: pd.DataFrame,
    issuer_bonds: pd.DataFrame,
    mmd_df: pd.DataFrame,
    benchmark_source_mode: str,
    benchmark_priority: str,
    benchmark_conflict_policy: str,
) -> dict:
    issuer_base = _add_workflow_spread_bps(issuer_trades)
    market_base = _add_workflow_spread_bps(market_df)
    cusip_summary = _focused_summary_with_peer_gaps(_build_workflow_cusip_summary(issuer_base))
    saved_watchlist = _focused_watchlist_dataframe(cusip_summary)
    top_opportunities = cusip_summary.head(5).copy()
    methodology = _focused_methodology_appendix(
        mmd_df=mmd_df,
        benchmark_source_mode=benchmark_source_mode,
        benchmark_priority=benchmark_priority,
        benchmark_conflict_policy=benchmark_conflict_policy,
    )
    chart_explanations = _focused_core_chart_explanations(selected_issuer, selected_sector, benchmark_source_mode)
    warning_cards = _build_snapshot_methodology_cards(issuer_base, mmd_df, benchmark_source_mode)
    takeaway_bullets, takeaway_labels = _build_snapshot_takeaway(
        issuer_df=issuer_base,
        market_df=market_base,
        cusip_summary=cusip_summary,
        selected_issuer=selected_issuer,
        benchmark_source_mode=benchmark_source_mode,
    )

    spread_series = pd.to_numeric(issuer_base.get("spread_bps"), errors="coerce") if "spread_bps" in issuer_base.columns else pd.Series(dtype="float64")
    liq_series = pd.to_numeric(cusip_summary.get("liquidity_score"), errors="coerce") if "liquidity_score" in cusip_summary.columns else pd.Series(dtype="float64")
    latest_trade = "No trades"
    if not issuer_base.empty and "trade_date" in issuer_base.columns:
        latest_date = pd.to_datetime(issuer_base["trade_date"], errors="coerce").dropna().max()
        if pd.notna(latest_date):
            latest_trade = latest_date.strftime("%m/%d/%Y")
    top_candidate = cusip_summary.iloc[0] if not cusip_summary.empty else None
    top_candidate_label = "N/A" if top_candidate is None else str(top_candidate.get("cusip", "N/A"))
    top_candidate_note = ""
    if top_candidate is not None:
        top_candidate_note = (
            f"{top_candidate.get('signal', 'Monitor')} | spread {_fmt_bps(top_candidate.get('current_spread_bps'))} | "
            f"liquidity {_fmt_num(top_candidate.get('liquidity_score'))} | RV {_fmt_num(top_candidate.get('rv_score'))}"
        )

    metrics = pd.DataFrame(
        [
            {"Metric": "Issuer", "Value": selected_issuer},
            {"Metric": "Sector", "Value": selected_sector or "Unknown"},
            {"Metric": "Trade date range", "Value": _workflow_date_range_text(issuer_base)},
            {"Metric": "Latest trade", "Value": latest_trade},
            {"Metric": "Trade rows", "Value": f"{len(issuer_base):,}"},
            {"Metric": "Security rows", "Value": f"{len(issuer_bonds):,}"},
            {"Metric": "CUSIPs", "Value": f"{issuer_base['cusip'].nunique() if 'cusip' in issuer_base.columns else 0:,}"},
            {"Metric": "Median spread", "Value": "N/A" if spread_series.dropna().empty else f"{spread_series.median():.1f} bps"},
            {"Metric": "Median liquidity", "Value": "N/A" if liq_series.dropna().empty else f"{liq_series.median():.1f}"},
            {"Metric": "Top candidate", "Value": top_candidate_label},
            {"Metric": "Saved watchlist", "Value": f"{len(saved_watchlist):,}"},
            {"Metric": "Benchmark source", "Value": benchmark_source_mode},
        ]
    )

    warning_rows = pd.DataFrame(
        [
            {
                "Area": card.get("kicker", ""),
                "Status": _status_label(card.get("status", "neutral")),
                "Value": card.get("value", ""),
                "Detail": card.get("detail", ""),
            }
            for card in warning_cards
        ]
    )

    return {
        "title": report_title or f"{selected_issuer} Secondary Market Desk Report",
        "prepared_for": prepared_for,
        "analyst_note": analyst_note,
        "generated": pd.Timestamp.now().strftime("%Y-%m-%d %H:%M"),
        "issuer": selected_issuer,
        "sector": selected_sector,
        "metrics": metrics,
        "takeaway_bullets": takeaway_bullets,
        "takeaway_labels": takeaway_labels,
        "top_opportunities": top_opportunities,
        "saved_watchlist": saved_watchlist,
        "methodology": methodology,
        "chart_explanations": chart_explanations,
        "warning_rows": warning_rows,
        "cusip_summary": cusip_summary,
        "benchmark_source_mode": benchmark_source_mode,
        "benchmark_priority": benchmark_priority,
        "benchmark_conflict_policy": benchmark_conflict_policy,
        "top_candidate_note": top_candidate_note,
    }


def _render_benchmark_methodology_block(mmd_df: pd.DataFrame, benchmark_source_mode: str, benchmark_priority: str, benchmark_conflict_policy: str):
    st.markdown(
        """
<div class="methodology-note">
<b>Benchmark policy:</b> uploaded MMD is treated as the AAA benchmark curve when it is the active benchmark source.
If trade-sheet Index / Index Rate is available, the app uses that trade-implied benchmark first and does not mix it with external MMD in the same run.
</div>
""",
        unsafe_allow_html=True,
    )
    safe_dataframe(
        pd.DataFrame(
            [
                {"Item": "Active benchmark source", "Value": benchmark_source_mode},
                {"Item": "Priority", "Value": benchmark_priority},
                {"Item": "Conflict policy", "Value": benchmark_conflict_policy},
                {"Item": "Benchmark row count", "Value": f"{len(mmd_df):,}"},
                {"Item": "Spread formula", "Value": "Issuer yield - active benchmark yield"},
                {"Item": "Attribution policy", "Value": "Rating, sector, callable, and liquidity effects stay separate from benchmark spread."},
            ]
        ),
        width="stretch",
        hide_index=True,
        auto_collapse=False,
    )
    with st.expander("Rating curve assumptions", expanded=False):
        safe_dataframe(rating_spread_table(), width="stretch", hide_index=True, auto_collapse=False)


def _build_snapshot_methodology_cards(
    issuer_df: pd.DataFrame,
    mmd_df: pd.DataFrame,
    benchmark_source_mode: str,
) -> list[dict]:
    """Build compact methodology warnings for the desk snapshot."""
    index_rate = _numeric_rate(issuer_df, "index_rate")
    spread_rate = _numeric_rate(issuer_df, "spread")
    cusip_rate = _nonnull_rate(issuer_df, "cusip")
    rating_col = _first_existing_col(issuer_df, ["ratings_m_s_f", "rating", "ratings", "benchmark_rating"])
    rating_rate = _nonnull_rate(issuer_df, rating_col) if rating_col else None
    benchmark_rows = len(mmd_df) if isinstance(mmd_df, pd.DataFrame) else 0

    if benchmark_source_mode == "Uploaded MMD fallback":
        mmd_status = "good"
        mmd_value = "Uploaded MMD active"
        mmd_detail = f"External MMD is being used as the AAA benchmark curve with {benchmark_rows:,} benchmark row(s)."
    elif benchmark_source_mode == "Trade Sheet Index / Index Rate":
        mmd_status = "warn"
        mmd_value = "Trade index active"
        mmd_detail = f"No external MMD is active in this run; trade-sheet Index Rate is the benchmark source with {benchmark_rows:,} benchmark row(s)."
    else:
        mmd_status = "bad"
        mmd_value = "No active benchmark"
        mmd_detail = "No trade index or uploaded MMD benchmark is available."

    if (index_rate or 0) >= 70:
        index_status = "good"
        index_value = f"{index_rate:.1f}% numeric"
        index_detail = "Trade-sheet Index Rate is available for benchmark spread analytics."
    elif benchmark_source_mode == "Uploaded MMD fallback":
        index_status = "warn"
        index_value = "Index Rate weak / absent"
        index_detail = "Uploaded MMD fallback is active; verify MMD date/tenor coverage before relying on spread outputs."
    elif (spread_rate or 0) >= 70:
        index_status = "warn"
        index_value = f"Spread field {spread_rate:.1f}% numeric"
        index_detail = "Spread field exists, but Index Rate is preferred for transparent benchmark governance."
    else:
        index_status = "bad"
        index_value = "No usable Index Rate"
        index_detail = "Spread-to-benchmark and RV outputs may be degraded."

    if rating_rate is None:
        rating_status = "warn"
        rating_value = "Missing"
        rating_detail = "Ratings are unavailable; peer grouping should fall back to sector and maturity."
    elif rating_rate >= 80:
        rating_status = "good"
        rating_value = f"{rating_rate:.1f}% populated"
        rating_detail = f"Rating field detected: {rating_col}."
    elif rating_rate >= 30:
        rating_status = "warn"
        rating_value = f"{rating_rate:.1f}% populated"
        rating_detail = "Partial ratings coverage; attribution should disclose fallback logic."
    else:
        rating_status = "warn"
        rating_value = f"{rating_rate:.1f}% populated"
        rating_detail = "Ratings are sparse; use sector/maturity fallback for peer grouping."

    return [
        {
            "status": mmd_status,
            "kicker": "MMD / benchmark",
            "title": "AAA curve availability",
            "value": mmd_value,
            "detail": mmd_detail,
        },
        {
            "status": index_status,
            "kicker": "Index Rate",
            "title": "Benchmark input",
            "value": index_value,
            "detail": index_detail,
        },
        {
            "status": rating_status,
            "kicker": "Ratings",
            "title": "Peer grouping input",
            "value": rating_value,
            "detail": rating_detail,
        },
        {
            "status": _rate_status(cusip_rate, good_threshold=95, warn_threshold=80),
            "kicker": "CUSIP quality",
            "title": "CUSIP-level reliability",
            "value": "N/A" if cusip_rate is None else f"{cusip_rate:.1f}% valid",
            "detail": "Low CUSIP quality weakens drilldown, watchlist, and opportunity ranking.",
        },
    ]


def _build_snapshot_takeaway(
    issuer_df: pd.DataFrame,
    market_df: pd.DataFrame,
    cusip_summary: pd.DataFrame,
    selected_issuer: str,
    benchmark_source_mode: str,
) -> tuple[list[str], dict]:
    """Return deterministic analyst takeaway bullets and supporting labels."""
    issuer_base = _add_workflow_spread_bps(issuer_df)
    universe_base = _add_workflow_spread_bps(market_df)
    issuer_spreads = pd.to_numeric(issuer_base.get("spread_bps"), errors="coerce").dropna()
    universe_spreads = pd.to_numeric(universe_base.get("spread_bps"), errors="coerce").dropna()

    spread_label = "Spread unavailable"
    spread_detail = "No usable spread/index-rate data was found."
    if not issuer_spreads.empty and not universe_spreads.empty:
        issuer_median = float(issuer_spreads.median())
        universe_median = float(universe_spreads.median())
        p25 = float(universe_spreads.quantile(0.25))
        p75 = float(universe_spreads.quantile(0.75))
        if issuer_median >= p75:
            spread_label = "Wide / cheaper"
            spread_detail = f"{selected_issuer} median spread is {issuer_median:.1f} bps vs universe median {universe_median:.1f} bps."
        elif issuer_median <= p25:
            spread_label = "Tight / richer"
            spread_detail = f"{selected_issuer} median spread is {issuer_median:.1f} bps vs universe median {universe_median:.1f} bps."
        else:
            spread_label = "Near uploaded universe"
            spread_detail = f"{selected_issuer} median spread is {issuer_median:.1f} bps vs universe median {universe_median:.1f} bps."

    liquidity_label = "Liquidity unavailable"
    liquidity_detail = "CUSIP-level liquidity could not be scored."
    if not cusip_summary.empty and "liquidity_score" in cusip_summary.columns:
        liquidity_scores = pd.to_numeric(cusip_summary["liquidity_score"], errors="coerce").dropna()
        if not liquidity_scores.empty:
            median_liq = float(liquidity_scores.median())
            top_liq = float(liquidity_scores.max())
            if median_liq >= 70:
                liquidity_label = "Liquidity strong"
            elif median_liq >= 45:
                liquidity_label = "Liquidity mixed"
            else:
                liquidity_label = "Liquidity thin"
            liquidity_detail = f"Median CUSIP liquidity score is {median_liq:.1f}; top score is {top_liq:.1f}."

    top_label = "No top CUSIP"
    top_detail = "No CUSIP summary is available."
    if not cusip_summary.empty:
        top = cusip_summary.iloc[0]
        top_label = str(top.get("cusip", "N/A"))
        top_detail = (
            f"{top.get('signal', 'Monitor')} with RV score {_fmt_num(top.get('rv_score'))} "
            f"and liquidity score {_fmt_num(top.get('liquidity_score'))}."
        )

    if benchmark_source_mode == "Trade Sheet Index / Index Rate":
        benchmark_label = "Benchmark OK"
        benchmark_detail = "Using trade-sheet Index / Index Rate; external MMD is not mixed into this run."
    elif benchmark_source_mode == "Uploaded MMD fallback":
        benchmark_label = "MMD fallback active"
        benchmark_detail = "Uploaded MMD is being used as the AAA benchmark curve."
    else:
        benchmark_label = "Benchmark warning"
        benchmark_detail = "No active benchmark source; spread/RV conclusions should be treated as incomplete."

    bullets = [
        f"Spread posture: {spread_label}. {spread_detail}",
        f"Liquidity posture: {liquidity_label}. {liquidity_detail}",
        f"Top CUSIP: {top_label}. {top_detail}",
        f"Benchmark: {benchmark_label}. {benchmark_detail}",
    ]
    labels = {
        "spread_label": spread_label,
        "liquidity_label": liquidity_label,
        "top_label": top_label,
        "benchmark_label": benchmark_label,
    }
    return bullets, labels


def render_focused_upload_audit(
    trade_reports: list[dict],
    bond_report: dict | None,
    mmd_report: dict | None,
    market_df: pd.DataFrame,
    bonds_df: pd.DataFrame,
    issuer_master: pd.DataFrame,
    mmd_df: pd.DataFrame,
    trade_payloads: list[tuple[str, bytes]],
    failed_files: list[str],
    duplicates_removed: int,
    benchmark_source_mode: str,
    benchmark_priority: str,
    benchmark_conflict_policy: str,
    use_external_mmd_fallback: bool,
    mmd_file_provided: bool,
):
    section_anchor("workflow-upload-audit", "Upload / Data Audit")
    st.markdown(
        "<div class='focus-band'>Start here. Confirm files, row counts, field mapping, date coverage, benchmark source, and blocking issues before reading any charts.</div>",
        unsafe_allow_html=True,
    )

    c1, c2, c3, c4, c5 = st.columns(5)
    with c1:
        clean_metric_card("Trade Files", f"{len(trade_payloads):,}", size="small")
    with c2:
        clean_metric_card("Trade Rows", f"{len(market_df):,}", size="small")
    with c3:
        clean_metric_card("Issuers", f"{market_df['issuer'].nunique() if 'issuer' in market_df.columns else 0:,}", size="small")
    with c4:
        clean_metric_card("CUSIPs", f"{market_df['cusip'].nunique() if 'cusip' in market_df.columns else 0:,}", size="small")
    with c5:
        clean_metric_card("Duplicates Removed", f"{duplicates_removed:,}", size="small")

    audit_cards, readiness = build_upload_audit_cards(
        trade_reports=trade_reports,
        market_df=market_df,
        benchmark_source_mode=benchmark_source_mode,
        use_external_mmd_fallback=use_external_mmd_fallback,
        mmd_file_provided=mmd_file_provided,
    )
    render_ready_to_analyze_card(readiness)

    st.subheader("Data Audit Status")
    _render_card_grid(audit_cards, "status-card-grid")

    audit_rows = []
    for report in trade_reports:
        audit_rows.append(
            {
                "File": report.get("dataset"),
                "Rows": report.get("row_count"),
                "Columns": report.get("column_count"),
                "Ready": "Yes" if report.get("can_run") else "No",
                "Missing Required": ", ".join(report.get("missing_required", [])) or "None",
                "Missing Recommended": ", ".join(report.get("missing_recommended", [])) or "None",
            }
        )
    if bond_report:
        audit_rows.append(
            {
                "File": "Optional bond reference",
                "Rows": bond_report.get("row_count", 0),
                "Columns": bond_report.get("column_count", 0),
                "Ready": "Yes" if bond_report.get("can_run", True) else "No",
                "Missing Required": ", ".join(bond_report.get("missing_required", [])) or "None",
                "Missing Recommended": ", ".join(bond_report.get("missing_recommended", [])) or "None",
            }
        )
    if mmd_report:
        audit_rows.append(
            {
                "File": "Optional MMD / benchmark curve",
                "Rows": mmd_report.get("row_count", 0),
                "Columns": mmd_report.get("column_count", 0),
                "Ready": "Yes" if mmd_report.get("can_run", True) else "No",
                "Missing Required": ", ".join(mmd_report.get("missing_required", [])) or "None",
                "Missing Recommended": ", ".join(mmd_report.get("missing_recommended", [])) or "None",
            }
        )
    if audit_rows:
        st.subheader("File Readiness Summary")
        safe_dataframe(pd.DataFrame(audit_rows), width="stretch", hide_index=True, auto_collapse=False)

    coverage_rows = [
        {"Metric": "Trade date coverage", "Value": _workflow_date_range_text(market_df)},
        {"Metric": "Security reference rows", "Value": f"{len(bonds_df):,}"},
        {"Metric": "Issuer master rows", "Value": f"{len(issuer_master):,}"},
        {"Metric": "Benchmark rows", "Value": f"{len(mmd_df):,}"},
        {"Metric": "Failed files", "Value": ", ".join(map(str, failed_files)) if failed_files else "None"},
    ]
    st.subheader("Data Coverage")
    safe_dataframe(pd.DataFrame(coverage_rows), width="stretch", hide_index=True, auto_collapse=False)

    st.subheader("Benchmark Priority / MMD Logic")
    mmd_logic_rows = [
        {"Policy Item": "Uploaded MMD role", "Current Setting": "AAA benchmark curve when external MMD fallback is active"},
        {"Policy Item": "Benchmark priority", "Current Setting": "Trade Sheet Index / Index Rate first; uploaded MMD fallback second"},
        {"Policy Item": "External MMD fallback toggle", "Current Setting": "Enabled" if use_external_mmd_fallback else "Disabled"},
        {"Policy Item": "MMD file provided", "Current Setting": "Yes" if mmd_file_provided else "No"},
        {"Policy Item": "Active benchmark source", "Current Setting": benchmark_source_mode},
        {"Policy Item": "Mixing policy", "Current Setting": "Never mix trade-sheet index rates and uploaded MMD in the same run"},
    ]
    safe_dataframe(pd.DataFrame(mmd_logic_rows), width="stretch", hide_index=True, auto_collapse=False)

    st.subheader("Fixed Benchmark / Methodology Audit")
    _render_benchmark_methodology_block(mmd_df, benchmark_source_mode, benchmark_priority, benchmark_conflict_policy)


def render_focused_snapshot(
    market_df: pd.DataFrame,
    bonds_df: pd.DataFrame,
    issuer_trades: pd.DataFrame,
    issuer_bonds: pd.DataFrame,
    mmd_df: pd.DataFrame,
    selected_issuer: str,
    selected_sector: str,
    benchmark_source_mode: str,
):
    section_anchor("workflow-desk-snapshot", "Desk Snapshot")
    st.markdown(
        "<div class='focus-band'>Decision-first view. Read this before opening detailed charts: coverage, current spread, liquidity, and the strongest CUSIP candidates.</div>",
        unsafe_allow_html=True,
    )
    issuer_base = _add_workflow_spread_bps(issuer_trades)
    cusip_summary = _build_workflow_cusip_summary(issuer_base)
    date_range = _workflow_date_range_text(issuer_base)
    spread_series = pd.to_numeric(issuer_base.get("spread_bps"), errors="coerce") if "spread_bps" in issuer_base.columns else pd.Series(dtype="float64")
    top_candidate = cusip_summary.iloc[0] if not cusip_summary.empty else None
    cusip_count = issuer_base["cusip"].nunique() if "cusip" in issuer_base.columns else 0
    top_liquidity = None if top_candidate is None else top_candidate.get("liquidity_score")
    top_candidate_text = "N/A" if top_candidate is None else str(top_candidate.get("cusip", "N/A"))
    top_candidate_note = None
    if top_candidate is not None:
        top_candidate_note = f"{top_candidate.get('signal', 'Monitor')} | liquidity {_fmt_num(top_liquidity)}"

    c1, c2, c3, c4, c5, c6 = st.columns(6)
    with c1:
        clean_metric_card("Issuer", selected_issuer, size="small", note=selected_sector)
    with c2:
        clean_metric_card("Date Range", date_range, size="small")
    with c3:
        clean_metric_card("Trade Rows", f"{len(issuer_trades):,}", size="small")
    with c4:
        clean_metric_card("CUSIPs", f"{cusip_count:,}", size="small")
    with c5:
        clean_metric_card("Median Spread", "N/A" if spread_series.dropna().empty else f"{spread_series.median():.1f} bps", size="small")
    with c6:
        clean_metric_card("Liquidity / Top", top_candidate_text, size="small", note=top_candidate_note)

    st.subheader("Analyst Takeaway")
    bullets, takeaway_labels = _build_snapshot_takeaway(
        issuer_df=issuer_base,
        market_df=market_df,
        cusip_summary=cusip_summary,
        selected_issuer=selected_issuer,
        benchmark_source_mode=benchmark_source_mode,
    )
    for bullet in bullets:
        st.markdown(f"- {bullet}")

    if not cusip_summary.empty:
        st.subheader("Top 5 Opportunities")
        opps = cusip_summary.head(5).copy()
        if "current_spread_bps" in opps.columns:
            opps["current_spread_bps"] = pd.to_numeric(opps["current_spread_bps"], errors="coerce").round(1)
        for col in ["liquidity_score", "rv_score"]:
            if col in opps.columns:
                opps[col] = pd.to_numeric(opps[col], errors="coerce").round(1)
        if "total_trade_amount" in opps.columns:
            opps["total_trade_amount"] = pd.to_numeric(opps["total_trade_amount"], errors="coerce").round(0)
        opps["snapshot_reason"] = opps.apply(
            lambda row: (
                f"{row.get('signal', 'Monitor')}; spread {_fmt_bps(row.get('current_spread_bps'))}; "
                f"liquidity {_fmt_num(row.get('liquidity_score'))}; RV {_fmt_num(row.get('rv_score'))}"
            ),
            axis=1,
        )
        display_cols = [
            "cusip", "signal", "maturity_bucket", "current_spread_bps", "liquidity_score",
            "rv_score", "trade_count", "total_trade_amount", "latest_trade", "snapshot_reason",
        ]
        safe_dataframe(opps[[c for c in display_cols if c in opps.columns]], hide_index=True, auto_collapse=False)
    else:
        st.info("No CUSIP-level opportunity table is available for the current filter.")

    st.subheader("Methodology Warnings")
    warning_cards = _build_snapshot_methodology_cards(
        issuer_df=issuer_base,
        mmd_df=mmd_df,
        benchmark_source_mode=benchmark_source_mode,
    )
    _render_card_grid(warning_cards, "status-card-grid")

    with st.expander("Snapshot calculation notes", expanded=False):
        st.markdown(
            f"""
- **Spread posture** compares the selected issuer median spread against the uploaded universe distribution.
- **Liquidity posture** uses the CUSIP-level liquidity score from trade count, par amount, and recency.
- **Top opportunity ranking** uses the focused workflow RV score, which combines spread rank and liquidity rank.
- **Benchmark source shown:** `{benchmark_source_mode}`.
- **Current snapshot labels:** spread = `{takeaway_labels.get('spread_label')}`, liquidity = `{takeaway_labels.get('liquidity_label')}`, benchmark = `{takeaway_labels.get('benchmark_label')}`.
            """
        )


def render_focused_core_charts(
    market_df: pd.DataFrame,
    issuer_trades: pd.DataFrame,
    mmd_df: pd.DataFrame,
    selected_issuer: str,
    comparison_issuers: list[str],
    selected_sector: str,
):
    section_anchor("workflow-core-charts", "Core Charts")
    st.markdown(
        "<div class='focus-band'>Core visual analysis only: spread trend, yield trend, trading volume, and maturity-year curve. Use this page when you want charts without the long audit/admin sections.</div>",
        unsafe_allow_html=True,
    )
    if market_df.empty or "issuer" not in market_df.columns:
        st.info("No uploaded market data is available for charts.")
        return

    all_issuers = sorted(market_df["issuer"].dropna().astype(str).unique().tolist())
    chart_issuers = [selected_issuer] + [x for x in comparison_issuers if x != selected_issuer]
    chart_issuers = [x for x in chart_issuers if x in all_issuers]
    if not chart_issuers:
        chart_issuers = [selected_issuer]

    chart_base_all = _add_workflow_spread_bps(market_df.copy())
    chart_base_all["trade_date"] = pd.to_datetime(chart_base_all.get("trade_date"), errors="coerce")
    if "trade_amount" in chart_base_all.columns:
        chart_base_all["trade_amount"] = pd.to_numeric(chart_base_all["trade_amount"], errors="coerce").fillna(0)
    else:
        chart_base_all["trade_amount"] = 0
    if "yield" in chart_base_all.columns:
        chart_base_all["yield"] = pd.to_numeric(chart_base_all["yield"], errors="coerce")

    chart_dates = chart_base_all["trade_date"].dropna() if "trade_date" in chart_base_all.columns else pd.Series(dtype="datetime64[ns]")
    if chart_dates.empty:
        st.warning("Core charts require valid trade_date values.")
        return

    ctrl1, ctrl2, ctrl3 = st.columns([1, 1, 1.4])
    with ctrl1:
        trend_frequency = st.selectbox(
            "Trend Frequency",
            ["Daily", "Weekly", "Monthly"],
            index=1,
            key="focused_core_trend_frequency",
        )
    with ctrl2:
        curve_benchmark_rating = st.selectbox(
            "Curve Benchmark",
            BENCHMARK_RATINGS,
            index=0,
            key="focused_core_curve_benchmark",
        )
    with ctrl3:
        reference_lines = st.multiselect(
            "Reference Lines",
            ["Sector median", "All uploaded median", "AAA/MMD baseline", "MMD benchmark curve"],
            default=["Sector median", "All uploaded median"],
            key="focused_core_reference_lines",
            help="MMD benchmark curve applies to the issuer curve when benchmark data is available. AAA/MMD baseline is 0 bps on spread charts.",
        )

    date_min = chart_dates.min().date()
    date_max = chart_dates.max().date()
    selected_chart_dates = st.date_input(
        "Chart Date Range",
        value=(date_min, date_max),
        min_value=date_min,
        max_value=date_max,
        key="focused_core_chart_date_range",
        help="Filters the focused core charts. Sidebar zoom only changes the visual viewport; this changes the chart data.",
    )
    chart_start_date, chart_end_date = date_min, date_max
    if isinstance(selected_chart_dates, tuple) and len(selected_chart_dates) == 2:
        chart_start_date, chart_end_date = selected_chart_dates
        chart_base_all = chart_base_all[
            (chart_base_all["trade_date"].dt.date >= chart_start_date)
            & (chart_base_all["trade_date"].dt.date <= chart_end_date)
        ].copy()

    freq_map = {"Daily": "D", "Weekly": "W", "Monthly": "M"}
    period_freq = freq_map.get(trend_frequency, "W")
    chart_base = chart_base_all[chart_base_all["issuer"].astype(str).isin(chart_issuers)].copy()

    st.subheader("Spread Trend")
    if not chart_base.empty and {"trade_date", "spread_bps", "issuer"}.issubset(chart_base.columns):
        spread_points = chart_base.dropna(subset=["trade_date", "spread_bps"]).copy()
        spread_trend = (
            spread_points.groupby([pd.Grouper(key="trade_date", freq=period_freq), "issuer"], as_index=False)
            .agg(
                spread_bps=("spread_bps", "median"),
                avg_yield=("yield", "mean") if "yield" in spread_points.columns else ("spread_bps", "count"),
                trade_count=("spread_bps", "count"),
                total_par=("trade_amount", "sum"),
            )
            .dropna(subset=["spread_bps"])
            .sort_values("trade_date")
        )
        fig = go.Figure()
        for issuer_name in chart_issuers:
            tmp = spread_trend[spread_trend["issuer"].astype(str) == str(issuer_name)]
            if tmp.empty:
                continue
            fig.add_trace(
                go.Scatter(
                    x=tmp["trade_date"],
                    y=tmp["spread_bps"],
                    mode="lines+markers",
                    name=issuer_name,
                    line=dict(width=3.2 if issuer_name == selected_issuer else 2.1),
                    customdata=np.stack(
                        [
                            tmp["trade_count"].fillna(0),
                            tmp["total_par"].fillna(0),
                            tmp["avg_yield"].fillna(np.nan),
                        ],
                        axis=-1,
                    ),
                    hovertemplate=(
                        "%{x|%m/%d/%Y}<br>"
                        "Spread: %{y:.1f} bps<br>"
                        "Trades: %{customdata[0]:,.0f}<br>"
                        "Par: $%{customdata[1]:,.0f}<br>"
                        "Avg yield: %{customdata[2]:.3f}%"
                        "<extra>%{fullData.name}</extra>"
                    ),
                )
            )

        reference_base = chart_base_all.dropna(subset=["trade_date", "spread_bps"]).copy()
        if "Sector median" in reference_lines and selected_sector and selected_sector != "Unknown" and "sector" in reference_base.columns:
            sector_ref = reference_base[reference_base["sector"].astype(str) == str(selected_sector)].copy()
            sector_ref = (
                sector_ref.groupby(pd.Grouper(key="trade_date", freq=period_freq), as_index=False)
                .agg(spread_bps=("spread_bps", "median"), trade_count=("spread_bps", "count"))
                .dropna(subset=["spread_bps"])
            )
            if not sector_ref.empty:
                fig.add_trace(
                    go.Scatter(
                        x=sector_ref["trade_date"],
                        y=sector_ref["spread_bps"],
                        mode="lines",
                        name=f"{selected_sector} median",
                        line=dict(width=2, dash="dash"),
                        hovertemplate="%{x|%m/%d/%Y}<br>Spread: %{y:.1f} bps<extra>%{fullData.name}</extra>",
                    )
                )
        if "All uploaded median" in reference_lines:
            all_ref = (
                reference_base.groupby(pd.Grouper(key="trade_date", freq=period_freq), as_index=False)
                .agg(spread_bps=("spread_bps", "median"), trade_count=("spread_bps", "count"))
                .dropna(subset=["spread_bps"])
            )
            if not all_ref.empty:
                fig.add_trace(
                    go.Scatter(
                        x=all_ref["trade_date"],
                        y=all_ref["spread_bps"],
                        mode="lines",
                        name="All uploaded median",
                        line=dict(width=2, dash="dot"),
                        hovertemplate="%{x|%m/%d/%Y}<br>Spread: %{y:.1f} bps<extra>%{fullData.name}</extra>",
                    )
                )
        if "AAA/MMD baseline" in reference_lines and not spread_trend.empty:
            fig.add_hline(y=0, line_dash="longdash", line_width=1.5, annotation_text="AAA/MMD baseline")

        if fig.data:
            fig.update_layout(
                title=f"{selected_issuer} Spread Trend with Reference Lines",
                xaxis_title="Trade Date",
                yaxis_title="Spread (bps)",
                hovermode="x unified",
                height=520,
                legend_title_text="Line Item",
                margin=dict(l=40, r=40, t=70, b=45),
            )
            safe_plotly_chart(fig, width="stretch")
            st.caption("Data table for spread trend")
            safe_dataframe(spread_trend, hide_index=True, top_rows=8)
        else:
            st.info("No spread trend traces were available for the selected filters.")
    else:
        st.info("Spread trend needs issuer, trade_date, and spread or index_rate/yield fields.")

    st.subheader("Volume & Activity")
    if not chart_base.empty and {"trade_date", "trade_amount", "issuer"}.issubset(chart_base.columns):
        vol = chart_base.copy()
        vol = vol.dropna(subset=["trade_date"])
        vol["period"] = vol["trade_date"].dt.to_period(period_freq).dt.to_timestamp()
        vol["volume_group"] = vol["issuer"].astype(str)
        volume_by_group = (
            vol.groupby(["period", "volume_group"], as_index=False)
            .agg(volume=("trade_amount", "sum"), trade_count=("trade_amount", "count"))
            .sort_values("period")
        )
        volume_total = (
            vol.groupby("period", as_index=False)
            .agg(total_volume=("trade_amount", "sum"), total_trade_count=("trade_amount", "count"))
            .sort_values("period")
        )
        selected_volume = (
            vol[vol["issuer"].astype(str) == str(selected_issuer)]
            .groupby("period", as_index=False)
            .agg(selected_volume=("trade_amount", "sum"))
        )
        volume_total = volume_total.merge(selected_volume, on="period", how="left")
        volume_total["selected_volume"] = volume_total["selected_volume"].fillna(0)
        volume_total["selected_issuer_share"] = np.where(
            volume_total["total_volume"] > 0,
            volume_total["selected_volume"] / volume_total["total_volume"] * 100,
            np.nan,
        )

        fig_vol = make_subplots(
            rows=2,
            cols=1,
            shared_xaxes=True,
            vertical_spacing=0.08,
            row_heights=[0.68, 0.32],
            specs=[[{}], [{"secondary_y": True}]],
            subplot_titles=("Trading Volume", "Activity & Selected Issuer Share"),
        )
        for group_name in chart_issuers:
            tmp = volume_by_group[volume_by_group["volume_group"].astype(str) == str(group_name)]
            if tmp.empty:
                continue
            fig_vol.add_trace(
                go.Bar(
                    x=tmp["period"],
                    y=tmp["volume"] / 1_000_000,
                    name=group_name,
                    customdata=np.stack([tmp["trade_count"].fillna(0)], axis=-1),
                    hovertemplate="%{x|%m/%d/%Y}<br>Volume: $%{y:,.1f}M<br>Trades: %{customdata[0]:,.0f}<extra>%{fullData.name}</extra>",
                ),
                row=1,
                col=1,
            )
        fig_vol.add_trace(
            go.Scatter(
                x=volume_total["period"],
                y=volume_total["total_trade_count"],
                mode="lines+markers",
                name="Total trade count",
                line=dict(width=2.4),
                hovertemplate="%{x|%m/%d/%Y}<br>Total trades: %{y:,.0f}<extra>Total trade count</extra>",
            ),
            row=2,
            col=1,
            secondary_y=False,
        )
        fig_vol.add_trace(
            go.Scatter(
                x=volume_total["period"],
                y=volume_total["selected_issuer_share"],
                mode="lines+markers",
                name=f"{selected_issuer} volume share",
                line=dict(width=2.4, dash="dash"),
                hovertemplate="%{x|%m/%d/%Y}<br>Share: %{y:.1f}%<extra>%{fullData.name}</extra>",
            ),
            row=2,
            col=1,
            secondary_y=True,
        )
        fig_vol.update_layout(
            title=f"{trend_frequency} Volume, Trade Count, and {selected_issuer} Share",
            barmode="stack",
            height=680,
            hovermode="x unified",
            legend_title_text="Series",
            margin=dict(l=40, r=50, t=85, b=45),
        )
        fig_vol.update_yaxes(title_text="Volume ($MM)", row=1, col=1)
        fig_vol.update_yaxes(title_text="Trade Count", row=2, col=1, secondary_y=False)
        fig_vol.update_yaxes(title_text=f"{selected_issuer} Share", ticksuffix="%", row=2, col=1, secondary_y=True)
        safe_plotly_chart(fig_vol, width="stretch")

        volume_table = volume_by_group.merge(volume_total[["period", "total_volume", "total_trade_count", "selected_issuer_share"]], on="period", how="left")
        st.caption("Data table for volume and activity")
        safe_dataframe(volume_table, hide_index=True, top_rows=8)
    else:
        st.info("Volume chart needs trade_date, trade_amount, and issuer fields.")

    st.subheader("Issuer Curve")
    curve_source = issuer_trades.copy()
    if curve_source.empty or not {"maturity_year", "yield"}.issubset(curve_source.columns):
        st.info("Issuer curve needs maturity_year and yield fields.")
        return

    curve_lookback = st.select_slider(
        "Issuer Curve Lookback",
        options=[7, 14, 30, 60, 90, 180, 365],
        value=60,
        format_func=lambda x: f"{x} days",
        key="focused_core_curve_lookback",
        help="Uses selected issuer trades inside this lookback window ending on the latest selected trade date.",
    )
    curve_source["trade_date"] = pd.to_datetime(curve_source.get("trade_date"), errors="coerce")
    if "trade_date" in curve_source.columns:
        curve_source = curve_source[
            (curve_source["trade_date"].dt.date >= chart_start_date)
            & (curve_source["trade_date"].dt.date <= chart_end_date)
        ].copy()
    latest_curve_date = curve_source["trade_date"].dropna().max()
    if pd.notna(latest_curve_date):
        curve_source = curve_source[curve_source["trade_date"] >= latest_curve_date - pd.Timedelta(days=int(curve_lookback))].copy()
    curve_source["maturity_year"] = pd.to_numeric(curve_source["maturity_year"], errors="coerce")
    curve_source["yield"] = pd.to_numeric(curve_source["yield"], errors="coerce")
    issuer_curve = (
        curve_source.dropna(subset=["maturity_year", "yield"])
        .groupby("maturity_year", as_index=False)
        .agg(avg_yield=("yield", "mean"), trade_count=("yield", "count"), total_par=("trade_amount", "sum"), latest_trade=("trade_date", "max"))
        .sort_values("maturity_year")
    )
    if issuer_curve.empty:
        st.info("No curve observations were available for the selected lookback.")
        return

    curve_fig = go.Figure()
    curve_fig.add_trace(
        go.Scatter(
            x=issuer_curve["maturity_year"],
            y=issuer_curve["avg_yield"],
            mode="lines+markers",
            name=f"{selected_issuer} issuer curve",
            line=dict(width=3.2),
            customdata=np.stack([issuer_curve["trade_count"].fillna(0), issuer_curve["total_par"].fillna(0)], axis=-1),
            hovertemplate=(
                "%{x:.0f}Y<br>Yield: %{y:.3f}%<br>"
                "Trades: %{customdata[0]:,.0f}<br>Par: $%{customdata[1]:,.0f}"
                "<extra>%{fullData.name}</extra>"
            ),
        )
    )

    curve_universe = chart_base_all.copy()
    if pd.notna(latest_curve_date):
        curve_universe = curve_universe[curve_universe["trade_date"] >= latest_curve_date - pd.Timedelta(days=int(curve_lookback))].copy()
    if {"sector", "maturity_year", "yield"}.issubset(curve_universe.columns) and "Sector median" in reference_lines and selected_sector != "Unknown":
        sector_curve = curve_universe[curve_universe["sector"].astype(str) == str(selected_sector)].copy()
        sector_curve["maturity_year"] = pd.to_numeric(sector_curve["maturity_year"], errors="coerce")
        sector_curve["yield"] = pd.to_numeric(sector_curve["yield"], errors="coerce")
        sector_curve = (
            sector_curve.dropna(subset=["maturity_year", "yield"])
            .groupby("maturity_year", as_index=False)
            .agg(avg_yield=("yield", "median"), trade_count=("yield", "count"))
            .sort_values("maturity_year")
        )
        if not sector_curve.empty:
            curve_fig.add_trace(
                go.Scatter(
                    x=sector_curve["maturity_year"],
                    y=sector_curve["avg_yield"],
                    mode="lines+markers",
                    name=f"{selected_sector} median curve",
                    line=dict(width=2, dash="dash"),
                    hovertemplate="%{x:.0f}Y<br>Yield: %{y:.3f}%<extra>%{fullData.name}</extra>",
                )
            )
    if {"maturity_year", "yield"}.issubset(curve_universe.columns) and "All uploaded median" in reference_lines:
        all_curve = curve_universe.copy()
        all_curve["maturity_year"] = pd.to_numeric(all_curve["maturity_year"], errors="coerce")
        all_curve["yield"] = pd.to_numeric(all_curve["yield"], errors="coerce")
        all_curve = (
            all_curve.dropna(subset=["maturity_year", "yield"])
            .groupby("maturity_year", as_index=False)
            .agg(avg_yield=("yield", "median"), trade_count=("yield", "count"))
            .sort_values("maturity_year")
        )
        if not all_curve.empty:
            curve_fig.add_trace(
                go.Scatter(
                    x=all_curve["maturity_year"],
                    y=all_curve["avg_yield"],
                    mode="lines+markers",
                    name="All uploaded median curve",
                    line=dict(width=2, dash="dot"),
                    hovertemplate="%{x:.0f}Y<br>Yield: %{y:.3f}%<extra>%{fullData.name}</extra>",
                )
            )

    benchmark_curve = pd.DataFrame()
    if "MMD benchmark curve" in reference_lines and isinstance(mmd_df, pd.DataFrame) and not mmd_df.empty:
        date_col = _detect_mmd_date_column(mmd_df)
        if date_col:
            mmd_work = mmd_df.copy()
            mmd_work[date_col] = pd.to_datetime(mmd_work[date_col], errors="coerce").dt.normalize()
            mmd_work = mmd_work.dropna(subset=[date_col])
            if pd.notna(latest_curve_date):
                mmd_work = mmd_work[mmd_work[date_col] <= latest_curve_date.normalize()].copy()
            if not mmd_work.empty:
                rows = []
                for year in sorted(issuer_curve["maturity_year"].dropna().astype(int).unique().tolist()):
                    bucket = f"{year}Y"
                    tenor = MMD_BUCKET_MAP.get(bucket, "10Y")
                    y, meta = get_benchmark_curve(mmd_work, tenor, curve_benchmark_rating)
                    if y is None:
                        continue
                    bench_tmp = pd.DataFrame({"benchmark_yield": pd.to_numeric(y, errors="coerce")}).dropna()
                    if bench_tmp.empty:
                        continue
                    rows.append(
                        {
                            "maturity_year": year,
                            "benchmark_yield": float(bench_tmp["benchmark_yield"].iloc[-1]),
                            "benchmark_rating": curve_benchmark_rating,
                            "mmd_tenor": tenor,
                            "benchmark_source": meta.get("benchmark_source"),
                            "source_column": meta.get("source_column"),
                            "rating_spread_bps": meta.get("rating_spread_bps"),
                        }
                    )
                benchmark_curve = pd.DataFrame(rows)
                if not benchmark_curve.empty:
                    curve_fig.add_trace(
                        go.Scatter(
                            x=benchmark_curve["maturity_year"],
                            y=benchmark_curve["benchmark_yield"],
                            mode="lines+markers",
                            name=f"{curve_benchmark_rating} benchmark curve",
                            line=dict(width=2.4, dash="longdash"),
                            hovertemplate="%{x:.0f}Y<br>Yield: %{y:.3f}%<extra>%{fullData.name}</extra>",
                        )
                    )

    curve_fig.update_layout(
        title=f"{selected_issuer} Issuer Curve vs References",
        xaxis_title="Maturity Year",
        yaxis_title="Yield (%)",
        hovermode="x unified",
        height=540,
        legend_title_text="Curve",
        margin=dict(l=40, r=40, t=70, b=45),
    )
    safe_plotly_chart(curve_fig, width="stretch")

    curve_table = issuer_curve.copy()
    if not benchmark_curve.empty:
        curve_table = curve_table.merge(benchmark_curve, on="maturity_year", how="left")
        curve_table["spread_to_benchmark_bps"] = (
            pd.to_numeric(curve_table["avg_yield"], errors="coerce")
            - pd.to_numeric(curve_table["benchmark_yield"], errors="coerce")
        ) * 100
    st.caption("Data table for issuer curve")
    display_cols = [
        "maturity_year", "avg_yield", "benchmark_rating", "benchmark_yield",
        "spread_to_benchmark_bps", "trade_count", "total_par", "latest_trade",
        "mmd_tenor", "benchmark_source", "source_column", "rating_spread_bps",
    ]
    safe_dataframe(curve_table[[c for c in display_cols if c in curve_table.columns]], hide_index=True, top_rows=12)


def render_focused_cusip_drilldown(issuer_trades: pd.DataFrame, selected_issuer: str):
    section_anchor("workflow-cusip-drilldown", "CUSIP Drilldown")
    st.markdown(
        "<div class='focus-band'>Security-level workflow: choose one CUSIP, review detail metrics, inspect trade path, then compare same-bucket peers.</div>",
        unsafe_allow_html=True,
    )
    summary = _build_workflow_cusip_summary(issuer_trades)
    if summary.empty:
        st.info("No CUSIP-level rows are available for the selected issuer/filter.")
        return

    selector_options = summary["cusip"].dropna().astype(str).tolist()
    selected_cusip = st.selectbox("Select CUSIP", selector_options)
    selected_row = summary[summary["cusip"].astype(str) == str(selected_cusip)].iloc[0]
    detail = _add_workflow_spread_bps(issuer_trades[issuer_trades["cusip"].astype(str) == str(selected_cusip)].copy())
    detail["trade_date"] = pd.to_datetime(detail.get("trade_date"), errors="coerce")
    for col in ["yield", "price", "trade_amount", "spread_bps"]:
        if col in detail.columns:
            detail[col] = pd.to_numeric(detail[col], errors="coerce")
    if "trade_amount" not in detail.columns:
        detail["trade_amount"] = 0.0
    if "yield" not in detail.columns:
        detail["yield"] = pd.NA
    if "price" not in detail.columns:
        detail["price"] = pd.NA

    detail_sorted = detail.sort_values("trade_date").copy()
    latest_trade_row = detail_sorted.dropna(subset=["trade_date"]).tail(1)
    latest_date = latest_trade_row["trade_date"].iloc[0] if not latest_trade_row.empty else pd.NaT
    latest_yield = latest_trade_row["yield"].iloc[0] if not latest_trade_row.empty and "yield" in latest_trade_row.columns else pd.NA
    latest_price = latest_trade_row["price"].iloc[0] if not latest_trade_row.empty and "price" in latest_trade_row.columns else pd.NA
    total_par = pd.to_numeric(detail_sorted["trade_amount"], errors="coerce").sum()
    avg_trade_size = pd.to_numeric(detail_sorted["trade_amount"], errors="coerce").mean()

    path = (
        detail_sorted.dropna(subset=["trade_date"])
        .groupby("trade_date", as_index=False)
        .agg(
            spread_bps=("spread_bps", "median"),
            avg_yield=("yield", "mean"),
            avg_price=("price", "mean"),
            par=("trade_amount", "sum"),
            trade_count=("trade_amount", "count"),
        )
        .sort_values("trade_date")
    )
    spread_change = pd.NA
    if not path.empty and pd.to_numeric(path["spread_bps"], errors="coerce").notna().sum() >= 2:
        clean_path_spread = path.dropna(subset=["spread_bps"])
        spread_change = float(clean_path_spread["spread_bps"].iloc[-1] - clean_path_spread["spread_bps"].iloc[0])

    side_col = _first_existing_col(detail, ["trade_type", "side", "buy_sell", "customer_side", "dealer_side"])
    if side_col:
        detail["flow_side"] = detail[side_col].map(_focused_trade_side)
    else:
        detail["flow_side"] = "Unknown"
    buy_count = int((detail["flow_side"] == "Buy").sum())
    sell_count = int((detail["flow_side"] == "Sell").sum())
    other_count = int((~detail["flow_side"].isin(["Buy", "Sell"])).sum())

    c1, c2, c3, c4, c5 = st.columns(5)
    with c1:
        clean_metric_card("CUSIP", selected_cusip, size="small")
    with c2:
        clean_metric_card("Signal", selected_row.get("signal"), size="small")
    with c3:
        clean_metric_card("Spread", _fmt_bps(selected_row.get("current_spread_bps")), size="small")
    with c4:
        clean_metric_card("Liquidity", _fmt_num(selected_row.get("liquidity_score")), size="small")
    with c5:
        clean_metric_card("Trades", f"{int(selected_row.get('trade_count', 0)):,}", size="small")

    d1, d2, d3, d4, d5 = st.columns(5)
    with d1:
        clean_metric_card("Latest Trade", _fmt_date(latest_date) if pd.notna(latest_date) else "N/A", size="small")
    with d2:
        clean_metric_card("Latest Yield", _fmt_pct(latest_yield), size="small")
    with d3:
        clean_metric_card("Latest Price", _fmt_num(latest_price), size="small")
    with d4:
        clean_metric_card("Total Par", _fmt_mm(total_par), size="small")
    with d5:
        clean_metric_card("Path Change", _fmt_bps(spread_change), size="small")

    st.subheader("Analyst Read-Through")
    readthrough = [
        f"{selected_cusip} has {len(detail):,} trade observation(s) in the current issuer/filter.",
        f"Latest spread screens at {_fmt_bps(selected_row.get('current_spread_bps'))}; liquidity score is {_fmt_num(selected_row.get('liquidity_score'))}.",
        f"Trade flow proxy: {buy_count:,} buy, {sell_count:,} sell, {other_count:,} other/unknown trade(s).",
    ]
    if pd.notna(spread_change):
        readthrough.append(f"Observed spread path moved {_fmt_bps(spread_change)} from first to latest trade in the selected window.")
    bucket = selected_row.get("maturity_bucket") if "maturity_bucket" in summary.columns else None
    if pd.notna(bucket):
        readthrough.append(f"Same-bucket comparison uses {bucket} peers from the selected issuer universe.")
    for line in readthrough:
        st.markdown(f"- {line}")

    records = _focused_watchlist_records()
    existing_note = records.get(str(selected_cusip), {}).get("note", "")
    note_col, action_col = st.columns([2.2, 1])
    with note_col:
        watch_note = st.text_area(
            "Watchlist note",
            value=existing_note,
            key=f"cusip_watch_note_{selected_cusip}",
            height=86,
            placeholder="Why this CUSIP is worth saving, what to verify, or how to frame it in the report.",
        )
    with action_col:
        st.caption("Save this CUSIP with the note so it carries into RV / Watchlist exports.")
        if st.button("Save / Update Watchlist", key=f"save_watch_{selected_cusip}"):
            _upsert_focused_watchlist(selected_cusip, selected_issuer, "CUSIP Drilldown", selected_row, watch_note)
            st.success(f"Saved {selected_cusip} to watchlist.")

    st.subheader("Trade Path")
    if not path.empty:
        fig_path = make_subplots(
            rows=3,
            cols=1,
            shared_xaxes=True,
            vertical_spacing=0.08,
            row_heights=[0.42, 0.30, 0.28],
            subplot_titles=("Spread Path", "Yield / Price", "Par Amount"),
            specs=[[{}], [{"secondary_y": True}], [{}]],
        )
        if pd.to_numeric(path["spread_bps"], errors="coerce").notna().any():
            fig_path.add_trace(
                go.Scatter(
                    x=path["trade_date"],
                    y=path["spread_bps"],
                    mode="lines+markers",
                    name="Spread",
                    line=dict(width=3),
                    customdata=np.stack([path["trade_count"].fillna(0), path["par"].fillna(0)], axis=-1),
                    hovertemplate="%{x|%m/%d/%Y}<br>Spread: %{y:.1f} bps<br>Trades: %{customdata[0]:,.0f}<br>Par: $%{customdata[1]:,.0f}<extra>Spread</extra>",
                ),
                row=1,
                col=1,
            )
        if pd.to_numeric(path["avg_yield"], errors="coerce").notna().any():
            fig_path.add_trace(
                go.Scatter(
                    x=path["trade_date"],
                    y=path["avg_yield"],
                    mode="lines+markers",
                    name="Yield",
                    line=dict(width=2.4),
                    hovertemplate="%{x|%m/%d/%Y}<br>Yield: %{y:.3f}%<extra>Yield</extra>",
                ),
                row=2,
                col=1,
                secondary_y=False,
            )
        if pd.to_numeric(path["avg_price"], errors="coerce").notna().any():
            fig_path.add_trace(
                go.Scatter(
                    x=path["trade_date"],
                    y=path["avg_price"],
                    mode="lines+markers",
                    name="Price",
                    line=dict(width=2.0, dash="dash"),
                    hovertemplate="%{x|%m/%d/%Y}<br>Price: %{y:.2f}<extra>Price</extra>",
                ),
                row=2,
                col=1,
                secondary_y=True,
            )
        fig_path.add_trace(
            go.Bar(
                x=path["trade_date"],
                y=path["par"],
                name="Par amount",
                hovertemplate="%{x|%m/%d/%Y}<br>Par: $%{y:,.0f}<extra>Par amount</extra>",
            ),
            row=3,
            col=1,
        )
        fig_path.update_layout(
            title=f"{selected_cusip} Trade Path",
            height=720,
            hovermode="x unified",
            legend_title_text="Series",
            margin=dict(l=40, r=50, t=85, b=45),
        )
        fig_path.update_yaxes(title_text="Spread (bps)", row=1, col=1)
        fig_path.update_yaxes(title_text="Yield (%)", row=2, col=1, secondary_y=False)
        fig_path.update_yaxes(title_text="Price", row=2, col=1, secondary_y=True)
        fig_path.update_yaxes(title_text="Par", row=3, col=1)
        safe_plotly_chart(fig_path, width="stretch")
        with st.expander("Trade path data", expanded=False):
            safe_dataframe(path, hide_index=True, auto_collapse=False)
    else:
        st.info("No dated trade path is available for the selected CUSIP.")

    st.subheader("Trade Detail")
    display_cols = ["trade_date", "trade_type", "yield", "price", "trade_amount", "spread_bps", "maturity_bucket", "description"]
    safe_dataframe(detail[[c for c in display_cols if c in detail.columns]].sort_values("trade_date", ascending=False), hide_index=True)

    if pd.notna(bucket):
        st.subheader("Same-Bucket Peers")
        peers = summary[summary["maturity_bucket"].astype(str) == str(bucket)].copy()
        peers["is_selected"] = peers["cusip"].astype(str).eq(str(selected_cusip))
        peer_median_spread = pd.to_numeric(peers["current_spread_bps"], errors="coerce").median()
        peers["peer_median_gap_bps"] = pd.to_numeric(peers["current_spread_bps"], errors="coerce") - peer_median_spread
        peers = peers.sort_values(["rv_score", "liquidity_score", "trade_count"], ascending=False)
        peer_cols = [
            "cusip", "is_selected", "signal", "current_spread_bps", "peer_median_gap_bps",
            "liquidity_score", "rv_score", "trade_count", "total_trade_amount", "latest_trade",
        ]
        safe_dataframe(peers[[c for c in peer_cols if c in peers.columns]].head(20), hide_index=True)

        if not peers.empty and pd.notna(peer_median_spread):
            selected_gap = peers.loc[peers["is_selected"], "peer_median_gap_bps"]
            selected_gap_val = selected_gap.iloc[0] if not selected_gap.empty else pd.NA
            st.info(
                f"Same-bucket median spread is {_fmt_bps(peer_median_spread)}. "
                f"{selected_cusip} screens {_fmt_bps(selected_gap_val)} versus that peer median."
            )


def render_focused_rv_watchlist(issuer_trades: pd.DataFrame, selected_issuer: str):
    section_anchor("workflow-rv-watchlist", "RV / Watchlist")
    st.markdown(
        "<div class='focus-band'>Ranking page for candidate discovery. Save CUSIPs during the session, then export the shortlist.</div>",
        unsafe_allow_html=True,
    )
    summary = _build_workflow_cusip_summary(issuer_trades)
    if summary.empty:
        st.info("No CUSIP-level rows are available for RV ranking.")
        return

    summary = summary.copy()
    if "maturity_bucket" in summary.columns:
        summary["peer_median_spread_bps"] = summary.groupby("maturity_bucket")["current_spread_bps"].transform("median")
        summary["peer_median_gap_bps"] = pd.to_numeric(summary["current_spread_bps"], errors="coerce") - pd.to_numeric(summary["peer_median_spread_bps"], errors="coerce")
    else:
        summary["peer_median_gap_bps"] = pd.NA

    filt1, filt2, filt3, filt4 = st.columns([1, 1, 1.2, 1.2])
    with filt1:
        min_liq = st.slider("Minimum liquidity score", 0, 100, 40)
    with filt2:
        min_trades = st.number_input("Minimum trade count", min_value=1, max_value=1000, value=2, step=1)
    with filt3:
        signal_options = sorted(summary["signal"].dropna().astype(str).unique().tolist()) if "signal" in summary.columns else []
        selected_signals = st.multiselect(
            "Signals",
            signal_options,
            default=signal_options,
            key="focused_rv_signal_filter",
        )
    with filt4:
        bucket_options = sorted(summary["maturity_bucket"].dropna().astype(str).unique().tolist()) if "maturity_bucket" in summary.columns else []
        selected_buckets = st.multiselect(
            "Maturity buckets",
            bucket_options,
            default=bucket_options,
            key="focused_rv_bucket_filter",
        )

    ranked = summary[
        (pd.to_numeric(summary["liquidity_score"], errors="coerce") >= min_liq)
        & (pd.to_numeric(summary["trade_count"], errors="coerce") >= min_trades)
    ].copy()
    if signal_options and "signal" in ranked.columns:
        ranked = ranked[ranked["signal"].astype(str).isin(selected_signals)].copy()
    if bucket_options and "maturity_bucket" in ranked.columns:
        ranked = ranked[ranked["maturity_bucket"].astype(str).isin(selected_buckets)].copy()
    ranked = ranked.sort_values(["rv_score", "liquidity_score", "trade_count"], ascending=False)

    display_cols = [
        "cusip", "signal", "maturity_bucket", "current_spread_bps", "peer_median_gap_bps",
        "liquidity_score", "rv_score", "trade_count", "total_trade_amount", "latest_trade",
    ]
    st.subheader("Opportunity Ranking")
    if ranked.empty:
        st.info("No candidates meet the current RV/watchlist filters.")
    else:
        r1, r2, r3, r4 = st.columns(4)
        with r1:
            clean_metric_card("Candidates", f"{len(ranked):,}", size="small")
        with r2:
            clean_metric_card("Top RV", _fmt_num(ranked["rv_score"].max()), size="small")
        with r3:
            clean_metric_card("Top Liquidity", _fmt_num(ranked["liquidity_score"].max()), size="small")
        with r4:
            clean_metric_card("Median Peer Gap", _fmt_bps(ranked["peer_median_gap_bps"].median()), size="small")
        safe_dataframe(ranked[[c for c in display_cols if c in ranked.columns]].head(50), hide_index=True)

    st.subheader("Watchlist")
    _focused_watchlist_records()
    add_options = ranked["cusip"].dropna().astype(str).head(150).tolist() if not ranked.empty else []
    add_col, note_col = st.columns([1.2, 2])
    with add_col:
        selected_add = st.multiselect("Add CUSIPs", add_options, key="focused_rv_add_cusips")
    with note_col:
        bulk_note = st.text_input(
            "Note for selected CUSIPs",
            key="focused_rv_bulk_note",
            placeholder="Why these belong on the shortlist, or what to verify next.",
        )
    add_button_col, clear_button_col = st.columns([1, 1])
    with add_button_col:
        if st.button("Add selected to watchlist", key="focused_rv_add_selected"):
            for item in selected_add:
                row_match = ranked[ranked["cusip"].astype(str) == str(item)]
                row = row_match.iloc[0] if not row_match.empty else {"cusip": item}
                _upsert_focused_watchlist(item, selected_issuer, "RV Ranking", row, bulk_note)
            st.success(f"Saved {len(selected_add):,} selected CUSIP(s).")
    with clear_button_col:
        if st.button("Clear full watchlist", key="focused_rv_clear_watchlist"):
            st.session_state["focused_watchlist_records"] = {}
            st.session_state["focused_watchlist"] = []
            st.info("Watchlist cleared.")

    saved = _focused_watchlist_dataframe(summary)
    if saved.empty:
        st.info("No saved CUSIPs yet.")
    else:
        st.caption(f"{len(saved):,} saved candidate(s). Notes are stored in this Streamlit session.")
        saved_display_cols = [
            "cusip", "issuer", "signal", "maturity_bucket", "current_spread_bps", "peer_median_gap_bps",
            "liquidity_score", "rv_score", "trade_count", "total_trade_amount", "latest_trade",
            "note", "source", "updated_at",
        ]
        safe_dataframe(saved[[c for c in saved_display_cols if c in saved.columns]], hide_index=True, auto_collapse=False)

        edit_col1, edit_col2 = st.columns([1, 2])
        with edit_col1:
            saved_cusips = saved["cusip"].dropna().astype(str).tolist()
            edit_cusip = st.selectbox("Edit saved CUSIP", saved_cusips, key="focused_watch_edit_cusip")
        with edit_col2:
            current_records = _focused_watchlist_records()
            current_note = current_records.get(str(edit_cusip), {}).get("note", "")
            edited_note = st.text_area(
                "Saved note",
                value=current_note,
                key=f"focused_watch_edit_note_{edit_cusip}",
                height=92,
            )
        update_col, remove_col = st.columns([1, 1])
        with update_col:
            if st.button("Update saved note", key="focused_watch_update_note"):
                current_records = _focused_watchlist_records()
                if str(edit_cusip) in current_records:
                    current_records[str(edit_cusip)]["note"] = edited_note
                    current_records[str(edit_cusip)]["updated_at"] = pd.Timestamp.now().strftime("%Y-%m-%d %H:%M")
                    st.success(f"Updated note for {edit_cusip}.")
        with remove_col:
            if st.button("Remove saved CUSIP", key="focused_watch_remove_cusip"):
                current_records = _focused_watchlist_records()
                current_records.pop(str(edit_cusip), None)
                st.session_state["focused_watchlist"] = sorted(current_records.keys())
                st.info(f"Removed {edit_cusip}.")

        export_col1, export_col2 = st.columns([1, 1])
        with export_col1:
            st.download_button(
                "Download Watchlist CSV",
                data=saved.to_csv(index=False).encode("utf-8"),
                file_name=f"{selected_issuer}_watchlist.csv".replace(" ", "_"),
                mime="text/csv",
            )
        with export_col2:
            watch_md = _focused_watchlist_markdown(saved, selected_issuer)
            st.download_button(
                "Download Watchlist Markdown",
                data=watch_md.encode("utf-8"),
                file_name=f"{selected_issuer}_watchlist.md".replace(" ", "_"),
                mime="text/markdown",
            )


def render_focused_workflow(
    workflow_view: str,
    trade_reports: list[dict],
    bond_report: dict | None,
    mmd_report: dict | None,
    market_df: pd.DataFrame,
    bonds_df: pd.DataFrame,
    issuer_master: pd.DataFrame,
    mmd_df: pd.DataFrame,
    trade_payloads: list[tuple[str, bytes]],
    failed_files: list[str],
    duplicates_removed: int,
    uploaded_issuers: list[str],
    selected_issuer: str,
    selected_sector: str,
    issuer_trades: pd.DataFrame,
    issuer_bonds: pd.DataFrame,
    comparison_issuers: list[str],
    benchmark_source_mode: str,
    benchmark_priority: str,
    benchmark_conflict_policy: str,
    use_external_mmd_fallback: bool,
    mmd_file_provided: bool,
):
    render_workflow_header(workflow_view, files_loaded=len(trade_payloads), issuers_loaded=len(uploaded_issuers))
    if workflow_view == "1. Upload / Data Audit":
        render_focused_upload_audit(
            trade_reports, bond_report, mmd_report, market_df, bonds_df, issuer_master, mmd_df,
            trade_payloads, failed_files, duplicates_removed, benchmark_source_mode, benchmark_priority,
            benchmark_conflict_policy, use_external_mmd_fallback, mmd_file_provided,
        )
    elif workflow_view == "2. Desk Snapshot":
        render_focused_snapshot(market_df, bonds_df, issuer_trades, issuer_bonds, mmd_df, selected_issuer, selected_sector, benchmark_source_mode)
    elif workflow_view == "3. Core Charts":
        render_focused_core_charts(market_df, issuer_trades, mmd_df, selected_issuer, comparison_issuers, selected_sector)
    elif workflow_view == "4. CUSIP Drilldown":
        render_focused_cusip_drilldown(issuer_trades, selected_issuer)
    elif workflow_view == "5. RV / Watchlist":
        render_focused_rv_watchlist(issuer_trades, selected_issuer)
    elif workflow_view == "6. Export / Methodology":
        render_focused_export_methodology(
            selected_issuer, selected_sector, market_df, issuer_trades, issuer_bonds, mmd_df,
            benchmark_source_mode, benchmark_priority, benchmark_conflict_policy,
            build_report_context=_build_focused_report_context,
            section_anchor=section_anchor,
            clean_metric_card=clean_metric_card,
            safe_dataframe=safe_dataframe,
            render_benchmark_methodology_block=_render_benchmark_methodology_block,
        )


def section_directory():
    """Compact workflow map for the main page.

    The full jump list lives in the sidebar. This main-page version is intentionally
    concise so it does not crowd the dashboard before users upload data.
    """
    with st.expander("Dashboard workflow map", expanded=False):
        st.markdown(
            """
<div class="nav-card">
<b>How to read this dashboard</b><br><br>

<b>1. Data readiness</b><br>
<a href="#file-readiness">File Readiness</a> ·
<a href="#executive-snapshot">Executive Snapshot</a><br><br>

<b>2. Benchmark & spread framework</b><br>
<a href="#yield-relative-value">Yield / RV Trend</a> ·
<a href="#issuer-curve">Issuer Curve</a> ·
<a href="#spread-level">Spread Level</a> ·
<a href="#spread-attribution">Spread Attribution</a><br><br>

<b>3. Relative value signals</b><br>
<a href="#peer-rv">Peer RV</a> ·
<a href="#cross-issuer-rv">Cross-Issuer RV</a> ·
<a href="#historical-spread">Historical Percentile</a> ·
<a href="#recommendation-engine">Rule-Based Narrative</a> ·
<a href="#ai-commentary-studio">AI Commentary Studio</a><br><br>

<b>4. Risk, flow & opportunity screening</b><br>
<a href="#curve-shape">Curve Shape</a> ·
<a href="#scenario-shock">Scenario Shock</a> ·
<a href="#dealer-proxy">Dealer Proxy</a> ·
<a href="#security-screener">Security Screener</a> ·
<a href="#watchlist">Watchlist</a><br><br>

<b>5. Outputs, methodology & raw detail</b><br>
<a href="#spread-movement">Spread Movement</a> ·
<a href="#cusip-drilldown">CUSIP Drilldown</a> ·
<a href="#report-export-center">Report Export Center</a> ·
<a href="#export-summary">Export Summary</a> ·
<a href="#admin-methodology">Admin Methodology</a> ·
<a href="#version-changelog">Version / Change Log</a> ·
<a href="#downloads">Downloads</a>
</div>
""",
            unsafe_allow_html=True,
        )

# Dashboard workflow map removed: the app now uses a desk-first sidebar index.


with st.expander("Instructions", expanded=False):
    st.markdown(
        """
<div style='font-size:15px; color:black; line-height:1.45;'>

<h4 style='margin-bottom:6px;'>Desk-First Workflow</h4>

This dashboard is designed to behave like a lightweight secondary-market desk tool. Start with the trade tape, then use optional reference files only when they improve the analysis.

<h5 style='margin-bottom:4px;'>1. Required Input: MuniPro Trade History</h5>
<div style='padding-left:18px;'>
<ul style='margin-top:2px; margin-bottom:6px;'>
<li>Upload one or more MuniPro trade-history files.</li>
<li>Name each file after the issuer, for example <code>LADWP_Trade.xlsx</code> or <code>State_of_California_Trade.csv</code>.</li>
<li>The app uses the file name as the issuer name. This avoids confusing bond-purpose text like <code>GO Various Purpose</code>, <code>Power</code>, or <code>Water</code> with the issuer.</li>
</ul>

<b>Minimum fields:</b><br>
CUSIP / CUSIP9, Trade Date, Yield, Maturity Date<br><br>

<b>Recommended fields:</b><br>
Trade Date/Time, Description, Coupon, Price, Trade Amount, Index, Index Rate, Spread, Trade Type, Ratings M/S/F
</div>

<h5 style='margin-top:10px; margin-bottom:4px;'>2. Optional Reference Files</h5>
<div style='padding-left:18px;'>
<ul style='margin-top:2px; margin-bottom:6px;'>
<li><b>Bond Reference:</b> use only for enrichment such as call date, call price, lien, tax status, and outstanding amount.</li>
<li><b>Issuer / Sector Mapping:</b> use when you want persistent sector labels instead of manual overrides.</li>
<li><b>MMD Curve:</b> use as the preferred benchmark only when you upload a clean, small curve file for the research period. Avoid oversized historical files.</li>
</ul>
</div>

<h5 style='margin-top:10px; margin-bottom:4px;'>3. Recommended Reading Order</h5>
<div style='padding-left:18px;'>
<ol style='margin-top:2px; margin-bottom:6px;'>
<li><b>Desk Market Snapshot</b>: spread trend, trading volume, curve snapshot, and top movers.</li>
<li><b>Issuer Curve vs Benchmark</b>: where the issuer curve sits versus the active benchmark.</li>
<li><b>Spread Movement Ladder</b>: use as a drilldown, not the first page.</li>
<li><b>Liquidity / Trading Frequency</b>: confirm whether the signal is supported by enough trading activity.</li>
<li><b>CUSIP Drilldown / Screener</b>: investigate specific bonds after the high-level view.</li>
</ol>
</div>

<h5 style='margin-top:10px; margin-bottom:4px;'>4. Performance Tips</h5>
<div style='padding-left:18px;'>
<ul style='margin-top:2px; margin-bottom:2px;'>
<li>Keep <b>Fast mode</b> on while exploring.</li>
<li>Do not show full raw tables unless auditing data.</li>
<li>Use MMD files covering only the years you are researching.</li>
<li>Use the sidebar index to jump to the most-used desk sections first.</li>
</ul>
</div>

</div>
""",
        unsafe_allow_html=True,
    )

# -----------------------------------------------------------------------------
# Team-readiness validation UI
# -----------------------------------------------------------------------------
# Pure readiness constants and validators live in engine/validation.py.

def display_validation_report(title: str, report: dict, warnings: list[str] | None = None):
    """Render a user-facing readiness card in Streamlit."""
    warnings = warnings or []
    status_icon = "✅" if report["can_run"] else "❌"
    with st.expander(f"{status_icon} {title} readiness check", expanded=not report["can_run"]):
        st.caption(f"Rows: {report['row_count']:,} · Columns: {report['column_count']:,}")

        c1, c2, c3 = st.columns(3)
        c1.metric("Required detected", f"{len(report['detected_required'])}/{len(report['detected_required']) + len(report['missing_required'])}")
        c2.metric("Recommended detected", f"{len(report['detected_recommended'])}/{len(report['detected_recommended']) + len(report['missing_recommended'])}")
        c3.metric("Ready to run", "Yes" if report["can_run"] else "No")

        if report["missing_required"]:
            st.error("Missing required fields: " + ", ".join(report["missing_required"]))
        if report["missing_recommended"]:
            st.warning("Missing recommended fields: " + ", ".join(report["missing_recommended"]))
        if warnings:
            for warning in warnings:
                st.warning(warning)

        mapping_rows = [
            {"Internal Field": key, "Uploaded Column Detected": value or "—"}
            for key, value in report["mapping"].items()
        ]
        safe_dataframe(pd.DataFrame(mapping_rows), width="stretch", hide_index=True)


def template_download_button(columns: list[str], label: str, filename: str):
    template = pd.DataFrame(columns=columns)
    st.download_button(
        label=label,
        data=template.to_csv(index=False).encode("utf-8"),
        file_name=filename,
        mime="text/csv",
    )




# Upload/data processing lives in engine/load_data.py.

def dataframe_download_button(df: pd.DataFrame, label: str, filename: str):
    if df.empty:
        return
    csv = df.to_csv(index=False).encode("utf-8")
    st.download_button(label=label, data=csv, file_name=filename, mime="text/csv")


with st.sidebar:
    st.header("Workflow")
    workflow_view = st.radio(
        "Workspace section",
        WORKFLOW_LABELS + [FULL_DASHBOARD_LABEL],
        index=0,
        help="Use the focused six-step flow for day-to-day work. Full Dashboard keeps the original long-form workstation available.",
    )
    st.markdown("---")
    st.header("Performance")
    PERFORMANCE_MODE = st.checkbox(
        "Fast mode",
        value=True,
        help="Caches heavy calculations, limits displayed rows, and keeps ladders readable.",
    )
    MAX_TABLE_ROWS = st.number_input(
        "Max table rows shown",
        min_value=500,
        max_value=20000,
        value=3000,
        step=500,
        help="Only limits displayed tables; underlying analytics still use the full filtered dataset.",
    )
    MAX_HEATMAP_ROWS = st.slider(
        "Max ladder maturity rows",
        min_value=8,
        max_value=40,
        value=18,
        help="Fast mode keeps the maturity years with the largest absolute signal.",
    )
    SHOW_FULL_RAW_TABLES = st.checkbox(
        "Show full raw tables",
        value=False,
        help="Usually keep this off. Full raw tables are one of the biggest Streamlit slowdowns.",
    )
    ENABLE_REPORT_EXPORTS = st.checkbox(
        "Enable report export builder",
        value=False,
        help="Keep off while exploring. Report export recomputes multiple charts and can slow the app.",
    )
    if st.button("Clear cached calculations"):
        st.cache_data.clear()
        st.rerun()

with st.expander(
    "Upload Center",
    expanded=(workflow_view == "1. Upload / Data Audit"),
):
    st.markdown(
        "<div class='focus-band'>Upload one or more MuniPro trade files here first. Optional reference files can enrich the analysis, but the trade file is the only required input.</div>",
        unsafe_allow_html=True,
    )
    upload_col1, upload_col2 = st.columns([1.15, 0.85])
    with upload_col1:
        trade_files = st.file_uploader(
            "Trade History File(s) — required",
            type=["csv", "xlsx", "xls"],
            accept_multiple_files=True,
            help="Required. Name each trade file after its issuer, e.g. State_of_California_Trade.csv or LADWP_Trade.xlsx.",
        )
        st.caption("Name each trade file after its issuer. The app uses the filename as the issuer name.")
        st.caption("Keep proprietary raw exports out of public GitHub. Upload them only during your own session.")

    with upload_col2:
        bond_file = st.file_uploader("Bond Reference File — optional enrichment", type=["csv", "xlsx", "xls"])
        issuer_mapping_file = st.file_uploader("Issuer / Sector Mapping — optional", type=["csv", "xlsx", "xls"])
        use_external_mmd_fallback = st.checkbox(
            "Enable External MMD Fallback",
            value=False,
            help=(
                "Off by default to prevent memory overload. The app uses Trade Sheet Index / Index Rate first. "
                "Only enable this if your trade files do not have usable Index Rate data."
            ),
        )
        mmd_file = st.file_uploader(
            "MMD Curve File — optional fallback",
            type=["csv", "xlsx", "xls"],
            disabled=not use_external_mmd_fallback,
            help="Loaded only when External MMD Fallback is enabled and trade-sheet Index Rate is unavailable.",
        )
        if not use_external_mmd_fallback:
            st.caption("External MMD loading is off. This avoids benchmark-source conflict and protects memory.")

    render_upload_file_cards(
        trade_file_names=[f.name for f in trade_files] if trade_files else [],
        bond_file_name=bond_file.name if bond_file else None,
        issuer_mapping_file_name=issuer_mapping_file.name if issuer_mapping_file else None,
        mmd_file_name=mmd_file.name if mmd_file else None,
        use_external_mmd_fallback=use_external_mmd_fallback,
    )

    with st.expander("Download blank templates", expanded=False):
        template_download_button(TRADE_REQUIRED + TRADE_RECOMMENDED + TRADE_OPTIONAL, "Trade template CSV", "trade_history_template.csv")
        template_download_button(BOND_REQUIRED + BOND_RECOMMENDED + BOND_OPTIONAL, "Optional bond reference template CSV", "bond_reference_template.csv")
        template_download_button(CURVE_TEMPLATE_COLUMNS, "Fallback MMD curve template CSV", "benchmark_curve_template.csv")

if not trade_files:
    render_workflow_header(workflow_view, files_loaded=0, issuers_loaded=0)
    st.info("Upload at least one MuniPro trade-history file to generate the dashboard. Bond reference data is optional enrichment.")
    with st.expander("Expected file logic"):
        st.write(
            "The app now uses a trade-first workflow: it standardizes CUSIP fields, uses each trade file name as the issuer name, "
            "builds maturity-year fields from trade maturity dates, and optionally enriches static fields from a bond reference file when provided."
        )
    st.stop()

bond_payload = (bond_file.name, bond_file.getvalue()) if bond_file else None
trade_payloads = [(f.name, f.getvalue()) for f in trade_files]
issuer_mapping_payload = (issuer_mapping_file.name, issuer_mapping_file.getvalue()) if issuer_mapping_file else None
mmd_payload = (mmd_file.name, mmd_file.getvalue()) if (use_external_mmd_fallback and mmd_file) else None
show_file_audit = workflow_view in {"1. Upload / Data Audit", FULL_DASHBOARD_LABEL}
show_methodology_audit = workflow_view in {"1. Upload / Data Audit", "6. Export / Methodology", FULL_DASHBOARD_LABEL}

# -----------------------------------------------------------------------------
# File-readiness gate: inspect the uploaded files before running full analytics.
# -----------------------------------------------------------------------------
if show_file_audit:
    section_anchor("file-readiness", "File Readiness Check")
if bond_payload is not None:
    raw_bonds_preview = read_uploaded_file(io.BytesIO(bond_payload[1]), bond_payload[0])
    bond_report = validate_dataset(raw_bonds_preview, bond_payload[0], ["cusip"], BOND_RECOMMENDED, BOND_OPTIONAL)
    bond_warnings = validate_basic_values(raw_bonds_preview, bond_report["mapping"], dataset_type="bond")
    if show_file_audit:
        display_validation_report("Optional Bond Reference File", bond_report, bond_warnings)
else:
    bond_report = {"can_run": True}
    if show_file_audit:
        st.info("No bond reference file uploaded. Running in trade-only mode; static bond metadata will be inferred from the trade tape where possible.")

trade_reports = []
trade_blocking_failures = []
for trade_name, trade_bytes in trade_payloads:
    try:
        raw_trade_preview = read_uploaded_file(io.BytesIO(trade_bytes), trade_name)
        report = validate_dataset(raw_trade_preview, trade_name, TRADE_REQUIRED, TRADE_RECOMMENDED, TRADE_OPTIONAL)
        warnings = validate_basic_values(raw_trade_preview, report["mapping"], dataset_type="trade")
        trade_reports.append(report)
        if show_file_audit:
            display_validation_report(f"Trade File — {trade_name}", report, warnings)
        if not report["can_run"]:
            trade_blocking_failures.append(trade_name)
    except Exception as exc:
        st.error(f"Could not read trade file {trade_name}: {exc}")
        trade_blocking_failures.append(trade_name)

if trade_blocking_failures:
    st.error(
        "The dashboard cannot run yet because at least one required trade file is missing minimum fields. "
        "Use the readiness tables above to rename/add columns, then upload again."
    )
    st.stop()

if mmd_payload is not None:
    try:
        mmd_name, mmd_bytes = mmd_payload
        raw_mmd_preview = read_uploaded_file(io.BytesIO(mmd_bytes), mmd_name)
        mmd_report = validate_dataset(raw_mmd_preview, mmd_name, MMD_REQUIRED, MMD_RECOMMENDED, [])
        if show_file_audit:
            display_validation_report("MMD Curve File", mmd_report)
        if show_file_audit and not mmd_report["can_run"]:
            st.warning("MMD comparison will be skipped unless the MMD file has a date column.")
    except Exception as exc:
        mmd_report = None
        if show_file_audit:
            st.warning(f"Could not validate MMD file. MMD comparison may be skipped: {exc}")
else:
    mmd_report = None

if show_file_audit:
    with st.expander("Methodology: how the app decides whether a file is usable", expanded=False):
        st.markdown(
            """
- **Required fields** are the minimum fields needed for the dashboard to run.
- **Recommended fields** improve liquidity, benchmark, tax, and relative-value analysis, but missing them should not break the app.
- **Column aliases** let the app recognize variants like `CUSIP9`, `Cusip`, or `CUSIP` as the same internal `cusip` field.
- **Warnings** flag data-quality issues, but the app only stops when a required trade field is missing.
        """
        )

(
    bonds_df,
    trades_df,
    issuer_master,
    market_df,
    mmd_df,
    failed_files,
    duplicates_removed,
) = process_uploads(
    trade_payloads=trade_payloads,
    issuer_mapping_payload=issuer_mapping_payload,
    mmd_payload=mmd_payload,
    bond_payload=bond_payload,
)

# Normalize legacy tenor-style bucket labels into presentation-friendly curve sectors.
# This keeps the dashboard compatible with existing data_utils output while making
# the user-facing terminology clearer.
for _df in [bonds_df, trades_df, market_df]:
    if isinstance(_df, pd.DataFrame) and "maturity_bucket" in _df.columns:
        _df["maturity_bucket"] = _df["maturity_bucket"].replace(MATURITY_BUCKET_RENAME)
    if isinstance(_df, pd.DataFrame) and "maturity_year" not in _df.columns and "years_to_maturity" in _df.columns:
        _y = pd.to_numeric(_df["years_to_maturity"], errors="coerce")
        _df["maturity_year"] = pd.Series(np.ceil(_y), index=_df.index).where(_y.notna()).astype("Int64")

if failed_files:
    with st.warning("Some trade files failed to process."):
        st.write(failed_files)

if market_df.empty:
    st.error("No usable trade rows found. Please check that trade files include CUSIP and trade date fields.")
    st.stop()

uploaded_issuers = sorted(market_df["issuer"].dropna().astype(str).unique().tolist())

if not uploaded_issuers:
    st.error("No issuer names were detected from the uploaded trade files. Please check Description, issuer mapping, or trade filenames.")
    st.stop()

if show_file_audit:
    st.success(
        f"Processed {len(market_df):,} trade rows and built {len(bonds_df):,} security-reference rows "
        f"from {len(trade_files):,} trade file(s). Detected {len(uploaded_issuers):,} issuer(s)."
    )

benchmark_source_mode = mmd_df.attrs.get("benchmark_source_mode", "None")
benchmark_priority = mmd_df.attrs.get("benchmark_source_priority", "None")
benchmark_conflict_policy = mmd_df.attrs.get("benchmark_conflict_policy", "No benchmark source selected")
uploaded_mmd_available = bool(mmd_df.attrs.get("uploaded_mmd_available", False))

if show_methodology_audit:
    if benchmark_source_mode == "Trade Sheet Index / Index Rate":
        st.info(
            "Benchmark source: using **Index / Index Rate from the uploaded trade sheet** as the primary benchmark universe. "
            "Any uploaded MMD file is treated as fallback only and is not mixed into the same analytics run."
        )
    elif benchmark_source_mode == "Uploaded MMD fallback":
        st.info(
            "Benchmark source: using the **uploaded MMD file as fallback** because the trade sheet did not contain usable Index / Index Rate data."
        )
    else:
        st.warning("No benchmark source detected. Upload trades with Index / Index Rate or provide an MMD file for benchmark analytics.")

    with st.expander("Benchmark source governance", expanded=False):
        st.markdown(
            """
This dashboard uses **one benchmark source at a time** to avoid benchmark-source conflict.

**Priority hierarchy**

1. **Trade Sheet Index / Index Rate — recommended primary source.**  
   This is preferred because it comes from the same uploaded trade tape and pricing context as the observed trades.
2. **Uploaded MMD file — fallback only.**  
   This is used only when the trade sheet does not include usable `Index` / `Index Rate` fields.
3. **No benchmark source.**  
   Yield-only and liquidity analytics can still run, but benchmark spread analytics are skipped or downgraded.

**Why not mix both?**

Trade-sheet index rates and an external MMD sheet may differ by date, tenor, rounding, provider convention, or interpolation method. Mixing them can shift spreads by several basis points and make relative-value signals inconsistent.
        """
        )
        safe_dataframe(
            pd.DataFrame([
                {"Item": "Active benchmark source", "Value": benchmark_source_mode},
                {"Item": "Priority", "Value": benchmark_priority},
                {"Item": "Conflict policy", "Value": benchmark_conflict_policy},
                {"Item": "Uploaded MMD detected", "Value": "Yes" if uploaded_mmd_available else "No / Not used"},
            ]),
            width="stretch",
            hide_index=True,
        )

with st.sidebar:
    st.markdown("---")
    st.header("2. Select From Uploaded Issuers")
    selected_issuer = st.selectbox(
        "Primary Issuer",
        uploaded_issuers,
        help="Main issuer shown first in desk snapshot and drilldown sections. Rename each trade file with the issuer name for best results."
    )

    # -----------------------------------------------------------------------------
    # Desk comparison controls
    # -----------------------------------------------------------------------------
    peer_options_sidebar = [x for x in uploaded_issuers if x != selected_issuer]
    default_peers_sidebar = peer_options_sidebar[:2] if peer_options_sidebar else []
    comparison_issuers_sidebar = st.multiselect(
        "Compare With Issuers",
        peer_options_sidebar,
        default=default_peers_sidebar,
        help="Optional peer issuers to plot beside the primary issuer in spread and volume charts.",
    )

    snapshot_reference_lines = st.multiselect(
        "Reference Lines",
        ["AAA / MMD Baseline", "Sector Average", "All Uploaded Issuers Average"],
        default=["Sector Average"] if peer_options_sidebar else [],
        help="Adds desk-style reference lines to the spread chart. AAA / MMD baseline is 0 bps when spread is measured vs AAA/MMD.",
    )

    volume_comparison_mode = st.radio(
        "Volume Chart Grouping",
        ["Primary vs Peers vs All Other", "Selected Issuers Only"],
        index=0,
        help="Controls whether the volume chart stacks selected issuers against the rest of the uploaded market universe.",
    )

    # -----------------------------------------------------------------------------
    # Manual issuer sector override
    # -----------------------------------------------------------------------------
    sector_options = [
        "Unknown", "General Government", "State GO", "Local Government",
        "Utilities", "Water / Sewer", "Power", "Transportation", "Airport",
        "Education", "School District", "Healthcare", "Housing",
        "Public Finance Authority", "Other",
    ]

    current_sector_sidebar = "Unknown"
    if "sector" in market_df.columns:
        vals = market_df.loc[market_df["issuer"] == selected_issuer, "sector"].dropna().astype(str).unique().tolist()
        vals = [v for v in vals if v and v.lower() != "nan"]
        if vals:
            current_sector_sidebar = vals[0]

    with st.expander("Issuer Sector Override", expanded=(current_sector_sidebar == "Unknown")):
        st.caption("Use this when the trade file has no sector field, or when the inferred sector is Unknown / wrong.")
        default_idx = sector_options.index(current_sector_sidebar) if current_sector_sidebar in sector_options else 0
        selected_sector_input = st.selectbox("Sector", sector_options, index=default_idx, key=f"sector_select_{selected_issuer}")
        custom_sector_input = st.text_input(
            "Custom sector",
            value="" if selected_sector_input != "Other" else current_sector_sidebar,
            key=f"sector_custom_{selected_issuer}",
        )
        final_sector_input = custom_sector_input.strip() if selected_sector_input == "Other" and custom_sector_input.strip() else selected_sector_input
        if st.button("Apply Sector to Current Issuer", key=f"apply_sector_{selected_issuer}"):
            st.session_state.setdefault("issuer_sector_overrides", {})[selected_issuer] = final_sector_input
            st.success(f"Applied: {selected_issuer} → {final_sector_input}")

    # -----------------------------------------------------------------------------
    # Maturity Year Selector
    # -----------------------------------------------------------------------------
    with st.expander("Maturity Year Methodology", expanded=False):
        st.markdown(
            """
### Maturity Year Definition

The dashboard now groups securities by **integer years to maturity** instead of broad ranges.

**Formula:**

`years_to_maturity = (maturity_date - trade_date) / 365.25`

**Bucket rule:**

- We use `ceil(years_to_maturity)` for the displayed maturity year.
- Example: 4.3 years to maturity → **5Y**.
- This is closer to curve-tenor convention than broad buckets like 1Y / 30Y.

This makes issuer-level analysis easier because you can compare 1Y, 2Y, 3Y, ... securities directly.
            """
        )

    issuer_year_values = []
    if "maturity_year" in market_df.columns:
        issuer_year_values = (
            market_df.loc[market_df["issuer"] == selected_issuer, "maturity_year"]
            .dropna()
            .astype(int)
            .sort_values()
            .unique()
            .tolist()
        )
    maturity_year_options = ["All"] + [f"{int(y)}Y" for y in issuer_year_values if int(y) >= 1]
    selected_maturity_year = st.selectbox(
        "Maturity Year",
        maturity_year_options,
        help="Filter securities by integer years to maturity. Example: 4.3 years to maturity is grouped as 5Y.",
    )

    # -----------------------------------------------------------------------------
    # Snapshot / Chart Period Selector
    # -----------------------------------------------------------------------------
    selected_trade_date_range = None
    trade_date_filter_enabled = False
    snapshot_period = st.selectbox(
        "Snapshot / Chart Period",
        ["All", "Last 3M", "Last 6M", "Last 1Y", "YTD", "Custom"],
        index=3,
        help="Choose Custom to freely select exact start/end trade dates. This is a data filter, unlike Plotly zoom which only changes the visual view.",
    )
    st.caption("Tip: use Custom here for a true date filter; chart zoom only changes the view and does not filter the dataset.")
    if "trade_date" in market_df.columns:
        _trade_dates = pd.to_datetime(market_df["trade_date"], errors="coerce").dropna()
        if not _trade_dates.empty:
            _min_ts = _trade_dates.min().normalize()
            _max_ts = _trade_dates.max().normalize()
            _min_date = _min_ts.date()
            _max_date = _max_ts.date()
            if snapshot_period == "All":
                trade_date_filter_enabled = False
                selected_trade_date_range = None
            elif snapshot_period == "Custom":
                trade_date_filter_enabled = True
                selected_trade_date_range = st.date_input(
                    "Custom Trade Date Range",
                    value=(_min_date, _max_date),
                    min_value=_min_date,
                    max_value=_max_date,
                    help="Choose the exact trade-date period to show in the spread and volume snapshot charts.",
                )
            else:
                trade_date_filter_enabled = True
                if snapshot_period == "Last 3M":
                    _start_ts = max(_min_ts, _max_ts - pd.DateOffset(months=3))
                elif snapshot_period == "Last 6M":
                    _start_ts = max(_min_ts, _max_ts - pd.DateOffset(months=6))
                elif snapshot_period == "Last 1Y":
                    _start_ts = max(_min_ts, _max_ts - pd.DateOffset(years=1))
                elif snapshot_period == "YTD":
                    _start_ts = max(_min_ts, pd.Timestamp(year=_max_ts.year, month=1, day=1))
                else:
                    _start_ts = _min_ts
                selected_trade_date_range = (_start_ts.date(), _max_date)
                st.caption(f"Active period: {_start_ts:%m/%d/%Y} → {_max_ts:%m/%d/%Y}")
    else:
        st.caption("Snapshot period unavailable until trade dates are loaded.")
    # Keep legacy variables available for older downstream chart blocks.
    maturity_bucket = selected_maturity_year
    trade_date_range = selected_trade_date_range

    # -----------------------------------------------------------------------------
    # Raw Table Toggle
    # -----------------------------------------------------------------------------
    show_raw_tables = st.checkbox(
        "Show Raw Tables",
        value=False,
        help="""
Display underlying trade-level and security-reference data tables.

Useful for:
- Audit review
- Data validation
- Trade-level investigation
- CUSIP drilldowns
"""
    )

    st.markdown("---")
    st.subheader("Index / Benchmark Source")
    st.caption(f"Active: {benchmark_source_mode}")
    st.caption(f"Policy: {benchmark_conflict_policy}")

    st.markdown("---")
    st.subheader("Desk Navigation")
    st.caption("Desk-first order: the most-used secondary-market views are listed first.")
    st.markdown(
        """
<div class="sidebar-nav-small">
<b>Primary Desk Views</b><br>
<a href="#yield-relative-value">1. Secondary Market Spreads</a><br>
<a href="#trading-volume">2. Secondary Market Trading Volume</a><br>
<a href="#issuer-curve">3. Issuer Curve vs Benchmark</a><br>
<a href="#spread-movement">4. Spread Movement Ladder</a><br>
<a href="#liquidity">5. Liquidity / Trading Frequency</a><br>
<a href="#cusip-drilldown">6. CUSIP Drilldown</a><br>
<a href="#security-screener">7. Security Screener</a><br><br>

<b>Relative Value Research</b><br>
<a href="#peer-rv">8. Peer RV Comparison</a><br>
<a href="#cross-issuer-rv">9. Cross-Issuer RV Analytics</a><br>
<a href="#spread-level">10. Current Spread Curve</a><br>
<a href="#spread-attribution">11. Spread Attribution</a><br>
<a href="#historical-spread">12. Historical Spread Percentile</a><br>
<a href="#curve-shape">13. Curve Shape Analytics</a><br><br>

<b>Advanced / Commentary</b><br>
<a href="#market-narrative">14. Market Narrative & Opportunity Map</a><br>
<a href="#dealer-proxy">15. Dealer Behavior Proxy</a><br>
<a href="#rv-positioning">16. RV Positioning Map</a><br>
<a href="#scenario-shock">17. Scenario Shock Analysis</a><br>
<a href="#watchlist">18. Watchlist / Saved Candidates</a><br>
<a href="#recommendation-engine">19. Recommendation Narrative</a><br>
<a href="#ai-commentary-studio">20. AI Commentary Studio</a><br><br>

<b>Data / Admin</b><br>
<a href="#file-readiness">File Readiness</a><br>
<a href="#executive-snapshot">Executive Snapshot</a><br>
<a href="#bond-master">Security Reference</a><br>
<a href="#trade-detail">Trade Detail</a><br>
<a href="#downloads">Downloads</a>
</div>
""",
        unsafe_allow_html=True,
    )

    with st.expander("Version / Change Log", expanded=False):
        st.markdown(
            """
**Current Version:** `v1.2-desk-order`

Recent additions:
- Cross-Issuer RV Analytics
- Scenario Shock Analysis
- Recommendation Narrative Engine
- Desk-first navigation and market snapshot
- Reordered primary desk views and sidebar index
- Faster exploration defaults
- Optional advanced/admin sections
            """
        )
    st.markdown("---")
    st.header("Data Health")

    if not market_df.empty and "trade_date" in market_df.columns:
        trade_dates = pd.to_datetime(market_df["trade_date"], errors="coerce").dropna()
        if not trade_dates.empty:
            earliest_trade = trade_dates.min()
            latest_trade = trade_dates.max()
            st.caption(
                f"📅 Data Coverage:\n"
                f"{earliest_trade:%m/%d/%Y} → {latest_trade:%m/%d/%Y}"
            )
        else:
            st.caption("📅 Data Coverage:\nNo valid trade dates detected")
    else:
        st.caption("📅 Data Coverage:\nNo trade data loaded")

    st.caption(
        f"📊 Trades Loaded:\n"
        f"{len(market_df):,}"
    )

    total_rows = len(market_df)
    if total_rows > 0 and "cusip" in market_df.columns:
        valid_cusip_count = market_df["cusip"].notna().sum()
        valid_cusip_rate = valid_cusip_count / total_rows * 100
    else:
        valid_cusip_rate = 0

    cusip_icon = "🟢" if valid_cusip_rate >= 95 else "🟡" if valid_cusip_rate >= 80 else "🔴"
    st.caption(
        f"{cusip_icon} Valid CUSIP Rate:\n"
        f"{valid_cusip_rate:.1f}%"
    )

    missing_issuers = market_df["issuer"].isna().sum() if "issuer" in market_df.columns else total_rows
    missing_issuer_rate = missing_issuers / total_rows * 100 if total_rows > 0 else 0
    missing_icon = "🟢" if missing_issuers == 0 else "🟡" if missing_issuer_rate <= 5 else "🔴"
    st.caption(
        f"{missing_icon} Missing Issuers:\n"
        f"{missing_issuers:,}"
    )

    st.caption(
        f"🧹 Duplicate Trades Removed:\n"
        f"{duplicates_removed:,}"
    )

    with st.expander("Data Health methodology", expanded=False):
        st.markdown(
            """
- **Data Coverage** uses the earliest and latest valid trade dates after standardization.
- **Trades Loaded** counts trade rows available for analytics.
- **Valid CUSIP Rate** is the share of trade rows with a usable CUSIP identifier.
- **Missing Issuers** counts rows without an issuer after trade-description inference and issuer-mapping logic.
- **Duplicate Trades Removed** counts exact duplicate standardized trade rows removed before analytics.
            """
        )


# Apply manual sector overrides selected in the sidebar.
issuer_sector_overrides = st.session_state.get("issuer_sector_overrides", {})
if issuer_sector_overrides and "issuer" in market_df.columns:
    for _issuer_name, _sector_value in issuer_sector_overrides.items():
        market_df.loc[market_df["issuer"] == _issuer_name, "sector"] = _sector_value
        if "issuer" in issuer_master.columns:
            issuer_master.loc[issuer_master["issuer"] == _issuer_name, "sector"] = _sector_value

issuer_bonds = bonds_df[bonds_df["issuer"] == selected_issuer].copy()
issuer_trades = market_df[market_df["issuer"] == selected_issuer].copy()

# -----------------------------------------------------------------------------
# Trade-only compatibility guard
# -----------------------------------------------------------------------------
# In trade-only mode, several downstream analytics sections may reference
# optional bond/security enrichment fields. Those fields should enhance the
# analysis when available, but they must never break the dashboard when absent.
OPTIONAL_SECURITY_FIELDS = {
    "maturity_bond": pd.NaT,
    "coupon_bond": pd.NA,
    "outstanding_amount": pd.NA,
    "call_date": pd.NaT,
    "call_price": pd.NA,
    "lien": pd.NA,
    "fed_tax": pd.NA,
    "amt": pd.NA,
    "secondary_credit": pd.NA,
    "description": pd.NA,
    "price": pd.NA,
    "trade_amount": 0,
}

for _df in [market_df, issuer_trades]:
    if isinstance(_df, pd.DataFrame):
        for _col, _default in OPTIONAL_SECURITY_FIELDS.items():
            if _col not in _df.columns:
                _df[_col] = _default


if issuer_sector_overrides:
    sector_download_df = pd.DataFrame(
        [{"issuer": k, "sector": v, "primary_type": pd.NA} for k, v in issuer_sector_overrides.items()]
    )
    with st.sidebar.expander("Download sector overrides", expanded=False):
        st.download_button(
            "Download issuer_sector_overrides.csv",
            data=sector_download_df.to_csv(index=False).encode("utf-8"),
            file_name="issuer_sector_overrides.csv",
            mime="text/csv",
        )

selected_sector = "Unknown"
if "sector" in market_df.columns:
    sector_values = issuer_trades["sector"].dropna().astype(str).unique().tolist()
    if sector_values:
        selected_sector = sector_values[0]
elif "sector" in issuer_master.columns:
    sector_values = issuer_master.loc[issuer_master["issuer"] == selected_issuer, "sector"].dropna().astype(str).unique().tolist()
    if sector_values:
        selected_sector = sector_values[0]

if not issuer_trades.empty and selected_maturity_year != "All" and "maturity_year" in issuer_trades.columns:
    _selected_year = int(str(selected_maturity_year).replace("Y", ""))
    issuer_trades = issuer_trades[issuer_trades["maturity_year"] == _selected_year].copy()

if not issuer_trades.empty and trade_date_filter_enabled and selected_trade_date_range:
    if isinstance(selected_trade_date_range, (tuple, list)) and len(selected_trade_date_range) == 2:
        _start_date, _end_date = selected_trade_date_range
        issuer_trades = issuer_trades[
            (pd.to_datetime(issuer_trades["trade_date"], errors="coerce").dt.date >= _start_date)
            & (pd.to_datetime(issuer_trades["trade_date"], errors="coerce").dt.date <= _end_date)
        ].copy()

if workflow_view != FULL_DASHBOARD_LABEL:
    render_focused_workflow(
        workflow_view=workflow_view,
        trade_reports=trade_reports,
        bond_report=bond_report,
        mmd_report=mmd_report,
        market_df=market_df,
        bonds_df=bonds_df,
        issuer_master=issuer_master,
        mmd_df=mmd_df,
        trade_payloads=trade_payloads,
        failed_files=failed_files,
        duplicates_removed=duplicates_removed,
        uploaded_issuers=uploaded_issuers,
        selected_issuer=selected_issuer,
        selected_sector=selected_sector,
        issuer_trades=issuer_trades,
        issuer_bonds=issuer_bonds,
        comparison_issuers=comparison_issuers_sidebar,
        benchmark_source_mode=benchmark_source_mode,
        benchmark_priority=benchmark_priority,
        benchmark_conflict_policy=benchmark_conflict_policy,
        use_external_mmd_fallback=use_external_mmd_fallback,
        mmd_file_provided=mmd_payload is not None,
    )
    st.stop()


# Data Quality Scorecard removed for trade-only workflow.
# The dashboard now relies on the File Readiness Check and Data Health sidebar metrics.


# -----------------------------------------------------------------------------
# Desk-first market snapshot
# -----------------------------------------------------------------------------
section_anchor("desk-market-snapshot", "Desk Market Snapshot")
st.caption(
    "Desk-style opening view: multi-issuer spread trend, stacked trading volume, curve snapshot, and top movers. "
    "Use the sidebar to add peer issuers and reference lines."
)

snapshot_issuers = [selected_issuer] + [x for x in comparison_issuers_sidebar if x != selected_issuer]
snapshot_issuers = [x for x in snapshot_issuers if x in uploaded_issuers]
if not snapshot_issuers:
    snapshot_issuers = [selected_issuer]

snapshot_base = market_df[market_df["issuer"].isin(snapshot_issuers)].copy()
if trade_date_filter_enabled and isinstance(trade_date_range, tuple) and len(trade_date_range) == 2:
    snapshot_base = snapshot_base[
        (pd.to_datetime(snapshot_base["trade_date"], errors="coerce").dt.date >= trade_date_range[0])
        & (pd.to_datetime(snapshot_base["trade_date"], errors="coerce").dt.date <= trade_date_range[1])
    ].copy()

# Stack the two desk snapshot charts vertically so each chart has full page width.
with st.container():
    st.subheader("Secondary Market Spreads")
    spread_universe = market_df.copy()
    if trade_date_filter_enabled and isinstance(trade_date_range, tuple) and len(trade_date_range) == 2:
        spread_universe = spread_universe[
            (pd.to_datetime(spread_universe["trade_date"], errors="coerce").dt.date >= trade_date_range[0])
            & (pd.to_datetime(spread_universe["trade_date"], errors="coerce").dt.date <= trade_date_range[1])
        ].copy()

    if not spread_universe.empty and {"trade_date", "yield", "issuer"}.issubset(spread_universe.columns):
        spread_universe["trade_date"] = pd.to_datetime(spread_universe["trade_date"], errors="coerce")
        spread_universe["yield"] = pd.to_numeric(spread_universe["yield"], errors="coerce")
        if "spread" in spread_universe.columns and spread_universe["spread"].notna().any():
            spread_universe["spread_bps"] = pd.to_numeric(spread_universe["spread"], errors="coerce") * 100
        elif "index_rate" in spread_universe.columns:
            spread_universe["spread_bps"] = (
                pd.to_numeric(spread_universe["yield"], errors="coerce")
                - pd.to_numeric(spread_universe["index_rate"], errors="coerce")
            ) * 100
        else:
            spread_universe["spread_bps"] = pd.NA

        spread_universe = spread_universe.dropna(subset=["trade_date", "spread_bps"])
        fig_spread_snapshot = go.Figure()

        selected_spread_df = spread_universe[spread_universe["issuer"].isin(snapshot_issuers)].copy()
        if not selected_spread_df.empty:
            issuer_daily = (
                selected_spread_df.groupby([pd.Grouper(key="trade_date", freq="D"), "issuer"], as_index=False)
                .agg(spread_bps=("spread_bps", "median"), trade_count=("spread_bps", "count"))
                .dropna(subset=["spread_bps"])
                .sort_values("trade_date")
            )
            for issuer_name in snapshot_issuers:
                tmp = issuer_daily[issuer_daily["issuer"] == issuer_name]
                if tmp.empty:
                    continue
                line_width = 3.2 if issuer_name == selected_issuer else 2.2
                fig_spread_snapshot.add_trace(
                    go.Scatter(
                        x=tmp["trade_date"],
                        y=tmp["spread_bps"],
                        mode="lines",
                        name=issuer_name,
                        line=dict(width=line_width),
                        hovertemplate="%{x|%m/%d/%Y}<br>%{y:.1f} bps<extra>%{fullData.name}</extra>",
                    )
                )

        if "Sector Average" in snapshot_reference_lines and "sector" in spread_universe.columns:
            selected_sector_for_ref = selected_sector if selected_sector and selected_sector != "Unknown" else None
            if selected_sector_for_ref:
                sector_df = spread_universe[spread_universe["sector"].astype(str) == str(selected_sector_for_ref)].copy()
                if not sector_df.empty:
                    sector_daily = (
                        sector_df.groupby(pd.Grouper(key="trade_date", freq="D"), as_index=False)
                        .agg(spread_bps=("spread_bps", "median"))
                        .dropna(subset=["spread_bps"])
                        .sort_values("trade_date")
                    )
                    if not sector_daily.empty:
                        fig_spread_snapshot.add_trace(
                            go.Scatter(
                                x=sector_daily["trade_date"],
                                y=sector_daily["spread_bps"],
                                mode="lines",
                                name=f"{selected_sector_for_ref} Avg",
                                line=dict(width=2, dash="dash"),
                                hovertemplate="%{x|%m/%d/%Y}<br>%{y:.1f} bps<extra>%{fullData.name}</extra>",
                            )
                        )

        if "All Uploaded Issuers Average" in snapshot_reference_lines:
            all_daily = (
                spread_universe.groupby(pd.Grouper(key="trade_date", freq="D"), as_index=False)
                .agg(spread_bps=("spread_bps", "median"))
                .dropna(subset=["spread_bps"])
                .sort_values("trade_date")
            )
            if not all_daily.empty:
                fig_spread_snapshot.add_trace(
                    go.Scatter(
                        x=all_daily["trade_date"],
                        y=all_daily["spread_bps"],
                        mode="lines",
                        name="All Uploaded Issuers Avg",
                        line=dict(width=2, dash="dot"),
                        hovertemplate="%{x|%m/%d/%Y}<br>%{y:.1f} bps<extra>%{fullData.name}</extra>",
                    )
                )

        if "AAA / MMD Baseline" in snapshot_reference_lines and not spread_universe.empty:
            min_dt = spread_universe["trade_date"].min()
            max_dt = spread_universe["trade_date"].max()
            if pd.notna(min_dt) and pd.notna(max_dt):
                fig_spread_snapshot.add_trace(
                    go.Scatter(
                        x=[min_dt, max_dt],
                        y=[0, 0],
                        mode="lines",
                        name="AAA / MMD Baseline",
                        line=dict(width=1.8, dash="longdash"),
                        hovertemplate="%{x|%m/%d/%Y}<br>0.0 bps<extra>AAA / MMD Baseline</extra>",
                    )
                )

        if fig_spread_snapshot.data:
            fig_spread_snapshot.update_layout(
                title="Multi-Issuer Spread Trend",
                xaxis_title="Trade Date",
                yaxis_title="Spread (bps)",
                height=560,
                margin=dict(l=40, r=40, t=70, b=50),
                legend_title_text="Line Item",
            )
            fig_spread_snapshot.update_xaxes(tickformat="%m/%d/%Y")
            safe_plotly_chart(fig_spread_snapshot, width="stretch")
            render_spread_trend_readthrough(spread_universe, selected_issuer, snapshot_issuers)
        else:
            st.info("No usable spread or index-rate data for the selected spread lines.")
    else:
        st.info("Upload trades with issuer, trade date, yield, and spread/index-rate fields to build spread trends.")

with st.container():
    st.markdown("<a id='trading-volume'></a>", unsafe_allow_html=True)
    st.subheader("Secondary Market Trading Volume")
    if not market_df.empty and {"trade_date", "trade_amount", "issuer"}.issubset(market_df.columns):
        vol_universe = market_df.copy()
        vol_universe["trade_date"] = pd.to_datetime(vol_universe["trade_date"], errors="coerce")
        vol_universe["trade_amount"] = pd.to_numeric(vol_universe["trade_amount"], errors="coerce")
        vol_universe = vol_universe.dropna(subset=["trade_date", "trade_amount", "issuer"])
        if trade_date_filter_enabled and isinstance(trade_date_range, tuple) and len(trade_date_range) == 2:
            vol_universe = vol_universe[
                (vol_universe["trade_date"].dt.date >= trade_date_range[0])
                & (vol_universe["trade_date"].dt.date <= trade_date_range[1])
            ].copy()

        if not vol_universe.empty:
            vol_universe["month"] = vol_universe["trade_date"].dt.to_period("M").dt.to_timestamp()
            if volume_comparison_mode == "Primary vs Peers vs All Other":
                def _volume_group(issuer_name: object) -> str:
                    issuer_name = str(issuer_name)
                    if issuer_name == selected_issuer:
                        return selected_issuer
                    if issuer_name in comparison_issuers_sidebar:
                        return issuer_name
                    return "All Other Uploaded Munis"
                vol_universe["volume_group"] = vol_universe["issuer"].apply(_volume_group)
            else:
                vol_universe = vol_universe[vol_universe["issuer"].isin(snapshot_issuers)].copy()
                vol_universe["volume_group"] = vol_universe["issuer"].astype(str)

            monthly_grouped = (
                vol_universe.groupby(["month", "volume_group"], as_index=False)
                .agg(monthly_volume=("trade_amount", "sum"), trade_count=("trade_amount", "count"))
                .sort_values("month")
            )
            monthly_total = (
                vol_universe.groupby("month", as_index=False)
                .agg(total_volume=("trade_amount", "sum"))
                .sort_values("month")
            )
            primary_monthly = (
                vol_universe[vol_universe["issuer"] == selected_issuer]
                .groupby("month", as_index=False)
                .agg(primary_volume=("trade_amount", "sum"))
            )
            pct_df = monthly_total.merge(primary_monthly, on="month", how="left")
            pct_df["primary_volume"] = pct_df["primary_volume"].fillna(0)
            pct_df["primary_pct"] = np.where(pct_df["total_volume"] > 0, pct_df["primary_volume"] / pct_df["total_volume"] * 100, np.nan)

            fig_vol_snapshot = make_subplots(specs=[[{"secondary_y": True}]])
            group_order = [selected_issuer] + [x for x in comparison_issuers_sidebar if x != selected_issuer]
            if volume_comparison_mode == "Primary vs Peers vs All Other":
                group_order += ["All Other Uploaded Munis"]
            else:
                group_order = snapshot_issuers
            for group_name in group_order:
                tmp = monthly_grouped[monthly_grouped["volume_group"] == group_name]
                if tmp.empty:
                    continue
                fig_vol_snapshot.add_trace(
                    go.Bar(
                        x=tmp["month"],
                        y=tmp["monthly_volume"] / 1_000_000,
                        name=group_name,
                        hovertemplate="%{x|%b %Y}<br>$%{y:,.1f}M<extra>%{fullData.name}</extra>",
                    ),
                    secondary_y=False,
                )
            if not pct_df.empty:
                fig_vol_snapshot.add_trace(
                    go.Scatter(
                        x=pct_df["month"],
                        y=pct_df["primary_pct"],
                        mode="lines+markers",
                        name=f"{selected_issuer} % of Volume",
                        line=dict(width=3),
                        hovertemplate="%{x|%b %Y}<br>%{y:.1f}%<extra>%{fullData.name}</extra>",
                    ),
                    secondary_y=True,
                )
            fig_vol_snapshot.update_layout(
                title="Monthly Trading Volume + Primary Issuer Share",
                barmode="stack",
                height=560,
                margin=dict(l=40, r=55, t=70, b=50),
                legend_title_text="Volume Item",
            )
            fig_vol_snapshot.update_xaxes(title_text="Trade Month", tickformat="%b %Y")
            fig_vol_snapshot.update_yaxes(title_text="Secondary Market Volume ($MM)", secondary_y=False)
            fig_vol_snapshot.update_yaxes(title_text=f"{selected_issuer} % of Total", ticksuffix="%", secondary_y=True)
            safe_plotly_chart(fig_vol_snapshot, width="stretch")
            render_volume_readthrough(vol_universe, selected_issuer)
        else:
            st.info("No usable trade amount data for monthly volume.")
    else:
        st.info("Upload trades with issuer, trade date, and trade amount to build trading-volume views.")

snap_col_a, snap_col_b, snap_col_c = st.columns(3)
with snap_col_a:
    _metric_source = market_df[market_df["issuer"] == selected_issuer].copy()
    if not _metric_source.empty and "spread" in _metric_source.columns:
        _spread_bps = pd.to_numeric(_metric_source["spread"], errors="coerce") * 100
        clean_metric_card("Primary Median Spread", "N/A" if _spread_bps.dropna().empty else f"{_spread_bps.median():.1f} bps", size="small")
    else:
        clean_metric_card("Primary Median Spread", "N/A", size="small")
with snap_col_b:
    if not snapshot_base.empty and "trade_amount" in snapshot_base.columns:
        _amt = pd.to_numeric(snapshot_base["trade_amount"], errors="coerce").sum()
        clean_metric_card("Selected Volume", f"${_amt/1_000_000:,.1f}M", size="small")
    else:
        clean_metric_card("Selected Volume", "N/A", size="small")
with snap_col_c:
    clean_metric_card("Compared Issuers", f"{len(snapshot_issuers):,}", size="small")

section_anchor("yield-relative-value", "Yield Trend / Relative Value Comparison")
with st.expander("Methodology: benchmark curve framework", expanded=False):
    st.markdown(
        """
This section groups uploaded trade rows by **trade date** and **issuer**, then plots average observed trade yield.

**Benchmark logic:**

- **Primary source = trade-sheet `Index` / `Index Rate`**, when available. This keeps benchmark spread analytics aligned with the same pricing environment as the uploaded MuniPro trades.
- **Uploaded MMD is fallback only** and is used only when trade-sheet index data is unavailable.
- The app intentionally uses **one benchmark universe at a time**; it does not mix trade-sheet index rates with external MMD rates in the same run.
- If explicit non-AAA curves are unavailable, the app can still use visible rating-spread assumptions as an analytical approximation.
- Units in the code are percentage points: `0.10 = 10 bps`.
- This is an internal analytical benchmark, not a live Bloomberg/BVAL/ICE curve. Replace assumptions with firm-approved or vendor curves when available.
        """
    )
    safe_dataframe(rating_spread_table(), width="stretch", hide_index=True)

issuer_choices = uploaded_issuers
default_compare = [selected_issuer] if selected_issuer in issuer_choices else issuer_choices[:1]
compare_issuers = st.multiselect("Compare Issuers", issuer_choices, default=default_compare)
compare_bucket = st.selectbox("Comparison Maturity Year", MATURITY_BUCKET_OPTIONS, key="compare_bucket")
benchmark_ratings = st.multiselect(
    "Benchmark Curve(s)",
    BENCHMARK_RATINGS,
    default=["AAA", "AA"],
    help="Priority: uploaded rating curve columns first; otherwise MMD/AAA plus the visible rating-spread assumptions above.",
)
show_spread_to_benchmark = st.checkbox(
    "Show issuer spread to selected benchmark",
    value=True,
    help="Calculates average issuer yield minus selected benchmark curve for dates where both are available.",
)

chart_df = market_df[market_df["issuer"].isin(compare_issuers)].copy()
if compare_bucket != "All":
    chart_df = chart_df[chart_df["maturity_bucket"] == compare_bucket].copy()

if chart_df.empty:
    st.warning("No trade data found for selected comparison filters.")
else:
    date_min = chart_df["trade_date"].min().date()
    date_max = chart_df["trade_date"].max().date()
    selected_dates = st.date_input("Trade Date Range", value=(date_min, date_max), min_value=date_min, max_value=date_max)
    if isinstance(selected_dates, tuple) and len(selected_dates) == 2:
        start_date, end_date = selected_dates
        chart_df = chart_df[(chart_df["trade_date"].dt.date >= start_date) & (chart_df["trade_date"].dt.date <= end_date)].copy()

    daily = (
        chart_df.groupby(["trade_date", "issuer"], as_index=False)
        .agg(avg_yield=("yield", "mean"), trade_count=("yield", "count"), total_trade_amount=("trade_amount", "sum"))
    )
    fig = px.line(
        daily.sort_values("trade_date"),
        x="trade_date",
        y="avg_yield",
        color="issuer",
        markers=True,
        hover_data=["trade_count", "total_trade_amount"],
        title="Average Trade Yield by Issuer",
    )

    benchmark_daily = pd.DataFrame()
    benchmark_ready = False
    if not mmd_df.empty and benchmark_ratings:
        date_col = _detect_mmd_date_column(mmd_df)
        mmd_col = MMD_BUCKET_MAP.get(compare_bucket, "10Y")
        if date_col:
            mmd_plot = mmd_df.copy()
            mmd_plot[date_col] = pd.to_datetime(mmd_plot[date_col], errors="coerce")
            mmd_plot = mmd_plot.dropna(subset=[date_col])
            if isinstance(selected_dates, tuple) and len(selected_dates) == 2:
                mmd_plot = mmd_plot[(mmd_plot[date_col].dt.date >= start_date) & (mmd_plot[date_col].dt.date <= end_date)]

            benchmark_frames = []
            unavailable_ratings = []
            for rating in benchmark_ratings:
                y, meta = get_benchmark_curve(mmd_plot, mmd_col, rating)
                if y is None:
                    unavailable_ratings.append(rating)
                    continue
                fig.add_scatter(
                    x=mmd_plot[date_col],
                    y=y,
                    mode="lines",
                    name=f"{rating} Curve ({mmd_col})",
                )
                benchmark_frames.append(
                    pd.DataFrame({
                        "trade_date": mmd_plot[date_col].dt.normalize(),
                        "benchmark_rating": rating,
                        "benchmark_yield": y,
                        "mmd_tenor": mmd_col,
                        "rating_spread_bps": meta.get("rating_spread_bps"),
                        "benchmark_source": meta.get("benchmark_source"),
                        "source_column": meta.get("source_column"),
                    })
                )
            benchmark_daily = pd.concat(benchmark_frames, ignore_index=True) if benchmark_frames else pd.DataFrame()
            benchmark_ready = not benchmark_daily.empty
            if unavailable_ratings:
                st.warning(
                    "Some benchmark curves could not be built because neither an uploaded curve column nor a usable AAA/MMD base tenor was found: "
                    + ", ".join(unavailable_ratings)
                )
        else:
            st.warning("Benchmark curves could not be plotted because the curve file does not contain a usable date column.")

    fig.update_layout(xaxis_title="Trade Date", yaxis_title="Yield (%)", hovermode="x unified")
    safe_plotly_chart(fig, width="stretch")

    if show_spread_to_benchmark and benchmark_ready and not daily.empty:
        spread_base = daily.copy()
        spread_base["trade_date"] = pd.to_datetime(spread_base["trade_date"], errors="coerce").dt.normalize()
        spread_to_benchmark = spread_base.merge(benchmark_daily, on="trade_date", how="inner")
        if spread_to_benchmark.empty:
            st.info("No overlapping dates were found between issuer trades and the selected benchmark curve.")
        else:
            spread_to_benchmark["spread_to_benchmark_bps"] = (
                spread_to_benchmark["avg_yield"] - spread_to_benchmark["benchmark_yield"]
            ) * 100
            spread_fig = px.line(
                spread_to_benchmark.sort_values("trade_date"),
                x="trade_date",
                y="spread_to_benchmark_bps",
                color="issuer",
                line_dash="benchmark_rating",
                markers=True,
                hover_data=["benchmark_rating", "mmd_tenor", "benchmark_source", "source_column", "rating_spread_bps", "trade_count", "total_trade_amount"],
                title="Issuer Spread to Selected Benchmark Curve(s)",
            )
            spread_fig.update_layout(xaxis_title="Trade Date", yaxis_title="Spread to Benchmark (bps)", hovermode="x unified")
            safe_plotly_chart(spread_fig, width="stretch")

            with st.expander("Spread-to-benchmark calculation details", expanded=False):
                st.markdown(
                    """
For each issuer/date/rating benchmark:

`Spread to Benchmark (bps) = (Average Issuer Trade Yield - Synthetic Benchmark Yield) × 100`

Where:

`Benchmark Yield = uploaded rating curve if available; otherwise MMD/AAA Tenor Yield + Rating Spread Assumption`
                    """
                )
                safe_dataframe(
                    spread_to_benchmark[[
                        "trade_date", "issuer", "benchmark_rating", "mmd_tenor", "avg_yield",
                        "benchmark_yield", "benchmark_source", "source_column", "rating_spread_bps", "spread_to_benchmark_bps",
                        "trade_count", "total_trade_amount",
                    ]].sort_values(["trade_date", "issuer", "benchmark_rating"], ascending=[False, True, True]).head(1000),
                    width="stretch",
                    hide_index=True,
                )
    elif show_spread_to_benchmark and mmd_df.empty:
        st.info("Upload an MMD curve file to enable AAA/AA/A/BBB benchmark curves and spread-to-benchmark analytics.")


section_anchor("issuer-curve", "Issuer Curve vs Benchmark Curve")
with st.expander("Methodology: issuer curve vs benchmark curve", expanded=False):
    st.markdown(
        """
This chart shows a **cross-sectional yield curve** by maturity year, rather than a time-series trend.

**Issuer curve logic:**

- The issuer curve is built from uploaded trade yields by maturity year: **1Y / 2Y / 3Y / ...**.
- Default aggregation uses **average yield over the latest selected window** ending on the curve date. This reduces noise from sparse municipal trading.
- You can also use **latest trade per bucket** when you want the most recent observation in each maturity year.

**Benchmark curve logic:**

- The benchmark curve uses uploaded rating curve columns when available.
- If an uploaded AA/A/BBB curve is missing, the app falls back to **MMD/AAA + transparent rating-spread assumptions**.
- The benchmark value uses the latest available curve observation at or before the selected curve date.

**How to read it:**

- If the issuer line is above the benchmark line, the issuer trades cheaper / wider for that bucket.
- If the issuer line is below the benchmark line, the issuer trades richer / tighter for that bucket.
- The accompanying table shows exact yields and spreads in basis points.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD / benchmark curve file to enable Issuer Curve vs Benchmark Curve analysis.")
else:
    selected_issuer_dates = pd.to_datetime(
        market_df.loc[market_df["issuer"] == selected_issuer, "trade_date"], errors="coerce"
    ).dropna()

    if selected_issuer_dates.empty:
        st.warning("No valid trade dates were found for the selected issuer, so the issuer curve cannot be built.")
    else:
        curve_min_date = selected_issuer_dates.min().date()
        curve_max_date = selected_issuer_dates.max().date()

        curve_ctrl1, curve_ctrl2, curve_ctrl3 = st.columns([1, 1, 1.4])
        with curve_ctrl1:
            curve_as_of_date = st.date_input(
                "Curve Date",
                value=curve_max_date,
                min_value=curve_min_date,
                max_value=curve_max_date,
                key="issuer_curve_as_of_date",
                help="The issuer and benchmark curves use observations available at or before this date.",
            )
        with curve_ctrl2:
            curve_aggregation = st.selectbox(
                "Issuer Curve Aggregation",
                ["Average last N days", "Latest trade per bucket"],
                index=0,
                key="issuer_curve_aggregation",
            )
        with curve_ctrl3:
            curve_benchmark_ratings = st.multiselect(
                "Benchmark Curve(s) for Curve Chart",
                BENCHMARK_RATINGS,
                default=[r for r in ["AAA", "AA"] if r in BENCHMARK_RATINGS],
                key="issuer_curve_benchmark_ratings",
                help="Priority: uploaded rating curve columns first; otherwise MMD/AAA plus the visible spread assumptions.",
            )

        curve_lookback_days = 30
        if curve_aggregation == "Average last N days":
            curve_lookback_days = st.select_slider(
                "Lookback Window for Issuer Curve",
                options=[7, 14, 30, 60, 90, 180],
                value=30,
                format_func=lambda x: f"{x} days",
                key="issuer_curve_lookback_days",
                help="Municipal trades can be sparse, so averaging over a window usually gives a more stable curve than using one day only.",
            )

        if not curve_benchmark_ratings:
            st.info("Select at least one benchmark curve to display the issuer curve comparison.")
        else:
            issuer_curve_plot_df, issuer_curve_audit = build_issuer_curve_snapshot(
                market_df=market_df,
                mmd_df=mmd_df,
                issuer=selected_issuer,
                ratings=curve_benchmark_ratings,
                as_of_date=pd.Timestamp(curve_as_of_date),
                lookback_days=curve_lookback_days,
                aggregation_method=curve_aggregation,
            )

            if issuer_curve_plot_df.empty or issuer_curve_audit.empty:
                st.warning(
                    "No overlapping issuer trades and benchmark curve observations were found for this curve setup. "
                    "Try a longer lookback window, a different curve date, or check that the benchmark file has usable 5Y/10Y/20Y/30Y columns."
                )
            else:
                curve_fig = px.line(
                    issuer_curve_plot_df,
                    x="maturity_bucket",
                    y="yield_value",
                    color="curve",
                    markers=True,
                    hover_data=["curve_type", "trade_count", "issuer_observation_date"],
                    title=f"{selected_issuer} Issuer Curve vs Benchmark Curve(s)",
                    labels={
                        "maturity_bucket": "Maturity Year",
                        "yield_value": "Yield (%)",
                        "curve": "Curve",
                    },
                )
                curve_fig.update_layout(hovermode="x unified", height=500)
                safe_plotly_chart(curve_fig, width="stretch")

                table_cols = [
                    "maturity_bucket", "benchmark_rating", "issuer_yield", "benchmark_yield",
                    "spread_to_benchmark_bps", "trade_count", "issuer_observation_date", "benchmark_date",
                    "mmd_tenor", "benchmark_source", "source_column", "rating_spread_bps",
                    "aggregation_method", "lookback_start", "lookback_end",
                ]
                curve_table = issuer_curve_audit[[c for c in table_cols if c in issuer_curve_audit.columns]].copy()
                for c in ["issuer_yield", "benchmark_yield", "spread_to_benchmark_bps", "rating_spread_bps"]:
                    if c in curve_table.columns:
                        curve_table[c] = pd.to_numeric(curve_table[c], errors="coerce").round(2)

                st.subheader("Curve Spread Table")
                safe_dataframe(curve_table, width="stretch", hide_index=True)

                primary_curve_rating = curve_benchmark_ratings[0]
                primary_rows = issuer_curve_audit[issuer_curve_audit["benchmark_rating"] == primary_curve_rating].copy()
                primary_rows = primary_rows.dropna(subset=["spread_to_benchmark_bps"])
                if not primary_rows.empty:
                    cheap_row = primary_rows.loc[primary_rows["spread_to_benchmark_bps"].idxmax()]
                    rich_row = primary_rows.loc[primary_rows["spread_to_benchmark_bps"].idxmin()]
                    st.info(
                        f"Curve read-through vs {primary_curve_rating}: "
                        f"{cheap_row['maturity_bucket']} is the widest bucket at {cheap_row['spread_to_benchmark_bps']:+.1f} bp, "
                        f"while {rich_row['maturity_bucket']} is the tightest bucket at {rich_row['spread_to_benchmark_bps']:+.1f} bp."
                    )



section_anchor("spread-movement", "Spread Movement Ladder")
with st.expander("Methodology: spread movement ladder", expanded=False):
    st.markdown(
        """
This section replaces the old ladder with a desk-readable **ranked ladder**.

**Calculation:**

`Issuer Spread = (Average Issuer Trade Yield - Benchmark Yield) × 100`

`Spread Movement = Latest Available Issuer Spread - Historical Issuer Spread`

**How to read it:**

- Positive = widening / cheaper versus the benchmark.
- Negative = tightening / richer versus the benchmark.
- The chart ranks the largest absolute movers first so the signal is visible without reading a dense grid.
        """
    )

if mmd_df.empty:
    st.info("Upload or enable a benchmark source to calculate spread movement ladders.")
else:
    ladder_col1, ladder_col2 = st.columns([1, 2])
    with ladder_col1:
        movement_rating = st.selectbox(
            "Movement Benchmark Curve",
            BENCHMARK_RATINGS,
            index=0,
            key="movement_ladder_rating",
        )
    with ladder_col2:
        movement_window = st.selectbox(
            "Lookback Window",
            ["1W", "1M", "3M", "6M", "1Y"],
            index=2,
            help="Ranks maturity years by spread movement over the selected lookback window.",
        )

    movement_spread_obs = build_spread_observations(
        market_df=market_df,
        mmd_df=mmd_df,
        issuer=selected_issuer,
        rating=movement_rating,
    )
    if movement_spread_obs.empty:
        st.info(
            "No overlapping issuer trade dates and benchmark dates were found for the spread movement ladder. "
            "Check whether the uploaded benchmark file covers the same dates as the trade data."
        )
    else:
        movement_table, movement_audit = build_spread_movement_ladder_data(movement_spread_obs)
        if movement_table.empty or movement_table.isna().all().all() or movement_window not in movement_table.columns:
            st.info("No spread movement observations are available for the selected lookback window.")
        else:
            ladder_df = movement_table[[movement_window]].reset_index().rename(
                columns={"index": "maturity_bucket", movement_window: "spread_movement_bps"}
            )
            if "maturity_bucket" not in ladder_df.columns:
                ladder_df = ladder_df.rename(columns={ladder_df.columns[0]: "maturity_bucket"})
            ladder_df["spread_movement_bps"] = pd.to_numeric(ladder_df["spread_movement_bps"], errors="coerce")
            ladder_df = ladder_df.dropna(subset=["spread_movement_bps"]).copy()
            ladder_df["abs_movement_bps"] = ladder_df["spread_movement_bps"].abs()
            ladder_df["maturity_zone"] = ladder_df["maturity_bucket"].apply(maturity_zone_label)
            ladder_df["signal"] = np.where(ladder_df["spread_movement_bps"] >= 0, "Widened / Cheaper", "Tightened / Richer")
            ladder_df = ladder_df.sort_values("abs_movement_bps", ascending=False).head(15)

            if ladder_df.empty:
                st.info("No valid spread movement values after filtering.")
            else:
                movement_fig = px.bar(
                    ladder_df.sort_values("spread_movement_bps"),
                    x="spread_movement_bps",
                    y="maturity_bucket",
                    orientation="h",
                    color="signal",
                    hover_data=["maturity_zone", "abs_movement_bps"],
                    title=f"{selected_issuer} Largest Spread Movers vs {movement_rating} ({movement_window})",
                    labels={
                        "spread_movement_bps": "Spread Movement (bps)",
                        "maturity_bucket": "Maturity Year",
                        "signal": "Signal",
                    },
                )
                movement_fig.add_vline(x=0, line_dash="dash", opacity=0.45)
                movement_fig.update_layout(height=max(420, 28 * len(ladder_df) + 160), legend_title_text="Direction")
                safe_plotly_chart(movement_fig, width="stretch")

                top_move = ladder_df.iloc[0]
                st.info(
                    f"Largest {movement_window} move: {top_move['maturity_bucket']} "
                    f"moved {top_move['spread_movement_bps']:+.1f} bps vs {movement_rating}."
                )

                with st.expander("Movement calculation audit table", expanded=False):
                    display_cols = [
                        "maturity_bucket", "window", "latest_date", "latest_spread_bps",
                        "target_date", "historical_date", "historical_spread_bps", "spread_movement_bps", "note",
                    ]
                    audit_display = movement_audit[[c for c in display_cols if c in movement_audit.columns]].copy()
                    for c in ["latest_spread_bps", "historical_spread_bps", "spread_movement_bps"]:
                        if c in audit_display.columns:
                            audit_display[c] = pd.to_numeric(audit_display[c], errors="coerce").round(2)
                    safe_dataframe(audit_display, width="stretch", hide_index=True)


section_anchor("liquidity", "Liquidity / Trading Frequency Analysis")
with st.expander("Methodology", expanded=False):
    st.write("Liquidity score is a transparent ranking measure: 35% trade count, 25% total trade amount, 25% recent 90-day trades, and 15% recency. It is a screening metric, not a credit rating or valuation recommendation.")
if issuer_trades.empty:
    st.warning("No trade rows found for this issuer and filter.")
else:
    today = pd.Timestamp.today().normalize()
    liq_base = issuer_trades.copy()
    liq_base["trade_month"] = liq_base["trade_date"].dt.to_period("M").astype(str)
    # Build aggregation dynamically so optional bond/security enrichment
    # columns do not trigger KeyError in trade-only mode.
    liquidity_agg = {
        "trade_count": ("trade_date", "count"),
        "first_trade": ("trade_date", "min"),
        "latest_trade": ("trade_date", "max"),
        "active_months": ("trade_month", "nunique"),
    }

    if "yield" in liq_base.columns:
        liquidity_agg.update({
            "avg_yield": ("yield", "mean"),
            "min_yield": ("yield", "min"),
            "max_yield": ("yield", "max"),
        })

    if "price" in liq_base.columns:
        liquidity_agg["avg_price"] = ("price", "mean")

    if "trade_amount" in liq_base.columns:
        liquidity_agg.update({
            "total_trade_amount": ("trade_amount", "sum"),
            "avg_trade_amount": ("trade_amount", "mean"),
            "median_trade_amount": ("trade_amount", "median"),
        })

    if "maturity_bond" in liq_base.columns:
        liquidity_agg["maturity"] = ("maturity_bond", "first")
    elif "maturity" in liq_base.columns:
        liquidity_agg["maturity"] = ("maturity", "first")

    if "coupon_bond" in liq_base.columns:
        liquidity_agg["coupon"] = ("coupon_bond", "first")
    elif "coupon" in liq_base.columns:
        liquidity_agg["coupon"] = ("coupon", "first")

    if "outstanding_amount" in liq_base.columns:
        liquidity_agg["outstanding_amount"] = ("outstanding_amount", "first")

    liq = (
        liq_base.groupby("cusip", dropna=False)
        .agg(**liquidity_agg)
        .reset_index()
    )

    # Ensure downstream formulas have safe defaults when optional columns are absent.
    if "total_trade_amount" not in liq.columns:
        liq["total_trade_amount"] = 0
    if "avg_trade_amount" not in liq.columns:
        liq["avg_trade_amount"] = pd.NA
    if "median_trade_amount" not in liq.columns:
        liq["median_trade_amount"] = pd.NA
    if "outstanding_amount" not in liq.columns:
        liq["outstanding_amount"] = pd.NA
    if "avg_yield" not in liq.columns:
        liq["avg_yield"] = pd.NA
    if "min_yield" not in liq.columns:
        liq["min_yield"] = pd.NA
    if "max_yield" not in liq.columns:
        liq["max_yield"] = pd.NA
    if "avg_price" not in liq.columns:
        liq["avg_price"] = pd.NA
    liq["days_since_last_trade"] = (today - liq["latest_trade"]).dt.days
    liq["trading_period_days"] = (liq["latest_trade"] - liq["first_trade"]).dt.days.clip(lower=1)
    liq["avg_days_between_trades"] = liq["trading_period_days"] / liq["trade_count"].clip(lower=1)
    liq["avg_trades_per_month"] = liq["trade_count"] / liq["active_months"].clip(lower=1)
    recent_cutoff = today - pd.DateOffset(days=90)
    recent = liq_base[liq_base["trade_date"] >= recent_cutoff].groupby("cusip").agg(recent_90d_trades=("trade_date", "count")).reset_index()
    liq = liq.merge(recent, on="cusip", how="left")
    liq["recent_90d_trades"] = liq["recent_90d_trades"].fillna(0).astype(int)
    liq["max_yield"] = pd.to_numeric(liq["max_yield"], errors="coerce")
    liq["min_yield"] = pd.to_numeric(liq["min_yield"], errors="coerce")
    liq["yield_range"] = liq["max_yield"] - liq["min_yield"]
    liq["total_trade_amount"] = pd.to_numeric(liq["total_trade_amount"], errors="coerce").fillna(0)
    liq["outstanding_amount"] = pd.to_numeric(liq["outstanding_amount"], errors="coerce")
    liq["turnover_ratio"] = liq["total_trade_amount"] / liq["outstanding_amount"].replace({0: pd.NA})
    liq["liquidity_score"] = (
        liq["trade_count"].rank(pct=True) * 35
        + liq["total_trade_amount"].rank(pct=True) * 25
        + liq["recent_90d_trades"].rank(pct=True) * 25
        + (1 - liq["days_since_last_trade"].rank(pct=True)) * 15
    )
    liq["liquidity_tier"] = pd.cut(
        liq["liquidity_score"], bins=[-1, 45, 75, 101], labels=["Low Liquidity", "Medium Liquidity", "High Liquidity"]
    ).astype(str)
    liq.loc[liq["days_since_last_trade"] > 365, "liquidity_tier"] = "Stale"
    liq = liq.sort_values(["liquidity_score", "trade_count", "total_trade_amount"], ascending=False)

    monthly = liq_base.groupby("trade_month", as_index=False).agg(trade_count=("trade_date", "count"), total_trade_amount=("trade_amount", "sum"), avg_yield=("yield", "mean"))
    st.subheader("1. Market Activity Over Time")
    safe_plotly_chart(px.line(monthly, x="trade_month", y="trade_count", markers=True, title="Monthly Trade Count"), width="stretch")
    render_monthly_activity_readthrough(monthly)

    st.subheader("2. Trade Size Distribution")
    with st.expander("Methodology: trade size distribution", expanded=False):
        st.markdown(
            """
This chart groups trades by par/trade amount to show whether activity is primarily retail-sized, institutional-sized, or block-oriented.

**Default buckets:**

- **< $100k**: odd-lot / retail-sized activity
- **$100k–$250k**: small institutional or advisor-sized activity
- **$250k–$1mm**: institutional-sized activity
- **$1mm+**: block trade / larger institutional flow

This is useful because trade count alone can overstate liquidity when most activity comes from small trades.
            """
        )

    if "trade_amount" not in liq_base.columns:
        st.info("Trade size distribution is unavailable because trade_amount is missing from the uploaded trade data.")
    else:
        trade_size_df = liq_base.copy()
        trade_size_df["trade_amount"] = pd.to_numeric(trade_size_df["trade_amount"], errors="coerce")
        trade_size_df = trade_size_df.dropna(subset=["trade_amount"])
        trade_size_df = trade_size_df[trade_size_df["trade_amount"] > 0]

        if trade_size_df.empty:
            st.info("Trade size distribution is unavailable because no positive trade_amount values were found.")
        else:
            trade_size_bins = [0, 100_000, 250_000, 1_000_000, float("inf")]
            trade_size_labels = ["< $100k", "$100k–$250k", "$250k–$1mm", "$1mm+"]
            trade_size_df["trade_size_bucket"] = pd.cut(
                trade_size_df["trade_amount"],
                bins=trade_size_bins,
                labels=trade_size_labels,
                include_lowest=True,
                right=False,
            )

            size_summary = (
                trade_size_df.groupby("trade_size_bucket", observed=False)
                .agg(
                    trade_count=("trade_amount", "count"),
                    total_trade_amount=("trade_amount", "sum"),
                    avg_trade_amount=("trade_amount", "mean"),
                    median_trade_amount=("trade_amount", "median"),
                )
                .reset_index()
            )
            size_summary["trade_size_bucket"] = size_summary["trade_size_bucket"].astype(str)
            size_summary["trade_count_share"] = size_summary["trade_count"] / size_summary["trade_count"].sum()
            size_summary["amount_share"] = size_summary["total_trade_amount"] / size_summary["total_trade_amount"].sum()

            size_fig = px.bar(
                size_summary,
                x="trade_size_bucket",
                y="trade_count",
                hover_data={
                    "total_trade_amount": ":,.0f",
                    "avg_trade_amount": ":,.0f",
                    "median_trade_amount": ":,.0f",
                    "trade_count_share": ":.1%",
                    "amount_share": ":.1%",
                },
                title="Trade Count by Size Bucket",
                labels={
                    "trade_size_bucket": "Trade Size Bucket",
                    "trade_count": "Number of Trades",
                    "total_trade_amount": "Total Trade Amount",
                    "trade_count_share": "Share of Trades",
                    "amount_share": "Share of Par Traded",
                },
            )
            size_fig.update_layout(height=430)
            safe_plotly_chart(size_fig, width="stretch")

            amount_fig = px.bar(
                size_summary,
                x="trade_size_bucket",
                y="total_trade_amount",
                hover_data={
                    "trade_count": ":,.0f",
                    "avg_trade_amount": ":,.0f",
                    "median_trade_amount": ":,.0f",
                    "trade_count_share": ":.1%",
                    "amount_share": ":.1%",
                },
                title="Total Par Traded by Size Bucket",
                labels={
                    "trade_size_bucket": "Trade Size Bucket",
                    "total_trade_amount": "Total Trade Amount",
                    "trade_count": "Number of Trades",
                    "trade_count_share": "Share of Trades",
                    "amount_share": "Share of Par Traded",
                },
            )
            amount_fig.update_layout(height=430)
            safe_plotly_chart(amount_fig, width="stretch")
            render_trade_size_readthrough(size_summary)

            retail_trade_share = size_summary.loc[
                size_summary["trade_size_bucket"] == "< $100k", "trade_count_share"
            ]
            block_amount_share = size_summary.loc[
                size_summary["trade_size_bucket"] == "$1mm+", "amount_share"
            ]

            retail_trade_share_val = float(retail_trade_share.iloc[0]) if not retail_trade_share.empty else 0.0
            block_amount_share_val = float(block_amount_share.iloc[0]) if not block_amount_share.empty else 0.0

            if retail_trade_share_val >= 0.60 and block_amount_share_val < 0.25:
                st.info(
                    f"Read-through: trading activity appears retail / odd-lot heavy. "
                    f"< $100k trades account for {retail_trade_share_val:.1%} of trades, "
                    f"while $1mm+ blocks account for {block_amount_share_val:.1%} of par traded."
                )
            elif block_amount_share_val >= 0.50:
                st.info(
                    f"Read-through: activity appears institutionally active. "
                    f"$1mm+ blocks account for {block_amount_share_val:.1%} of par traded."
                )
            else:
                st.info(
                    f"Read-through: trade activity is mixed across retail-sized and institutional-sized buckets. "
                    f"< $100k trades account for {retail_trade_share_val:.1%} of trades; "
                    f"$1mm+ blocks account for {block_amount_share_val:.1%} of par traded."
                )

            with st.expander("Trade size distribution table", expanded=False):
                table_display = size_summary.copy()
                for pct_col in ["trade_count_share", "amount_share"]:
                    table_display[pct_col] = table_display[pct_col].map(lambda x: f"{x:.1%}" if pd.notna(x) else "")
                for amt_col in ["total_trade_amount", "avg_trade_amount", "median_trade_amount"]:
                    table_display[amt_col] = pd.to_numeric(table_display[amt_col], errors="coerce").round(0)
                safe_dataframe(table_display, width="stretch", hide_index=True)

    st.subheader("3. Most Frequently Traded CUSIPs")
    safe_plotly_chart(px.bar(liq.head(25), x="cusip", y="trade_count", color="liquidity_tier", title="Top 25 Most Frequently Traded CUSIPs"), width="stretch")

    st.subheader("4. Trade Recency / Staleness")
    safe_plotly_chart(px.histogram(liq, x="days_since_last_trade", nbins=30, color="liquidity_tier", title="Distribution of Days Since Last Trade"), width="stretch")
    render_liquidity_readthrough(liq)

    st.subheader("5. Liquidity Ranking Table")
    display_cols = [
        "cusip", "liquidity_tier", "liquidity_score", "trade_count", "recent_90d_trades", "active_months",
        "avg_trades_per_month", "avg_days_between_trades", "days_since_last_trade", "first_trade", "latest_trade",
        "avg_yield", "yield_range", "avg_price", "total_trade_amount", "avg_trade_amount", "turnover_ratio",
        "maturity", "coupon", "outstanding_amount",
    ]
    safe_dataframe(liq[[c for c in display_cols if c in liq.columns]], width="stretch", height=500)

section_anchor("cusip-drilldown", "CUSIP Opportunity Drilldown")
with st.expander("Methodology: CUSIP opportunity drilldown", expanded=False):
    st.markdown(
        """
This section moves from issuer-level signals into **specific bond-level candidates**.

**Purpose:**

- Identify which CUSIPs are driving a maturity year's relative-value signal.
- Compare current CUSIP-level spread, recent movement, trade count, and liquidity.
- Help the team move from: *"30Y widened"* to *"which 30Y bonds should we look at?"*

**Calculation overview:**

- Current CUSIP yield is calculated over the selected lookback window.
- Current spread is calculated as:

`CUSIP Spread = (Average CUSIP Yield - Benchmark Yield) × 100`

- Historical spread uses the most recent observation at or before the lookback target date.
- Spread change is:

`Spread Change = Current Spread - Historical Spread`

Positive spread change means widening; negative spread change means tightening.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD curve file to enable CUSIP-level spread drilldown.")
elif issuer_trades.empty:
    st.warning("No trade rows found for the selected issuer and filters.")
else:
    dd_col1, dd_col2, dd_col3, dd_col4 = st.columns([1, 1, 1, 1])
    with dd_col1:
        dd_bucket = st.selectbox(
            "Drilldown Maturity Year",
            MATURITY_BUCKET_ORDER,
            index=3,
            help="Focus the drilldown on one maturity year.",
        )
    with dd_col2:
        dd_rating = st.selectbox(
            "Drilldown Benchmark",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            help="Priority: uploaded benchmark curve first; otherwise modeled from MMD + spread assumptions.",
        )
    with dd_col3:
        dd_lookback_label = st.selectbox(
            "Movement Lookback",
            ["1W", "1M", "3M", "6M", "1Y"],
            index=1,
        )
    with dd_col4:
        dd_min_trades = st.number_input(
            "Minimum Trades",
            min_value=1,
            max_value=50,
            value=1,
            step=1,
            help="Filter out CUSIPs with fewer trades in the current lookback window.",
        )

    dd_window_days = {"1W": 7, "1M": 30, "3M": 90, "6M": 180, "1Y": 365}[dd_lookback_label]
    dd_tenor = MMD_BUCKET_MAP.get(dd_bucket, "10Y")
    dd_date_col = _detect_mmd_date_column(mmd_df)

    if dd_date_col is None:
        st.warning("CUSIP drilldown cannot run because the benchmark file does not contain a usable date column.")
    else:
        dd_base = market_df[
            (market_df["issuer"] == selected_issuer)
            & (market_df["maturity_bucket"] == dd_bucket)
        ].copy()

        if dd_base.empty:
            st.warning(f"No {dd_bucket} CUSIP-level trade rows were found for {selected_issuer}.")
        else:
            dd_base["trade_date"] = pd.to_datetime(dd_base["trade_date"], errors="coerce").dt.normalize()
            dd_base["yield"] = pd.to_numeric(dd_base["yield"], errors="coerce")
            if "trade_amount" in dd_base.columns:
                dd_base["trade_amount"] = pd.to_numeric(dd_base["trade_amount"], errors="coerce").fillna(0)
            else:
                dd_base["trade_amount"] = 0.0
            if "price" in dd_base.columns:
                dd_base["price"] = pd.to_numeric(dd_base["price"], errors="coerce")
            else:
                dd_base["price"] = pd.NA

            dd_base = dd_base.dropna(subset=["trade_date", "yield", "cusip"])

            if dd_base.empty:
                st.warning("CUSIP drilldown cannot run because no valid CUSIP/date/yield rows remain after cleaning.")
            else:
                dd_latest_date = dd_base["trade_date"].max()
                dd_current_start = dd_latest_date - pd.Timedelta(days=dd_window_days)
                dd_hist_target = dd_current_start

                dd_current = dd_base[dd_base["trade_date"] >= dd_current_start].copy()
                if dd_current.empty:
                    st.warning("No CUSIP trades were found inside the selected current lookback window.")
                else:
                    # Benchmark curve on or before latest issuer trade date.
                    dd_mmd = mmd_df.copy()
                    dd_mmd[dd_date_col] = pd.to_datetime(dd_mmd[dd_date_col], errors="coerce")
                    dd_mmd = dd_mmd.dropna(subset=[dd_date_col])
                    dd_mmd = dd_mmd[dd_mmd[dd_date_col].dt.normalize() <= dd_latest_date].sort_values(dd_date_col)

                    if dd_mmd.empty:
                        st.warning("No benchmark curve observation was available on or before the latest issuer trade date.")
                    else:
                        dd_latest_mmd = dd_mmd.iloc[[-1]].copy()
                        dd_benchmark_date = dd_latest_mmd[dd_date_col].iloc[0]
                        dd_benchmark_yield_series, dd_meta = get_benchmark_curve(dd_latest_mmd, dd_tenor, dd_rating)

                        if dd_benchmark_yield_series is None or pd.isna(dd_benchmark_yield_series.iloc[0]):
                            st.warning(f"{dd_rating} {dd_tenor} benchmark could not be built for this drilldown.")
                        else:
                            dd_benchmark_yield = float(dd_benchmark_yield_series.iloc[0])

                            current_summary = (
                                dd_current.groupby("cusip", dropna=False)
                                .agg(
                                    current_avg_yield=("yield", "mean"),
                                    latest_yield=("yield", "last"),
                                    avg_price=("price", "mean"),
                                    trade_count=("trade_date", "count"),
                                    latest_trade=("trade_date", "max"),
                                    first_trade=("trade_date", "min"),
                                    total_trade_amount=("trade_amount", "sum"),
                                    avg_trade_amount=("trade_amount", "mean"),
                                    maturity=("maturity_bond", "first") if "maturity_bond" in dd_current.columns else ("trade_date", "max"),
                                    coupon=("coupon_bond", "first") if "coupon_bond" in dd_current.columns else ("yield", "count"),
                                    call_date=("call_date", "first") if "call_date" in dd_current.columns else ("trade_date", "max"),
                                    call_price=("call_price", "first") if "call_price" in dd_current.columns else ("yield", "count"),
                                    outstanding_amount=("outstanding_amount", "first") if "outstanding_amount" in dd_current.columns else ("trade_amount", "sum"),
                                )
                                .reset_index()
                            )

                            current_summary["current_spread_bps"] = (
                                current_summary["current_avg_yield"] - dd_benchmark_yield
                            ) * 100

                            # Historical CUSIP spread at or before the lookback target date.
                            hist_candidates = dd_base[dd_base["trade_date"] <= dd_hist_target].copy()
                            if not hist_candidates.empty:
                                hist_rows = (
                                    hist_candidates.sort_values("trade_date")
                                    .groupby("cusip", as_index=False)
                                    .tail(1)[["cusip", "trade_date", "yield", "price", "trade_amount"]]
                                    .rename(
                                        columns={
                                            "trade_date": "historical_trade_date",
                                            "yield": "historical_yield",
                                            "price": "historical_price",
                                            "trade_amount": "historical_trade_amount",
                                        }
                                    )
                                )

                                dd_hist_mmd = mmd_df.copy()
                                dd_hist_mmd[dd_date_col] = pd.to_datetime(dd_hist_mmd[dd_date_col], errors="coerce")
                                dd_hist_mmd = dd_hist_mmd.dropna(subset=[dd_date_col])
                                dd_hist_mmd = dd_hist_mmd[dd_hist_mmd[dd_date_col].dt.normalize() <= dd_hist_target].sort_values(dd_date_col)

                                if not dd_hist_mmd.empty:
                                    dd_hist_latest_mmd = dd_hist_mmd.iloc[[-1]].copy()
                                    dd_hist_benchmark_yield_series, dd_hist_meta = get_benchmark_curve(
                                        dd_hist_latest_mmd, dd_tenor, dd_rating
                                    )
                                    if dd_hist_benchmark_yield_series is not None and pd.notna(dd_hist_benchmark_yield_series.iloc[0]):
                                        dd_hist_benchmark_yield = float(dd_hist_benchmark_yield_series.iloc[0])
                                        hist_rows["historical_spread_bps"] = (
                                            hist_rows["historical_yield"] - dd_hist_benchmark_yield
                                        ) * 100
                                    else:
                                        hist_rows["historical_spread_bps"] = pd.NA
                                else:
                                    hist_rows["historical_spread_bps"] = pd.NA
                            else:
                                hist_rows = pd.DataFrame(columns=["cusip", "historical_trade_date", "historical_yield", "historical_price", "historical_trade_amount", "historical_spread_bps"])

                            dd_opps = current_summary.merge(hist_rows, on="cusip", how="left")
                            dd_opps["spread_change_bps"] = dd_opps["current_spread_bps"] - dd_opps["historical_spread_bps"]
                            dd_opps["yield_change_bps"] = (dd_opps["current_avg_yield"] - dd_opps["historical_yield"]) * 100

                            # Liquidity score proxy for current window.
                            dd_today = pd.Timestamp.today().normalize()
                            dd_opps["days_since_last_trade"] = (dd_today - dd_opps["latest_trade"]).dt.days
                            dd_opps["liquidity_score"] = (
                                dd_opps["trade_count"].rank(pct=True) * 40
                                + dd_opps["total_trade_amount"].rank(pct=True) * 35
                                + (1 - dd_opps["days_since_last_trade"].rank(pct=True)) * 25
                            )
                            dd_opps["liquidity_tier"] = pd.cut(
                                dd_opps["liquidity_score"],
                                bins=[-1, 45, 75, 101],
                                labels=["Low", "Medium", "High"],
                            ).astype(str)

                            dd_opps = dd_opps[dd_opps["trade_count"] >= dd_min_trades].copy()

                            if dd_opps.empty:
                                st.info("No CUSIPs met the selected minimum trade filter.")
                            else:
                                dd_sort_options = {
                                    "Current Spread": "current_spread_bps",
                                    "Spread Change": "spread_change_bps",
                                    "Liquidity Score": "liquidity_score",
                                    "Trade Count": "trade_count",
                                    "Total Trade Amount": "total_trade_amount",
                                }
                                dd_sort_label = st.selectbox(
                                    "Sort Opportunities By",
                                    list(dd_sort_options.keys()),
                                    index=0,
                                    key="dd_sort_opportunities",
                                )
                                dd_sort_col = dd_sort_options[dd_sort_label]
                                dd_opps = dd_opps.sort_values(dd_sort_col, ascending=False, na_position="last")

                                summary_c1, summary_c2, summary_c3, summary_c4 = st.columns(4)
                                summary_c1.metric("CUSIPs Found", f"{len(dd_opps):,}")
                                summary_c2.metric("Bucket", dd_bucket)
                                summary_c3.metric("Total Par Traded", f"{dd_opps['total_trade_amount'].sum():,.0f}")
                                summary_c4.metric("Benchmark", f"{dd_rating} {dd_tenor}")

                                top_row = dd_opps.iloc[0]
                                spread_change_text = ""
                                if pd.notna(top_row.get("spread_change_bps")):
                                    spread_change_text = f"{top_row.get('spread_change_bps'):+.1f} bp spread change, "

                                st.info(
                                    f"Top read-through by {dd_sort_label}: CUSIP {top_row['cusip']} shows "
                                    f"{top_row['current_spread_bps']:+.1f} bp current spread to {dd_rating}, "
                                    f"{spread_change_text}"
                                    f"{int(top_row['trade_count'])} trades, and {top_row['liquidity_tier']} liquidity in the selected window."
                                )

                                display_cols = [
                                    "cusip", "coupon", "maturity", "call_date", "call_price",
                                    "current_avg_yield", "current_spread_bps", "spread_change_bps",
                                    "yield_change_bps", "trade_count", "total_trade_amount", "avg_trade_amount",
                                    "latest_trade", "historical_trade_date", "historical_spread_bps",
                                    "avg_price", "historical_price", "liquidity_score", "liquidity_tier",
                                    "outstanding_amount",
                                ]
                                dd_display = dd_opps[[c for c in display_cols if c in dd_opps.columns]].copy()
                                for col in ["current_avg_yield", "current_spread_bps", "spread_change_bps", "yield_change_bps", "avg_price", "historical_price", "liquidity_score"]:
                                    if col in dd_display.columns:
                                        dd_display[col] = pd.to_numeric(dd_display[col], errors="coerce").round(2)

                                st.subheader("CUSIP Opportunity Table")
                                safe_dataframe(dd_display, width="stretch", hide_index=True, height=420)

                                st.subheader("Security Detail")
                                selected_cusip = st.selectbox(
                                    "Select CUSIP for detail",
                                    dd_opps["cusip"].astype(str).tolist(),
                                    index=0,
                                    key="selected_cusip_drilldown",
                                )

                                sec_trades = dd_base[dd_base["cusip"].astype(str) == str(selected_cusip)].copy()
                                sec_trades = sec_trades.sort_values("trade_date")
                                if sec_trades.empty:
                                    st.warning("No trade rows found for the selected CUSIP.")
                                else:
                                    sec_daily = (
                                        sec_trades.groupby("trade_date", as_index=False)
                                        .agg(
                                            avg_yield=("yield", "mean"),
                                            trade_count=("yield", "count"),
                                            total_trade_amount=("trade_amount", "sum"),
                                            avg_price=("price", "mean"),
                                        )
                                    )

                                    # Build benchmark series for selected security dates.
                                    bench_long = make_benchmark_long(mmd_df, dd_rating)
                                    if not bench_long.empty:
                                        sec_daily = sec_daily.merge(
                                            bench_long[bench_long["maturity_bucket"] == dd_bucket][["trade_date", "benchmark_yield", "benchmark_source", "source_column"]],
                                            on="trade_date",
                                            how="left",
                                        )
                                        sec_daily["spread_to_benchmark_bps"] = (
                                            sec_daily["avg_yield"] - sec_daily["benchmark_yield"]
                                        ) * 100

                                    detail_col1, detail_col2 = st.columns(2)
                                    with detail_col1:
                                        sec_yield_fig = px.line(
                                            sec_daily,
                                            x="trade_date",
                                            y="avg_yield",
                                            markers=True,
                                            hover_data=["trade_count", "total_trade_amount", "avg_price"],
                                            title=f"{selected_cusip} Yield History",
                                            labels={"trade_date": "Trade Date", "avg_yield": "Average Yield (%)"},
                                        )
                                        sec_yield_fig.update_layout(height=380)
                                        safe_plotly_chart(sec_yield_fig, width="stretch")

                                    with detail_col2:
                                        if "spread_to_benchmark_bps" in sec_daily.columns and sec_daily["spread_to_benchmark_bps"].notna().any():
                                            sec_spread_fig = px.line(
                                                sec_daily,
                                                x="trade_date",
                                                y="spread_to_benchmark_bps",
                                                markers=True,
                                                hover_data=["trade_count", "total_trade_amount", "benchmark_source", "source_column"],
                                                title=f"{selected_cusip} Spread to {dd_rating} Benchmark",
                                                labels={"trade_date": "Trade Date", "spread_to_benchmark_bps": "Spread (bps)"},
                                            )
                                            sec_spread_fig.update_layout(height=380)
                                            safe_plotly_chart(sec_spread_fig, width="stretch")
                                            render_security_detail_readthrough(sec_daily, selected_cusip, dd_rating)
                                        else:
                                            sec_amt_fig = px.bar(
                                                sec_daily,
                                                x="trade_date",
                                                y="total_trade_amount",
                                                hover_data=["trade_count", "avg_yield", "avg_price"],
                                                title=f"{selected_cusip} Trade Amount History",
                                                labels={"trade_date": "Trade Date", "total_trade_amount": "Total Trade Amount"},
                                            )
                                            sec_amt_fig.update_layout(height=380)
                                            safe_plotly_chart(sec_amt_fig, width="stretch")

                                    with st.expander("Latest trades for selected CUSIP", expanded=False):
                                        latest_trade_cols = [
                                            "trade_datetime", "trade_date", "cusip", "description", "maturity_trade",
                                            "maturity_bond", "coupon_trade", "coupon_bond", "yield", "price",
                                            "trade_amount", "spread", "trade_type", "ratings_m_s_f",
                                        ]
                                        safe_dataframe(
                                            sec_trades[[c for c in latest_trade_cols if c in sec_trades.columns]]
                                            .sort_values("trade_date", ascending=False)
                                            .head(500),
                                            width="stretch",
                                            hide_index=True,
                                        )

                            with st.expander("Drilldown benchmark/audit details", expanded=False):
                                st.markdown(
                                    f"""
- Latest CUSIP trade date used: **{dd_latest_date.strftime('%Y-%m-%d')}**
- Current window start: **{dd_current_start.strftime('%Y-%m-%d')}**
- Historical target date: **{dd_hist_target.strftime('%Y-%m-%d')}**
- Benchmark date: **{dd_benchmark_date.strftime('%Y-%m-%d')}**
- Benchmark source: **{dd_meta.get('benchmark_source')}**
- Source column: **{dd_meta.get('source_column')}**
- Benchmark yield: **{dd_benchmark_yield:.4f}%**
                                    """
                                )


section_anchor("security-screener", "Security Screener — Top Relative Value Candidates")
with st.expander("Methodology: security screener", expanded=False):
    st.markdown(
        """
This section turns the dashboard into a practical **find me bonds** workflow.

**Goal:**

Screen uploaded bonds/trades for securities that are both relatively cheap and sufficiently liquid.

**Core fields used:**

- Sector / issuer / maturity year
- Spread to benchmark
- Liquidity score
- Trade count
- Total trade amount
- Days since last trade

**Core spread calculation:**

`Spread to Benchmark = (Average CUSIP Yield - Benchmark Yield) × 100`

**Important limitation:**

This is a screening tool. It does not replace credit review, call analysis, structure review, tax status review, or PM/trader judgment.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable security screener spread calculations.")
else:
    screen_col1, screen_col2, screen_col3, screen_col4 = st.columns([1, 1, 1, 1])
    with screen_col1:
        screen_sector_options = ["All"]
        if "sector" in market_df.columns:
            screen_sector_options += sorted(market_df["sector"].dropna().astype(str).unique().tolist())
        screen_sector = st.selectbox("Screener Sector", screen_sector_options, index=0, key="screen_sector")
    with screen_col2:
        screen_bucket = st.selectbox(
            "Screener Maturity Year",
            MATURITY_BUCKET_OPTIONS,
            index=0,
            key="screen_bucket",
        )
    with screen_col3:
        screen_rating = st.selectbox(
            "Screener Benchmark",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="screen_rating",
        )
    with screen_col4:
        screen_window = st.selectbox(
            "Screener Lookback",
            ["Latest 30D", "Latest 60D", "Latest 90D", "All"],
            index=2,
            key="screen_window",
        )

    filt_col1, filt_col2, filt_col3, filt_col4 = st.columns([1, 1, 1, 1])
    with filt_col1:
        min_spread = st.number_input("Min Spread (bps)", value=40.0, step=5.0, key="min_screen_spread")
    with filt_col2:
        min_liquidity = st.number_input("Min Liquidity Score", value=50.0, min_value=0.0, max_value=100.0, step=5.0, key="min_screen_liq")
    with filt_col3:
        min_trades_screen = st.number_input("Min Trades", value=2, min_value=1, max_value=100, step=1, key="min_screen_trades")
    with filt_col4:
        min_trade_amount = st.number_input("Min Total Trade Amount", value=0.0, step=100000.0, key="min_screen_amount")

    screener_df = market_df.copy()
    screener_df["trade_date"] = pd.to_datetime(screener_df["trade_date"], errors="coerce").dt.normalize()
    screener_df["yield"] = pd.to_numeric(screener_df["yield"], errors="coerce")
    if "trade_amount" in screener_df.columns:
        screener_df["trade_amount"] = pd.to_numeric(screener_df["trade_amount"], errors="coerce").fillna(0)
    else:
        screener_df["trade_amount"] = 0.0
    if "price" in screener_df.columns:
        screener_df["price"] = pd.to_numeric(screener_df["price"], errors="coerce")
    else:
        screener_df["price"] = pd.NA

    screener_df = screener_df.dropna(subset=["trade_date", "yield", "cusip"])
    if screen_sector != "All" and "sector" in screener_df.columns:
        screener_df = screener_df[screener_df["sector"].astype(str) == str(screen_sector)].copy()
    if screen_bucket != "All" and "maturity_bucket" in screener_df.columns:
        screener_df = screener_df[screener_df["maturity_bucket"] == screen_bucket].copy()

    screen_days = {"Latest 30D": 30, "Latest 60D": 60, "Latest 90D": 90, "All": None}[screen_window]
    if not screener_df.empty and screen_days is not None:
        latest_screen_date = screener_df["trade_date"].max()
        screener_df = screener_df[screener_df["trade_date"] >= latest_screen_date - pd.Timedelta(days=screen_days)].copy()

    if screener_df.empty:
        st.warning("No securities remain after the selected screener universe filters.")
    else:
        screen_summary = (
            screener_df.groupby("cusip", dropna=False)
            .agg(
                issuer=("issuer", "first"),
                sector=("sector", "first") if "sector" in screener_df.columns else ("issuer", "first"),
                maturity_bucket=("maturity_bucket", "first") if "maturity_bucket" in screener_df.columns else ("issuer", "first"),
                maturity=("maturity_bond", "first") if "maturity_bond" in screener_df.columns else ("trade_date", "max"),
                coupon=("coupon_bond", "first") if "coupon_bond" in screener_df.columns else ("yield", "count"),
                call_date=("call_date", "first") if "call_date" in screener_df.columns else ("trade_date", "max"),
                avg_yield=("yield", "mean"),
                latest_yield=("yield", "last"),
                avg_price=("price", "mean"),
                trade_count=("trade_date", "count"),
                latest_trade=("trade_date", "max"),
                first_trade=("trade_date", "min"),
                total_trade_amount=("trade_amount", "sum"),
                avg_trade_amount=("trade_amount", "mean"),
                outstanding_amount=("outstanding_amount", "first") if "outstanding_amount" in screener_df.columns else ("trade_amount", "sum"),
            )
            .reset_index()
        )

        # Latest benchmark by maturity year.
        screen_date_col = _detect_mmd_date_column(mmd_df)
        if screen_date_col is None:
            st.warning("Security screener cannot calculate spreads because the benchmark file has no usable date column.")
        else:
            screen_mmd = mmd_df.copy()
            screen_mmd[screen_date_col] = pd.to_datetime(screen_mmd[screen_date_col], errors="coerce")
            screen_mmd = screen_mmd.dropna(subset=[screen_date_col])
            latest_trade_for_screen = screener_df["trade_date"].max()
            screen_mmd = screen_mmd[screen_mmd[screen_date_col].dt.normalize() <= latest_trade_for_screen].sort_values(screen_date_col)

            if screen_mmd.empty:
                st.warning("No benchmark curve observation was available on or before the latest screener trade date.")
            else:
                screen_latest_mmd = screen_mmd.iloc[[-1]].copy()
                screen_benchmark_date = screen_latest_mmd[screen_date_col].iloc[0]

                bench_rows = []
                for bucket in MATURITY_BUCKET_ORDER:
                    tenor = MMD_BUCKET_MAP.get(bucket, "10Y")
                    y, meta = get_benchmark_curve(screen_latest_mmd, tenor, screen_rating)
                    if y is not None and pd.notna(y.iloc[0]):
                        bench_rows.append(
                            {
                                "maturity_bucket": bucket,
                                "mmd_tenor": tenor,
                                "benchmark_yield": float(y.iloc[0]),
                                "benchmark_source": meta.get("benchmark_source"),
                                "source_column": meta.get("source_column"),
                                "rating_spread_bps": meta.get("rating_spread_bps"),
                            }
                        )
                screen_bench = pd.DataFrame(bench_rows)

                if screen_bench.empty:
                    st.warning("Selected benchmark curve could not be built for screener.")
                else:
                    screen_summary["maturity_bucket"] = screen_summary["maturity_bucket"].astype(str)
                    screen_summary = screen_summary.merge(screen_bench, on="maturity_bucket", how="left")
                    screen_summary["spread_to_benchmark_bps"] = (
                        screen_summary["avg_yield"] - screen_summary["benchmark_yield"]
                    ) * 100

                    # Liquidity score proxy.
                    today_screen = pd.Timestamp.today().normalize()
                    screen_summary["days_since_last_trade"] = (today_screen - screen_summary["latest_trade"]).dt.days
                    screen_summary["liquidity_score"] = (
                        screen_summary["trade_count"].rank(pct=True) * 35
                        + screen_summary["total_trade_amount"].rank(pct=True) * 35
                        + (1 - screen_summary["days_since_last_trade"].rank(pct=True)) * 30
                    )
                    screen_summary["turnover_ratio"] = (
                        screen_summary["total_trade_amount"]
                        / pd.to_numeric(screen_summary["outstanding_amount"], errors="coerce").replace({0: pd.NA})
                    )

                    # Apply screen filters.
                    candidates = screen_summary[
                        (screen_summary["spread_to_benchmark_bps"] >= min_spread)
                        & (screen_summary["liquidity_score"] >= min_liquidity)
                        & (screen_summary["trade_count"] >= min_trades_screen)
                        & (screen_summary["total_trade_amount"] >= min_trade_amount)
                    ].copy()

                    candidates["rv_score"] = (
                        candidates["spread_to_benchmark_bps"].rank(pct=True) * 45
                        + candidates["liquidity_score"].rank(pct=True) * 35
                        + candidates["trade_count"].rank(pct=True) * 10
                        + candidates["total_trade_amount"].rank(pct=True) * 10
                    )

                    candidates = candidates.sort_values("rv_score", ascending=False, na_position="last")

                    s1, s2, s3, s4 = st.columns(4)
                    s1.metric("Candidates", f"{len(candidates):,}")
                    s2.metric("Benchmark", f"{screen_rating}")
                    s3.metric("Benchmark Date", screen_benchmark_date.strftime("%Y-%m-%d"))
                    s4.metric("Universe CUSIPs", f"{len(screen_summary):,}")

                    if candidates.empty:
                        st.info(
                            "No securities met the current screener filters. Try lowering minimum spread, liquidity score, or trade count."
                        )
                    else:
                        top_candidate = candidates.iloc[0]
                        st.info(
                            f"Top candidate by RV score: CUSIP {top_candidate['cusip']} "
                            f"({top_candidate['issuer']}) screens at {top_candidate['spread_to_benchmark_bps']:+.1f} bp "
                            f"to {screen_rating}, liquidity score {top_candidate['liquidity_score']:.1f}, "
                            f"and {int(top_candidate['trade_count'])} trades in the selected window."
                        )

                        display_cols = [
                            "rv_score",
                            "cusip",
                            "issuer",
                            "sector",
                            "maturity_bucket",
                            "maturity",
                            "coupon",
                            "call_date",
                            "avg_yield",
                            "benchmark_yield",
                            "spread_to_benchmark_bps",
                            "liquidity_score",
                            "trade_count",
                            "total_trade_amount",
                            "avg_trade_amount",
                            "days_since_last_trade",
                            "avg_price",
                            "outstanding_amount",
                            "turnover_ratio",
                            "benchmark_source",
                            "source_column",
                        ]
                        display_candidates = candidates[[c for c in display_cols if c in candidates.columns]].copy()
                        for c in ["rv_score", "avg_yield", "benchmark_yield", "spread_to_benchmark_bps", "liquidity_score", "avg_price", "turnover_ratio"]:
                            if c in display_candidates.columns:
                                display_candidates[c] = pd.to_numeric(display_candidates[c], errors="coerce").round(2)

                        st.caption("Showing top 15 candidates. Expand for a larger preview.")
                        safe_dataframe(
                            display_candidates.head(15),
                            width="stretch",
                            hide_index=True,
                            height=420,
                            auto_collapse=False,
                        )
                        if len(display_candidates) > 15:
                            with st.expander(f"View broader candidate table ({len(display_candidates):,} rows)", expanded=False):
                                safe_dataframe(display_candidates, width="stretch", hide_index=True, height=480, max_rows=1000, auto_collapse=False)

                        # Desk-friendly replacement for the old bubble scatter:
                        # show the clearest Top-N ranked opportunities first.
                        candidates_labeled = add_security_label(candidates)
                        screener_fig = ranked_bar_chart(
                            candidates_labeled,
                            value_col="spread_to_benchmark_bps",
                            label_col="security_label",
                            title="Top Cheap / Wide Bonds vs Benchmark",
                            x_title="Spread to Benchmark (bps)",
                            top_n=15,
                            color_col="maturity_bucket",
                            hover_cols=[
                                "issuer", "sector", "avg_yield", "benchmark_yield", "trade_count",
                                "days_since_last_trade", "rv_score", "liquidity_score", "total_trade_amount",
                            ],
                        )
                        if screener_fig is not None:
                            safe_plotly_chart(screener_fig, width="stretch")
                            render_ladder_readthrough(candidates_labeled, "spread_to_benchmark_bps", "security_label", "top cheap / wide bonds vs benchmark")

                        # Secondary read-through table: cheap + liquid / rich / review buckets.
                        quadrant = candidates_labeled.copy()
                        quadrant["desk_signal"] = np.select(
                            [
                                (quadrant["spread_to_benchmark_bps"] >= min_spread) & (quadrant["liquidity_score"] >= min_liquidity),
                                (quadrant["spread_to_benchmark_bps"] >= min_spread) & (quadrant["liquidity_score"] < min_liquidity),
                                (quadrant["spread_to_benchmark_bps"] < 0) & (quadrant["liquidity_score"] >= min_liquidity),
                            ],
                            ["Cheap + Liquid", "Cheap / Needs Liquidity Check", "Rich + Liquid"],
                            default="Review",
                        )
                        st.subheader("Top Opportunity Read-Through")
                        q_cols = [
                            "desk_signal", "security_label", "issuer", "sector", "maturity_bucket",
                            "spread_to_benchmark_bps", "liquidity_score", "trade_count", "total_trade_amount", "rv_score",
                        ]
                        q_display = quadrant[[c for c in q_cols if c in quadrant.columns]].head(15).copy()
                        for c in ["spread_to_benchmark_bps", "liquidity_score", "rv_score"]:
                            if c in q_display.columns:
                                q_display[c] = pd.to_numeric(q_display[c], errors="coerce").round(2)
                        safe_dataframe(q_display, width="stretch", hide_index=True, auto_collapse=False, height=420)

                        csv_candidates = candidates.to_csv(index=False).encode("utf-8")
                        st.download_button(
                            label="Download Top Relative Value Candidates CSV",
                            data=csv_candidates,
                            file_name="top_relative_value_candidates.csv",
                            mime="text/csv",
                        )

                    with st.expander("Screener universe audit table", expanded=False):
                        audit_cols = [
                            "cusip",
                            "issuer",
                            "sector",
                            "maturity_bucket",
                            "avg_yield",
                            "benchmark_yield",
                            "spread_to_benchmark_bps",
                            "liquidity_score",
                            "trade_count",
                            "total_trade_amount",
                            "latest_trade",
                            "benchmark_source",
                            "source_column",
                        ]
                        audit_screen = screen_summary[[c for c in audit_cols if c in screen_summary.columns]].copy()
                        for c in ["avg_yield", "benchmark_yield", "spread_to_benchmark_bps", "liquidity_score"]:
                            if c in audit_screen.columns:
                                audit_screen[c] = pd.to_numeric(audit_screen[c], errors="coerce").round(2)
                        safe_dataframe(audit_screen.head(5000), width="stretch", hide_index=True)




section_anchor("peer-rv", "Peer Relative Value Comparison")
with st.expander("Methodology: peer relative value comparison", expanded=False):
    st.markdown(
        """
This module is intentionally **optional**. It only becomes fully useful when the uploaded dataset contains multiple issuers.

**Purpose:**

- Compare the selected issuer against uploaded peers.
- Prefer same-sector peers when available.
- Allow manual cross-sector comparison when needed, with a warning that sector differences may dominate issuer-level relative value.

**Core calculation:**

`Issuer Spread to Benchmark = (Average Issuer Yield - Benchmark Yield) × 100`

**Default framework:**

- Peer universe = uploaded issuers in the same sector as the selected issuer.
- Benchmark = user-selected AAA / AA / A / BBB curve.
- Window = recent average, because municipal trading can be sparse.
- If fewer than two issuers are uploaded, the module shows an information message instead of failing.
        """
    )

if len(uploaded_issuers) < 2:
    st.info(
        "Peer comparison is unavailable. Upload trade data for at least two issuers "
        "to compare relative value across peers."
    )
elif mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable peer spread comparison.")
else:
    peer_base = market_df.copy()
    peer_base["issuer"] = peer_base["issuer"].astype(str)

    # Determine same-sector peer universe from uploaded data.
    selected_sector_for_peers = selected_sector if selected_sector and selected_sector != "Unknown" else None

    if selected_sector_for_peers and "sector" in peer_base.columns:
        same_sector_issuers = sorted(
            peer_base.loc[
                peer_base["sector"].astype(str) == str(selected_sector_for_peers),
                "issuer",
            ].dropna().astype(str).unique().tolist()
        )
    else:
        same_sector_issuers = []

    peer_col1, peer_col2, peer_col3 = st.columns([1.2, 1, 1])
    with peer_col1:
        peer_mode = st.radio(
            "Peer Universe",
            ["Same-sector uploaded peers", "Manual selection"],
            index=0 if len(same_sector_issuers) >= 2 else 1,
            horizontal=False,
            help="Same-sector comparison is usually cleaner. Manual selection is useful when sector data is missing or when the team wants a custom comp set.",
        )
    with peer_col2:
        peer_rating = st.selectbox(
            "Peer Benchmark Curve",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="peer_benchmark_rating",
        )
    with peer_col3:
        peer_window_label = st.selectbox(
            "Peer Lookback Window",
            ["Latest 30D", "Latest 60D", "Latest 90D", "All"],
            index=0,
            key="peer_lookback_window",
        )

    if peer_mode == "Same-sector uploaded peers":
        if len(same_sector_issuers) < 2:
            st.info(
                "No same-sector peer set with at least two issuers was detected. "
                "Use Manual selection to compare across uploaded issuers, but interpret cross-sector comparisons carefully."
            )
            peer_issuers = [selected_issuer]
        else:
            st.success(
                f"Detected {len(same_sector_issuers):,} uploaded issuer(s) in sector: {selected_sector_for_peers}."
            )
            default_peer_issuers = same_sector_issuers
            peer_issuers = st.multiselect(
                "Same-sector peers",
                same_sector_issuers,
                default=default_peer_issuers,
                key="same_sector_peer_issuers",
            )
    else:
        st.warning(
            "Manual peer selection can compare issuers across sectors. Cross-sector spreads may reflect sector risk, "
            "not just issuer-level relative value."
        )
        default_manual = [selected_issuer]
        # add up to 3 additional issuers by default
        for issuer in uploaded_issuers:
            if issuer != selected_issuer and len(default_manual) < 4:
                default_manual.append(issuer)
        peer_issuers = st.multiselect(
            "Manual peers from uploaded issuers",
            uploaded_issuers,
            default=default_manual,
            key="manual_peer_issuers",
        )

    peer_days = {"Latest 30D": 30, "Latest 60D": 60, "Latest 90D": 90, "All": None}[peer_window_label]

    if len(peer_issuers) < 2:
        st.info("Select at least two issuers to generate peer comparison charts.")
    else:
        peer_work = peer_base[peer_base["issuer"].isin(peer_issuers)].copy()
        peer_work["trade_date"] = pd.to_datetime(peer_work["trade_date"], errors="coerce").dt.normalize()
        peer_work["yield"] = pd.to_numeric(peer_work["yield"], errors="coerce")
        if "trade_amount" in peer_work.columns:
            peer_work["trade_amount"] = pd.to_numeric(peer_work["trade_amount"], errors="coerce").fillna(0)
        else:
            peer_work["trade_amount"] = 0.0

        peer_work = peer_work.dropna(subset=["trade_date", "yield", "issuer", "maturity_bucket"])
        peer_work = peer_work[peer_work["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)].copy()

        if peer_work.empty:
            st.warning("No usable peer trade rows were found after cleaning.")
        else:
            latest_peer_date = peer_work["trade_date"].max()
            if peer_days is not None:
                peer_cutoff = latest_peer_date - pd.Timedelta(days=peer_days)
                peer_work = peer_work[peer_work["trade_date"] >= peer_cutoff].copy()

            if peer_work.empty:
                st.info("No peer observations remain inside the selected lookback window.")
            else:
                # Build issuer/bucket average yields over lookback.
                peer_summary = (
                    peer_work.groupby(["issuer", "maturity_bucket"], as_index=False)
                    .agg(
                        avg_yield=("yield", "mean"),
                        trade_count=("yield", "count"),
                        latest_trade=("trade_date", "max"),
                        total_trade_amount=("trade_amount", "sum"),
                    )
                )

                # Build benchmark curves and use latest benchmark observation at/before latest peer date.
                peer_date_col = _detect_mmd_date_column(mmd_df)
                if peer_date_col is None:
                    st.warning("Peer comparison cannot run because benchmark curve file has no usable date column.")
                else:
                    peer_mmd = mmd_df.copy()
                    peer_mmd[peer_date_col] = pd.to_datetime(peer_mmd[peer_date_col], errors="coerce")
                    peer_mmd = peer_mmd.dropna(subset=[peer_date_col])
                    peer_mmd = peer_mmd[peer_mmd[peer_date_col].dt.normalize() <= latest_peer_date].sort_values(peer_date_col)

                    if peer_mmd.empty:
                        st.warning("No benchmark curve observation was available on or before the latest peer trade date.")
                    else:
                        peer_latest_mmd = peer_mmd.iloc[[-1]].copy()
                        peer_benchmark_date = peer_latest_mmd[peer_date_col].iloc[0]

                        bench_rows = []
                        for bucket in MATURITY_BUCKET_ORDER:
                            tenor = MMD_BUCKET_MAP.get(bucket, "10Y")
                            y, meta = get_benchmark_curve(peer_latest_mmd, tenor, peer_rating)
                            if y is not None and pd.notna(y.iloc[0]):
                                bench_rows.append(
                                    {
                                        "maturity_bucket": bucket,
                                        "mmd_tenor": tenor,
                                        "benchmark_yield": float(y.iloc[0]),
                                        "benchmark_source": meta.get("benchmark_source"),
                                        "source_column": meta.get("source_column"),
                                        "rating_spread_bps": meta.get("rating_spread_bps"),
                                    }
                                )
                        bench_df = pd.DataFrame(bench_rows)

                        if bench_df.empty:
                            st.warning("Selected benchmark curve could not be built for any maturity year.")
                        else:
                            peer_summary = peer_summary.merge(bench_df, on="maturity_bucket", how="left")
                            peer_summary["spread_to_benchmark_bps"] = (
                                peer_summary["avg_yield"] - peer_summary["benchmark_yield"]
                            ) * 100
                            peer_summary = peer_summary.dropna(subset=["spread_to_benchmark_bps"])

                            if peer_summary.empty:
                                st.info("No overlapping peer observations and benchmark tenors were available.")
                            else:
                                maturity_order = MATURITY_BUCKET_ORDER
                                peer_summary["maturity_bucket"] = pd.Categorical(
                                    peer_summary["maturity_bucket"],
                                    categories=maturity_order,
                                    ordered=True,
                                )
                                peer_summary = peer_summary.sort_values(["issuer", "maturity_bucket"])

                                st.subheader("1. Peer Spread Curve Comparison")
                                peer_curve_fig = px.line(
                                    peer_summary,
                                    x="maturity_bucket",
                                    y="spread_to_benchmark_bps",
                                    color="issuer",
                                    markers=True,
                                    hover_data=[
                                        "avg_yield",
                                        "benchmark_yield",
                                        "trade_count",
                                        "total_trade_amount",
                                        "latest_trade",
                                        "benchmark_source",
                                        "source_column",
                                    ],
                                    title=f"Peer Spread Curve Comparison vs {peer_rating} Benchmark",
                                    labels={
                                        "maturity_bucket": "Maturity Year",
                                        "spread_to_benchmark_bps": "Spread to Benchmark (bps)",
                                        "issuer": "Issuer",
                                    },
                                )
                                peer_curve_fig.add_hline(y=0, line_dash="dash", opacity=0.45)
                                peer_curve_fig.update_layout(height=520, hovermode="x unified")
                                safe_plotly_chart(peer_curve_fig, width="stretch")
                                render_ladder_readthrough(peer_summary, "spread_to_benchmark_bps", "issuer", "peer spread curve comparison")

                                st.subheader("2. Peer Spread Ladder")
                                peer_ladder = peer_summary.copy()
                                peer_ladder["spread_to_benchmark_bps"] = pd.to_numeric(peer_ladder["spread_to_benchmark_bps"], errors="coerce")
                                peer_ladder = peer_ladder.dropna(subset=["spread_to_benchmark_bps"]).copy()
                                if peer_ladder.empty:
                                    st.info("No peer spread values available for the ladder view.")
                                else:
                                    peer_ladder["security_bucket"] = peer_ladder["issuer"].astype(str) + " " + peer_ladder["maturity_bucket"].astype(str)
                                    peer_ladder["abs_spread_bps"] = peer_ladder["spread_to_benchmark_bps"].abs()
                                    peer_ladder["maturity_zone"] = peer_ladder["maturity_bucket"].apply(maturity_zone_label)
                                    peer_ladder = peer_ladder.sort_values("abs_spread_bps", ascending=False).head(20)
                                    peer_ladder_fig = px.bar(
                                        peer_ladder.sort_values("spread_to_benchmark_bps"),
                                        x="spread_to_benchmark_bps",
                                        y="security_bucket",
                                        orientation="h",
                                        color="issuer",
                                        hover_data=["maturity_zone", "avg_yield", "benchmark_yield", "trade_count", "total_trade_amount"],
                                        title=f"Largest Peer Spreads vs {peer_rating}",
                                        labels={
                                            "spread_to_benchmark_bps": "Spread to Benchmark (bps)",
                                            "security_bucket": "Issuer / Maturity",
                                            "issuer": "Issuer",
                                        },
                                    )
                                    peer_ladder_fig.add_vline(x=0, line_dash="dash", opacity=0.45)
                                    peer_ladder_fig.update_layout(height=max(420, 28 * len(peer_ladder) + 160), legend_title_text="Issuer")
                                    safe_plotly_chart(peer_ladder_fig, width="stretch")
                                    render_ladder_readthrough(peer_ladder, "spread_to_benchmark_bps", "security_bucket", "peer spread ladder")

                                st.subheader("3. Peer Ranking Table")

                                # -------------------------------------------------------------
                                # Peer ranking methodology
                                # -------------------------------------------------------------
                                # Average spread is a simple arithmetic average across maturity
                                # buckets. Weighted average spread gives more influence to buckets
                                # with larger uploaded trade amount / par volume.
                                #
                                # Spread-to-benchmark is calculated as:
                                #   (Issuer average trade yield - benchmark yield) * 100
                                # where yields are in percentage-point terms and the output is bps.
                                # -------------------------------------------------------------

                                def _weighted_avg_spread_by_amount(group: pd.DataFrame) -> float:
                                    spreads = pd.to_numeric(group["spread_to_benchmark_bps"], errors="coerce")
                                    weights = pd.to_numeric(group["total_trade_amount"], errors="coerce").fillna(0)
                                    valid = spreads.notna() & weights.gt(0)
                                    if valid.any():
                                        return float((spreads[valid] * weights[valid]).sum() / weights[valid].sum())
                                    return float(spreads.mean())

                                weighted_spread = (
                                    peer_summary.groupby("issuer")
                                    .apply(_weighted_avg_spread_by_amount)
                                    .reset_index(name="weighted_avg_spread_bps")
                                )

                                weight_basis = (
                                    peer_summary.groupby("issuer", as_index=False)
                                    .agg(weighted_trade_amount_used=("total_trade_amount", "sum"))
                                )
                                weight_basis["weighting_method"] = weight_basis["weighted_trade_amount_used"].apply(
                                    lambda x: "Trade-amount weighted" if pd.notna(x) and x > 0 else "Unweighted fallback"
                                )

                                peer_rank = (
                                    peer_summary.groupby("issuer", as_index=False)
                                    .agg(
                                        avg_spread_bps=("spread_to_benchmark_bps", "mean"),
                                        max_spread_bps=("spread_to_benchmark_bps", "max"),
                                        trade_count=("trade_count", "sum"),
                                        total_trade_amount=("total_trade_amount", "sum"),
                                        latest_trade=("latest_trade", "max"),
                                    )
                                    .merge(weighted_spread, on="issuer", how="left")
                                    .merge(weight_basis[["issuer", "weighting_method"]], on="issuer", how="left")
                                )

                                peer_rank["rank_by_weighted_spread"] = (
                                    peer_rank["weighted_avg_spread_bps"]
                                    .rank(method="dense", ascending=False)
                                    .astype(int)
                                )
                                peer_rank["rank_by_avg_spread"] = (
                                    peer_rank["avg_spread_bps"]
                                    .rank(method="dense", ascending=False)
                                    .astype(int)
                                )
                                peer_rank = peer_rank.sort_values(
                                    ["rank_by_weighted_spread", "rank_by_avg_spread", "issuer"],
                                    ascending=[True, True, True],
                                )

                                with st.expander("Peer ranking methodology", expanded=False):
                                    st.markdown(f"""
### How the Peer Ranking Is Calculated

This table ranks issuers by **weighted average spread-to-benchmark**, measured in basis points.

**Spread-to-benchmark formula:**

`Spread (bps) = (Issuer Average Trade Yield - Benchmark Yield) × 100`

**Weighted average spread formula:**

`Weighted Avg Spread = Σ(Spread × Trade Amount) / Σ(Trade Amount)`

### Current Ranking Logic

- Benchmark rating selected: **{peer_rating}**
- Lookback window selected: **{peer_window_label}**
- Maturity years included: selected peer maturity scope shown in this section
- Issuers are ranked from **widest** to **tightest** based on `weighted_avg_spread_bps`
- If trade amount is unavailable or zero, the dashboard falls back to the simple average spread

### Interpretation

- **Positive spread**: issuer is trading wider than the selected benchmark
- **Negative spread**: issuer is trading tighter than the selected benchmark
- **Higher rank**: wider screened spread, potentially more yield compensation versus benchmark

### Important Caveats

This is a **relative-value screening tool**, not a standalone trading recommendation. Rankings can be affected by:

- Trade size and odd-lot effects
- Liquidity differences across CUSIPs
- Callable structures and coupon differences
- Credit quality differences not fully captured by benchmark rating
- Limited observations in the selected time window
- Stale or uneven trading activity across issuers

Use this ranking as a first-pass screen, then review the CUSIP-level drilldown and liquidity metrics before drawing trading conclusions.
""")

                                safe_dataframe(
                                    peer_rank[
                                        [
                                            "rank_by_weighted_spread",
                                            "rank_by_avg_spread",
                                            "issuer",
                                            "weighted_avg_spread_bps",
                                            "avg_spread_bps",
                                            "max_spread_bps",
                                            "trade_count",
                                            "total_trade_amount",
                                            "weighting_method",
                                            "latest_trade",
                                        ]
                                    ],
                                    width="stretch",
                                    hide_index=True,
                                )

                                # Read-through
                                if not peer_rank.empty:
                                    widest = peer_rank.iloc[0]
                                    tightest = peer_rank.iloc[-1]
                                    st.info(
                                        f"Peer read-through: {widest['issuer']} screens widest on a trade-amount weighted basis "
                                        f"at {widest['weighted_avg_spread_bps']:+.1f} bp versus {peer_rating}, while "
                                        f"{tightest['issuer']} screens tightest at {tightest['weighted_avg_spread_bps']:+.1f} bp. "
                                        f"The simple average spread is also shown for comparison. Use this as a screening signal "
                                        f"and review CUSIP-level liquidity before drawing trading conclusions."
                                    )

                                with st.expander("Peer comparison audit table", expanded=False):
                                    audit_cols = [
                                        "issuer",
                                        "maturity_bucket",
                                        "avg_yield",
                                        "benchmark_yield",
                                        "spread_to_benchmark_bps",
                                        "trade_count",
                                        "total_trade_amount",
                                        "latest_trade",
                                        "mmd_tenor",
                                        "benchmark_source",
                                        "source_column",
                                        "rating_spread_bps",
                                    ]
                                    audit_df = peer_summary[[c for c in audit_cols if c in peer_summary.columns]].copy()
                                    for c in ["avg_yield", "benchmark_yield", "spread_to_benchmark_bps", "rating_spread_bps"]:
                                        if c in audit_df.columns:
                                            audit_df[c] = pd.to_numeric(audit_df[c], errors="coerce").round(2)
                                    st.caption(
                                        f"Benchmark date used: {peer_benchmark_date.strftime('%Y-%m-%d')}. "
                                        f"Lookback window: {peer_window_label}."
                                    )
                                    safe_dataframe(audit_df, width="stretch", hide_index=True)




section_anchor("cross-issuer-rv", "Cross-Issuer Relative Value Analytics")
with st.expander("Methodology: cross-issuer relative value analytics", expanded=False):
    st.markdown(
        """
This section upgrades peer comparison from **visual comparison** into a systematic **issuer-bucket ranking framework**.

**Core purpose:**

- Identify which issuer / maturity year screens cheap or rich versus the uploaded peer group.
- Convert peer spreads into **peer gaps**, **z-scores**, and **relative value scores**.
- Keep the module optional: it only becomes meaningful when at least two issuers are uploaded.

**Core calculations:**

`Issuer Spread = (Average Issuer Yield - Benchmark Yield) × 100`

`Peer Gap = Issuer Spread - Peer Median Spread within the same maturity year`

`Peer Z-Score = (Issuer Spread - Bucket Mean Spread) / Bucket Spread Std`

**Simplified RV Score:**

`RV Score = 45% Spread Percentile + 35% Liquidity Percentile + 20% Trade Activity Percentile`

**How to read it:**

- Positive peer gap: issuer/bucket is wider than peers, potentially cheaper.
- Negative peer gap: issuer/bucket is tighter than peers, potentially richer.
- Higher RV score: screens cheaper while retaining better liquidity/trading support.
- This is a screening tool, not a final trade recommendation.
        """
    )

if len(uploaded_issuers) < 2:
    st.info(
        "Cross-issuer RV analytics is unavailable. Upload at least two issuers to compare issuer-bucket relative value."
    )
elif mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable cross-issuer spread analytics.")
else:
    xrv_col1, xrv_col2, xrv_col3, xrv_col4 = st.columns([1.2, 1, 1, 1])
    with xrv_col1:
        xrv_mode = st.radio(
            "Cross-Issuer Universe",
            ["Same-sector uploaded issuers", "Manual issuer set"],
            index=0 if "sector" in market_df.columns and selected_sector != "Unknown" else 1,
            key="xrv_universe_mode",
            help="Same-sector is cleaner. Manual set is useful when sector data is missing or when the desk wants a custom comp set.",
        )
    with xrv_col2:
        xrv_rating = st.selectbox(
            "X-Issuer Benchmark",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="xrv_benchmark",
        )
    with xrv_col3:
        xrv_window = st.selectbox(
            "X-Issuer Lookback",
            ["Latest 30D", "Latest 60D", "Latest 90D", "All"],
            index=1,
            key="xrv_window",
        )
    with xrv_col4:
        xrv_min_trades = st.number_input(
            "Min Trades / Bucket",
            min_value=1,
            max_value=100,
            value=1,
            step=1,
            key="xrv_min_trades",
        )

    if xrv_mode == "Same-sector uploaded issuers":
        if "sector" not in market_df.columns or selected_sector == "Unknown":
            st.info("Same-sector universe is unavailable because sector data is missing or unknown. Use Manual issuer set.")
            xrv_issuers = [selected_issuer]
        else:
            xrv_issuers = sorted(
                market_df.loc[
                    market_df["sector"].astype(str) == str(selected_sector),
                    "issuer",
                ].dropna().astype(str).unique().tolist()
            )
            if len(xrv_issuers) < 2:
                st.info(
                    "Fewer than two same-sector issuers were detected. Use Manual issuer set or upload more peer data."
                )
    else:
        default_xrv = [selected_issuer]
        for issuer in uploaded_issuers:
            if issuer != selected_issuer and len(default_xrv) < 5:
                default_xrv.append(issuer)
        xrv_issuers = st.multiselect(
            "Manual cross-issuer set",
            uploaded_issuers,
            default=default_xrv,
            key="xrv_manual_issuers",
        )
        st.warning(
            "Manual cross-issuer comparison may include different sectors. Interpret peer gaps carefully because sector risk can dominate issuer-level RV."
        )

    if len(xrv_issuers) < 2:
        st.info("Select or upload at least two issuers to generate cross-issuer RV analytics.")
    else:
        xrv_days = {"Latest 30D": 30, "Latest 60D": 60, "Latest 90D": 90, "All": None}[xrv_window]

        xrv_df = market_df[market_df["issuer"].astype(str).isin([str(x) for x in xrv_issuers])].copy()
        xrv_df["trade_date"] = pd.to_datetime(xrv_df["trade_date"], errors="coerce").dt.normalize()
        xrv_df["yield"] = pd.to_numeric(xrv_df["yield"], errors="coerce")
        if "trade_amount" in xrv_df.columns:
            xrv_df["trade_amount"] = pd.to_numeric(xrv_df["trade_amount"], errors="coerce").fillna(0)
        else:
            xrv_df["trade_amount"] = 0.0

        xrv_df = xrv_df.dropna(subset=["trade_date", "yield", "issuer", "maturity_bucket"])
        xrv_df = xrv_df[xrv_df["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)].copy()

        if xrv_df.empty:
            st.warning("No usable cross-issuer observations remained after cleaning.")
        else:
            xrv_latest_date = xrv_df["trade_date"].max()
            if xrv_days is not None:
                xrv_cutoff = xrv_latest_date - pd.Timedelta(days=xrv_days)
                xrv_df = xrv_df[xrv_df["trade_date"] >= xrv_cutoff].copy()

            if xrv_df.empty:
                st.info("No cross-issuer observations remain inside the selected lookback window.")
            else:
                xrv_summary = (
                    xrv_df.groupby(["issuer", "maturity_bucket"], as_index=False)
                    .agg(
                        avg_yield=("yield", "mean"),
                        trade_count=("yield", "count"),
                        total_trade_amount=("trade_amount", "sum"),
                        latest_trade=("trade_date", "max"),
                        sector=("sector", "first") if "sector" in xrv_df.columns else ("issuer", "first"),
                    )
                )
                xrv_summary = xrv_summary[xrv_summary["trade_count"] >= xrv_min_trades].copy()

                if xrv_summary.empty:
                    st.info("No issuer-bucket observations met the minimum trade filter.")
                else:
                    date_col = _detect_mmd_date_column(mmd_df)
                    if date_col is None:
                        st.warning("Cross-issuer RV cannot run because the benchmark file has no usable date column.")
                    else:
                        xrv_mmd = mmd_df.copy()
                        xrv_mmd[date_col] = pd.to_datetime(xrv_mmd[date_col], errors="coerce")
                        xrv_mmd = xrv_mmd.dropna(subset=[date_col])
                        xrv_mmd = xrv_mmd[xrv_mmd[date_col].dt.normalize() <= xrv_latest_date].sort_values(date_col)

                        if xrv_mmd.empty:
                            st.warning("No benchmark curve observation was available on or before the latest cross-issuer trade date.")
                        else:
                            xrv_latest_mmd = xrv_mmd.iloc[[-1]].copy()
                            xrv_benchmark_date = xrv_latest_mmd[date_col].iloc[0]

                            bench_rows = []
                            for bucket in MATURITY_BUCKET_ORDER:
                                tenor = MMD_BUCKET_MAP.get(bucket, "10Y")
                                y, meta = get_benchmark_curve(xrv_latest_mmd, tenor, xrv_rating)
                                if y is not None and pd.notna(y.iloc[0]):
                                    bench_rows.append(
                                        {
                                            "maturity_bucket": bucket,
                                            "mmd_tenor": tenor,
                                            "benchmark_yield": float(y.iloc[0]),
                                            "benchmark_source": meta.get("benchmark_source"),
                                            "source_column": meta.get("source_column"),
                                            "rating_spread_bps": meta.get("rating_spread_bps"),
                                        }
                                    )
                            xrv_bench = pd.DataFrame(bench_rows)

                            if xrv_bench.empty:
                                st.warning("Selected benchmark curve could not be built for cross-issuer RV.")
                            else:
                                xrv_summary = xrv_summary.merge(xrv_bench, on="maturity_bucket", how="left")
                                xrv_summary["spread_to_benchmark_bps"] = (
                                    xrv_summary["avg_yield"] - xrv_summary["benchmark_yield"]
                                ) * 100
                                xrv_summary = xrv_summary.dropna(subset=["spread_to_benchmark_bps"])

                                if xrv_summary.empty:
                                    st.info("No issuer-bucket observations had usable benchmark spreads.")
                                else:
                                    # Peer-relative metrics by maturity year.
                                    xrv_summary["bucket_peer_median_bps"] = (
                                        xrv_summary.groupby("maturity_bucket")["spread_to_benchmark_bps"].transform("median")
                                    )
                                    xrv_summary["bucket_peer_mean_bps"] = (
                                        xrv_summary.groupby("maturity_bucket")["spread_to_benchmark_bps"].transform("mean")
                                    )
                                    xrv_summary["bucket_peer_std_bps"] = (
                                        xrv_summary.groupby("maturity_bucket")["spread_to_benchmark_bps"].transform("std")
                                    )
                                    xrv_summary["peer_gap_bps"] = (
                                        xrv_summary["spread_to_benchmark_bps"] - xrv_summary["bucket_peer_median_bps"]
                                    )
                                    xrv_summary["peer_z_score"] = (
                                        (xrv_summary["spread_to_benchmark_bps"] - xrv_summary["bucket_peer_mean_bps"])
                                        / xrv_summary["bucket_peer_std_bps"].replace({0: pd.NA})
                                    )

                                    # Liquidity proxy within universe.
                                    today_xrv = pd.Timestamp.today().normalize()
                                    xrv_summary["days_since_last_trade"] = (today_xrv - xrv_summary["latest_trade"]).dt.days
                                    xrv_summary["liquidity_score"] = (
                                        xrv_summary["trade_count"].rank(pct=True) * 35
                                        + xrv_summary["total_trade_amount"].rank(pct=True) * 35
                                        + (1 - xrv_summary["days_since_last_trade"].rank(pct=True)) * 30
                                    )

                                    xrv_summary["spread_percentile"] = (
                                        xrv_summary.groupby("maturity_bucket")["spread_to_benchmark_bps"]
                                        .rank(pct=True)
                                    )
                                    xrv_summary["liquidity_percentile"] = xrv_summary["liquidity_score"].rank(pct=True)
                                    xrv_summary["trade_activity_percentile"] = xrv_summary["trade_count"].rank(pct=True)
                                    xrv_summary["x_issuer_rv_score"] = (
                                        xrv_summary["spread_percentile"] * 45
                                        + xrv_summary["liquidity_percentile"] * 35
                                        + xrv_summary["trade_activity_percentile"] * 20
                                    )

                                    def classify_xrv(row):
                                        if row["peer_gap_bps"] >= 10 and row["liquidity_score"] >= 60:
                                            return "Cheap + Liquid"
                                        if row["peer_gap_bps"] >= 10:
                                            return "Cheap / Needs Liquidity Check"
                                        if row["peer_gap_bps"] <= -10:
                                            return "Rich vs Peers"
                                        return "In Line"

                                    xrv_summary["x_issuer_signal"] = xrv_summary.apply(classify_xrv, axis=1)

                                    st.subheader("1. Peer Gap Ladder")
                                    gap_ladder = xrv_summary.copy()
                                    gap_ladder["peer_gap_bps"] = pd.to_numeric(gap_ladder["peer_gap_bps"], errors="coerce")
                                    gap_ladder = gap_ladder.dropna(subset=["peer_gap_bps"]).copy()
                                    if gap_ladder.empty:
                                        st.info("No peer-gap values available for the ladder view.")
                                    else:
                                        gap_ladder["issuer_bucket"] = gap_ladder["issuer"].astype(str) + " " + gap_ladder["maturity_bucket"].astype(str)
                                        gap_ladder["abs_peer_gap_bps"] = gap_ladder["peer_gap_bps"].abs()
                                        gap_ladder["maturity_zone"] = gap_ladder["maturity_bucket"].apply(maturity_zone_label)
                                        gap_ladder = gap_ladder.sort_values("abs_peer_gap_bps", ascending=False).head(20)
                                        gap_fig = px.bar(
                                            gap_ladder.sort_values("peer_gap_bps"),
                                            x="peer_gap_bps",
                                            y="issuer_bucket",
                                            orientation="h",
                                            color="issuer",
                                            hover_data=["maturity_zone", "spread_to_benchmark_bps", "bucket_peer_median_bps", "liquidity_score"],
                                            title=f"Largest Peer Gaps vs {xrv_rating} Benchmark",
                                            labels={
                                                "peer_gap_bps": "Peer Gap (bps)",
                                                "issuer_bucket": "Issuer / Maturity",
                                                "issuer": "Issuer",
                                            },
                                        )
                                        gap_fig.add_vline(x=0, line_dash="dash", opacity=0.45)
                                        gap_fig.update_layout(height=max(420, 28 * len(gap_ladder) + 160), legend_title_text="Issuer")
                                        safe_plotly_chart(gap_fig, width="stretch")
                                        render_ladder_readthrough(gap_ladder, "peer_gap_bps", "issuer_bucket", "peer gap ladder")

                                    st.subheader("2. Cross-Issuer RV Ranking")
                                    ranking = xrv_summary.sort_values("x_issuer_rv_score", ascending=False, na_position="last").copy()
                                    rank_cols = [
                                        "issuer",
                                        "sector",
                                        "maturity_bucket",
                                        "spread_to_benchmark_bps",
                                        "bucket_peer_median_bps",
                                        "peer_gap_bps",
                                        "peer_z_score",
                                        "liquidity_score",
                                        "trade_count",
                                        "total_trade_amount",
                                        "days_since_last_trade",
                                        "x_issuer_rv_score",
                                        "x_issuer_signal",
                                    ]
                                    ranking_display = ranking[[c for c in rank_cols if c in ranking.columns]].copy()
                                    for c in [
                                        "spread_to_benchmark_bps",
                                        "bucket_peer_median_bps",
                                        "peer_gap_bps",
                                        "peer_z_score",
                                        "liquidity_score",
                                        "x_issuer_rv_score",
                                    ]:
                                        if c in ranking_display.columns:
                                            ranking_display[c] = pd.to_numeric(ranking_display[c], errors="coerce").round(2)

                                    ranking_display_top = ranking_display.head(15).copy()
                                    st.caption("Showing top 15 cross-issuer RV candidates. Expand the table for a larger preview.")
                                    safe_dataframe(
                                        ranking_display_top,
                                        width="stretch",
                                        hide_index=True,
                                        height=430,
                                        auto_collapse=False,
                                    )
                                    if len(ranking_display) > len(ranking_display_top):
                                        with st.expander(f"View broader Cross-Issuer RV Ranking ({len(ranking_display):,} rows)", expanded=False):
                                            safe_dataframe(
                                                ranking_display,
                                                width="stretch",
                                                hide_index=True,
                                                height=430,
                                                max_rows=min(MAX_TABLE_ROWS, 1000),
                                                auto_collapse=False,
                                            )

                                    st.subheader("3. Cross-Issuer Opportunity Ranking")
                                    # Bubble maps become unreadable with many issuer/maturity combinations.
                                    # Use a ranked bar chart plus a decision table instead.
                                    ranking_labeled = ranking.copy()
                                    ranking_labeled["issuer_bucket"] = ranking_labeled["issuer"].astype(str) + " " + ranking_labeled["maturity_bucket"].astype(str)
                                    xrv_bar = ranked_bar_chart(
                                        ranking_labeled,
                                        value_col="peer_gap_bps",
                                        label_col="issuer_bucket",
                                        title="Top Cross-Issuer Cheapness vs Peer Median",
                                        x_title="Peer Gap (bps)",
                                        top_n=18,
                                        color_col="x_issuer_signal",
                                        hover_cols=["issuer", "sector", "maturity_bucket", "spread_to_benchmark_bps", "bucket_peer_median_bps", "liquidity_score", "trade_count", "x_issuer_rv_score"],
                                    )
                                    if xrv_bar is not None:
                                        safe_plotly_chart(xrv_bar, width="stretch")
                                        render_ladder_readthrough(ranking_labeled, "peer_gap_bps", "issuer_bucket", "cross-issuer opportunity ranking")

                                    st.subheader("Cross-Issuer Decision Table")
                                    decision_cols = [
                                        "x_issuer_signal", "issuer", "sector", "maturity_bucket", "peer_gap_bps",
                                        "spread_to_benchmark_bps", "bucket_peer_median_bps", "liquidity_score",
                                        "trade_count", "total_trade_amount", "x_issuer_rv_score",
                                    ]
                                    decision_display = ranking_labeled[[c for c in decision_cols if c in ranking_labeled.columns]].head(20).copy()
                                    for c in ["peer_gap_bps", "spread_to_benchmark_bps", "bucket_peer_median_bps", "liquidity_score", "x_issuer_rv_score"]:
                                        if c in decision_display.columns:
                                            decision_display[c] = pd.to_numeric(decision_display[c], errors="coerce").round(2)
                                    safe_dataframe(decision_display, width="stretch", hide_index=True, auto_collapse=False, height=460)

                                    if not ranking.empty:
                                        top = ranking.iloc[0]
                                        st.info(
                                            f"Cross-issuer read-through: {top['issuer']} / {top['maturity_bucket']} screens highest by RV score. "
                                            f"It is {top['peer_gap_bps']:+.1f} bp versus the bucket peer median, "
                                            f"with liquidity score {top['liquidity_score']:.1f} and signal: {top['x_issuer_signal']}."
                                        )

                                    with st.expander("Cross-issuer RV audit table", expanded=False):
                                        audit_cols = [
                                            "issuer",
                                            "sector",
                                            "maturity_bucket",
                                            "avg_yield",
                                            "benchmark_yield",
                                            "spread_to_benchmark_bps",
                                            "bucket_peer_median_bps",
                                            "peer_gap_bps",
                                            "peer_z_score",
                                            "trade_count",
                                            "total_trade_amount",
                                            "latest_trade",
                                            "mmd_tenor",
                                            "benchmark_source",
                                            "source_column",
                                            "rating_spread_bps",
                                        ]
                                        audit_xrv = xrv_summary[[c for c in audit_cols if c in xrv_summary.columns]].copy()
                                        for c in [
                                            "avg_yield",
                                            "benchmark_yield",
                                            "spread_to_benchmark_bps",
                                            "bucket_peer_median_bps",
                                            "peer_gap_bps",
                                            "peer_z_score",
                                            "rating_spread_bps",
                                        ]:
                                            if c in audit_xrv.columns:
                                                audit_xrv[c] = pd.to_numeric(audit_xrv[c], errors="coerce").round(2)
                                        st.caption(
                                            f"Benchmark date used: {xrv_benchmark_date.strftime('%Y-%m-%d')}. "
                                            f"Lookback window: {xrv_window}."
                                        )
                                        safe_dataframe(audit_xrv, width="stretch", hide_index=True)


section_anchor("spread-level", "Current Spread Level Framework")
with st.expander("Methodology: current spread level", expanded=False):
    st.markdown(
        """
This section shows where the selected issuer is trading **now** versus transparent benchmark curves.

**Calculation:**

`Current Spread Level = (Average Issuer Trade Yield - Benchmark Yield) × 100`

Where:

`Benchmark Yield = uploaded rating curve if available; otherwise MMD/AAA Curve + Rating Spread Assumption`

**How to read it:**

- **Positive spread**: issuer yield is above the selected benchmark curve; the issuer/bucket is trading cheaper than that benchmark.
- **Negative spread**: issuer yield is below the selected benchmark curve; the issuer/bucket is trading richer than that benchmark.
- Rows are maturity years. Columns are benchmark curves.
- This is a **level** view, not a movement view. Level answers: *is it cheap or rich right now?* Movement answers: *did it widen or tighten recently?*
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD curve file to enable current spread level analytics.")
else:
    level_col1, level_col2 = st.columns([1, 2])
    with level_col1:
        level_ratings = st.multiselect(
            "Spread Level Benchmark Curves",
            BENCHMARK_RATINGS,
            default=[r for r in ["AAA", "AA", "A", "BBB"] if r in BENCHMARK_RATINGS],
            help="Priority: uploaded rating curve columns first; otherwise MMD/AAA plus the visible rating-spread assumptions.",
        )
    with level_col2:
        st.caption(
            "Cells show latest available issuer spread to each benchmark curve, in basis points. "
            "Higher positive values generally indicate cheaper relative value versus that benchmark."
        )

    if not level_ratings:
        st.info("Select at least one benchmark curve to display current spread levels.")
    else:
        level_table, level_audit = build_spread_level_data(
            market_df=market_df,
            mmd_df=mmd_df,
            issuer=selected_issuer,
            ratings=level_ratings,
        )
        if level_table.isna().all().all():
            st.warning(
                "No overlapping issuer trade dates and benchmark dates were found for current spread levels. "
                "Check that the curve file has a Date column plus either 5Y/10Y/20Y/30Y base columns or explicit rating curve columns such as AA_10Y, and that trade dates overlap with the curve history."
            )
        else:
            level_table = aggregate_maturity_rows_for_display(level_table) if len(level_table.index) > 10 else level_table
            level_text = level_table.map(lambda x: "" if pd.isna(x) else f"{x:+.1f} bp")

            # 1) Spread level curve: one line per selected benchmark rating.
            # Defensive schema handling: level_table may have maturity labels as the index
            # or under maturity_bucket/maturity_year depending on earlier transformations.
            curve_df = level_table.copy()
            curve_long = safe_melt_by_maturity(
                curve_df,
                value_name="spread_to_benchmark_bps",
                var_name="benchmark_rating",
            )

            st.subheader("1. Current Spread Curve")
            curve_long = sanitize_curve_long_for_plot(
                curve_long,
                x_col="maturity_bucket",
                y_col="spread_to_benchmark_bps",
                color_col="benchmark_rating",
            )
            if curve_long.empty:
                st.warning(
                    "Current Spread Curve skipped safely: no valid maturity + spread rows were available. "
                    "This usually means benchmark dates did not overlap with issuer trades, or spread values were all blank after the benchmark merge."
                )
                with st.expander("Debug current spread curve data", expanded=False):
                    safe_dataframe(curve_data_audit(curve_df), width="stretch")
            else:
                level_curve_fig = px.line(
                    curve_long,
                    x="maturity_bucket",
                    y="spread_to_benchmark_bps",
                    color="benchmark_rating",
                    markers=True,
                    title=f"{selected_issuer} Current Spread Curve vs Selected Benchmarks",
                    labels={
                        "maturity_bucket": "Maturity Year / Zone",
                        "spread_to_benchmark_bps": "Spread to Benchmark (bps)",
                        "benchmark_rating": "Benchmark Curve",
                    },
                )
                level_curve_fig.add_hline(y=0, line_dash="dash", opacity=0.5)
                level_curve_fig.update_layout(hovermode="x unified")
                safe_plotly_chart(level_curve_fig, width="stretch")
                render_ladder_readthrough(curve_long, "spread_to_benchmark_bps", "maturity_bucket", "current spread curve")

            # 2) Spread level ladder: maturity year x benchmark rating, shown as ranked bars instead of a ladder.
            st.subheader("2. Current Spread Level Ladder")
            level_ladder = safe_melt_by_maturity(
                level_table.reset_index().rename(columns={level_table.index.name or "index": "maturity_bucket"}),
                value_vars=[c for c in level_table.columns if c in level_ratings],
                id_vars="maturity_bucket",
                var_name="benchmark_rating",
                value_name="spread_to_benchmark_bps",
            )
            if level_ladder.empty:
                st.info("No current spread level values available for the ladder view.")
            else:
                level_ladder["spread_to_benchmark_bps"] = pd.to_numeric(level_ladder["spread_to_benchmark_bps"], errors="coerce")
                level_ladder = level_ladder.dropna(subset=["spread_to_benchmark_bps"]).copy()
                level_ladder["security_bucket"] = level_ladder["benchmark_rating"].astype(str) + " / " + level_ladder["maturity_bucket"].astype(str)
                level_ladder["abs_spread_bps"] = level_ladder["spread_to_benchmark_bps"].abs()
                level_ladder["maturity_zone"] = level_ladder["maturity_bucket"].apply(maturity_zone_label)
                level_ladder = level_ladder.sort_values("abs_spread_bps", ascending=False).head(20)
                level_ladder_fig = px.bar(
                    level_ladder.sort_values("spread_to_benchmark_bps"),
                    x="spread_to_benchmark_bps",
                    y="security_bucket",
                    orientation="h",
                    color="benchmark_rating",
                    hover_data=["maturity_zone", "abs_spread_bps"],
                    title=f"{selected_issuer} Largest Current Spread Levels",
                    labels={
                        "spread_to_benchmark_bps": "Current Spread (bps)",
                        "security_bucket": "Benchmark / Maturity",
                        "benchmark_rating": "Benchmark Curve",
                    },
                )
                level_ladder_fig.add_vline(x=0, line_dash="dash", opacity=0.45)
                level_ladder_fig.update_layout(height=max(420, 28 * len(level_ladder) + 160), legend_title_text="Benchmark")
                safe_plotly_chart(level_ladder_fig, width="stretch")
                render_ladder_readthrough(level_ladder, "spread_to_benchmark_bps", "security_bucket", "current spread level ladder")

            # 3) Quick signal: identify the cheapest bucket vs the first selected benchmark.
            primary_rating = level_ratings[0]
            if primary_rating in level_table.columns and level_table[primary_rating].notna().any():
                cheapest_bucket = level_table[primary_rating].astype(float).idxmax()
                cheapest_spread = level_table.loc[cheapest_bucket, primary_rating]
                richest_bucket = level_table[primary_rating].astype(float).idxmin()
                richest_spread = level_table.loc[richest_bucket, primary_rating]
                st.info(
                    f"Relative value read-through vs {primary_rating}: "
                    f"{cheapest_bucket} appears cheapest at {cheapest_spread:+.1f} bp, "
                    f"while {richest_bucket} appears richest at {richest_spread:+.1f} bp."
                )

            with st.expander("Current spread level audit table", expanded=False):
                display_cols = [
                    "maturity_bucket", "benchmark_rating", "latest_date", "avg_yield", "benchmark_yield",
                    "spread_to_benchmark_bps", "mmd_tenor", "benchmark_source", "source_column", "rating_spread_bps", "trade_count",
                    "total_trade_amount", "note",
                ]
                audit_display = level_audit[[c for c in display_cols if c in level_audit.columns]].copy()
                for c in ["avg_yield", "benchmark_yield", "spread_to_benchmark_bps", "rating_spread_bps"]:
                    if c in audit_display.columns:
                        audit_display[c] = pd.to_numeric(audit_display[c], errors="coerce").round(2)
                safe_dataframe(audit_display, width="stretch", hide_index=True)


section_anchor("spread-attribution", "Spread Attribution Waterfall")
with st.expander("Methodology: spread attribution waterfall", expanded=False):
    st.markdown(
        """
This section decomposes the selected issuer's spread versus the **AAA/MMD curve** into transparent, reviewable components.

**Framework:**

`Issuer Spread vs AAA = Rating Premium + Liquidity Premium + Callable Adjustment + Residual / Issuer-Specific Premium`

**Important notes:**

- This is a **modeled attribution**, not a vendor curve or investment recommendation.
- **Rating Premium** uses the visible maturity-adjusted rating-spread assumptions already shown in the benchmark methodology.
- **Liquidity Premium** is estimated from the issuer/bucket liquidity score. Less liquid buckets receive a larger modeled premium.
- **Callable Adjustment** is a simple proxy based on whether bonds in the bucket appear callable.
- **Residual / Issuer-Specific Premium** captures what remains after the modeled components. This can reflect credit, sector, structure, supply/demand, data noise, or model misspecification.
- The purpose is pitchbook-style explanation and screening, not final pricing.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD curve file to enable spread attribution waterfall analytics.")
else:
    wf_col1, wf_col2, wf_col3 = st.columns([1, 1, 1])
    with wf_col1:
        wf_bucket = st.selectbox(
            "Waterfall Maturity Year",
            MATURITY_BUCKET_ORDER,
            index=1,
            help="The issuer spread will be attributed for this maturity year.",
        )
    with wf_col2:
        wf_rating = st.selectbox(
            "Modeled Rating Premium",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AA") if "AA" in BENCHMARK_RATINGS else 0,
            help="Used only for the rating-premium component. The total spread is measured versus AAA/MMD.",
        )
    with wf_col3:
        wf_lookback_days = st.selectbox(
            "Issuer Yield Lookback",
            [7, 30, 60, 90, 180],
            index=1,
            format_func=lambda x: f"Latest {x} days",
            help="Average issuer yield is calculated from trades in this lookback window to reduce muni trading noise.",
        )

    wf_tenor = MMD_BUCKET_MAP.get(wf_bucket, "10Y")
    date_col = _detect_mmd_date_column(mmd_df)

    if date_col is None:
        st.warning("Waterfall cannot run because the MMD/curve file does not contain a usable date column.")
    else:
        issuer_bucket_trades = market_df[
            (market_df["issuer"] == selected_issuer)
            & (market_df["maturity_bucket"] == wf_bucket)
        ].copy()

        if issuer_bucket_trades.empty:
            st.warning(f"No {wf_bucket} trade rows were found for {selected_issuer}.")
        else:
            issuer_bucket_trades["trade_date"] = pd.to_datetime(issuer_bucket_trades["trade_date"], errors="coerce")
            issuer_bucket_trades["yield"] = pd.to_numeric(issuer_bucket_trades["yield"], errors="coerce")
            issuer_bucket_trades = issuer_bucket_trades.dropna(subset=["trade_date", "yield"])

            if issuer_bucket_trades.empty:
                st.warning("Waterfall cannot run because no valid trade dates/yields remain after cleaning.")
            else:
                wf_latest_trade_date = issuer_bucket_trades["trade_date"].max().normalize()
                wf_start_date = wf_latest_trade_date - pd.Timedelta(days=int(wf_lookback_days))
                wf_window_trades = issuer_bucket_trades[issuer_bucket_trades["trade_date"] >= wf_start_date].copy()

                if wf_window_trades.empty:
                    st.warning("No trades were found inside the selected lookback window.")
                else:
                    wf_avg_issuer_yield = wf_window_trades["yield"].mean()
                    wf_trade_count = len(wf_window_trades)

                    if "trade_amount" in wf_window_trades.columns:
                        wf_total_trade_amount = pd.to_numeric(
                            wf_window_trades["trade_amount"], errors="coerce"
                        ).fillna(0).sum()
                    else:
                        wf_total_trade_amount = 0.0

                    wf_mmd = mmd_df.copy()
                    wf_mmd[date_col] = pd.to_datetime(wf_mmd[date_col], errors="coerce")
                    wf_mmd = wf_mmd.dropna(subset=[date_col])
                    wf_mmd = wf_mmd[wf_mmd[date_col].dt.normalize() <= wf_latest_trade_date].copy()

                    if wf_mmd.empty:
                        st.warning("No benchmark curve observations were available on or before the latest issuer trade date.")
                    else:
                        wf_mmd = wf_mmd.sort_values(date_col)
                        wf_latest_mmd = wf_mmd.iloc[[-1]].copy()
                        wf_benchmark_date = wf_latest_mmd[date_col].iloc[0]

                        wf_aaa_yield_series, wf_aaa_meta = get_benchmark_curve(wf_latest_mmd, wf_tenor, "AAA")
                        if wf_aaa_yield_series is None or pd.isna(wf_aaa_yield_series.iloc[0]):
                            st.warning(f"AAA/MMD {wf_tenor} curve could not be built for the waterfall.")
                        else:
                            wf_aaa_yield = float(wf_aaa_yield_series.iloc[0])
                            wf_total_spread_bps = (wf_avg_issuer_yield - wf_aaa_yield) * 100

                            wf_rating_premium_bps = (
                                RATING_SPREADS.get(wf_rating, RATING_SPREADS["AAA"]).get(wf_tenor, 0.00) * 100
                            )

                            wf_liq_source = wf_window_trades.copy()
                            wf_liq_source["trade_month"] = wf_liq_source["trade_date"].dt.to_period("M").astype(str)
                            wf_today = pd.Timestamp.today().normalize()

                            if "trade_amount" not in wf_liq_source.columns:
                                wf_liq_source["trade_amount"] = 0.0

                            wf_liq_by_cusip = (
                                wf_liq_source.groupby("cusip", dropna=False)
                                .agg(
                                    trade_count=("trade_date", "count"),
                                    latest_trade=("trade_date", "max"),
                                    active_months=("trade_month", "nunique"),
                                    total_trade_amount=("trade_amount", "sum"),
                                )
                                .reset_index()
                            )

                            if wf_liq_by_cusip.empty:
                                wf_avg_liquidity_score = pd.NA
                                wf_liquidity_premium_bps = 10.0
                                wf_liquidity_note = "No CUSIP-level liquidity score available; default proxy used."
                            else:
                                wf_liq_by_cusip["days_since_last_trade"] = (
                                    wf_today - wf_liq_by_cusip["latest_trade"]
                                ).dt.days
                                wf_liq_by_cusip["recent_90d_trades"] = wf_liq_by_cusip["trade_count"]
                                wf_liq_by_cusip["liquidity_score"] = (
                                    wf_liq_by_cusip["trade_count"].rank(pct=True) * 35
                                    + wf_liq_by_cusip["total_trade_amount"].rank(pct=True) * 25
                                    + wf_liq_by_cusip["recent_90d_trades"].rank(pct=True) * 25
                                    + (1 - wf_liq_by_cusip["days_since_last_trade"].rank(pct=True)) * 15
                                )
                                wf_avg_liquidity_score = wf_liq_by_cusip["liquidity_score"].mean()

                                if pd.isna(wf_avg_liquidity_score):
                                    wf_liquidity_premium_bps = 10.0
                                    wf_liquidity_note = "Liquidity score unavailable; default proxy used."
                                elif wf_avg_liquidity_score < 45:
                                    wf_liquidity_premium_bps = 15.0
                                    wf_liquidity_note = "Low liquidity bucket proxy."
                                elif wf_avg_liquidity_score < 75:
                                    wf_liquidity_premium_bps = 7.5
                                    wf_liquidity_note = "Medium liquidity bucket proxy."
                                else:
                                    wf_liquidity_premium_bps = 2.5
                                    wf_liquidity_note = "High liquidity bucket proxy."

                            callable_cols = [c for c in ["call_date", "call_date_bond"] if c in wf_window_trades.columns]
                            wf_callable_share = 0.0
                            if callable_cols:
                                call_col = callable_cols[0]
                                parsed_calls = pd.to_datetime(wf_window_trades[call_col], errors="coerce")
                                wf_callable_share = parsed_calls.notna().mean()
                            wf_callable_adjustment_bps = 5.0 if wf_callable_share >= 0.50 else 0.0

                            wf_residual_bps = (
                                wf_total_spread_bps
                                - wf_rating_premium_bps
                                - wf_liquidity_premium_bps
                                - wf_callable_adjustment_bps
                            )

                            waterfall_df = pd.DataFrame(
                                {
                                    "Component": [
                                        "AAA / MMD Base",
                                        "Rating Premium",
                                        "Liquidity Premium",
                                        "Callable Adjustment",
                                        "Residual / Issuer-Specific Premium",
                                        "Implied Issuer Yield",
                                    ],
                                    "Value": [
                                        wf_aaa_yield,
                                        wf_rating_premium_bps / 100,
                                        wf_liquidity_premium_bps / 100,
                                        wf_callable_adjustment_bps / 100,
                                        wf_residual_bps / 100,
                                        wf_avg_issuer_yield,
                                    ],
                                    "Display": [
                                        f"{wf_aaa_yield:.2f}%",
                                        f"{wf_rating_premium_bps:+.1f} bp",
                                        f"{wf_liquidity_premium_bps:+.1f} bp",
                                        f"{wf_callable_adjustment_bps:+.1f} bp",
                                        f"{wf_residual_bps:+.1f} bp",
                                        f"{wf_avg_issuer_yield:.2f}%",
                                    ],
                                }
                            )

                            wf_fig = go.Figure(
                                go.Waterfall(
                                    name="Spread Attribution",
                                    orientation="v",
                                    measure=["absolute", "relative", "relative", "relative", "relative", "total"],
                                    x=waterfall_df["Component"],
                                    y=waterfall_df["Value"],
                                    text=waterfall_df["Display"],
                                    textposition="outside",
                                    connector={"line": {"width": 1}},
                                )
                            )
                            wf_fig.update_layout(
                                title=f"{selected_issuer} Spread Attribution Waterfall ({wf_bucket}, vs AAA/MMD)",
                                yaxis_title="Yield / Spread Contribution (%)",
                                height=540,
                                showlegend=False,
                            )
                            safe_plotly_chart(wf_fig, width="stretch")

                            wf_metric1, wf_metric2, wf_metric3, wf_metric4 = st.columns(4)
                            wf_metric1.metric("Issuer Yield", f"{wf_avg_issuer_yield:.2f}%")
                            wf_metric2.metric("AAA / MMD Yield", f"{wf_aaa_yield:.2f}%")
                            wf_metric3.metric("Total Spread vs AAA", f"{wf_total_spread_bps:+.1f} bp")
                            wf_metric4.metric("Residual Premium", f"{wf_residual_bps:+.1f} bp")

                            if wf_residual_bps > 15:
                                st.info(
                                    f"Read-through: after modeled rating, liquidity, and callable components, "
                                    f"{selected_issuer}'s {wf_bucket} bucket still shows a positive residual premium "
                                    f"of {wf_residual_bps:+.1f} bp. This may indicate issuer-specific cheapness, "
                                    f"sector/supply pressure, data noise, or a component assumption that should be reviewed."
                                )
                            elif wf_residual_bps < -15:
                                st.info(
                                    f"Read-through: the modeled components exceed the observed spread by "
                                    f"{abs(wf_residual_bps):.1f} bp. This may indicate rich trading, stronger demand, "
                                    f"or overly conservative component assumptions."
                                )
                            else:
                                st.info(
                                    "Read-through: modeled components broadly explain the observed spread versus AAA/MMD. "
                                    "Residual premium is relatively modest."
                                )

                            with st.expander("Waterfall calculation audit table", expanded=False):
                                audit_df = pd.DataFrame(
                                    [
                                        {"Metric": "Selected issuer", "Value": selected_issuer},
                                        {"Metric": "Maturity year", "Value": wf_bucket},
                                        {"Metric": "MMD tenor", "Value": wf_tenor},
                                        {"Metric": "Latest issuer trade date", "Value": wf_latest_trade_date.strftime("%Y-%m-%d")},
                                        {"Metric": "Benchmark curve date", "Value": wf_benchmark_date.strftime("%Y-%m-%d")},
                                        {"Metric": "Issuer yield lookback", "Value": f"{wf_lookback_days} days"},
                                        {"Metric": "Trade count in lookback", "Value": f"{wf_trade_count:,}"},
                                        {"Metric": "Total trade amount in lookback", "Value": f"{wf_total_trade_amount:,.0f}"},
                                        {"Metric": "Average issuer yield", "Value": f"{wf_avg_issuer_yield:.4f}%"},
                                        {"Metric": "AAA/MMD yield", "Value": f"{wf_aaa_yield:.4f}%"},
                                        {"Metric": "Total spread vs AAA", "Value": f"{wf_total_spread_bps:+.2f} bp"},
                                        {"Metric": "Rating premium assumption", "Value": f"{wf_rating} / {wf_tenor}: {wf_rating_premium_bps:+.2f} bp"},
                                        {"Metric": "Average liquidity score", "Value": "" if pd.isna(wf_avg_liquidity_score) else f"{wf_avg_liquidity_score:.2f}"},
                                        {"Metric": "Liquidity premium proxy", "Value": f"{wf_liquidity_premium_bps:+.2f} bp — {wf_liquidity_note}"},
                                        {"Metric": "Callable share proxy", "Value": f"{wf_callable_share:.1%}"},
                                        {"Metric": "Callable adjustment proxy", "Value": f"{wf_callable_adjustment_bps:+.2f} bp"},
                                        {"Metric": "Residual / issuer-specific premium", "Value": f"{wf_residual_bps:+.2f} bp"},
                                        {"Metric": "Benchmark source", "Value": wf_aaa_meta.get("benchmark_source")},
                                        {"Metric": "Benchmark source column", "Value": wf_aaa_meta.get("source_column")},
                                    ]
                                )
                                safe_dataframe(audit_df, width="stretch", hide_index=True)



section_anchor("historical-spread", "Historical Spread Range & Percentile")
with st.expander("Methodology: historical spread range and percentile", expanded=False):
    st.markdown(
        """
This section compares the selected issuer's current spread against its own historical spread range.

**Core question:**

> Is the current spread historically wide, tight, or normal?

**Calculation:**

`Spread to Benchmark = (Average Issuer Yield - Benchmark Yield) × 100`

The module then calculates:

- **Current Spread**: latest available spread observation in the selected bucket.
- **Historical Median**: median spread over the selected history window.
- **Percentile**: where the current spread ranks versus the historical spread distribution.
- **Z-Score**: current spread distance from historical mean in standard deviation units.

**Signal guide:**

- **> 90th percentile**: very wide / potentially cheap
- **75th–90th percentile**: wide
- **25th–75th percentile**: normal range
- **10th–25th percentile**: tight
- **< 10th percentile**: very tight / potentially rich

Municipal trading can be sparse, so this is a screening signal rather than a final trade recommendation.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable historical spread percentile analysis.")
else:
    hist_col1, hist_col2, hist_col3, hist_col4 = st.columns([1, 1, 1, 1])
    with hist_col1:
        hist_bucket = st.selectbox(
            "Historical Maturity Year",
            MATURITY_BUCKET_ORDER,
            index=1,
            key="hist_spread_bucket",
        )
    with hist_col2:
        hist_rating = st.selectbox(
            "Historical Benchmark Curve",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="hist_spread_benchmark",
        )
    with hist_col3:
        hist_window_label = st.selectbox(
            "History Window",
            ["1M", "3M", "6M", "1Y", "All"],
            index=3,
            key="hist_spread_window",
        )
    with hist_col4:
        hist_smoothing = st.selectbox(
            "Smoothing",
            ["Daily", "7D average", "30D average"],
            index=0,
            key="hist_spread_smoothing",
        )

    hist_window_days = {"1M": 30, "3M": 90, "6M": 180, "1Y": 365, "All": None}[hist_window_label]

    hist_obs = build_spread_observations(
        market_df=market_df,
        mmd_df=mmd_df,
        issuer=selected_issuer,
        rating=hist_rating,
    )

    if hist_obs.empty:
        st.warning(
            "No overlapping issuer trade dates and benchmark dates were found for historical spread analysis. "
            "Check that the curve file date range overlaps with trade dates."
        )
    else:
        hist_obs = hist_obs.copy()
        hist_obs["trade_date"] = pd.to_datetime(hist_obs["trade_date"], errors="coerce").dt.normalize()
        hist_obs = hist_obs[
            (hist_obs["maturity_bucket"] == hist_bucket)
            & hist_obs["spread_to_benchmark_bps"].notna()
        ].sort_values("trade_date")

        if hist_obs.empty:
            st.warning(f"No historical spread observations found for {hist_bucket}.")
        else:
            latest_hist_date = hist_obs["trade_date"].max()
            if hist_window_days is not None:
                hist_cutoff = latest_hist_date - pd.Timedelta(days=hist_window_days)
                hist_plot = hist_obs[hist_obs["trade_date"] >= hist_cutoff].copy()
            else:
                hist_plot = hist_obs.copy()

            if hist_plot.empty or len(hist_plot) < 2:
                st.info("Not enough historical spread observations to calculate a percentile for the selected window.")
            else:
                # Aggregate to one observation per date first.
                hist_daily = (
                    hist_plot.groupby("trade_date", as_index=False)
                    .agg(
                        spread_to_benchmark_bps=("spread_to_benchmark_bps", "mean"),
                        avg_yield=("avg_yield", "mean"),
                        benchmark_yield=("benchmark_yield", "mean"),
                        trade_count=("trade_count", "sum"),
                        total_trade_amount=("total_trade_amount", "sum"),
                        benchmark_source=("benchmark_source", "first"),
                        source_column=("source_column", "first"),
                    )
                    .sort_values("trade_date")
                )

                if hist_smoothing == "7D average":
                    hist_daily["display_spread_bps"] = (
                        hist_daily["spread_to_benchmark_bps"].rolling(window=7, min_periods=1).mean()
                    )
                    display_label = "7D Average Spread"
                elif hist_smoothing == "30D average":
                    hist_daily["display_spread_bps"] = (
                        hist_daily["spread_to_benchmark_bps"].rolling(window=30, min_periods=1).mean()
                    )
                    display_label = "30D Average Spread"
                else:
                    hist_daily["display_spread_bps"] = hist_daily["spread_to_benchmark_bps"]
                    display_label = "Daily Spread"

                current_spread = float(hist_daily["display_spread_bps"].iloc[-1])
                hist_values = pd.to_numeric(hist_daily["display_spread_bps"], errors="coerce").dropna()

                hist_median = float(hist_values.median())
                hist_mean = float(hist_values.mean())
                hist_std = float(hist_values.std()) if len(hist_values) > 1 else 0.0
                hist_p25 = float(hist_values.quantile(0.25))
                hist_p75 = float(hist_values.quantile(0.75))
                hist_p10 = float(hist_values.quantile(0.10))
                hist_p90 = float(hist_values.quantile(0.90))
                hist_min = float(hist_values.min())
                hist_max = float(hist_values.max())

                percentile = float((hist_values <= current_spread).mean() * 100)
                z_score = (current_spread - hist_mean) / hist_std if hist_std and hist_std > 0 else pd.NA

                # Signal label
                if percentile >= 90:
                    signal = "Very Wide / Potentially Cheap"
                elif percentile >= 75:
                    signal = "Wide"
                elif percentile >= 25:
                    signal = "Normal Range"
                elif percentile >= 10:
                    signal = "Tight"
                else:
                    signal = "Very Tight / Potentially Rich"

                card1, card2, card3, card4 = st.columns(4)
                card1.metric("Current Spread", f"{current_spread:+.1f} bp")
                card2.metric("Historical Percentile", f"{percentile:.0f}th")
                card3.metric("Historical Median", f"{hist_median:+.1f} bp")
                card4.metric("Z-Score", "N/A" if pd.isna(z_score) else f"{z_score:+.2f}")

                st.info(
                    f"Read-through: {selected_issuer}'s {hist_bucket} spread to {hist_rating} is currently "
                    f"at the {percentile:.0f}th percentile versus the selected {hist_window_label} history. "
                    f"Signal: **{signal}**."
                )

                hist_fig = px.line(
                    hist_daily,
                    x="trade_date",
                    y="display_spread_bps",
                    markers=True,
                    hover_data={
                        "spread_to_benchmark_bps": ":.1f",
                        "avg_yield": ":.2f",
                        "benchmark_yield": ":.2f",
                        "trade_count": ":,.0f",
                        "total_trade_amount": ":,.0f",
                        "benchmark_source": True,
                        "source_column": True,
                    },
                    title=f"{selected_issuer} {hist_bucket} Spread to {hist_rating} Benchmark Over Time",
                    labels={
                        "trade_date": "Trade Date",
                        "display_spread_bps": f"{display_label} (bps)",
                        "spread_to_benchmark_bps": "Raw Daily Spread (bps)",
                    },
                )
                hist_fig.add_hline(y=current_spread, line_dash="solid", opacity=0.65, annotation_text="Current")
                hist_fig.add_hline(y=hist_median, line_dash="dash", opacity=0.65, annotation_text="Median")
                hist_fig.add_hline(y=hist_p25, line_dash="dot", opacity=0.45, annotation_text="25th")
                hist_fig.add_hline(y=hist_p75, line_dash="dot", opacity=0.45, annotation_text="75th")
                hist_fig.update_layout(height=520, hovermode="x unified")
                safe_plotly_chart(hist_fig, width="stretch")

                dist_col1, dist_col2 = st.columns([1.2, 1])
                with dist_col1:
                    hist_dist_fig = px.histogram(
                        hist_daily,
                        x="display_spread_bps",
                        nbins=25,
                        title=f"Historical Spread Distribution ({hist_window_label})",
                        labels={"display_spread_bps": f"{display_label} (bps)", "count": "Observation Count"},
                    )
                    hist_dist_fig.add_vline(x=current_spread, line_dash="solid", annotation_text="Current")
                    hist_dist_fig.add_vline(x=hist_median, line_dash="dash", annotation_text="Median")
                    hist_dist_fig.update_layout(height=430)
                    safe_plotly_chart(hist_dist_fig, width="stretch")

                with dist_col2:
                    pct_table = pd.DataFrame(
                        [
                            {"Statistic": "Minimum", "Spread (bps)": hist_min},
                            {"Statistic": "10th Percentile", "Spread (bps)": hist_p10},
                            {"Statistic": "25th Percentile", "Spread (bps)": hist_p25},
                            {"Statistic": "Median", "Spread (bps)": hist_median},
                            {"Statistic": "75th Percentile", "Spread (bps)": hist_p75},
                            {"Statistic": "90th Percentile", "Spread (bps)": hist_p90},
                            {"Statistic": "Maximum", "Spread (bps)": hist_max},
                            {"Statistic": "Current", "Spread (bps)": current_spread},
                        ]
                    )
                    pct_table["Spread (bps)"] = pct_table["Spread (bps)"].round(2)
                    safe_dataframe(pct_table, width="stretch", hide_index=True)

                with st.expander("Historical spread audit table", expanded=False):
                    audit_cols = [
                        "trade_date",
                        "spread_to_benchmark_bps",
                        "display_spread_bps",
                        "avg_yield",
                        "benchmark_yield",
                        "trade_count",
                        "total_trade_amount",
                        "benchmark_source",
                        "source_column",
                    ]
                    audit_hist = hist_daily[[c for c in audit_cols if c in hist_daily.columns]].copy()
                    for c in ["spread_to_benchmark_bps", "display_spread_bps", "avg_yield", "benchmark_yield"]:
                        if c in audit_hist.columns:
                            audit_hist[c] = pd.to_numeric(audit_hist[c], errors="coerce").round(2)
                    safe_dataframe(
                        audit_hist.sort_values("trade_date", ascending=False).head(1000),
                        width="stretch",
                        hide_index=True,
                    )



section_anchor("curve-shape", "Curve Shape Analytics")
with st.expander("Methodology: curve shape analytics", expanded=False):
    st.markdown(
        """
This section turns the issuer curve into **curve mathematics**, similar to what rates / muni desks monitor.

**Metrics:**

- **5s10s Slope** = 10Y yield − 5Y yield
- **10s30s Slope** = 30Y yield − 10Y yield
- **5s30s Slope** = 30Y yield − 5Y yield
- **5s10s30s Butterfly** = 10Y yield − average(5Y yield, 30Y yield)

**How to read it:**

- Higher positive slope = steeper curve.
- Lower or negative slope = flatter / inverted curve shape.
- Positive butterfly = 10Y “belly” is high/cheap versus wings.
- Negative butterfly = 10Y “belly” is low/rich versus wings.

The issuer curve uses uploaded trade data over the selected lookback window. The benchmark curve uses uploaded rating curves when available; otherwise it falls back to MMD/AAA plus transparent rating-spread assumptions.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable curve shape analytics.")
else:
    cs_col1, cs_col2, cs_col3 = st.columns([1, 1, 1])
    with cs_col1:
        cs_rating = st.selectbox(
            "Curve Shape Benchmark",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="curve_shape_benchmark",
        )
    with cs_col2:
        cs_lookback = st.selectbox(
            "Issuer Curve Lookback",
            [7, 30, 60, 90, 180],
            index=1,
            format_func=lambda x: f"Latest {x} days",
            key="curve_shape_lookback",
        )
    with cs_col3:
        cs_curve_basis = st.selectbox(
            "Curve Basis",
            ["Yield Curve", "Spread Curve"],
            index=0,
            key="curve_shape_basis",
            help="Yield Curve uses issuer yields. Spread Curve uses issuer spread to selected benchmark.",
        )

    cs_base = market_df[market_df["issuer"] == selected_issuer].copy()
    cs_base["trade_date"] = pd.to_datetime(cs_base["trade_date"], errors="coerce").dt.normalize()
    cs_base["yield"] = pd.to_numeric(cs_base["yield"], errors="coerce")
    cs_base = cs_base.dropna(subset=["trade_date", "yield", "maturity_bucket"])
    cs_base = cs_base[cs_base["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)].copy()

    if cs_base.empty:
        st.warning("No usable issuer trade rows were available for curve shape analytics.")
    else:
        cs_latest_date = cs_base["trade_date"].max()
        cs_start_date = cs_latest_date - pd.Timedelta(days=int(cs_lookback))
        cs_window = cs_base[cs_base["trade_date"] >= cs_start_date].copy()

        if cs_window.empty:
            st.warning("No issuer trades were found inside the selected lookback window.")
        else:
            issuer_curve = (
                cs_window.groupby("maturity_bucket", as_index=False)
                .agg(
                    issuer_yield=("yield", "mean"),
                    trade_count=("yield", "count"),
                    total_trade_amount=("trade_amount", "sum") if "trade_amount" in cs_window.columns else ("yield", "count"),
                    latest_trade=("trade_date", "max"),
                )
            )

            date_col = _detect_mmd_date_column(mmd_df)
            if date_col is None:
                st.warning("Curve shape analytics cannot build benchmark curve because the curve file has no usable date column.")
            else:
                cs_mmd = mmd_df.copy()
                cs_mmd[date_col] = pd.to_datetime(cs_mmd[date_col], errors="coerce")
                cs_mmd = cs_mmd.dropna(subset=[date_col])
                cs_mmd = cs_mmd[cs_mmd[date_col].dt.normalize() <= cs_latest_date].sort_values(date_col)

                if cs_mmd.empty:
                    st.warning("No benchmark curve observation was available on or before the latest issuer trade date.")
                else:
                    cs_latest_mmd = cs_mmd.iloc[[-1]].copy()
                    cs_benchmark_date = cs_latest_mmd[date_col].iloc[0]

                    bench_rows = []
                    for bucket in MATURITY_BUCKET_ORDER:
                        tenor = MMD_BUCKET_MAP.get(bucket, "10Y")
                        y, meta = get_benchmark_curve(cs_latest_mmd, tenor, cs_rating)
                        if y is not None and pd.notna(y.iloc[0]):
                            bench_rows.append(
                                {
                                    "maturity_bucket": bucket,
                                    "mmd_tenor": tenor,
                                    "benchmark_yield": float(y.iloc[0]),
                                    "benchmark_source": meta.get("benchmark_source"),
                                    "source_column": meta.get("source_column"),
                                    "rating_spread_bps": meta.get("rating_spread_bps"),
                                }
                            )
                    bench_curve = pd.DataFrame(bench_rows)

                    if bench_curve.empty:
                        st.warning("Selected benchmark curve could not be built for curve shape analytics.")
                    else:
                        curve_shape_df = issuer_curve.merge(bench_curve, on="maturity_bucket", how="outer")
                        curve_shape_df["spread_to_benchmark_bps"] = (
                            curve_shape_df["issuer_yield"] - curve_shape_df["benchmark_yield"]
                        ) * 100

                        maturity_order = MATURITY_BUCKET_ORDER
                        curve_shape_df["maturity_bucket"] = pd.Categorical(
                            curve_shape_df["maturity_bucket"],
                            categories=maturity_order,
                            ordered=True,
                        )
                        curve_shape_df = curve_shape_df.sort_values("maturity_bucket")

                        metric_col = "issuer_yield" if cs_curve_basis == "Yield Curve" else "spread_to_benchmark_bps"
                        metric_label = "Issuer Yield (%)" if cs_curve_basis == "Yield Curve" else f"Spread to {cs_rating} (bps)"

                        curve_plot = curve_shape_df.dropna(subset=[metric_col]).copy()
                        if curve_plot.empty:
                            st.warning("Not enough curve points were available to calculate curve shape metrics.")
                        else:
                            fig_curve_shape = px.line(
                                curve_plot,
                                x="maturity_bucket",
                                y=metric_col,
                                markers=True,
                                hover_data=[
                                    c for c in [
                                        "issuer_yield", "benchmark_yield", "spread_to_benchmark_bps",
                                        "trade_count", "total_trade_amount", "latest_trade",
                                        "benchmark_source", "source_column"
                                    ] if c in curve_plot.columns
                                ],
                                title=f"{selected_issuer} {cs_curve_basis} Shape",
                                labels={
                                    "maturity_bucket": "Maturity Year",
                                    metric_col: metric_label,
                                },
                            )
                            fig_curve_shape.update_layout(height=450, hovermode="x unified")
                            safe_plotly_chart(fig_curve_shape, width="stretch")

                            curve_values = (
                                curve_shape_df.set_index("maturity_bucket")[metric_col]
                                .astype(float)
                                .to_dict()
                            )

                            def get_curve_value(bucket: str):
                                value = curve_values.get(bucket)
                                return value if pd.notna(value) else pd.NA

                            v_short = get_curve_value("5Y")
                            v_10 = get_curve_value("10Y")
                            v_20 = get_curve_value("20Y")
                            v_30 = get_curve_value("30Y")

                            # -------------------------
                            # Dynamic curve diagnostics
                            # -------------------------
                            available_points = []
                            missing_points = []

                            for bucket_name, bucket_value in {
                                "5Y": v_short,
                                "10Y": v_10,
                                "20Y": v_20,
                                "30Y": v_30,
                            }.items():
                                if pd.notna(bucket_value):
                                    available_points.append(bucket_name)
                                else:
                                    missing_points.append(bucket_name)

                            diag_col1, diag_col2 = st.columns(2)

                            with diag_col1:
                                st.success(
                                    "Available Curve Points:\n\n"
                                    + ", ".join(available_points)
                                    if available_points
                                    else "No usable curve points detected."
                                )

                            with diag_col2:
                                if missing_points:
                                    st.warning(
                                        "Missing Curve Points:\n\n"
                                        + ", ".join(missing_points)
                                    )
                                else:
                                    st.success("All core curve points detected.")

                            metrics_rows = []
                            analytics_status = []

                            # 5s10s
                            if pd.notna(v_short) and pd.notna(v_10):
                                metrics_rows.append({
                                    "Metric": "5s10s Slope",
                                    "Value": v_10 - v_short,
                                })
                                analytics_status.append({
                                    "Analytics": "5s10s Slope",
                                    "Status": "Available",
                                    "Requirement": "5Y + 10Y",
                                })
                            else:
                                analytics_status.append({
                                    "Analytics": "5s10s Slope",
                                    "Status": "Missing Required Points",
                                    "Requirement": "5Y + 10Y",
                                })

                            # 10s30s
                            if pd.notna(v_10) and pd.notna(v_30):
                                metrics_rows.append({
                                    "Metric": "10s30s Slope",
                                    "Value": v_30 - v_10,
                                })
                                analytics_status.append({
                                    "Analytics": "10s30s Slope",
                                    "Status": "Available",
                                    "Requirement": "10Y + 30Y",
                                })
                            else:
                                analytics_status.append({
                                    "Analytics": "10s30s Slope",
                                    "Status": "Missing Required Points",
                                    "Requirement": "10Y + 30Y",
                                })

                            # 5s30s
                            if pd.notna(v_short) and pd.notna(v_30):
                                metrics_rows.append({
                                    "Metric": "5s30s Slope",
                                    "Value": v_30 - v_short,
                                })
                                analytics_status.append({
                                    "Analytics": "5s30s Slope",
                                    "Status": "Available",
                                    "Requirement": "5Y + 30Y",
                                })
                            else:
                                analytics_status.append({
                                    "Analytics": "5s30s Slope",
                                    "Status": "Missing Required Points",
                                    "Requirement": "5Y + 30Y",
                                })

                            # Butterfly
                            if pd.notna(v_short) and pd.notna(v_10) and pd.notna(v_30):
                                metrics_rows.append({
                                    "Metric": "5s10s30s Butterfly",
                                    "Value": v_10 - ((v_short + v_30) / 2),
                                })
                                analytics_status.append({
                                    "Analytics": "5s10s30s Butterfly",
                                    "Status": "Available",
                                    "Requirement": "5Y + 10Y + 30Y",
                                })
                            else:
                                analytics_status.append({
                                    "Analytics": "5s10s30s Butterfly",
                                    "Status": "Missing Required Points",
                                    "Requirement": "5Y + 10Y + 30Y",
                                })

                            # Steepness score
                            if pd.notna(v_short) and pd.notna(v_10) and pd.notna(v_20) and pd.notna(v_30):
                                metrics_rows.append({
                                    "Metric": "Steepness Score",
                                    "Value": (
                                        ((v_10 - v_short)
                                        + (v_30 - v_10)
                                        + (v_30 - v_short)) / 3
                                    ),
                                })
                                analytics_status.append({
                                    "Analytics": "Steepness Score",
                                    "Status": "Available",
                                    "Requirement": "5Y + 10Y + 20Y + 30Y",
                                })
                            else:
                                analytics_status.append({
                                    "Analytics": "Steepness Score",
                                    "Status": "Missing Required Points",
                                    "Requirement": "5Y + 10Y + 20Y + 30Y",
                                })

                            analytics_status_df = pd.DataFrame(analytics_status)

                            with st.expander("Curve Analytics Availability", expanded=False):
                                safe_dataframe(
                                    analytics_status_df,
                                    width="stretch",
                                    hide_index=True,
                                )

                            if not metrics_rows:
                                st.info("At least two compatible curve points are needed to calculate curve shape metrics.")
                            else:
                                metrics_df = pd.DataFrame(metrics_rows)
                                unit = "bp" if cs_curve_basis == "Spread Curve" else "%"
                                metrics_df["Display"] = metrics_df["Value"].map(
                                    lambda x: f"{x:+.1f} {unit}" if cs_curve_basis == "Spread Curve" else f"{x:+.2f}%"
                                )

                                mcols = st.columns(min(4, len(metrics_df)))
                                for idx, (_, row) in enumerate(metrics_df.head(4).iterrows()):
                                    mcols[idx % len(mcols)].metric(row["Metric"], row["Display"])

                                safe_dataframe(metrics_df, width="stretch", hide_index=True)

                                # Read-through
                                slope_1030 = metrics_df.loc[metrics_df["Metric"] == "10s30s Slope", "Value"]
                                butterfly = metrics_df.loc[metrics_df["Metric"] == "5s10s30s Butterfly", "Value"]

                                notes = []
                                if not slope_1030.empty:
                                    s_val = float(slope_1030.iloc[0])
                                    if s_val > (20 if cs_curve_basis == "Spread Curve" else 0.20):
                                        notes.append("long end screens steep versus the 10Y point")
                                    elif s_val < (-10 if cs_curve_basis == "Spread Curve" else -0.10):
                                        notes.append("long end screens flat/inverted versus the 10Y point")
                                    else:
                                        notes.append("10s30s slope appears relatively contained")

                                if not butterfly.empty:
                                    b_val = float(butterfly.iloc[0])
                                    if b_val > (10 if cs_curve_basis == "Spread Curve" else 0.10):
                                        notes.append("10Y belly appears cheap/high versus wings")
                                    elif b_val < (-10 if cs_curve_basis == "Spread Curve" else -0.10):
                                        notes.append("10Y belly appears rich/low versus wings")

                                if notes:
                                    st.info("Curve read-through: " + "; ".join(notes) + ".")

                                with st.expander("Curve shape audit table", expanded=False):
                                    audit_cols = [
                                        "maturity_bucket",
                                        "issuer_yield",
                                        "benchmark_yield",
                                        "spread_to_benchmark_bps",
                                        "trade_count",
                                        "total_trade_amount",
                                        "latest_trade",
                                        "mmd_tenor",
                                        "benchmark_source",
                                        "source_column",
                                        "rating_spread_bps",
                                    ]
                                    audit_curve = curve_shape_df[[c for c in audit_cols if c in curve_shape_df.columns]].copy()
                                    for c in ["issuer_yield", "benchmark_yield", "spread_to_benchmark_bps", "rating_spread_bps"]:
                                        if c in audit_curve.columns:
                                            audit_curve[c] = pd.to_numeric(audit_curve[c], errors="coerce").round(2)
                                    st.caption(
                                        f"Issuer lookback: latest {cs_lookback} days. "
                                        f"Benchmark date: {cs_benchmark_date.strftime('%Y-%m-%d')}."
                                    )
                                    safe_dataframe(audit_curve, width="stretch", hide_index=True)


section_anchor("market-narrative", "Market Narrative & Opportunity Map")
with st.expander("Methodology: market narrative and opportunity map", expanded=False):
    st.markdown(
        """
This section combines recent trading behavior with relative-value positioning.

**Trading Activity Timeline**

- Aggregates selected issuer trades by day.
- Shows daily trade count and total par traded.
- Adds a daily average yield line.
- Generates simple narrative events when activity, volume, or yield moves are unusually high relative to the issuer's recent history.

**Rich / Cheap Quadrant**

- Uses security-level observations.
- `x = Liquidity Score`
- `y = Spread to Benchmark` when benchmark data is available; otherwise `y = Average Yield`
- Vertical and horizontal median lines divide the map into four desk-style zones:
    - **Cheap + Liquid**: potential buy candidate
    - **Cheap + Illiquid**: opportunistic / liquidity premium required
    - **Rich + Liquid**: benchmark-like / monitor
    - **Rich + Illiquid**: low priority / avoid
        """
    )

mn_tab1, mn_tab2 = st.tabs(["Trading Activity Timeline", "Rich / Cheap Quadrant"])

with mn_tab1:
    if issuer_trades.empty:
        st.warning("No trade rows found for the selected issuer and filters.")
    else:
        timeline_df = issuer_trades.copy()
        timeline_df["trade_date"] = pd.to_datetime(timeline_df["trade_date"], errors="coerce")
        timeline_df["yield"] = pd.to_numeric(timeline_df["yield"], errors="coerce")
        if "trade_amount" in timeline_df.columns:
            timeline_df["trade_amount"] = pd.to_numeric(timeline_df["trade_amount"], errors="coerce").fillna(0)
        else:
            timeline_df["trade_amount"] = 0.0

        timeline_df = timeline_df.dropna(subset=["trade_date"])
        if timeline_df.empty:
            st.warning("Timeline cannot be generated because no valid trade dates were found.")
        else:
            timeline_daily = (
                timeline_df.groupby(timeline_df["trade_date"].dt.normalize(), as_index=False)
                .agg(
                    trade_date=("trade_date", "first"),
                    trade_count=("cusip", "count") if "cusip" in timeline_df.columns else ("yield", "count"),
                    total_trade_amount=("trade_amount", "sum"),
                    avg_yield=("yield", "mean"),
                )
                .sort_values("trade_date")
            )
            timeline_daily["yield_change_bp"] = timeline_daily["avg_yield"].diff() * 100

            lookback_options = {
                "Latest 30D": 30,
                "Latest 60D": 60,
                "Latest 90D": 90,
                "All": None,
            }
            timeline_window_label = st.selectbox(
                "Timeline Window",
                list(lookback_options.keys()),
                index=2,
                key="timeline_window",
            )
            timeline_days = lookback_options[timeline_window_label]
            timeline_plot = timeline_daily.copy()
            if timeline_days is not None and not timeline_plot.empty:
                cutoff = timeline_plot["trade_date"].max() - pd.Timedelta(days=timeline_days)
                timeline_plot = timeline_plot[timeline_plot["trade_date"] >= cutoff].copy()

            if timeline_plot.empty:
                st.info("No timeline observations are available for the selected window.")
            else:
                tl_fig = px.bar(
                    timeline_plot,
                    x="trade_date",
                    y="trade_count",
                    hover_data={
                        "total_trade_amount": ":,.0f",
                        "avg_yield": ":.2f",
                        "yield_change_bp": ":.1f",
                    },
                    title=f"{selected_issuer} Trading Activity Timeline",
                    labels={
                        "trade_date": "Trade Date",
                        "trade_count": "Trade Count",
                        "total_trade_amount": "Total Trade Amount",
                        "avg_yield": "Average Yield",
                        "yield_change_bp": "Daily Yield Change (bp)",
                    },
                )
                tl_fig.add_scatter(
                    x=timeline_plot["trade_date"],
                    y=timeline_plot["avg_yield"],
                    mode="lines+markers",
                    name="Average Yield",
                    yaxis="y2",
                )
                tl_fig.update_layout(
                    height=500,
                    yaxis=dict(title="Trade Count"),
                    yaxis2=dict(title="Average Yield (%)", overlaying="y", side="right"),
                    hovermode="x unified",
                )
                safe_plotly_chart(tl_fig, width="stretch")

                amount_fig = px.bar(
                    timeline_plot,
                    x="trade_date",
                    y="total_trade_amount",
                    title="Daily Total Par Traded",
                    labels={
                        "trade_date": "Trade Date",
                        "total_trade_amount": "Total Trade Amount",
                    },
                    hover_data={"trade_count": ":,.0f", "avg_yield": ":.2f"},
                )
                amount_fig.update_layout(height=380)
                safe_plotly_chart(amount_fig, width="stretch")

                # Simple narrative event detection
                event_rows = []
                tc_mean = timeline_daily["trade_count"].mean()
                tc_std = timeline_daily["trade_count"].std()
                amt_mean = timeline_daily["total_trade_amount"].mean()
                amt_std = timeline_daily["total_trade_amount"].std()
                yield_abs_threshold = 10.0

                for _, row in timeline_plot.iterrows():
                    notes = []
                    if pd.notna(tc_std) and tc_std > 0 and row["trade_count"] >= tc_mean + tc_std:
                        notes.append("Trade count above recent norm")
                    if pd.notna(amt_std) and amt_std > 0 and row["total_trade_amount"] >= amt_mean + amt_std:
                        notes.append("Heavy par traded")
                    if pd.notna(row.get("yield_change_bp")) and abs(row["yield_change_bp"]) >= yield_abs_threshold:
                        direction = "moved higher" if row["yield_change_bp"] > 0 else "moved lower"
                        notes.append(f"Average yield {direction} by {row['yield_change_bp']:+.1f} bp")
                    if notes:
                        event_rows.append(
                            {
                                "Date": row["trade_date"].strftime("%Y-%m-%d"),
                                "Narrative Signal": "; ".join(notes),
                                "Trade Count": int(row["trade_count"]),
                                "Total Trade Amount": row["total_trade_amount"],
                                "Average Yield": row["avg_yield"],
                                "Yield Change (bp)": row.get("yield_change_bp"),
                            }
                        )

                if event_rows:
                    st.subheader("Narrative Signals")
                    events_df = pd.DataFrame(event_rows)
                    safe_dataframe(events_df, width="stretch", hide_index=True)
                else:
                    st.info("No unusually large activity/volume/yield-move events were detected in the selected timeline window.")

with mn_tab2:
    if "liq" not in locals() or liq.empty:
        st.info("The quadrant map will be available after liquidity metrics are calculated for the selected issuer.")
    else:
        quadrant_df = liq.copy()

        # Choose Y-axis: spread to benchmark when available, else avg yield.
        quadrant_y_col = "avg_yield"
        quadrant_y_label = "Average Yield (%)"
        quadrant_source_note = "Using Average Yield because benchmark spread is not available inside this liquidity section."

        if not mmd_df.empty:
            try:
                q_rating = st.selectbox(
                    "Quadrant Benchmark Curve",
                    BENCHMARK_RATINGS,
                    index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
                    key="quadrant_benchmark_rating",
                )
                q_spread_obs = build_spread_observations(
                    market_df=market_df,
                    mmd_df=mmd_df,
                    issuer=selected_issuer,
                    rating=q_rating,
                )
                if not q_spread_obs.empty and "cusip" in issuer_trades.columns:
                    # Approximate CUSIP-level benchmark spread by merging the latest bucket-level spread to each CUSIP bucket.
                    latest_bucket_spread = (
                        q_spread_obs.sort_values("trade_date")
                        .groupby("maturity_bucket", as_index=False)
                        .tail(1)[["maturity_bucket", "spread_to_benchmark_bps"]]
                    )
                    if "maturity" in quadrant_df.columns:
                        pass
                    if "maturity_bucket" not in quadrant_df.columns and "maturity" in quadrant_df.columns:
                        # If liquidity table does not retain bucket, infer from maturity when possible.
                        maturity_tmp = pd.to_datetime(quadrant_df["maturity"], errors="coerce")
                        years_tmp = (maturity_tmp - pd.Timestamp.today()).dt.days / 365.25
                        quadrant_df["maturity_bucket"] = pd.cut(
                            years_tmp,
                            bins=[-float("inf"), 7, 15, 25, float("inf")],
                            labels=MATURITY_BUCKET_ORDER,
                        ).astype("string")
                    if "maturity_bucket" in quadrant_df.columns:
                        quadrant_df = quadrant_df.merge(latest_bucket_spread, on="maturity_bucket", how="left")
                        if quadrant_df["spread_to_benchmark_bps"].notna().any():
                            quadrant_y_col = "spread_to_benchmark_bps"
                            quadrant_y_label = f"Spread to {q_rating} Benchmark (bps)"
                            quadrant_source_note = (
                                "Using latest available bucket-level spread to benchmark. "
                                "Higher values generally indicate cheaper relative value."
                            )
            except Exception as exc:
                st.warning(f"Benchmark spread overlay was unavailable, so the quadrant uses average yield. Details: {exc}")

        required_cols = ["liquidity_score", quadrant_y_col]
        for col in required_cols:
            if col in quadrant_df.columns:
                quadrant_df[col] = pd.to_numeric(quadrant_df[col], errors="coerce")
        quadrant_df = quadrant_df.dropna(subset=[c for c in required_cols if c in quadrant_df.columns])

        if quadrant_df.empty:
            st.warning("No usable observations were available for the rich/cheap quadrant.")
        else:
            valid_buckets = MATURITY_BUCKET_ORDER
            if "maturity_bucket" not in quadrant_df.columns:
                quadrant_df["maturity_bucket"] = "Unknown"
            quadrant_df["maturity_bucket"] = quadrant_df["maturity_bucket"].astype(str)

            if "total_trade_amount" not in quadrant_df.columns:
                quadrant_df["total_trade_amount"] = 1
            quadrant_df["total_trade_amount"] = pd.to_numeric(
                quadrant_df["total_trade_amount"], errors="coerce"
            ).fillna(0).clip(lower=0)
            if quadrant_df["total_trade_amount"].sum() <= 0:
                quadrant_df["point_size"] = 10
                q_size_col = "point_size"
            else:
                q_size_col = "total_trade_amount"

            q_median_liq = quadrant_df["liquidity_score"].median()
            q_median_y = quadrant_df[quadrant_y_col].median()

            quadrant_df["Quadrant"] = "Unclassified"
            quadrant_df.loc[
                (quadrant_df["liquidity_score"] >= q_median_liq) & (quadrant_df[quadrant_y_col] >= q_median_y),
                "Quadrant",
            ] = "Cheap + Liquid"
            quadrant_df.loc[
                (quadrant_df["liquidity_score"] < q_median_liq) & (quadrant_df[quadrant_y_col] >= q_median_y),
                "Quadrant",
            ] = "Cheap + Illiquid"
            quadrant_df.loc[
                (quadrant_df["liquidity_score"] >= q_median_liq) & (quadrant_df[quadrant_y_col] < q_median_y),
                "Quadrant",
            ] = "Rich + Liquid"
            quadrant_df.loc[
                (quadrant_df["liquidity_score"] < q_median_liq) & (quadrant_df[quadrant_y_col] < q_median_y),
                "Quadrant",
            ] = "Rich + Illiquid"

            st.caption(quadrant_source_note)

            q_fig = px.scatter(
                quadrant_df,
                x="liquidity_score",
                y=quadrant_y_col,
                size=q_size_col,
                size_max=38,
                color="Quadrant",
                hover_name="cusip" if "cusip" in quadrant_df.columns else None,
                hover_data=[c for c in [
                    "maturity_bucket", "liquidity_tier", "trade_count", "recent_90d_trades",
                    "days_since_last_trade", "avg_yield", "spread_to_benchmark_bps",
                    "total_trade_amount", "outstanding_amount"
                ] if c in quadrant_df.columns],
                title=f"{selected_issuer} Rich / Cheap Liquidity Quadrant",
                labels={
                    "liquidity_score": "Liquidity Score",
                    quadrant_y_col: quadrant_y_label,
                    "total_trade_amount": "Total Trade Amount",
                },
            )
            q_fig.add_vline(x=q_median_liq, line_dash="dash", opacity=0.45)
            q_fig.add_hline(y=q_median_y, line_dash="dash", opacity=0.45)
            q_fig.add_annotation(
                x=0.98, y=0.98, xref="paper", yref="paper",
                text="Cheap + Liquid<br>Buy candidate",
                showarrow=False, align="right",
                bgcolor="rgba(255,255,255,0.75)",
            )
            q_fig.add_annotation(
                x=0.02, y=0.98, xref="paper", yref="paper",
                text="Cheap + Illiquid<br>Opportunistic",
                showarrow=False, align="left",
                bgcolor="rgba(255,255,255,0.75)",
            )
            q_fig.add_annotation(
                x=0.98, y=0.02, xref="paper", yref="paper",
                text="Rich + Liquid<br>Benchmark / monitor",
                showarrow=False, align="right",
                bgcolor="rgba(255,255,255,0.75)",
            )
            q_fig.add_annotation(
                x=0.02, y=0.02, xref="paper", yref="paper",
                text="Rich + Illiquid<br>Low priority",
                showarrow=False, align="left",
                bgcolor="rgba(255,255,255,0.75)",
            )
            q_fig.update_layout(height=560, hovermode="closest")
            safe_plotly_chart(q_fig, width="stretch")

            q_summary = (
                quadrant_df.groupby("Quadrant", as_index=False)
                .agg(
                    cusip_count=("cusip", "count") if "cusip" in quadrant_df.columns else ("liquidity_score", "count"),
                    avg_liquidity_score=("liquidity_score", "mean"),
                    avg_y_axis=(quadrant_y_col, "mean"),
                    total_trade_amount=("total_trade_amount", "sum") if "total_trade_amount" in quadrant_df.columns else ("liquidity_score", "count"),
                )
            )
            safe_dataframe(q_summary, width="stretch", hide_index=True)

            with st.expander("Quadrant security-level table", expanded=False):
                display_cols = [
                    "Quadrant", "cusip", "maturity_bucket", "liquidity_score", quadrant_y_col,
                    "avg_yield", "trade_count", "recent_90d_trades", "days_since_last_trade",
                    "total_trade_amount", "outstanding_amount", "liquidity_tier",
                ]
                safe_dataframe(
                    quadrant_df[[c for c in display_cols if c in quadrant_df.columns]]
                    .sort_values(["Quadrant", "liquidity_score"], ascending=[True, False])
                    .head(5000),
                    width="stretch",
                    hide_index=True,
                )



section_anchor("dealer-proxy", "Bid / Ask & Dealer Behavior Proxy")
with st.expander("Methodology: dealer behavior proxy", expanded=False):
    st.markdown(
        """
This section is intentionally **data-dependent**. It becomes most useful when uploaded trade files include a buy/sell indicator, dealer side, customer side, or trade type.

**Realistic implementation:**

- If a usable `trade_type` / side field exists, the app classifies trades into:
    - **Customer Buy / Buy**
    - **Customer Sell / Sell**
    - **Other / Unknown**
- It estimates a simple imbalance proxy:

`Selling Imbalance = (Sell Amount - Buy Amount) / (Sell Amount + Buy Amount)`

**How to read it:**

- Positive value: more selling pressure / potential dealer inventory pressure.
- Negative value: more buying demand.
- Near zero: balanced flow.

**Important limitation:**

This is a **proxy**, not true dealer inventory. True dealer behavior would require richer MSRB/dealer-side fields, bid-wanted data, or inventory data.
        """
    )

dealer_trade_type_col = None
for candidate_col in ["trade_type", "side", "buy_sell", "customer_side", "dealer_side"]:
    if candidate_col in market_df.columns:
        dealer_trade_type_col = candidate_col
        break

if dealer_trade_type_col is None:
    st.info(
        "Dealer behavior proxy is unavailable because no trade side / trade type field was detected. "
        "Upload trade data with fields such as Trade Type, Side, Customer Buy/Sell, or Dealer Side to enable this module."
    )
else:
    dealer_col1, dealer_col2, dealer_col3 = st.columns([1, 1, 1])
    with dealer_col1:
        dealer_scope = st.selectbox(
            "Dealer Proxy Scope",
            ["Selected issuer", "All uploaded issuers"],
            index=0,
            key="dealer_proxy_scope",
        )
    with dealer_col2:
        dealer_bucket = st.selectbox(
            "Dealer Proxy Maturity Year",
            MATURITY_BUCKET_OPTIONS,
            index=0,
            key="dealer_proxy_bucket",
        )
    with dealer_col3:
        dealer_window = st.selectbox(
            "Dealer Proxy Window",
            ["Latest 30D", "Latest 60D", "Latest 90D", "All"],
            index=2,
            key="dealer_proxy_window",
        )

    dealer_df = market_df.copy()
    dealer_df["trade_date"] = pd.to_datetime(dealer_df["trade_date"], errors="coerce").dt.normalize()
    dealer_df["yield"] = pd.to_numeric(dealer_df["yield"], errors="coerce")
    if "trade_amount" in dealer_df.columns:
        dealer_df["trade_amount"] = pd.to_numeric(dealer_df["trade_amount"], errors="coerce").fillna(0)
    else:
        dealer_df["trade_amount"] = 0.0

    dealer_df = dealer_df.dropna(subset=["trade_date"])
    if dealer_scope == "Selected issuer":
        dealer_df = dealer_df[dealer_df["issuer"] == selected_issuer].copy()
    if dealer_bucket != "All" and "maturity_bucket" in dealer_df.columns:
        dealer_df = dealer_df[dealer_df["maturity_bucket"] == dealer_bucket].copy()

    dealer_days = {"Latest 30D": 30, "Latest 60D": 60, "Latest 90D": 90, "All": None}[dealer_window]
    if not dealer_df.empty and dealer_days is not None:
        dealer_latest = dealer_df["trade_date"].max()
        dealer_df = dealer_df[dealer_df["trade_date"] >= dealer_latest - pd.Timedelta(days=dealer_days)].copy()

    if dealer_df.empty:
        st.warning("No trades remain for the selected dealer-proxy filters.")
    else:
        def classify_trade_side(value: object) -> str:
            text_val = str(value).strip().lower()
            if not text_val or text_val in {"nan", "none"}:
                return "Unknown"
            # Common variants from trade exports. These are intentionally broad.
            if any(token in text_val for token in ["sell", "sold", "sld", "customer sell", "cust sell", "cs"]):
                return "Sell"
            if any(token in text_val for token in ["buy", "bought", "purchase", "customer buy", "cust buy", "cb"]):
                return "Buy"
            if text_val in {"s"}:
                return "Sell"
            if text_val in {"b"}:
                return "Buy"
            return "Other / Unknown"

        dealer_df["flow_side"] = dealer_df[dealer_trade_type_col].map(classify_trade_side)

        flow_summary = (
            dealer_df.groupby("flow_side", as_index=False)
            .agg(
                trade_count=("trade_date", "count"),
                total_trade_amount=("trade_amount", "sum"),
                avg_yield=("yield", "mean"),
            )
        )

        buy_amount = flow_summary.loc[flow_summary["flow_side"] == "Buy", "total_trade_amount"].sum()
        sell_amount = flow_summary.loc[flow_summary["flow_side"] == "Sell", "total_trade_amount"].sum()
        buy_count = flow_summary.loc[flow_summary["flow_side"] == "Buy", "trade_count"].sum()
        sell_count = flow_summary.loc[flow_summary["flow_side"] == "Sell", "trade_count"].sum()
        denom_amount = buy_amount + sell_amount
        denom_count = buy_count + sell_count
        amount_imbalance = (sell_amount - buy_amount) / denom_amount if denom_amount > 0 else pd.NA
        count_imbalance = (sell_count - buy_count) / denom_count if denom_count > 0 else pd.NA

        dp1, dp2, dp3, dp4 = st.columns(4)
        dp1.metric("Buy Amount", f"{buy_amount:,.0f}")
        dp2.metric("Sell Amount", f"{sell_amount:,.0f}")
        dp3.metric("Amount Imbalance", "N/A" if pd.isna(amount_imbalance) else f"{amount_imbalance:+.1%}")
        dp4.metric("Count Imbalance", "N/A" if pd.isna(count_imbalance) else f"{count_imbalance:+.1%}")

        if pd.notna(amount_imbalance):
            if amount_imbalance >= 0.25:
                st.info("Read-through: flow appears sell-heavy, which may indicate customer selling pressure or dealer balance-sheet pressure.")
            elif amount_imbalance <= -0.25:
                st.info("Read-through: flow appears buy-heavy, which may indicate stronger demand or dealer distribution.")
            else:
                st.info("Read-through: buy/sell flow appears relatively balanced in the selected window.")
        else:
            st.info("Read-through: buy/sell imbalance cannot be calculated because buy/sell amount data is unavailable or zero.")

        flow_fig = px.bar(
            flow_summary,
            x="flow_side",
            y="total_trade_amount",
            hover_data={"trade_count": ":,.0f", "avg_yield": ":.2f"},
            title="Trade Amount by Flow Side",
            labels={
                "flow_side": "Flow Side",
                "total_trade_amount": "Total Trade Amount",
                "trade_count": "Trade Count",
                "avg_yield": "Average Yield",
            },
        )
        flow_fig.update_layout(height=430)
        safe_plotly_chart(flow_fig, width="stretch")

        dealer_daily = (
            dealer_df.groupby(["trade_date", "flow_side"], as_index=False)
            .agg(
                trade_count=("trade_date", "count"),
                total_trade_amount=("trade_amount", "sum"),
            )
        )
        daily_pivot = (
            dealer_daily.pivot_table(
                index="trade_date",
                columns="flow_side",
                values="total_trade_amount",
                aggfunc="sum",
                fill_value=0,
                observed=False,
            )
            .reset_index()
        )
        if "Buy" not in daily_pivot.columns:
            daily_pivot["Buy"] = 0
        if "Sell" not in daily_pivot.columns:
            daily_pivot["Sell"] = 0
        daily_pivot["net_selling_pressure"] = daily_pivot["Sell"] - daily_pivot["Buy"]

        pressure_fig = px.bar(
            daily_pivot,
            x="trade_date",
            y="net_selling_pressure",
            title="Daily Net Selling Pressure Proxy",
            labels={
                "trade_date": "Trade Date",
                "net_selling_pressure": "Sell Amount - Buy Amount",
            },
        )
        pressure_fig.add_hline(y=0, line_dash="dash", opacity=0.45)
        pressure_fig.update_layout(height=430)
        safe_plotly_chart(pressure_fig, width="stretch")

        with st.expander("Dealer proxy audit table", expanded=False):
            audit_cols = [
                "trade_date",
                "issuer",
                "cusip",
                "maturity_bucket",
                dealer_trade_type_col,
                "flow_side",
                "yield",
                "price",
                "trade_amount",
                "trade_type",
            ]

            # Remove missing columns and duplicate column names while preserving order.
            # This matters when dealer_trade_type_col itself is "trade_type".
            audit_cols = list(dict.fromkeys([c for c in audit_cols if c in dealer_df.columns]))

            dealer_audit_display = dealer_df[audit_cols].copy()
            dealer_audit_display = dealer_audit_display.loc[:, ~dealer_audit_display.columns.duplicated()].copy()

            st.caption(f"Side classification source column: `{dealer_trade_type_col}`.")
            safe_dataframe(
                dealer_audit_display.sort_values("trade_date", ascending=False).head(5000),
                width="stretch",
                hide_index=True,
            )



section_anchor("rv-positioning", "Relative Value Positioning Map")
with st.expander("Methodology: relative value positioning map", expanded=False):
    st.markdown(
        """
This scatter plot maps individual CUSIPs by **tradability** and **relative value**.

**Default interpretation:**

- **X-axis = Liquidity Score**: higher means more actively traded, larger traded amount, more recent activity, and less staleness.
- **Y-axis = Spread to Benchmark**: higher means the bond is trading cheaper versus the selected benchmark curve.
- **Bubble size = Total Trade Amount**: larger dots indicate more secondary-market trading volume.
- **Color = Maturity Year**: 1Y / 2Y / 3Y / ....

**Quadrants:**

- **Upper-right:** cheap and liquid; often the first area to investigate.
- **Upper-left:** cheap but illiquid; may require a liquidity premium.
- **Lower-right:** liquid but rich; useful benchmark-like bonds.
- **Lower-left:** illiquid and rich; usually less attractive from a relative-value screen.

This is a **screening view**, not an investment recommendation. It helps analysts identify bonds worth deeper review.
        """
    )

if issuer_trades.empty:
    st.warning("No trade rows found for this issuer and filter.")
else:
    rv_controls = st.columns([1, 1, 1, 1])
    with rv_controls[0]:
        rv_benchmark_rating = st.selectbox(
            "RV Benchmark Curve",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="rv_benchmark_rating",
            help="Used only when Y-axis is spread to benchmark. Uploaded curve columns are used first; otherwise MMD + assumption spread.",
        )
    with rv_controls[1]:
        rv_y_axis = st.selectbox(
            "Y-axis",
            ["Spread to Benchmark (bps)", "Average Yield (%)"],
            index=0,
            key="rv_y_axis",
        )
    with rv_controls[2]:
        rv_size_by = st.selectbox(
            "Bubble size",
            ["Total Trade Amount", "Outstanding Amount", "Trade Count"],
            index=0,
            key="rv_size_by",
        )
    with rv_controls[3]:
        rv_min_trades = st.number_input(
            "Minimum Trades",
            min_value=1,
            max_value=100,
            value=1,
            step=1,
            key="rv_min_trades",
        )

    rv_base = issuer_trades.copy()
    rv_base["trade_date"] = pd.to_datetime(rv_base["trade_date"], errors="coerce").dt.normalize()
    rv_base["yield"] = pd.to_numeric(rv_base["yield"], errors="coerce")
    if "trade_amount" in rv_base.columns:
        rv_base["trade_amount"] = pd.to_numeric(rv_base["trade_amount"], errors="coerce")
    else:
        rv_base["trade_amount"] = pd.NA
    if "price" in rv_base.columns:
        rv_base["price"] = pd.to_numeric(rv_base["price"], errors="coerce")
    else:
        rv_base["price"] = pd.NA

    rv_base = rv_base.dropna(subset=["cusip", "trade_date", "yield"])

    if rv_base.empty:
        st.warning("No usable CUSIP-level trade rows are available for the positioning map.")
    else:
        today_rv = pd.Timestamp.today().normalize()
        rv_base["trade_month"] = rv_base["trade_date"].dt.to_period("M").astype(str)
        rv_agg_dict = {
            "avg_yield": ("yield", "mean"),
            "latest_yield": ("yield", "last"),
            "avg_price": ("price", "mean"),
            "trade_count": ("trade_date", "count"),
            "first_trade": ("trade_date", "min"),
            "latest_trade": ("trade_date", "max"),
            "active_months": ("trade_month", "nunique"),
            "total_trade_amount": ("trade_amount", "sum"),
            "avg_trade_amount": ("trade_amount", "mean"),
        }
        optional_first_cols = {
            "maturity_bucket": "maturity_bucket",
            "maturity": "maturity_bond",
            "coupon": "coupon_bond",
            "outstanding_amount": "outstanding_amount",
            "description": "description",
        }
        for output_col, source_col in optional_first_cols.items():
            if source_col in rv_base.columns:
                rv_agg_dict[output_col] = (source_col, "first")

        rv_summary = (
            rv_base.groupby("cusip", dropna=False)
            .agg(**rv_agg_dict)
            .reset_index()
        )

        for required_col in ["maturity_bucket", "maturity", "coupon", "outstanding_amount", "description"]:
            if required_col not in rv_summary.columns:
                rv_summary[required_col] = pd.NA
        rv_summary["days_since_last_trade"] = (today_rv - rv_summary["latest_trade"]).dt.days
        rv_summary["trading_period_days"] = (rv_summary["latest_trade"] - rv_summary["first_trade"]).dt.days.clip(lower=1)
        rv_summary["avg_days_between_trades"] = rv_summary["trading_period_days"] / rv_summary["trade_count"].clip(lower=1)
        rv_summary["avg_trades_per_month"] = rv_summary["trade_count"] / rv_summary["active_months"].clip(lower=1)

        recent_cutoff_rv = today_rv - pd.DateOffset(days=90)
        rv_recent = (
            rv_base[rv_base["trade_date"] >= recent_cutoff_rv]
            .groupby("cusip")
            .agg(recent_90d_trades=("trade_date", "count"))
            .reset_index()
        )
        rv_summary = rv_summary.merge(rv_recent, on="cusip", how="left")
        rv_summary["recent_90d_trades"] = rv_summary["recent_90d_trades"].fillna(0).astype(int)

        for numeric_col in ["total_trade_amount", "outstanding_amount", "avg_trade_amount"]:
            if numeric_col in rv_summary.columns:
                rv_summary[numeric_col] = pd.to_numeric(rv_summary[numeric_col], errors="coerce")
        rv_summary["turnover_ratio"] = rv_summary["total_trade_amount"] / rv_summary["outstanding_amount"].replace({0: pd.NA})
        rv_summary["liquidity_score"] = (
            rv_summary["trade_count"].rank(pct=True) * 35
            + rv_summary["total_trade_amount"].fillna(0).rank(pct=True) * 25
            + rv_summary["recent_90d_trades"].rank(pct=True) * 25
            + (1 - rv_summary["days_since_last_trade"].rank(pct=True)) * 15
        )
        rv_summary["liquidity_tier"] = pd.cut(
            rv_summary["liquidity_score"],
            bins=[-1, 45, 75, 101],
            labels=["Low Liquidity", "Medium Liquidity", "High Liquidity"],
        ).astype(str)
        rv_summary.loc[rv_summary["days_since_last_trade"] > 365, "liquidity_tier"] = "Stale"

        rv_summary = rv_summary[rv_summary["trade_count"] >= rv_min_trades].copy()

        # Add benchmark spread at each CUSIP's latest trade date and maturity year.
        if rv_y_axis == "Spread to Benchmark (bps)":
            if mmd_df.empty:
                st.info("Upload an MMD / benchmark curve file to use Spread to Benchmark. Showing Average Yield instead.")
                rv_y_axis_col = "avg_yield"
                rv_y_axis_label = "Average Yield (%)"
            else:
                benchmark_long_rv = make_benchmark_long(mmd_df, rv_benchmark_rating)
                if benchmark_long_rv.empty:
                    st.info("No usable benchmark curve was found for the selected rating. Showing Average Yield instead.")
                    rv_y_axis_col = "avg_yield"
                    rv_y_axis_label = "Average Yield (%)"
                else:
                    benchmark_long_rv = benchmark_long_rv.sort_values(["maturity_bucket", "trade_date"])
                    merge_frames = []
                    for bucket in MATURITY_BUCKET_ORDER:
                        left = rv_summary[rv_summary["maturity_bucket"] == bucket].sort_values("latest_trade")
                        right = benchmark_long_rv[benchmark_long_rv["maturity_bucket"] == bucket].sort_values("trade_date")
                        if left.empty or right.empty:
                            continue
                        merged_bucket = pd.merge_asof(
                            left,
                            right,
                            left_on="latest_trade",
                            right_on="trade_date",
                            direction="backward",
                            tolerance=pd.Timedelta(days=14),
                        )
                        merge_frames.append(merged_bucket)
                    if merge_frames:
                        rv_summary = pd.concat(merge_frames, ignore_index=True)
                        rv_summary["spread_to_benchmark_bps"] = (
                            rv_summary["avg_yield"] - rv_summary["benchmark_yield"]
                        ) * 100
                        rv_y_axis_col = "spread_to_benchmark_bps"
                        rv_y_axis_label = "Spread to Benchmark (bps)"
                    else:
                        st.info("No overlapping CUSIP latest-trade dates and benchmark dates were found. Showing Average Yield instead.")
                        rv_y_axis_col = "avg_yield"
                        rv_y_axis_label = "Average Yield (%)"
        else:
            rv_y_axis_col = "avg_yield"
            rv_y_axis_label = "Average Yield (%)"

        rv_summary = rv_summary.dropna(subset=["liquidity_score", rv_y_axis_col])

        if rv_summary.empty:
            st.warning("No CUSIPs meet the selected filters for the positioning map.")
        else:
            size_map = {
                "Total Trade Amount": "total_trade_amount",
                "Outstanding Amount": "outstanding_amount",
                "Trade Count": "trade_count",
            }
            size_col = size_map.get(rv_size_by, "total_trade_amount")

            # Defensive plotting layer -------------------------------------------------
            # Plotly scatter is sensitive to missing/non-numeric/negative values in
            # size, x, and y columns. Muni exports often have blank outstanding amount,
            # missing trade amount, missing maturity dates, or unmatched benchmark values.
            #
            # We handle this in two steps:
            #   1) Clean numeric plotting inputs so the chart does not crash.
            #   2) Split known maturity years from unknown maturity years so
            #      "Unknown" does not dominate or pollute the main positioning map.
            rv_plot = rv_summary.copy()

            for numeric_col in ["liquidity_score", rv_y_axis_col, size_col]:
                if numeric_col in rv_plot.columns:
                    rv_plot[numeric_col] = pd.to_numeric(rv_plot[numeric_col], errors="coerce")
                    rv_plot[numeric_col] = rv_plot[numeric_col].replace([float("inf"), -float("inf")], pd.NA)

            required_plot_cols = ["liquidity_score", rv_y_axis_col]
            rv_plot = rv_plot.dropna(subset=[c for c in required_plot_cols if c in rv_plot.columns])

            # Resolve maturity year from common merge variants. If the bucket still
            # cannot be determined, keep the row for audit but exclude it from the
            # main scatter chart.
            valid_buckets = MATURITY_BUCKET_ORDER

            if "maturity_bucket" not in rv_plot.columns:
                possible_bucket_cols = [
                    "maturity_bucket_x",
                    "maturity_bucket_y",
                    "maturity_bucket_trade",
                    "maturity_bucket_bond",
                ]
                found_bucket_col = next((c for c in possible_bucket_cols if c in rv_plot.columns), None)
                if found_bucket_col:
                    rv_plot["maturity_bucket"] = rv_plot[found_bucket_col]
                else:
                    rv_plot["maturity_bucket"] = pd.NA

            rv_plot["maturity_bucket"] = rv_plot["maturity_bucket"].astype("string")

            if "cusip" not in rv_plot.columns:
                rv_plot["cusip"] = rv_plot.index.astype(str)
            else:
                rv_plot["cusip"] = rv_plot["cusip"].fillna("Unknown").astype(str)

            rv_known = rv_plot[rv_plot["maturity_bucket"].isin(valid_buckets)].copy()
            rv_unknown = rv_plot[~rv_plot["maturity_bucket"].isin(valid_buckets)].copy()

            # Clean the bubble-size column only on the known-bucket plotting set.
            if size_col not in rv_known.columns:
                rv_known["point_size"] = 10
                size_col = "point_size"
            else:
                rv_known[size_col] = pd.to_numeric(rv_known[size_col], errors="coerce")
                rv_known[size_col] = rv_known[size_col].replace([float("inf"), -float("inf")], pd.NA)
                rv_known[size_col] = rv_known[size_col].fillna(0).clip(lower=0)

                if rv_known[size_col].sum() <= 0:
                    rv_known["point_size"] = 10
                    size_col = "point_size"

            hover_cols = [
                "cusip", "maturity_bucket", "maturity", "coupon", "avg_yield", "avg_price",
                "trade_count", "recent_90d_trades", "days_since_last_trade", "total_trade_amount",
                "outstanding_amount", "turnover_ratio", "liquidity_tier",
            ]
            if "spread_to_benchmark_bps" in rv_known.columns:
                hover_cols.extend(["spread_to_benchmark_bps", "benchmark_yield", "benchmark_source", "source_column"])
            hover_cols = [c for c in hover_cols if c in rv_known.columns]

            if rv_plot.empty:
                st.warning(
                    "No valid observations remain after cleaning the positioning-map inputs. "
                    "Try lowering the minimum trade filter or using Average Yield instead of Spread to Benchmark."
                )
                rv_summary = rv_plot
                median_liquidity = pd.NA
                median_y = pd.NA

            elif rv_known.empty:
                st.warning(
                    "No bonds with known maturity years were available for the main positioning map. "
                    "Unknown-maturity bonds are listed below for audit."
                )
                rv_summary = rv_known
                median_liquidity = pd.NA
                median_y = pd.NA

            else:
                try:
                    median_liquidity = rv_known["liquidity_score"].median()
                    median_y = rv_known[rv_y_axis_col].median()

                    rv_known_labeled = add_security_label(rv_known)
                    # Replace crowded bubble scatter with a ranked horizontal bar view.
                    rv_fig = ranked_bar_chart(
                        rv_known_labeled,
                        value_col=rv_y_axis_col,
                        label_col="security_label",
                        title=f"{selected_issuer} Ranked Relative Value Candidates",
                        x_title=rv_y_axis_label,
                        top_n=18,
                        color_col="maturity_bucket",
                        hover_cols=hover_cols,
                    )
                    if rv_fig is not None:
                        safe_plotly_chart(rv_fig, width="stretch")

                    # Add a compact decision table so users do not need to interpret a dense bubble map.
                    rv_read = rv_known_labeled.copy()
                    rv_read["desk_signal"] = np.select(
                        [
                            (rv_read["liquidity_score"] >= median_liquidity) & (rv_read[rv_y_axis_col] >= median_y),
                            (rv_read["liquidity_score"] < median_liquidity) & (rv_read[rv_y_axis_col] >= median_y),
                            (rv_read["liquidity_score"] >= median_liquidity) & (rv_read[rv_y_axis_col] < median_y),
                        ],
                        ["Cheap + Liquid", "Cheap / Needs Liquidity Check", "Rich + Liquid"],
                        default="In Line / Review",
                    )
                    read_cols = ["desk_signal", "security_label", "maturity_bucket", rv_y_axis_col, "liquidity_score", "trade_count", "total_trade_amount"]
                    read_display = rv_read[[c for c in read_cols if c in rv_read.columns]].sort_values(rv_y_axis_col, ascending=False).head(15)
                    st.subheader("Positioning Read-Through")
                    safe_dataframe(read_display, width="stretch", hide_index=True, auto_collapse=False, height=400)
                except Exception as exc:
                    st.warning(
                        "The positioning map could not be plotted because the scatter inputs were not usable. "
                        f"The cleaned known-maturity data table is shown below for review. Error: {exc}"
                    )
                    safe_dataframe(rv_known.head(1000), width="stretch", hide_index=True)
                    median_liquidity = rv_known["liquidity_score"].median() if "liquidity_score" in rv_known.columns else pd.NA
                    median_y = rv_known[rv_y_axis_col].median() if rv_y_axis_col in rv_known.columns else pd.NA

                # Use the cleaned known-bucket plotting data for quadrant/read-through logic.
                rv_summary = rv_known

            # Unknown maturity year audit ------------------------------------------
            # These rows are not bad data; they are simply excluded from the main map
            # because the maturity year could not be determined from the uploaded
            # bond/trade data. Keeping them visible makes the dashboard transparent
            # without letting Unknown dominate the legend.
            if not rv_unknown.empty:
                with st.expander(
                    f"Unknown maturity year bonds excluded from main map ({len(rv_unknown):,})",
                    expanded=False,
                ):
                    st.caption(
                        "These CUSIPs were excluded from the main positioning map because their maturity year "
                        "could not be determined from the uploaded bond/trade data. They are retained here for audit."
                    )
                    unknown_display_cols = [
                        "cusip",
                        "avg_yield",
                        "spread_to_benchmark_bps",
                        "liquidity_score",
                        "trade_count",
                        "recent_90d_trades",
                        "days_since_last_trade",
                        "total_trade_amount",
                        "outstanding_amount",
                    ]
                    unknown_existing_cols = [c for c in unknown_display_cols if c in rv_unknown.columns]
                    safe_dataframe(
                        rv_unknown[unknown_existing_cols].head(5000),
                        width="stretch",
                        hide_index=True,
                    )

            if (
                rv_y_axis_col == "spread_to_benchmark_bps"
                and not rv_summary.empty
                and pd.notna(median_liquidity)
                and pd.notna(median_y)
            ):
                candidates = rv_summary[
                    (rv_summary["liquidity_score"] >= median_liquidity)
                    & (rv_summary["spread_to_benchmark_bps"] >= median_y)
                ].sort_values(["spread_to_benchmark_bps", "liquidity_score"], ascending=False)
                if not candidates.empty:
                    top = candidates.iloc[0]
                    st.info(
                        f"Positioning read-through: {top['cusip']} screens as relatively cheap and liquid "
                        f"at {top['spread_to_benchmark_bps']:+.1f} bp versus {rv_benchmark_rating}, "
                        f"with a liquidity score of {top['liquidity_score']:.1f}."
                    )

            with st.expander("Positioning map data table", expanded=False):
                display_cols = [
                    "cusip", "maturity_bucket", "liquidity_score", "liquidity_tier", rv_y_axis_col,
                    "avg_yield", "benchmark_yield", "benchmark_source", "source_column", "trade_count",
                    "recent_90d_trades", "days_since_last_trade", "total_trade_amount", "outstanding_amount",
                    "turnover_ratio", "maturity", "coupon", "avg_price",
                ]
                display_cols = [c for c in display_cols if c in rv_summary.columns]
                rv_display = rv_summary[display_cols].copy()
                for c in ["liquidity_score", rv_y_axis_col, "avg_yield", "benchmark_yield", "turnover_ratio", "avg_price"]:
                    if c in rv_display.columns:
                        rv_display[c] = pd.to_numeric(rv_display[c], errors="coerce").round(2)
                safe_dataframe(
                    rv_display.sort_values([rv_y_axis_col, "liquidity_score"], ascending=False),
                    width="stretch",
                    hide_index=True,
                    height=420,
                )

section_anchor("scenario-shock", "Scenario Shock Analysis")
with st.expander("Methodology: scenario shock analysis", expanded=False):
    st.markdown(
        """
This section estimates how the selected issuer or uploaded securities may react under simple interest-rate shock scenarios.

**Purpose:**

- Estimate approximate price impact under rate shocks.
- Identify which maturity years or CUSIPs are most exposed to parallel moves, steepening, or flattening.
- Provide a first-pass risk lens for secondary trading and pitchbook discussion.

**Version 1 approximation:**

This is a **duration-proxy model**, not a full bond pricing engine.

Core formula:

`Approximate Price Impact ≈ -Duration × Yield Shock`

Where:

- Duration is proxied by maturity year unless a duration field is uploaded.
- Yield shock is expressed in decimal form. Example: `+25 bp = +0.0025`.
- Price impact is an approximate percentage price move.

**Default proxy durations:**

The model now assigns proxy duration by annual maturity year. As a simple first-pass approximation, duration is capped at 18 years:

`proxy_duration = min(maturity_year, 18)`

**Important limitations:**

- This does not model callable optionality, convexity, OAS, amortization, tax effects, or full cash flows.
- Callable / premium bonds may behave differently from this simple duration approximation.
- Treat this as a screening and risk-discussion tool, not a final valuation model.
        """
    )

DURATION_PROXY = {f"{y}Y": min(float(y), 18.0) for y in range(1, MAX_MATURITY_YEAR + 1)}

SHOCK_SCENARIOS_BPS = {
    "+25bp Parallel": {f"{y}Y": 25 for y in range(1, MAX_MATURITY_YEAR + 1)},
    "+50bp Parallel": {f"{y}Y": 50 for y in range(1, MAX_MATURITY_YEAR + 1)},
    "-25bp Parallel": {f"{y}Y": -25 for y in range(1, MAX_MATURITY_YEAR + 1)},
    "Bear Steepening": {f"{y}Y": min(5 + y, 35) for y in range(1, MAX_MATURITY_YEAR + 1)},
    "Bull Flattening": {f"{y}Y": max(-25 + int(y / 3), -5) for y in range(1, MAX_MATURITY_YEAR + 1)},
    "Front-End Selloff": {f"{y}Y": max(35 - y, 5) for y in range(1, MAX_MATURITY_YEAR + 1)},
}

shock_col1, shock_col2, shock_col3 = st.columns([1, 1, 1])
with shock_col1:
    shock_scope = st.selectbox(
        "Shock Scope",
        ["Selected issuer", "All uploaded issuers"],
        index=0,
        key="shock_scope",
    )
with shock_col2:
    shock_scenario = st.selectbox(
        "Rate Shock Scenario",
        list(SHOCK_SCENARIOS_BPS.keys()) + ["Custom"],
        index=0,
        key="shock_scenario",
    )
with shock_col3:
    shock_view = st.selectbox(
        "Shock View",
        ["Maturity year summary", "CUSIP-level detail"],
        index=0,
        key="shock_view",
    )

if shock_scenario == "Custom":
    custom_col1, custom_col2, custom_col3, custom_col4 = st.columns(4)
    with custom_col1:
        shock_1 = st.number_input("1Y Shock (bp)", value=25.0, step=5.0, key="shock_1")
    with custom_col2:
        shock_5 = st.number_input("5Y Shock (bp)", value=25.0, step=5.0, key="shock_5")
    with custom_col3:
        shock_10 = st.number_input("10Y Shock (bp)", value=25.0, step=5.0, key="shock_10")
    with custom_col4:
        shock_30 = st.number_input("30Y Shock (bp)", value=25.0, step=5.0, key="shock_30")
    # Interpolate a simple custom shock path across annual maturity years.
    selected_shocks = {}
    anchor_points = [(1, float(shock_1)), (5, float(shock_5)), (10, float(shock_10)), (30, float(shock_30))]
    for y in range(1, MAX_MATURITY_YEAR + 1):
        if y <= 1:
            val = anchor_points[0][1]
        elif y >= 30:
            val = anchor_points[-1][1]
        else:
            left, right = anchor_points[0], anchor_points[-1]
            for a, b in zip(anchor_points, anchor_points[1:]):
                if a[0] <= y <= b[0]:
                    left, right = a, b
                    break
            val = left[1] + (right[1] - left[1]) * (y - left[0]) / (right[0] - left[0])
        selected_shocks[f"{y}Y"] = float(val)
else:
    selected_shocks = SHOCK_SCENARIOS_BPS[shock_scenario]

shock_base = market_df.copy()
if shock_scope == "Selected issuer":
    shock_base = shock_base[shock_base["issuer"] == selected_issuer].copy()

if shock_base.empty:
    st.warning("No trade rows are available for the selected shock scope.")
else:
    shock_base["trade_date"] = pd.to_datetime(shock_base["trade_date"], errors="coerce").dt.normalize()
    shock_base["yield"] = pd.to_numeric(shock_base["yield"], errors="coerce")
    if "price" in shock_base.columns:
        shock_base["price"] = pd.to_numeric(shock_base["price"], errors="coerce")
    else:
        shock_base["price"] = pd.NA
    if "trade_amount" in shock_base.columns:
        shock_base["trade_amount"] = pd.to_numeric(shock_base["trade_amount"], errors="coerce").fillna(0)
    else:
        shock_base["trade_amount"] = 0.0

    shock_base = shock_base.dropna(subset=["trade_date", "yield", "maturity_bucket"])
    shock_base = shock_base[shock_base["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)].copy()

    if shock_base.empty:
        st.warning("No usable rows with maturity years were available for scenario shock analysis.")
    else:
        # Build bucket summary from latest available selected-scope data.
        latest_shock_date = shock_base["trade_date"].max()
        shock_recent = shock_base[shock_base["trade_date"] >= latest_shock_date - pd.Timedelta(days=90)].copy()
        if shock_recent.empty:
            shock_recent = shock_base.copy()

        bucket_summary = (
            shock_recent.groupby("maturity_bucket", as_index=False)
            .agg(
                avg_yield=("yield", "mean"),
                avg_price=("price", "mean"),
                trade_count=("yield", "count"),
                total_trade_amount=("trade_amount", "sum"),
                latest_trade=("trade_date", "max"),
            )
        )
        bucket_summary["proxy_duration"] = bucket_summary["maturity_bucket"].map(DURATION_PROXY)
        bucket_summary["shock_bps"] = bucket_summary["maturity_bucket"].map(selected_shocks)
        bucket_summary["shock_decimal"] = bucket_summary["shock_bps"] / 10000
        bucket_summary["approx_price_impact_pct"] = -bucket_summary["proxy_duration"] * bucket_summary["shock_decimal"] * 100
        bucket_summary["shocked_yield"] = bucket_summary["avg_yield"] + (bucket_summary["shock_bps"] / 100)
        bucket_summary["impact_direction"] = bucket_summary["approx_price_impact_pct"].map(
            lambda x: "Price Down" if x < 0 else "Price Up" if x > 0 else "Flat"
        )

        shock_m1, shock_m2, shock_m3, shock_m4 = st.columns(4)
        worst_bucket_row = bucket_summary.sort_values("approx_price_impact_pct").iloc[0]
        best_bucket_row = bucket_summary.sort_values("approx_price_impact_pct", ascending=False).iloc[0]
        weighted_impact = (
            (bucket_summary["approx_price_impact_pct"] * bucket_summary["total_trade_amount"]).sum()
            / bucket_summary["total_trade_amount"].sum()
            if bucket_summary["total_trade_amount"].sum() > 0
            else bucket_summary["approx_price_impact_pct"].mean()
        )

        shock_m1.metric("Scenario", shock_scenario)
        shock_m2.metric("Worst Bucket", f"{worst_bucket_row['maturity_bucket']}")
        shock_m3.metric("Worst Approx Impact", f"{worst_bucket_row['approx_price_impact_pct']:+.2f}%")
        shock_m4.metric("Weighted Impact", f"{weighted_impact:+.2f}%")

        st.info(
            f"Read-through: under **{shock_scenario}**, the most rate-sensitive bucket is "
            f"{worst_bucket_row['maturity_bucket']} with an approximate price impact of "
            f"{worst_bucket_row['approx_price_impact_pct']:+.2f}%. "
            f"This is based on proxy duration and should be treated as a first-pass risk estimate."
        )

        st.subheader("1. Shock Impact by Maturity Year")
        shock_bar = px.bar(
            bucket_summary,
            x="maturity_bucket",
            y="approx_price_impact_pct",
            hover_data={
                "avg_yield": ":.2f",
                "shocked_yield": ":.2f",
                "proxy_duration": ":.1f",
                "shock_bps": ":.0f",
                "trade_count": ":,.0f",
                "total_trade_amount": ":,.0f",
            },
            title=f"Approximate Price Impact by Bucket — {shock_scenario}",
            labels={
                "maturity_bucket": "Maturity Year",
                "approx_price_impact_pct": "Approximate Price Impact (%)",
                "avg_yield": "Current Avg Yield",
                "shocked_yield": "Shocked Yield",
                "proxy_duration": "Proxy Duration",
                "shock_bps": "Shock (bp)",
            },
        )
        shock_bar.add_hline(y=0, line_dash="dash", opacity=0.45)
        shock_bar.update_layout(height=440)
        safe_plotly_chart(shock_bar, width="stretch")

        st.subheader("2. Shock Summary Table")
        bucket_display = bucket_summary.copy()
        for col in ["avg_yield", "avg_price", "proxy_duration", "shocked_yield", "approx_price_impact_pct"]:
            if col in bucket_display.columns:
                bucket_display[col] = pd.to_numeric(bucket_display[col], errors="coerce").round(2)
        safe_dataframe(
            bucket_display[
                [
                    "maturity_bucket",
                    "avg_yield",
                    "shocked_yield",
                    "shock_bps",
                    "proxy_duration",
                    "approx_price_impact_pct",
                    "trade_count",
                    "total_trade_amount",
                    "latest_trade",
                    "impact_direction",
                ]
            ],
            width="stretch",
            hide_index=True,
        )

        # CUSIP-level shock detail.
        if shock_view == "CUSIP-level detail":
            st.subheader("3. CUSIP-Level Shock Detail")
            cusip_shock = (
                shock_recent.groupby("cusip", dropna=False)
                .agg(
                    issuer=("issuer", "first"),
                    sector=("sector", "first") if "sector" in shock_recent.columns else ("issuer", "first"),
                    maturity_bucket=("maturity_bucket", "first"),
                    maturity=("maturity_bond", "first") if "maturity_bond" in shock_recent.columns else ("trade_date", "max"),
                    coupon=("coupon_bond", "first") if "coupon_bond" in shock_recent.columns else ("yield", "count"),
                    avg_yield=("yield", "mean"),
                    avg_price=("price", "mean"),
                    trade_count=("yield", "count"),
                    total_trade_amount=("trade_amount", "sum"),
                    latest_trade=("trade_date", "max"),
                )
                .reset_index()
            )
            cusip_shock["proxy_duration"] = cusip_shock["maturity_bucket"].map(DURATION_PROXY)
            cusip_shock["shock_bps"] = cusip_shock["maturity_bucket"].map(selected_shocks)
            cusip_shock["shock_decimal"] = cusip_shock["shock_bps"] / 10000
            cusip_shock["approx_price_impact_pct"] = -cusip_shock["proxy_duration"] * cusip_shock["shock_decimal"] * 100
            cusip_shock["shocked_yield"] = cusip_shock["avg_yield"] + (cusip_shock["shock_bps"] / 100)
            cusip_shock = cusip_shock.sort_values("approx_price_impact_pct")

            detail_cols = [
                "cusip",
                "issuer",
                "sector",
                "maturity_bucket",
                "maturity",
                "coupon",
                "avg_yield",
                "shocked_yield",
                "shock_bps",
                "proxy_duration",
                "approx_price_impact_pct",
                "trade_count",
                "total_trade_amount",
                "latest_trade",
            ]
            detail_display = cusip_shock[[c for c in detail_cols if c in cusip_shock.columns]].copy()
            for col in ["avg_yield", "shocked_yield", "proxy_duration", "approx_price_impact_pct"]:
                if col in detail_display.columns:
                    detail_display[col] = pd.to_numeric(detail_display[col], errors="coerce").round(2)

            safe_dataframe(detail_display.head(5000), width="stretch", hide_index=True, height=480)

            shock_scatter = px.scatter(
                cusip_shock,
                x="proxy_duration",
                y="approx_price_impact_pct",
                size="total_trade_amount",
                color="maturity_bucket",
                hover_name="cusip",
                hover_data=[
                    c for c in ["issuer", "avg_yield", "shocked_yield", "shock_bps", "trade_count", "latest_trade"]
                    if c in cusip_shock.columns
                ],
                title="CUSIP Shock Exposure Map",
                labels={
                    "proxy_duration": "Proxy Duration",
                    "approx_price_impact_pct": "Approx Price Impact (%)",
                    "maturity_bucket": "Maturity Year",
                },
            )
            shock_scatter.add_hline(y=0, line_dash="dash", opacity=0.45)
            shock_scatter.update_layout(height=500)
            safe_plotly_chart(shock_scatter, width="stretch")

        with st.expander("Scenario shock assumptions and audit", expanded=False):
            shock_assumption_df = pd.DataFrame(
                [
                    {
                        "Maturity Year": bucket,
                        "Shock (bp)": selected_shocks.get(bucket),
                        "Proxy Duration": DURATION_PROXY.get(bucket),
                        "Formula": "Approx Price Impact ≈ -Duration × Shock",
                    }
                    for bucket in MATURITY_BUCKET_ORDER
                ]
            )
            safe_dataframe(shock_assumption_df, width="stretch", hide_index=True)

            st.download_button(
                label="Download Scenario Shock Results CSV",
                data=bucket_summary.to_csv(index=False).encode("utf-8"),
                file_name="scenario_shock_results.csv",
                mime="text/csv",
            )


section_anchor("watchlist", "Watchlist / Saved Candidates")
with st.expander("Methodology: watchlist / saved candidates", expanded=False):
    st.markdown(
        """
This section lets users save CUSIPs for later review during the current session.

**Why it matters:**

A trading workflow often moves from screening → shortlist → detailed review. The watchlist keeps promising CUSIPs visible without forcing users to re-filter every time.

**Current implementation:**

- Uses Streamlit session state, so it persists while the app session is active.
- Users can add CUSIPs from the selected issuer or all uploaded data.
- Users can download the watchlist as CSV.
- This is not a database-backed permanent watchlist yet; that would be a future production feature.
        """
    )

if "watchlist_cusips" not in st.session_state:
    st.session_state["watchlist_cusips"] = []

watch_col1, watch_col2 = st.columns([1, 1])
with watch_col1:
    watch_scope = st.selectbox(
        "Watchlist Add Scope",
        ["Selected issuer", "All uploaded issuers"],
        index=0,
        key="watchlist_scope",
    )
with watch_col2:
    watch_universe = market_df.copy()
    if watch_scope == "Selected issuer":
        watch_universe = watch_universe[watch_universe["issuer"] == selected_issuer].copy()
    watch_options = sorted(watch_universe["cusip"].dropna().astype(str).unique().tolist()) if "cusip" in watch_universe.columns else []
    cusips_to_add = st.multiselect(
        "Add CUSIPs to Watchlist",
        watch_options,
        default=[],
        key="watchlist_add_cusips",
    )

add_col, clear_col = st.columns([1, 1])
with add_col:
    if st.button("Add selected CUSIPs", key="add_watchlist_cusips"):
        current = set(st.session_state.get("watchlist_cusips", []))
        current.update(cusips_to_add)
        st.session_state["watchlist_cusips"] = sorted(current)
        st.success(f"Added {len(cusips_to_add):,} CUSIP(s) to watchlist.")
with clear_col:
    if st.button("Clear watchlist", key="clear_watchlist_cusips"):
        st.session_state["watchlist_cusips"] = []
        st.info("Watchlist cleared.")

watchlist_cusips = st.session_state.get("watchlist_cusips", [])
if not watchlist_cusips:
    st.info("No CUSIPs saved yet. Add candidates from the selector above.")
else:
    watchlist_df = market_df[market_df["cusip"].astype(str).isin(watchlist_cusips)].copy()
    if watchlist_df.empty:
        st.warning("Saved CUSIPs were not found in the current uploaded dataset.")
    else:
        watchlist_df["trade_date"] = pd.to_datetime(watchlist_df["trade_date"], errors="coerce")
        if "trade_amount" in watchlist_df.columns:
            watchlist_df["trade_amount"] = pd.to_numeric(watchlist_df["trade_amount"], errors="coerce").fillna(0)
        else:
            watchlist_df["trade_amount"] = 0.0
        watchlist_summary = (
            watchlist_df.groupby("cusip", dropna=False)
            .agg(
                issuer=("issuer", "first"),
                sector=("sector", "first") if "sector" in watchlist_df.columns else ("issuer", "first"),
                maturity_bucket=("maturity_bucket", "first") if "maturity_bucket" in watchlist_df.columns else ("issuer", "first"),
                maturity=("maturity_bond", "first") if "maturity_bond" in watchlist_df.columns else ("trade_date", "max"),
                coupon=("coupon_bond", "first") if "coupon_bond" in watchlist_df.columns else ("trade_date", "count"),
                avg_yield=("yield", "mean"),
                avg_price=("price", "mean") if "price" in watchlist_df.columns else ("yield", "mean"),
                trade_count=("trade_date", "count"),
                total_trade_amount=("trade_amount", "sum"),
                latest_trade=("trade_date", "max"),
            )
            .reset_index()
        )
        st.metric("Saved CUSIPs", f"{len(watchlist_summary):,}")
        safe_dataframe(watchlist_summary, width="stretch", hide_index=True)

        st.download_button(
            label="Download Watchlist CSV",
            data=watchlist_summary.to_csv(index=False).encode("utf-8"),
            file_name="watchlist_saved_candidates.csv",
            mime="text/csv",
        )


section_anchor("recommendation-engine", "Trade Recommendation Narrative Engine")
with st.expander("Methodology: rule-based recommendation narrative", expanded=False):
    st.markdown(
        """
This section generates a **rule-driven market commentary** from the analytics already shown in the dashboard.

It is intentionally not an AI black box. Each phrase is triggered by a transparent data rule.

**Signals used when available:**

- **Spread movement:** whether the selected maturity year widened/tightened over the chosen lookback window.
- **Historical percentile:** whether the current spread is wide/tight versus its own history.
- **Liquidity:** whether the bucket remains tradable based on trade count, amount traded, and recency.
- **Peer comparison:** whether the selected issuer screens wide/tight versus uploaded peers.
- **Dealer flow proxy:** whether observed trade type/side suggests buy-heavy or sell-heavy activity.

**Example rule mapping:**

| Data trigger | Narrative phrase |
|---|---|
| Spread change > +15 bp | widened materially |
| Historical percentile > 90th | near the upper end of its historical range |
| Liquidity score > 70 | while maintaining above-average liquidity |
| Wider than peer median by > 10 bp | screens wide versus uploaded peers |
| Sell imbalance > 25% | flow appears sell-heavy |

**Important limitation:**

This is a screening commentary, not an investment recommendation. It should be reviewed alongside credit fundamentals, call structure, tax status, and actual executable market levels.
        """
    )

if mmd_df.empty:
    st.info("Upload an MMD/benchmark curve file to enable rule-based recommendation commentary.")
else:
    rec_col1, rec_col2, rec_col3 = st.columns([1, 1, 1])
    with rec_col1:
        rec_bucket = st.selectbox(
            "Narrative Maturity Year",
            MATURITY_BUCKET_ORDER,
            index=3,
            key="rec_bucket",
        )
    with rec_col2:
        rec_rating = st.selectbox(
            "Narrative Benchmark",
            BENCHMARK_RATINGS,
            index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
            key="rec_rating",
        )
    with rec_col3:
        rec_window_label = st.selectbox(
            "Narrative Lookback",
            ["1W", "1M", "3M", "6M", "1Y"],
            index=1,
            key="rec_window",
        )

    rec_window_days = {"1W": 7, "1M": 30, "3M": 90, "6M": 180, "1Y": 365}[rec_window_label]

    # -----------------------------
    # Signal 1: spread movement + current spread
    # -----------------------------
    rec_obs = build_spread_observations(
        market_df=market_df,
        mmd_df=mmd_df,
        issuer=selected_issuer,
        rating=rec_rating,
    )

    narrative_lines = []
    evidence_rows = []

    current_spread = pd.NA
    spread_change = pd.NA
    historical_percentile = pd.NA
    liquidity_score = pd.NA
    peer_gap = pd.NA
    amount_imbalance = pd.NA

    if rec_obs.empty:
        st.warning("No overlapping issuer trade dates and benchmark dates were found for narrative generation.")
    else:
        rec_obs = rec_obs.copy()
        rec_obs["trade_date"] = pd.to_datetime(rec_obs["trade_date"], errors="coerce").dt.normalize()
        rec_obs = rec_obs[
            (rec_obs["maturity_bucket"] == rec_bucket)
            & rec_obs["spread_to_benchmark_bps"].notna()
        ].sort_values("trade_date")

        if rec_obs.empty:
            st.warning(f"No spread observations found for {selected_issuer} / {rec_bucket}.")
        else:
            latest_row = rec_obs.iloc[-1]
            latest_date = latest_row["trade_date"]
            current_spread = float(latest_row["spread_to_benchmark_bps"])
            target_date = latest_date - pd.Timedelta(days=rec_window_days)
            hist_candidates = rec_obs[rec_obs["trade_date"] <= target_date]

            if not hist_candidates.empty:
                hist_row = hist_candidates.iloc[-1]
                historical_spread = float(hist_row["spread_to_benchmark_bps"])
                spread_change = current_spread - historical_spread
            else:
                historical_spread = pd.NA

            # Historical percentile uses up to 1Y of history when available.
            hist_1y = rec_obs[rec_obs["trade_date"] >= latest_date - pd.Timedelta(days=365)].copy()
            hist_values = pd.to_numeric(hist_1y["spread_to_benchmark_bps"], errors="coerce").dropna()
            if len(hist_values) >= 2:
                historical_percentile = float((hist_values <= current_spread).mean() * 100)

            evidence_rows.append({
                "Signal": "Current spread",
                "Value": f"{current_spread:+.1f} bp",
                "Rule / Source": f"Latest {rec_bucket} spread to {rec_rating}",
            })
            if pd.notna(spread_change):
                evidence_rows.append({
                    "Signal": f"{rec_window_label} spread movement",
                    "Value": f"{spread_change:+.1f} bp",
                    "Rule / Source": "Latest spread minus historical spread at/before lookback target date",
                })
            if pd.notna(historical_percentile):
                evidence_rows.append({
                    "Signal": "Historical percentile",
                    "Value": f"{historical_percentile:.0f}th",
                    "Rule / Source": "Current spread percentile versus latest 1Y spread observations",
                })

            # Movement phrase.
            if pd.notna(spread_change):
                if spread_change >= 15:
                    movement_phrase = f"{rec_bucket} widened materially by {spread_change:+.1f} bp over {rec_window_label}"
                elif spread_change >= 5:
                    movement_phrase = f"{rec_bucket} widened modestly by {spread_change:+.1f} bp over {rec_window_label}"
                elif spread_change <= -15:
                    movement_phrase = f"{rec_bucket} tightened materially by {spread_change:+.1f} bp over {rec_window_label}"
                elif spread_change <= -5:
                    movement_phrase = f"{rec_bucket} tightened modestly by {spread_change:+.1f} bp over {rec_window_label}"
                else:
                    movement_phrase = f"{rec_bucket} was broadly stable over {rec_window_label}"
                narrative_lines.append(f"{selected_issuer} {movement_phrase} versus the {rec_rating} benchmark.")
            else:
                narrative_lines.append(
                    f"{selected_issuer} {rec_bucket} currently screens at {current_spread:+.1f} bp versus the {rec_rating} benchmark."
                )

            # Historical phrase.
            if pd.notna(historical_percentile):
                if historical_percentile >= 90:
                    narrative_lines.append(
                        f"Current spread sits near the upper end of its recent historical range ({historical_percentile:.0f}th percentile)."
                    )
                elif historical_percentile >= 75:
                    narrative_lines.append(
                        f"Current spread screens wide versus recent history ({historical_percentile:.0f}th percentile)."
                    )
                elif historical_percentile <= 10:
                    narrative_lines.append(
                        f"Current spread sits near the tight end of its recent historical range ({historical_percentile:.0f}th percentile)."
                    )
                elif historical_percentile <= 25:
                    narrative_lines.append(
                        f"Current spread screens tight versus recent history ({historical_percentile:.0f}th percentile)."
                    )
                else:
                    narrative_lines.append(
                        f"Current spread is close to its normal recent historical range ({historical_percentile:.0f}th percentile)."
                    )

    # -----------------------------
    # Signal 2: liquidity proxy
    # -----------------------------
    rec_trades = market_df[
        (market_df["issuer"] == selected_issuer)
        & (market_df["maturity_bucket"] == rec_bucket)
    ].copy()

    if not rec_trades.empty:
        rec_trades["trade_date"] = pd.to_datetime(rec_trades["trade_date"], errors="coerce").dt.normalize()
        if "trade_amount" in rec_trades.columns:
            rec_trades["trade_amount"] = pd.to_numeric(rec_trades["trade_amount"], errors="coerce").fillna(0)
        else:
            rec_trades["trade_amount"] = 0.0

        rec_trades = rec_trades.dropna(subset=["trade_date"])
        if not rec_trades.empty:
            latest_trade_date = rec_trades["trade_date"].max()
            rec_trade_window = rec_trades[rec_trades["trade_date"] >= latest_trade_date - pd.Timedelta(days=rec_window_days)].copy()

            if not rec_trade_window.empty:
                trade_count = len(rec_trade_window)
                total_trade_amount = rec_trade_window["trade_amount"].sum()
                days_since_last = (pd.Timestamp.today().normalize() - latest_trade_date).days

                # Simple bounded liquidity proxy.
                trade_count_score = min(trade_count / 10, 1) * 35
                amount_score = min(total_trade_amount / 5_000_000, 1) * 35
                recency_score = max(0, 1 - min(days_since_last / 180, 1)) * 30
                liquidity_score = trade_count_score + amount_score + recency_score

                evidence_rows.append({
                    "Signal": "Liquidity score",
                    "Value": f"{liquidity_score:.1f}",
                    "Rule / Source": "35% trade count + 35% par traded + 30% recency proxy",
                })
                evidence_rows.append({
                    "Signal": "Window trade count",
                    "Value": f"{trade_count:,}",
                    "Rule / Source": f"Trades in selected {rec_window_label} window",
                })
                evidence_rows.append({
                    "Signal": "Window par traded",
                    "Value": f"{total_trade_amount:,.0f}",
                    "Rule / Source": f"Total par traded in selected {rec_window_label} window",
                })

                if liquidity_score >= 70:
                    narrative_lines.append("The bucket maintains above-average liquidity based on recent trade count, par traded, and recency.")
                elif liquidity_score >= 45:
                    narrative_lines.append("Liquidity appears moderate; execution quality should still be checked at the CUSIP level.")
                else:
                    narrative_lines.append("Liquidity appears limited, so any apparent cheapness may include a meaningful liquidity premium.")

    # -----------------------------
    # Signal 3: peer comparison proxy
    # -----------------------------
    try:
        if "sector" in market_df.columns and selected_sector and selected_sector != "Unknown":
            peer_universe = market_df[
                market_df["sector"].astype(str) == str(selected_sector)
            ].copy()
        else:
            peer_universe = market_df.copy()

        peer_universe = peer_universe[
            (peer_universe["issuer"].astype(str) != str(selected_issuer))
            & (peer_universe["maturity_bucket"] == rec_bucket)
        ].copy()

        if not peer_universe.empty and pd.notna(current_spread):
            peer_universe["trade_date"] = pd.to_datetime(peer_universe["trade_date"], errors="coerce").dt.normalize()
            peer_universe["yield"] = pd.to_numeric(peer_universe["yield"], errors="coerce")
            peer_universe = peer_universe.dropna(subset=["trade_date", "yield"])
            if not peer_universe.empty:
                latest_peer_date = peer_universe["trade_date"].max()
                peer_universe = peer_universe[
                    peer_universe["trade_date"] >= latest_peer_date - pd.Timedelta(days=rec_window_days)
                ].copy()

                if not peer_universe.empty:
                    date_col = _detect_mmd_date_column(mmd_df)
                    if date_col is not None:
                        peer_mmd = mmd_df.copy()
                        peer_mmd[date_col] = pd.to_datetime(peer_mmd[date_col], errors="coerce")
                        peer_mmd = peer_mmd.dropna(subset=[date_col])
                        peer_mmd = peer_mmd[peer_mmd[date_col].dt.normalize() <= latest_peer_date].sort_values(date_col)
                        if not peer_mmd.empty:
                            peer_latest_mmd = peer_mmd.iloc[[-1]].copy()
                            peer_tenor = MMD_BUCKET_MAP.get(rec_bucket, "10Y")
                            peer_bench, _peer_meta = get_benchmark_curve(peer_latest_mmd, peer_tenor, rec_rating)
                            if peer_bench is not None and pd.notna(peer_bench.iloc[0]):
                                peer_benchmark_yield = float(peer_bench.iloc[0])
                                peer_summary_for_narrative = (
                                    peer_universe.groupby("issuer", as_index=False)
                                    .agg(avg_yield=("yield", "mean"), trade_count=("yield", "count"))
                                )
                                peer_summary_for_narrative["peer_spread_bps"] = (
                                    peer_summary_for_narrative["avg_yield"] - peer_benchmark_yield
                                ) * 100
                                peer_median = peer_summary_for_narrative["peer_spread_bps"].median()
                                peer_gap = current_spread - peer_median

                                evidence_rows.append({
                                    "Signal": "Peer gap",
                                    "Value": f"{peer_gap:+.1f} bp",
                                    "Rule / Source": f"Selected issuer spread minus uploaded peer median in {rec_bucket}",
                                })

                                if peer_gap >= 10:
                                    narrative_lines.append(
                                        f"The issuer screens {peer_gap:+.1f} bp wide to the uploaded peer median in the same bucket."
                                    )
                                elif peer_gap <= -10:
                                    narrative_lines.append(
                                        f"The issuer screens {abs(peer_gap):.1f} bp tight to the uploaded peer median in the same bucket."
                                    )
                                else:
                                    narrative_lines.append(
                                        "The issuer screens broadly in line with the uploaded peer median."
                                    )
    except Exception as exc:
        evidence_rows.append({
            "Signal": "Peer comparison",
            "Value": "Unavailable",
            "Rule / Source": f"Peer narrative skipped: {exc}",
        })

    # -----------------------------
    # Signal 4: dealer flow proxy
    # -----------------------------
    dealer_col = None
    for candidate_col in ["trade_type", "side", "buy_sell", "customer_side", "dealer_side"]:
        if candidate_col in market_df.columns:
            dealer_col = candidate_col
            break

    if dealer_col is not None:
        try:
            flow_df = market_df[
                (market_df["issuer"] == selected_issuer)
                & (market_df["maturity_bucket"] == rec_bucket)
            ].copy()
            flow_df["trade_date"] = pd.to_datetime(flow_df["trade_date"], errors="coerce").dt.normalize()
            if "trade_amount" in flow_df.columns:
                flow_df["trade_amount"] = pd.to_numeric(flow_df["trade_amount"], errors="coerce").fillna(0)
            else:
                flow_df["trade_amount"] = 0.0
            flow_df = flow_df.dropna(subset=["trade_date"])
            if not flow_df.empty:
                latest_flow_date = flow_df["trade_date"].max()
                flow_df = flow_df[flow_df["trade_date"] >= latest_flow_date - pd.Timedelta(days=rec_window_days)].copy()

            def rec_classify_side(value: object) -> str:
                t = str(value).strip().lower()
                if any(token in t for token in ["sell", "sold", "sld", "customer sell", "cust sell", "cs"]):
                    return "Sell"
                if any(token in t for token in ["buy", "bought", "purchase", "customer buy", "cust buy", "cb"]):
                    return "Buy"
                if t == "s":
                    return "Sell"
                if t == "b":
                    return "Buy"
                return "Other"

            if not flow_df.empty:
                flow_df["flow_side"] = flow_df[dealer_col].map(rec_classify_side)
                buy_amt = flow_df.loc[flow_df["flow_side"] == "Buy", "trade_amount"].sum()
                sell_amt = flow_df.loc[flow_df["flow_side"] == "Sell", "trade_amount"].sum()
                denom = buy_amt + sell_amt
                if denom > 0:
                    amount_imbalance = (sell_amt - buy_amt) / denom
                    evidence_rows.append({
                        "Signal": "Dealer / flow imbalance",
                        "Value": f"{amount_imbalance:+.1%}",
                        "Rule / Source": f"(Sell amount - Buy amount) / total classified amount from `{dealer_col}`",
                    })
                    if amount_imbalance >= 0.25:
                        narrative_lines.append("Classified flow appears sell-heavy, which may indicate customer selling pressure.")
                    elif amount_imbalance <= -0.25:
                        narrative_lines.append("Classified flow appears buy-heavy, which may indicate stronger demand.")
        except Exception as exc:
            evidence_rows.append({
                "Signal": "Dealer / flow imbalance",
                "Value": "Unavailable",
                "Rule / Source": f"Flow narrative skipped: {exc}",
            })

    # -----------------------------
    # Recommendation label
    # -----------------------------
    score = 0
    if pd.notna(spread_change) and spread_change >= 15:
        score += 1
    if pd.notna(historical_percentile) and historical_percentile >= 75:
        score += 1
    if pd.notna(liquidity_score) and liquidity_score >= 60:
        score += 1
    if pd.notna(peer_gap) and peer_gap >= 10:
        score += 1
    if pd.notna(amount_imbalance) and amount_imbalance >= 0.25:
        score += 0.5

    if score >= 3:
        recommendation_label = "Potential Relative Value Candidate"
    elif score >= 2:
        recommendation_label = "Watchlist Candidate"
    elif score <= 0.5 and pd.notna(historical_percentile) and historical_percentile <= 25:
        recommendation_label = "Potentially Rich / Lower Priority"
    else:
        recommendation_label = "Neutral / Needs More Evidence"

    rec_m1, rec_m2, rec_m3, rec_m4 = st.columns(4)
    rec_m1.metric("Narrative Signal", recommendation_label)
    rec_m2.metric("Current Spread", "N/A" if pd.isna(current_spread) else f"{current_spread:+.1f} bp")
    rec_m3.metric("Spread Movement", "N/A" if pd.isna(spread_change) else f"{spread_change:+.1f} bp")
    rec_m4.metric("Historical Percentile", "N/A" if pd.isna(historical_percentile) else f"{historical_percentile:.0f}th")

    if narrative_lines:
        st.markdown("### Generated Commentary")
        commentary_text = " ".join(narrative_lines)
        st.info(commentary_text)
    else:
        st.info("Not enough data was available to generate a recommendation narrative for the selected inputs.")

    if evidence_rows:
        st.markdown("### Evidence Trail")
        evidence_df = pd.DataFrame(evidence_rows)
        safe_dataframe(evidence_df, width="stretch", hide_index=True)

    with st.expander("Rule thresholds used in this narrative", expanded=False):
        rule_df = pd.DataFrame(
            [
                {"Rule": "Material widening", "Threshold": "Spread movement >= +15 bp"},
                {"Rule": "Material tightening", "Threshold": "Spread movement <= -15 bp"},
                {"Rule": "Historically wide", "Threshold": "Historical percentile >= 75th"},
                {"Rule": "Very historically wide", "Threshold": "Historical percentile >= 90th"},
                {"Rule": "Above-average liquidity", "Threshold": "Liquidity score >= 70"},
                {"Rule": "Moderate liquidity", "Threshold": "Liquidity score 45–70"},
                {"Rule": "Wide versus peers", "Threshold": "Peer gap >= +10 bp"},
                {"Rule": "Tight versus peers", "Threshold": "Peer gap <= -10 bp"},
                {"Rule": "Sell-heavy flow", "Threshold": "Flow imbalance >= +25%"},
                {"Rule": "Buy-heavy flow", "Threshold": "Flow imbalance <= -25%"},
            ]
        )
        safe_dataframe(rule_df, width="stretch", hide_index=True)




section_anchor("ai-commentary-studio", "AI Commentary Studio")
with st.expander("Methodology: AI commentary studio", expanded=False):
    st.markdown(
        """
This section adds an **AI writing layer** on top of the dashboard analytics.

**Design principle:**

The AI should not invent the analysis. The dashboard first creates a structured evidence package from spreads, liquidity, historical percentile, peer gap, curve shape, and selected trade activity. The AI then converts that package into polished institutional commentary.

**Recommended use:**

- Use **Dashboard Analytics Only** when you want controlled commentary from internal numbers.
- Add **manual market context** when you already know what happened during the period.
- Enable **web search** only when you want current public market/news context. Review the output carefully.

**Important limitation:**

This commentary is for market discussion and screening. It is not an investment recommendation, credit opinion, or substitute for trader/PM judgment.
        """
    )

ai_col1, ai_col2, ai_col3, ai_col4 = st.columns([1, 1, 1, 1])
with ai_col1:
    ai_bucket = st.selectbox(
        "AI Commentary Bucket",
        MATURITY_BUCKET_ORDER,
        index=3,
        key="ai_commentary_bucket",
    )
with ai_col2:
    ai_rating = st.selectbox(
        "AI Benchmark",
        BENCHMARK_RATINGS,
        index=BENCHMARK_RATINGS.index("AAA") if "AAA" in BENCHMARK_RATINGS else 0,
        key="ai_commentary_rating",
    )
with ai_col3:
    ai_period = st.selectbox(
        "AI Commentary Period",
        ["1W", "1M", "3M", "6M", "1Y"],
        index=1,
        key="ai_commentary_period",
    )
with ai_col4:
    ai_model = st.selectbox(
        "AI Model",
        ["gpt-4.1-mini", "gpt-4o-mini"],
        index=0,
        key="ai_commentary_model",
    )

ai_period_days = {"1W": 7, "1M": 30, "3M": 90, "6M": 180, "1Y": 365}[ai_period]

st.markdown("### Controlled Market Context Retrieval")
st.caption(
    "Recommended workflow: first build dashboard signals, then retrieve/review public market context, "
    "then generate final commentary. This keeps the AI evidence-linked instead of free-form."
)

context_col1, context_col2 = st.columns([1.2, 1])
with context_col1:
    market_context_query = st.text_area(
        "Market / sector context search focus",
        value=(
            f"Municipal bond market context for {selected_sector} sector and {selected_issuer}; "
            f"focus on Treasury curve, muni market tone, fund flows, and sector headlines over the last {ai_period_days} days."
        ),
        height=110,
        key="ai_market_context_query",
        help="This tells the AI web search what public context to retrieve. Avoid confidential internal details.",
    )
with context_col2:
    ai_context_mode = st.radio(
        "AI Context Mode",
        [
            "Dashboard analytics only",
            "Manual context only",
            "Retrieve market context first",
            "Manual + retrieved context",
        ],
        index=2,
        key="ai_context_mode",
    )

manual_market_context = st.text_area(
    "Optional manual market / sector context",
    placeholder=(
        "Example: Treasury curve steepened during the period; long-duration munis underperformed; "
        "utility sector saw weaker tone after fund outflows..."
    ),
    height=120,
    key="manual_ai_market_context",
)

direct_web_search_in_commentary = st.checkbox(
    "Allow direct web search during final commentary generation",
    value=False,
    key="direct_web_search_in_commentary",
    help=(
        "Usually keep this off. Preferred workflow is: Retrieve Market Context → review it → Generate Commentary. "
        "Turn this on only if you want the final commentary call to do additional web search."
    ),
)

# -----------------------------
# Build structured context package
# -----------------------------
ai_context = {
    "issuer": selected_issuer,
    "sector": selected_sector,
    "bucket": ai_bucket,
    "benchmark": ai_rating,
    "period": ai_period,
    "analytics_as_of": pd.Timestamp.now().strftime("%Y-%m-%d %H:%M"),
    "executive_snapshot": {
        "securities": len(issuer_bonds),
        "trades_current_filter": len(issuer_trades),
        "latest_trade": issuer_trades["trade_date"].max().strftime("%m/%d/%Y") if not issuer_trades.empty else None,
    },
    "signals": {},
    "market_context_request": {
        "sector": selected_sector,
        "issuer": selected_issuer,
        "period_days": ai_period_days,
        "requested_public_context": [
            "Treasury curve / rates move",
            "municipal market tone",
            "fund flows if available",
            "sector-specific headlines",
            "issuer-specific public news only if available",
        ],
    },
}

# Spread movement and historical percentile.
try:
    ai_obs = build_spread_observations(
        market_df=market_df,
        mmd_df=mmd_df,
        issuer=selected_issuer,
        rating=ai_rating,
    )
    ai_obs = ai_obs.copy()
    ai_obs["trade_date"] = pd.to_datetime(ai_obs["trade_date"], errors="coerce").dt.normalize()
    ai_obs = ai_obs[
        (ai_obs["maturity_bucket"] == ai_bucket)
        & ai_obs["spread_to_benchmark_bps"].notna()
    ].sort_values("trade_date")

    if not ai_obs.empty:
        ai_latest = ai_obs.iloc[-1]
        ai_latest_date = ai_latest["trade_date"]
        ai_current_spread = float(ai_latest["spread_to_benchmark_bps"])
        ai_target_date = ai_latest_date - pd.Timedelta(days=ai_period_days)
        ai_hist_candidates = ai_obs[ai_obs["trade_date"] <= ai_target_date]

        if not ai_hist_candidates.empty:
            ai_hist_row = ai_hist_candidates.iloc[-1]
            ai_spread_change = ai_current_spread - float(ai_hist_row["spread_to_benchmark_bps"])
        else:
            ai_spread_change = None

        ai_1y = ai_obs[ai_obs["trade_date"] >= ai_latest_date - pd.Timedelta(days=365)]
        ai_values = pd.to_numeric(ai_1y["spread_to_benchmark_bps"], errors="coerce").dropna()
        ai_percentile = float((ai_values <= ai_current_spread).mean() * 100) if len(ai_values) >= 2 else None

        ai_context["signals"]["spread"] = {
            "current_spread_bps": round(ai_current_spread, 2),
            "spread_change_bps": None if ai_spread_change is None else round(ai_spread_change, 2),
            "historical_percentile_1y": None if ai_percentile is None else round(ai_percentile, 1),
            "latest_spread_date": ai_latest_date.strftime("%Y-%m-%d"),
            "benchmark_source": ai_latest.get("benchmark_source"),
            "benchmark_column": ai_latest.get("source_column"),
        }
except Exception as exc:
    ai_context["signals"]["spread_error"] = str(exc)

# Liquidity proxy.
try:
    ai_trades = market_df[
        (market_df["issuer"] == selected_issuer)
        & (market_df["maturity_bucket"] == ai_bucket)
    ].copy()
    ai_trades["trade_date"] = pd.to_datetime(ai_trades["trade_date"], errors="coerce").dt.normalize()
    if "trade_amount" in ai_trades.columns:
        ai_trades["trade_amount"] = pd.to_numeric(ai_trades["trade_amount"], errors="coerce").fillna(0)
    else:
        ai_trades["trade_amount"] = 0.0
    ai_trades = ai_trades.dropna(subset=["trade_date"])
    if not ai_trades.empty:
        ai_latest_trade = ai_trades["trade_date"].max()
        ai_trade_window = ai_trades[ai_trades["trade_date"] >= ai_latest_trade - pd.Timedelta(days=ai_period_days)].copy()
        ai_trade_count = len(ai_trade_window)
        ai_total_amount = float(ai_trade_window["trade_amount"].sum())
        ai_days_since_last = int((pd.Timestamp.today().normalize() - ai_latest_trade).days)
        ai_liq_score = (
            min(ai_trade_count / 10, 1) * 35
            + min(ai_total_amount / 5_000_000, 1) * 35
            + max(0, 1 - min(ai_days_since_last / 180, 1)) * 30
        )

        ai_context["signals"]["liquidity"] = {
            "liquidity_score_proxy": round(ai_liq_score, 1),
            "trade_count": int(ai_trade_count),
            "total_trade_amount": round(ai_total_amount, 0),
            "days_since_last_trade": ai_days_since_last,
        }
except Exception as exc:
    ai_context["signals"]["liquidity_error"] = str(exc)

# Peer gap proxy.
try:
    if "sector" in market_df.columns and selected_sector and selected_sector != "Unknown":
        ai_peer_universe = market_df[
            (market_df["sector"].astype(str) == str(selected_sector))
            & (market_df["issuer"].astype(str) != str(selected_issuer))
            & (market_df["maturity_bucket"] == ai_bucket)
        ].copy()
    else:
        ai_peer_universe = market_df[
            (market_df["issuer"].astype(str) != str(selected_issuer))
            & (market_df["maturity_bucket"] == ai_bucket)
        ].copy()

    if not ai_peer_universe.empty and "spread" in ai_context["signals"]:
        ai_peer_universe["trade_date"] = pd.to_datetime(ai_peer_universe["trade_date"], errors="coerce").dt.normalize()
        ai_peer_universe["yield"] = pd.to_numeric(ai_peer_universe["yield"], errors="coerce")
        ai_peer_universe = ai_peer_universe.dropna(subset=["trade_date", "yield"])
        if not ai_peer_universe.empty:
            ai_peer_latest = ai_peer_universe["trade_date"].max()
            ai_peer_universe = ai_peer_universe[
                ai_peer_universe["trade_date"] >= ai_peer_latest - pd.Timedelta(days=ai_period_days)
            ].copy()

            ai_date_col = _detect_mmd_date_column(mmd_df)
            if ai_date_col is not None and not ai_peer_universe.empty:
                ai_peer_mmd = mmd_df.copy()
                ai_peer_mmd[ai_date_col] = pd.to_datetime(ai_peer_mmd[ai_date_col], errors="coerce")
                ai_peer_mmd = ai_peer_mmd.dropna(subset=[ai_date_col])
                ai_peer_mmd = ai_peer_mmd[
                    ai_peer_mmd[ai_date_col].dt.normalize() <= ai_peer_latest
                ].sort_values(ai_date_col)

                if not ai_peer_mmd.empty:
                    ai_peer_tenor = MMD_BUCKET_MAP.get(ai_bucket, "10Y")
                    ai_peer_bench, _ = get_benchmark_curve(ai_peer_mmd.iloc[[-1]], ai_peer_tenor, ai_rating)
                    if ai_peer_bench is not None and pd.notna(ai_peer_bench.iloc[0]):
                        ai_peer_benchmark_yield = float(ai_peer_bench.iloc[0])
                        ai_peer_summary = (
                            ai_peer_universe.groupby("issuer", as_index=False)
                            .agg(avg_yield=("yield", "mean"), trade_count=("yield", "count"))
                        )
                        ai_peer_summary["peer_spread_bps"] = (
                            ai_peer_summary["avg_yield"] - ai_peer_benchmark_yield
                        ) * 100
                        ai_peer_median = float(ai_peer_summary["peer_spread_bps"].median())
                        ai_peer_gap = ai_context["signals"]["spread"]["current_spread_bps"] - ai_peer_median

                        ai_context["signals"]["peer_comparison"] = {
                            "peer_median_spread_bps": round(ai_peer_median, 2),
                            "peer_gap_bps": round(ai_peer_gap, 2),
                            "peer_count": int(len(ai_peer_summary)),
                        }
except Exception as exc:
    ai_context["signals"]["peer_error"] = str(exc)

# Curve shape snapshot if enough points.
try:
    ai_curve = market_df[market_df["issuer"] == selected_issuer].copy()
    ai_curve["trade_date"] = pd.to_datetime(ai_curve["trade_date"], errors="coerce").dt.normalize()
    ai_curve["yield"] = pd.to_numeric(ai_curve["yield"], errors="coerce")
    ai_curve = ai_curve.dropna(subset=["trade_date", "yield", "maturity_bucket"])
    ai_curve = ai_curve[ai_curve["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)]
    if not ai_curve.empty:
        ai_curve_latest = ai_curve["trade_date"].max()
        ai_curve_window = ai_curve[ai_curve["trade_date"] >= ai_curve_latest - pd.Timedelta(days=30)]
        ai_curve_summary = ai_curve_window.groupby("maturity_bucket")["yield"].mean().to_dict()
        ai_context["signals"]["curve_shape"] = {
            "short_yield": ai_curve_summary.get("5Y"),
            "ten_yield": ai_curve_summary.get("10Y"),
            "twenty_yield": ai_curve_summary.get("20Y"),
            "thirty_yield": ai_curve_summary.get("30Y"),
        }
        if "10Y" in ai_curve_summary and "30Y" in ai_curve_summary:
            ai_context["signals"]["curve_shape"]["10s30s_slope_pct"] = round(
                ai_curve_summary["30Y"] - ai_curve_summary["10Y"], 4
            )
except Exception as exc:
    ai_context["signals"]["curve_error"] = str(exc)

st.subheader("AI Context Package")
with st.expander("Review structured evidence before sending to AI", expanded=False):
    st.json(ai_context)

st.subheader("Ask the Dashboard Analyst")
st.caption(
    "Use this as one centralized AI copilot for questions about the selected issuer, benchmark, bucket, and dashboard signals."
)
ask_ai_question = st.text_area(
    "Ask a context question",
    placeholder="Example: What should I say about this issuer in a slide? Is this cheapness more likely spread-driven or liquidity-driven?",
    key="dashboard_ai_question",
)
ask_ai_model = st.selectbox(
    "Dashboard assistant model",
    ["gpt-4.1-mini", "gpt-4.1", "gpt-4o-mini"],
    index=0,
    key="dashboard_ai_question_model",
)
if st.button("Ask AI Analyst", key="ask_dashboard_ai_button"):
    if not ask_ai_question.strip():
        st.warning("Type a question first.")
    else:
        with st.spinner("Reading dashboard context..."):
            st.session_state["dashboard_ai_answer"] = answer_dashboard_question_with_ai(
                context_package=ai_context,
                user_question=ask_ai_question,
                model=ask_ai_model,
            )
if "dashboard_ai_answer" in st.session_state:
    st.markdown(st.session_state["dashboard_ai_answer"])

# -----------------------------
# Controlled retrieval + synthesis workflow
# -----------------------------
retrieve_enabled = ai_context_mode in ["Retrieve market context first", "Manual + retrieved context"]
manual_enabled = ai_context_mode in ["Manual context only", "Manual + retrieved context"]

if retrieve_enabled or direct_web_search_in_commentary:
    st.warning(
        "Public web context may be retrieved. Do not include confidential, proprietary, or client-sensitive information "
        "in the market context query or manual context field."
    )

action_col1, action_col2 = st.columns([1, 1])

with action_col1:
    if st.button("Retrieve Market / Sector Context", key="retrieve_ai_market_context", disabled=not retrieve_enabled):
        with st.spinner("Retrieving controlled public market context..."):
            retrieved_context = retrieve_market_context_with_openai(
                context_package=ai_context,
                market_context_query=market_context_query,
                model=ai_model,
            )
            st.session_state["latest_retrieved_market_context"] = retrieved_context

with action_col2:
    if st.button("Generate AI Institutional Commentary", key="generate_ai_institutional_commentary"):
        retrieved_context_for_commentary = (
            st.session_state.get("latest_retrieved_market_context", "")
            if ai_context_mode in ["Retrieve market context first", "Manual + retrieved context"]
            else ""
        )
        manual_context_for_commentary = manual_market_context if manual_enabled else ""

        with st.spinner("Generating evidence-linked commentary..."):
            commentary = generate_ai_market_commentary(
                context_package=ai_context,
                manual_market_context=manual_context_for_commentary,
                retrieved_market_context=retrieved_context_for_commentary,
                use_web_search=direct_web_search_in_commentary,
                market_context_query=market_context_query,
                model=ai_model,
            )
            st.session_state["latest_ai_commentary"] = commentary

if "latest_retrieved_market_context" in st.session_state:
    st.subheader("Retrieved Market / Sector Context")
    st.markdown(st.session_state["latest_retrieved_market_context"])

    st.download_button(
        label="Download Retrieved Market Context Markdown",
        data=st.session_state["latest_retrieved_market_context"].encode("utf-8"),
        file_name=f"{selected_issuer}_retrieved_market_context.md".replace(" ", "_"),
        mime="text/markdown",
    )

if "latest_ai_commentary" in st.session_state:
    st.subheader("Generated Institutional Commentary")
    st.markdown(st.session_state["latest_ai_commentary"])

    st.download_button(
        label="Download AI Commentary Markdown",
        data=st.session_state["latest_ai_commentary"].encode("utf-8"),
        file_name=f"{selected_issuer}_ai_market_commentary.md".replace(" ", "_"),
        mime="text/markdown",
    )

with st.expander("Recommended AI architecture for this dashboard", expanded=False):
    st.markdown(
        """
**Best practice: use a centralized AI Commentary Studio, not one AI button per section.**

Why:
1. Lower API cost.
2. Less duplicated commentary.
3. Lower risk of conflicting explanations.
4. Easier review by analysts / managers.
5. Cleaner evidence trail.

Recommended workflow:
1. Dashboard computes deterministic analytics.
2. AI Context Package captures the evidence.
3. Controlled Market Context Retrieval gathers public market/sector context.
4. User reviews retrieved context.
5. Final AI commentary synthesizes data + context into institutional language.
        """
    )

with st.expander("Where should AI commentary live?", expanded=False):
    st.markdown(
        """
I recommend **not** putting an AI button inside every single section.

A cleaner institutional structure is:

1. **Section-level rule commentary** stays deterministic and auditable.
2. **AI Commentary Studio** synthesizes across multiple sections.
3. Specific section AI can be added later only for the highest-value areas:
   - Historical Spread Percentile
   - Peer / Cross-Issuer RV
   - Security Screener
   - Recommendation Narrative

This keeps cost lower, avoids repeated/conflicting narratives, and makes the output easier for a team to review.
        """
    )


section_anchor("executive-snapshot", "Executive Snapshot")

latest_trade_display = (
    issuer_trades["trade_date"].max().strftime("%m/%d/%Y")
    if not issuer_trades.empty
    else "No trades"
)

# Custom cards give long sector/issuer names enough horizontal room, while keeping numeric fields quieter.
snap_col1, snap_col2, snap_col3, snap_col4, snap_col5 = st.columns([1.55, 2.15, 0.75, 0.9, 1.1])
with snap_col1:
    clean_metric_card("Sector", selected_sector, size="large")
with snap_col2:
    clean_metric_card("Issuer", selected_issuer, size="large")
with snap_col3:
    clean_metric_card("Securities", f"{len(issuer_bonds):,}", size="small")
with snap_col4:
    clean_metric_card("Trades", f"{len(issuer_trades):,}", size="small")
with snap_col5:
    clean_metric_card("Latest Trade", latest_trade_display, size="small")

section_anchor("bond-master", "Security Reference / Optional Bond Enrichment")
bond_cols = ["issuer", "sector", "primary_type", "election", "series", "cusip", "secondary_credit", "term", "maturity", "par_amount", "outstanding_amount", "coupon", "call_date", "call_price", "fed_tax", "amt"]
safe_dataframe(issuer_bonds[[c for c in bond_cols if c in issuer_bonds.columns]].sort_values([c for c in ["maturity", "cusip"] if c in issuer_bonds.columns]), width="stretch")

section_anchor("trade-detail", "Underlying Trade Detail")
trade_cols = ["trade_datetime", "cusip", "description", "maturity_trade", "maturity_bond", "maturity_bucket", "coupon_trade", "yield", "price", "trade_amount", "spread", "trade_type", "ratings_m_s_f"]
safe_dataframe(issuer_trades[[c for c in trade_cols if c in issuer_trades.columns]].sort_values("trade_datetime", ascending=False).head(20000), width="stretch")



if ENABLE_REPORT_EXPORTS:
    section_anchor("report-export-center", "Report Export Center")
    with st.expander("Methodology: report export center", expanded=False):
        st.markdown(
            """
    This section creates exportable reporting packages from the current dashboard state.

    **Export options:**

    - **Interactive HTML report:** includes selected charts, summary metrics, methodology notes, and key tables.
    - **PDF summary:** creates a lightweight PDF using `reportlab` when available. For a full visual PDF, download the HTML report and use browser **Print → Save as PDF**.
    - **PowerPoint slides:** creates a simple presentation using `python-pptx` when available.
    - **Chart HTML bundle:** exports selected Plotly charts as standalone HTML files inside a ZIP.
    - **Chart data CSV bundle:** exports the underlying data used to generate selected charts.

    **Important limitation:**

    A Streamlit app cannot reliably export the exact live browser page, all expanded/collapsed states, and all interactive widgets into a perfect PDF/PPTX without a browser automation service. This module therefore exports a clean, reproducible report built from the uploaded data and current issuer selection.
            """
        )

    # -----------------------------
    # Build exportable chart/data objects
    # -----------------------------
    export_chart_items = []
    export_data_items = {}

    def add_export_chart(name: str, fig, data: pd.DataFrame | None = None):
        """Collect charts for HTML/ZIP/PPT export without breaking if one chart fails."""
        try:
            export_chart_items.append((name, fig))
            if data is not None and not data.empty:
                export_data_items[name] = data.copy()
        except Exception:
            pass

    # 1) Yield trend chart
    try:
        export_yield_df = market_df[market_df["issuer"] == selected_issuer].copy()
        export_yield_df["trade_date"] = pd.to_datetime(export_yield_df["trade_date"], errors="coerce")
        export_yield_df["yield"] = pd.to_numeric(export_yield_df["yield"], errors="coerce")
        export_yield_df = export_yield_df.dropna(subset=["trade_date", "yield"])
        if not export_yield_df.empty:
            export_yield_daily = (
                export_yield_df.groupby("trade_date", as_index=False)
                .agg(avg_yield=("yield", "mean"), trade_count=("yield", "count"), total_trade_amount=("trade_amount", "sum") if "trade_amount" in export_yield_df.columns else ("yield", "count"))
                .sort_values("trade_date")
            )
            export_yield_fig = px.line(
                export_yield_daily,
                x="trade_date",
                y="avg_yield",
                markers=True,
                title=f"{selected_issuer} Average Trade Yield",
                labels={"trade_date": "Trade Date", "avg_yield": "Average Yield (%)"},
            )
            add_export_chart("yield_trend", export_yield_fig, export_yield_daily)
    except Exception:
        pass

    # 2) Issuer curve chart
    try:
        export_curve_df = market_df[market_df["issuer"] == selected_issuer].copy()
        export_curve_df["trade_date"] = pd.to_datetime(export_curve_df["trade_date"], errors="coerce").dt.normalize()
        export_curve_df["yield"] = pd.to_numeric(export_curve_df["yield"], errors="coerce")
        export_curve_df = export_curve_df.dropna(subset=["trade_date", "yield", "maturity_bucket"])
        export_curve_df = export_curve_df[export_curve_df["maturity_bucket"].isin(MATURITY_BUCKET_ORDER)].copy()
        if not export_curve_df.empty:
            latest_curve_date = export_curve_df["trade_date"].max()
            export_curve_df = export_curve_df[export_curve_df["trade_date"] >= latest_curve_date - pd.Timedelta(days=30)]
            export_curve_summary = (
                export_curve_df.groupby("maturity_bucket", as_index=False)
                .agg(avg_yield=("yield", "mean"), trade_count=("yield", "count"))
            )
            export_curve_summary["maturity_bucket"] = pd.Categorical(
                export_curve_summary["maturity_bucket"],
                categories=MATURITY_BUCKET_ORDER,
                ordered=True,
            )
            export_curve_summary = export_curve_summary.sort_values("maturity_bucket")
            export_curve_fig = px.line(
                export_curve_summary,
                x="maturity_bucket",
                y="avg_yield",
                markers=True,
                title=f"{selected_issuer} Issuer Curve — Latest 30D Average",
                labels={"maturity_bucket": "Maturity Year", "avg_yield": "Average Yield (%)"},
            )
            add_export_chart("issuer_curve_latest_30d", export_curve_fig, export_curve_summary)
    except Exception:
        pass

    # 3) Current spread level ladder
    try:
        if not mmd_df.empty:
            export_level_table, export_level_audit = build_spread_level_data(
                market_df=market_df,
                mmd_df=mmd_df,
                issuer=selected_issuer,
                ratings=["AAA", "AA", "A", "BBB"],
            )
            if not export_level_table.empty and not export_level_table.isna().all().all():
                export_level_ladder = safe_melt_by_maturity(
                    export_level_table.reset_index().rename(columns={export_level_table.index.name or "index": "maturity_bucket"}),
                    value_vars=[c for c in export_level_table.columns if c in ["AAA", "AA", "A", "BBB"]],
                    id_vars="maturity_bucket",
                    var_name="benchmark_rating",
                    value_name="spread_to_benchmark_bps",
                )
                if not export_level_ladder.empty:
                    export_level_ladder["spread_to_benchmark_bps"] = pd.to_numeric(export_level_ladder["spread_to_benchmark_bps"], errors="coerce")
                    export_level_ladder = export_level_ladder.dropna(subset=["spread_to_benchmark_bps"]).copy()
                    export_level_ladder["label"] = export_level_ladder["benchmark_rating"].astype(str) + " / " + export_level_ladder["maturity_bucket"].astype(str)
                    export_level_ladder["abs_spread_bps"] = export_level_ladder["spread_to_benchmark_bps"].abs()
                    export_level_ladder = export_level_ladder.sort_values("abs_spread_bps", ascending=False).head(20)
                    export_level_fig = px.bar(
                        export_level_ladder.sort_values("spread_to_benchmark_bps"),
                        x="spread_to_benchmark_bps",
                        y="label",
                        orientation="h",
                        color="benchmark_rating",
                        title=f"{selected_issuer} Current Spread Level Ladder",
                        labels={"spread_to_benchmark_bps": "Spread (bps)", "label": "Benchmark / Maturity"},
                    )
                    export_level_fig.add_vline(x=0, line_dash="dash", opacity=0.45)
                    add_export_chart("current_spread_level_ladder", export_level_fig, export_level_audit)
    except Exception:
        pass

    # 4) Liquidity monthly activity chart
    try:
        export_liq_df = market_df[market_df["issuer"] == selected_issuer].copy()
        export_liq_df["trade_date"] = pd.to_datetime(export_liq_df["trade_date"], errors="coerce")
        export_liq_df = export_liq_df.dropna(subset=["trade_date"])
        if "trade_amount" in export_liq_df.columns:
            export_liq_df["trade_amount"] = pd.to_numeric(export_liq_df["trade_amount"], errors="coerce").fillna(0)
        else:
            export_liq_df["trade_amount"] = 0.0
        if not export_liq_df.empty:
            export_liq_df["trade_month"] = export_liq_df["trade_date"].dt.to_period("M").astype(str)
            export_monthly = (
                export_liq_df.groupby("trade_month", as_index=False)
                .agg(trade_count=("trade_date", "count"), total_trade_amount=("trade_amount", "sum"))
            )
            export_monthly_fig = px.line(
                export_monthly,
                x="trade_month",
                y="trade_count",
                markers=True,
                title=f"{selected_issuer} Monthly Trade Count",
                labels={"trade_month": "Trade Month", "trade_count": "Trade Count"},
            )
            add_export_chart("monthly_trade_count", export_monthly_fig, export_monthly)
    except Exception:
        pass

    # -----------------------------
    # Export controls
    # -----------------------------
    export_options = st.multiselect(
        "Select charts to include",
        [name for name, _fig in export_chart_items],
        default=[name for name, _fig in export_chart_items],
        help="These are reconstructed export charts based on current selected issuer and uploaded data.",
    )

    selected_export_charts = [(name, fig) for name, fig in export_chart_items if name in export_options]

    export_meta = {
        "Generated": pd.Timestamp.now().strftime("%Y-%m-%d %H:%M"),
        "Selected Issuer": selected_issuer,
        "Sector": selected_sector,
        "Securities": f"{len(issuer_bonds):,}",
        "Trades in Current Filter": f"{len(issuer_trades):,}",
        "Latest Trade": issuer_trades["trade_date"].max().strftime("%m/%d/%Y") if not issuer_trades.empty else "No trades",
    }

    report_html_parts = [
        "<html><head><meta charset='utf-8'><title>Municipal Secondary Market Dashboard Report</title>",
        "<style>body{font-family:Arial,sans-serif;margin:32px;color:#111827;} h1,h2{color:#111827;} table{border-collapse:collapse;width:100%;margin:16px 0;} td,th{border:1px solid #e5e7eb;padding:8px;text-align:left;} .note{color:#64748b;font-size:13px;}</style>",
        "</head><body>",
        "<h1>Municipal Secondary Market Dashboard Report</h1>",
        "<h2>Executive Summary</h2>",
        "<table>",
    ]
    for k, v in export_meta.items():
        report_html_parts.append(f"<tr><th>{k}</th><td>{v}</td></tr>")
    report_html_parts.extend([
        "</table>",
        "<p class='note'>Benchmark curves use uploaded rating curves when available; otherwise the app falls back to MMD/AAA plus transparent spread assumptions. Screening outputs are not investment recommendations.</p>",
    ])

    for name, fig in selected_export_charts:
        report_html_parts.append(f"<h2>{name.replace('_', ' ').title()}</h2>")
        report_html_parts.append(fig.to_html(full_html=False, include_plotlyjs="cdn"))

    report_html_parts.append("</body></html>")
    full_report_html = "\n".join(report_html_parts)

    export_col1, export_col2, export_col3 = st.columns(3)

    with export_col1:
        st.download_button(
            label="Download Interactive HTML Report",
            data=full_report_html.encode("utf-8"),
            file_name=f"{selected_issuer}_dashboard_report.html".replace(" ", "_"),
            mime="text/html",
            help="Open this file in a browser. For a visual PDF, use browser Print → Save as PDF.",
        )

    with export_col2:
        # Lightweight PDF summary via reportlab, if installed.
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors

            pdf_buffer = io.BytesIO()
            doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
            styles = getSampleStyleSheet()
            story = [
                Paragraph("Municipal Secondary Market Dashboard Summary", styles["Title"]),
                Spacer(1, 12),
            ]
            meta_table = Table([[k, str(v)] for k, v in export_meta.items()])
            meta_table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.lightgrey),
                ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(meta_table)
            story.append(Spacer(1, 14))
            story.append(Paragraph("Included Charts", styles["Heading2"]))
            for name, _fig in selected_export_charts:
                story.append(Paragraph(f"• {name.replace('_', ' ').title()}", styles["BodyText"]))
            story.append(Spacer(1, 14))
            story.append(Paragraph("Methodology Note", styles["Heading2"]))
            story.append(Paragraph(
                "This PDF is a lightweight summary. For interactive charts and fuller visual output, use the HTML report and browser Print → Save as PDF.",
                styles["BodyText"],
            ))
            doc.build(story)
            pdf_bytes = pdf_buffer.getvalue()

            st.download_button(
                label="Download PDF Summary",
                data=pdf_bytes,
                file_name=f"{selected_issuer}_dashboard_summary.pdf".replace(" ", "_"),
                mime="application/pdf",
            )
        except Exception:
            st.info("PDF summary export requires `reportlab`. Add `reportlab` to requirements.txt, or download HTML and print/save as PDF.")

    with export_col3:
        # PPTX summary via python-pptx, if installed.
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Municipal Secondary Market Dashboard"
            slide.placeholders[1].text = f"{selected_issuer} | {selected_sector} | Generated {export_meta['Generated']}"

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Executive Snapshot"
            body = slide.placeholders[1].text_frame
            body.clear()
            for k, v in export_meta.items():
                p = body.add_paragraph()
                p.text = f"{k}: {v}"
                p.font.size = Pt(18)

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Included Dashboard Charts"
            body = slide.placeholders[1].text_frame
            body.clear()
            for name, _fig in selected_export_charts:
                p = body.add_paragraph()
                p.text = name.replace("_", " ").title()
                p.level = 0
                p.font.size = Pt(18)

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Methodology Notes"
            body = slide.placeholders[1].text_frame
            body.clear()
            notes = [
                "Benchmark curves use uploaded rating curves when available.",
                "Fallback curves use MMD/AAA plus visible rating-spread assumptions.",
                "Liquidity and RV scores are screening tools, not trade recommendations.",
                "Scenario shock uses duration proxies, not full cash-flow pricing.",
            ]
            for note in notes:
                p = body.add_paragraph()
                p.text = note
                p.font.size = Pt(16)

            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_bytes = pptx_buffer.getvalue()

            st.download_button(
                label="Download PowerPoint Slides",
                data=pptx_bytes,
                file_name=f"{selected_issuer}_dashboard_slides.pptx".replace(" ", "_"),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        except Exception:
            st.info("PowerPoint export requires `python-pptx`. Add `python-pptx` to requirements.txt to enable slide export.")

    # Chart HTML bundle and chart data bundle
    bundle_col1, bundle_col2 = st.columns(2)

    with bundle_col1:
        try:
            import zipfile
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "w", compression=zipfile.ZIP_DEFLATED) as zf:
                for name, fig in selected_export_charts:
                    zf.writestr(f"{name}.html", fig.to_html(full_html=True, include_plotlyjs="cdn"))
            st.download_button(
                label="Download Chart HTML Bundle",
                data=zip_buffer.getvalue(),
                file_name=f"{selected_issuer}_chart_html_bundle.zip".replace(" ", "_"),
                mime="application/zip",
            )
        except Exception as exc:
            st.info(f"Chart bundle export unavailable: {exc}")

    with bundle_col2:
        try:
            import zipfile
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "w", compression=zipfile.ZIP_DEFLATED) as zf:
                for name, df in export_data_items.items():
                    if name in export_options and df is not None and not df.empty:
                        zf.writestr(f"{name}_data.csv", df.to_csv(index=False))
            st.download_button(
                label="Download Chart Data CSV Bundle",
                data=zip_buffer.getvalue(),
                file_name=f"{selected_issuer}_chart_data_bundle.zip".replace(" ", "_"),
                mime="application/zip",
            )
        except Exception as exc:
            st.info(f"Chart data bundle export unavailable: {exc}")

    with st.expander("How to export the full live webpage as PDF", expanded=False):
        st.markdown(
            """
    For the exact live Streamlit page:

    1. Open the dashboard in your browser.
    2. Expand the sections you want included.
    3. Press **Cmd+P** on Mac or **Ctrl+P** on Windows.
    4. Choose **Save as PDF**.
    5. Set scale to 70–85% if charts are too wide.

    For a cleaner report with reproducible charts, use **Download Interactive HTML Report** above.
            """
        )


else:
    section_anchor("report-export-center", "Report Export Center")
    st.info("Report export builder is disabled for speed. Enable it in the sidebar Performance panel when you need HTML/PDF/PPTX exports.")
    export_chart_items = []
    export_data_items = {}
section_anchor("export-summary", "Export Summary Package")
with st.expander("Methodology: export summary package", expanded=False):
    st.markdown(
        """
This section generates a lightweight export package that can be copied into internal updates, pitchbook drafts, or meeting notes.

**Current implementation:**

- Generates a Markdown summary and an HTML summary.
- Uses the selected issuer, selected sector, counts, latest trade date, and available dashboard outputs.
- For PDF export, open the HTML file in a browser and print/save as PDF.
- For PowerPoint, copy the HTML/Markdown summary into your deck and insert key charts from the dashboard.

This avoids adding fragile report-generation dependencies while keeping the workflow practical for internal use.
        """
    )

latest_trade_text = issuer_trades["trade_date"].max().strftime("%m/%d/%Y") if not issuer_trades.empty else "No trades"
summary_timestamp = pd.Timestamp.now().strftime("%Y-%m-%d %H:%M")
summary_lines = [
    f"# Municipal Secondary Market Dashboard Summary",
    "",
    f"**Generated:** {summary_timestamp}",
    f"**Selected Issuer:** {selected_issuer}",
    f"**Sector:** {selected_sector}",
    f"**Securities:** {len(issuer_bonds):,}",
    f"**Trades in Current Filter:** {len(issuer_trades):,}",
    f"**Latest Trade:** {latest_trade_text}",
    "",
    "## Key Dashboard Modules",
    "- Yield Trend / Relative Value Comparison",
    "- Issuer Curve vs Benchmark Curve",
    "- Current Spread Level Framework",
    "- Peer and Cross-Issuer Relative Value",
    "- Historical Spread Percentile",
    "- Security Screener",
    "- Recommendation Narrative",
    "- Scenario Shock Analysis",
    "- CUSIP Opportunity Drilldown",
    "",
    "## Notes",
    "- Benchmark curves use uploaded rating curves when available; otherwise the dashboard falls back to MMD/AAA plus visible spread assumptions.",
    "- Liquidity, RV score, dealer proxy, and scenario shock outputs are screening tools, not final trade recommendations.",
]

if "data_quality_score" in locals():
    summary_lines.extend([
        "",
        "## Data Quality",
        f"- Data Quality Score: {data_quality_score:.1f}/100",
        f"- Valid CUSIP Rate: {cusip_match_rate:.1f}%",
        f"- Known Maturity Year Rate: {known_bucket_rate:.1f}%",
        f"- Duplicates Removed: {duplicates_removed:,}",
    ])

summary_md = "\n".join(summary_lines)
summary_html = summary_md.replace("\n", "<br>")

export_c1, export_c2 = st.columns(2)
with export_c1:
    st.download_button(
        label="Download Markdown Summary",
        data=summary_md.encode("utf-8"),
        file_name=f"{selected_issuer}_dashboard_summary.md".replace(" ", "_"),
        mime="text/markdown",
    )
with export_c2:
    st.download_button(
        label="Download HTML Summary",
        data=f"<html><body>{summary_html}</body></html>".encode("utf-8"),
        file_name=f"{selected_issuer}_dashboard_summary.html".replace(" ", "_"),
        mime="text/html",
    )

with st.expander("Preview export summary", expanded=True):
    st.markdown(summary_md)

section_anchor("admin-methodology", "Admin Methodology Page")
st.markdown(
    """
This page centralizes the assumptions used throughout the dashboard so the tool is easier to hand off and maintain.

### Benchmark Curves
- AAA = uploaded MMD / AAA curve.
- Non-AAA curves use uploaded rating-specific columns when available.
- If not available, non-AAA curves are modeled as MMD/AAA + visible rating-spread assumptions.

### Liquidity Score
Uses trade count, total trade amount, recent activity, and recency. It is a screening score, not a credit rating.

### Relative Value Score
Combines spread percentile, liquidity percentile, and trade activity percentile to identify cheap + tradable candidates.

### Dealer Behavior Proxy
Only enabled when trade side/type data exists. It estimates buy/sell imbalance but does not represent true dealer inventory.

### Scenario Shock
Uses duration proxies by maturity year. It does not model full cash flows, convexity, OAS, tax status, or callable optionality.

### Recommendation Narrative
Rule-based and explainable. Each phrase is triggered by spread movement, historical percentile, liquidity, peer gap, or flow proxy thresholds.
"""
)

with st.expander("Rating spread assumptions", expanded=False):
    safe_dataframe(rating_spread_table(), width="stretch", hide_index=True)

with st.expander("Duration proxy assumptions", expanded=False):
    duration_proxy_df = pd.DataFrame(
        [{"Maturity Year": f"{y}Y", "Proxy Duration": DURATION_PROXY[f"{y}Y"]} for y in range(1, MAX_MATURITY_YEAR + 1)]
    )
    safe_dataframe(duration_proxy_df, width="stretch", hide_index=True)

section_anchor("version-changelog", "Version / Change Log")
version_rows = [
    {"Version": "v1.0-team-ready", "Change": "Stabilized data validation, benchmark framework, relative value analytics, and team-readiness modules."},
    {"Version": "v1.1", "Change": "Added Cross-Issuer RV Analytics, Scenario Shock, Recommendation Narrative, and CUSIP Drilldown."},
    {"Version": "v1.2", "Change": "Added Data Quality Scorecard, Export Summary Package, Admin Methodology Page, and Watchlist."},
]
safe_dataframe(pd.DataFrame(version_rows), width="stretch", hide_index=True)
st.caption("Update this changelog whenever the team changes methodology, assumptions, or major modules.")


section_anchor("downloads", "Download Outputs")
d1, d2, d3 = st.columns(3)
with d1:
    dataframe_download_button(market_df, "Download Merged Market Data CSV", "merged_market_data.csv")
with d2:
    dataframe_download_button(issuer_master, "Download Issuer Master CSV", "issuer_master.csv")
with d3:
    dataframe_download_button(bonds_df, "Download Security Reference CSV", "security_reference.csv")

if show_raw_tables:
    st.header("Raw / Processed Tables")
    st.subheader("Issuer Master")
    safe_dataframe(issuer_master, width="stretch")
    st.subheader("Security Reference")
    safe_dataframe(bonds_df, width="stretch")
    st.subheader("All Trades")
    safe_dataframe(trades_df.head(20000), width="stretch")
    st.subheader("Merged Market Data")
    safe_dataframe(market_df.head(20000), width="stretch")
